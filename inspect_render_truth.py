"""inspect_render_truth.py — ask real PowerPoint where it actually broke the lines.

Everything else in this repo MODELS PowerPoint's layout (Pillow/metrics.json glyph
widths, a 1.2x line pitch, a 3pt paragraph gap), and every checker we have
measures with the same measurer the packer used -- so the two agree by
construction and a clean `inspect_pptx.py` run is not evidence that PowerPoint
agrees. The only real oracle was measure_boundheight.bas, pasted into the VBE by
hand and read back by eye. That manual loop is why the capacity investigation ran
15 rounds instead of 3.

This is the automated replacement. On Windows it drives PowerPoint over COM,
pulls back what its layout engine actually produced, and scores the model
against it.

The ground truth
----------------
    TextRange.BoundHeight          the height PowerPoint gave the text
    Paragraphs(k).Lines.Count      how many lines it broke paragraph k into

The per-paragraph count is the sharp instrument. A whole-shape height can be
right for the wrong reasons -- two errors cancelling is exactly what a mixed
slot produces -- but a paragraph PowerPoint broke into 4 lines where we
predicted 3 cannot be.

What it scores
--------------
For every "■ key - ..." bullet it computes three predictions and asks which one
PowerPoint agrees with:

    A  what the packer does today  regular-weight widths, and the label is the
                                   mapping_key (_account_cost_key), which is NOT
                                   the text the renderer writes
    B  A + the rendered label      the display_name/display_name_zh actually on
                                   the slide (_rendered_bullet_label already
                                   exists for this; 5 of its 7 call sites in
                                   gen_packing.py still use _account_cost_key)
    C  B + a bold-aware key run    generation.py sets run_key.font.bold = True,
                                   but text_metrics.Measurer holds one
                                   regular-weight table and has no weight
                                   parameter

If C scores no better than A on real content, the fixes are not worth making and
the remaining error is somewhere else. That is the point: decide from measured
truth, before writing anything into the packer.

It also reports
---------------
  * real fill % = BoundHeight / usable box height. The honest answer to "why do
    the pages look under-filled" -- unlike inspect_pptx.py's fill%, this
    numerator came from PowerPoint. Commentary slots left ENTIRELY empty are
    listed too, at a measured 0%: an unused column is the largest under-fill
    there is and filtering it out made a wasted half-page look like a page
    that simply had fewer columns.
  * pitch per shape = BoundHeight with the paragraph gaps taken back out, over
    the real line count. Below the nominal 10.8pt means PowerPoint re-ran its
    own autofit shrink and ignored the fontScale we wrote. (Not BoundHeight/
    Lines: that bundles the gaps in and always reads high, so it could never
    show a shrink at all.)
  * pitch and gap least-squares-fitted out of the real numbers, both ways round
    on whether the last paragraph's space_after counts -- so 10.8/3.0 gets
    re-measured rather than re-argued. (3.0 was already restored once, 63e4120,
    after 2.2 was back-solved and proved wrong. Re-fit, don't re-litigate.)

Usage
-----
    python inspect_render_truth.py exported_deck.pptx
    python inspect_render_truth.py exported_deck.pptx --lines --csv truth.csv
    python inspect_render_truth.py exported_deck.pptx --model-only   # no PowerPoint

Requires pywin32 (`pip install pywin32`) and PowerPoint. The deck is opened
READ-ONLY and closed without saving; PowerPoint is only quit if this script was
the one that started it.
"""

from __future__ import annotations

import argparse
import csv
import os
import sys
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Sequence, Tuple

from pptx import Presentation

# The SAME measurer production uses. Imported, never reimplemented -- a fourth
# copy of the formula is exactly the drift this tool exists to detect.
from fdd_utils.financial_common import load_yaml_file
from fdd_utils.text_metrics import (POWERPOINT_LINE_PITCH_FACTOR, get_measurer,
                                    text_box_from_shape)
from fdd_utils.pptx.helpers import (
    _measurer_family,
    _real_font_size_pt,
    _real_line_spacing,
    _real_para_gap_pt,
    _resolve_font_metrics_path,
)
from fdd_utils.pptx.payloads import _load_pptx_settings
from inspect_pptx import _is_chinese_text, _slot_of

# generation.py's p_key: left_indent 0.15" / first_line_indent -0.15", so line 1
# spans the full box and wrapped continuation lines sit 10.8pt narrower.
BULLET_HANGING_INDENT_PT = 10.8

BULLET_MARKER = "■"
# _explanation_render_text prefixes a post-table explanation with "➢ " in a
# Chinese deck and "- " in an English one.
EXPLAIN_MARKERS = ("➢", "- ", "• ")

# Bold faces for variant C. Windows first -- that is where this runs.
_BOLD_FONTS = {
    False: [r"C:\Windows\Fonts\arialbd.ttf",
            "/System/Library/Fonts/Supplemental/Arial Bold.ttf",
            "/usr/share/fonts/truetype/msttcorefonts/Arial_Bold.ttf"],
    True: [r"C:\Windows\Fonts\msyhbd.ttc",
           r"C:\Windows\Fonts\msyh.ttc",
           "/Library/Fonts/Microsoft/Microsoft YaHei Bold.ttf"],
}


# ---------------------------------------------------------------------------
# Records
# ---------------------------------------------------------------------------

@dataclass
class ParaRow:
    slide: int
    shape: str
    slot: str
    index: int              # 1-based, matches PowerPoint's Paragraphs(k)
    kind: str               # bullet | continuation | category | explain | blank
    has_bold: bool
    chars: int
    lines_a: int            # today's packer
    lines_b: int            # + rendered label
    lines_c: int            # + bold-aware key run
    label: str = ""         # the text actually drawn in the bold run
    mapping_key: str = ""   # what the packer would have charged instead
    real_lines: Optional[int] = None
    space_after_pt: float = 0.0
    space_before_pt: float = 0.0
    selfcheck_ok: bool = True
    text: str = ""
    font_pt: float = 9.0    # this paragraph's own size, NOT the deck default
    real_line_texts: List[str] = field(default_factory=list)

    def delta(self, variant: str) -> Optional[int]:
        if self.real_lines is None:
            return None
        return self.real_lines - getattr(self, f"lines_{variant}")


@dataclass
class ShapeRow:
    slide: int
    shape: str
    slot: str
    is_chinese: bool
    box_w_pt: float
    box_h_pt: float
    model_lines: int
    model_pt: float
    real_lines: Optional[int] = None
    real_bound_pt: Optional[float] = None
    is_empty: bool = False
    paras: List[ParaRow] = field(default_factory=list)
    measurer: object = None      # the one this shape was measured with

    @property
    def real_pitch(self) -> Optional[float]:
        """PowerPoint's real baseline-to-baseline pitch, normalised to a 9pt
        line, with the paragraph gaps taken back out.

        Two things have to come out before this number means anything.

        The gaps: BoundHeight/Lines bundles every inter-paragraph gap into the
        average, so it always reads ABOVE the nominal line height (11.1-11.7pt
        against a real 10.8pt) and could never show the one thing this column
        exists for -- a normAutofit shrink makes the pitch SMALLER.

        The font size: a table slot's blank spacer paragraphs are deliberately
        sized down (real exports carry 1.0, 3.66, 4.33, 6.33 and 7.0pt runs),
        so counting them as 9pt lines drags the average below 10.8 and cries
        "autofit shrink" on a deck with noAutofit set. Each line is weighted by
        its own size instead, which is why this reads 10.80 on a table slide
        and only moves when PowerPoint really did shrink something.
        """
        if not self.real_lines or self.real_bound_pt is None:
            return None
        equiv = self.equivalent_9pt_lines
        if equiv <= 0:
            return None
        return (self.real_bound_pt - self.gap_total_pt) / equiv

    @property
    def gap_total_pt(self) -> float:
        """Every point of paragraph spacing inside this shape, summed from the
        paragraphs themselves.

        Not `count x _real_para_gap_pt`: that approximation is what has now put
        this file wrong about the model three times. model_pt accumulates the
        real sa/sb of each paragraph and refunds the last one's space_after, so
        the pitch and the fit have to subtract exactly that, not a nominal
        stand-in for it.
        """
        if not self.paras:
            return 0.0
        total = sum(p.space_after_pt + p.space_before_pt for p in self.paras)
        return total - self.paras[-1].space_after_pt

    @property
    def equivalent_9pt_lines(self) -> float:
        """Real line count with each line weighted by its paragraph's own font
        size, so a 1pt spacer line counts as 1/9th of a 9pt line."""
        base = _real_font_size_pt(self.is_chinese) or 9.0
        total = 0.0
        for p in self.paras:
            n = p.real_lines if p.real_lines is not None else p.lines_b
            total += n * (p.font_pt / base)
        return total

    @property
    def real_fill(self) -> Optional[float]:
        if self.real_bound_pt is None or self.box_h_pt <= 0:
            return None
        return self.real_bound_pt / self.box_h_pt


# ---------------------------------------------------------------------------
# Width sources
# ---------------------------------------------------------------------------

def _bold_font(is_cjk: bool, size_pt: float):
    """Pillow face for the BOLD key run. Pillow rather than a MetricsTable
    because there is no bold metrics.json and fontTools is not in
    requirements.txt -- but for Arial the JSON was dumped from the same file,
    so the two agree. Returns None when no bold face is installed, in which
    case variant C is simply not scored rather than silently faked."""
    from PIL import ImageFont
    for path in _BOLD_FONTS[is_cjk]:
        if os.path.exists(path):
            try:
                return ImageFont.truetype(path, size=size_pt)
            except Exception:
                continue
    return None


def _is_cjk_char(ch: str) -> bool:
    cp = ord(ch)
    return (0x3000 <= cp <= 0x303F or 0x3400 <= cp <= 0x4DBF or 0x4E00 <= cp <= 0x9FFF
            or 0xF900 <= cp <= 0xFAFF or 0xFF00 <= cp <= 0xFFEF)


def _atomize(text: str) -> List[str]:
    """Same atoms text_metrics.wrap_paragraph uses: one CJK char, or a maximal
    run of non-CJK non-space."""
    out, i, n = [], 0, len(text)
    while i < n:
        if text[i].isspace():
            i += 1
            continue
        if _is_cjk_char(text[i]):
            out.append(text[i]); i += 1
            continue
        j = i
        while j < n and not text[j].isspace() and not _is_cjk_char(text[j]):
            j += 1
        out.append(text[i:j]); i = j
    return out


def _wrap_runs(runs: Sequence[Tuple[str, bool]], measurer, bold_font,
               first_width_pt: float, width_pt: float) -> int:
    """Line count for a paragraph whose runs have DIFFERENT weights.

    A tool-local greedy wrapper, because the production `measurer.wrap()` takes
    one string at one weight and cannot express "this run is bold". It uses the
    production measurer for every regular run, so the widths are the same
    source; only the loop is local. `_selfcheck` below proves that loop agrees
    with `measurer.wrap()` whenever nothing is bold -- if it ever stops
    agreeing, the tool says so and variant C is not to be trusted.
    """
    def w(text: str, is_bold: bool) -> float:
        if is_bold and bold_font is not None:
            return float(bold_font.getlength(text))
        return measurer.text_width_pt(text)

    lines, cur, limit = 1, 0.0, first_width_pt
    prev_atom = ""
    for text, is_bold in runs:
        for atom in _atomize(text):
            sep = 0.0
            if cur > 0 and prev_atom and not _is_cjk_char(atom[0]) and not _is_cjk_char(prev_atom[-1]):
                sep = w(" ", is_bold)
            aw = w(atom, is_bold)
            if cur > 0 and cur + sep + aw > limit:
                lines += 1
                cur, limit = aw, width_pt
            else:
                cur += sep + aw
            prev_atom = atom
    return lines


# ---------------------------------------------------------------------------
# mapping_key recovery
# ---------------------------------------------------------------------------

def _alias_index() -> Dict[str, List[str]]:
    """rendered label -> the mapping_key(s) it could have come from.

    The deck only carries the label the renderer wrote; the packer charged the
    mapping_key. mappings.yml is what links them, and it is the same file the
    payload builder used, so this recovers variant A's input exactly rather
    than approximating it.
    """
    raw = load_yaml_file("fdd_utils/mappings.yml") or {}
    raw = raw.get("mappings", raw)
    idx: Dict[str, set] = {}
    for key, cfg in raw.items():
        if not isinstance(cfg, dict):
            continue
        for alias in [key] + list(cfg.get("aliases") or []):
            idx.setdefault(str(alias).strip(), set()).add(key)
    return {k: sorted(v) for k, v in idx.items()}


# ---------------------------------------------------------------------------
# Do the model's assumptions still hold in this deck?
# ---------------------------------------------------------------------------
#
# Every number the packer produces rests on constants that describe the CURRENT
# renderer and the CURRENT template: 9pt, line spacing 1.0, a 3pt paragraph gap,
# a 0.15" hanging indent, Arial + Microsoft YaHei. None of them is read from the
# document. Change the shape or the font and the model does not fail -- it keeps
# answering, wrongly and silently.
#
# So check them. This is the difference between a model that happens to be right
# and one that is reliable: when an assumption stops holding, something says so.
#
# The indent is the reason this section exists. generation.py sets
#     p_key.left_indent = Inches(0.15); p_key.first_line_indent = Inches(-0.15)
# inside a try/except -- but python-pptx's _Paragraph has NO such properties
# (checked: 1.0.2 exposes alignment/level/line_spacing/space_after/space_before
# and nothing else). Assigning them creates two ordinary Python attributes,
# writes no XML, and raises nothing for the except to catch. On this machine's
# template the resulting paragraphs carry an empty <a:pPr>, an empty lstStyle,
# and are not placeholders -- so marL/indent resolve to 0 and the hanging indent
# the model subtracts 10.8pt for does not exist. The model therefore gives every
# WRAPPED line 10.8pt less room than it really has, over-counts lines, and
# under-fills slots. The user's template is a different file, which is exactly
# why this is measured per-deck rather than asserted here.

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _indent_source(shape, paragraph, prs) -> Tuple[float, float, str]:
    """(marL_pt, indent_pt, where it came from) for this paragraph.

    Walks the same inheritance chain PowerPoint does, nearest first: the
    paragraph's own pPr, then the shape's lstStyle lvl1pPr, then the
    presentation-level defaultTextStyle. Placeholders can inherit further from
    the layout and master, so those are reported as such rather than guessed at.
    """
    def _read(el):
        if el is None:
            return None
        marL, ind = el.get("marL"), el.get("indent")
        if marL is None and ind is None:
            return None
        return (int(marL or 0) / 12700.0, int(ind or 0) / 12700.0)

    got = _read(paragraph._p.find(_A + "pPr"))
    if got:
        return got[0], got[1], "paragraph pPr"

    try:
        lst = shape.text_frame._txBody.find(_A + "lstStyle")
        got = _read(None if lst is None else lst.find(_A + "lvl1pPr"))
        if got:
            return got[0], got[1], "shape lstStyle"
    except Exception:
        pass

    try:
        dts = prs.part._element.find(_A.replace("drawingml/2006/main", "drawingml/2006/main")
                                     + "defaultTextStyle")
    except Exception:
        dts = None
    got = _read(None if dts is None else dts.find(_A + "lvl1pPr"))
    if got:
        return got[0], got[1], "presentation defaultTextStyle"

    where = "nothing declared (placeholder — may still inherit from the layout/master)" \
        if getattr(shape, "is_placeholder", False) else "nothing declared anywhere"
    return 0.0, 0.0, where


def _check_assumptions(prs, rows: List[ShapeRow], packing: Dict) -> List[str]:
    """One line per assumption, with what the deck actually says."""
    out: List[str] = []
    shapes_by_key = {(r.slide, r.shape): r for r in rows}
    sizes, spacings, gaps, latins, eas, autofits = set(), set(), set(), set(), set(), set()
    indents: Dict[Tuple[float, float, str], int] = {}
    placeholders = set()

    for s_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            key = (s_idx, str(getattr(shape, "name", "") or "") or "(unnamed)")
            if key not in shapes_by_key or shapes_by_key[key].is_empty:
                continue
            tf = shape.text_frame
            if getattr(shape, "is_placeholder", False):
                placeholders.add(key[1])
            try:
                autofits.add(",".join(e.tag.split("}")[-1] for e in tf._txBody.bodyPr) or "(none)")
            except Exception:
                pass
            for para in tf.paragraphs:
                if not (para.text or "").strip():
                    continue
                if para.line_spacing is not None:
                    spacings.add(round(float(para.line_spacing), 3))
                for attr in ("space_after", "space_before"):
                    v = getattr(para, attr)
                    if v is not None and v.pt > 0:
                        gaps.add(round(v.pt, 2))
                if (para.text or "").lstrip().startswith(BULLET_MARKER):
                    ind = _indent_source(shape, para, prs)
                    indents[ind] = indents.get(ind, 0) + 1
                for run in para.runs:
                    if run.font.size is not None:
                        sizes.add(round(run.font.size.pt, 2))
                    if run.font.name:
                        latins.add(run.font.name)
                    try:
                        rPr = run._r.find(_A + "rPr")
                        ea = None if rPr is None else rPr.find(_A + "ea")
                        if ea is not None and ea.get("typeface"):
                            eas.add(ea.get("typeface"))
                    except Exception:
                        pass

    def _cmp(label: str, expected, found, note: str = "") -> None:
        found_s = ", ".join(str(f) for f in sorted(found)) if found else "(nothing declared)"
        ok = (len(found) == 1 and next(iter(found)) == expected) if found else False
        out.append(f"  {'OK  ' if ok else '  !!'} {label:<26} model assumes {expected!r:<20} "
                   f"deck says {found_s}{('  — ' + note) if note and not ok else ''}")

    _cmp("font size (pt)", _real_font_size_pt(False), sizes)
    _cmp("line spacing", _real_line_spacing(False), spacings)
    _cmp("paragraph gap (pt)", _real_para_gap_pt(False), gaps)
    _cmp("latin typeface", _metrics_family(False, packing), latins,
         "the metrics table measuring this text belongs to a different family")
    _cmp("east-asian typeface", _metrics_family(True, packing), eas,
         "CJK glyph widths come from <a:ea>; if it is absent PowerPoint picks its own fallback")

    for (marL, ind, src), n in sorted(indents.items(), key=lambda kv: -kv[1]):
        effective_hang = max(0.0, marL) if ind < 0 else 0.0
        ok = abs(effective_hang - BULLET_HANGING_INDENT_PT) < 0.5
        out.append(f"  {'OK  ' if ok else '  !!'} {'bullet hanging indent':<26} "
                   f"model assumes {BULLET_HANGING_INDENT_PT}pt        "
                   f"deck says marL={marL:.1f}pt indent={ind:.1f}pt "
                   f"-> {effective_hang:.1f}pt  [{src}, {n} bullets]")
    if not indents:
        out.append("       (no ■ bullets found to check the indent against)")

    out.append(f"  {'OK  ' if autofits <= {'noAutofit'} else '  !!'} {'autofit':<26} "
               f"model assumes {'noAutofit'!r:<20} deck says "
               f"{', '.join(sorted(autofits)) or '(none)'}")
    if placeholders:
        out.append(f"       note: {', '.join(sorted(placeholders))} are PLACEHOLDERS — they can "
                   f"inherit paragraph properties from the layout/master that are not visible here")
    return out


def _metrics_family(is_chi: bool, packing: Dict) -> str:
    """Family name of the metrics table actually in use, so the check compares
    against the ruler rather than against the config's fallback family."""
    path = _resolve_font_metrics_path(is_chi, packing)
    if path:
        try:
            import json
            with open(path, encoding="utf-8") as fh:
                fam = json.load(fh).get("family")
            if fam:
                return str(fam)
        except Exception:
            pass
    return _measurer_family(is_chi, packing)


# ---------------------------------------------------------------------------
# The model half (runs anywhere -- this is what --model-only exercises)
# ---------------------------------------------------------------------------

def _para_kind(text: str, space_after_pt: float, starts_bold: bool) -> str:
    """Classify a paragraph the way the renderer built it.

    A category header is the only non-blank paragraph the renderer gives
    space_after = Pt(0) (generation.py's p_category); every bullet and
    continuation gets Pt(3). Read from the shape rather than pattern-matched,
    because prose continuations look exactly like category text.
    """
    stripped = text.strip()
    if not stripped:
        return "blank"
    if stripped.startswith(BULLET_MARKER):
        return "bullet"
    if stripped.startswith(EXPLAIN_MARKERS):
        return "explain"
    if space_after_pt == 0.0 and not starts_bold:
        return "category"
    return "continuation"


def _collect_model(deck_path: str, want_slide: Optional[int], want_shape: Optional[str]):
    prs = Presentation(deck_path)
    packing = (_load_pptx_settings() or {})
    packing = packing.get("packing", packing)
    alias_idx = _alias_index()

    measurers, bold_fonts, env = {}, {}, {}
    for is_chi in (False, True):
        mpath = _resolve_font_metrics_path(is_chi, packing)
        m = get_measurer(
            _measurer_family(is_chi, packing), _real_font_size_pt(is_chi),
            is_cjk=is_chi, line_spacing=_real_line_spacing(is_chi), metrics_path=mpath,
        )
        measurers[is_chi] = m
        bold_fonts[is_chi] = _bold_font(is_chi, _real_font_size_pt(is_chi))
        tag = "CHI" if is_chi else "ENG"
        bold_note = "bold face NOT FOUND — variant C not scored"
        if bold_fonts[is_chi] is not None:
            bold_note = f"bold face {os.path.basename(getattr(bold_fonts[is_chi], 'path', '?'))}"
        env[tag] = (f"{m.source:<14} size={m.size_pt}pt  spacing={m.line_spacing}  "
                    f"line_h={m.line_height_pt():.2f}pt  gap={_real_para_gap_pt(is_chi):.2f}pt"
                    f"\n                 {mpath or '(no metrics.json — system font)'}"
                    f"\n                 {bold_note}")

    rows: List[ShapeRow] = []
    warnings: List[str] = []
    for s_idx, slide in enumerate(prs.slides, start=1):
        if want_slide and s_idx != want_slide:
            continue
        for shape in slide.shapes:
            name = str(getattr(shape, "name", "") or "")
            if not getattr(shape, "has_text_frame", False):
                continue
            tf = shape.text_frame
            text = tf.text or ""
            is_slot = name.startswith("textMainBullets")
            # An EMPTY commentary slot is the largest under-fill there is, so it
            # has to appear in the fill picture rather than be filtered out of
            # it. Skipping empties made a page whose whole second column went
            # unused look like a page that simply had fewer columns.
            if not text.strip() and not is_slot:
                continue
            # Commentary slots, plus the unnamed textboxes
            # _render_table_accounts_stack drops beside a table.
            if not (is_slot or BULLET_MARKER in text):
                continue
            if want_shape and want_shape.lower() not in name.lower():
                continue

            box = text_box_from_shape(shape)
            is_chi = _is_chinese_text(text)
            measurer = measurers[is_chi]
            bold_font = bold_fonts[is_chi]
            line_h = measurer.line_height_pt()
            hang_w = max(10.0, box.width_pt - BULLET_HANGING_INDENT_PT)

            row = ShapeRow(slide=s_idx, shape=name or "(unnamed)", slot=_slot_of(name),
                           is_chinese=is_chi, box_w_pt=box.width_pt,
                           box_h_pt=box.height_pt, model_lines=0, model_pt=0.0,
                           is_empty=not text.strip(), measurer=measurer)
            if row.is_empty:
                # No paragraphs to walk, and PowerPoint reports HasText false so
                # ground truth will never match it -- record it as a real,
                # measured 0% and move on.
                row.real_lines, row.real_bound_pt = 0, 0.0
                rows.append(row)
                continue

            for p_idx, para in enumerate(tf.paragraphs, start=1):
                p_text = para.text or ""
                sa = para.space_after.pt if para.space_after is not None else 0.0
                sb = para.space_before.pt if para.space_before is not None else 0.0
                runs = list(para.runs)
                # _render_table_accounts_stack reserves a table's vertical
                # space as BLANK paragraphs inside the same frame, and sizes
                # the last one to a FRACTION of a line -- real exports carry
                # runs at 1.0, 3.66, 4.33, 6.33 and 7.0pt. Pricing every
                # paragraph at the deck's 9pt made model_pt overstate every
                # table slide (by 14.0pt on one real slot) and dragged the
                # fitted pitch down to 10.11pt, i.e. this file blaming the
                # model for its own assumption. Same mistake as the
                # space_before one; inspect_pptx.py already reads the size.
                _sizes = [r.font.size.pt for r in runs if r.font.size is not None]
                font_pt = max(_sizes) if _sizes else _real_font_size_pt(is_chi)
                has_bold = any(bool(r.font.bold) for r in runs)
                starts_bold = bool(runs and runs[0].font.bold)
                kind = _para_kind(p_text, sa, starts_bold)

                label = mapping_key = ""
                selfcheck_ok = True
                if kind == "blank":
                    n_a = n_b = n_c = 1
                else:
                    # Production call, unmodified. This is variant B: the text
                    # that is really on the slide, regular weight throughout.
                    # Both markers hang: "\u25a0 " on a lead-in and "\u27a2 " on a
                    # post-table explanation. Line 1 spans the box, wrapped
                    # lines sit one indent in. A continuation or category
                    # paragraph is narrow on every line.
                    first_w = box.width_pt if kind in ("bullet", "explain") else None
                    n_b = max(1, len(measurer.wrap(p_text, hang_w, first_line_width_pt=first_w)))
                    n_a = n_c = n_b

                    if kind == "bullet":
                        bold_run = next((r.text for r in runs if r.font.bold), "")
                        label = (bold_run or "").strip()
                        head, sep_found, tail = p_text.partition(bold_run) if bold_run else ("", "", "")
                        first_w = box.width_pt

                        # Self-check: the tool-local mixed wrapper must agree
                        # with the production wrapper when nothing is bold.
                        if sep_found:
                            n_plain = _wrap_runs([(head, False), (bold_run, False), (tail, False)],
                                                 measurer, bold_font, first_w, hang_w)
                            selfcheck_ok = (n_plain == n_b)
                            if not selfcheck_ok:
                                # Print the evidence, not just the verdict. The two
                                # wrappers tokenise differently -- wrap_text_with_
                                # metrics (the metrics.json backend production uses
                                # here) keeps whitespace as its own token and charges
                                # its width, while _atomize (the Pillow backend, and
                                # what the local loop mirrors) drops whitespace and
                                # only re-inserts a space between two LATIN atoms. A
                                # stray space inside CJK prose therefore costs one of
                                # them nothing and the other ~3pt. Which is right is a
                                # question for PowerPoint, so show the break-up.
                                _prod = measurer.wrap(p_text, hang_w,
                                                      first_line_width_pt=first_w)
                                warnings.append(
                                    f"slide {s_idx} {name} para {p_idx}: local wrapper says "
                                    f"{n_plain} lines where measurer.wrap says {n_b}"
                                    + f"\n       text: {p_text[:90]}"
                                    + "".join(f"\n       measurer.wrap L{i}: {ln}"
                                              for i, ln in enumerate(_prod, 1)))

                            # Variant C: same text, key run measured bold.
                            if bold_font is not None:
                                n_c = _wrap_runs([(head, False), (bold_run, True), (tail, False)],
                                                 measurer, bold_font, first_w, hang_w)

                            # Variant A: what the packer charged -- mapping_key
                            # in place of the rendered label.
                            cands = alias_idx.get(label, [])
                            if len(cands) == 1:
                                mapping_key = cands[0]
                                if mapping_key != label:
                                    n_a = max(1, len(measurer.wrap(
                                        head + mapping_key + tail, hang_w,
                                        first_line_width_pt=first_w)))

                row.paras.append(ParaRow(
                    slide=s_idx, shape=row.shape, slot=row.slot, index=p_idx,
                    kind=kind, has_bold=has_bold, chars=len(p_text),
                    lines_a=n_a, lines_b=n_b, lines_c=n_c, label=label,
                    mapping_key=mapping_key, space_after_pt=sa, space_before_pt=sb,
                    selfcheck_ok=selfcheck_ok, text=p_text,
                ))
                row.model_lines += n_b
                row.model_pt += n_b * (font_pt * POWERPOINT_LINE_PITCH_FACTOR) + sa + sb
            # The final paragraph's space_after is invisible padding at the
            # bottom of the frame, not occupied height -- the same correction
            # _calculate_content_lines makes. Whether BoundHeight agrees is one
            # of the things the CALIBRATION fit settles.
            if row.paras:
                row.model_pt -= row.paras[-1].space_after_pt
            rows.append(row)
    return rows, env, warnings, prs, packing


# ---------------------------------------------------------------------------
# The ground-truth half (Windows + PowerPoint only)
# ---------------------------------------------------------------------------

def _attach_powerpoint() -> Tuple[object, bool, str]:
    """Return (app, started_by_us, binding).

    EARLY binding first (gencache builds a wrapper from PowerPoint's type
    library, so Lines/Paragraphs are known to be methods and return real range
    objects). Late binding is the fallback and works too -- see _sub_range,
    whose call ladder handles both, and the record above it of the two forms
    that failed on a real Windows run.

    Reuses a running PowerPoint when there is one, so this never quits an
    instance the user already had open with their own work in it.
    """
    import win32com.client as win32
    from win32com.client import gencache

    def _early(target):
        try:
            return gencache.EnsureDispatch(target)
        except Exception:
            return None   # stale/unwritable gen_py cache — fall back, don't die

    try:
        raw = win32.GetActiveObject("PowerPoint.Application")
        started = False
    except Exception:
        raw, started = None, True

    if raw is not None:
        app = _early(raw)
        return (app or raw), False, ("early" if app else "late")
    app = _early("PowerPoint.Application")
    return (app or win32.Dispatch("PowerPoint.Application")), True, ("early" if app else "late")


# How TextRange2.Lines / .Paragraphs(k) actually have to be called.
#
# measure_boundheight.bas does `shp.TextFrame2.TextRange.Lines.Count` and that
# works, so the member exists. Getting there from Python took three attempts,
# all recorded here because each one FAILED in a way that looked like the
# previous fix:
#
#   1. `tr.Lines.Count` -- pywin32 late binding resolved Lines as a plain
#      property get, PowerPoint returned the range, and the dynamic wrapper
#      collapsed it to its default member (Text). Result: a **str**, no
#      exception, `.Count` became str.count. Silent wrong answer.
#   2. `_FlagAsMethod("Lines")` then `tr.Lines()` -- this is the documented
#      late-binding fix for members that are methods, and it made things worse:
#      Lines is declared PROPGET in the type library, so invoking it with
#      DISPATCH_METHOD alone no longer matches any member and COM answers
#      "Member not found" (-2147352573). Do NOT put _FlagAsMethod back.
#   3. What works: invoke it explicitly with DISPATCH_PROPERTYGET |
#      DISPATCH_METHOD and the arguments inline -- a parameterised property
#      get, which is exactly what VBA is doing when it writes `.Lines` and
#      `.Paragraphs(k)` with the optional arguments omitted.
#
# The ladder below tries the VBA-equivalent form first and keeps the other
# forms as fallbacks, because early binding (when gencache succeeds) exposes
# these as ordinary callables instead. Whichever one works is reported, so a
# future failure says which rung broke rather than just "it broke".

_STRATEGY_USED: List[str] = []


def _via_propget(obj, member: str, args: Tuple):
    """Parameterised property get -- the VBA-equivalent call."""
    import pythoncom
    from win32com.client import Dispatch
    dispid = obj._oleobj_.GetIDsOfNames(0, member)
    res = obj._oleobj_.Invoke(
        dispid, 0, pythoncom.DISPATCH_PROPERTYGET | pythoncom.DISPATCH_METHOD,
        True, *args)
    try:
        return Dispatch(res)
    except Exception:
        return res


def _via_call(obj, member: str, args: Tuple):
    """Early-bound: gencache generated an ordinary method."""
    return getattr(obj, member)(*args)


def _via_attr(obj, member: str, args: Tuple):
    """Whatever plain attribute access gives, if it happens to be a range."""
    if args:
        raise TypeError("attribute access cannot take an index")
    return getattr(obj, member)


def _sub_range(text_range, member: str, index: Optional[int] = None):
    """Get TextRange2.Lines / .Paragraphs(k) as a real range object.

    Every ground-truth number this tool prints comes through here, so the
    late-binding traps above can only bite in one place -- and so that a bite
    RAISES rather than quietly reporting a number that came from somewhere
    else. `.Count` is touched before returning for exactly that reason: it is
    the cheapest proof that what came back is a range and not a string.
    """
    args: Tuple = () if index is None else (index,)
    failures: List[str] = []
    for name, fn in (("propget", _via_propget), ("call", _via_call), ("attr", _via_attr)):
        try:
            res = fn(text_range, member, args)
            _ = res.Count          # must really be a range, not a str or a method
            if name not in _STRATEGY_USED:
                _STRATEGY_USED.append(name)
            return res
        except Exception as exc:
            failures.append(f"{name}: {type(exc).__name__}: {exc}")
    raise RuntimeError(
        f"Could not read TextRange2.{member}"
        f"{'' if index is None else f'({index})'} through any call form.\n  "
        + "\n  ".join(failures)
        + "\n\nThis is the one thing the tool cannot work around. Paste the above "
          "back;\nthe VBA in measure_boundheight.bas reads the same member "
          "successfully,\nso the member exists and only the Python call form is wrong.")


def _normalize_com_text(text: str) -> str:
    """COM returns \\r for paragraph breaks and \\x0b for soft line breaks."""
    return str(text or "").replace("\r\n", "\n").replace("\r", "\n").replace("\x0b", "\n").strip()


def _fill_ground_truth(deck_path: str, rows: List[ShapeRow], warnings: List[str]) -> str:
    """Open the deck in real PowerPoint and record what its layout engine did.
    Read-only, closed without saving."""
    app, started_by_us, binding = _attach_powerpoint()
    version, pres = "", None
    try:
        try:
            app.Visible = True   # PowerPoint refuses invisible automation
        except Exception:
            pass
        version = (f"PowerPoint {getattr(app, 'Version', '?')} "
                   f"build {getattr(app, 'Build', '?')} ({binding}-bound)")
        # Positional, not keyword: late-bound Dispatch resolves named arguments
        # through GetIDsOfNames and that is not reliable across Office builds.
        # Open(FileName, ReadOnly, Untitled, WithWindow); msoTrue = -1, and
        # WithWindow MUST be true or PowerPoint never lays the text out and
        # BoundHeight comes back as 0.
        pres = app.Presentations.Open(os.path.abspath(deck_path), -1, 0, -1)

        by_slide: Dict[int, List[ShapeRow]] = {}
        for r in rows:
            if r.is_empty:
                continue   # TextFrame2.HasText is false; already recorded as a real 0%
            by_slide.setdefault(r.slide, []).append(r)

        for s_idx, wanted in by_slide.items():
            if s_idx > pres.Slides.Count:
                continue
            sld = pres.Slides(s_idx)
            unmatched = list(wanted)
            for j in range(1, sld.Shapes.Count + 1):
                shp = sld.Shapes(j)
                try:
                    if not shp.HasTextFrame or not shp.TextFrame2.HasText:
                        continue
                except Exception:
                    continue
                com_name = str(shp.Name or "")
                tr = shp.TextFrame2.TextRange
                com_text = _normalize_com_text(tr.Text)

                target = next((r for r in unmatched if r.shape == com_name), None)
                if target is None:
                    target = next(
                        (r for r in unmatched
                         if _normalize_com_text("\n".join(p.text for p in r.paras)) == com_text),
                        None)
                if target is None:
                    continue
                unmatched.remove(target)

                target.real_lines = int(_sub_range(tr, "Lines").Count)
                target.real_bound_pt = float(tr.BoundHeight)
                n_com = int(_sub_range(tr, "Paragraphs").Count)
                for p in target.paras:
                    if p.index > n_com:
                        continue
                    try:
                        para_range = _sub_range(tr, "Paragraphs", p.index)
                        lines = _sub_range(para_range, "Lines")
                        p.real_lines = int(lines.Count)
                    except Exception:
                        p.real_lines = None
                        continue
                    # Only where we got it WRONG, because this is one extra COM
                    # round trip per line and it is only ever read for those.
                    # Lines(i).Text is the literal text PowerPoint drew on line
                    # i -- it does not say how wide the model thinks that line
                    # is, it says where PowerPoint actually broke, which is the
                    # only thing that can settle WHY it broke there. Guessing at
                    # the rule from aggregate counts has now been wrong twice
                    # (mixed-script widths, then punctuation compression).
                    if p.real_lines != p.lines_c:
                        try:
                            p.real_line_texts = [
                                str(_sub_range(para_range, "Lines", i).Text)
                                for i in range(1, p.real_lines + 1)]
                        except Exception as exc:
                            warnings.append(f"slide {s_idx} {com_name} para {p.index}: "
                                            f"could not read the real line breaks ({exc})")
                if n_com != len(target.paras):
                    # Not cosmetic: the two sides disagree about what a
                    # paragraph even is, so every per-paragraph delta below it
                    # compares different things.
                    warnings.append(f"slide {s_idx} {com_name}: PowerPoint sees {n_com} "
                                    f"paragraphs, python-pptx sees {len(target.paras)} — "
                                    f"per-paragraph rows may be misaligned")
            for leftover in unmatched:
                warnings.append(f"slide {s_idx} {leftover.shape}: no matching shape in "
                                f"PowerPoint — not scored")
    finally:
        try:
            if pres is not None:
                pres.Close()
        except Exception:
            pass
        if started_by_us:
            try:
                app.Quit()
            except Exception:
                pass
    if _STRATEGY_USED:
        version += f", range access via {'/'.join(_STRATEGY_USED)}"
    return version


# ---------------------------------------------------------------------------
# Calibration
# ---------------------------------------------------------------------------

def _gap_count(row: ShapeRow, *, drop_last: bool) -> int:
    """How many 3pt inter-paragraph gaps this shape's text really contains.

    space_BEFORE counts too. Leaving it out is not a rounding matter: the
    renderer puts space_before = Pt(3) on every category header after the first
    (generation.py's p_category), so a slot with three category groups carries
    two gaps this used to miss entirely. On the first real Windows run that
    alone pulled the fitted pitch to 10.88pt / gap 3.08pt with an rmse of
    0.53pt, making the model look ~1% off when the SAME measurements resolve to
    lines x 10.80 + gaps x 3.00 with a residual of exactly zero on every shape.
    A calibration that flatters or maligns the thing it calibrates is worse
    than none.
    """
    n = sum(1 for p in row.paras if p.space_after_pt > 0)
    if drop_last and row.paras and row.paras[-1].space_after_pt > 0:
        n -= 1
    n += sum(1 for p in row.paras if p.space_before_pt > 0)
    return n


def _fit_pitch_and_gap(rows: List[ShapeRow]) -> Optional[Dict[str, float]]:
    """Least-squares-fit pitch and gap out of the real measurements:

        BoundHeight = lines * pitch + gap_count * gap

    Fitted twice -- charging a gap for every paragraph, and for all but the last
    -- because which of those BoundHeight includes is precisely the correction
    _calculate_content_lines makes on faith.
    """
    try:
        import numpy as np
    except ImportError:
        return None
    usable = [r for r in rows if r.real_lines and r.real_bound_pt]
    if len(usable) < 3:
        return None

    out: Dict[str, float] = {"n": float(len(usable))}
    for label, drop_last in (("gap_per_para", False), ("gap_between_paras", True)):
        # Weighted, not raw: see ShapeRow.real_pitch. A raw count makes a
        # table slide's 1pt spacers look like short 9pt lines and pulls the
        # fitted pitch to 10.11pt on a deck that renders at exactly 10.80.
        A = np.array([[r.equivalent_9pt_lines, _gap_count(r, drop_last=drop_last)]
                      for r in usable], float)
        b = np.array([r.real_bound_pt for r in usable], float)
        # Reported alongside: the pitch each shape implies once its OWN
        # measured spacing is removed. If the fit and these disagree, the
        # gap model is what is wrong, not the pitch.
        implied = [r.real_pitch for r in usable if r.real_pitch]
        if implied:
            out['implied_pitch_min'] = min(implied)
            out['implied_pitch_max'] = max(implied)
            out['implied_pitch_mean'] = sum(implied) / len(implied)
        sol, *_ = np.linalg.lstsq(A, b, rcond=None)
        out[f"{label}_pitch"] = float(sol[0])
        out[f"{label}_gap"] = float(sol[1])
        out[f"{label}_rmse"] = float(np.sqrt(np.mean((A @ sol - b) ** 2)))
    return out


# ---------------------------------------------------------------------------
# Report
# ---------------------------------------------------------------------------

VARIANTS = (("a", "today's packer (mapping_key, regular)"),
            ("b", "+ rendered label"),
            ("c", "+ bold-aware key run"))


def _report(rows: List[ShapeRow], env: Dict[str, str], version: str,
            warnings: List[str], show_lines: bool, model_only: bool,
            assumptions: Optional[List[str]] = None) -> int:
    print("=" * 96)
    print("RENDER TRUTH — real PowerPoint layout vs this repo's model")
    print("=" * 96)
    print("\nMeasurement source (the ruler production itself uses):")
    for tag in ("ENG", "CHI"):
        print(f"  [{tag}] {env.get(tag, '?')}")
    print(f"  ground truth: {version or 'SKIPPED (--model-only)'}")

    if assumptions:
        print("\n" + "-" * 96)
        print("MODEL ASSUMPTIONS vs THIS DECK  — every number below depends on these holding.")
        print("                                 A `!!` means the model is measuring something")
        print("                                 the deck does not actually contain.")
        print("-" * 96)
        for line in assumptions:
            print(line)

    if warnings:
        print("\n  !! " + "\n  !! ".join(warnings[:12]))
        if len(warnings) > 12:
            print(f"  !! ... and {len(warnings) - 12} more")

    # ---- per shape -------------------------------------------------------
    print("\n" + "-" * 96)
    print("PER SHAPE   fill% is BoundHeight/box — the numerator came from PowerPoint, not us.")
    print("            pitch is BoundHeight less the paragraph gaps, over the real line count")
    print("            weighted by each paragraph's own font size (a table slot's spacers are")
    print(f"            sized down); below the nominal {_real_font_size_pt(False) * 1.2:.2f}pt "
          f"means PowerPoint re-ran its own autofit.")
    print("-" * 96)
    print(f"{'sl':>3} {'shape':<22}{'lang':>5}{'box_w':>8}{'box_h':>8}{'mLines':>7}{'rLines':>7}"
          f"{'d':>4}{'model_pt':>10}{'BoundH':>9}{'pitch':>7}{'fill%':>8}")
    n_shape_bad = 0
    for r in sorted(rows, key=lambda x: (x.slide, x.shape)):
        d = "" if r.real_lines is None else f"{r.real_lines - r.model_lines:+d}"
        if d not in ("", "+0"):
            n_shape_bad += 1
        print(f"{r.slide:>3} {r.shape[:22]:<22}{'CHI' if r.is_chinese else 'ENG':>5}"
              f"{r.box_w_pt:8.1f}{r.box_h_pt:8.1f}{r.model_lines:7d}"
              f"{(r.real_lines if r.real_lines is not None else ''):>7}{d:>4}{r.model_pt:10.1f}"
              f"{(f'{r.real_bound_pt:.1f}' if r.real_bound_pt is not None else ''):>9}"
              f"{(f'{r.real_pitch:.2f}' if r.real_pitch else ''):>7}"
              f"{(f'{r.real_fill*100:.1f}' if r.real_fill is not None else ''):>8}"
              + ('  EMPTY' if r.is_empty else ''))

    all_paras = [p for r in rows for p in r.paras]
    scored = [p for p in all_paras if p.real_lines is not None]

    # The one thing that IS verifiable without PowerPoint: that the tool-local
    # mixed-weight wrapper reproduces the production wrapper exactly when no run
    # is bold. If it does not, variant C is measuring the tool's own bug rather
    # than the bold key run, so this has to be checked before the numbers mean
    # anything -- and it is checked on every run, not just --model-only.
    checked = [p for p in all_paras if p.kind == "bullet"]
    failed = [p for p in checked if not p.selfcheck_ok]
    print("\n" + "-" * 96)
    print(f"SELF-CHECK  local mixed-weight wrapper vs production measurer.wrap(): "
          f"{len(checked) - len(failed)}/{len(checked)} agree")
    print("-" * 96)
    if failed:
        print("  variant C is NOT trustworthy on the rows listed in the warnings above.")

    if model_only:
        print("\n" + "=" * 96)
        print("MODEL-ONLY RUN — no verdict on accuracy (that needs PowerPoint).")
        print("Re-run on Windows without --model-only.")
        print("=" * 96)
        return 0

    # ---- the scoreboard: which variant does PowerPoint agree with? -------
    print("\n" + "-" * 96)
    print("VARIANT SCOREBOARD  — how often each prediction matches PowerPoint, per paragraph")
    print("-" * 96)
    bullets = [p for p in scored if p.kind == "bullet"]
    if not bullets:
        print("  no ■ bullets scored")
    else:
        print(f"{'variant':<40}{'exact':>8}{'rate':>8}{'net lines':>11}{'worst':>7}")
        for v, name in VARIANTS:
            deltas = [p.delta(v) for p in bullets]
            exact = sum(1 for d in deltas if d == 0)
            print(f"{name:<40}{exact:>8}{exact/len(bullets)*100:7.1f}%"
                  f"{sum(deltas):>+11d}{max(abs(d) for d in deltas):>7}")
        print(f"\n  ({len(bullets)} bullets scored; 'net lines' is real minus predicted — "
              f"positive means\n   the model UNDER-counts and the slot will overflow, negative "
              f"means it over-counts\n   and the slot is left under-filled)")
        best = max("abc", key=lambda v: sum(1 for p in bullets if p.delta(v) == 0))
        print(f"\n  -> PowerPoint agrees most with variant '{best.upper()}'.")
        if best == "a":
            print("     The proposed fixes do NOT help on this deck; the error is elsewhere.")

    # non-bullet paragraphs use one prediction for all three variants
    others = [p for p in scored if p.kind != "bullet"]
    if others:
        bad = [p for p in others if p.delta("b") != 0]
        print(f"\n  non-bullet paragraphs: {len(bad)}/{len(others)} mis-wrapped, "
              f"net {sum(p.delta('b') for p in others):+d} lines")

    # ---- mis-wrapped paragraphs -----------------------------------------
    wrong = [p for p in scored if p.delta("c") != 0]
    print("\n" + "-" * 96)
    print("STILL MIS-WRAPPED UNDER THE BEST VARIANT (C)  — this is what remains unexplained")
    print("-" * 96)
    if not wrong:
        print(f"  none — all {len(scored)} paragraphs wrapped exactly as variant C predicts")
    else:
        print(f"{'sl':>3} {'shape':<20}{'#':>4}{'kind':>13}{'A':>3}{'B':>3}{'C':>3}"
              f"{'real':>5}{'d':>4}  text")
        for p in wrong:
            print(f"{p.slide:>3} {p.shape[:20]:<20}{p.index:>4}{p.kind:>13}"
                  f"{p.lines_a:>3}{p.lines_b:>3}{p.lines_c:>3}{p.real_lines:>5}"
                  f"{p.delta('c'):>+4}  {p.text[:30]}")

    if any(not p.selfcheck_ok for p in all_paras):
        print("\n  !! the tool-local mixed wrapper disagreed with measurer.wrap() on some "
              "paragraphs;\n     variant C is NOT trustworthy on those rows (listed in the "
              "warnings above).")

    # ---- calibration -----------------------------------------------------
    fit = _fit_pitch_and_gap(rows)
    print("\n" + "-" * 96)
    print("CALIBRATION  (fitted from the real measurements, not assumed)")
    print("-" * 96)
    if fit is None:
        print("  (needs numpy and at least 3 measured shapes)")
    else:
        print(f"  model constants:  pitch {_real_font_size_pt(False) * 1.2:.2f}pt   "
              f"gap {_real_para_gap_pt(False):.2f}pt")
        for label in ("gap_per_para", "gap_between_paras"):
            print(f"  fit [{label:<18}] pitch {fit[label+'_pitch']:6.2f}pt   "
                  f"gap {fit[label+'_gap']:5.2f}pt   rmse {fit[label+'_rmse']:5.2f}pt")
        best_fit = min(("gap_per_para", "gap_between_paras"), key=lambda k: fit[k + "_rmse"])
        print(f"  -> BoundHeight here is best explained by '{best_fit}' "
              f"(n={int(fit['n'])} shapes)")
        if "implied_pitch_mean" in fit:
            print(f"  per-shape implied pitch (BoundHeight less that shape's OWN measured "
                  f"spacing,\n  over its size-weighted line count): "
                  f"{fit['implied_pitch_min']:.3f} .. {fit['implied_pitch_max']:.3f}pt, "
                  f"mean {fit['implied_pitch_mean']:.3f}pt")

        nominal = _real_font_size_pt(False) * POWERPOINT_LINE_PITCH_FACTOR
        odd = [r for r in rows if r.real_pitch and abs(r.real_pitch - nominal) > 0.02]
        if odd:
            print(f"\n  {len(odd)} shape(s) do not land on {nominal:.2f}pt. Breaking each one")
            print("  down rather than reasoning about it -- this is the intermediate:")
            for r in odd:
                print(f"    slide {r.slide} {r.shape}: BoundH={r.real_bound_pt:.1f} "
                      f"gap_total={r.gap_total_pt:.1f} equiv_lines={r.equivalent_9pt_lines:.3f} "
                      f"-> implied {r.real_pitch:.3f}pt   (model_pt={r.model_pt:.1f})")
                sizes: Dict[float, List[int]] = {}
                for q in r.paras:
                    n = q.real_lines if q.real_lines is not None else q.lines_b
                    sizes.setdefault(q.font_pt, []).append(n)
                for size in sorted(sizes):
                    ns = sizes[size]
                    print(f"        {len(ns):>3} para(s) at {size:>5.2f}pt "
                          f"holding {sum(ns):>3} line(s)")

    if show_lines and wrong:
        print("\n" + "-" * 96)
        print("WHERE POWERPOINT ACTUALLY BROKE THE LINE")
        print("-" * 96)
        print("  `our_w` is this repo's own width for the text PowerPoint put on that line, and")
        print("  `limit` is the width we believe it had. A line whose our_w EXCEEDS its limit is")
        print("  PowerPoint fitting text our widths say cannot fit -- the `over` column is by how")
        print("  much, and the last characters of that line are what bought the room. That is the")
        print("  measurement; the rule follows from it. (our_w ignores the bold key run, so a")
        print("  small positive `over` on line 1 of a bullet is expected and not the finding.)")
        by_shape = {(r.slide, r.shape): r for r in rows}
        for p in wrong:
            print(f"\n  slide {p.slide} {p.shape} para {p.index} [{p.kind}] "
                  f"A={p.lines_a} B={p.lines_b} C={p.lines_c} real={p.real_lines}"
                  + (f"\n  label drawn: {p.label!r}   packer charged: {p.mapping_key!r}"
                     if p.label else ""))
            shape_row = by_shape.get((p.slide, p.shape))
            measurer = getattr(shape_row, "measurer", None)
            if not p.real_line_texts or measurer is None:
                print(f"    (no line breaks captured)\n    {p.text}")
                continue
            box_w = shape_row.box_w_pt
            hang_w = max(10.0, box_w - BULLET_HANGING_INDENT_PT)
            print(f"    box_w={box_w:.1f}pt   hanging continuation width={hang_w:.1f}pt")
            print(f"    {'#':>3}{'our_w':>9}{'limit':>8}{'over':>8}  line as PowerPoint drew it")
            for i, line in enumerate(p.real_line_texts, start=1):
                limit = box_w if (i == 1 and p.kind in ("bullet", "explain")) else hang_w
                our_w = measurer.text_width_pt(line.rstrip("\r\n\x0b"))
                over = our_w - limit
                print(f"    {i:>3}{our_w:>9.1f}{limit:>8.1f}{over:>+8.1f}  "
                      f"{'!! ' if over > 0.5 else '   '}{line.rstrip()}")

    print("\n" + "=" * 96)
    net = sum(p.delta("c") for p in scored)
    print(f"VERDICT: under the best variant, {len(wrong)} of {len(scored)} paragraphs are still "
          f"mis-wrapped\n         (net {net:+d} lines deck-wide); {n_shape_bad} of {len(rows)} "
          f"shapes off on total lines.")
    print("=" * 96)
    return 1 if wrong else 0


def _write_csv(path: str, rows: List[ShapeRow]) -> None:
    with open(path, "w", newline="", encoding="utf-8-sig") as fh:
        w = csv.writer(fh)
        w.writerow(["slide", "shape", "slot", "para", "kind", "label", "mapping_key",
                    "chars", "lines_A", "lines_B", "lines_C", "real_lines",
                    "delta_A", "delta_B", "delta_C", "space_before_pt",
                    "space_after_pt", "selfcheck_ok", "text"])
        for r in rows:
            for p in r.paras:
                w.writerow([p.slide, p.shape, p.slot, p.index, p.kind, p.label,
                            p.mapping_key, p.chars, p.lines_a, p.lines_b, p.lines_c,
                            "" if p.real_lines is None else p.real_lines,
                            *["" if p.delta(v) is None else p.delta(v) for v in "abc"],
                            f"{p.space_before_pt:.2f}", f"{p.space_after_pt:.2f}",
                            int(p.selfcheck_ok), p.text])
    print(f"\nwrote {path}")


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.split("\n")[0])
    ap.add_argument("deck", help="an EXPORTED .pptx (not the template — it has no commentary)")
    ap.add_argument("--lines", action="store_true",
                    help="print the full text of every still-mis-wrapped paragraph")
    ap.add_argument("--slide", type=int, help="restrict to one slide (1-based)")
    ap.add_argument("--shape", help="restrict to shapes whose name contains this")
    ap.add_argument("--csv", help="write the per-paragraph table here")
    ap.add_argument("--model-only", action="store_true",
                    help="skip PowerPoint entirely (exercises the prediction half off Windows)")
    args = ap.parse_args()

    if not os.path.exists(args.deck):
        print(f"No such file: {args.deck}", file=sys.stderr)
        return 2

    rows, env, warnings, prs, packing = _collect_model(args.deck, args.slide, args.shape)
    if not rows:
        print("No commentary shapes found. Is this an exported deck rather than the template?",
              file=sys.stderr)
        return 2

    version = ""
    if not args.model_only:
        if sys.platform != "win32":
            print(f"This needs Windows + PowerPoint (running on {sys.platform}). "
                  f"Use --model-only to exercise the prediction half.", file=sys.stderr)
            return 2
        try:
            import win32com.client  # noqa: F401
        except ImportError:
            print("pywin32 is not installed. Run:  pip install pywin32", file=sys.stderr)
            return 2
        version = _fill_ground_truth(args.deck, rows, warnings)

    rc = _report(rows, env, version, warnings, args.lines, args.model_only,
                 assumptions=_check_assumptions(prs, rows, packing))
    if args.csv:
        _write_csv(args.csv, rows)
    return rc


if __name__ == "__main__":
    raise SystemExit(main())
