#!/usr/bin/env python3
"""
Debug script to inspect PowerPoint template shapes
"""

import os
import sys
sys.path.append('/Users/ytsang/Desktop/Github/python-pptx')

def inspect_template():
    """Inspect the PowerPoint template to see available shapes"""
    try:
        from pptx import Presentation

        template_path = 'fdd_utils/template.pptx'
        if not os.path.exists(template_path):
            print(f"❌ Template not found: {template_path}")
            return

        prs = Presentation(template_path)
        print(f"📊 Template inspection: {len(prs.slides)} slides found")
        print("=" * 80)

        for i, slide in enumerate(prs.slides):
            print(f"\n📄 SLIDE {i}:")
            print("-" * 40)

            # Count shapes by type
            text_shapes = []
            other_shapes = []

            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    text_shapes.append(shape)
                else:
                    other_shapes.append(shape)

            print(f"Total shapes: {len(slide.shapes)}")
            print(f"Text shapes: {len(text_shapes)}")
            print(f"Other shapes: {len(other_shapes)}")

            # Inspect text shapes
            for shape in text_shapes:
                has_text = bool(shape.text_frame.text.strip())
                print(f"  📝 Shape: '{shape.name}' (has text: {has_text})")
                if has_text:
                    preview = shape.text_frame.text[:100].replace('\n', ' ')
                    print(f"      Preview: {preview}{'...' if len(shape.text_frame.text) > 100 else ''}")

            # Look for specific shapes we're interested in
            target_shapes = ['textMainBullets', 'textMainBullets_L', 'textMainBullets_R', 'coSummaryShape']
            found_shapes = [s.name for s in text_shapes if s.name in target_shapes]

            if found_shapes:
                print(f"  ✅ Target shapes found: {found_shapes}")
            else:
                print("  ⚠️  No target shapes found")
        print("\n" + "=" * 80)
        print("🎯 TARGET SHAPES TO LOOK FOR:")
        print("  - textMainBullets (main content)")
        print("  - textMainBullets_L (left column)")
        print("  - textMainBullets_R (right column)")
        print("  - coSummaryShape (summary section)")
        print("=" * 80)

    except ImportError as e:
        print(f"❌ Import error: {e}")
        print("💡 Make sure python-pptx is installed")
    except Exception as e:
        print(f"❌ Error inspecting template: {e}")

if __name__ == "__main__":
    inspect_template()
