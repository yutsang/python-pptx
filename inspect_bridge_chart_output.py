#!/usr/bin/env python3
"""Inspects a GENERATED bridge/waterfall chart .pptx's embedded chart object
data directly -- shows exactly what ended up in the RENDERED chart
(categories + every series' actual values), as opposed to the raw source-
side block detection (bridge_chart_prototype.py's own printed item list,
already checked against a manual calculation and found to match within
rounding). Use this specifically to check whether the CHART RENDERING step
(_compute_waterfall_series / build_waterfall_chart, the invisible-base-
series stacked-bar technique) introduced any discrepancy on the way from
detected item -> chart series -> what actually displays on the slide.

Reconstructs one "effective value" per category (exactly one of the
Total/Increase/Decrease series is non-zero per category in this chart
technique; Base is the invisible offset, not a real displayed value) so
it's directly comparable to a manual calculation, the same way the raw
block dump already was.

Usage:
    python inspect_bridge_chart_output.py "bridge_chart_prototype_output.pptx"
    python inspect_bridge_chart_output.py "some_deck.pptx" --slide 3
"""
import argparse
import sys

from pptx import Presentation


def _chart_title(chart) -> str:
    try:
        if chart.has_title:
            return chart.chart_title.text_frame.text
    except Exception:
        pass
    return ""


def inspect_pptx_charts(path: str, only_slide: "int | None" = None) -> int:
    prs = Presentation(path)
    chart_count = 0
    for slide_idx, slide in enumerate(prs.slides, 1):
        if only_slide is not None and slide_idx != only_slide:
            continue
        for shape in slide.shapes:
            if not getattr(shape, "has_chart", False):
                continue
            chart_count += 1
            chart = shape.chart
            title = _chart_title(chart)
            print(f"\n{'=' * 78}")
            print(f"Slide {slide_idx}, chart: {title!r}")
            print(f"{'=' * 78}")

            plot = chart.plots[0]
            categories = [str(c) for c in plot.categories]
            series_values = {}
            for series in plot.series:
                series_values[series.name] = [float(v) if v is not None else None for v in series.values]

            print(f"Categories ({len(categories)}): {categories}")
            for name, values in series_values.items():
                print(f"  Series {name!r}: {values}")

            print("\n  Reconstructed per-category effective value (what the chart actually shows):")
            for i, cat in enumerate(categories):
                total_v = series_values.get("Total", [None] * len(categories))[i]
                inc_v = series_values.get("Increase", [None] * len(categories))[i]
                dec_v = series_values.get("Decrease", [None] * len(categories))[i]
                if total_v:
                    print(f"    {cat}: {total_v:,.2f}  [Total]")
                elif inc_v:
                    print(f"    {cat}: {inc_v:,.2f}  [Increase]")
                elif dec_v:
                    print(f"    {cat}: -{dec_v:,.2f}  [Decrease]")
                else:
                    print(f"    {cat}: 0.00  [all series zero at this category]")

    if chart_count == 0:
        print(f"⚠️  No chart object found in {path!r}"
              + (f" on slide {only_slide}" if only_slide is not None else "")
              + " -- either the wrong file/slide, or the chart wasn't a native python-pptx chart object.")
        return 1
    print(f"\n{chart_count} chart(s) inspected.")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("path", help="path to the generated bridge chart .pptx")
    ap.add_argument("--slide", type=int, default=None, help="only inspect this one slide number (1-indexed)")
    args = ap.parse_args()
    return inspect_pptx_charts(args.path, only_slide=args.slide)


if __name__ == "__main__":
    sys.exit(main())
