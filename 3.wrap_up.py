import argparse, time, os
from pathlib import Path
from pptx import Presentation
from datetime import datetime

def find_shape_by_name(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
    return None

def replace_text_preserve_formatting(shape, replacements):
    if not shape.has_text_frame:
        return
    
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)
                    
def update_project_titles(presentation_path, project_name, output_path=None):
    # Load the presentation
    prs = Presentation(presentation_path)
    total_slides = len(prs.slides)
    
    # Process each slide
    for slide_index, slide in enumerate(prs.slides):
        current_slide_number = slide_index + 1
        
        # find the projTitle shape
        projTitle_shape = find_shape_by_name(slide.shapes, "projTitle")
        
        if projTitle_shape:
            replacements = {
                "[PROJECT]": project_name,
                "[Current]": str(current_slide_number),
                "[Total]": str(total_slides)
            }
            
            # Repalce text while preserving formatting
            replace_text_preserve_formatting(projTitle_shape, replacements)
        
        else: 
            print(f"Warning: No 'projTitle' shape found on slide {current_slide_number}")
            
    # Generate the output filename with date and time
    datetime_string = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    input_path = Path(presentation_path)
    output_file = input_path.parent / f"{input_path.stem}_{project_name}_{datetime_string}{input_path.suffix}"
    
    # Save the presentation
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")
    
    # Optionally delete the original file
    try:
        os.remove(input_path)
    except Exception as e:
        print(f"Error deleting original file: {e}")
        
def main():
    parser = argparse.ArgumentParser(
        description="Update Project Title in PowerPoint Presentation",
        epilog="Example: python 3.wrap_up.py final_report.pptx 'Project_Name'"
    )
    parser.add_argument("input_file", help="Path to input PowerPoint file (.pptx)")
    parser.add_argument("project_name", help="Project name of the PowerPoint Presentation")
    parser.add_argument("-o", "--output", help="Output file path (optional)")
    
    args = parser.parse_args()
    
    # validate the input file 
    input_path = Path(args.input_file)
    if not input_path.exists():
        print(f"Error: Input file '{args.input_file}' not found")
        return 1
    try:
        update_project_titles(args.input_file, args.project_name, args.output)
        print("Project Presentation updated successfully!")
        return 0
    except Exception as e:
        print(f"Error processing presentation: {e}")
        return 1
    
if __name__ == "__main__":
    exit(main())
    
# python 3.wrap_up.py final_report.pptx 'Haining'