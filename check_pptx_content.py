#!/usr/bin/env python3
"""
Check what content is actually in the generated PowerPoint file
"""

from pptx import Presentation
import os

def check_pptx_content(pptx_path):
    """Check what content is actually in the PowerPoint file"""
    
    if not os.path.exists(pptx_path):
        print(f"❌ PowerPoint file not found: {pptx_path}")
        return
    
    print(f"🔍 Checking PowerPoint content: {pptx_path}")
    print("=" * 60)
    
    prs = Presentation(pptx_path)
    
    print(f"📊 PowerPoint has {len(prs.slides)} slides")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 Slide {slide_idx + 1}:")
        print(f"   Layout: {slide.slide_layout.name}")
        print(f"   Shapes: {len(slide.shapes)}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, 'name', 'No name')
            shape_type = type(shape).__name__
            
            print(f"   Shape {shape_idx + 1}: {shape_name} ({shape_type})")
            
            # If it's a text shape, show the content
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"      Content: {text[:200]}{'...' if len(text) > 200 else ''}")
                    print(f"      Content length: {len(text)} characters")
                    
                    # Show paragraph count
                    paragraphs = shape.text_frame.paragraphs
                    print(f"      Paragraphs: {len(paragraphs)}")
                    
                    # Show first few paragraphs
                    for i, para in enumerate(paragraphs[:3]):
                        para_text = para.text.strip()
                        if para_text:
                            print(f"         Para {i+1}: {para_text[:100]}{'...' if len(para_text) > 100 else ''}")
                else:
                    print(f"      Content: (empty)")
            
            # If it's a table, show table info
            if hasattr(shape, 'table'):
                table = shape.table
                print(f"      Table: {len(table.rows)} rows x {len(table.columns)} columns")
                
                # Show table content
                for row_idx in range(min(3, len(table.rows))):
                    row_data = []
                    for col_idx in range(min(3, len(table.columns))):
                        cell_text = table.cell(row_idx, col_idx).text.strip()
                        row_data.append(cell_text[:20])
                    print(f"         Row {row_idx+1}: {row_data}")

if __name__ == "__main__":
    pptx_path = "fdd_utils/output/test_BS_export.pptx"
    check_pptx_content(pptx_path)
