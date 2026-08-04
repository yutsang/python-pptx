#!/usr/bin/env python3
"""Read-only diagnostic for the "real spare capacity exceeds computed
capacity" gap reported on both table-lead-in and plain commentary boxes.

Dumps every internal number _calculate_max_lines_for_textbox's own
"exact by construction" identity (capacity_units * std_lh == box.height_pt)
depends on for named real boxes in a real exported .pptx, PLUS the real
paragraph count of each box's own content -- so a per-box back-solve can
be checked against the "does the residual gap scale with paragraph
count" hypothesis (each real paragraph carries one para_gap; a box with
more real paragraphs should show a bigger gap if para_gap is still
overstated, and little/no gap if it isn't).

Usage:
    # Dump every matched box (no back-solve):
    python diagnose_capacity_gap.py "Crescent_....pptx"

    # Back-solve specific boxes against your own real PowerPoint spare-
    # capacity measurements -- repeat --target as many times as needed,
    # one per box you tested:
    python diagnose_capacity_gap.py "Crescent_....pptx" \
        --target "1:textMainBullets:5" \
        --target "2:textMainBullets_L:2.15"
    # (slide number is 1-indexed, matching what inspect_pptx.py prints)

No PPTX is modified. Safe to run any time.
"""
import argparse
import sys

from pptx import Presentation
from pptx.util import Emu

from fdd_utils.text_metrics import get_measurer, text_box_from_shape

FONT_SIZE_PT = 9.0
LINE_SPACING = 1.0
PARA_GAP_PT = 3.0  # keep in lockstep with fdd_utils/pptx.py's _real_para_gap_pt


def _is_chinese(text: str) -> bool:
    return any('一' <= c <= '鿿' for c in text)


def _load_config():
    try:
        import yaml
        with open("fdd_utils/config.yml", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}
    except Exception:
        return {}


def _parse_target(spec: str):
    parts = spec.split(":")
    if len(parts) != 3:
        raise ValueError(f"--target must be SLIDE:SHAPE_NAME:EXTRA_LINES, got {spec!r}")
    slide_no, name, extra = parts
    return int(slide_no), name, float(extra)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path")
    ap.add_argument("--target", action="append", default=[],
                     help="SLIDE:SHAPE_NAME:EXTRA_LINES -- repeatable, one per box you "
                          "empirically measured spare capacity on.")
    args = ap.parse_args()
    targets = {}
    for spec in args.target:
        slide_no, name, extra = _parse_target(spec)
        targets[(slide_no, name)] = extra

    config = _load_config()
    packing_cfg = ((config.get("pptx") or {}).get("commentary_packing") or {})
    metrics_eng = packing_cfg.get("font_metrics_path_eng") or "fdd_utils/font_metrics/arial_eng.json"
    metrics_chi = packing_cfg.get("font_metrics_path_chi") or "fdd_utils/font_metrics/msyh_chi.json"
    family_eng = packing_cfg.get("font_family_eng") or "Arial"
    family_chi = packing_cfg.get("font_family_chi") or "Microsoft YaHei"

    eng_measurer = get_measurer(family_eng, FONT_SIZE_PT, is_cjk=False,
                                 line_spacing=LINE_SPACING, metrics_path=metrics_eng)
    chi_measurer = get_measurer(family_chi, FONT_SIZE_PT, is_cjk=True,
                                 line_spacing=LINE_SPACING, metrics_path=metrics_chi)
    print(f"Measurement source: ENG={eng_measurer.source}  CHI={chi_measurer.source}")
    print(f"  eng line_height_pt() = {eng_measurer.line_height_pt():.4f}")
    print(f"  chi line_height_pt() = {chi_measurer.line_height_pt():.4f}")
    print()

    prs = Presentation(args.pptx_path)
    checked = 0
    backsolved = []  # (slide_no, name, n_paras, implied_para_gap)

    for slide_idx, slide in enumerate(prs.slides):
        slide_no = slide_idx + 1
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            name = (getattr(shape, "name", "") or "")
            is_target = ("textmainbullets" in name.lower() or name.startswith("TextBox"))
            if not is_target:
                continue
            text = shape.text_frame.text or ""
            if len(text.strip()) < 15:
                continue  # skip stray short labels (e.g. a leftover "评述" band)

            checked += 1
            is_chi = _is_chinese(text)
            measurer = chi_measurer if is_chi else eng_measurer

            raw_height_pt = Emu(shape.height).pt
            box = text_box_from_shape(shape)
            tIns_bIns_pt = raw_height_pt - box.height_pt

            line_h = measurer.line_height_pt()
            std_lh = line_h + PARA_GAP_PT
            capacity_units = max(1.0, box.height_pt / std_lh) if std_lh > 0 else 0.0

            # Real paragraph count -- one per XML <a:p>, which is exactly
            # what _calculate_content_lines charges one para_gap against
            # (category header, each "■ key - ..." bullet, each
            # continuation line). If the residual gap scales with THIS
            # number rather than being a flat per-box amount, that's
            # direct evidence para_gap itself (not something else) is
            # still overstated, and by how much.
            n_paras = sum(1 for p in shape.text_frame.paragraphs if (p.text or "").strip())

            print(f"Slide {slide_no} | {name!r} | {'CHI' if is_chi else 'ENG'} | "
                  f"{len(text)} chars | {n_paras} real paragraphs")
            print(f"  raw shape.height       = {raw_height_pt:.3f}pt ({raw_height_pt/72:.4f}in)")
            print(f"  tIns+bIns (margins)    = {tIns_bIns_pt:.3f}pt")
            print(f"  box.height_pt (usable) = {box.height_pt:.3f}pt")
            print(f"  measurer.line_height_pt() = {line_h:.4f}pt   (source={measurer.source})")
            print(f"  para_gap_pt (current)  = {PARA_GAP_PT:.3f}pt")
            print(f"  std_lh (line_h+gap)    = {std_lh:.4f}pt")
            print(f"  capacity_units         = {capacity_units:.4f}L")
            print(f"  IDENTITY CHECK: capacity_units * std_lh = {capacity_units*std_lh:.3f}pt "
                  f"(should equal box.height_pt {box.height_pt:.3f}pt)")

            extra = targets.get((slide_no, name))
            if extra is not None:
                real_capacity_units = capacity_units + extra
                implied_std_lh = box.height_pt / real_capacity_units
                implied_line_h_if_paragap_fixed = implied_std_lh - PARA_GAP_PT
                # Alternative attribution: hold line_h fixed (it's the
                # separately-proven-correct 1.2x factor) and solve for
                # what para_gap WOULD need to be, per real paragraph, to
                # absorb the WHOLE gap alone.
                implied_para_gap = implied_std_lh - line_h
                print()
                print(f"  >>> Back-solved from your reported {extra} extra real lines:")
                print(f"      implied REAL std_lh   = {implied_std_lh:.4f}pt "
                      f"(currently: {std_lh:.4f}pt, delta={std_lh-implied_std_lh:+.4f}pt)")
                print(f"      implied REAL para_gap = {implied_para_gap:.4f}pt "
                      f"(currently: {PARA_GAP_PT:.3f}pt, delta={PARA_GAP_PT-implied_para_gap:+.4f}pt)"
                      f"  [{n_paras} real paragraphs in this box]")
                backsolved.append((slide_no, name, n_paras, implied_para_gap))
            print()

            if checked >= 12:
                print("(stopping after 12 shapes -- enough for comparison)")
                print()
                break
        if checked >= 12:
            break

    if len(backsolved) >= 2:
        print("=" * 70)
        print("CROSS-CHECK: does the implied para_gap correlate with paragraph count?")
        for slide_no, name, n_paras, implied_para_gap in backsolved:
            print(f"  slide {slide_no} {name!r}: {n_paras} paragraphs -> implied para_gap {implied_para_gap:.4f}pt")
        print("If these numbers are CLOSE regardless of paragraph count, para_gap is the")
        print("right lever and this is the value to use. If they're spread out (e.g. one")
        print("strongly negative), a flat para_gap correction can't explain the full gap")
        print("on its own -- something else (possibly per-paragraph, possibly not) is")
        print("also contributing, and needs separate investigation before changing anything.")

    if checked == 0:
        print("No textMainBullets*/TextBox shapes with text found.")
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
