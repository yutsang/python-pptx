#!/usr/bin/env python3
"""Reads the tables that are actually IN a reference deck and traces each one
back to the databook block it came from.

This is the ground truth for "which table belongs in the report". Deriving it
from the databook alone means inventing rules about which of a sheet's blocks
looks report-worthy -- and a real scan showed that second block is variously a
bank-account listing, a fee-rate workpaper, a rollforward or a set of operating
KPIs, so such rules can only ever approximate. The finished deck already
contains the decision; this reads it back and reports the observed rule.

Handles the three ways a table can reach a deck, since an UpSlide-produced one
often is not a native PowerPoint table:
  1. native table shape -- read directly;
  2. embedded/linked OLE object (an Excel range pasted in) -- reported, with
     inspect_pptx_tables.py named as the tool that extracts the blob;
  3. picture / flattened image -- reported as not structurally readable, so a
     table that exists only as an image is never silently counted as absent.

Read-only.

Usage:
    python trace_report_tables.py "for_test/Crescent-report.pptx"
    python trace_report_tables.py "for_test/Crescent-report.pptx" --databook "for_test/Crescent-databook.xlsx"
"""
import argparse
import re
import sys
import warnings

warnings.filterwarnings("ignore")
sys.path.insert(0, ".")

import pandas as pd
from pptx import Presentation

_SKIP_SHEET_HINTS = ("-->", ">>", "_tm_", "upslide", "封面", "目录", "目錄")
_NUMERIC_RE = re.compile(r"^[-(]?[\d,]+(\.\d+)?\)?%?$")


def _walk(shapes, slide_no):
    """(slide, shape, kind) for every shape, descending into groups. kind is
    'table', 'ole', 'picture' or 'other'."""
    for shape in shapes:
        if shape.shape_type == 6 and hasattr(shape, "shapes"):
            yield from _walk(shape.shapes, slide_no)
            continue
        if getattr(shape, "has_table", False):
            kind = "table"
        elif shape.shape_type == 7:  # MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT
            kind = "ole"
        elif shape.shape_type == 13:  # PICTURE
            kind = "picture"
        else:
            kind = "other"
        yield slide_no, shape, kind


def _cell(t, r, c):
    try:
        return t.cell(r, c).text.strip()
    except Exception:
        return ""


_XLSX_MAGIC = b"PK\x03\x04"


def _read_ole_table(shape, slide_no, tmp_dir):
    """An embedded Excel range as {labels, rows} so it can be traced like a
    native table.

    A real reference deck turned out to carry every one of its tables this
    way -- 7 OLE embeds, 0 native tables -- because the analyst pastes an
    Excel range through UpSlide. Reporting them as merely 'present' would
    have left the actual question unanswered."""
    import os
    try:
        blob = shape.ole_format.blob
    except Exception:
        return None  # linked rather than embedded: the source file is external
    if not blob or blob[:4] != _XLSX_MAGIC:
        return None  # legacy .xls or not a workbook
    path = os.path.join(tmp_dir, f"ole_slide{slide_no}_{shape.shape_id}.xlsx")
    try:
        with open(path, "wb") as f:
            f.write(blob)
        frames = pd.read_excel(path, sheet_name=None, header=None)
    except Exception:
        return None
    best = None
    for name, df in frames.items():
        if df is None or df.empty:
            continue
        labels, rows = [], []
        for r in range(len(df)):
            first = str(df.iloc[r, 0]).strip()
            if first in ("", "nan"):
                first = ""
            values = [str(v).strip() for v in df.iloc[r].tolist()[1:6]
                      if str(v).strip() not in ("", "nan")]
            if first and not _NUMERIC_RE.match(first):
                labels.append(first)
            if first or values:
                rows.append({"label": first, "values": values})
        if labels and (best is None or len(labels) > len(best["labels"])):
            best = {"slide": slide_no, "shape": f"{shape.name} [OLE:{name}]",
                    "header": [str(v).strip() for v in df.iloc[0].tolist()[:6]],
                    "labels": labels, "rows": rows,
                    "n_rows": len(df), "n_cols": len(df.columns),
                    "source_file": path}
    return best


def extract_deck_tables(pptx_path, tmp_dir=None):
    import tempfile
    tmp_dir = tmp_dir or tempfile.mkdtemp(prefix="ole_tables_")
    prs = Presentation(pptx_path)
    tables, ole, pictures = [], [], []
    for idx, slide in enumerate(prs.slides, 1):
        for slide_no, shape, kind in _walk(slide.shapes, idx):
            if kind == "ole":
                ole.append((slide_no, shape.name))
                parsed = _read_ole_table(shape, slide_no, tmp_dir)
                if parsed:
                    tables.append(parsed)
                continue
            if kind == "picture":
                pictures.append((slide_no, shape.name))
                continue
            if kind != "table":
                continue
            table = shape.table
            n_rows, n_cols = len(table.rows), len(table.columns)
            if n_rows < 2 or n_cols < 2:
                continue
            header = [_cell(table, 0, c) for c in range(n_cols)]
            labels, rows = [], []
            for r in range(1, n_rows):
                first = _cell(table, r, 0)
                values = [_cell(table, r, c) for c in range(1, n_cols)]
                if first and not _NUMERIC_RE.match(first):
                    labels.append(first)
                rows.append({"label": first, "values": values})
            tables.append({"slide": slide_no, "shape": shape.name, "header": header,
                           "labels": labels, "rows": rows,
                           "n_rows": n_rows, "n_cols": n_cols})
    return tables, ole, pictures


def find_source_block(xl_path, labels, min_hits=2, _cache={}):
    """Databook sheets containing these row labels, best match first."""
    wanted = [l for l in labels if l and not _NUMERIC_RE.match(l)]
    if not wanted:
        return []
    if xl_path not in _cache:
        xl = pd.ExcelFile(xl_path)
        frames = {}
        for sheet in xl.sheet_names:
            if any(h in sheet.lower() for h in _SKIP_SHEET_HINTS):
                continue
            try:
                frames[sheet] = pd.read_excel(xl_path, sheet_name=sheet, header=None)
            except Exception:
                continue
        _cache[xl_path] = frames
    results = []
    for sheet, df in _cache[xl_path].items():
        positions = {}
        for r in range(len(df)):
            for c in range(min(4, len(df.columns))):
                text = str(df.iloc[r, c]).strip()
                if text in wanted and text not in positions:
                    positions[text] = (r, c)
        if len(positions) >= min_hits:
            results.append((sheet, min(p[0] for p in positions.values()),
                            len(positions), len(wanted), positions))
    results.sort(key=lambda t: (-t[2], t[1]))
    return results


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx", help="the reference deck")
    ap.add_argument("--databook", default=None, help="trace each table back to its source sheet")
    ap.add_argument("--max-rows", type=int, default=14, help="rows to print per table")
    args = ap.parse_args()

    tables, ole, pictures = extract_deck_tables(args.pptx)
    n_from_ole = sum(1 for t in tables if "[OLE:" in t["shape"])
    print(f"Deck: {args.pptx!r}")
    print(f"  readable tables: {len(tables)}"
          + (f"   ({n_from_ole} of them extracted from OLE embeds)" if n_from_ole else ""))
    print(f"  OLE embeds    : {len(ole)}"
          + (f"   ({len(ole) - n_from_ole} not readable -- linked, or legacy .xls)"
             if len(ole) > n_from_ole else "   (all read)"))
    print(f"  pictures      : {len(pictures)}"
          + ("   <-- a table flattened to an image cannot be read structurally" if pictures else ""))
    print()
    if not tables:
        print("No native tables. If the deck visibly HAS tables they are OLE embeds or")
        print("images -- see the counts above; this tool can only trace native ones.")
        if ole:
            for s, n in ole[:10]:
                print(f"    OLE on slide {s}: {n!r}")
        if pictures:
            for s, n in pictures[:10]:
                print(f"    picture on slide {s}: {n!r}")
        return 1

    for i, t in enumerate(tables, 1):
        print("=" * 78)
        print(f"TABLE {i}  slide {t['slide']}  shape {t['shape']!r}  {t['n_rows']}x{t['n_cols']}")
        print("=" * 78)
        print(f"  header: {t['header']}")
        for row in t["rows"][:args.max_rows]:
            print(f"    {row['label'][:26]:28s} {row['values'][:5]}")
        if len(t["rows"]) > args.max_rows:
            print(f"    ... and {len(t['rows']) - args.max_rows} more row(s)")
        if args.databook:
            matches = find_source_block(args.databook, t["labels"])
            print(f"\n  SOURCE in the databook:")
            if not matches:
                print(f"    ⚠️ no sheet carries these row labels -- typed straight into the")
                print(f"       deck, or relabelled on the way in")
            for sheet, first_row, hits, total, positions in matches[:3]:
                exact = "  EXACT" if hits == total else ""
                print(f"    {sheet!r}: {hits}/{total} labels matched, first at row {first_row}{exact}")
                sample = list(positions.items())[:4]
                print(f"      {', '.join(f'{k}@r{v[0]}' for k, v in sample)}")
        print()

    if args.databook:
        print("=" * 78)
        print("OBSERVED RULE: which accounts got a table")
        print("=" * 78)
        traced = {}
        for t in tables:
            for sheet, _fr, hits, total, _pos in find_source_block(args.databook, t["labels"])[:1]:
                traced[sheet] = (hits, total, t["slide"], len(t["rows"]), t["n_cols"])
        for sheet, (hits, total, slide, n_rows, n_cols) in sorted(traced.items()):
            print(f"  ✅ {sheet:22s} slide {slide:>2}  {n_rows} rows x {n_cols} cols  "
                  f"({hits}/{total} labels traced)")
        print(f"\n  {len(traced)} databook sheet(s) supplied a table to this deck.")
        print(f"  Implement THIS list, not a guess about which block looks report-worthy.")
        print(f"  Cross-check it against extract_presentation_detail_table's own picks --")
        print(f"  where the two disagree, the deck is right and the heuristic needs work.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
