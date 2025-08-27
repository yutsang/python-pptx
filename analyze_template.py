#!/usr/bin/env python3
"""
Analyze PowerPoint template to see available shapes
"""

from pptx import Presentation
import os

def analyze_template(template_path):
    """Analyze the PowerPoint template and list all shapes"""
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return
    
    print(f"🔍 Analyzing template: {template_path}")
    print("=" * 60)
    
    prs = Presentation(template_path)
    
    print(f"📊 Template has {len(prs.slides)} slides")
    print(f"📋 Template has {len(prs.slide_layouts)} slide layouts")
    
    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n📄 Slide {slide_idx + 1}:")
        print(f"   Layout: {slide.slide_layout.name}")
        print(f"   Shapes: {len(slide.shapes)}")
        
        for shape_idx, shape in enumerate(slide.shapes):
            shape_name = getattr(shape, 'name', 'No name')
            shape_type = type(shape).__name__
            
            print(f"   Shape {shape_idx + 1}: {shape_name} ({shape_type})")
            
            # If it's a text shape, show some content
            if hasattr(shape, 'text_frame') and shape.text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    print(f"      Text: {text[:50]}{'...' if len(text) > 50 else ''}")
            
            # If it's a table, show table info
            if hasattr(shape, 'table'):
                table = shape.table
                print(f"      Table: {len(table.rows)} rows x {len(table.columns)} columns")

if __name__ == "__main__":
    template_path = "fdd_utils/template.pptx"
    analyze_template(template_path) 