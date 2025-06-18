from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from dataclasses import dataclass
from typing import List, Tuple
import textwrap
import logging
from pptx.oxml.ns import qn  # Required import for XML namespace handling
from pptx.oxml.xmlchemy import OxmlElement
import re
import argparse

logging.basicConfig(level=logging.INFO)

@dataclass
class FinancialItem:
    accounting_type: str
    account_title: str
    descriptions: List[str]
    layer1_continued: bool = False
    layer2_continued: bool = False
    is_table: bool = False  # Added missing is_table parameter
    
class ReportGenerator: 
    def __init__(self, template_path, markdown_file, output_path):
        self.template_path = template_path
        self.markdown_file = markdown_file
        self.output_path = output_path
        
    def extract_location_from_path(self, file_path):
        match = re.match(r"utils/template_(.+)\.pptx", file_path)
        if match:
            return match.group(1)
        else: 
            raise ValueError("The file path does not match the expected format.")
    
    def generate(self):
        path = self.extract_location_from_path(path)
        location = self.extract_location_from_path(path)
        if location == "Haining":
            summary_content = """
            ## Summary
            The company demonstrates strong financial health with total assets of $180 million, liabilities of $75 million, and shareholder's equity of $105 million. Current assets including $45 million cash and $30 million receivables provide ample liquidity to cover short-term obligations of $50 milliion. Long-term invesmtnets in property and equipment total $90 million, supported by conservative debt levels with a debt-to-equity ratio of 0.71. Retained earnings of $80 million reflect consistent profitability and prudent dividend policies. The balance sheet structure shows optimal asset allocation with 60% long-term investments and 40% working capital. Financial ratios indicate robust solvency with current ratio of 2.4 and quick ratio of 1.8. Equity growth of 12% year-over-year demonstrates sustainable value creation. Conservative accounting practives ensure asset valutations remain realistic, while liability management maintains healthy interest coverage. Overall, the balance sheet positions the company for strategic investments while maintaining financial stability.
            """
        elif location == "Ningbo":
            summary_content = """
            ## Summary
            The company demonstrates strong financial health with total assets of $180 million, liabilities of $75 million, and shareholder's equity of $105 million. Current assets including $45 million cash and $30 million receivables provide ample liquidity to cover short-term obligations of $50 milliion. Long-term invesmtnets in property and equipment total $90 million, supported by conservative debt levels with a debt-to-equity ratio of 0.71. Retained earnings of $80 million reflect consistent profitability and prudent dividend policies. The balance sheet structure shows optimal asset allocation with 60% long-term investments and 40% working capital. Financial ratios indicate robust solvency with current ratio of 2.4 and quick ratio of 1.8. Equity growth of 12% year-over-year demonstrates sustainable value creation. Conservative accounting practives ensure asset valutations remain realistic, while liability management maintains healthy interest coverage. Overall, the balance sheet positions the company for strategic investments while maintaining financial stability.
            """
        elif location == "Nanjing":
            summary_content = """
            ## Summary
            The company demonstrates strong financial health with total assets of $180 million, liabilities of $75 million, and shareholder's equity of $105 million. Current assets including $45 million cash and $30 million receivables provide ample liquidity to cover short-term obligations of $50 milliion. Long-term invesmtnets in property and equipment total $90 million, supported by conservative debt levels with a debt-to-equity ratio of 0.71. Retained earnings of $80 million reflect consistent profitability and prudent dividend policies. The balance sheet structure shows optimal asset allocation with 60% long-term investments and 40% working capital. Financial ratios indicate robust solvency with current ratio of 2.4 and quick ratio of 1.8. Equity growth of 12% year-over-year demonstrates sustainable value creation. Conservative accounting practives ensure asset valutations remain realistic, while liability management maintains healthy interest coverage. Overall, the balance sheet positions the company for strategic investments while maintaining financial stability.
            """
        generator = PowerPointGenerator(self.template_path)
        try: 
            generator.generate_full_report(self.markdown_file, summary_content, self.output_path)
        except Exception as e:
            print(f"Generation failed: {str(e)}")

class PowerPointGenerator:
    def __init__(self, template_path: str):
        self.prs = Presentation(template_path)
        #self.log_template_shapes()
        #self._validate_template()
        self.current_slide_index = 0
        self.LINE_HEIGHT = Pt(12)
        self.ROWS_PER_SECTION = self._calculate_max_rows()
        self.CHARS_PER_ROW = 70
        self.BULLET_CHAR = '■ '
        self.DARK_BLUE = RGBColor(0, 50, 150)
        self.DARK_GREY = RGBColor(169, 169, 169)
        self.prev_layer1 = None
        self.prev_layer2 = None
        
    def _load_markdown(self, file_path: str) -> str:
        """Load markdown content from file."""
        try: 
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            return content
        except FileNotFoundError:
            print(f"The file at {file_path} was not found.")
            return ""
        except Exception as e:
            print(f"An error occured: {e}")
            return ""
        
    def _calculate_max_rows(self):
        slide = self.prs.slides[0]
        shape = next(s for s in slide.shapes if s.name == "textMainBullets")
        return int(shape.height / self.LINE_HEIGHT) - 3
    
    def _apply_paragraph_formatting(self, paragraph, is_layer2_3=False):
        """Version-safe paragraph formatting with layer-specific settings"""
        try:
            # Modern versions (>=0.6.18) with paragraph_format
            pf = paragraph.paragraph_format
            if is_layer2_3:
                pf.left_indent = Inches(0.21)
                pf.first_line_indent = Inches(-0.19)  # Hanging indent
                pf.space_before = Pt(3)  # Layer 2/3 spacing
            else:
                pf.left_indent = Inches(0.3)
                pf.first_line_indent = Inches(-0.3)
                pf.space_before = Pt(6)  # Layer 1 spacing
            
            pf.space_after = Pt(0)
            pf.line_spacing = 1.0
            pf.alignment = PP_ALIGN.LEFT
            
        except AttributeError:
            # Legacy version handling
            if is_layer2_3:
                try: paragraph.left_indent = Inches(0.21)
                except: pass
                try: paragraph.first_line_indent = Inches(-0.19)
                except: pass
                try: paragraph.space_before = Pt(3)
                except: pass
            else:
                try: paragraph.left_indent = Inches(0.3)
                except: pass
                try: paragraph.first_line_indent = Inches(-0.3)
                except: pass
                try: paragraph.space_before = Pt(6)
                except: pass
            
            try: paragraph.space_after = Pt(0)
            except: pass
            try: paragraph.line_spacing = 1.0
            except: pass
            try: paragraph.alignment = PP_ALIGN.LEFT
            except: self._handle_legacy_alignment(paragraph)

    def _handle_legacy_alignment(self, paragraph):
        """XML-based left alignment for legacy versions"""
        try:
            pPr = paragraph._element.get_or_add_pPr()
            for child in list(pPr):
                if child.tag.endswith('jc'):
                    pPr.remove(child)
            align = OxmlElement('a:jc')
            align.set('val', 'left')
            pPr.append(align)
        except Exception as e:
            logging.warning(f"Legacy alignment failed: {str(e)}")


    def _validate_template(self):
        """Main version validation logic"""
        required_shapes = {
            0: ["textMainBullets"],
            1: ["textMainBullets_L", "textMainBullets_R"]
        }
        for slide_idx, slide in enumerate(self.prs.slides):
            if slide_idx in required_shapes:
                missing = [name for name in required_shapes[slide_idx]
                          if not any(s.name == name for s in slide.shapes)]
                if missing:
                    raise ValueError(f"Missing shapes on slide {slide_idx+1}: {', '.join(missing)}")
                
    def log_template_shapes(self):
        """Log all shapes in the template for debugging"""
        logging.info("=== Template Shape Audit ===")
        for slide_idx, slide in enumerate(self.prs.slides):
            logging.info(f"Slide {slide_idx + 1} has {len(slide.shapes)} shapes:")
            for shape in slide.shapes:
                logging.info(f"  - Name: '{shape.name}' | Type: {shape.shape_type}")
        logging.info("=== End Shape Audit ===")


    def _validate_content_placement(self, distribution):
        """Anti-duplication safeguard"""
        seen = set()
        for slide_idx, section, items in distribution:
            for item in items:
                key = (item.accounting_type, item.account_title, tuple(item.descriptions))
                if key in seen:
                    raise ValueError(f"Duplicate content detected: {key}")
                seen.add(key)

    def parse_markdown(self, md_content: str) -> List[FinancialItem]:
        items = []
        current_type = ""
        current_title = ""
        current_descs = []
        is_table = False

        for line in md_content.strip().split('\n'):
            stripped = line.strip()
            
            if stripped.startswith('## '):
                if current_descs:
                    items.append(FinancialItem(
                        current_type, current_title, current_descs, 
                        is_table=is_table
                    ))
                current_type = stripped[3:]
                current_title = ""
                current_descs = []
                is_table = False
            elif stripped.startswith('### '):
                if current_descs:
                    items.append(FinancialItem(
                        current_type, current_title, current_descs,
                        is_table=is_table
                    ))
                current_title = stripped[4:]
                current_descs = []
                is_table = "taxes payables" in current_title.lower()
            elif re.match(r'^={20,}', stripped):
                is_table = True
                current_descs.append("="*40)
            else:
                if stripped:
                    current_descs.append(stripped)

        if current_descs:
            items.append(FinancialItem(current_type, current_title, current_descs, is_table=is_table))
        return items

    def _plan_content_distribution(self, items: List[FinancialItem]):
        distribution = []
        content_queue = items.copy()
        slide_idx = 0

        while content_queue:
            sections = ['c'] if slide_idx == 0 else ['b', 'c']
            
            for section in sections:
                if not content_queue:
                    break
                
                section_items = []
                lines_used = 0
                
                while content_queue and lines_used < self.ROWS_PER_SECTION:
                    item = content_queue[0]
                    item_lines = self._calculate_item_lines(item)
                    
                    if lines_used + item_lines <= self.ROWS_PER_SECTION:
                        section_items.append(item)
                        content_queue.pop(0)
                        lines_used += item_lines
                    else:
                        # Handle continuation logic
                        remaining_lines = self.ROWS_PER_SECTION - lines_used
                        if remaining_lines > 1:
                            split_item, remaining_item = self._split_item(item, remaining_lines)
                            section_items.append(split_item)
                            
                            # Only set remaining item if there's actual remaining content
                            if remaining_item and remaining_item.descriptions and remaining_item.descriptions[0]:
                                content_queue[0] = remaining_item
                            else:
                                content_queue.pop(0)  # Remove if no remaining content
                
                if section_items:
                    distribution.append((slide_idx, section, section_items))
            
            slide_idx += 1
            if slide_idx >= len(self.prs.slides):
                break
                
        return distribution
    
    def _split_item(self, item: FinancialItem, max_lines: int) -> tuple[FinancialItem, FinancialItem]:
        available_chars = max_lines * self.CHARS_PER_ROW
        header = f"{self.BULLET_CHAR}{item.account_title} - "
        available_for_desc = available_chars - len(header)
        
        desc_text = ' '.join(item.descriptions)
        
        # Find split point at word boundary
        if len(desc_text) <= available_for_desc:
            split_item = FinancialItem(
                item.accounting_type,
                item.account_title,
                item.descriptions,
                layer1_continued=item.layer1_continued,
                layer2_continued=item.layer2_continued
            )
            remaining_item = None
        else:
            split_pos = desc_text.rfind(' ', 0, available_for_desc)
            if split_pos == -1:
                split_pos = available_for_desc
            
            first_part = desc_text[:split_pos].strip()
            remaining_part = desc_text[split_pos:].strip()
            
            split_item = FinancialItem(
                item.accounting_type,
                item.account_title,
                [first_part] if first_part else [],
                layer1_continued=item.layer1_continued,
                layer2_continued=item.layer2_continued
            )
            
            remaining_item = FinancialItem(
                item.accounting_type,
                item.account_title,
                [remaining_part] if remaining_part else [],
                layer1_continued=True,
                layer2_continued=True
            ) if remaining_part else None
        
        return split_item, remaining_item

        
        
    def _split_text_at_boundary(self, text: str, max_chars: int) -> str:
        """Split text at word boundary within character limit"""
        if len(text) <= max_chars:
            return text
        
        # Find last space within limit
        split_pos = text.rfind(' ', 0, max_chars)
        if split_pos == -1:  # No space found, split at character limit
            split_pos = max_chars
        
        return text[:split_pos]

    def _wrap_text(self, text: str) -> List[str]:
        return textwrap.wrap(text, width=self.CHARS_PER_ROW, break_long_words=True)

    def _calculate_wrapped_lines(self, text: str) -> int:
        """Calculate actual wrapped lines using textwrap"""
        wrapper = textwrap.TextWrapper(
            width=self.CHARS_PER_ROW,
            break_long_words=True,
            replace_whitespace=False
        )
        return len(wrapper.wrap(text))
    
    def _calculate_effective_width_for_description(self, account_title: str) -> int:
        """Calculate available character width for description after accounting for bullet and title"""
        bullet_overhead = len(self.BULLET_CHAR)  # "■ " = 2 characters
        separator_overhead = len(" - ")  # " - " = 3 characters
        title_overhead = len(account_title)
        
        total_overhead = bullet_overhead + title_overhead + separator_overhead
        effective_width = self.CHARS_PER_ROW - total_overhead
        
        # Ensure minimum width for description
        return max(effective_width, 10)  # At least 10 chars for description

    def _calculate_item_lines(self, item: FinancialItem) -> int:
        lines = 0
        
        # Layer 1 lines
        lines += len(textwrap.wrap(f"{item.accounting_type} (continued)" if item.layer1_continued 
                                  else item.accounting_type, width=self.CHARS_PER_ROW))
        
        # Combined Layer 2+3 lines
        combined_text = f"{self.BULLET_CHAR}{item.account_title} - {' '.join(item.descriptions)}"
        if item.is_table:
            combined_text = self._format_table_text(combined_text)
            
        lines += len(textwrap.wrap(combined_text, width=self.CHARS_PER_ROW))
        
        return lines

    def _get_section_shape(self, slide, section: str):
        """Direct shape access without layout assumptions"""
        if self.current_slide_index == 0:
            if section == 'c':
                return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
            return None  # No 'b' section on first slide
        
        # For subsequent slides
        target_name = "textMainBullets_L" if section == 'b' else "textMainBullets_R"
        return next((s for s in slide.shapes if s.name == target_name), None)

    def _populate_section(self, shape, items: List[FinancialItem]):
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        self.prev_layer1 = None  # Reset for each new section

        for item in items:
            # Layer 1 Header
            if (item.accounting_type != self.prev_layer1) or (self.prev_layer1 is None):
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=False)
                run = p.add_run()
                cont_text = " (continued)" if item.layer1_continued else ""
                run.text = f"{item.accounting_type}{cont_text}"
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.name = 'Arial'
                run.font.color.rgb = self.DARK_BLUE
                self.prev_layer1 = item.accounting_type



            # Handle taxes payables table
            if 'taxes payables' in item.account_title.lower():
                self._create_taxes_table(shape, item)
            else:
                desc_text = " ".join(item.descriptions)
                bullet_lines = self._detect_bullet_content(desc_text)
                
                if bullet_lines:
                    self._populate_bullet_content(shape, item, bullet_lines)
                else:
                    p = tf.add_paragraph()
                    self._apply_paragraph_formatting(p, is_layer2_3=True)
                    
                    # Bullet symbol
                    bullet_run = p.add_run()
                    bullet_run.text = self.BULLET_CHAR
                    bullet_run.font.color.rgb = self.DARK_GREY
                    bullet_run.font.name = 'Arial'
                    bullet_run.font.size = Pt(9)
                    
                    # Layer 2 (bold)
                    title_run = p.add_run()
                    cont_text = " (continued)" if item.layer2_continued else ""
                    title_run.text = f"{item.account_title}{cont_text}"
                    title_run.font.bold = True
                    title_run.font.name = 'Arial'
                    title_run.font.size = Pt(9)
                    
                    # Separator and Layer 3 (not bold)
                    desc_run = p.add_run()
                    desc_run.text = f" - {desc_text}"
                    desc_run.font.bold = False
                    desc_run.font.name = 'Arial'
                    desc_run.font.size = Pt(9)

    def _detect_bullet_content(self, text):
        """
        Enhanced bullet detection that handles multi-level bullet formats up to 4 levels deep
        with special handling for layer 3 bullet points that might contain dashes indicating layer 4
        """
        lines = text.split('\n')
        bullet_lines = []
        has_bullets = False

        # Check if we have any lines starting with "- "
        for line in lines:
            if line.strip().startswith('- '):
                has_bullets = True
                break
        
        # If we have bullets, process the content
        if has_bullets:
            main_content = []
            
            for line in lines:
                stripped = line.strip()
                original_line = line
                
                # Detect bullet lines with '- ' prefix
                if original_line.lstrip().startswith('- '):
                    # If we have accumulated main content, add it first
                    if main_content:
                        main_text = ' '.join(main_content).strip()
                        if main_text:
                            bullet_lines.append((0, main_text))
                        main_content = []
                    
                    # Calculate indentation level (based on spaces/tabs before the bullet)
                    indent_spaces = len(original_line) - len(original_line.lstrip())
                    
                    # Determine bullet level based on indentation
                    # We divide by 2 assuming 2 spaces per indentation level
                    level = min(4, (indent_spaces // 2) + 1)  # Cap at 4 levels
                    
                    # Clean and store bullet line
                    clean_line = stripped[2:]  # Remove '- '
                    
                    # Special handling for level 3 bullets that contain a dash indicating level 4
                    if level == 3 and " - " in clean_line:
                        # Split at the first occurrence of " - "
                        parts = clean_line.split(" - ", 1)
                        if len(parts) > 1:
                            # Add level 3 content
                            bullet_lines.append((level, parts[0].strip()))
                            # Add level 4 content
                            bullet_lines.append((4, parts[1].strip()))
                        else:
                            bullet_lines.append((level, clean_line))
                    else:
                        bullet_lines.append((level, clean_line))
                
                elif stripped:
                    # This is regular content - could be main content or continuation
                    main_content.append(stripped)
            
            # Add any remaining main content
            if main_content:
                main_text = ' '.join(main_content).strip()
                if main_text:
                    bullet_lines.insert(0, (0, main_text))
        
        return bullet_lines if bullet_lines else None


    def _create_taxes_table(self, shape, item):
        """Formats taxes payables content as a structured table"""
        tf = shape.text_frame
        
        # Header
        p = tf.add_paragraph()
        self._apply_paragraph_formatting(p, is_layer2_3=True)
        header_run = p.add_run()
        cont_text = " (continued)" if item.layer2_continued else ""
        header_run.text = f"{self.BULLET_CHAR}{item.account_title}{cont_text}"
        header_run.font.bold = True
        header_run.font.name = 'Arial'
        header_run.font.size = Pt(9)

        # Process table content
        content = ' '.join(item.descriptions)
        current_category = None
        
        for line in content.split('\n'):
            line = line.strip()
            if not line or '===' in line:
                continue
                
            if 'Tax:' in line or 'Payable:' in line:
                current_category = line.replace(':', '')
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=True)
                p.paragraph_format.left_indent = Inches(0.4)
                run = p.add_run()
                run.text = f"• {current_category}"
                run.font.name = 'Arial'
                run.font.bold = True
                run.font.size = Pt(9)
            elif line.startswith('- '):
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=True)
                p.paragraph_format.left_indent = Inches(0.6)
                run = p.add_run()
                run.text = f"◦ {line[2:]}"
                run.font.name = 'Arial'
                run.font.size = Pt(9)

    def _populate_bullet_content(self, shape, item, bullet_lines):
        """Formats detected bullet points with proper indentation"""
        tf = shape.text_frame
        
        p = tf.add_paragraph()
        self._apply_paragraph_formatting(p, is_layer2_3=True)
        
        # Layer 2 header
        bullet_run = p.add_run()
        bullet_run.text = self.BULLET_CHAR
        bullet_run.font.color.rgb = self.DARK_GREY
        bullet_run.font.name = 'Arial'
        bullet_run.font.size = Pt(9)
        
        title_run = p.add_run()
        cont_text = " (continued)" if item.layer2_continued else ""
        title_run.text = f"{item.account_title}{cont_text}"
        title_run.font.bold = True
        title_run.font.name = 'Arial'
        title_run.font.size = Pt(9)
        
        # Bullet items
        for line in bullet_lines:
            p = tf.add_paragraph()
            self._apply_paragraph_formatting(p, is_layer2_3=True)
            p.paragraph_format.left_indent = Inches(0.4)
            p.paragraph_format.first_line_indent = Inches(-0.19)
            
            run = p.add_run()
            clean_line = line.lstrip('- ').strip()
            run.text = f"• {clean_line}"
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.bold = False

    def _format_table_text(self, text: str) -> str:
        lines = []
        current_line = []
        for part in text.split():
            if part == self.BULLET_CHAR.strip():
                if current_line:
                    lines.append(" ".join(current_line))
                current_line = ["  " + part]
            else:
                current_line.append(part)
        lines.append(" ".join(current_line))
        return "\n".join(lines)
                
    def _handle_paragraph_spacing(self, paragraph, space_before=None, space_after=None):
        """Universal spacing handler"""
        try:
            # Modern versions (>=0.6.18)
            pf = paragraph.paragraph_format
            if space_before is not None:
                pf.space_before = space_before
            if space_after is not None:
                pf.space_after = space_after
        except AttributeError:
            # Legacy version fallback
            if space_before is not None:
                paragraph.space_before = space_before
            if space_after is not None:
                paragraph.space_after = space_after


    def _handle_paragraph_indent(self, paragraph, indent):
        """Version-safe indentation handling"""
        try:
            paragraph.paragraph_format.left_indent = indent
        except AttributeError:
            paragraph.left_indent = indent
                
    def _handle_alignment(self, paragraph):
        """Enhanced justification handling with XML fallback"""
        try:
            # Modern versions (>=0.6.18)
            paragraph.paragraph_format.alignment = PP_ALIGN.JUSTIFY
        except AttributeError:
            try:
                # Legacy version fallback
                paragraph.alignment = PP_ALIGN.JUSTIFY
            except:
                # Direct XML manipulation for stubborn cases
                pPr = paragraph._element.get_or_add_pPr()
                align = OxmlElement('a:jc')
                align.set('val', 'dist')
                pPr.append(align)

    
    def _detect_unused_slides(self, distribution):
        """Adjusted slide retention logic with content-aware detection"""
        used_slides = set()
        content_slides = set()
        
        # Track slides that actually contain content
        for slide_idx, section, items in distribution:
            if items:  # Only consider slides with actual content
                content_slides.add(slide_idx)
            used_slides.add(slide_idx)
        
        # Calculate maximum content-bearing slide
        max_content_slide = max(content_slides) if content_slides else 0
        
        # Determine minimum slides to keep (content slides + buffer)
        min_slides_to_keep = max(2, max_content_slide + 1)
        
        # Preserve all slides up to min_slides_to_keep
        remove_slides = []
        for slide_idx in range(len(self.prs.slides)-1, min_slides_to_keep-1, -1):
            remove_slides.append(slide_idx)
            
        return sorted(remove_slides, reverse=True)
    
    def _remove_slides(self, slide_indices):
        """Remove slides by indices (must be in descending order)"""
        for slide_idx in slide_indices:
            if slide_idx < len(self.prs.slides):
                # Method 1: Using _sldIdLst (more reliable)
                xml_slides = self.prs.slides._sldIdLst
                slides = list(xml_slides)
                
                # Remove relationship
                rId = slides[slide_idx].rId
                self.prs.part.drop_rel(rId)
                
                # Remove from slide list
                xml_slides.remove(slides[slide_idx])
                
                #logging.info(f"Removed slide {slide_idx + 1}")
                
    #################### NEW SECTION ####################
    def parse_summary_markdown(self, md_content: str) -> str:
        """Enhanced summary extraction with better parsing"""
        summary_lines = []
        in_summary = False
        
        for line in md_content.strip().split('\n'):
            stripped = line.strip()
            if stripped.startswith('## Summary'):
                in_summary = True
                continue
            if in_summary:
                if stripped.startswith('##'):
                    break  # Next section found
                if stripped:  # Skip empty lines
                    summary_lines.append(stripped)
        
        summary_text = " ".join(summary_lines)
        #logging.info(f"Extracted summary text: {summary_text[:100]}...")  # Debug logging
        return summary_text

    def _populate_summary_section(self, shape, chunks: List[str]):
        """Populate summary section with bold Arial text"""
        tf = shape.text_frame
        tf.clear()
        
        for chunk in chunks:
            p = tf.add_paragraph()
            run = p.add_run()
            run.text = chunk
            run.font.bold = True
            run.font.name = 'Arial'
            run.font.size = Pt(9)
            run.font.color.rgb = RGBColor(0, 0, 0)  # Black
            
            # Set paragraph formatting
            try:
                p.paragraph_format.space_after = Pt(6)
            except AttributeError:
                pass

    def generate_full_report(self, md_content: str, summary_md: str, output_path: str):
        """Enhanced generation with proper summary population"""
        try:
            # Process main content
            items = self.parse_markdown(md_content)
            distribution = self._plan_content_distribution(items)
            self._validate_content_placement(distribution)
            
            # Process summary content
            summary_text = self.parse_summary_markdown(summary_md)
            
            # Calculate slides needed from distribution
            max_slide_used = max((slide_idx for slide_idx, _, _ in distribution), default=0)
            total_slides_needed = max_slide_used + 1
            
            # Split summary content appropriately
            wrapper = textwrap.TextWrapper(
                width=self.CHARS_PER_ROW, 
                break_long_words=True,
                replace_whitespace=False
            )
            summary_chunks = wrapper.wrap(summary_text)
            
            # Distribute summary chunks across slides
            chunks_per_slide = len(summary_chunks) // total_slides_needed + 1
            slide_summary_chunks = [
                summary_chunks[i:i+chunks_per_slide] 
                for i in range(0, len(summary_chunks), chunks_per_slide)
            ]
            
            # Ensure we have enough slides
            while len(self.prs.slides) < total_slides_needed:
                # Use layout index 1 for additional slides (2-column layout)
                self.prs.slides.add_slide(self.prs.slide_layouts[1])
            
            # Populate main content sections
            for slide_idx, section, section_items in distribution:
                if slide_idx >= len(self.prs.slides):
                    raise ValueError("Insufficient slides in template")
                slide = self.prs.slides[slide_idx]
                self.current_slide_index = slide_idx
                shape = self._get_section_shape(slide, section)
                if shape:
                    self._populate_section(shape, section_items)
            
            # Populate summary content on all slides
            for slide_idx in range(total_slides_needed):
                slide = self.prs.slides[slide_idx]
                summary_shape = next((s for s in slide.shapes if s.name == "coSummaryShape"), None)
                
                if summary_shape:
                    self._populate_summary_section_safe(
                        summary_shape, 
                        slide_summary_chunks[slide_idx] if slide_idx < len(slide_summary_chunks) else []
                    )
            
            # Remove unused slides
            unused_slides = self._detect_unused_slides(distribution)
            if unused_slides:
                logging.info(f"Removing unused slides: {[idx+1 for idx in unused_slides]}")
                self._remove_slides(unused_slides)
            
            self.prs.save(output_path)
            logging.info(f"Successfully generated PowerPoint with {len(self.prs.slides)} slides and summary content")
            
        except Exception as e:
            logging.error(f"Generation failed: {str(e)}")
            raise
        
    def _populate_summary_section_safe(self, shape, chunks: List[str]):
        """Summary section with updated formatting"""
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        
        tf.margin_left = Inches(0.07)
        tf.margin_right = Inches(0.07)
        
        p = tf.add_paragraph()
        full_text = " ".join(chunks)
        
        run = p.add_run()
        run.text = full_text
        run.font.size = Pt(9)  # All text pt 9
        run.font.bold = True
        run.font.name = 'Arial'
        run.font.color.rgb = RGBColor(255, 255, 255)

        # Set LEFT alignment
        try:
            p.paragraph_format.alignment = PP_ALIGN.LEFT
        except AttributeError:
            try:
                p.alignment = PP_ALIGN.LEFT
            except AttributeError:
                self._handle_legacy_alignment(p)
        
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP


# Usage
if __name__ == "__main__":
    parser =argparse.ArgumentParser(description='Generate PowerPoint presentation from Markdown')
    parser.add_argument('--template', required=True, help='Path to thhe template PowerPoint file')
    parser.add_argument('--markdown', required=True, help='Path to the markdown file')
    parser.add_argument('--output', required=True, help='Path to the output PowerPoint file')
    
    args = parser.parse_args()
    
    report_generator  = ReportGenerator(args.template, args.markdown, args.output)
    report_generator.generate()

# python 2.pptx.py --template template.pptx --markdown utils/bs_content.md --output final_report.pptx