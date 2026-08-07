#!/usr/bin/env python3
"""Why is coSummaryShape blank?

The executive-summary band shipping empty has survived several rounds of
fixes because every failure mode looks identical from the outside: the shape
is missing, or there is no source text, or the generator returned nothing, or
a blank pre-generated summary won -- each one quietly does nothing.

This runs the two checks that actually separate them, using files that
already exist on this machine. No databook, no AI, no 12-minute export.

    python diagnose_summary.py                          # template only
    python diagnose_summary.py path/to/exported.pptx    # both checks

CHECK A -- does the fill work with THIS machine's template?
    Runs the real export entry point over synthetic Chinese commentary. If
    the band fills here, the template and the fill logic are both fine and
    the problem is upstream, in what the real run feeds them. If it does NOT
    fill, the difference is the template, which is gitignored and per-machine
    -- so this is the only place that can be found.

CHECK B -- does the summary generator work on THIS deck's real text?
    Pulls the actual commentary off the exported deck's first slide and runs
    _generate_page_summary on it. Real Chinese, real sentence structure, no
    AI call. If this returns nothing, the generator is the culprit and the
    text that broke it is printed.
"""

import sys
from pathlib import Path

from pptx import Presentation

from fdd_utils.pptx import PowerPointGenerator, export_pptx_from_structured_data_combined

TEMPLATE = Path("fdd_utils/template.pptx")


def _bar(title):
    print("\n" + "=" * 74)
    print(f"  {title}")
    print("=" * 74)


def _report_summary_shapes(prs, label):
    found_any = False
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if "summary" not in (shape.name or "").lower():
                continue
            found_any = True
            text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            state = f"{len(text)} chars: {text[:80]!r}" if text else "*** EMPTY ***"
            exact = "" if shape.name == "coSummaryShape" else "  ❌ NAME IS NOT EXACTLY 'coSummaryShape'"
            print(f"  {label} slide {i + 1}  [{shape.name}]  {state}{exact}")
    if not found_any:
        print(f"  {label}: no summary-named shape on any slide")


def _warn_if_stale(deck: Path):
    """Is this deck older than the exporter that supposedly fixed it?

    A .pptx sitting in pptx_previews/ carries no version stamp, so a deck
    exported before a fix looks exactly like one exported after it and
    still broken. Diagnosing the stale one wastes a 12-minute re-run at
    best and sends the investigation somewhere false at worst.
    """
    import subprocess
    from datetime import datetime

    stamp = "%Y-%m-%d %H:%M"
    exported = datetime.fromtimestamp(deck.stat().st_mtime)
    print(f"  exported:         {exported:{stamp}}")

    # The comparison is against git where available, and the exporter file's
    # own mtime otherwise -- a machine without git on PATH still gets a real
    # answer. Whichever is used is named, and a failure to compare is stated
    # rather than swallowed: the first version of this check returned
    # silently on any exception, which is indistinguishable from "checked,
    # deck is current" -- the exact silent-no-op this whole investigation has
    # been chasing.
    exporter = Path("fdd_utils/pptx.py")
    try:
        changed = datetime.fromtimestamp(int(subprocess.run(
            ["git", "log", "-1", "--format=%ct", "--", str(exporter)],
            capture_output=True, text=True, check=True,
        ).stdout.strip()))
        source = "last commit touching fdd_utils/pptx.py"
    except Exception:
        if not exporter.exists():
            print("  ⚠️  cannot tell if this deck is current: no git and no "
                  "fdd_utils/pptx.py. Run from the repo root.")
            return
        changed = datetime.fromtimestamp(exporter.stat().st_mtime)
        source = "mtime of fdd_utils/pptx.py (git unavailable)"

    print(f"  exporter changed: {changed:{stamp}}  ({source})")
    if exported < changed:
        print("  ⚠️  THIS DECK PREDATES THE CURRENT EXPORTER — whatever it shows may "
              "already be fixed.\n      Re-export before drawing conclusions from it.")
    else:
        print("  ✅ deck is newer than the exporter — what it shows is current behaviour.")


def check_a():
    _bar("CHECK A — fill the band using THIS machine's template")
    if not TEMPLATE.exists():
        print(f"  template not found: {TEMPLATE.resolve()}")
        return
    print(f"  template: {TEMPLATE.resolve()}")

    def rows(names, category):
        return [{
            "mapping_key": n, "account_name": n, "category": category, "is_chinese": True,
            "commentary": (f"截至2026年03月31日，{n}余额为人民币1,234千元，"
                           f"较2025年12月31日增加人民币200千元，主要由于业务规模扩大所致。"),
        } for n in names]

    out = Path("diagnose_summary_output.pptx")
    export_pptx_from_structured_data_combined(
        template_path=str(TEMPLATE),
        bs_data=rows(["货币资金", "应收账款", "预付款项"], "流动资产"),
        is_data=rows(["营业收入", "营业成本"], "收入"),
        output_path=str(out),
        project_name="Diagnostic", language="Chinese",
        temp_path=None, selected_sheet=None,
        is_chinese_databook=True, bs_is_results=None,
    )
    print()
    _report_summary_shapes(Presentation(str(out)), "result")
    print(f"\n  (wrote {out} — delete it when done)")


def check_b(deck_path):
    _bar("CHECK B — run the summary generator on THIS deck's real commentary")
    # A mistyped path otherwise surfaces as a 30-line python-pptx traceback
    # ending in "Package not found", which reads like the deck is corrupt.
    path = Path(deck_path)
    if not path.exists():
        print(f"  ❌ no such file: {path.resolve()}")
        near = sorted(path.parent.glob("*.pptx")) if path.parent.exists() else []
        if near:
            print("  .pptx files in that folder:")
            for p in near:
                print(f"    {p}")
        else:
            print(f"  (folder {path.parent} has no .pptx files, or does not exist)")
        return
    prs = Presentation(str(path))
    print(f"  deck: {path.resolve()}")
    _warn_if_stale(path)
    print()
    _report_summary_shapes(prs, "as shipped:")

    gen = PowerPointGenerator(str(TEMPLATE))
    for i, slide in enumerate(prs.slides):
        parts = [
            s.text_frame.text.strip()
            for s in slide.shapes
            if (s.name or "").startswith("textMainBullets")
            and s.has_text_frame and s.text_frame.text.strip()
        ]
        if not parts:
            continue
        source = "\n\n".join(parts)
        result = gen._generate_page_summary(source, True)
        print(f"\n  slide {i + 1}: {len(source)} chars of commentary "
              f"-> summary {len(result or '')} chars")
        if result:
            print(f"    {result[:200]!r}")
        else:
            print("    ❌ RETURNED NOTHING — this is the failure. Source began:")
            print(f"    {source[:200]!r}")


if __name__ == "__main__":
    check_a()
    if len(sys.argv) > 1:
        check_b(sys.argv[1])
    else:
        print("\n(no deck path given — skipping CHECK B. Pass the exported "
              ".pptx to run it, e.g. for_test/pptx_previews/<name>.preview.pptx)")
