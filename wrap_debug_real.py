#!/usr/bin/env python3
"""Read-only: reproduces the wrap-line-count prediction for slide 1's real
textMainBullets text using THIS machine's REAL client-metrics measurer,
line by line, so it can be compared directly against a real PowerPoint
BoundHeight measurement (see diagnose_capacity_gap.py / the VBA macro
that produced it).

Tries 3 width variants for the wrapped (non-first) lines of each bullet,
since a Mac-only test already showed the hanging-indent width reduction
alone doesn't explain the full gap -- this checks whether it's a bigger
contributor on REAL client-metrics fonts, or whether the character-width
measurement itself is the larger factor.

Usage:
    python wrap_debug_real.py
"""
import sys

from fdd_utils.text_metrics import get_measurer, text_box_from_shape
from pptx import Presentation


def _load_config():
    try:
        import yaml
        with open("fdd_utils/config.yml", encoding="utf-8") as f:
            return yaml.safe_load(f) or {}
    except Exception:
        return {}


PARA_GAP = 2.2
REAL_BOUND_HEIGHT_PT = 288.0

PARAS_RAW = """■ 货币资金 - 截至2026年03月31日，货币资金余额为1.23亿元，均为银行存款，无使用受限资金。我方已核对截至2026年03月31日的银行对账单，未发现显著差异。
■ 预付款项 - 截至2026年03月31日，预付款项无余额。2025年12月31日余额为16.6万元，主要为预付保险费；2023年12月31日及2024年12月31日余额分别为23.6万元、9.9万元。2026年03月31日余额归零，推测主要系前期预付保险费随期间摊销或结转所致，具体原因仍待管理层确认。
■ 其他流动资产 - 截至2026年03月31日，其他流动资产余额为589.4万元，主要为应交增值税借方余额，其中已认证进项税额529.2万元、计提的进项税额60.2万元。该科目由2023年末289.2万元增至2024年末346.9万元，并于2025年末进一步增至529.2万元。我方已将截至2026年03月31日的余额核对至增值税申报表，未发现重大异常。
■ 固定资产 - 截至2026年03月31日，固定资产账面净值为6.65亿元，主要为房屋及土地。截至同日，固定资产原值为7.75亿元，累计折旧为1.10亿元。固定资产原值自2024年末起保持稳定；折旧采用年限平均法，土地一期及二期折旧月数分别为572个月和582个月，房屋及房屋改造按240个月折旧，电子设备及办公家具按3至5年折旧。
■ 长期待摊费用 - 截至2026年03月31日，长期待摊费用余额为355.8万元，主要为项目管理费和绿化改造工程费摊销后净值。相关原值系支付给上海熙麦企业管理有限公司的租赁管理费及绿化改造工程服务费，项目管理费和绿化改造工程均按5年摊销、残值为零。余额于2024年末由2023年末216.1万元降至199.6万元，2025年末增至364.2万元，增幅约82.5%，2026年03月31日回落至355.8万元。
■ 其他非流动资产 - 截至2026年03月31日，其他非流动资产余额为111.3万元，系待抵扣进项税余额。该科目由2023年12月31日的639.4万元降至2024年12月31日的442.5万元，并于2025年12月31日进一步降至193.3万元。2026年03月31日较2025年12月31日下降约42.4%。结合该余额性质，推测下降主要系待抵扣进项税持续抵扣或结转所致。""".split("\n")

ALL_PARAS = ["流动资产", PARAS_RAW[0], PARAS_RAW[1], PARAS_RAW[2],
             "非流动资产", PARAS_RAW[3], PARAS_RAW[4], PARAS_RAW[5]]


def run(measurer, box_width_pt, line_h, label, wrap_w, first_w, verbose=False):
    total_lines = 0
    total_pt = 0.0
    for para in ALL_PARAS:
        is_cat = para in ("流动资产", "非流动资产")
        if is_cat:
            wrapped = measurer.wrap(para, wrap_w)
            n = len(wrapped)
            total_lines += n
            total_pt += n * line_h
        else:
            wrapped = measurer.wrap(para, wrap_w, first_line_width_pt=first_w)
            n = len(wrapped)
            total_lines += n
            total_pt += n * line_h + PARA_GAP
            if verbose:
                print(f"  ({len(para)} chars) -> {n} lines:")
                for j, ln in enumerate(wrapped):
                    print(f"    L{j} ({len(ln)} chars): {ln!r}")
    diff = total_pt - REAL_BOUND_HEIGHT_PT
    print(f"{label}: lines={total_lines}  predicted_pt={total_pt:.2f}  "
          f"real_pt={REAL_BOUND_HEIGHT_PT}  diff={diff:+.2f}pt")
    return total_lines, total_pt


def main() -> int:
    config = _load_config()
    packing_cfg = ((config.get("pptx") or {}).get("commentary_packing") or {})
    metrics_chi = packing_cfg.get("font_metrics_path_chi") or "fdd_utils/font_metrics/msyh_chi.json"
    family_chi = packing_cfg.get("font_family_chi") or "Microsoft YaHei"

    measurer = get_measurer(family_chi, 9.0, is_cjk=True, line_spacing=1.0, metrics_path=metrics_chi)
    print(f"Measurement source: {measurer.source}  (must say 'client-metrics' for this to be meaningful)")
    line_h = measurer.line_height_pt()
    print(f"line_height_pt() = {line_h:.4f}")

    prs = Presentation("fdd_utils/template.pptx")
    shape = [s for s in prs.slides[0].shapes if s.name == "textMainBullets"][0]
    box = text_box_from_shape(shape)
    print(f"box.width_pt={box.width_pt:.3f}  box.height_pt={box.height_pt:.3f}")
    print()

    hang = 10.8
    print("--- Summary across width variants ---")
    run(measurer, box.width_pt, line_h, "A) WITH hanging-indent reduction (current formula)",
        box.width_pt - hang, box.width_pt)
    run(measurer, box.width_pt, line_h, "B) NO hanging-indent reduction (full width both lines)",
        box.width_pt, box.width_pt)
    run(measurer, box.width_pt, line_h, "C) Half hanging-indent (5.4pt)",
        box.width_pt - hang / 2, box.width_pt)
    print()
    print("--- Full line-by-line breakdown (variant A, the current formula) ---")
    run(measurer, box.width_pt, line_h, "A) WITH hanging-indent reduction",
        box.width_pt - hang, box.width_pt, verbose=True)

    print()
    print("Paste this ENTIRE output back -- the line-by-line breakdown under")
    print("variant A is the important part: if real PowerPoint fits ONE MORE")
    print("word/character on any specific line than shown here before wrapping,")
    print("that pinpoints the exact character-width measurement gap.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
