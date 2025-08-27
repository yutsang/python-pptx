#!/usr/bin/env python3
"""
Check what slide layouts are available in the template
"""

from pptx import Presentation
import os

def check_slide_layouts(template_path):
    """Check what slide layouts are available"""
    
    if not os.path.exists(template_path):
        print(f"❌ Template not found: {template_path}")
        return
    
    print(f"🔍 Checking slide layouts: {template_path}")
    print("=" * 60)
    
    prs = Presentation(template_path)
    
    print(f"📊 Template has {len(prs.slide_layouts)} slide layouts")
    
    for i, layout in enumerate(prs.slide_layouts):
        print(f"\n📋 Layout {i}: {layout.name}")
        print(f"   Type: {type(layout).__name__}")
        
        # Check what shapes are in this layout
        print(f"   Shapes: {len(layout.shapes)}")
        for j, shape in enumerate(layout.shapes):
            shape_name = getattr(shape, 'name', 'No name')
            shape_type = type(shape).__name__
            print(f"      Shape {j+1}: {shape_name} ({shape_type})")

if __name__ == "__main__":
    template_path = "fdd_utils/template.pptx"
    check_slide_layouts(template_path)
