#!/usr/bin/env python3
"""Pairs each account's SOURCE material in a databook with the COMMENTARY
that actually came out in the exported deck, so the writing style can be
judged against what the model was given rather than in isolation.

inspect_is_variance.py looks only at the databook; inspect_pptx.py looks
only at the deck. Neither answers the question this one is for: given
these remarks, is the output using them, and how is it phrasing them?

Per account it prints the substantive source remarks, the generated
bullet, whether each remark visibly reached the output, and any style
flags. Then a summary covering remark uptake and the style issues found
across the whole deck.

Style flags (each is a worklist item, not an automatic fail):
  * source meta-reference -- '根据备注信息...' cites the databook's own
    remarks column; a consultant states the fact, or attributes it to
    management ('管理层表示...'), which is an established convention here.
  * repeated 人民币 in one sentence.
  * tie-out ('Check') text leaking into the commentary -- a working
    artefact that should never appear in a deliverable.
  * a bullet that reproduces a remark almost verbatim rather than
    summarising it into report prose.

Usage:
    python inspect_databook_to_pptx.py --databook "for_test/x.xlsx" --pptx "for_test/x.pptx"
    python inspect_databook_to_pptx.py --databook ... --pptx ... --account 管理费用
"""
import argparse
import difflib
import re
import sys
import warnings

warnings.filterwarnings("ignore")
sys.path.insert(0, ".")

from pptx import Presentation

from fdd_utils.workbook import process_workbook_data
from inspect_is_variance import (
    _entity_from_filename, _substantive_notes, _substantive_linked, _note_prose,
)

_SOURCE_META_RE = re.compile(
    r"(根据|依据|按照)\s*(补充|管理层)?(备注|说明|备注信息)(信息|内容|说明)?|"
    r"(备注|说明)(显示|表明|提到|中提及|仅说明|未说明|未进一步说明)|"
    r"according\s+to\s+the\s+(supplementary\s+)?(remarks?|notes?)\b|"
    r"(as|per)\s+(stated|noted)\s+in\s+the\s+(supplementary\s+)?remarks?\b",
    re.IGNORECASE,
)
# A bare "Check"/"对账单" is NOT enough to call something an artefact: the
# reference deck legitimately writes 我方获取并核对了截至2026年3月31日的银行
# 对账单 (we obtained and checked the bank statement), which is exactly the
# work-performed language the report is supposed to contain. Only flag a
# marker that stands alone as a label, i.e. followed by a separator and
# figures rather than embedded in a sentence.
_TIEOUT_LEAK_RE = re.compile(
    r"(?:^|[\s。；;])(?:check|对帐单)\s*[|｜:：]\s*-?[\d,]+", re.IGNORECASE
)
_CHI_SENTENCE_SPLIT_RE = re.compile(r"[。；;]")


def _iter_text_shapes(shapes):
    """Yields every text-bearing shape, descending into groups -- a shape
    inside a group is not reachable from slide.shapes directly."""
    for shape in shapes:
        if shape.shape_type == 6 and hasattr(shape, "shapes"):  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_text_shapes(shape.shapes)
            continue
        if getattr(shape, "has_text_frame", False):
            yield shape


# The generator writes bullets as "■ <key> - <text>" (see pptx.py's
# key_prefix), but a deck can arrive with the marker stripped, replaced by
# PowerPoint's own list formatting, or using a different dash -- so match a
# family of markers and both dash characters rather than one literal.
_BULLET_MARKERS = ("■", "●", "▪", "•", "◆", "-")
_LABEL_SPLIT_RE = re.compile(r"\s+[-–—]\s+")


def extract_bullets(pptx_path: str, debug: bool = False):
    """{account_label: full bullet text} for every account bullet in the
    deck. A long account can be split across slots with a '(cont'd)' marker,
    so fragments are stitched back together under one label -- otherwise a
    continuation looks like a separate, truncated account and its remark
    uptake is judged against half the text.

    Returns (bullets, order, diagnostics). diagnostics is only populated
    when nothing matched, so the caller can show what the deck actually
    contains instead of silently reporting zero accounts."""
    prs = Presentation(pptx_path)
    bullets, order = {}, []
    sample_lines, shape_names = [], []

    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in _iter_text_shapes(slide.shapes):
            text = shape.text_frame.text or ""
            if text.strip():
                shape_names.append(f"slide {slide_idx} / {shape.name!r} ({len(text)} chars)")
            for line in text.split("\n"):
                s = line.strip()
                if not s:
                    continue
                if len(sample_lines) < 12:
                    sample_lines.append(f"slide {slide_idx} [{shape.name}] {s[:110]}")
                marker = next((m for m in _BULLET_MARKERS if s.startswith(m)), None)
                if marker is None:
                    continue
                body = s[len(marker):].strip()
                parts = _LABEL_SPLIT_RE.split(body, maxsplit=1)
                if len(parts) == 2:
                    label, body_text = parts
                else:
                    # No dash separator: treat a short leading line as a
                    # label-only bullet rather than dropping it entirely.
                    label, body_text = body, ""
                    if len(label) > 30:
                        continue
                label = re.sub(r"\s*[(（](cont'd|续|續)[)）]\s*$", "", label.strip()).strip()
                if not label:
                    continue
                if label in bullets:
                    bullets[label] += " " + body_text.strip()
                else:
                    bullets[label] = body_text.strip()
                    order.append(label)

    diagnostics = {}
    if not bullets:
        diagnostics = {"shapes": shape_names[:20], "sample_lines": sample_lines}
    return bullets, order, diagnostics


def _distinctive_tokens(prose: str):
    """CJK runs and long latin words from a remark -- used to ask whether
    that remark visibly reached the commentary. Short tokens are dropped
    because generic words ('费用', 'the') match anything."""
    toks = re.findall(r"[一-鿿]{2,}|[A-Za-z]{4,}", prose or "")
    return [t for t in toks if len(t) >= 2]


def _remark_reached_output(prose: str, output: str) -> bool:
    toks = _distinctive_tokens(prose)
    if not toks:
        return False
    hits = sum(1 for t in toks if t in output)
    return hits >= max(1, len(toks) // 3)


def _style_flags(text: str):
    """[(kind, detail)] -- kind is a stable short label so the summary can
    group by issue type rather than by the varying quoted text."""
    flags = []
    for m in _SOURCE_META_RE.finditer(text):
        flags.append(("source meta-reference",
                      f"{m.group(0)!r} -- state the fact directly"))
    for sentence in _CHI_SENTENCE_SPLIT_RE.split(text):
        n = sentence.count("人民币")
        if n >= 2:
            flags.append(("repeated 人民币 in one sentence",
                          f"x{n}: {sentence.strip()[:90]!r}"))
    for m in _TIEOUT_LEAK_RE.finditer(text):
        flags.append(("tie-out artefact leaked into the deliverable", repr(m.group(0))))
    return flags


def _verbatim_ratio(remark_prose: str, output: str) -> float:
    """Longest common run between a remark and the commentary, as a share
    of the remark. High values mean the note was pasted rather than
    rewritten into report prose."""
    if not remark_prose or not output:
        return 0.0
    sm = difflib.SequenceMatcher(None, remark_prose, output)
    match = sm.find_longest_match(0, len(remark_prose), 0, len(output))
    return match.size / max(1, len(remark_prose))


def extract_paragraphs(pptx_path: str):
    """[(slide, shape, text)] for every substantial paragraph in a deck that
    wasn't produced by this tool -- e.g. a real analyst deck built in
    PowerPoint/UpSlide, whose commentary sits in ordinary text placeholders
    with no '■ <account> - ' structure to key off. Boilerplate (titles, nav
    chrome, source lines) is short, so a length floor removes most of it."""
    prs = Presentation(pptx_path)
    out = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in _iter_text_shapes(slide.shapes):
            for para in (shape.text_frame.text or "").split("\n"):
                s = para.strip()
                if len(s) >= 40:
                    out.append((slide_idx, shape.name, s))
    return out


# A paragraph that hands off to a table instead of narrating the detail.
_TABLE_HANDOFF_RE = re.compile(r"明细如下|明细详见|详见下表|如下表|汇总如下|列示如下|详见附件")
# Work-performed language: what the consultant did, as opposed to what the
# numbers are. This is the register the reference deck uses constantly.
_WORK_PERFORMED_RE = re.compile(
    r"我(?:方|们)(?:已)?(?:获取|取得|核对|检查|查看|复核|比对)|未(?:发现|见)(?:明显|显著|重大)?(?:差异|异常)|"
    r"无重大(?:差异|异常)|未见异常|已核对至|核对至"
)
_MGMT_ATTRIB_RE = re.compile(r"管理层(?:表示|称|提供|确认|解释)")
_ADJUSTMENT_RE = re.compile(r"示意性调整|补计提|重分类")
_AMOUNT_RE = re.compile(r"(人民币\s*)?([\d,]+(?:\.\d+)?)\s*(万元|亿元|元)")
_OPENING_RE = re.compile(r"^(.{2,12}?)\s*[–—-]\s*(.{0,20})")


def style_profile(args, dfs) -> int:
    """Measures the conventions a human-written reference deck actually
    follows, instead of inferring them by eye.

    The point is calibration: the project's own prompts encode assumptions
    about currency prefixes, opening sentences and how much narrative an
    account should carry. This counts what the real deliverable does, so
    those assumptions can be checked against evidence rather than memory."""
    paragraphs = extract_paragraphs(args.pptx)
    if not paragraphs:
        print("❌ no substantial paragraphs found in this deck.")
        return 1
    texts = [t for _s, _n, t in paragraphs]
    joined = "\n".join(texts)

    print("=" * 78)
    print(f"STYLE PROFILE -- {args.pptx}")
    print("=" * 78)
    print(f"  paragraphs analysed: {len(texts)}")
    lens = sorted(len(t) for t in texts)
    print(f"  paragraph length: min {lens[0]}, median {lens[len(lens) // 2]}, max {lens[-1]} chars")

    amounts = _AMOUNT_RE.findall(joined)
    with_prefix = sum(1 for pfx, _n, _u in amounts if pfx)
    print(f"\n  AMOUNTS: {len(amounts)} figure(s) found")
    if amounts:
        print(f"    with a 人民币 prefix : {with_prefix} ({with_prefix / len(amounts) * 100:.0f}%)")
        print(f"    bare (no prefix)    : {len(amounts) - with_prefix} "
              f"({(len(amounts) - with_prefix) / len(amounts) * 100:.0f}%)")
        units = {}
        for _p, _n, u in amounts:
            units[u] = units.get(u, 0) + 1
        print(f"    units used          : {', '.join(f'{u} x{n}' for u, n in sorted(units.items(), key=lambda kv: -kv[1]))}")
        if with_prefix / len(amounts) < 0.25:
            print("    => the reference deck states amounts BARE by default and reserves")
            print("       人民币 for disambiguation (e.g. a figure converted from USD).")

    per_sentence = []
    for t in texts:
        for sent in re.split(r"[。；;]", t):
            n = sent.count("人民币")
            if n >= 2:
                per_sentence.append((n, sent.strip()[:80]))
    print(f"    sentences repeating 人民币 2+ times: {len(per_sentence)}")
    for n, s in per_sentence[:3]:
        print(f"      x{n}: {s!r}")

    # Bare-by-default is right, but the one case the convention keeps 人民币
    # for is a sentence that also names a foreign currency -- '注册资本为
    # 7,000万美元（折合人民币4.88亿元）'. A bare figure there is genuinely
    # ambiguous, so flag it rather than treat 0% as unambiguously ideal.
    ambiguous = []
    for t in texts:
        for sent in re.split(r"[。；;]", t):
            if not re.search(r"美元|港元|港币|欧元|日元|USD|HKD|EUR", sent):
                continue
            bare = [m for m in _AMOUNT_RE.finditer(sent) if not m.group(1)]
            if bare:
                ambiguous.append(sent.strip()[:110])
    if ambiguous:
        print(f"\n    ⚠️ {len(ambiguous)} sentence(s) name a FOREIGN currency yet leave a figure")
        print(f"       bare -- this is the case where 人民币 should be kept, for clarity:")
        for s in ambiguous[:3]:
            print(f"         {s!r}")
    elif amounts:
        print(f"    (no sentence mixes a foreign currency with a bare figure)")

    openings = []
    for t in texts:
        m = _OPENING_RE.match(t)
        if m:
            openings.append(m.group(2).strip()[:18])
    print(f"\n  OPENINGS: {len(openings)} paragraph(s) use the '<account> – <text>' form")
    starts = {}
    for o in openings:
        key = o[:6]
        starts[key] = starts.get(key, 0) + 1
    for k, n in sorted(starts.items(), key=lambda kv: -kv[1])[:6]:
        print(f"    {n:>3d} x  starts '{k}...'")

    def _count(rx, label):
        hits = [t for t in texts if rx.search(t)]
        print(f"    {len(hits):>3d} paragraph(s) ({len(hits) / len(texts) * 100:.0f}%)  {label}")
        return hits

    print(f"\n  REGISTER:")
    table_paras = _count(_TABLE_HANDOFF_RE, "hand off to a table ('明细如下' / '详见下表')")
    _count(_WORK_PERFORMED_RE, "state work performed ('我方核对了...未发现差异')")
    _count(_MGMT_ATTRIB_RE, "attribute to management ('管理层表示...')")
    _count(_ADJUSTMENT_RE, "mention an adjustment ('示意性调整' / '补计提')")

    if table_paras:
        print(f"\n  TABLE-INSTEAD-OF-PROSE (the POC you asked about):")
        print(f"    {len(table_paras)} paragraph(s) stop and defer to a table. Their length:")
        tl = sorted(len(t) for t in table_paras)
        print(f"      min {tl[0]}, median {tl[len(tl) // 2]}, max {tl[-1]} chars")
        others = [len(t) for t in texts if t not in table_paras]
        if others:
            others.sort()
            print(f"    vs paragraphs that narrate in full: median {others[len(others) // 2]} chars")
            print(f"    => a table-backed account gets a SHORT lead-in, not a long narrative.")
        print(f"    Examples:")
        for t in table_paras[:4]:
            print(f"      · {t[:150]}")

    # How close is the deck's prose to the databook's own remark text?
    print(f"\n  REMARK REUSE (databook remark -> deck sentence):")
    verbatim, paraphrased, absent = 0, 0, 0
    examples = []
    for key in sorted(dfs):
        df = dfs[key]
        for note in _substantive_notes(df.attrs.get("supporting_notes") or []):
            prose = _note_prose(note)
            if len(prose) < 12:
                continue
            best = max((_verbatim_ratio(prose, t) for t in texts), default=0.0)
            if best >= 0.6:
                verbatim += 1
                if len(examples) < 4:
                    examples.append((key, best, prose[:90]))
            elif best >= 0.2:
                paraphrased += 1
            else:
                absent += 1
    total_notes = verbatim + paraphrased + absent
    if total_notes:
        print(f"    {verbatim} reproduced near-verbatim ({verbatim / total_notes * 100:.0f}%)")
        print(f"    {paraphrased} partially reused ({paraphrased / total_notes * 100:.0f}%)")
        print(f"    {absent} not visible in the deck ({absent / total_notes * 100:.0f}%)")
        for key, ratio, prose in examples:
            print(f"      {ratio:.0%}  {key}: {prose!r}")
        if verbatim and verbatim >= paraphrased:
            print("    => the analyst writes the finished sentence INTO the databook remark")
            print("       column and the deck reproduces it. Near-verbatim reuse is the")
            print("       intended workflow here, not a defect to rewrite away.")
    return 0


def prose_mode(args, dfs, diagnostics) -> int:
    """Reports how a human-written deck talks about each account: which
    paragraphs mention it, and what the databook had available to say."""
    paragraphs = extract_paragraphs(args.pptx)
    print(f"{len(paragraphs)} substantial paragraph(s) in the deck "
          f"(>=40 chars, excludes titles/navigation).\n")
    if not paragraphs:
        print("  Text-bearing shapes found:")
        for s in diagnostics.get("shapes") or ["    (none)"]:
            print(f"    {s}")
        return 1

    matched_accounts, unmatched, used_paras = 0, [], set()
    for key in sorted(dfs):
        if args.account and key != args.account:
            continue
        df = dfs[key]
        display = str(df.attrs.get("display_key") or "").strip()
        needles = {n for n in (key, display) if n}
        hits = [(i, p) for i, p in enumerate(paragraphs)
                if any(n in p[2] for n in needles)]
        notes = _substantive_notes(df.attrs.get("supporting_notes") or [])
        rhs = df.attrs.get("adjacent_detail_rows") or []

        if not hits:
            unmatched.append(key)
            continue
        matched_accounts += 1
        print("=" * 78)
        print(f"{key}")
        print("=" * 78)
        print(f"  DATABOOK source material: {len(notes)} substantive note(s), {len(rhs)} RHS row(s)")
        for n in notes[:4]:
            print(f"    · {str(n)[:180]}")
        print(f"\n  HOW THE DECK WRITES IT ({len(hits)} paragraph(s) mentioning this account):")
        for i, (slide_no, shape_name, text) in hits[:3]:
            used_paras.add(i)
            print(f"    [slide {slide_no} / {shape_name}]")
            for j in range(0, min(len(text), 700), 108):
                print(f"      {text[j:j + 108]}")
            if len(text) > 700:
                print(f"      ... (+{len(text) - 700} chars)")
        flags = _style_flags(" ".join(p[2] for _i, p in hits))
        if flags:
            print(f"\n  STYLE FLAGS in the reference deck ({len(flags)}):")
            for kind, detail in flags:
                print(f"    ⚠️ {kind}: {detail}")
        print()

    print("=" * 78)
    print("SUMMARY")
    print("=" * 78)
    print(f"  accounts named somewhere in the deck : {matched_accounts}")
    print(f"  accounts never named                 : {len(unmatched)}")
    if unmatched:
        print(f"    {', '.join(unmatched[:18])}{' ...' if len(unmatched) > 18 else ''}")
    print(f"  deck paragraphs matched to an account: {len(used_paras)}/{len(paragraphs)}")
    print("\n  NOTE: this is a REFERENCE deck (human-written), so the value here is the")
    print("  target style, not a defect list -- read the paragraphs above to see how the")
    print("  team actually phrases an account given the same databook.")
    return 0


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--databook", required=True, help="path to the source databook .xlsx")
    ap.add_argument("--pptx", required=True, help="path to the exported deck .pptx")
    ap.add_argument("--entity", default=None, help="entity name (default: derived from the filename)")
    ap.add_argument("--sheet", default=None, help="specific sheet, if needed")
    ap.add_argument("--account", default=None, help="only report this one account")
    ap.add_argument("--style-profile", action="store_true",
                     help="measure the conventions a reference deck actually follows (currency "
                          "prefix rate, opening forms, work-performed / management-attribution / "
                          "table-handoff frequency, and how much databook remark text is reused "
                          "verbatim) -- use this on a human-written deck to calibrate the prompts")
    ap.add_argument("--verbatim-threshold", type=float, default=0.5,
                     help="share of a remark appearing as one unbroken run in the output above "
                          "which it's called near-verbatim (default 0.5)")
    args = ap.parse_args()

    entity = args.entity or _entity_from_filename(args.databook)
    print(f"Databook : {args.databook!r} (entity={entity!r})")
    print(f"Deck     : {args.pptx!r}\n")

    result = process_workbook_data(temp_path=args.databook, entity_name=entity,
                                    selected_sheet=args.sheet)
    dfs = result["dfs"]
    print(f"{len(dfs)} account(s) in the databook. Language: {result.get('language')}")

    if args.style_profile:
        return style_profile(args, dfs)

    bullets, order, diagnostics = extract_bullets(args.pptx)
    print(f"{len(bullets)} account bullet(s) in the deck: {', '.join(order[:12])}"
          f"{' ...' if len(order) > 12 else ''}\n")

    if not bullets:
        print("No '■ <account> - ' bullets found -- this deck was not produced by this tool.")
        print("Falling back to PROSE mode: matching account names inside the deck's own")
        print("paragraphs. A human-written reference deck is more useful this way anyway,")
        print("since it shows the target style rather than our own output.\n")
        return prose_mode(args, dfs, diagnostics)

    # Deck labels are display names; databook keys are mapping keys. Match
    # exactly first, then by containment either way, before giving up --
    # a label like '应付账款' can appear in the deck while the key is 'AP'.
    def find_bullet(key: str, df):
        display = str(df.attrs.get("display_key") or "").strip()
        for cand in (key, display):
            if cand and cand in bullets:
                return bullets[cand], cand
        for label, text in bullets.items():
            if key and (key in label or label in key):
                return text, label
            if display and (display in label or label in display):
                return text, label
        return None, None

    total, matched, with_remarks, used_remarks = 0, 0, 0, 0
    all_flags, unused, verbatim_hits = [], [], []

    for key in sorted(dfs):
        if args.account and key != args.account:
            continue
        df = dfs[key]
        notes = _substantive_notes(df.attrs.get("supporting_notes") or [])
        linked = _substantive_linked(df.attrs.get("table_linked_remarks") or [])
        rhs = df.attrs.get("adjacent_detail_rows") or []
        total += 1

        output, label = find_bullet(key, df)
        print("=" * 78)
        print(f"{key}" + (f"   [deck label: {label!r}]" if label and label != key else ""))
        print("=" * 78)
        if output is None:
            print("  ⚠️  no matching bullet found in the deck (account not written up, or the")
            print("      deck label differs too much to match automatically)\n")
            continue
        matched += 1

        print(f"  SOURCE remarks ({len(notes)} note(s), {len(rhs)} RHS row(s), {len(linked)} linked):")
        if not notes and not rhs and not linked:
            print("    (none substantive -- commentary can only restate the figures)")
        for n in notes[:6]:
            print(f"    · {str(n)[:200]}")
        for r in rhs[:4]:
            desc = r.get("Description") if isinstance(r, dict) else None
            extra = [f"{k}={str(v)[:60]}" for k, v in (r.items() if isinstance(r, dict) else [])
                     if "Detail" in str(k) and str(v).strip()]
            if desc or extra:
                print(f"    · [RHS] {desc}: {'; '.join(extra[:2])}"[:200])

        print(f"\n  GENERATED commentary ({len(output)} chars):")
        for i in range(0, min(len(output), 900), 110):
            print(f"    {output[i:i + 110]}")
        if len(output) > 900:
            print(f"    ... (+{len(output) - 900} more chars)")

        if notes or rhs:
            with_remarks += 1
            reached, missed = [], []
            for n in notes:
                prose = _note_prose(n)
                (reached if _remark_reached_output(prose, output) else missed).append(prose[:70])
            if reached:
                used_remarks += 1
            print(f"\n  remark uptake: {len(reached)}/{len(notes)} note(s) visibly reached the output")
            for m in missed[:4]:
                print(f"    ✗ not visible: {m!r}")
                unused.append((key, m))
            for n in notes:
                prose = _note_prose(n)
                ratio = _verbatim_ratio(prose, output)
                if ratio >= args.verbatim_threshold and len(prose) >= 12:
                    print(f"    ⚠️ near-verbatim ({ratio:.0%} of the note is one unbroken run in the "
                          f"output) -- summarise rather than paste: {prose[:60]!r}")
                    verbatim_hits.append((key, prose[:60], ratio))

        flags = _style_flags(output)
        if flags:
            print(f"\n  STYLE FLAGS ({len(flags)}):")
            for kind, detail in flags:
                print(f"    ⚠️ {kind}: {detail}")
            all_flags.extend((key, kind) for kind, _detail in flags)
        print()

    print("=" * 78)
    print("SUMMARY")
    print("=" * 78)
    print(f"  accounts examined                 : {total}")
    print(f"  matched to a deck bullet          : {matched}")
    print(f"  had substantive source material   : {with_remarks}")
    print(f"  ... of which visibly used any     : {used_remarks}")
    if with_remarks:
        print(f"  => remark uptake rate: {used_remarks / with_remarks * 100:.0f}%")
    print(f"  style flags across the deck       : {len(all_flags)}")
    if all_flags:
        by_kind = {}
        for _key, kind in all_flags:
            by_kind[kind] = by_kind.get(kind, 0) + 1
        for kind, n in sorted(by_kind.items(), key=lambda kv: -kv[1]):
            print(f"    {n:>3d} x  {kind}")
    if verbatim_hits:
        print(f"\n  near-verbatim remark reuse: {len(verbatim_hits)}")
        for key, prose, ratio in verbatim_hits[:8]:
            print(f"    {ratio:.0%}  {key}: {prose!r}")
    if unused:
        print(f"\n  source remarks that did NOT visibly reach the deck: {len(unused)}")
        for key, m in unused[:10]:
            print(f"    {key}: {m!r}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
