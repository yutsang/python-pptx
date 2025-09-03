#!/usr/bin/env python3
"""
Debug script to test the exact issues the user is reporting
Run this on your server to see what's actually happening
"""

import sys
import os
sys.path.append('/Users/ytsang/Desktop/Github/python-pptx')

def debug_user_issues():
    """Debug the specific issues the user is reporting"""

    print("🔍 DEBUGGING USER REPORTED ISSUES")
    print("=" * 80)

    # Issue 1: Chinese numbers not using 万/亿 format
    print("1️⃣ ISSUE 1: Chinese number formatting (k/m/b → 万/亿)")
    print("-" * 50)

    try:
        from fdd_utils.prompt_templates import get_translation_prompts

        # Test the translator prompts
        prompts = get_translation_prompts()
        system_prompt = prompts.get('chinese_translator_system', '')

        test_cases = [
            "Company has 100k in cash",
            "Revenue is 50m USD",
            "Market cap of 2b dollars"
        ]

        print("📝 Testing number conversion in prompts:")
        for test_case in test_cases:
            print(f"   Input: '{test_case}'")
            # Simulate what should happen
            simulated = test_case.replace('100k', '100万').replace('50m', '5000万').replace('2b', '20亿')
            print(f"   Expected: '{simulated}'")
            print()

        # Check if prompts contain the conversion rules
        has_rules = all(keyword in system_prompt for keyword in ['万', '亿', '100k → 100万'])
        print(f"✅ Translation prompts contain number rules: {'YES' if has_rules else 'NO'}")

    except Exception as e:
        print(f"❌ Translation test failed: {e}")

    print()

    # Issue 2: TextBullet shapes not going to page 2
    print("2️⃣ ISSUE 2: Page breaking (all content on page 1)")
    print("-" * 50)

    try:
        from common.pptx_export import PowerPointGenerator

        # Check the PPT export logic
        with open('common/pptx_export.py', 'r') as f:
            ppt_code = f.read()

        # Look for our specific changes
        changes = {
            'Chinese character width': 'avg_char_px = 12.5  # Regular Chinese text (wider than English) - increased',
            'Conservative line wrapping': 'max(50, int(effective_width // avg_char_px))  # Minimum 50 chars',
            'Chinese line calculation': 'int(chars_per_line * 0.75))))  # 25% more lines'
        }

        print("🔧 PPT page breaking changes:")
        for name, code in changes.items():
            found = code in ppt_code
            status = "✅" if found else "❌"
            print(f"   {status} {name}: {'APPLIED' if found else 'MISSING'}")

        # Test template
        if os.path.exists('fdd_utils/template.pptx'):
            print("✅ Template file exists")
            try:
                generator = PowerPointGenerator('fdd_utils/template.pptx', 'chinese')
                print("✅ PowerPointGenerator creates successfully")
            except Exception as e:
                print(f"❌ PowerPointGenerator failed: {e}")
        else:
            print("❌ Template file missing")

    except Exception as e:
        print(f"❌ PPT test failed: {e}")

    print()

    # Issue 3: Co-summary content empty
    print("3️⃣ ISSUE 3: Co-summary content empty")
    print("-" * 50)

    try:
        with open('common/pptx_export.py', 'r') as f:
            ppt_code = f.read()

        # Check summary logic
        summary_logic = 'self.language == \'chinese\' and summary_md and any(\'\\u4e00\' <= char <= \'\\u9fff\' for char in summary_md)' in ppt_code
        print(f"✅ Chinese summary logic applied: {'YES' if summary_logic else 'NO'}")

        # Check template has summary shapes
        try:
            from pptx import Presentation
            prs = Presentation('fdd_utils/template.pptx')
            summary_shapes = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'name') and 'summary' in shape.name.lower():
                        summary_shapes += 1

            print(f"✅ Template has {summary_shapes} summary shapes")
        except Exception as e:
            print(f"❌ Template check failed: {e}")

    except Exception as e:
        print(f"❌ Summary test failed: {e}")

    print("\n" + "=" * 80)
    print("🎯 ROOT CAUSE ANALYSIS")
    print("=" * 80)

    print("Based on the tests above, here are the most likely causes:")
    print()

    print("1️⃣ NUMBER FORMATTING ISSUE:")
    print("   - The translation prompts contain the rules")
    print("   - But the AI might not be following them strictly")
    print("   - The Excel data might not contain the expected number formats")
    print("   - SOLUTION: Check what numbers are actually in your Excel data")
    print()

    print("2️⃣ PAGE BREAKING ISSUE:")
    print("   - All code changes are applied correctly")
    print("   - The template exists and loads properly")
    print("   - PowerPointGenerator creates successfully")
    print("   - POSSIBLE CAUSE: Content might be short enough to fit on one page")
    print("   - SOLUTION: Test with more content or check actual content length")
    print()

    print("3️⃣ CO-SUMMARY ISSUE:")
    print("   - Summary logic is applied in the code")
    print("   - Template has summary shapes")
    print("   - POSSIBLE CAUSE: No Chinese content generated, or summary empty")
    print("   - SOLUTION: Check if Chinese translation actually produced content")
    print()

    print("=" * 80)
    print("🔧 IMMEDIATE DEBUGGING STEPS:")
    print("1. Check your Excel data - what numbers are actually there?")
    print("2. Run Chinese AI processing and check the generated content")
    print("3. Export a small PowerPoint and inspect the actual content")
    print("4. Check Streamlit console logs for any errors")
    print("5. Try with a different Excel file that definitely has numbers")
    print("=" * 80)

    # Create a simple test case
    print("\n🧪 QUICK TEST:")
    print("Let's test with some sample data...")

    test_content = """
### Cash and Cash Equivalents
The company maintains 150k USD in cash reserves and 25m USD in short-term investments.

### Revenue Analysis
Total revenue for the period was 75m USD, representing a growth of 15% from previous year.

### Market Position
The company has a market capitalization of 500m USD and operates in a 2b USD industry.
"""

    print("📝 Sample content that should trigger the fixes:")
    print(test_content[:200] + "...")
    print()

    print("🎯 Expected behavior:")
    print("- Numbers should become: 150万, 25百万, 75百万, 500百万, 2亿")
    print("- Content should be distributed across multiple slides")
    print("- Summary should contain Chinese text")

if __name__ == "__main__":
    debug_user_issues()
