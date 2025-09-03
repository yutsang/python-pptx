#!/usr/bin/env python3
"""
Comprehensive debug script to test ALL three issues systematically
"""

import sys
import os
sys.path.append('/Users/ytsang/Desktop/Github/python-pptx')

def test_issue_1_chinese_numbers():
    """Test Issue 1: Chinese number formatting"""
    print("🔍 TESTING ISSUE 1: Chinese Number Formatting")
    print("=" * 60)
    
    try:
        from fdd_utils.prompt_templates import get_translation_prompts
        
        prompts = get_translation_prompts()
        system_prompt = prompts.get('chinese_translator_system', '')
        user_prompt = prompts.get('chinese_translator_user', '')
        
        # Check if prompts contain the number conversion rules
        required_elements = [
            '万', '亿', '100k → 100万', '50m → 5000万', '2b → 20亿'
        ]
        
        print("📋 Checking translation prompts:")
        all_found = True
        for element in required_elements:
            found = element in system_prompt and element in user_prompt
            status = "✅" if found else "❌"
            print(f"   {status} '{element}': {'FOUND' if found else 'MISSING'}")
            if not found:
                all_found = False
        
        print()
        print(f"🎯 RESULT: {'PASS' if all_found else 'FAIL'}")
        return all_found
        
    except Exception as e:
        print(f"❌ Test failed: {e}")
        return False

def test_issue_2_page_breaking():
    """Test Issue 2: Page breaking functionality"""
    print("\n🔍 TESTING ISSUE 2: Page Breaking")
    print("=" * 60)

    try:
        # Try to read the file with proper encoding
        try:
            with open('common/pptx_export.py', 'r', encoding='utf-8') as f:
                ppt_code = f.read()
        except UnicodeDecodeError:
            # If UTF-8 fails, try with errors='ignore'
            print("⚠️  UTF-8 decoding failed, trying alternative...")
            with open('common/pptx_export.py', 'r', encoding='utf-8', errors='ignore') as f:
                ppt_code = f.read()

        # Check for all the page breaking changes
        required_changes = [
            'avg_char_px = 12.5',
            'max(50, int(effective_width',
            '0.75))))  # 25% more lines',
            '0.78))))  # 22% more lines',
            '0.82))))  # 18% more lines'
        ]

        print("📋 Checking page breaking changes:")
        all_found = True
        for change in required_changes:
            found = change in ppt_code
            status = "✅" if found else "❌"
            print(f"   {status} Page breaking change: {'APPLIED' if found else 'MISSING'}")
            if not found:
                all_found = False

        # Test template loading
        if os.path.exists('fdd_utils/template.pptx'):
            print("✅ Template file exists")
            try:
                from common.pptx_export import PowerPointGenerator
                generator = PowerPointGenerator('fdd_utils/template.pptx', 'chinese')
                print("✅ PowerPointGenerator created successfully")
            except Exception as e:
                print(f"❌ PowerPointGenerator failed: {e}")
                all_found = False
        else:
            print("❌ Template file missing")
            all_found = False

        print()
        print(f"🎯 RESULT: {'PASS' if all_found else 'FAIL'}")
        return all_found

    except Exception as e:
        print(f"❌ Test failed: {e}")
        print("💡 This might be a Windows encoding issue.")
        print("   Try running: chcp 65001 && python debug_comprehensive_test.py")
        return False

def test_issue_3_co_summary():
    """Test Issue 3: Co-summary content"""
    print("\n🔍 TESTING ISSUE 3: Co-Summary Content")
    print("=" * 60)

    try:
        # Try to read the file with proper encoding
        try:
            with open('common/pptx_export.py', 'r', encoding='utf-8') as f:
                ppt_code = f.read()
        except UnicodeDecodeError:
            # If UTF-8 fails, try with errors='ignore'
            print("⚠️  UTF-8 decoding failed, trying alternative...")
            with open('common/pptx_export.py', 'r', encoding='utf-8', errors='ignore') as f:
                ppt_code = f.read()

        # Check for summary logic
        summary_logic = 'self.language == \'chinese\' and summary_md and any(\'\\u4e00\' <= char <= \'\\u9fff\' for char in summary_md)' in ppt_code
        print(f"✅ Chinese summary logic: {'APPLIED' if summary_logic else 'MISSING'}")

        # Check template summary shapes
        try:
            from pptx import Presentation
            prs = Presentation('fdd_utils/template.pptx')

            summary_shapes = []
            for i, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, 'name') and 'summary' in shape.name.lower():
                        summary_shapes.append((i, shape.name))

            print(f"✅ Template has {len(summary_shapes)} summary shapes:")
            for slide_idx, shape_name in summary_shapes:
                print(f"   - Slide {slide_idx}: {shape_name}")

            if summary_shapes:
                print("✅ Summary shapes found in template")
                return True
            else:
                print("❌ No summary shapes found in template")
                return False

        except Exception as e:
            print(f"❌ Template check failed: {e}")
            return False

    except Exception as e:
        print(f"❌ Test failed: {e}")
        print("💡 This might be a Windows encoding issue.")
        print("   Try running: chcp 65001 && python debug_comprehensive_test.py")
        return False

def test_actual_workflow():
    """Test the actual workflow to see what happens"""
    print("\n🔍 TESTING ACTUAL WORKFLOW")
    print("=" * 60)
    
    try:
        # Test with sample data that should trigger the fixes
        test_content = """
### Cash and Cash Equivalents
The company maintains 150k USD in cash reserves and 25m USD in short-term investments.

### Revenue Analysis  
Total revenue for the period was 75m USD, representing a growth of 15% from previous year.

### Market Position
The company has a market capitalization of 500m USD and operates in a 2b USD industry.

### Additional Financial Metrics
Operating expenses totaled 30m USD while net income reached 45m USD.
"""
        
        print("📝 Sample content that should trigger fixes:")
        print(test_content[:200] + "...")
        print()
        
        # Check if content has the right format to trigger fixes
        has_numbers = any(char in test_content for char in ['k', 'm', 'b'])
        print(f"✅ Content contains k/m/b numbers: {'YES' if has_numbers else 'NO'}")
        
        # Check content length
        content_length = len(test_content)
        print(f"✅ Content length: {content_length} characters")
        
        # Simulate what should happen
        print("\n🎯 EXPECTED BEHAVIOR:")
        print("1. Numbers: 150k → 150万, 25m → 2500万, 75m → 7500万, 500m → 5亿, 2b → 20亿")
        print("2. Page breaking: Content distributed across multiple slides")
        print("3. Co-summary: Chinese content in summary sections")
        
        return True
        
    except Exception as e:
        print(f"❌ Workflow test failed: {e}")
        return False

def main():
    """Run all tests"""
    print("🧪 COMPREHENSIVE DEBUG TESTING")
    print("=" * 80)
    
    results = []
    
    # Test all three issues
    results.append(("Chinese Number Formatting", test_issue_1_chinese_numbers()))
    results.append(("Page Breaking", test_issue_2_page_breaking()))
    results.append(("Co-Summary Content", test_issue_3_co_summary()))
    
    # Test actual workflow
    results.append(("Actual Workflow", test_actual_workflow()))
    
    # Summary
    print("\n" + "=" * 80)
    print("📊 FINAL RESULTS:")
    print("=" * 80)
    
    all_pass = True
    for test_name, result in results:
        status = "✅" if result else "❌"
        print(f"{status} {test_name}: {'PASS' if result else 'FAIL'}")
        if not result:
            all_pass = False
    
    print("\n" + "=" * 80)
    
    if all_pass:
        print("🎉 ALL TESTS PASSED!")
        print("   The code changes are correctly applied.")
        print("   If you're still seeing issues, the problem is likely:")
        print("   1. Streamlit caching (restart needed)")
        print("   2. Browser cache (clear needed)")
        print("   3. Multiple instances running")
        print("   4. Testing with wrong data format")
    else:
        print("❌ SOME TESTS FAILED!")
        print("   The code changes are NOT fully applied.")
        print("   Please check the failed tests above.")
    
    print("\n🔧 IMMEDIATE ACTIONS:")
    print("1. Fix Windows encoding (if on Windows):")
    print("   chcp 65001")
    print()
    print("2. Clear Python cache:")
    print("   find . -name '*.pyc' -delete")
    print("   (or on Windows: for /r %i in (*.pyc) do del \"%i\")")
    print()
    print("3. Restart Streamlit:")
    print("   pkill -f streamlit && streamlit run fdd_app.py")
    print("   (or on Windows: taskkill /f /im streamlit.exe && streamlit run fdd_app.py)")
    print()
    print("4. Clear browser cache completely")
    print("5. Test with Excel data containing k/m/b numbers")
    print("=" * 80)
    
    return all_pass

if __name__ == "__main__":
    main()
