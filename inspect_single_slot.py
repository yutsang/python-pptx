"""Deep-dive diagnostic for ONE textMainBullets shape -- dumps geometry,
margins, capacity math, and a per-paragraph wrap breakdown, so a visible
"box looks empty but inspect_pptx.py says it's 95% full" mismatch can be
pinned down to a specific number instead of a screenshot.

Usage:
    python inspect_single_slot.py "output.pptx" --slide 1 --shape textMainBullets
    python inspect_single_slot.py "output.pptx" --slide 2 --shape textMainBullets_L
"""
import argparse
import sys

from pptx import Presentation
from pptx.util import Emu

from fdd_utils.financial_common import load_yaml_file
from fdd_utils.text_metrics import get_measurer, text_box_from_shape

DEFAULT_CONFIG_CANDIDATES = ["fdd_utils/config.yml", "fdd_utils/config.example.yml"]


def _load_config():
    for candidate in DEFAULT_CONFIG_CANDIDATES:
        try:
            cfg = load_yaml_file(candidate)
            if cfg:
                return cfg
        except (FileNotFoundError, OSError):
            continue
    return {}


def _is_chinese_text(text: str) -> bool:
    chinese_chars = sum(1 for ch in text if "一" <= ch <= "鿿")
    return len(text) > 0 and (chinese_chars / len(text)) > 0.3


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path")
    ap.add_argument("--slide", type=int, required=True, help="1-based slide number")
    ap.add_argument("--shape", required=True, help="Shape name, e.g. textMainBullets or textMainBullets_L")
    args = ap.parse_args()

    config = _load_config()
    packing_cfg = ((config.get("pptx") or {}).get("commentary_packing") or {})
    metrics_eng = packing_cfg.get("font_metrics_path_eng") or "fdd_utils/font_metrics/arial_eng.json"
    metrics_chi = packing_cfg.get("font_metrics_path_chi") or "fdd_utils/font_metrics/msyh_chi.json"
    family_eng = packing_cfg.get("font_family_eng") or "Arial"
    family_chi = packing_cfg.get("font_family_chi") or "Microsoft YaHei"

    # line_spacing=1.0 for BOTH languages, and a flat 3pt inter-paragraph gap
    # below -- matches _fill_text_main_bullets_with_category_and_key, the
    # function that ACTUALLY sets a textMainBullets run's formatting
    # (hardcodes space_before=Pt(0)/space_after=Pt(3)/line_spacing=1.0
    # unconditionally). get_line_spacing_for_text/get_space_after_for_text/
    # get_space_before_for_text's language-dependent values belong to a
    # separate, legacy code path never reached by the live commentary
    # renderer -- assuming those inflated "capacity used" by roughly 30%.
    eng_measurer = get_measurer(family_eng, 9.0, is_cjk=False, line_spacing=1.0, metrics_path=metrics_eng)
    chi_measurer = get_measurer(family_chi, 9.0, is_cjk=True, line_spacing=1.0, metrics_path=metrics_chi)
    print(f"Measurement source: ENG={eng_measurer.source}  CHI={chi_measurer.source}")

    prs = Presentation(args.pptx_path)
    slide = prs.slides[args.slide - 1]
    shape = None
    for s in slide.shapes:
        if s.name == args.shape:
            shape = s
            break
    if shape is None:
        print(f"Shape {args.shape!r} not found on slide {args.slide}. Shapes on this slide:")
        for s in slide.shapes:
            print(f"  - {s.name!r}")
        return 1

    print(f"\nShape: {shape.name!r}")
    print(f"  left={Emu(shape.left).inches:.3f}in top={Emu(shape.top).inches:.3f}in "
          f"width={Emu(shape.width).inches:.3f}in height={Emu(shape.height).inches:.3f}in")

    body_pr = shape.text_frame._txBody.bodyPr
    print(f"  bodyPr attrs: {dict(body_pr.attrib)}")
    for tag in ("a:noAutofit", "a:normAutofit", "a:spAutoFit"):
        el = body_pr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}" + tag.split(":")[1])
        if el is not None:
            print(f"  autofit element: {tag} attrs={dict(el.attrib)}")

    full_text = shape.text_frame.text
    is_chi = _is_chinese_text(full_text)
    measurer = chi_measurer if is_chi else eng_measurer
    print(f"  whole-shape language guess: {'Chinese' if is_chi else 'English'} (measurer={measurer.source})")

    box = text_box_from_shape(shape)
    line_h = measurer.line_height_pt()
    para_gap = 3.0
    std_lh = line_h + para_gap
    # Float, not int(...) floored -- matches fdd_utils/pptx.py's
    # _calculate_max_lines_for_textbox (fixed in 5bbec43). This diagnostic
    # script had its own standalone copy of the old int()-floored formula
    # that never got updated alongside that fix, so it kept reporting the
    # pre-fix capacity/percentages even against freshly re-exported files.
    capacity = box.height_pt / std_lh
    print(f"\n  box.width_pt={box.width_pt:.2f} box.height_pt={box.height_pt:.2f}")
    print(f"  line_h={line_h:.2f}pt para_gap={para_gap:.2f}pt std_lh={std_lh:.2f}pt")
    print(f"  capacity = {box.height_pt:.2f} / {std_lh:.2f} = {capacity:.2f} lines")
    print(f"  capacity in inches = {capacity * std_lh / 72:.3f}in  (box height = {box.height_pt/72:.3f}in)")

    print("\n  Per-paragraph breakdown:")
    total_units = 0.0
    total_pt = 0.0
    for i, p in enumerate(shape.text_frame.paragraphs):
        text = "".join(r.text for r in p.runs)
        if not text.strip():
            print(f"    para[{i}]: BLANK (still consumes real vertical space in PowerPoint)")
            continue
        p_is_chi = _is_chinese_text(text)
        p_measurer = chi_measurer if p_is_chi else eng_measurer
        p_line_h = p_measurer.line_height_pt()
        p_para_gap = 3.0
        wrapped = p_measurer.wrap(text, box.width_pt)
        p_pt = len(wrapped) * p_line_h + p_para_gap
        units = p_pt / std_lh
        total_units += units
        total_pt += p_pt
        print(f"    para[{i}]: {len(wrapped)} wrapped lines, {p_pt:.1f}pt -> {units:.2f} units | {text[:50]!r}")

    print(f"\n  TOTAL: {total_units:.2f} units / {capacity:.2f} capacity = {100*total_units/capacity:.1f}% (line-unit basis)")
    print(f"  TOTAL: {total_pt:.1f}pt / {box.height_pt:.1f}pt box = {100*total_pt/box.height_pt:.1f}% (raw-height basis)")
    print("\n  If these two percentages differ noticeably, the capacity formula's")
    print("  std_lh unit conversion doesn't track raw physical height 1:1 --")
    print("  that gap is worth investigating further.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
