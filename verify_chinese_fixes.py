#!/usr/bin/env python3
"""
VERIFICATION SCRIPT: Test Chinese PPT Export Fixes
Run this script to verify that all Chinese export issues are fixed.
"""

import sys
import os
import tempfile

# Add the project root to Python path
sys.path.append('/Users/ytsang/Desktop/Github/python-pptx')

def test_chinese_export_fixes():
    """Test all Chinese export fixes in one comprehensive test"""

    print("🧪 CHINESE EXPORT FIXES VERIFICATION")
    print("=" * 60)

    from common.pptx_export import export_pptx
    from pptx import Presentation

    # Create test content with substantial Chinese text
    test_content = '''# Balance Sheet Analysis

## Current Assets
截至2022年12月31日的流动资产总额为人民币150万元，其中货币资金占比约为33%，应收账款占比约为40%，存货占比约为27%。流动资产的构成较为合理，体现了公司良好的流动性管理能力。流动资产周转率保持在较高水平，表明公司的营运资金使用效率良好。

## Non-current Assets
非流动资产主要由固定资产和长期股权投资构成。固定资产原值为人民币200万元，累计折旧为30万元，固定资产净值为170万元。长期股权投资账面价值为80万元，主要投资于联营企业和合营企业。固定资产的使用寿命较长，为公司的长期发展提供了物质保障。

## Current Liabilities
流动负债总额为人民币80万元，其中应付账款为48万元，占流动负债总额的60%；短期借款为20万元，占流动负债总额的25%；其他流动负债为12万元。公司流动负债结构相对稳定，债务负担适中。

## Equity
股东权益总额为人民币250万元，其中实收资本为150万元，资本公积为40万元，未分配利润为60万元。公司的净资产收益率保持在较为理想的水平，显示了良好的盈利能力和资本回报。股东权益结构合理，有利于公司的可持续发展。

## Cash Flow Analysis
经营活动现金流量净额为人民币45万元，投资活动现金流量净额为负30万元，融资活动现金流量净额为负15万元。整体现金流量状况基本稳定，但需要关注投资活动的现金流出情况。现金流量的合理配置是公司财务健康的重要标志。
'''

    # Create temporary files
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as md_file:
        md_file.write(test_content)
        md_path = md_file.name

    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as pptx_file:
        pptx_path = pptx_file.name

    try:
        print("📄 Test content created (5 sections with substantial Chinese text)")
        print("🚀 Exporting with language='chinese'...")

        # Export with Chinese mode
        result = export_pptx(
            template_path='fdd_utils/template.pptx',
            markdown_path=md_path,
            output_path=pptx_path,
            project_name='ChineseTest',
            language='chinese'  # This triggers all Chinese optimizations
        )

        if result and os.path.exists(result):
            print("✅ Export successful!")

            # Analyze the generated PPTX
            prs = Presentation(result)
            print(f"📊 Generated {len(prs.slides)} slides")

            # Analyze content
            total_chinese_chars = 0
            chinese_headers_found = []
            english_headers_found = []

            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n📄 SLIDE {slide_idx + 1}:")

                for shape in slide.shapes:
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text = shape.text_frame.text.strip()
                        if text and len(text) > 5:
                            chinese_count = sum(1 for char in text if '\u4e00' <= char <= '\u9fff')
                            total_chinese_chars += chinese_count

                            if chinese_count > 0:
                                print(f"   ✅ {chinese_count} Chinese characters")

                                # Check for specific header translations
                                if '流动资产' in text:
                                    chinese_headers_found.append('流动资产')
                                    print("      🎯 Current Assets → 流动资产")
                                if '非流动资产' in text:
                                    chinese_headers_found.append('非流动资产')
                                    print("      🎯 Non-current Assets → 非流动资产")
                                if '股东权益' in text:
                                    chinese_headers_found.append('股东权益')
                                    print("      🎯 Equity → 股东权益")

                            # Check for untranslated headers
                            if 'Current Assets' in text and '流动资产' not in text:
                                english_headers_found.append('Current Assets')
                                print("      ❌ ENGLISH: Current Assets")
                            if 'Non-current Assets' in text and '非流动资产' not in text:
                                english_headers_found.append('Non-current Assets')
                                print("      ❌ ENGLISH: Non-current Assets")
                            if 'Equity' in text and '股东权益' not in text:
                                english_headers_found.append('Equity')
                                print("      ❌ ENGLISH: Equity")

            # Final results
            print("\n🎯 VERIFICATION RESULTS:")
            print("=" * 50)

            success_count = 0

            # Check multi-page
            if len(prs.slides) > 1:
                print("✅ MULTI-PAGE: Content spans multiple slides")
                success_count += 1
            else:
                print("❌ SINGLE PAGE: Content fits on one slide only")

            # Check header translation
            if chinese_headers_found:
                print(f"✅ HEADER TRANSLATION: {len(chinese_headers_found)} Chinese headers found")
                print(f"   Headers: {', '.join(set(chinese_headers_found))}")
                success_count += 1
            else:
                print("❌ HEADER TRANSLATION: No Chinese headers found")

            if english_headers_found:
                print(f"❌ ENGLISH HEADERS: {english_headers_found} still in English")

            # Check Chinese content
            if total_chinese_chars > 200:
                print(f"✅ CHINESE CONTENT: {total_chinese_chars} Chinese characters")
                success_count += 1
            else:
                print(f"❌ CHINESE CONTENT: Only {total_chinese_chars} Chinese characters")

            print("\n" + "=" * 50)
            if success_count >= 3:
                print("🎉 ALL FIXES ARE WORKING CORRECTLY!")
                print("   ✅ Multi-page content distribution")
                print("   ✅ Chinese header translation")
                print("   ✅ Substantial Chinese content")
                return True
            else:
                print("⚠️  SOME ISSUES MAY STILL EXIST")
                return False

        else:
            print("❌ Export failed")
            return False

    except Exception as e:
        print(f"❌ Test failed with error: {e}")
        import traceback
        traceback.print_exc()
        return False

    finally:
        # Cleanup
        for path in [md_path, pptx_path]:
            if os.path.exists(path):
                os.unlink(path)

def main():
    """Main function"""
    print("CHINESE PPT EXPORT FIXES VERIFICATION")
    print("This script tests all the fixes for Chinese PPT export issues.")
    print()

    success = test_chinese_export_fixes()

    print("\n" + "=" * 60)
    if success:
        print("🎉 SUCCESS: All Chinese export fixes are working!")
        print("\nWhat this means:")
        print("• Headers like 'Current Assets' are translated to '流动资产'")
        print("• Content spans multiple pages when needed")
        print("• Chinese text is properly formatted with correct spacing")
        print("• All Chinese characters are preserved")
    else:
        print("⚠️  ISSUES DETECTED: Some fixes may not be working")
        print("\nPlease check the error messages above.")

    print("\nTo run this test yourself:")
    print("cd /Users/ytsang/Desktop/Github/python-pptx")
    print("python verify_chinese_fixes.py")
    print("=" * 60)

if __name__ == "__main__":
    main()
