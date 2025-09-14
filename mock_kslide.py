"""
Mock KSlide Implementation
This simulates KSlide functionality for testing PowerPoint generation
"""

import os
import datetime
from typing import Dict, List, Any, Optional
from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

class MockPresentation:
    """Mock KSlide Presentation class"""
    
    def __init__(self):
        self.slides = []
        self.slide_layouts = {
            'title': 0,
            'content': 1,
            'two_content': 2,
            'comparison': 3,
            'title_only': 5,
            'blank': 6
        }
    
    def add_slide(self, layout_type='content'):
        """Add a new slide to the presentation"""
        slide = MockSlide(layout_type)
        self.slides.append(slide)
        return slide
    
    def save(self, filename):
        """Save the presentation to a file"""
        # Create a real PowerPoint presentation using python-pptx with template
        template_path = "fdd_utils/template.pptx"
        if os.path.exists(template_path):
            prs = PptxPresentation(template_path)
        else:
            prs = PptxPresentation()
        
        for slide_data in self.slides:
            if slide_data.layout_type == 'title':
                slide_layout = prs.slide_layouts[0]  # Title slide
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                subtitle = slide.placeholders[1]
                
                if slide_data.title:
                    title.text = slide_data.title
                if slide_data.subtitle:
                    subtitle.text = slide_data.subtitle
                    
            elif slide_data.layout_type == 'content':
                slide_layout = prs.slide_layouts[1]  # Title and content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                content = slide.placeholders[1]
                
                if slide_data.title:
                    title.text = slide_data.title
                if slide_data.content:
                    content.text = slide_data.content
                    
            elif slide_data.layout_type == 'two_content':
                slide_layout = prs.slide_layouts[2]  # Two content
                slide = prs.slides.add_slide(slide_layout)
                title = slide.shapes.title
                left_content = slide.placeholders[1]
                right_content = slide.placeholders[2]
                
                if slide_data.title:
                    title.text = slide_data.title
                if slide_data.left_content:
                    left_content.text = slide_data.left_content
                if slide_data.right_content:
                    right_content.text = slide_data.right_content
        
        # Ensure output directory exists
        output_dir = os.path.dirname(filename)
        if output_dir:  # Only create directory if there's a directory path
            os.makedirs(output_dir, exist_ok=True)
        prs.save(filename)
        return True  # Return True to indicate success

class MockSlide:
    """Mock KSlide Slide class"""
    
    def __init__(self, layout_type='content'):
        self.layout_type = layout_type
        self.title = ""
        self.subtitle = ""
        self.content = ""
        self.left_content = ""
        self.right_content = ""
        self.shapes = []
    
    def add_textbox(self, x, y, width, height, text="", font_size=18):
        """Add a text box to the slide"""
        textbox = MockTextBox(x, y, width, height, text, font_size)
        self.shapes.append(textbox)
        return textbox
    
    def add_table(self, x, y, width, height, rows, cols):
        """Add a table to the slide"""
        table = MockTable(x, y, width, height, rows, cols)
        self.shapes.append(table)
        return table

class MockTextBox:
    """Mock KSlide TextBox class"""
    
    def __init__(self, x, y, width, height, text="", font_size=18):
        self.x = x
        self.y = y
        self.width = width
        self.height = height
        self.text = text
        self.font_size = font_size
        self.alignment = 'left'
        self.bold = False
        self.italic = False
        self.color = (0, 0, 0)  # RGB black
    
    def set_text(self, text):
        """Set the text content"""
        self.text = text
    
    def set_font_size(self, size):
        """Set the font size"""
        self.font_size = size
    
    def set_alignment(self, alignment):
        """Set text alignment"""
        self.alignment = alignment
    
    def set_bold(self, bold=True):
        """Set bold formatting"""
        self.bold = bold
    
    def set_italic(self, italic=True):
        """Set italic formatting"""
        self.italic = italic
    
    def set_color(self, r, g, b):
        """Set text color"""
        self.color = (r, g, b)

class MockTable:
    """Mock KSlide Table class"""
    
    def __init__(self, x, y, width, height, rows, cols):
        self.x = x
        self.y = y
        self.width = width
        self.height = height
        self.rows = rows
        self.cols = cols
        self.data = [["" for _ in range(cols)] for _ in range(rows)]
    
    def set_cell(self, row, col, value):
        """Set cell value"""
        if 0 <= row < self.rows and 0 <= col < self.cols:
            self.data[row][col] = str(value)
    
    def set_header_row(self, row, values):
        """Set header row values"""
        for col, value in enumerate(values):
            if col < self.cols:
                self.set_cell(row, col, value)

# Mock KSlide functions
def Presentation():
    """Create a new presentation"""
    return MockPresentation()

def Slide(presentation, layout_type='content'):
    """Add a slide to presentation"""
    return presentation.add_slide(layout_type)

def TextBox(slide, x, y, width, height, text="", font_size=18):
    """Add a text box to slide"""
    return slide.add_textbox(x, y, width, height, text, font_size)

def Table(slide, x, y, width, height, rows, cols):
    """Add a table to slide"""
    return slide.add_table(x, y, width, height, rows, cols)

# Mock Chart class (placeholder)
class MockChart:
    def __init__(self, *args, **kwargs):
        pass

def Chart(*args, **kwargs):
    """Mock chart function"""
    return MockChart(*args, **kwargs)

# Export the mock classes and functions
__all__ = [
    'Presentation', 'Slide', 'TextBox', 'Table', 'Chart',
    'MockPresentation', 'MockSlide', 'MockTextBox', 'MockTable', 'MockChart'
]
