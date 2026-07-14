#!/usr/bin/env python3
"""Dump every shape's name/type/position/size from a PPTX template, one
line per shape, across every slide.

fdd_utils/template.pptx is gitignored (*template*.pptx in .gitignore) --
it's a per-machine local binary that never syncs via git, so two machines'
copies can silently drift apart (e.g. one has textMainBullets resized to
4.6in and a new "Table Placeholder" shape, the other still has the
original dimensions). Run this on EACH machine and diff the output instead
of describing shapes verbally back and forth -- that's what caused the
drift to go unnoticed in the first place.

Usage:
    python inspect_template.py [path/to/template.pptx]

Defaults to fdd_utils/template.pptx if no path is given.
"""
import sys

from pptx import Presentation


def _emu_to_in(value):
    return round(value / 914400, 3) if value is not None else None


def main() -> int:
    path = sys.argv[1] if len(sys.argv) > 1 else "fdd_utils/template.pptx"
    try:
        prs = Presentation(path)
    except Exception as exc:
        print(f"Failed to open {path!r}: {exc}")
        return 1

    print(f"Template: {path}")
    print(f"Slide size: {_emu_to_in(prs.slide_width)}in x {_emu_to_in(prs.slide_height)}in")
    print(f"Slides: {len(prs.slides)}")

    for idx, slide in enumerate(prs.slides):
        print(f"\n=== Slide {idx} ===")
        for shape in slide.shapes:
            try:
                left, top = _emu_to_in(shape.left), _emu_to_in(shape.top)
                width, height = _emu_to_in(shape.width), _emu_to_in(shape.height)
            except Exception:
                left = top = width = height = None

            is_ph = shape.is_placeholder
            ph_type = ph_idx = None
            if is_ph:
                try:
                    ph_type = str(shape.placeholder_format.type)
                    ph_idx = shape.placeholder_format.idx
                except Exception:
                    pass

            print(
                f"  name={shape.name!r:35s} type={str(shape.shape_type):28s} "
                f"placeholder={is_ph!s:5s} ph_type={ph_type} ph_idx={ph_idx} "
                f"left={left} top={top} width={width} height={height}"
            )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
