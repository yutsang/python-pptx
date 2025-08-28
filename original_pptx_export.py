import os
import textwrap
import logging
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from dataclasses import dataclass
from typing import List, Tuple
from datetime import datetime
import re
from pptx.oxml.ns import qn  # Required import for XML namespace handling

logging.basicConfig(level=logging.INFO)

def get_tab_name(project_name):
    """Get tab name based on project name - more flexible approach"""
    try:
        # Hardcoded mappings for known entities
        if project_name == 'Haining':
            return "BSHN"
        elif project_name == 'Nanjing':
            return "BSNJ"
        elif project_name == 'Ningbo':
            return "BSNB"
        else:
            # Try to find a sheet that contains the project name
            return None
    except Exception:
        return None

def clean_content_quotes(content):
    """
    Remove outermost quotation marks from content while preserving legitimate internal quotes.
    Supports both straight quotes (") and curly quotes (" ").
    """
    if not content:
        return content
    
    # Handle straight quotes
    content = re.sub(r'^"([^"]*)"$', r'\1', content)
    
    # Handle curly quotes
    content = re.sub(r'^"([^"]*)"$', r'\1', content)
    content = re.sub(r'^"([^"]*)"$', r'\1', content)
    
    return content

@dataclass
class FinancialItem:
    accounting_type: str
    account_title: str
    descriptions: List[str]
    layer1_continued: bool = False
    layer2_continued: bool = False
    is_table: bool = False

class PowerPointGenerator:
    def __init__(self, template_path: str):
        self.prs = Presentation(template_path)
        self.current_slide_index = 0
        self.LINE_HEIGHT = Pt(12)
        self.ROWS_PER_SECTION = 30  # Use the same value for all sections
        
        # Find a slide with textMainBullets shape
        shape = None
        for slide in self.prs.slides:
            try:
                shape = next(s for s in slide.shapes if s.name == "textMainBullets")
                break
            except StopIteration:
                continue
        
        if not shape:
            # If no textMainBullets found, try to find any text shape
            for slide in self.prs.slides:
                for s in slide.shapes:
                    if hasattr(s, 'text_frame'):
                        shape = s
                        break
                if shape:
                    break
        
        if not shape:
            raise ValueError("No suitable text shape found in template")
        
        self.CHARS_PER_ROW = 50
        self.BULLET_CHAR = '■ '
        self.DARK_BLUE = RGBColor(0, 50, 150)
        self.DARK_GREY = RGBColor(169, 169, 169)
        self.prev_layer1 = None
        self.prev_layer2 = None

    def log_template_shapes(self):
        # Removed shape audit logging to reduce debug output
        pass

    def _calculate_max_rows_for_shape(self, shape):
        """Calculate the actual maximum number of rows that can fit in a shape"""
        if not shape or not hasattr(shape, 'height'):
            return 25  # Default fallback
        
        # Get the actual shape height in EMU (English Metric Units)
        shape_height_emu = shape.height
        
        # Convert EMU to points (1 EMU = 1/914400 inches, 1 inch = 72 points)
        shape_height_pt = shape_height_emu * 72 / 914400
        
        # Account for margins and padding - moderate space for shorter summary
        effective_height_pt = shape_height_pt * 0.85  # Moderate height utilization for shorter summary
        
        # Calculate line height based on font size and line spacing
        # Use dynamic font size detection or default to 10pt for summary shapes
        font_size_pt = 10  # Use 10pt for summary shapes to match actual usage
        line_spacing = 1.15  # Comfortable spacing for shorter, more readable summary
        line_height_pt = font_size_pt * line_spacing
        
        # Calculate maximum rows that can fit
        max_rows = int(effective_height_pt / line_height_pt)
        
        # Use moderate space for shorter, more readable summary
        max_rows = max(20, max_rows)  # Minimum 20 rows for compact summary
        
        return max_rows

    def _calculate_max_rows(self):
        # Find a slide with textMainBullets shape
        shape = None
        for slide in self.prs.slides:
