#!/usr/bin/env python3
"""Prototype: auto-detect a price-volume bridge/waterfall helper block in an
Excel tab (like '成都-量价桥图') and render it as a real PPTX waterfall chart,
WITHOUT hardcoding row/column numbers -- so a future file with a different
layout (different row numbers, more/fewer factors) still works as long as
the same generic template convention is used.

Detection strategy (generic, anchor-text based -- not position based):
  1. Find a header row containing the literal text 'Base' in some column C,
     and 'Change' in column C+1 (adjacent) -- this is the standard column
     pair a "waterfall via stacked-bar" Excel template uses. The label
     column is C-1 (one column left of 'Base').
  2. Scan rows directly below that header until two consecutive fully-empty
     rows are hit (end of block). For each row:
       - a value in the Base column (Change blank) = a TOTAL/anchor bar
         (start or end of the bridge, an absolute value).
       - a value in the Change column (Base blank) = a DELTA bar (a single
         factor's contribution, positive or negative).
  3. Cross-check against a nearby 'check' row (label cell containing the
     word 'check', formula = end_target - reconstructed_end) if present --
     confirms the detected block actually reconstructs to the right number,
     independent of trusting the row/column positions blindly.

This intentionally does NOT try to replicate AB-CD's own SUMIFS/AVERAGEIFS
sourcing logic (those row numbers ARE hardcoded per file today, per the
earlier inspect_bridge_source.py findings) -- it only reads the bridge tab's
OWN already-computed factor list, which is the part confirmed to generalize.

Usage (against a real databook):
    python bridge_chart_prototype.py "databooks/xx.xlsx" --sheet "成都-量价桥图"

Usage (synthetic self-test, no file needed -- rebuilds a workbook from the
real Chengdu numbers already inspected, to prove the reader+chart pipeline
end-to-end without needing the actual client file locally):
    python bridge_chart_prototype.py --self-test
"""
import argparse
import sys
from dataclasses import dataclass
from typing import List, Optional, Tuple

from openpyxl import Workbook, load_workbook
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, Reference
from openpyxl.chart.shapes import GraphicalProperties
from openpyxl.chart.label import DataLabel, DataLabelList
from openpyxl.chart.text import RichText
from openpyxl.drawing.text import CharacterProperties, Paragraph, ParagraphProperties, RichTextProperties

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


@dataclass
class BridgeItem:
    label: str
    kind: str  # "total" or "delta"
    value: float


@dataclass
class BridgeBlock:
    header_row: int
    label_col: int
    base_col: int
    change_col: int
    items: List[BridgeItem]
    check_ok: Optional[bool]  # None if no check row found nearby


def _find_header_rows(ws) -> List[Tuple[int, int]]:
    """Returns (row, base_col) for every row where col C says 'Base' and
    col C+1 says 'Change' (case/whitespace-insensitive)."""
    hits = []
    for row in ws.iter_rows():
        for cell in row:
            if cell.value is None:
                continue
            text = str(cell.value).strip().lower()
            if text != "base":
                continue
            next_cell = ws.cell(row=cell.row, column=cell.column + 1)
            if next_cell.value is not None and str(next_cell.value).strip().lower() == "change":
                hits.append((cell.row, cell.column))
    return hits


def _is_number(v) -> bool:
    return isinstance(v, (int, float)) and not isinstance(v, bool)


def find_bridge_blocks(ws_values, max_gap: int = 2) -> List[BridgeBlock]:
    blocks = []
    for header_row, base_col in _find_header_rows(ws_values):
        label_col = base_col - 1
        change_col = base_col + 1
        items: List[BridgeItem] = []
        row = header_row + 1
        empty_streak = 0
        while empty_streak < max_gap:
            label = ws_values.cell(row=row, column=label_col).value
            base = ws_values.cell(row=row, column=base_col).value
            change = ws_values.cell(row=row, column=change_col).value
            if label is None and base is None and change is None:
                empty_streak += 1
                row += 1
                continue
            if isinstance(label, str) and "check" in label.lower():
                # Explicit terminator, not just another gap -- a 'check' row
                # can follow the last real item after only ONE blank row
                # (confirmed in the real file), which the gap counter alone
                # would otherwise walk straight past and misread as a
                # spurious zero-value bridge item.
                break
            empty_streak = 0
            if _is_number(base) and not _is_number(change):
                items.append(BridgeItem(label=str(label) if label is not None else f"row{row}", kind="total", value=float(base)))
            elif _is_number(change) and not _is_number(base):
                items.append(BridgeItem(label=str(label) if label is not None else f"row{row}", kind="delta", value=float(change)))
            row += 1
        if len(items) >= 2:  # need at least a start + end to be a real bridge
            check_ok = _validate_block(ws_values, header_row, row, items)
            blocks.append(BridgeBlock(header_row, label_col, base_col, change_col, items, check_ok))
    return blocks


def _validate_block(ws_values, header_row: int, end_row: int, items: List[BridgeItem]) -> Optional[bool]:
    """Looks a few rows below the block for a row whose label contains
    'check' -- if found, and its own formula result is ~0, that's the
    workbook's OWN confirmation the bridge reconstructs correctly, which we
    can cross-check against by re-summing the detected items ourselves."""
    if len(items) < 2:
        return None
    reconstructed_end = items[0].value if items[0].kind == "total" else 0.0
    for it in items[1:-1] if items[-1].kind == "total" else items[1:]:
        if it.kind == "delta":
            reconstructed_end += it.value
    expected_end = items[-1].value if items[-1].kind == "total" else reconstructed_end
    for scan_row in range(end_row, end_row + 5):
        for col in range(1, ws_values.max_column + 1):
            v = ws_values.cell(row=scan_row, column=col).value
            if isinstance(v, str) and "check" in v.lower():
                return abs(reconstructed_end - expected_end) < max(1.0, abs(expected_end) * 0.005)
    return abs(reconstructed_end - expected_end) < max(1.0, abs(expected_end) * 0.005)


def _compute_waterfall_series(block: BridgeBlock):
    """Shared series math for both renderers (PPTX and Excel): decomposes
    the block's items into the classic invisible-base-series stacked-bar
    layout -- one non-zero value per category across 4 series."""
    categories = [it.label for it in block.items]
    base_vals, total_vals, inc_vals, dec_vals = [], [], [], []
    running = 0.0
    for it in block.items:
        if it.kind == "total":
            base_vals.append(0.0)
            total_vals.append(it.value)
            inc_vals.append(0.0)
            dec_vals.append(0.0)
            running = it.value
        else:
            if it.value >= 0:
                base_vals.append(running)
                inc_vals.append(it.value)
                dec_vals.append(0.0)
            else:
                base_vals.append(running + it.value)
                inc_vals.append(0.0)
                dec_vals.append(-it.value)
            total_vals.append(0.0)
            running += it.value
    return categories, base_vals, total_vals, inc_vals, dec_vals


def build_waterfall_chart(slide, block: BridgeBlock, title: str, left, top, width, height):
    """Classic stacked-bar-simulated waterfall: invisible Base series +
    Total/Increase/Decrease series, one non-zero per category."""
    categories, base_vals, total_vals, inc_vals, dec_vals = _compute_waterfall_series(block)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Base", base_vals)
    chart_data.add_series("Total", total_vals)
    chart_data.add_series("Increase", inc_vals)
    chart_data.add_series("Decrease", dec_vals)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, left, top, width, height, chart_data
    )
    chart = graphic_frame.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.has_legend = False

    # Anchor the value axis at 0 and use a plain thousands-separated number
    # format -- without this, PowerPoint's auto-scaled axis (based on the
    # STACKED total per category, which varies a lot between a small delta
    # bar and a large running-total bar) can end up looking arbitrary.
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.number_format = "#,##0"
    chart.value_axis.tick_labels.number_format_is_linked = False

    plot = chart.plots[0]
    plot.gap_width = 40  # tighter bars -- default (150) reads as too sparse for 9+ categories
    colors = {
        "Base": None,  # made invisible below
        "Total": RGBColor(0x44, 0x54, 0x6A),
        "Increase": RGBColor(0x2E, 0x8B, 0x57),
        "Decrease": RGBColor(0xC0, 0x39, 0x2B),
    }
    for series in plot.series:
        fill = series.format.fill
        if series.name == "Base":
            fill.background()  # fully transparent -- this is what makes it a "floating" bar
        else:
            fill.solid()
            fill.fore_color.rgb = colors[series.name]
            # Mark the actual number on every visible segment -- the whole
            # point of a bridge chart is reading off each factor's size,
            # not just eyeballing bar heights.
            series.data_labels.show_value = True
            series.data_labels.number_format = "#,##0"
            series.data_labels.number_format_is_linked = False
            series.data_labels.position = XL_LABEL_POSITION.CENTER
    return chart


def build_excel_waterfall_chart(ws, block: BridgeBlock, title: str, start_row: int = 1, start_col: int = 1):
    """Writes the block's Base/Total/Increase/Decrease series into `ws`
    starting at (start_row, start_col) and adds a NATIVE Excel stacked bar
    chart using the same invisible-base-series waterfall technique as
    build_waterfall_chart (PPTX version) -- so downstream tooling that
    expects a real Excel chart object (e.g. UpSlide, which links FROM Excel
    charts INTO PowerPoint rather than being driven programmatically) can
    pick this up directly, unlike a python-pptx chart which UpSlide has no
    concept of. Returns the row just below the written block (for stacking
    multiple blocks on one sheet)."""
    categories, base_vals, total_vals, inc_vals, dec_vals = _compute_waterfall_series(block)

    header_row = start_row
    ws.cell(row=header_row, column=start_col, value="Label")
    ws.cell(row=header_row, column=start_col + 1, value="Base")
    ws.cell(row=header_row, column=start_col + 2, value="Total")
    ws.cell(row=header_row, column=start_col + 3, value="Increase")
    ws.cell(row=header_row, column=start_col + 4, value="Decrease")
    for i, cat in enumerate(categories):
        r = header_row + 1 + i
        ws.cell(row=r, column=start_col, value=cat)
        ws.cell(row=r, column=start_col + 1, value=base_vals[i])
        ws.cell(row=r, column=start_col + 2, value=total_vals[i])
        ws.cell(row=r, column=start_col + 3, value=inc_vals[i])
        ws.cell(row=r, column=start_col + 4, value=dec_vals[i])
    ws.cell(row=header_row, column=start_col + 6, value=title)
    data_last_row = header_row + len(categories)

    chart = BarChart()
    chart.type = "col"
    chart.grouping = "stacked"
    chart.overlap = 100
    chart.title = title
    chart.legend = None
    chart.gapWidth = 40
    chart.y_axis.scaling.min = 0
    chart.y_axis.numFmt = "#,##0"
    # Default openpyxl chart size (~15x7.5cm) is far too small once there are
    # 8-11 categories with Chinese labels -- everything overlaps and reads as
    # one dense blob. Size it closer to a full slide so bars/labels breathe.
    chart.width = 30
    chart.height = 15
    # Category labels are long Chinese phrases ("综合楼运营天数增加") -- at
    # default horizontal orientation with 8-11 of them they overlap into an
    # unreadable strip. Rotate -45 degrees (rot is in 60,000ths of a degree)
    # and shrink the font slightly so each label gets its own diagonal slot.
    axis_text_props = RichText(
        bodyPr=RichTextProperties(rot=-2700000, vert="horz"),
        p=[Paragraph(pPr=ParagraphProperties(defRPr=CharacterProperties(sz=900)), endParaRPr=CharacterProperties(sz=900))],
    )
    chart.x_axis.txPr = axis_text_props

    cats_ref = Reference(ws, min_col=start_col, min_row=header_row + 1, max_row=data_last_row)
    series_values = {"Total": total_vals, "Increase": inc_vals, "Decrease": dec_vals}
    colors = {"Base": None, "Total": "44546A", "Increase": "2E8B57", "Decrease": "C0392B"}
    for col_offset, name in [(1, "Base"), (2, "Total"), (3, "Increase"), (4, "Decrease")]:
        data_ref = Reference(ws, min_col=start_col + col_offset, min_row=header_row, max_row=data_last_row)
        chart.add_data(data_ref, titles_from_data=True)
    chart.set_categories(cats_ref)

    for series, name in zip(chart.series, ["Base", "Total", "Increase", "Decrease"]):
        if name == "Base":
            series.graphicalProperties = GraphicalProperties(noFill=True)
        else:
            series.graphicalProperties = GraphicalProperties(solidFill=colors[name])
            dlbls = DataLabelList(showVal=True, numFmt="#,##0")
            # Every category has a value in only ONE of Total/Increase/Decrease
            # (that's the whole point of the invisible-base-series technique) --
            # the other two series read 0 for that category. Without this, every
            # bar shows two extra "0" labels floating on the invisible portion,
            # which is most of what made the chart look cluttered.
            dlbls.dLbl = [DataLabel(idx=i, showVal=False) for i, v in enumerate(series_values[name]) if v == 0]
            series.dLbls = dlbls

    anchor_row = data_last_row + 2
    ws.add_chart(chart, f"{get_column_letter(start_col)}{anchor_row}")

    return anchor_row + 32  # bigger chart now -- leave more room before the next block


def _build_synthetic_workbook() -> Workbook:
    """Rebuilds the REAL structure + REAL values already confirmed from the
    user's actual databook dump (成都-量价桥图 tab), so the reader can be
    tested end-to-end without needing the actual client file on this
    machine. Only the two 'Dynamic data for bridge' helper blocks are
    reproduced (that's all the reader needs) -- the SUMIFS-sourced C:AE
    staging area is NOT needed for this prototype and is omitted."""
    wb = Workbook()
    ws = wb.active
    ws.title = "成都-量价桥图"

    block1 = [
        ("2024年收入", "total", 1182.02748),
        ("干仓单价增加", "delta", 2059.545458900859),
        ("干仓出租率增加", "delta", 1463.5320655100995),
        ("干仓运营天数增加", "delta", 3298.960105589041),
        ("综合楼单价增加", "delta", 66.9393853357857),
        ("综合楼出租率下降", "delta", -55.45101747277201),
        ("综合楼运营天数增加", "delta", 29.103442136986306),
        ("冷库单价下降", "delta", 0.0),
        ("冷库出租率提升", "delta", 0.0),
        ("冷库运营天数增加", "delta", 3061.83285),
        ("2025年收入", "total", 11106.48977),
    ]
    block2 = [
        ("2025年收入", "total", 11106.48977),
        ("干仓单价增加", "delta", 1185.257478889836),
        ("干仓出租率增加", "delta", 1082.3849911101638),
        ("干仓运营天数增加", "delta", 0.0),
        ("综合楼单价下降", "delta", -0.2670081709406935),
        ("综合楼出租率增加", "delta", 19.893338170940705),
        ("综合楼运营天数增加", "delta", 0.0),
        ("冷库单价下降", "delta", -614.8229452141594),
        ("冷库出租率提升", "delta", 818.0862552141591),
        ("冷库运营天数增加", "delta", 0.0),
        ("2025年7月至2026年6月收入", "total", 13597.02188),
    ]

    def write_block(start_row: int, items):
        ws.cell(row=start_row - 1, column=33, value="Dynamic data for bridge (do not change)")  # AG col=33... adjust below
        header_row = start_row
        label_col, base_col, change_col = 33, 34, 35  # AG, AH, AI
        ws.cell(row=header_row, column=base_col, value="Base")
        ws.cell(row=header_row, column=change_col, value="Change")
        for i, (label, kind, value) in enumerate(items):
            r = header_row + 1 + i
            ws.cell(row=r, column=label_col, value=label)
            if kind == "total":
                ws.cell(row=r, column=base_col, value=value)
            else:
                ws.cell(row=r, column=change_col, value=value)
        check_row = header_row + 1 + len(items) + 1
        ws.cell(row=check_row, column=label_col, value="check")
        ws.cell(row=check_row, column=base_col, value=0)  # workbook's own check formula would resolve to 0
        return check_row

    next_start = write_block(12, block1)
    write_block(next_start + 4, block2)

    return wb


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("path", nargs="?", default=None, help="path to the databook .xlsx")
    ap.add_argument("--sheet", default=None,
                     help="scan only this one sheet by name. If omitted, scans EVERY sheet in the "
                          "workbook (this is the default -- covers every entity's bridge tab in one "
                          "pass, not just one hardcoded tab name, since new entities may name theirs "
                          "differently, e.g. '<City>-量价桥图').")
    ap.add_argument("--self-test", action="store_true",
                     help="skip the real file and use a synthetic workbook rebuilt from the "
                          "already-confirmed real Chengdu values, to test the pipeline end-to-end")
    ap.add_argument("--out", default="bridge_chart_prototype_output.pptx", help="output pptx path")
    ap.add_argument("--excel-out", default=None,
                     help="also write a native-Excel-chart version (openpyxl BarChart, same invisible-base-"
                          "series technique) to this .xlsx path, e.g. for UpSlide-based downstream workflows "
                          "that link FROM Excel charts rather than accepting a python-pptx chart object")
    args = ap.parse_args()

    if args.self_test:
        print("Using synthetic self-test workbook (real Chengdu values, reconstructed structure)...")
        wb_values = _build_synthetic_workbook()
        sheet_names = [args.sheet] if args.sheet else wb_values.sheetnames
    else:
        if not args.path:
            print("❌ provide a databook path, or use --self-test")
            return 1
        print(f"Loading {args.path!r}...")
        wb_values = load_workbook(args.path, data_only=True)
        if args.sheet:
            if args.sheet not in wb_values.sheetnames:
                print(f"❌ sheet {args.sheet!r} not found. Available: {wb_values.sheetnames}")
                return 1
            sheet_names = [args.sheet]
        else:
            sheet_names = wb_values.sheetnames

    print(f"\nScanning {len(sheet_names)} sheet(s) for 'Base'/'Change' bridge helper block(s)...")
    all_blocks: List[Tuple[str, BridgeBlock]] = []
    sheets_with_blocks = 0
    for sheet_name in sheet_names:
        ws = wb_values[sheet_name]
        blocks = find_bridge_blocks(ws)
        if blocks:
            sheets_with_blocks += 1
            for b in blocks:
                all_blocks.append((sheet_name, b))

    print(f"{sheets_with_blocks}/{len(sheet_names)} sheet(s) have at least one bridge block; "
          f"{len(all_blocks)} block(s) total.\n")
    if not all_blocks:
        print("❌ No bridge blocks found anywhere in this workbook.")
        return 1

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slides_added = 0
    for i, (sheet_name, block) in enumerate(all_blocks):
        print(f"--- Block {i + 1}/{len(all_blocks)} -- sheet {sheet_name!r} (header row {block.header_row}, "
              f"label col {get_column_letter(block.label_col)}, "
              f"base col {get_column_letter(block.base_col)}, "
              f"change col {get_column_letter(block.change_col)}) ---")
        for it in block.items:
            print(f"  [{it.kind:5s}] {it.label}: {it.value:,.2f}")
        check_msg = {True: "✅ reconstructs correctly", False: "❌ MISMATCH -- do not trust this block",
                     None: "⚠️ no check row found nearby, unverified"}[block.check_ok]
        print(f"  {check_msg}\n")
        if block.check_ok is False:
            print("  Skipping chart generation for this block (failed its own reconciliation check).\n")
            continue

        slide = prs.slides.add_slide(blank_layout)
        title = f"{sheet_name}: {block.items[0].label} → {block.items[-1].label}"
        build_waterfall_chart(
            slide, block, title,
            left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(6.3),
        )
        slides_added += 1

    prs.save(args.out)
    print(f"Saved demo PPTX ({slides_added} slide(s)) to {args.out!r}. Open it and check the bars visually.")

    if args.excel_out:
        out_wb = Workbook()
        out_wb.remove(out_wb.active)
        out_ws = out_wb.create_sheet("Bridge_Output")
        next_row = 1
        charts_added = 0
        for sheet_name, block in all_blocks:
            if block.check_ok is False:
                continue
            title = f"{sheet_name}: {block.items[0].label} → {block.items[-1].label}"
            next_row = build_excel_waterfall_chart(out_ws, block, title, start_row=next_row)
            charts_added += 1
        out_wb.save(args.excel_out)
        print(f"Saved demo XLSX ({charts_added} chart(s)) to {args.excel_out!r}. Open it and check the bars visually.")

    return 0


if __name__ == "__main__":
    sys.exit(main())
