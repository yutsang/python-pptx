"""
Clean, modular FDD Application
Main application file - properly modularized and lightweight
"""

import streamlit as st
import pandas as pd
import json
import warnings
import os
import datetime
import time
from pathlib import Path
import threading
import tempfile
import uuid

# Suppress warnings and bytecode generation
os.environ['PYTHONDONTWRITEBYTECODE'] = '1'
import urllib3
urllib3.disable_warnings()
warnings.filterwarnings('ignore', category=UserWarning, module='openpyxl')
warnings.simplefilter(action='ignore', category=UserWarning)

# Import all required modules
from fdd_utils.mappings import KEY_TO_SECTION_MAPPING, KEY_TERMS_BY_KEY
from fdd_utils.category_config import DISPLAY_NAME_MAPPING_DEFAULT, DISPLAY_NAME_MAPPING_NB_NJ
from fdd_utils.excel_processing import get_worksheet_sections_by_keys
from fdd_utils.data_utils import (
    get_tab_name, get_financial_keys, get_key_display_name,
    format_date_to_dd_mmm_yyyy, load_config_files
)
from fdd_utils.content_utils import (
    clean_content_quotes, get_content_from_json,
    generate_content_from_session_storage
)
from common.ui import configure_streamlit_page
from common.ui_sections import (
    render_balance_sheet_sections, render_income_statement_sections,
    render_combined_sections
)
from common.pptx_export import export_pptx, merge_presentations

from common.assistant import (
    process_keys, QualityAssuranceAgent, DataValidationAgent,
    PatternValidationAgent, find_financial_figures_with_context_check,
    get_financial_figure, ProofreadingAgent
)

# Import logging
import logging
logging.getLogger('streamlit.watcher.event_based_path_watcher').setLevel(logging.ERROR)
logging.getLogger('streamlit.watcher.util').setLevel(logging.ERROR)


class MockUploadedFile:
    """Mock uploaded file object for default file handling"""
    def __init__(self, file_path):
        self.name = file_path
        self.file_path = file_path
        self._file = None
    
    def read(self, size=-1):
        if self._file is None:
            self._file = open(self.file_path, 'rb')
        return self._file.read(size)
    
    def getbuffer(self):
        with open(self.file_path, 'rb') as f:
            return f.read()
    
    def getvalue(self):
        return self.getbuffer()
    
    def seek(self, offset, whence=0):
        if self._file is None:
            self._file = open(self.file_path, 'rb')
        return self._file.seek(offset, whence)
    
    def tell(self):
        if self._file is None:
            return 0
        return self._file.tell()
    
    def seekable(self):
        return True
    
    def close(self):
        if self._file:
            self._file.close()
            self._file = None


def initialize_session_state():
    """Initialize session state variables"""
    if 'processing_started' not in st.session_state:
        st.session_state['processing_started'] = False
    
    if 'agent_states' not in st.session_state:
        st.session_state['agent_states'] = {
            'agent1_completed': False,
            'agent2_completed': False, 
            'agent3_completed': False,
            'agent1_results': {},
            'agent2_results': {},
            'agent3_results': {},
            'agent1_success': False,
            'agent2_success': False,
            'agent3_success': False
        }
    
    if 'ai_logger' not in st.session_state:
        from fdd_utils.app_helpers import AIAgentLogger
        st.session_state.ai_logger = AIAgentLogger()


def get_available_ai_providers(config):
    """Get available AI providers from config"""
    providers = []
    
    if not config:
        return ["Offline"]  # Fallback if no config
    
    # Check Local AI
    if config.get('LOCAL_AI_ENABLED') and config.get('LOCAL_AI_API_BASE'):
        providers.append("Local AI")
    
    # Check OpenAI
    if config.get('OPENAI_API_KEY') and config.get('OPENAI_API_KEY') != "placeholder-openai-api-key":
        providers.append("Open AI")
    
    # Check DeepSeek
    if config.get('DEEPSEEK_API_KEY') and config.get('DEEPSEEK_API_KEY') != "placeholder-deepseek-api-key":
        providers.append("DeepSeek")
    
    # Always include Offline as fallback
    providers.append("Offline")
    
    # Use default provider from config if available
    default_provider = config.get('DEFAULT_AI_PROVIDER')
    if default_provider and default_provider in providers:
        # Move default to front
        providers.remove(default_provider)
        providers.insert(0, default_provider)
    
    return providers


def detect_entity_mode_automatically(uploaded_file, selected_entity, entity_keywords):
    """Automatically detect if the Excel file contains single or multiple entities"""
    try:
        if not uploaded_file:
            # If no file uploaded, check default file
            if os.path.exists('databook.xlsx'):
                file_to_check = 'databook.xlsx'
            else:
                return 'single'  # Default to single if no file
        else:
            file_to_check = uploaded_file
        
        # Read Excel file and check for multiple entities
        xl = pd.ExcelFile(file_to_check)
        
        # Known entity patterns to look for (including the selected entity)
        entity_patterns = [
            'ningbo wanchen', 'haining wanpu', 'nanjing jingya',
            '宁波万晨', '海宁万普', '南京京亚'
        ]
        
        # Add the selected entity and its keywords to the patterns
        if selected_entity:
            entity_patterns.append(selected_entity.lower())
        if entity_keywords:
            entity_patterns.extend([kw.lower() for kw in entity_keywords])
        
        entities_found = set()
        selected_entity_found = False
        sheets_checked = 0
        
        # Check first few sheets for entity patterns
        for sheet_name in xl.sheet_names[:8]:  # Check first 8 sheets
            try:
                df = xl.parse(sheet_name)
                if df.empty:
                    continue
                    
                sheets_checked += 1
                
                # Convert all data to text for searching
                all_text = ""
                for row_idx in range(min(25, len(df))):  # Check first 25 rows
                    for col_idx in range(len(df.columns)):
                        try:
                            cell_value = str(df.iloc[row_idx, col_idx]).lower()
                            if cell_value and cell_value != 'nan':
                                all_text += cell_value + " "
                        except:
                            continue
                
                # Look for entity patterns
                sheet_entities = set()
                for pattern in entity_patterns:
                    if pattern in all_text:
                        sheet_entities.add(pattern)
                        entities_found.add(pattern)
                        
                        # Check if this is the selected entity
                        if (selected_entity and pattern.lower() == selected_entity.lower()) or \
                           (entity_keywords and pattern in [kw.lower() for kw in entity_keywords]):
                            selected_entity_found = True
                
                if sheet_entities:
                    #print(f"🔍 AUTO-DETECT: Sheet '{sheet_name}' contains entities: {sheet_entities}")
                    pass    
            except Exception as e:
                print(f"Error checking sheet {sheet_name}: {e}")
                continue
        
        # Decision logic
        unique_entities = len(entities_found)
        #print(f"🔍 AUTO-DETECT: Total unique entities found: {unique_entities}")
        print(f"🔍 AUTO-DETECT: Entities found: {entities_found}")
        print(f"🔍 AUTO-DETECT: Selected entity found: {selected_entity_found}")
        
        if unique_entities > 1:
            print(f"🔍 AUTO-DETECT: Multiple entities detected -> MULTIPLE mode")
            return 'multiple'
        elif unique_entities == 1 and selected_entity_found:
            print(f"🔍 AUTO-DETECT: Only selected entity found -> SINGLE mode")
            return 'single'
        elif unique_entities >= 1:
            print(f"🔍 AUTO-DETECT: Other entities found but not selected entity -> MULTIPLE mode (for filtering)")
            return 'multiple'
        else:
            print(f"🔍 AUTO-DETECT: No specific entities detected -> SINGLE mode (default)")
            return 'single'
                
    except Exception as e:
        print(f"Error in auto-detection: {e}")
        return 'single'  # Default to single on error


def generate_entity_keywords(entity_input):
    """Generate comprehensive entity keywords from input"""
    if not entity_input:
        return [], None, []
    
    words = entity_input.split()
    entity_keywords = [words[0]] if words else []
    
    # Generate combinations
    if len(words) >= 2:
        for i in range(1, len(words)):
            entity_keywords.append(f"{words[0]} {words[i]}")
    
    # Generate suffixes
    if len(words) > 1:
        entity_suffixes = [s.strip() for s in " ".join(words[1:]).split(',') if s.strip()]
    else:
        entity_suffixes = ["Limited"]
    
    return entity_keywords, entity_input, entity_suffixes


def detect_language_from_data(sections_by_key):
    """Auto-detect language from 'Indicative adjusted' vs '示意性调整后' columns in Excel data"""
    chinese_indicators = ['示意性调整后', '示意性調整後']
    english_indicators = ['Indicative adjusted']
    
    chinese_count = 0
    english_count = 0
    
    # Silently scan all data for language indicators
    for key, sections in sections_by_key.items():
        if not sections:
            continue
            
        for section in sections:
            if 'parsed_data' in section and section['parsed_data']:
                metadata = section['parsed_data'].get('metadata', {})
                table_name = metadata.get('table_name', '')
                
                # Check table name for language indicators
                if any(indicator in table_name for indicator in english_indicators):
                    english_count += 1
                elif any(indicator in table_name for indicator in chinese_indicators):
                    chinese_count += 1
                
                # Check the raw Excel data for language indicators
                raw_data = section.get('data', None)
                if raw_data is not None:
                    # Convert DataFrame to list of lists for processing
                    if hasattr(raw_data, 'values'):
                        data_rows = raw_data.values.tolist()
                    elif isinstance(raw_data, list):
                        data_rows = raw_data
                    else:
                        data_rows = []
                    
                    for row_idx, row in enumerate(data_rows[:10]):  # Check first 10 rows of raw data
                        if isinstance(row, list):
                            for cell_idx, cell in enumerate(row):
                                if isinstance(cell, str):
                                    # Check for the specific "Indicative adjusted" vs "示意性调整后" indicators
                                    if "Indicative adjusted" in cell:
                                        english_count += 1
                                    elif "示意性调整后" in cell or "示意性調整後" in cell:
                                        chinese_count += 1
                                    # Also check for other language indicators
                                    elif any(indicator in cell for indicator in english_indicators):
                                        english_count += 1
                                    elif any(indicator in cell for indicator in chinese_indicators):
                                        chinese_count += 1
                
                # Also check parsed data content for language indicators
                data_rows = section['parsed_data'].get('data', [])
                
                for row_idx, row in enumerate(data_rows[:5]):  # Check first 5 rows
                    if isinstance(row, list):
                        for cell_idx, cell in enumerate(row):
                                if isinstance(cell, str):
                                    if any(indicator in cell for indicator in english_indicators):
                                        english_count += 1
                                    elif any(indicator in cell for indicator in chinese_indicators):
                                        chinese_count += 1
    
    print(f"🌏 LANGUAGE DETECTION: Chinese: {chinese_count}, English: {english_count}")
    
    # If no indicators found, try to detect from any Chinese characters in the data
    if chinese_count == 0 and english_count == 0:
        print("🌏 LANGUAGE DETECTION: No specific indicators found, checking for Chinese characters...")
        chinese_char_count = 0
        total_char_count = 0
        
        for key, sections in sections_by_key.items():
            if not sections:
                continue
                
            for section in sections:
                # Check raw data first (stored in 'data' field)
                raw_data = section.get('data', None)
                if raw_data is not None:
                    # Convert DataFrame to list of lists for processing
                    if hasattr(raw_data, 'values'):
                        data_rows = raw_data.values.tolist()
                    elif isinstance(raw_data, list):
                        data_rows = raw_data
                    else:
                        data_rows = []
                    
                    for row in data_rows[:10]:  # Check first 10 rows
                        if isinstance(row, list):
                            for cell in row:
                                if isinstance(cell, str):
                                    total_char_count += len(cell)
                                    # Count Chinese characters (CJK Unified Ideographs)
                                    for char in cell:
                                        if '\u4e00' <= char <= '\u9fff':
                                            chinese_char_count += 1
                
                # Also check parsed data
                if 'parsed_data' in section and section['parsed_data']:
                    data_rows = section['parsed_data'].get('data', [])
                    for row in data_rows[:10]:  # Check first 10 rows
                        if isinstance(row, list):
                            for cell in row:
                                if isinstance(cell, str):
                                    total_char_count += len(cell)
                                    # Count Chinese characters (CJK Unified Ideographs)
                                    for char in cell:
                                        if '\u4e00' <= char <= '\u9fff':
                                            chinese_char_count += 1
        
        if total_char_count > 0:
            chinese_ratio = chinese_char_count / total_char_count
            print(f"🌏 LANGUAGE DETECTION: Chinese character ratio: {chinese_ratio:.2f} ({chinese_char_count}/{total_char_count})")
            
            if chinese_ratio > 0.1:  # If more than 10% Chinese characters
                detected_language = 'chinese'
                print(f"🌏 LANGUAGE DETECTED: Chinese (based on character ratio: {chinese_ratio:.2f})")
            else:
                detected_language = 'english'
                print(f"🌏 LANGUAGE DETECTED: English (based on character ratio: {chinese_ratio:.2f})")
        else:
            # Default to English if no characters found
            detected_language = 'english'
            print("🌏 LANGUAGE DETECTED: English (default - no characters found)")
    else:
        # Determine language based on counts - prioritize Chinese if any Chinese indicators found
        if chinese_count > 0:
            detected_language = 'chinese'
            print(f"🌏 LANGUAGE DETECTED: Chinese ({chinese_count} indicators found)")
        elif english_count > 0:
            detected_language = 'english'
            print(f"🌏 LANGUAGE DETECTED: English ({english_count} indicators found)")
        else:
            # This shouldn't happen if the logic above is correct, but just in case
            detected_language = 'english'
            print(f"🌏 LANGUAGE DETECTED: English (default fallback)")
    
    return detected_language


def process_excel_with_timeout(uploaded_file, mapping, selected_entity, entity_suffixes, entity_keywords, entity_mode, timeout=30):
    """Process Excel file with timeout protection"""
    result_container = {}
    exception_container = {}

    def excel_worker():
        try:
            result = get_worksheet_sections_by_keys(
                uploaded_file=uploaded_file,
                tab_name_mapping=mapping,
                entity_name=selected_entity,
                entity_suffixes=entity_suffixes,
                entity_keywords=entity_keywords,
                entity_mode=entity_mode,
                debug=True
            )
            result_container['result'] = result
        except Exception as e:
            exception_container['exception'] = e

    processing_thread = threading.Thread(target=excel_worker)
    processing_thread.daemon = True
    processing_thread.start()
    processing_thread.join(timeout=timeout)

    if processing_thread.is_alive():
        return None, "timeout"
    elif 'exception' in exception_container:
        raise exception_container['exception']
    else:
        return result_container['result'], "success"


def run_ai_processing(filtered_keys, ai_data, language='english', progress_callback=None):
    """Run AI processing for content generation"""
    # AI processing started
    try:
        # Create temporary file for processing
        temp_file_path = None
        if 'uploaded_file_data' in st.session_state:
            unique_filename = f"databook_{uuid.uuid4().hex[:8]}.xlsx"
            temp_file_path = os.path.join(tempfile.gettempdir(), unique_filename)
            with open(temp_file_path, 'wb') as tmp_file:
                tmp_file.write(st.session_state['uploaded_file_data'])
        elif os.path.exists('databook.xlsx'):
            temp_file_path = 'databook.xlsx'
        
        if not temp_file_path:
            st.error("❌ No databook available for processing")
            return {}

        # Get entity information
        entity_name = ai_data.get('entity_name', '')
        entity_keywords = ai_data.get('entity_keywords', [])

        # Process keys using the assistant
        results = process_keys(
            keys=filtered_keys,
            entity_name=entity_name,
            entity_helpers=entity_keywords,
            input_file=temp_file_path,
            mapping_file="fdd_utils/mapping.json",
            pattern_file="fdd_utils/pattern.json",
            config_file='fdd_utils/config.json',
            prompts_file='fdd_utils/prompts.json',
            use_ai=True,
            processed_table_data=ai_data.get('sections_by_key', {}),
            use_local_ai=st.session_state.get('use_local_ai', False),
            use_openai=st.session_state.get('use_openai', False),
            language=language,
            progress_callback=progress_callback
        )

        return results or {}
        
    except Exception as e:
        st.error(f"❌ AI processing failed: {e}")
        return {}


def run_simple_chinese_translation(english_results, ai_data, progress_callback=None):
    """Simple Chinese translation function that fixes the translator issues"""
    try:
        from common.assistant import initialize_ai_services, generate_response, get_chat_model, load_config
        
        # Load AI configuration
        config = load_config('fdd_utils/config.json')
        oai_client, _ = initialize_ai_services(config, 
                                             use_local=st.session_state.get('use_local_ai', False),
                                             use_openai=st.session_state.get('use_openai', False))
        
        model_name = get_chat_model(config)
        entity_name = ai_data.get('entity_name', '')
        
        # Chinese translation prompt
        system_prompt = """您是中国财务报告翻译专家。您的任务是将英文财务分析内容完整翻译成简体中文。

重要要求：
1. 必须将所有英文内容翻译成简体中文
2. 保留所有数字、百分比、货币符号和技术术语（如VAT、CIT、WHT、Surtax、IPO）不变
3. 保持专业的财务报告语气和格式结构
4. 确保最终输出100%是中文内容，除了上述保留的数字和技术术语
5. 不要添加任何解释或额外文本
6. 翻译必须准确、专业，适合中国财务报告使用
7. 禁止在翻译结果中保留任何英文句子或短语
8. 直接返回翻译后的完整中文内容"""
        
        translated_results = {}
        
        for key, result in english_results.items():
            if progress_callback:
                progress_callback(0.1, f"Translating {key}")
            
            content = result.get('content', str(result)) if isinstance(result, dict) else str(result)
            
            if content:
                user_prompt = f"""请将以下英文财务分析内容翻译成简体中文：

{content}

请提供准确的中文翻译，保持专业财务报告的语气和格式。"""
                
                try:
                    translated_content = generate_response(
                        user_query=user_prompt,
                        system_prompt=system_prompt,
                        oai_client=oai_client,
                        context_content="",
                        openai_chat_model=model_name,
                        entity_name=entity_name,
                        use_local_ai=st.session_state.get('use_local_ai', False)
                    )
                    
                    translated_results[key] = {
                        'content': translated_content,
                        'translated_content': translated_content,
                        'is_chinese': True,
                        'original_english': content
                    }
                    
                except Exception as e:
                    print(f"Translation error for {key}: {e}")
                    # Fallback: return original content with error note
                    translated_results[key] = {
                        'content': f"[翻译失败] {content}",
                        'translated_content': f"[翻译失败] {content}",
                        'is_chinese': False,
                        'original_english': content,
                        'error': str(e)
                    }
            else:
                translated_results[key] = {
                    'content': '',
                    'translated_content': '',
                    'is_chinese': True,
                    'original_english': ''
                }
        
        return translated_results
        
    except Exception as e:
        st.error(f"❌ Chinese translation failed: {e}")
        return {}


def run_simple_proofreader(english_results, ai_data, progress_callback=None):
    """Enhanced proofreader function that provides detailed feedback"""
    try:
        from common.assistant import ProofreadingAgent, load_config
        
        # Load AI configuration
        config = load_config('fdd_utils/config.json')
        entity_name = ai_data.get('entity_name', '')
        
        # Create proofreader agent
        proofreader = ProofreadingAgent(
            use_local_ai=st.session_state.get('use_local_ai', False),
            use_openai=st.session_state.get('use_openai', False),
            language='English'
        )
        
        proofread_results = {}
        
        for key, result in english_results.items():
            if progress_callback:
                progress_callback(0.1, f"Proofreading {key}")
            
            content = result.get('content', str(result)) if isinstance(result, dict) else str(result)
            
            if content:
                try:
                    # Get tables markdown for context
                    tables_markdown = ""
                    if 'sections_by_key' in ai_data and key in ai_data['sections_by_key']:
                        tables_data = ai_data['sections_by_key'][key]
                        # Ensure tables_markdown is a string, not a list
                        if isinstance(tables_data, list):
                            tables_markdown = "\n".join(str(item) for item in tables_data)
                        else:
                            tables_markdown = str(tables_data)
                    
                    # Use the full proofreader implementation
                    proofread_result = proofreader.proofread(
                        content=content,
                        key=key,
                        tables_markdown=tables_markdown,
                        entity=entity_name
                    )
                    
                    
                    # Ensure issues is always a list - with robust handling
                    issues = proofread_result.get('issues', [])
                    if not isinstance(issues, list):
                        print(f"⚠️ WARNING: Issues field is not a list in proofread result, converting from {type(issues)}")
                        if isinstance(issues, str):
                            # If it's a string, try to split it if it looks like multiple issues
                            issues_str = issues.strip()
                            if ',' in issues_str or '\n' in issues_str:
                                # Try to split by comma or newline
                                issues = [issue.strip() for issue in issues_str.replace('\n', ',').split(',') if issue.strip()]
                            else:
                                issues = [issues_str] if issues_str else []
                        elif isinstance(issues, (int, float)):
                            # If it's a number, convert to string first
                            issues = [str(issues)]
                        else:
                            # For any other type, try to convert to string and wrap in list
                            try:
                                issues = [str(issues)] if issues else []
                            except:
                                issues = []
                    
                    # Ensure all list fields are properly handled
                    def ensure_list_field(field_value, field_name):
                        if not isinstance(field_value, list):
                            print(f"⚠️ WARNING: {field_name} field is not a list, converting from {type(field_value)}")
                            if isinstance(field_value, str):
                                field_value_str = field_value.strip()
                                if ',' in field_value_str or '\n' in field_value_str:
                                    field_value = [item.strip() for item in field_value_str.replace('\n', ',').split(',') if item.strip()]
                                else:
                                    field_value = [field_value_str] if field_value_str else []
                            elif isinstance(field_value, (int, float)):
                                field_value = [str(field_value)]
                            else:
                                try:
                                    field_value = [str(field_value)] if field_value else []
                                except:
                                    field_value = []
                        return field_value
                    
                    figure_checks = ensure_list_field(proofread_result.get('figure_checks', []), 'figure_checks')
                    entity_checks = ensure_list_field(proofread_result.get('entity_checks', []), 'entity_checks')
                    grammar_notes = ensure_list_field(proofread_result.get('grammar_notes', []), 'grammar_notes')
                    
                    proofread_results[key] = {
                        'content': proofread_result.get('corrected_content', content),
                        'original_content': content,
                        'is_compliant': proofread_result.get('is_compliant', True),
                        'issues': issues,
                        'figure_checks': figure_checks,
                        'entity_checks': entity_checks,
                        'grammar_notes': grammar_notes,
                        'pattern_used': proofread_result.get('pattern_used', ''),
                        'translation_runs': proofread_result.get('translation_runs', 0)
                    }
                    
                except Exception as e:
                    print(f"Proofreading error for {key}: {e}")
                    # Fallback: return original content
                    proofread_results[key] = {
                        'content': content,
                        'original_content': content,
                        'error': str(e),
                        'is_compliant': False,
                        'issues': [f"Proofreading failed: {str(e)}"],
                        'figure_checks': [],
                        'entity_checks': [],
                        'grammar_notes': [],
                        'pattern_used': '',
                        'translation_runs': 0
                    }
            else:
                proofread_results[key] = {
                    'content': '',
                    'original_content': '',
                    'is_compliant': True,
                    'issues': [],
                    'figure_checks': [],
                    'entity_checks': [],
                    'grammar_notes': [],
                    'pattern_used': '',
                    'translation_runs': 0
                }
        
        return proofread_results
        
    except Exception as e:
        st.error(f"❌ Proofreading failed: {e}")
        return english_results  # Return original results if proofreading fails



def cleanup_old_files():
    """Clean up any leftover temporary files from previous runs"""
    try:
        temp_files = [
            "fdd_utils/bs_content_temp.md", "fdd_utils/is_content_temp.md",
            "fdd_utils/bs_content_backup.md", "fdd_utils/is_content_backup.md"
        ]

        for temp_file in temp_files:
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                    print(f"🧹 Cleaned up leftover temp file: {temp_file}")
                except:
                    pass

        # Clean up very old output files (keep only last 10)
        try:
            output_dir = "fdd_utils/output"
            if os.path.exists(output_dir):
                import glob
                pptx_files = glob.glob(f"{output_dir}/*.pptx")
                if len(pptx_files) > 10:
                    # Sort by creation time (oldest first)
                    pptx_files.sort(key=os.path.getctime)
                    # Remove oldest files
                    files_to_remove = pptx_files[:-10]
                    for old_file in files_to_remove:
                        try:
                            os.remove(old_file)
                            print(f"🧹 Cleaned up old output file: {old_file}")
                        except:
                            pass
        except:
            pass
    except:
        pass

def main():
    """Main application function"""
    # Clean up old files on startup
    cleanup_old_files()

    # Configure Streamlit
    configure_streamlit_page()
    st.title("🏢 Real Estate DD Report Writer")
    
    # Initialize session state
    initialize_session_state()

    # Navigation description
    if not st.session_state.get('processing_started', False):
        st.info("📋 **Welcome!** Please navigate to the left sidebar to upload your databook and configure input data.")

    # Sidebar controls
    with st.sidebar:
        # File uploader
        uploaded_file = st.file_uploader(
            "Upload Excel File (Optional)",
            type=['xlsx', 'xls'],
            help="Upload your financial data Excel file or use the default databook.xlsx"
        )

        # Use default file if none uploaded
        if uploaded_file is None:
            default_file_path = "databook.xlsx"
            if os.path.exists(default_file_path):
                st.caption(f"Using default file: {default_file_path}")
                uploaded_file = MockUploadedFile(default_file_path)
            else:
                st.error("❌ Default file not found: databook.xlsx")
                st.info("Please upload an Excel file to get started.")
                st.stop()
        
        # Store uploaded file in session state
        st.session_state['uploaded_file'] = uploaded_file
        
        # Check if file has changed and clear AI content if so
        current_file_name = uploaded_file.name if hasattr(uploaded_file, 'name') else 'default'
        previous_file_name = st.session_state.get('previous_uploaded_file_name', '')
        
        if current_file_name != previous_file_name and previous_file_name != '':
            # File has changed, clear AI content and entity cache
            if 'ai_content_store' in st.session_state:
                del st.session_state['ai_content_store']
            if 'ai_data' in st.session_state:
                del st.session_state['ai_data']
            if 'uploaded_file_data' in st.session_state:
                del st.session_state['uploaded_file_data']
            # Clear any cached entity names and related data
            keys_to_clear = [k for k in st.session_state.keys() if any(term in k.lower() for term in ['entity', 'entities', 'databook'])]
            for key in keys_to_clear:
                del st.session_state[key]

            # Clean up old content files when switching databooks
            old_content_files = [
                "fdd_utils/bs_content.md", "fdd_utils/is_content.md",
                "fdd_utils/bs_content_backup.md", "fdd_utils/is_content_backup.md",
                "fdd_utils/bs_content_temp.md", "fdd_utils/is_content_temp.md",
                "fdd_utils/output/"  # Clean up old output files too
            ]

            for file_path in old_content_files:
                if file_path.endswith('/'):
                    # Clean up old files in output directory
                    try:
                        import glob
                        pattern = file_path + "*.pptx"
                        old_files = glob.glob(pattern)
                        for old_file in old_files:
                            try:
                                os.remove(old_file)
                                print(f"🧹 Cleaned up old output file: {old_file}")
                            except:
                                pass
                    except:
                        pass
                else:
                    # Clean up specific content files
                    try:
                        if os.path.exists(file_path):
                            os.remove(file_path)
                            print(f"🧹 Cleaned up old content file: {file_path}")
                    except:
                        pass

            # Don't show reminder message - just silently reset
            print(f"🔄 New databook uploaded: {current_file_name} - cleaned up old files")
            st.session_state['previous_uploaded_file_name'] = current_file_name
        elif previous_file_name == '':
            st.session_state['previous_uploaded_file_name'] = current_file_name

        # Entity input with dropdown
        from fdd_utils.data_utils import extract_entity_names_from_databook

        # Get entity names from the uploaded file
        uploaded_file_path = None
        if uploaded_file and hasattr(uploaded_file, 'name'):
            uploaded_file_path = uploaded_file.name

        entity_names = extract_entity_names_from_databook(uploaded_file_path or 'databook.xlsx')
        
        # Multi-entity mode: Use dropdown with option to type custom entity
        if entity_names:
            # Create selectbox with custom option
            entity_options = entity_names + ["[Custom Entity]"]
            
            selected_entity_option = st.selectbox(
                "Select Entity Name", 
                options=entity_options,
                help="Select entity from databook or choose '[Custom Entity]' to type your own"
            )
            
            if selected_entity_option == "[Custom Entity]":
                entity_input = st.text_input(
                    "Enter Custom Entity Name",
                    value="",
                    placeholder="e.g., Company Name Limited, Entity Name Corp",
                    help="Enter the full entity name to configure processing"
                )
            else:
                entity_input = selected_entity_option
        else:
            # Fallback to text input if no entities found
            entity_input = st.text_input(
                "Enter Entity Name",
                value="",
                placeholder="e.g., Company Name Limited, Entity Name Corp",
                help="Enter the full entity name to configure processing"
            )
        
        # Row limit is hardcoded in settings
        row_limit = 20  # Fixed limit for all slides
        
        # Clear session state when entity changes
        if 'last_entity_input' in st.session_state:
            if st.session_state['last_entity_input'] != entity_input:
                if 'ai_data' in st.session_state:
                    del st.session_state['ai_data']
                if 'filtered_keys_for_ai' in st.session_state:
                    del st.session_state['filtered_keys_for_ai']
                if 'processing_started' in st.session_state:
                    del st.session_state['processing_started']
        
        st.session_state['last_entity_input'] = entity_input
        
        # Generate entity configuration first
        entity_keywords, selected_entity, entity_suffixes = generate_entity_keywords(entity_input)
        
        if not selected_entity:
            st.warning("⚠️ Please enter an entity name to start processing")
            st.stop()

        # Automatic entity mode detection
        st.markdown("---")
        entity_mode_internal = detect_entity_mode_automatically(uploaded_file, selected_entity, entity_keywords)
        st.session_state['entity_mode'] = entity_mode_internal

        # Show entity info
        # Statement type selection
        statement_type_display = st.radio(
            "Financial Statement Type",
            ["Balance Sheet", "Income Statement", "All"],
            help="Select the type of financial statement to process"
        )
        
        statement_type_mapping = {
            "Balance Sheet": "BS",
            "Income Statement": "IS", 
            "All": "ALL"
        }
        statement_type = statement_type_mapping[statement_type_display]
        
        # Store statement type in session state for excel processing
        st.session_state['current_statement_type'] = statement_type
        
        # Financial Data Selection (moved here from PowerPoint section)
        st.markdown("**Financial Data Source:**")
        
        # Get available Excel sheets for selection
        uploaded_file = st.session_state.get('uploaded_file')
        
        if uploaded_file:
            try:
                import pandas as pd
                if hasattr(uploaded_file, 'file_path'):
                    file_path = uploaded_file.file_path
                else:
                    # For uploaded files, we need to save them temporarily
                    import tempfile
                    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
                    temp_file.write(uploaded_file.getvalue())
                    temp_file.close()
                    file_path = temp_file.name
                
                # Read Excel file to get sheet names
                excel_file = pd.ExcelFile(file_path)
                available_sheets = excel_file.sheet_names
            except Exception as e:
                print(f"Error reading Excel sheets: {e}")
                available_sheets = ["BS", "IS", "BSHN", "Cash", "AR", "AP"]
        else:
            # Default sheets if no file uploaded
            available_sheets = ["BS", "IS", "BSHN", "Cash", "AR", "AP"]
        
        # Show dropdown(s) based on statement type selection
        if statement_type == "BS":
            financial_statement_tab = st.selectbox(
                "Select Excel tab for Balance Sheet data:",
                options=available_sheets,
                index=0 if available_sheets else 0,
                help="Choose which Excel sheet contains the Balance Sheet data"
            )
        elif statement_type == "IS":
            financial_statement_tab = st.selectbox(
                "Select Excel tab for Income Statement data:",
                options=available_sheets,
                index=0 if available_sheets else 0,
                help="Choose which Excel sheet contains the Income Statement data"
            )
        else:  # statement_type == "ALL"
            # For "All" mode, use single dropdown since scraping logic handles both BS and IS
            financial_statement_tab = st.selectbox(
                "Select Excel tab for Financial Data:",
                options=available_sheets,
                index=0 if available_sheets else 0,
                help="Choose Excel sheet containing both Balance Sheet and Income Statement data"
            )
            
        # AI Provider Selection - Load from config
        config, _, _, _ = load_config_files()
        ai_providers = get_available_ai_providers(config)
        
        mode_display = st.selectbox(
            "Select Mode", 
            ai_providers,
            index=0,
            help="Choose AI provider. Models are taken from fdd_utils/config.json"
        )
            
        # Show API configuration status
        if config:
            # AI provider configuration status (simplified)
            if mode_display == "Open AI":
                if config.get('OPENAI_API_KEY') and config.get('OPENAI_API_BASE'):
                    model = config.get('OPENAI_CHAT_MODEL', 'Not configured')
                    st.info(f"🤖 Model: {model}")
            elif mode_display == "DeepSeek":
                if config.get('DEEPSEEK_API_KEY') and config.get('DEEPSEEK_API_BASE'):
                    model = config.get('DEEPSEEK_CHAT_MODEL', 'Not configured')
                    st.info(f"🤖 Model: {model}")
            elif mode_display == "Local AI":
                if config.get('LOCAL_AI_API_BASE') and config.get('LOCAL_AI_ENABLED'):
                    model = config.get('LOCAL_AI_CHAT_MODEL', 'Not configured')
                    st.info(f"🏠 Model: {model}")

        # Store mode configuration
        st.session_state['selected_mode'] = f"AI Mode - {mode_display}"
        st.session_state['ai_model'] = mode_display
        st.session_state['selected_provider'] = mode_display
        st.session_state['use_local_ai'] = (mode_display == "Local AI")
        st.session_state['use_openai'] = (mode_display == "Open AI")
            
    # Main processing area
    if uploaded_file is not None:
        # Start processing button
        if not st.session_state.get('processing_started', False):
            st.markdown("### 🎯 Ready to Process")
            st.info("📋 Configuration loaded. Click 'Start Processing' to begin data analysis and AI processing.")

            col1, col2, col3 = st.columns([1, 2, 1])
            with col2:
                start_processing = st.button(
                    "🚀 Start Processing",
                    type="primary",
                    use_container_width=True,
                    key="btn_start_processing",
                    help="Begin data processing and AI analysis"
                )

            if start_processing:
                st.session_state['processing_started'] = True
                st.rerun()

            st.stop()

        # Load configuration
        config, mapping, pattern, prompts = load_config_files()

        
        # Process Excel data
        entity_changed = st.session_state.get('last_processed_entity') != selected_entity
        needs_processing = 'ai_data' not in st.session_state or 'sections_by_key' not in st.session_state['ai_data'] or entity_changed

        if needs_processing:
            with st.spinner("🔄 Processing Excel file..."):
                print(f"🔄 Processing Excel for {selected_entity}")
                start_time = time.time()

                sections_by_key, status = process_excel_with_timeout(
                    uploaded_file=uploaded_file,
                    mapping=mapping,
                    selected_entity=selected_entity,
                    entity_suffixes=entity_suffixes,
                    entity_keywords=entity_keywords,
                    entity_mode=entity_mode_internal,
                    timeout=30
                )

                if status == "timeout":
                    if st.button("⚠️ Continue Without Excel Data", key="continue_without_excel"):
                        st.warning("⚠️ Continuing without Excel data. Some features may be limited.")
                        sections_by_key = {}
                        st.session_state['excel_processing_skipped'] = True
                    else:
                        st.error("❌ Excel processing timed out. Click 'Continue Without Excel Data' to proceed.")
                        st.stop()

                processing_time = time.time() - start_time
                print(f"✅ Excel processed in {processing_time:.2f}s")
                

                # Auto-detect language from data
                detected_language = detect_language_from_data(sections_by_key)

                # Store processed data
                if 'ai_data' not in st.session_state:
                    st.session_state['ai_data'] = {}
                st.session_state['ai_data']['sections_by_key'] = sections_by_key
                st.session_state['ai_data']['entity_name'] = selected_entity
                st.session_state['ai_data']['entity_keywords'] = entity_keywords
                st.session_state['ai_data']['detected_language'] = detected_language
                st.session_state['last_processed_entity'] = selected_entity
        else:
            sections_by_key = st.session_state.get('ai_data', {}).get('sections_by_key', {})

        # Filter sections by statement type to avoid showing BS content when IS is selected
        # But skip filtering for "ALL" mode to show all data
        if statement_type in ["BS", "IS"] and statement_type != "ALL":
            # Get the original unfiltered data from session state
            original_sections_by_key = st.session_state.get('ai_data', {}).get('sections_by_key', {})

            # Define BS and IS keys (allow partial matches for Chinese databooks)
            bs_keywords = ["Cash", "AR", "Prepayments", "OR", "CA", "NCA", "IP", "NCA",
                          "AP", "payable", "OP", "Capital", "Reserve"]
            is_keywords = ["OI", "OC", "Tax", "GA", "Fin", "Loss", "Income", "operating", "LT DTA"]

            # Filter the original sections_by_key based on statement type
            filtered_sections_by_key = {}
            if statement_type == "BS":
                for key in original_sections_by_key:
                    if any(bs_kw.lower() in key.lower() for bs_kw in bs_keywords):
                        filtered_sections_by_key[key] = original_sections_by_key[key]
                print(f"🔍 BS: {len(original_sections_by_key)}→{len(filtered_sections_by_key)} keys")
                print(f"🔍 BS keys found: {list(filtered_sections_by_key.keys())}")
            elif statement_type == "IS":
                print(f"🔍 IS DEBUG: Original keys: {list(original_sections_by_key.keys())}")
                print(f"🔍 IS DEBUG: IS keywords: {is_keywords}")
                print(f"🔍 IS DEBUG: BS keywords: {bs_keywords}")

                for key in original_sections_by_key:
                    is_match = any(is_kw.lower() in key.lower() for is_kw in is_keywords)
                    bs_match = any(bs_kw.lower() in key.lower() for bs_kw in bs_keywords)
                    include_key = is_match or not bs_match

                    print(f"🔍 IS DEBUG: Key '{key}' - IS match: {is_match}, BS match: {bs_match}, include: {include_key}")

                    if include_key:
                        filtered_sections_by_key[key] = original_sections_by_key[key]

                print(f"🔍 IS: {len(original_sections_by_key)}→{len(filtered_sections_by_key)} keys")
                print(f"🔍 IS keys found: {list(filtered_sections_by_key.keys())}")

            sections_by_key = filtered_sections_by_key
        else:
            print(f"🔍 ALL: {len(sections_by_key)} keys")

        # Display financial statements
        
        if statement_type == "BS":
            render_balance_sheet_sections(
                sections_by_key, get_key_display_name, selected_entity, format_date_to_dd_mmm_yyyy
            )
        elif statement_type == "IS":
            render_income_statement_sections(
                sections_by_key, get_key_display_name, selected_entity, format_date_to_dd_mmm_yyyy
            )
        elif statement_type == "ALL":
            render_combined_sections(
                sections_by_key, get_key_display_name, selected_entity, format_date_to_dd_mmm_yyyy
            )

        # AI Processing Section
        st.markdown("---")
        st.markdown("## 🤖 AI Report Generation")
        
        # Prepare AI data
        keys_with_data = [key for key, sections in sections_by_key.items() if sections]

        print(f"📊 SUMMARY: {len(keys_with_data)} financial keys have data: {keys_with_data[:5]}{'...' if len(keys_with_data) > 5 else ''}")
        print(f"📊 DEBUG: All sections_by_key keys: {list(sections_by_key.keys())}")
        print(f"📊 DEBUG: Sections with data: {[(k, len(v)) for k, v in sections_by_key.items() if v]}")

        # Filter keys by statement type
        bs_keys = ["Cash", "AR", "Prepayments", "OR", "Other CA", "Other NCA", "IP", "NCA",
                   "AP", "Taxes payable", "OP", "Capital", "Reserve"]
        is_keys = ["OI", "OC", "Tax and Surcharges", "GA", "Fin Exp", "Cr Loss", "Other Income",
                   "Non-operating Income", "Non-operating Exp", "Income tax", "LT DTA"]

        print(f"📊 DEBUG: Expected BS keys: {bs_keys}")
        print(f"📊 DEBUG: Expected IS keys: {is_keys}")

        if statement_type == "BS":
            filtered_keys_for_ai = [key for key in keys_with_data if key in bs_keys]
            print(f"📊 {statement_type} MODE: Processing {len(filtered_keys_for_ai)} keys {filtered_keys_for_ai[:3]}{'...' if len(filtered_keys_for_ai) > 3 else ''}")
            if not filtered_keys_for_ai:
                filtered_keys_for_ai = keys_with_data
                print(f"⚠️ {statement_type} MODE: No specific keys found, using all {len(filtered_keys_for_ai)} available keys")
        elif statement_type == "IS":
            filtered_keys_for_ai = [key for key in keys_with_data if key in is_keys]
            print(f"📊 {statement_type} MODE: Processing {len(filtered_keys_for_ai)} keys {filtered_keys_for_ai[:3]}{'...' if len(filtered_keys_for_ai) > 3 else ''}")
            if not filtered_keys_for_ai:
                filtered_keys_for_ai = keys_with_data
                print(f"⚠️ {statement_type} MODE: No specific keys found, using all {len(filtered_keys_for_ai)} available keys")
        else:  # ALL
            filtered_keys_for_ai = keys_with_data
            print(f"📊 {statement_type} MODE: Processing {len(filtered_keys_for_ai)} keys for combined report")

        st.session_state['filtered_keys_for_ai'] = filtered_keys_for_ai
        
        # Store uploaded file data for AI processing
        if hasattr(uploaded_file, 'getbuffer'):
            st.session_state['uploaded_file_data'] = uploaded_file.getbuffer()
        elif hasattr(uploaded_file, 'getvalue'):
            st.session_state['uploaded_file_data'] = uploaded_file.getvalue()
                
        # Get detected language before preparing AI data
        detected_language = st.session_state.get('ai_data', {}).get('detected_language', 'english')
        print(f"🔍 DEBUG: Retrieved detected_language from session state: '{detected_language}'")
        
        # Prepare AI data (preserve the detected_language)
        temp_ai_data = {
            'entity_name': selected_entity,
            'entity_keywords': entity_keywords,
            'sections_by_key': sections_by_key,
            'pattern': pattern,
            'mapping': mapping,
            'config': config,
            'detected_language': detected_language  # Preserve the detected language
        }
        st.session_state['ai_data'] = temp_ai_data
        language_display = "🇨🇳 Chinese" if detected_language == 'chinese' else "🇺🇸 English"
        
        
        # Show detected language prominently
        if detected_language == 'chinese':
            st.info(f"🌏 **Language Detected**: Chinese databook - AI will generate content in Chinese")
        else:
            st.info(f"🌏 **Language Detected**: English databook - AI will generate content in English")
        
        # BSHN Sheet Options (default enabled)
        include_bshn = True  # Always include BSHN sheet by default
        
        # Simplified AI Processing Buttons
        
        col1, col2 = st.columns(2)

        with col1:
            generate_report_clicked = st.button(
                f"🚀 Generate & Download Report ({language_display})",
                type="primary",
                use_container_width=True,
                key="btn_generate_report",
                help=f"Generate AI content and download PowerPoint in {detected_language}"
            )

        with col2:
            # Download button will be populated here after generation
            pass

        # Store the column reference for later use
        st.session_state['download_button_column'] = col2

        
        # Check if AI processing has completed (use same logic as result display)
        agent_states = st.session_state.get('agent_states', {})
        ai_completed = any([
            agent_states.get('agent1_completed', False),
            agent_states.get('agent2_completed', False),
            agent_states.get('agent3_completed', False)
        ])

        # Check if PowerPoint file exists for download
        output_dir = "fdd_utils/output"
        pptx_file_exists = False
        latest_file = None

        if os.path.exists(output_dir):
            pptx_files = [f for f in os.listdir(output_dir) if f.endswith('.pptx')]
            if pptx_files:
                pptx_file_exists = True
                latest_file = max(pptx_files, key=lambda x: os.path.getctime(os.path.join(output_dir, x)))
                print(f"📥 DOWNLOAD: Found {len(pptx_files)} PowerPoint files, latest: {latest_file}")
        else:
            print(f"📥 DOWNLOAD: Output directory does not exist: {output_dir}")

        # Show download button that directly downloads the file
        if pptx_file_exists and ai_completed:
            file_path = os.path.join(output_dir, latest_file)
            with open(file_path, 'rb') as f:
                file_data = f.read()

            # Use the stored column reference for proper alignment
            download_col = st.session_state.get('download_button_column')
            if download_col:
                with download_col:
                    st.download_button(
                        label="📥 Download Full Report",
                        data=file_data,
                        file_name=latest_file,
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="btn_download_full_pptx",
                        help="Download the full PowerPoint report with AI content",
                        use_container_width=True
                    )
            else:
                st.download_button(
                    label="📥 Download Full Report",
                    data=file_data,
                    file_name=latest_file,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    key="btn_download_full_pptx",
                    help="Download the full PowerPoint report with AI content",
                    use_container_width=True
                )
        else:
            if not ai_completed:
                print(f"📥 DOWNLOAD: Waiting for AI processing to complete")
            elif not pptx_file_exists:
                print(f"📥 DOWNLOAD: No PowerPoint file found yet")
            download_col = st.session_state.get('download_button_column')
            if download_col:
                with download_col:
                    st.button(
                        "📥 Download Full Report",
                        type="secondary",
                        use_container_width=True,
                        key="btn_download_full_pptx",
                        help="Download the full PowerPoint report with AI content",
                        disabled=True
                    )
                    if not ai_completed:
                        st.info("💡 Complete AI processing first to enable full report download")
                    elif not pptx_file_exists:
                        st.info("💡 Generate a report first to enable download")
            else:
                st.button(
                    "📥 Download Full Report",
                    type="secondary",
                    use_container_width=True,
                    key="btn_download_full_pptx",
                    help="Download the full PowerPoint report with AI content",
                    disabled=True
                )
                if not ai_completed:
                    st.info("💡 Complete AI processing first to enable full report download")
                elif not pptx_file_exists:
                    st.info("💡 Generate a report first to enable download")

        # Handle combined AI processing and PowerPoint export
        if generate_report_clicked:
            progress_bar = st.progress(0)
            status_text = st.empty()
            
            try:
                if detected_language == 'chinese':
                    status_text.text("🤖 初始化中文AI处理...")
                    progress_bar.progress(0.1)
                    
                    # First generate English content, then proofread, then translate
                    total_keys = len(filtered_keys_for_ai)
                    status_text.text(f"📊 生成英文内容... (0/{total_keys} keys)")
                    progress_bar.progress(0.2)
                    
                    # Initialize timing for proper ETA calculation
                    if 'processing_start_time' not in st.session_state:
                        st.session_state['processing_start_time'] = time.time()
                    
                    # Track current key index for proper progress display
                    current_key_index = 0
                    current_key = "Processing"

                    def progress_callback_eng(p, msg):
                        nonlocal current_key_index, current_key

                        # Store progress info in session state for display
                        if 'debug_progress' not in st.session_state:
                            st.session_state['debug_progress'] = []
                        st.session_state['debug_progress'].append(f"p={p}, msg='{msg}'")

                        # Parse the detailed message from the AI processing
                        current_key = "Processing"  # Default fallback

                        if msg and isinstance(msg, str):
                            # Format: "🔄 Processing Cash • OpenAI • Key 1/9"
                            if 'Processing' in msg and '•' in msg:
                                parts = msg.split('•')
                                if len(parts) >= 3 and 'Key' in parts[2]:
                                    # Extract key index from format like "Key 1/9"
                                    try:
                                        key_info = parts[2].strip()
                                        if '/' in key_info:
                                            current_key_index = int(key_info.split('/')[0].replace('Key ', ''))
                                            current_key = parts[0].replace('🔄 Processing', '').strip()
                                    except:
                                        pass
                                elif len(parts) >= 1:
                                    key_part = parts[0].replace('🔄 Processing', '').strip()
                                    if key_part:
                                        current_key = key_part
                            # Format: "🔄 Processing Cash" (without bullet points)
                            elif 'Processing' in msg:
                                key_part = msg.replace('🔄 Processing', '').strip()
                                if key_part:
                                    current_key = key_part
                            # Format: Direct key name
                            elif len(msg.strip()) < 50 and not '•' in msg and not 'Processing' in msg:
                                current_key = msg.strip()
                            # Format: Check if it's just a key name without "Processing"
                            elif msg.strip() in filtered_keys_for_ai:
                                current_key = msg.strip()

                        # Calculate progress based on current key index
                        if current_key_index == 0:
                            # Fallback: estimate from progress value
                            current_key_index = max(1, int(p * total_keys))

                        # Calculate proper ETA based on actual processing time
                        if p > 0 and current_key_index > 0:
                            elapsed_time = time.time() - st.session_state['processing_start_time']
                            avg_time_per_key = elapsed_time / current_key_index
                            remaining_keys = max(0, total_keys - current_key_index)
                            if remaining_keys > 0:
                                eta_seconds = int(remaining_keys * avg_time_per_key)

                                if eta_seconds > 0:
                                    eta_minutes = eta_seconds // 60
                                    eta_seconds = eta_seconds % 60
                                    eta_text = f"⏱️ ETA: {eta_minutes}m {eta_seconds}s" if eta_minutes > 0 else f"⏱️ ETA: {eta_seconds}s"
                                else:
                                    eta_text = "⏱️ Almost done!"
                            else:
                                eta_text = "⏱️ Almost done!"
                        else:
                            eta_text = ""

                        # Enhanced status display with ETA on same line
                        status_display = f"📊 生成英文内容... ({current_key_index}/{total_keys} keys) - {current_key}"
                        if eta_text:
                            status_display += f" {eta_text}"
                        status_text.text(status_display)
                        progress_bar.progress(0.1 + p * 0.2)
                    
                        # Starting AI processing
                    english_results = run_ai_processing(filtered_keys_for_ai, temp_ai_data, language=detected_language, progress_callback=progress_callback_eng)

                    if not english_results:
                        print(f"❌ AI PROCESSING FAILED: No results generated for {len(filtered_keys_for_ai)} keys")
                        if detected_language == 'chinese':
                            st.error("❌ 中文内容生成失败")
                        else:
                            st.error("❌ 英文内容生成失败，无法进行中文翻译")
                        return
                    else:
                        print(f"✅ AI PROCESSING: Generated content for {len(english_results)} keys")
                    
                    # Proofread content
                    if detected_language == 'chinese':
                        status_text.text(f"🧐 校对中文内容... (0/{total_keys} keys)")
                    else:
                        status_text.text(f"🧐 校对英文内容... (0/{total_keys} keys)")
                    progress_bar.progress(0.3)
                    
                    def progress_callback_proof(p, msg):
                        nonlocal current_key_index, current_key

                        # Parse the detailed message from the proofreader
                        current_key = "Proofreading"  # Default fallback

                        if msg and isinstance(msg, str):
                            # Format: "🔄 Processing Cash • OpenAI • Key 1/9"
                            if 'Processing' in msg and '•' in msg:
                                parts = msg.split('•')
                                if len(parts) >= 1:
                                    key_part = parts[0].replace('🔄 Processing', '').strip()
                                    if key_part:
                                        current_key = key_part
                            # Format: "🔄 Processing Cash" (without bullet points)
                            elif 'Processing' in msg:
                                key_part = msg.replace('🔄 Processing', '').strip()
                                if key_part:
                                    current_key = key_part
                            # Format: Direct key name
                            elif len(msg.strip()) < 50 and not '•' in msg and not 'Processing' in msg:
                                current_key = msg.strip()
                            # Format: Check if it's just a key name without "Processing"
                            elif msg.strip() in filtered_keys_for_ai:
                                current_key = msg.strip()

                        key_index = int(p * total_keys) if p > 0 else 0
                        status_text.text(f"🧐 校对英文内容... ({key_index}/{total_keys} keys) - {current_key}")
                        progress_bar.progress(0.3 + p * 0.1)
                    
                    proofread_results = run_simple_proofreader(english_results, temp_ai_data, progress_callback=progress_callback_proof)
                    
                    # Then translate to Chinese
                    status_text.text(f"🌐 翻译为中文... (0/{total_keys} keys)")
                    progress_bar.progress(0.5)
                    
                    def progress_callback_trans(p, msg):
                        nonlocal current_key_index, current_key

                        # Store progress info in session state for display
                        if 'debug_progress' not in st.session_state:
                            st.session_state['debug_progress'] = []
                        st.session_state['debug_progress'].append(f"p={p}, msg='{msg}'")

                        # Parse the detailed message from the AI processing
                        current_key = "Translating"  # Default fallback

                        if msg and isinstance(msg, str):
                            # Format: "🔄 Processing Cash • OpenAI • Key 1/9"
                            if 'Processing' in msg and '•' in msg:
                                parts = msg.split('•')
                                if len(parts) >= 3 and 'Key' in parts[2]:
                                    # Extract key index from format like "Key 1/9"
                                    try:
                                        key_info = parts[2].strip()
                                        if '/' in key_info:
                                            current_key_index = int(key_info.split('/')[0].replace('Key ', ''))
                                            current_key = parts[0].replace('🔄 Processing', '').strip()
                                    except:
                                        pass
                                elif len(parts) >= 1:
                                    key_part = parts[0].replace('🔄 Processing', '').strip()
                                    if key_part:
                                        current_key = key_part
                            # Format: "🔄 Processing Cash" (without bullet points)
                            elif 'Processing' in msg:
                                key_part = msg.replace('🔄 Processing', '').strip()
                                if key_part:
                                    current_key = key_part
                            # Format: Direct key name
                            elif len(msg.strip()) < 50 and not '•' in msg and not 'Processing' in msg:
                                current_key = msg.strip()
                            # Format: Check if it's just a key name without "Processing"
                            elif msg.strip() in filtered_keys_for_ai:
                                current_key = msg.strip()

                        # Calculate progress based on current key index
                        if current_key_index == 0:
                            # Fallback: estimate from progress value
                            current_key_index = max(1, int(p * total_keys))

                        # Calculate proper ETA based on actual processing time
                        if p > 0 and current_key_index > 0:
                            elapsed_time = time.time() - st.session_state['processing_start_time']
                            avg_time_per_key = elapsed_time / current_key_index
                            remaining_keys = max(0, total_keys - current_key_index)
                            if remaining_keys > 0:
                                eta_seconds = int(remaining_keys * avg_time_per_key)

                                if eta_seconds > 0:
                                    eta_minutes = eta_seconds // 60
                                    eta_seconds = eta_seconds % 60
                                    eta_text = f"⏱️ ETA: {eta_minutes}m {eta_seconds}s" if eta_minutes > 0 else f"⏱️ ETA: {eta_seconds}s"
                                else:
                                    eta_text = "⏱️ Almost done!"
                            else:
                                eta_text = "⏱️ Almost done!"
                        else:
                            eta_text = ""

                        # Enhanced status display with ETA on same line
                        status_display = f"🌐 翻译为中文... ({current_key_index}/{total_keys} keys) - {current_key}"
                        if eta_text:
                            status_display += f" {eta_text}"
                        status_text.text(status_display)
                        progress_bar.progress(0.5 + p * 0.3)
                    
                    final_results = run_simple_chinese_translation(proofread_results, temp_ai_data, progress_callback=progress_callback_trans)
                    
                    # Store results
                    if 'ai_content_store' not in st.session_state:
                        st.session_state['ai_content_store'] = {}

                    for key, result in final_results.items():
                        if key not in st.session_state['ai_content_store']:
                            st.session_state['ai_content_store'][key] = {}
                        content = result.get('content', str(result)) if isinstance(result, dict) else str(result)
                        st.session_state['ai_content_store'][key]['agent3_content'] = content
                        st.session_state['ai_content_store'][key]['current_content'] = content
                        st.session_state['ai_content_store'][key]['agent3_timestamp'] = time.time()
                        
                        # Store proofread content for preview (store the full result object, not just content)
                        if key in proofread_results:
                            proofread_result = proofread_results[key]
                            st.session_state['ai_content_store'][key]['agent2_content'] = proofread_result
                        

                    st.session_state['agent_states']['agent3_results'] = final_results
                    st.session_state['agent_states']['agent3_completed'] = True
                    st.session_state['agent_states']['agent3_success'] = True
                    
                else:
                    total_keys = len(filtered_keys_for_ai)
                    status_text.text(f"🤖 Generating English content... (0/{total_keys} keys)")
                    progress_bar.progress(0.3)
                    
                    # Initialize timing for proper ETA calculation
                    if 'processing_start_time' not in st.session_state:
                        st.session_state['processing_start_time'] = time.time()

                    # Track current key index for proper progress display
                    current_key_index = 0
                    current_key = "Processing"

                    def progress_callback_eng_simple(p, msg):
                        nonlocal current_key_index, current_key

                        # Parse the detailed message from the AI processing
                        current_key = "Processing"  # Default fallback

                        if msg and isinstance(msg, str):
                            # Format: "🔄 Processing Cash • OpenAI • Key 1/9"
                            if 'Processing' in msg and '•' in msg:
                                parts = msg.split('•')
                                if len(parts) >= 1:
                                    key_part = parts[0].replace('🔄 Processing', '').strip()
                                    if key_part:
                                        current_key = key_part
                            # Format: "🔄 Processing Cash" (without bullet points)
                            elif 'Processing' in msg:
                                key_part = msg.replace('🔄 Processing', '').strip()
                                if key_part:
                                    current_key = key_part
                            # Format: Direct key name
                            elif len(msg.strip()) < 50 and not '•' in msg and not 'Processing' in msg:
                                current_key = msg.strip()
                            # Format: Check if it's just a key name without "Processing"
                            elif msg.strip() in filtered_keys_for_ai:
                                current_key = msg.strip()

                        key_index = int(p * total_keys) if p > 0 else 0

                        # Calculate proper ETA based on actual processing time
                        if p > 0 and key_index > 0:
                            elapsed_time = time.time() - st.session_state['processing_start_time']
                            avg_time_per_key = elapsed_time / key_index
                            remaining_keys = total_keys - key_index
                            eta_seconds = int(remaining_keys * avg_time_per_key)
                            
                            if eta_seconds > 0:
                                eta_minutes = eta_seconds // 60
                                eta_seconds = eta_seconds % 60
                                eta_text = f"⏱️ ETA: {eta_minutes}m {eta_seconds}s" if eta_minutes > 0 else f"⏱️ ETA: {eta_seconds}s"
                            else:
                                eta_text = "⏱️ Almost done!"
                        else:
                            eta_text = ""
                        
                        # Enhanced status display with ETA on same line
                        status_display = f"🤖 Generating English content... ({key_index}/{total_keys} keys) - {current_key}"
                        if eta_text:
                            status_display += f" {eta_text}"
                        status_text.text(status_display)
                        progress_bar.progress(p)
                    
                        # Starting AI processing
                    english_results = run_ai_processing(filtered_keys_for_ai, temp_ai_data, language=detected_language, progress_callback=progress_callback_eng_simple)
                    
                    if not english_results:
                        if detected_language == 'chinese':
                            st.error("❌ 中文内容生成失败")
                        else:
                            st.error("❌ English content generation failed")
                        return
                    
                    # Proofread content
                    if detected_language == 'chinese':
                        status_text.text(f"🧐 校对中文内容... (0/{total_keys} keys)")
                    else:
                        status_text.text(f"🧐 Proofreading English content... (0/{total_keys} keys)")
                    progress_bar.progress(0.6)
                    
                    def progress_callback_proof_eng(p, msg):
                        nonlocal current_key_index, current_key

                        # Store progress info in session state for display
                        if 'debug_progress' not in st.session_state:
                            st.session_state['debug_progress'] = []
                        st.session_state['debug_progress'].append(f"p={p}, msg='{msg}'")

                        # Parse the detailed message from the AI processing
                        current_key = "Proofreading"  # Default fallback

                        if msg and isinstance(msg, str):
                            # Format: "🔄 Processing Cash • OpenAI • Key 1/9"
                            if 'Processing' in msg and '•' in msg:
                                parts = msg.split('•')
                                if len(parts) >= 3 and 'Key' in parts[2]:
                                    # Extract key index from format like "Key 1/9"
                                    try:
                                        key_info = parts[2].strip()
                                        if '/' in key_info:
                                            current_key_index = int(key_info.split('/')[0].replace('Key ', ''))
                                            current_key = parts[0].replace('🔄 Processing', '').strip()
                                    except:
                                        pass
                                elif len(parts) >= 1:
                                    key_part = parts[0].replace('🔄 Processing', '').strip()
                                    if key_part:
                                        current_key = key_part
                            # Format: "🔄 Processing Cash" (without bullet points)
                            elif 'Processing' in msg:
                                key_part = msg.replace('🔄 Processing', '').strip()
                                if key_part:
                                    current_key = key_part
                            # Format: Direct key name
                            elif len(msg.strip()) < 50 and not '•' in msg and not 'Processing' in msg:
                                current_key = msg.strip()
                            # Format: Check if it's just a key name without "Processing"
                            elif msg.strip() in filtered_keys_for_ai:
                                current_key = msg.strip()

                        # Calculate progress based on current key index
                        if current_key_index == 0:
                            # Fallback: estimate from progress value
                            current_key_index = max(1, int(p * total_keys))

                        # Calculate proper ETA based on actual processing time
                        if p > 0 and current_key_index > 0:
                            elapsed_time = time.time() - st.session_state['processing_start_time']
                            avg_time_per_key = elapsed_time / current_key_index
                            remaining_keys = max(0, total_keys - current_key_index)
                            if remaining_keys > 0:
                                eta_seconds = int(remaining_keys * avg_time_per_key)

                                if eta_seconds > 0:
                                    eta_minutes = eta_seconds // 60
                                    eta_seconds = eta_seconds % 60
                                    eta_text = f"⏱️ ETA: {eta_minutes}m {eta_seconds}s" if eta_minutes > 0 else f"⏱️ ETA: {eta_seconds}s"
                                else:
                                    eta_text = "⏱️ Almost done!"
                            else:
                                eta_text = "⏱️ Almost done!"
                        else:
                            eta_text = ""

                        # Enhanced status display with ETA on same line
                        status_display = f"🧐 校对英文内容... ({current_key_index}/{total_keys} keys) - {current_key}"
                        if eta_text:
                            status_display += f" {eta_text}"
                        status_text.text(status_display)
                        progress_bar.progress(0.6 + p * 0.1)
                    
                    proofread_results = run_simple_proofreader(english_results, temp_ai_data, progress_callback=progress_callback_proof_eng)
                    
                    if proofread_results:
                        # Store results
                        if 'ai_content_store' not in st.session_state:
                            st.session_state['ai_content_store'] = {}

                        for key, result in proofread_results.items():
                            if key not in st.session_state['ai_content_store']:
                                st.session_state['ai_content_store'][key] = {}
                            content = result.get('content', str(result)) if isinstance(result, dict) else str(result)
                            st.session_state['ai_content_store'][key]['agent1_content'] = content
                            st.session_state['ai_content_store'][key]['current_content'] = content
                            st.session_state['ai_content_store'][key]['agent1_timestamp'] = datetime.datetime.now().isoformat()
                            
                            # Store proofread content for preview (store the full result object, not just content)
                            st.session_state['ai_content_store'][key]['agent2_content'] = result
                            

                        st.session_state['agent_states']['agent1_results'] = proofread_results
                        st.session_state['agent_states']['agent1_completed'] = True
                        st.session_state['agent_states']['agent1_success'] = True

                        print(f"💾 STORED: {len(proofread_results)} AI results in session state")
                
                # Generate content files
                status_text.text("📝 Generating content files...")
                progress_bar.progress(0.8)

                print(f"📝 CONTENT GEN: Processing {statement_type} mode with {len(st.session_state.get('ai_content_store', {}))} stored results")

                # Handle "All" case - generate both BS and IS content
                if statement_type == "ALL":
                    print(f"🔄 ALL mode: Processing content generation...")

                    # Backup existing files BEFORE generating new ones
                    if os.path.exists("fdd_utils/bs_content.md"):
                        try:
                            os.rename("fdd_utils/bs_content.md", "fdd_utils/bs_content_backup.md")
                        except:
                            pass
                    if os.path.exists("fdd_utils/is_content.md"):
                        try:
                            os.rename("fdd_utils/is_content.md", "fdd_utils/is_content_backup.md")
                        except:
                            pass

                    try:
                        # Generate BS content
                        st.session_state['current_statement_type'] = 'BS'
                        generate_content_from_session_storage(selected_entity)

                        # Generate IS content
                        st.session_state['current_statement_type'] = 'IS'
                        generate_content_from_session_storage(selected_entity)

                        # Clean up backup files (no longer needed)
                        if os.path.exists("fdd_utils/bs_content_backup.md"):
                            try:
                                os.remove("fdd_utils/bs_content_backup.md")
                            except:
                                pass
                        if os.path.exists("fdd_utils/is_content_backup.md"):
                            try:
                                os.remove("fdd_utils/is_content_backup.md")
                            except:
                                pass

                        print(f"✅ ALL mode content generation completed")

                        # Verify content files were created for ALL mode
                        bs_content_file = "fdd_utils/bs_content.md"
                        is_content_file = "fdd_utils/is_content.md"
                        if os.path.exists(bs_content_file):
                            with open(bs_content_file, 'r', encoding='utf-8') as f:
                                content = f.read()
                            print(f"✅ ALL mode BS content file: {bs_content_file}, size: {len(content)}")
                        else:
                            print(f"❌ ALL mode BS content file not found: {bs_content_file}")

                        if os.path.exists(is_content_file):
                            with open(is_content_file, 'r', encoding='utf-8') as f:
                                content = f.read()
                            print(f"✅ ALL mode IS content file: {is_content_file}, size: {len(content)}")
                        else:
                            print(f"❌ ALL mode IS content file not found: {is_content_file}")

                    except Exception as e:
                        # Restore from backup if generation failed
                        if os.path.exists("fdd_utils/bs_content_backup.md"):
                            try:
                                os.rename("fdd_utils/bs_content_backup.md", "fdd_utils/bs_content.md")
                            except:
                                pass
                        if os.path.exists("fdd_utils/is_content_backup.md"):
                            try:
                                os.rename("fdd_utils/is_content_backup.md", "fdd_utils/is_content.md")
                            except:
                                pass
                        st.error(f"❌ Content generation failed: {str(e)}")
                        raise e

                    # Restore current statement type
                    st.session_state['current_statement_type'] = 'ALL'

                    # Mark all agents as completed for "All" mode
                    st.session_state['agent_states']['agent1_completed'] = True
                    st.session_state['agent_states']['agent1_success'] = True
                    st.session_state['agent_states']['agent2_completed'] = True
                    st.session_state['agent_states']['agent2_success'] = True
                    st.session_state['agent_states']['agent3_completed'] = True
                    st.session_state['agent_states']['agent3_success'] = True
                    print(f"✅ ALL mode agents completed")
                else:
                    # For individual BS/IS modes, mark agents as completed
                    if statement_type == "BS":
                        # Mark all agents as completed for BS mode
                        st.session_state['agent_states']['agent1_completed'] = True
                        st.session_state['agent_states']['agent1_success'] = True
                        st.session_state['agent_states']['agent2_completed'] = True
                        st.session_state['agent_states']['agent2_success'] = True
                        st.session_state['agent_states']['agent3_completed'] = True
                        st.session_state['agent_states']['agent3_success'] = True
                        print(f"✅ BS mode agents completed")
                        # Verify BS content files were created
                        bs_content_file = "fdd_utils/bs_content.md"
                        if os.path.exists(bs_content_file):
                            with open(bs_content_file, 'r', encoding='utf-8') as f:
                                content = f.read()
                            print(f"✅ BS content file created: {bs_content_file}, size: {len(content)}")
                        else:
                            print(f"❌ BS content file not found: {bs_content_file}")
                    elif statement_type == "IS":
                        print(f"🔄 Processing IS mode content generation...")
                        # Set current statement type for IS mode
                        st.session_state['current_statement_type'] = 'IS'
                        # Mark all agents as completed for IS mode
                        st.session_state['agent_states']['agent1_completed'] = True
                        st.session_state['agent_states']['agent1_success'] = True
                        st.session_state['agent_states']['agent2_completed'] = True
                        st.session_state['agent_states']['agent2_success'] = True
                        st.session_state['agent_states']['agent3_completed'] = True
                        st.session_state['agent_states']['agent3_success'] = True
                        print(f"✅ All agents marked as completed for IS mode")
                        # Generate content for IS mode
                        generate_content_from_session_storage(selected_entity)
                        print(f"✅ IS mode content generation completed")

                        # Verify IS content files were created
                        is_content_file = "fdd_utils/is_content.md"
                        if os.path.exists(is_content_file):
                            with open(is_content_file, 'r', encoding='utf-8') as f:
                                content = f.read()
                            print(f"✅ IS content file created: {is_content_file}, size: {len(content)}")
                        else:
                            print(f"❌ IS content file not found: {is_content_file}")
                    else:
                        generate_content_from_session_storage(selected_entity)

                # Export PowerPoint
                status_text.text("📊 Exporting PowerPoint...")
                progress_bar.progress(0.9)
                
                try:
                    # Export PowerPoint and automatically show download
                    output_path = export_enhanced_pptx(selected_entity, statement_type, language=detected_language,
                                       financial_statement_tab=financial_statement_tab, include_bshn=include_bshn, row_limit=row_limit)

                    # Verify the file was created
                    if os.path.exists(output_path):
                        print(f"✅ PowerPoint file created successfully: {output_path}")
                        progress_bar.progress(1.0)
                        status_text.text(f"✅ Report generation and export completed ({language_display})")

                        # Show success message with download info
                        st.success(f"🎉 Report generated successfully! The download button should appear above.")
                    else:
                        print(f"❌ PowerPoint file was not created: {output_path}")
                        progress_bar.progress(1.0)
                        status_text.text(f"⚠️ Report generation completed but PowerPoint export failed")
                        st.error(f"❌ PowerPoint file was not created at: {output_path}")

                except Exception as export_error:
                    progress_bar.progress(1.0)
                    status_text.text(f"⚠️ Report generated but export failed: {str(export_error)}")
                    st.error(f"❌ PowerPoint export failed: {str(export_error)}")
                    st.info("💡 Content has been generated successfully. You can try the export again.")
                    print(f"❌ PowerPoint export error: {export_error}")

                time.sleep(1)
                st.rerun()
                
            except Exception as e:
                progress_bar.progress(1.0)
                status_text.text(f"❌ Report generation failed: {e}")
                time.sleep(1)
                st.rerun()


        # Display AI Results
        agent_states = st.session_state.get('agent_states', {})
        any_agent_completed = any([
            agent_states.get('agent1_completed', False),
            agent_states.get('agent2_completed', False),
            agent_states.get('agent3_completed', False)
        ])
        
        if any_agent_completed:
            filtered_keys = st.session_state.get('filtered_keys_for_ai', [])
            
            if filtered_keys:
                # Create tabs for each key
                tab_labels = [get_key_display_name(key) for key in filtered_keys]
                key_tabs = st.tabs(tab_labels)
                
                for i, key in enumerate(filtered_keys):
                    with key_tabs[i]:
                        # Get all available content
                        ai_content_store = st.session_state.get('ai_content_store', {})
                        agent1_results = agent_states.get('agent1_results', {}) or {}
                        agent3_results_all = agent_states.get('agent3_results', {}) or {}
                        
                        # Show final content (Agent 3 if available, otherwise Agent 1)
                        final_content = None
                        if key in agent3_results_all:
                            pr = agent3_results_all[key]
                            translated_content = pr.get('translated_content', '')
                            corrected_content = pr.get('corrected_content', '') or pr.get('content', '')
                            final_content = translated_content if translated_content and pr.get('is_chinese', False) else corrected_content
                        elif key in agent1_results and agent1_results[key]:
                            content = agent1_results[key]
                            final_content = content.get('content', str(content)) if isinstance(content, dict) else str(content)
                        
                        if final_content:
                            st.markdown("**Final Content:**")
                            st.markdown(final_content)
                        
                        # Show Agent 1 results (AI Generation)
                        if key in agent1_results and agent1_results[key]:
                            content = agent1_results[key]
                            content_str = content.get('content', str(content)) if isinstance(content, dict) else str(content)
                            
                            with st.expander("📝 AI Generation (Original)", expanded=False):
                                st.markdown(content_str)
                        
                        # Show Agent 2 results (Proofreading) with detailed feedback
                        if key in ai_content_store:
                            proofread_result = ai_content_store[key].get('agent2_content', '')
                            if proofread_result:
                                with st.expander("🧐 Proofread Content", expanded=False):
                                    # Display the corrected content
                                    if isinstance(proofread_result, dict):
                                        corrected_content = proofread_result.get('content', '')
                                        st.markdown(corrected_content)
                                        
                                        # Show detailed proofreader feedback
                                        st.markdown("**🔍 Proofreader Analysis:**")
                                        
                                        # Show compliance status
                                        is_compliant = proofread_result.get('is_compliant', True)
                                        if is_compliant:
                                            st.success("✅ Content is compliant")
                                        else:
                                            st.warning("⚠️ Content has compliance issues")
                                        
                                        # Show issues found - with robust type checking
                                        issues = proofread_result.get('issues', [])
                                        if not isinstance(issues, list):
                                            print(f"⚠️ WARNING: Issues field is not a list in display, converting from {type(issues)}")
                                            if isinstance(issues, str):
                                                issues_str = issues.strip()
                                                if ',' in issues_str or '\n' in issues_str:
                                                    issues = [issue.strip() for issue in issues_str.replace('\n', ',').split(',') if issue.strip()]
                                                else:
                                                    issues = [issues_str] if issues_str else []
                                            elif isinstance(issues, (int, float)):
                                                issues = [str(issues)]
                                            else:
                                                try:
                                                    issues = [str(issues)] if issues else []
                                                except:
                                                    issues = []
                                        
                                        if issues:
                                            st.markdown("**Issues Found:**")
                                            for issue in issues:
                                                st.markdown(f"• {issue}")
                                        
                                        # Show figure checks - with robust type checking
                                        figure_checks = proofread_result.get('figure_checks', [])
                                        if not isinstance(figure_checks, list):
                                            print(f"⚠️ WARNING: figure_checks field is not a list in display, converting from {type(figure_checks)}")
                                            if isinstance(figure_checks, str):
                                                figure_checks_str = figure_checks.strip()
                                                if ',' in figure_checks_str or '\n' in figure_checks_str:
                                                    figure_checks = [check.strip() for check in figure_checks_str.replace('\n', ',').split(',') if check.strip()]
                                                else:
                                                    figure_checks = [figure_checks_str] if figure_checks_str else []
                                            elif isinstance(figure_checks, (int, float)):
                                                figure_checks = [str(figure_checks)]
                                            else:
                                                try:
                                                    figure_checks = [str(figure_checks)] if figure_checks else []
                                                except:
                                                    figure_checks = []
                                        
                                        if figure_checks:
                                            st.markdown("**Figure Validation:**")
                                            for check in figure_checks:
                                                st.markdown(f"• {check}")
                                        
                                        # Show entity checks - with robust type checking
                                        entity_checks = proofread_result.get('entity_checks', [])
                                        if not isinstance(entity_checks, list):
                                            print(f"⚠️ WARNING: entity_checks field is not a list in display, converting from {type(entity_checks)}")
                                            if isinstance(entity_checks, str):
                                                entity_checks_str = entity_checks.strip()
                                                if ',' in entity_checks_str or '\n' in entity_checks_str:
                                                    entity_checks = [check.strip() for check in entity_checks_str.replace('\n', ',').split(',') if check.strip()]
                                                else:
                                                    entity_checks = [entity_checks_str] if entity_checks_str else []
                                            elif isinstance(entity_checks, (int, float)):
                                                entity_checks = [str(entity_checks)]
                                            else:
                                                try:
                                                    entity_checks = [str(entity_checks)] if entity_checks else []
                                                except:
                                                    entity_checks = []
                                        
                                        if entity_checks:
                                            st.markdown("**Entity Validation:**")
                                            for check in entity_checks:
                                                st.markdown(f"• {check}")
                                        
                                        # Show grammar notes - with robust type checking
                                        grammar_notes = proofread_result.get('grammar_notes', [])
                                        if not isinstance(grammar_notes, list):
                                            print(f"⚠️ WARNING: grammar_notes field is not a list in display, converting from {type(grammar_notes)}")
                                            if isinstance(grammar_notes, str):
                                                grammar_notes_str = grammar_notes.strip()
                                                if ',' in grammar_notes_str or '\n' in grammar_notes_str:
                                                    grammar_notes = [note.strip() for note in grammar_notes_str.replace('\n', ',').split(',') if note.strip()]
                                                else:
                                                    grammar_notes = [grammar_notes_str] if grammar_notes_str else []
                                            elif isinstance(grammar_notes, (int, float)):
                                                grammar_notes = [str(grammar_notes)]
                                            else:
                                                try:
                                                    grammar_notes = [str(grammar_notes)] if grammar_notes else []
                                                except:
                                                    grammar_notes = []
                                        
                                        if grammar_notes:
                                            st.markdown("**Grammar & Style:**")
                                            for note in grammar_notes:
                                                st.markdown(f"• {note}")
                                        
                                        # Show pattern used
                                        pattern_used = proofread_result.get('pattern_used', '')
                                        if pattern_used:
                                            st.markdown(f"**Pattern Used:** {pattern_used}")
                                        
                                        # Show translation runs
                                        translation_runs = proofread_result.get('translation_runs', 0)
                                        if translation_runs > 0:
                                            st.markdown(f"**Heuristic Translation:** {translation_runs} run(s) applied")
                                    else:
                                        # Fallback for string content
                                        st.markdown(proofread_result)
                        
                        # Show Agent 3 results (Translation) if different from final content
                        if key in agent3_results_all and key not in agent1_results:
                            pr = agent3_results_all[key]
                            translated_content = pr.get('translated_content', '')
                            if translated_content and translated_content != final_content:
                                with st.expander("🌐 Translated Content", expanded=False):
                                    st.markdown(translated_content)
            else:
                st.info("No financial keys available for results display.")
        else:
            pass  # AI processing section ready
        
        # Debug section for progress messages
        if 'debug_progress' in st.session_state and st.session_state['debug_progress']:
            with st.expander("🔍 Debug: Progress Messages", expanded=False):
                for i, debug_msg in enumerate(st.session_state['debug_progress'][-10:]):  # Show last 10 messages
                    st.text(f"{i+1}. {debug_msg}")
        
        





def embed_bshn_data_simple(presentation_path, excel_file_path, sheet_name, project_name, language='english'):
    """Add BSHN data table to the first slide (BS1)"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
        import pandas as pd
        import os
        
        # Load the presentation
        prs = Presentation(presentation_path)
        
        # Read Excel data and apply column filtering
        df = pd.read_excel(excel_file_path, sheet_name=sheet_name)
        print(f"\n{'='*80}")
        print(f"🎯 POWERPOINT TABLE GENERATION - BALANCE SHEET MODE")
        print(f"{'='*80}")
        print(f"📊 Processing sheet: {sheet_name}")
        print(f"📊 Original data shape: {df.shape}")
        
        # Apply BS/IS separation first
        from fdd_utils.excel_processing import separate_balance_sheet_and_income_statement_tables, filter_to_indicative_adjusted_columns
        bs_data, is_data, separation_metadata = separate_balance_sheet_and_income_statement_tables(df, [project_name])
        
        # Use only Balance Sheet data (since this function is for BS mode)
        if bs_data:
            df = bs_data['data']
            print(f"✅ USING BALANCE SHEET DATA ONLY")
            print(f"   📊 BS data shape: {df.shape}")
            print(f"   📋 BS row range: rows {bs_data.get('data_range', (0, len(df)-1))[0]} to {bs_data.get('data_range', (0, len(df)-1))[1]} in original databook")
            print(f"   📋 Selected columns: {list(df.columns)}")
        else:
            print(f"⚠️ NO BS DATA FOUND - using filtered original")
            df = filter_to_indicative_adjusted_columns(df)
        
        # Filter out zero rows (all values are exactly 0, not NaN)
        original_row_count = len(df)
        if len(df.columns) > 1:
            desc_col = df.columns[0]
            value_cols = df.columns[1:]
            
            print(f"🔍 ZERO ROW FILTERING: {original_row_count} rows")
            
            mask = []
            kept_count = 0
            filtered_count = 0
            
            for idx, row in df.iterrows():
                desc_value = str(row[desc_col]).strip()
                has_description = desc_value not in ['', 'nan', 'None', 'NaN']
                
                # Check if ALL value columns are exactly zero (not NaN)
                all_values_zero = True
                has_any_data = False
                for col in value_cols:
                    val = pd.to_numeric(row[col], errors='coerce')
                    if pd.notna(val):
                        has_any_data = True
                        if val != 0:
                            all_values_zero = False
                            break
                
                # Keep row if has description AND (has non-zero values OR is header with NaN)
                if has_description:
                    if not has_any_data or not all_values_zero:
                        mask.append(True)
                        kept_count += 1
                    else:
                        mask.append(False)
                        filtered_count += 1
                else:
                    mask.append(False)
                    filtered_count += 1
            
            if any(mask):
                df = df[mask]
            
            print(f"   Kept: {kept_count}, Filtered: {filtered_count}")
        print(f"📊 FINAL PPT TABLE DATA: {df.shape}")
        print(f"   📋 Rows after zero filtering: {len(df)}")
        print(f"   📋 Columns: {list(df.columns)}")
        if len(df) > 0:
            print(f"   📋 Sample rows:")
            for i in range(min(3, len(df))):
                desc = df.iloc[i, 0] if len(df.columns) > 0 else 'N/A'
                val = df.iloc[i, 1] if len(df.columns) > 1 else 'N/A'
                print(f"      Row {i}: \"{desc}\" | {val}")
        print(f"{'='*80}")
        
        # Get the first slide (BS1)
        first_slide = prs.slides[0]
        
        # Create table with proper header structure
        rows = len(df) + 1  # +1 for header only (no separate currency row)
        cols = len(df.columns)
        
        # Convert cm to inches (1 inch = 2.54 cm)
        table_width = 12.14 / 2.54  # 4.78 inches
        table_height = 10.49 / 2.54  # 4.13 inches
        
        # Position table at specific coordinates (5.01cm from top, 0.36cm from left)
        table_x = 0.36 / 2.54  # 0.14 inches from left
        table_y = 5.01 / 2.54  # 1.97 inches from top
        
        # Add table to the first slide with specific dimensions and position
        table_shape = first_slide.shapes.add_table(rows, cols, Inches(table_x), Inches(table_y), Inches(table_width), Inches(table_height))
        
        # Force the table to maintain exact dimensions
        table_shape.width = Inches(table_width)
        table_shape.height = Inches(table_height)
        table = table_shape.table
        
        # Set table properties to prevent auto-sizing
        table.autofit = False
        
        # Reduce cell margins and padding for more compact table
        for row in table.rows:
            for cell in row.cells:
                # Set minimal margins
                cell.margin_left = Inches(0.05)
                cell.margin_right = Inches(0.05)
                cell.margin_top = Inches(0.02)
                cell.margin_bottom = Inches(0.02)
        
        # Set table name for future reference
        table_shape.name = "financialData"
        
        # Extract first word from project name for header
        entity_first_word = project_name.split()[0] if project_name else "Entity"
        
        # Fill header row (row 0) - language-aware title
        cell = table.cell(0, 0)
        if language.lower() == 'chinese':
            cell.text = f"經示意性調整後資產負債表 - {entity_first_word}"
        else:
            cell.text = f"Indicative adjusted balance sheet - Project {entity_first_word}"
        # Format header with highlighting
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(9)  # Font size 9
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT  # Left align
        cell.text_frame.word_wrap = False  # Disable text wrapping
        cell.text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # Vertically center
        # Add background color for header - RGB(42, 72, 121)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(42, 72, 121)  # Custom blue background
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
        
        # Merge header cells across all columns
        if cols > 1:
            table.cell(0, 0).merge(table.cell(0, cols-1))
        
        # Set row height for header row - 0.64cm
        table.rows[0].height = Inches(0.64 / 2.54)  # 0.25 inches (0.64cm)
        
        # Adjust column widths - make first column 35-40% of the table width
        if cols >= 4:
            # First column 37.5% of the table width (middle of 35-40% range)
            first_col_width = table_width * 0.375  # 37.5% of the table width
            remaining_width = table_width * 0.625
            other_col_width = remaining_width / (cols - 1)  # Distribute remaining width
            
            table.columns[0].width = Inches(first_col_width)  # First column 37.5% width
            for i in range(1, cols):
                table.columns[i].width = Inches(other_col_width)  # Other columns share remaining 62.5%
        else:
            # If less than 4 columns, distribute evenly
            for i in range(cols):
                table.columns[i].width = Inches(table_width / cols)
        
        # No separate currency row - integrate into header
        
        # Fill data rows with special formatting for totals/subtotals
        for row_idx, (_, row_data) in enumerate(df.iterrows()):
            for col_idx, value in enumerate(row_data):
                cell = table.cell(row_idx + 1, col_idx)  # +1 because we only have header row
                
                # Check if this is a total or subtotal row (case-insensitive)
                row_text = str(row_data.iloc[0]).lower() if len(row_data) > 0 else ""
                is_total_row = any(keyword in row_text for keyword in ['total', 'subtotal', 'sum', '合计', '小计'])
                
                # Check if this is a date column and format accordingly
                if col_idx > 0:
                    # Try to detect and convert date format
                    try:
                        date_str = str(value)
                        if ' ' in date_str:
                            date_part = date_str.split(' ')[0]  # Get date part only
                        else:
                            date_part = date_str
                        
                        # Parse and reformat date
                        from datetime import datetime
                        if '-' in date_part and len(date_part) >= 8:
                            # Try different date formats
                            for fmt in ['%Y-%m-%d', '%d-%m-%Y', '%m-%d-%Y']:
                                try:
                                    date_obj = datetime.strptime(date_part, fmt)
                                    cell_text = date_obj.strftime('%d-%b-%Y')
                                    break
                                except:
                                    continue
                            else:
                                cell_text = date_part  # Keep original if can't parse
                        elif '/' in date_part and len(date_part) >= 8:
                            # Try different date formats with /
                            for fmt in ['%Y/%m/%d', '%d/%m/%Y', '%m/%d/%Y']:
                                try:
                                    date_obj = datetime.strptime(date_part, fmt)
                                    cell_text = date_obj.strftime('%d-%b-%Y')
                                    break
                                except:
                                    continue
                            else:
                                cell_text = date_part  # Keep original if can't parse
                        else:
                            # Regular value formatting
                            if pd.isna(value):
                                cell_text = ""
                            elif isinstance(value, (int, float)):
                                cell_text = f"{value:,.0f}"
                            else:
                                cell_text = str(value)
                    except:
                        # Regular value formatting
                        if pd.isna(value):
                            cell_text = ""
                        elif isinstance(value, (int, float)):
                            cell_text = f"{value:,.0f}"
                        else:
                            cell_text = str(value)
                else:
                    # Regular value formatting for first column
                    if pd.isna(value):
                        cell_text = ""
                    elif isinstance(value, (int, float)):
                        cell_text = f"{value:,.0f}"
                    else:
                        cell_text = str(value)
                
                cell.text = cell_text
                
                if is_total_row:
                    # Special formatting for total/subtotal rows
                    cell.text_frame.paragraphs[0].font.bold = True
                    cell.text_frame.paragraphs[0].font.size = Pt(7)
                    cell.text_frame.word_wrap = False
                    # Light gray background for totals
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(217, 217, 217)  # Light gray
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
                else:
                    # Regular data formatting
                    cell.text_frame.paragraphs[0].font.bold = False
                    cell.text_frame.paragraphs[0].font.size = Pt(7)
                    cell.text_frame.word_wrap = False
                    # White background for regular data
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
                    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text
                
                # Right-align numbers
                if col_idx > 0:
                    cell.text_frame.paragraphs[0].alignment = PP_ALIGN.RIGHT
        
        # Set row heights for all data rows to be smaller
        for row_idx in range(1, len(table.rows)):
            table.rows[row_idx].height = Inches(0.15)  # Extremely small row height
        
        # Note: Table borders are complex in PowerPoint and may not always work reliably
        # The table will still function correctly without borders
        print("ℹ️ Table created successfully (borders skipped for compatibility)")
        
        # Save the updated presentation
        prs.save(presentation_path)
        print(f"✅ BSHN data successfully embedded from sheet '{sheet_name}' in first slide (BS1)")
        
    except Exception as e:
        print(f"❌ Error embedding BSHN data: {e}")
        raise


def export_enhanced_pptx(selected_entity, statement_type, language='english', financial_statement_tab=None, include_bshn=True, row_limit=20):
    """Enhanced PowerPoint export function with BSHN sheet and page designer using template"""
    try:
        if language == 'chinese':
            st.info("📊 开始生成中文 PowerPoint 演示文稿...")
        else:
            st.info("📊 Generating English PowerPoint presentation...")

        # Get project name
        words = selected_entity.split() if selected_entity else ['Project']
        project_name = ' '.join(words[:2]) if len(words) >= 2 else words[0] if words else 'Project'

        # Find template
        template_path = None
        for template in ["fdd_utils/template.pptx", "template.pptx"]:
            if os.path.exists(template):
                template_path = template
                break

        if not template_path:
            st.error("❌ PowerPoint template not found")
            return

        # Create output filename (sanitize project name)
        import re
        import uuid
        sanitized_project_name = re.sub(r'[^\w\-_]', '_', project_name).strip('_')
        language_suffix = "_CN" if language == 'chinese' else "_EN"
        unique_id = str(uuid.uuid4())[:8]  # Add unique identifier to prevent caching
        timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
        output_filename = f"{sanitized_project_name}_{statement_type.upper()}_{timestamp}_{unique_id}{language_suffix}.pptx"
        output_path = f"fdd_utils/output/{output_filename}"

        # Ensure output directory exists
        os.makedirs("fdd_utils/output", exist_ok=True)

        # Clean up old files (keep only last 5 files to prevent clutter)
        try:
            existing_files = [f for f in os.listdir("fdd_utils/output") if f.endswith('.pptx') and f.startswith(sanitized_project_name)]
            if len(existing_files) > 5:
                existing_files.sort(key=lambda x: os.path.getctime(os.path.join("fdd_utils/output", x)))
                files_to_remove = existing_files[:-5]  # Keep only the 5 most recent
                for old_file in files_to_remove:
                    try:
                        os.remove(os.path.join("fdd_utils/output", old_file))
                        print(f"🧹 Cleaned up old file: {old_file}")
                    except:
                        pass  # Ignore if file can't be removed
        except:
            pass  # Ignore cleanup errors

        # Get content file
        if statement_type == "IS":
            markdown_path = "fdd_utils/is_content.md"
            print(f"🔍 EXPORT: Looking for IS content file: {markdown_path}")
            print(f"🔍 EXPORT: Current working directory: {os.getcwd()}")
            print(f"🔍 EXPORT: File exists check: {os.path.exists(markdown_path)}")

            if os.path.exists(markdown_path):
                with open(markdown_path, 'r', encoding='utf-8') as f:
                    is_content = f.read()
                print(f"✅ EXPORT: IS content file exists, length: {len(is_content)}")
                if len(is_content) == 0:
                    print(f"❌ EXPORT: IS content file is empty!")
                    st.error("❌ IS content file is empty")
                    return
            else:
                print(f"❌ EXPORT: IS content file does not exist: {markdown_path}")
                # List files in fdd_utils directory to see what's there
                try:
                    files = os.listdir("fdd_utils")
                    print(f"🔍 EXPORT: Files in fdd_utils: {files}")
                except:
                    print("❌ EXPORT: Cannot list fdd_utils directory")
                st.error(f"❌ IS content file not found: {markdown_path}")
                return
        elif statement_type == "BS":
            markdown_path = "fdd_utils/bs_content.md"
            print(f"🔍 EXPORT: Looking for BS content file: {markdown_path}")
            print(f"🔍 EXPORT: File exists check: {os.path.exists(markdown_path)}")

            if not os.path.exists(markdown_path):
                print(f"❌ EXPORT: BS content file does not exist: {markdown_path}")
                # List files in fdd_utils directory to see what's there
                try:
                    files = os.listdir("fdd_utils")
                    print(f"🔍 EXPORT: Files in fdd_utils: {files}")
                except:
                    print("❌ EXPORT: Cannot list fdd_utils directory")
                st.error(f"❌ BS content file not found: {markdown_path}")
                return
            else:
                print(f"✅ EXPORT: BS content file exists")
        else:  # ALL
            st.info("🔄 Generating combined presentation...")
            # For combined, create both BS and IS then merge
            bs_path = "fdd_utils/bs_content.md"
            is_path = "fdd_utils/is_content.md"

            # Check if content files exist and have content
            bs_content = ""
            is_content = ""

            if os.path.exists(bs_path):
                with open(bs_path, 'r', encoding='utf-8') as f:
                    bs_content = f.read()
                print(f"🔍 ALL MODE: BS content file exists, length: {len(bs_content)}")
            else:
                print(f"🔍 ALL MODE: BS content file does not exist: {bs_path}")

            if os.path.exists(is_path):
                with open(is_path, 'r', encoding='utf-8') as f:
                    is_content = f.read()
                print(f"🔍 ALL MODE: IS content file exists, length: {len(is_content)}")
            else:
                print(f"🔍 ALL MODE: IS content file does not exist: {is_path}")

            if not os.path.exists(bs_path) or not os.path.exists(is_path):
                st.error("❌ Content files not found. Please run AI processing first.")
                return

            if not bs_content or not is_content:
                st.error("❌ Content files are empty. Please run AI processing first.")
                return

            with tempfile.TemporaryDirectory() as temp_dir:
                bs_temp = os.path.join(temp_dir, "bs_temp.pptx")
                is_temp = os.path.join(temp_dir, "is_temp.pptx")

                try:
                    # Generate BS and IS presentations using template
                    export_pptx(template_path, bs_path, bs_temp, project_name, language=language, row_limit=row_limit)
                    export_pptx(template_path, is_path, is_temp, project_name, language=language, row_limit=row_limit)

                    # Merge presentations
                    merge_presentations(bs_temp, is_temp, output_path)
                    print(f"✅ ALL mode export completed")

                except Exception as merge_error:
                    st.error(f"❌ PowerPoint generation failed: {str(merge_error)}")
                    raise merge_error

            if language == 'chinese':
                st.success("✅ 中文组合演示文稿生成成功!")
            else:
                st.success("✅ Combined presentation generated successfully!")

        if statement_type in ["IS", "BS"]:
            if not os.path.exists(markdown_path):
                st.error(f"❌ Content file not found: {markdown_path}")
                st.info("💡 Please run AI processing first.")
                return

            # Get Excel file path for BSHN integration
            excel_file_path = None
            if include_bshn and statement_type == "BS" and financial_statement_tab:
                # Get the uploaded file path
                uploaded_file = st.session_state.get('uploaded_file')
                if uploaded_file:
                    if hasattr(uploaded_file, 'file_path'):
                        excel_file_path = uploaded_file.file_path
                    else:
                        # For uploaded files, save temporarily
                        import tempfile
                        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx')
                        temp_file.write(uploaded_file.getvalue())
                        temp_file.close()
                        excel_file_path = temp_file.name
                        st.info(f"📊 BSHN sheet will be included in BS1 from Excel tab: {financial_statement_tab}")
            
            # Use the template with the original export_pptx function (without automatic Excel embedding)
            try:
                print(f"🔍 EXPORT: Calling export_pptx with:")
                print(f"  template_path: {template_path}")
                print(f"  markdown_path: {markdown_path}")
                print(f"  output_path: {output_path}")
                print(f"  project_name: {project_name}")
                print(f"  statement_type: {statement_type}")

                # Check if markdown content is valid before calling export
                if statement_type == "IS":
                    print(f"🔍 EXPORT: IS markdown content preview: {is_content[:200]}...")
                elif statement_type == "BS":
                    with open(markdown_path, 'r', encoding='utf-8') as f:
                        bs_content = f.read()
                    print(f"🔍 EXPORT: BS markdown content preview: {bs_content[:200]}...")

                export_pptx(template_path, markdown_path, output_path, project_name, excel_file_path=None, language=language, statement_type=statement_type, row_limit=row_limit)

                # Verify the file was created
                if os.path.exists(output_path):
                    print(f"✅ EXPORT: PowerPoint file created successfully: {output_path}")
                    file_size = os.path.getsize(output_path)
                    print(f"✅ EXPORT: PowerPoint file size: {file_size} bytes")
                else:
                    print(f"❌ EXPORT: PowerPoint file was not created: {output_path}")
                    st.error(f"❌ PowerPoint file was not created at: {output_path}")
                    return

            except Exception as export_error:
                print(f"❌ EXPORT: PowerPoint generation failed: {str(export_error)}")
                import traceback
                print(f"❌ EXPORT: Full traceback: {traceback.format_exc()}")
                st.error(f"❌ PowerPoint generation failed: {str(export_error)}")
                st.info(f"💡 Check if content file exists: {markdown_path}")
                st.info(f"💡 Check if template exists: {template_path}")
                raise
            
            # Add BSHN data if requested
            if include_bshn and statement_type == "BS" and excel_file_path and financial_statement_tab:
                try:
                    embed_bshn_data_simple(output_path, excel_file_path, financial_statement_tab, project_name, language)
                    st.success(f"✅ BSHN data from '{financial_statement_tab}' sheet added to BS1")
                except Exception as e:
                    st.warning(f"⚠️ Could not add BSHN data: {str(e)}")
                    st.info("💡 The presentation was created but BSHN data could not be embedded")

        # Show download button
        if os.path.exists(output_path):
            with open(output_path, "rb") as file:
                download_label = f"📥 下载中文 PowerPoint: {output_filename}" if language == 'chinese' else f"📥 Download English PowerPoint: {output_filename}"

                st.download_button(
                    label=download_label,
                    data=file.read(),
                    file_name=output_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True
                )

        success_msg = f"✅ 中文 PowerPoint 生成完成: {output_filename}" if language == 'chinese' else f"✅ English PowerPoint generated successfully: {output_filename}"
        st.success(success_msg)

        # Return the output path for the main app to verify
        return output_path

    except Exception as e:
        st.error(f"❌ Export failed: {e}")
        return None




if __name__ == "__main__":
    main() 