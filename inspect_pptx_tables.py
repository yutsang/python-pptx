#!/usr/bin/env python3
"""Dump the exact cell-level formatting of every table-like object in a .pptx.

Use this against a real, correctly-formatted past report to extract a
concrete spec (table header style, borders, total-row style, content-cell
style) that fdd_utils/pptx.py's embed_financial_tables() should be
matching. Also works against our own generated output, so the same tool
can diff "what we produce" against "what the company format is".

Handles three cases, since UpSlide-produced decks often don't use a
native PowerPoint table shape:
  1. Native pptx table (GraphicFrame with has_table) -- dumped directly.
  2. Embedded/linked OLE object (e.g. an Excel range pasted via UpSlide) --
     the embedded blob is extracted; if it's an Excel file, it's saved
     next to the input and its cell formatting is dumped with openpyxl
     (same font/fill/border/alignment/number-format detail as the pptx
     table dump).
  3. Picture / flattened image -- flagged as not structurally inspectable
     (formatting can only be read visually), so it's not silently missed.
  4. Anything else -- listed in the per-slide shape inventory so nothing
     is silently skipped.

Shapes are walked recursively into groups, since UpSlide/manually-grouped
content often nests the real object inside one or more GroupShapes.

Usage:
    python inspect_pptx_tables.py <path.pptx> [--slide N] [--save-ole DIR]
"""
from __future__ import annotations

import argparse
import os
import sys
import zipfile
from typing import Any, Iterator, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu


# --------------------------------------------------------------------------
# shared formatting helpers
# --------------------------------------------------------------------------

def _emu_to_in(value: Optional[int]) -> Optional[float]:
    if value is None:
        return None
    return round(Emu(value).inches, 3)


def _color_str(color_format) -> str:
    try:
        if color_format is None or color_format.type is None:
            return "inherited"
        from pptx.enum.dml import MSO_THEME_COLOR
        ctype = color_format.type
        if str(ctype).startswith("MSO_COLOR_TYPE.THEME") or ctype == 2:
            try:
                return f"theme:{MSO_THEME_COLOR(color_format.theme_color).name}"
            except Exception:
                return f"theme:{color_format.theme_color}"
        return f"#{color_format.rgb}"
    except Exception:
        return "?"


def _fill_str(fill) -> str:
    try:
        ftype = fill.type
        if ftype is None:
            return "inherited"
        name = str(ftype).split(".")[-1] if ftype is not None else "NONE"
        if "SOLID" in name.upper():
            return f"solid {_color_str(fill.fore_color)}"
        return name
    except Exception:
        return "?"


def _run_font_desc(run) -> str:
    f = run.font
    try:
        name = f.name or "(inherit)"
        size = f"{f.size.pt:g}pt" if f.size else "(inherit)"
        bold = "B" if f.bold else ""
        italic = "I" if f.italic else ""
        style = "".join(x for x in (bold, italic) if x) or "-"
        color = _color_str(f.color) if f.color and f.color.type is not None else "inherited"
        return f"{name} {size} {style} color={color}"
    except Exception as exc:
        return f"(font read error: {exc})"


_BORDER_TAGS = (("L", "a:lnL"), ("R", "a:lnR"), ("T", "a:lnT"), ("B", "a:lnB"))


def _cell_borders(cell) -> str:
    tc = cell._tc
    tcPr = tc.find(qn("a:tcPr"))
    if tcPr is None:
        return "no tcPr"
    parts = []
    for label, tag in _BORDER_TAGS:
        ln = tcPr.find(qn(tag))
        if ln is None:
            continue
        w = ln.get("w")
        w_pt = round(int(w) / 12700, 2) if w else None
        if ln.find(qn("a:noFill")) is not None:
            parts.append(f"{label}=none")
            continue
        solid = ln.find(qn("a:solidFill"))
        color = "?"
        if solid is not None:
            srgb = solid.find(qn("a:srgbClr"))
            if srgb is not None:
                color = f"#{srgb.get('val')}"
            else:
                scheme = solid.find(qn("a:schemeClr"))
                if scheme is not None:
                    color = f"theme:{scheme.get('val')}"
        parts.append(f"{label}={w_pt}pt {color}")
    return "; ".join(parts) if parts else "no explicit borders (inherits table style)"


def _cell_fill(cell) -> str:
    try:
        return _fill_str(cell.fill)
    except Exception as exc:
        return f"(fill read error: {exc})"


def _is_total_like(text: str) -> bool:
    lowered = text.lower()
    return any(kw in lowered or kw in text for kw in ("total", "合计", "总计", "小计", "subtotal"))


# --------------------------------------------------------------------------
# shape tree walking (recurses into groups so nested/UpSlide content isn't missed)
# --------------------------------------------------------------------------

def _iter_shapes_recursive(shapes) -> Iterator[Any]:
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_shapes_recursive(shape.shapes)


def _shape_type_name(shape) -> str:
    try:
        return str(shape.shape_type)
    except Exception:
        return "UNKNOWN"


# --------------------------------------------------------------------------
# native pptx table dump
# --------------------------------------------------------------------------

def dump_pptx_table(slide_idx: int, shape) -> None:
    tbl = shape.table
    n_rows = len(tbl.rows)
    n_cols = len(tbl.columns)
    print("=" * 100)
    print(
        f"SLIDE {slide_idx + 1}  NATIVE PPTX TABLE  shape_name={shape.name!r}  "
        f"pos=({_emu_to_in(shape.left)}in, {_emu_to_in(shape.top)}in)  "
        f"size=({_emu_to_in(shape.width)}in x {_emu_to_in(shape.height)}in)  "
        f"rows={n_rows} cols={n_cols}"
    )
    tblPr = tbl._tbl.find(qn("a:tblPr"))
    if tblPr is not None:
        style_id_elem = tblPr.find(qn("a:tableStyleId"))
        style_id = style_id_elem.text if style_id_elem is not None else None
        print(
            f"  tableStyleId={style_id}  firstRow={tblPr.get('firstRow')}  "
            f"firstCol={tblPr.get('firstCol')}  bandRow={tblPr.get('bandRow')}  "
            f"bandCol={tblPr.get('bandCol')}"
        )
    print("  column widths (in):", [_emu_to_in(col.width) for col in tbl.columns])
    print("  row heights   (in):", [_emu_to_in(row.height) for row in tbl.rows])
    print("-" * 100)
    for r in range(n_rows):
        for c in range(n_cols):
            cell = tbl.cell(r, c)
            if cell.is_spanned:
                continue
            text = cell.text
            span = f"  [merged {cell.span_height}x{cell.span_width}]" if cell.is_merge_origin else ""
            fonts, alignments = [], []
            for p in cell.text_frame.paragraphs:
                alignments.append(str(p.alignment) if p.alignment is not None else "inherited")
                for run in p.runs:
                    fonts.append(_run_font_desc(run))
            flag = "  *** TOTAL/SUBTOTAL-LIKE ***" if _is_total_like(text) and text.strip() else ""
            print(f"  R{r}C{c}{span}: text={text!r}")
            print(f"        font={fonts[0] if fonts else '(no runs / empty)'}  "
                  f"align={alignments[0] if alignments else 'inherited'}  fill={_cell_fill(cell)}{flag}")
            print(f"        borders: {_cell_borders(cell)}")
    print()


# --------------------------------------------------------------------------
# OLE object handling (UpSlide-style embedded/linked Excel range)
# --------------------------------------------------------------------------

_XLSX_MAGIC = b"PK\x03\x04"
_OLE_COMPOUND_MAGIC = b"\xd0\xcf\x11\xe0"


def dump_ole_object(slide_idx: int, shape, save_dir: Optional[str]) -> None:
    print("=" * 100)
    kind = "EMBEDDED" if shape.shape_type == MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT else "LINKED"
    print(
        f"SLIDE {slide_idx + 1}  {kind} OLE OBJECT  shape_name={shape.name!r}  "
        f"pos=({_emu_to_in(shape.left)}in, {_emu_to_in(shape.top)}in)  "
        f"size=({_emu_to_in(shape.width)}in x {_emu_to_in(shape.height)}in)"
    )
    try:
        ole = shape.ole_format
    except Exception as exc:
        print(f"  Could not read ole_format: {exc}")
        return
    print(f"  prog_id={ole.prog_id!r}  show_as_icon={ole.show_as_icon}")
    try:
        blob = ole.blob
    except Exception as exc:
        print(f"  No embedded blob available (likely a LINKED, not EMBEDDED, object): {exc}")
        print("  This references an external Excel file by path/link, not stored inside the pptx --")
        print("  the source .xlsx would need to be provided separately to inspect its formatting.")
        return

    if blob[:4] == _XLSX_MAGIC:
        ext = ".xlsx"
    elif blob[:4] == _OLE_COMPOUND_MAGIC:
        ext = ".xls"
    else:
        ext = ".bin"
    print(f"  Embedded blob: {len(blob)} bytes, detected format={ext}")

    if ext == ".bin":
        print("  Not a recognisable Excel file signature -- can't inspect cell formatting directly.")
        return

    out_dir = save_dir or os.path.dirname(os.path.abspath(shape.part.package.filename if hasattr(shape.part.package, 'filename') else ".")) or "."
    out_path = os.path.join(out_dir, f"extracted_ole_slide{slide_idx + 1}_{shape.shape_id}{ext}")
    try:
        with open(out_path, "wb") as f:
            f.write(blob)
        print(f"  Saved to: {out_path}")
    except Exception as exc:
        print(f"  Could not save extracted blob: {exc}")
        return

    if ext == ".xlsx":
        try:
            dump_xlsx_formatting(out_path)
        except Exception as exc:
            print(f"  openpyxl formatting dump failed: {exc}")
    else:
        print("  Legacy .xls (OLE compound) format -- open manually in Excel to inspect formatting,")
        print("  or re-save as .xlsx and re-run this script's --xlsx mode.")
    print()


def dump_xlsx_formatting(
    xlsx_path: str,
    max_rows: int = 100,
    max_cols: int = 30,
    sheet_name: Optional[str] = None,
) -> None:
    import openpyxl

    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    worksheets = wb.worksheets
    if sheet_name is not None:
        if sheet_name not in wb.sheetnames:
            print(f"  Sheet {sheet_name!r} not found. Available sheets: {wb.sheetnames}")
            return
        worksheets = [wb[sheet_name]]

    for ws in worksheets:
        used_rows = min(ws.max_row or 0, max_rows)
        used_cols = min(ws.max_column or 0, max_cols)
        if not used_rows or not used_cols:
            continue
        print(f"  --- worksheet {ws.title!r} ({ws.max_row}x{ws.max_column}, showing top-left "
              f"{used_rows}x{used_cols}) ---")
        for row in ws.iter_rows(min_row=1, max_row=used_rows, min_col=1, max_col=used_cols):
            for cell in row:
                if cell.value in (None, ""):
                    continue
                font = cell.font
                fill = cell.fill
                border = cell.border
                text = str(cell.value)
                flag = "  *** TOTAL/SUBTOTAL-LIKE ***" if _is_total_like(text) else ""

                def _rgb_str(color) -> Optional[str]:
                    # openpyxl's Color.rgb can be a real hex string, an indexed/theme
                    # int, or an internal enum sentinel depending on how the cell was
                    # styled -- only format it if it's actually an 8-char ARGB hex string.
                    try:
                        rgb = color.rgb if color else None
                    except Exception:
                        return None
                    if isinstance(rgb, str) and len(rgb) == 8:
                        return f"#{rgb}"
                    return None

                fill_color = None
                try:
                    if fill and fill.fgColor:
                        rgb = _rgb_str(fill.fgColor)
                        if rgb and rgb != "#00000000":
                            fill_color = rgb
                except Exception:
                    pass

                def _side(side):
                    if side is None or side.style is None:
                        return "none"
                    return f"{side.style} {_rgb_str(side.color) or 'inherited'}"

                font_name = font.name or "(inherited)"
                font_size = f"{font.size}pt" if font.size else "(inherited)"
                font_color = _rgb_str(font.color) or "inherited"
                print(f"  {cell.coordinate}: value={text!r}  number_format={cell.number_format!r}{flag}")
                print(
                    f"        font={font_name} {font_size} "
                    f"{'B' if font.bold else ''}{'I' if font.italic else ''} "
                    f"color={font_color}  "
                    f"fill={fill_color or 'none'}  align={cell.alignment.horizontal}"
                )
                print(
                    f"        borders: L={_side(border.left)}; R={_side(border.right)}; "
                    f"T={_side(border.top)}; B={_side(border.bottom)}"
                )
    print()


# --------------------------------------------------------------------------
# source-file breadcrumb scan (UpSlide/Excel links often leave a path
# reference behind even when the pasted content itself is a flattened
# picture, e.g. in customXml parts or a picture's alt-text/description)
# --------------------------------------------------------------------------

def scan_for_source_link_breadcrumbs(pptx_path: str) -> None:
    print("\n### Scanning for a source-Excel-file breadcrumb (customXml parts, picture alt-text) ###")
    found_any = False
    try:
        with zipfile.ZipFile(pptx_path) as z:
            candidate_parts = [
                n for n in z.namelist()
                if "customxml" in n.lower() or "upslide" in n.lower() or "externallink" in n.lower()
            ]
            for name in candidate_parts:
                if not name.lower().endswith((".xml", ".rels")):
                    continue
                data = z.read(name)
                if len(data) > 4000:
                    print(f"  {name}: {len(data)} bytes (too large to print in full; open manually if relevant)")
                    found_any = True
                    continue
                text = data.decode("utf-8", errors="replace")
                if any(hint in text.lower() for hint in (".xlsx", ".xls", "upslide", "target=")):
                    print(f"  {name}:")
                    print("    " + text.replace("\n", "\n    "))
                    found_any = True
    except Exception as exc:
        print(f"  Could not scan zip parts: {exc}")

    try:
        prs = Presentation(pptx_path)
        for slide_idx, slide in enumerate(prs.slides):
            for shape in _iter_shapes_recursive(slide.shapes):
                cNvPr = shape._element.find(f".//{qn('p:cNvPr')}")
                if cNvPr is None:
                    continue
                descr = cNvPr.get("descr")
                title = cNvPr.get("title")
                if descr or title:
                    print(f"  Slide {slide_idx + 1} shape {shape.name!r}: descr={descr!r} title={title!r}")
                    found_any = True
    except Exception as exc:
        print(f"  Could not scan shape metadata: {exc}")

    if not found_any:
        print("  Nothing found. UpSlide's link-back reference (if any) isn't stored inside this pptx in a "
              "readable form -- you'll need to locate the source .xlsx directly (e.g. via the UpSlide panel "
              "in PowerPoint's ribbon, which usually shows 'linked to: <path>' for each pasted object, or "
              "by asking whoever prepared this deck which workbook it was built from).")


# --------------------------------------------------------------------------
# main
# --------------------------------------------------------------------------

def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path", help="Path to a .pptx file, OR (with --xlsx) a raw .xlsx to inspect directly")
    ap.add_argument("--slide", type=int, default=None, help="1-indexed slide number to restrict to")
    ap.add_argument("--save-ole", default=None, help="Directory to save any extracted OLE/Excel blobs into (default: alongside the pptx)")
    ap.add_argument("--xlsx", action="store_true", help="Treat pptx_path as a raw .xlsx file and dump its cell formatting directly (skip pptx parsing)")
    ap.add_argument("--sheet", default=None, help="With --xlsx, restrict the dump to this worksheet name only")
    args = ap.parse_args()

    if args.xlsx:
        dump_xlsx_formatting(args.pptx_path, sheet_name=args.sheet)
        return 0

    prs = Presentation(args.pptx_path)
    found_table_or_ole = False

    for slide_idx, slide in enumerate(prs.slides):
        if args.slide is not None and (slide_idx + 1) != args.slide:
            continue

        all_shapes = list(_iter_shapes_recursive(slide.shapes))
        print(f"\n### Slide {slide_idx + 1}: {len(all_shapes)} shape(s) (including nested in groups) ###")
        for shape in all_shapes:
            print(f"    - name={shape.name!r:30} type={_shape_type_name(shape):35} "
                  f"pos=({_emu_to_in(shape.left)}in, {_emu_to_in(shape.top)}in) "
                  f"size=({_emu_to_in(shape.width)}in x {_emu_to_in(shape.height)}in)")

        for shape in all_shapes:
            if getattr(shape, "has_table", False):
                found_table_or_ole = True
                dump_pptx_table(slide_idx, shape)
            elif shape.shape_type in (MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT):
                found_table_or_ole = True
                dump_ole_object(slide_idx, shape, args.save_ole)
            elif shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
                print(f"    [picture shape {shape.name!r} on slide {slide_idx + 1} -- flattened image, "
                      f"no structural cell data available; formatting can only be read visually]")

    if not found_table_or_ole:
        print(
            "\nNo native table AND no OLE object found anywhere (including inside groups). "
            "The shape inventory above lists everything actually on each slide -- if the "
            "'table' shows up as PICTURE/LINKED_PICTURE, it's a flattened image and structural "
            "formatting can't be extracted programmatically; describe it visually instead, or "
            "check if UpSlide kept a linked source .xlsx file you can share directly."
        )
        scan_for_source_link_breadcrumbs(args.pptx_path)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
