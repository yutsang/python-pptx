#!/usr/bin/env python3
"""Dump the vertical band of every paragraph in a commentary frame, next to
the tables floated over it.

The table-account layout reserves a table's height as BLANK paragraphs
inside the column's own text frame and floats the table on top of them
(fdd_utils/pptx.py's _render_table_accounts_stack). That only works while
the renderer's cursor arithmetic and the frame's real rendered paragraph
heights agree. When they drift, the table lands on a paragraph that has
actual text in it and inspect_pptx.py reports TABLE OVERLAPS REAL TEXT --
but not by how much, or which paragraph, which is what you need to fix it.

This prints, per slide and per commentary shape:
    * every paragraph, its band in absolute points, and whether it is a
      blank spacer or real text
    * every table's own band
    * the exact overlap, in points, for any collision

Read-only. No AI, no export -- runs against an already-exported .pptx.

Usage:
    python inspect_table_bands.py for_test/pptx_previews/Crescent-databook.preview.pptx
    python inspect_table_bands.py <file.pptx> --slide 4
"""
from __future__ import annotations

import argparse
import sys

from pptx import Presentation
from pptx.util import Emu

from fdd_utils.financial_common import load_yaml_file
from fdd_utils.pptx.helpers import _resolve_font_metrics_path
from fdd_utils.text_metrics import get_measurer, text_box_from_shape

FONT_PT = 9.0
BULLET_HANG_PT = 10.8


def _is_chinese(text: str) -> bool:
    if not text:
        return False
    return sum(1 for c in text if "一" <= c <= "鿿") / len(text) > 0.3


def _load_cfg():
    for c in ("fdd_utils/config.yml", "fdd_utils/config.example.yml"):
        try:
            cfg = load_yaml_file(c)
            if cfg:
                return cfg
        except (FileNotFoundError, OSError):
            continue
    return {}


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path")
    ap.add_argument("--slide", type=int, default=0,
                    help="1-indexed slide to dump; 0 (default) = every slide that has a table")
    args = ap.parse_args()

    packing = ((_load_cfg().get("pptx") or {}).get("commentary_packing") or {})
    measurers = {
        True: get_measurer(packing.get("font_family_chi") or "Microsoft YaHei", FONT_PT,
                           is_cjk=True, line_spacing=1.0,
                           metrics_path=packing.get("font_metrics_path_chi")
                           or _resolve_font_metrics_path(True, packing)),
        False: get_measurer(packing.get("font_family_eng") or "Arial", FONT_PT,
                            is_cjk=False, line_spacing=1.0,
                            metrics_path=packing.get("font_metrics_path_eng")
                            or _resolve_font_metrics_path(False, packing)),
    }
    print(f"Measurement source: CHI={measurers[True].source}  ENG={measurers[False].source}\n")

    prs = Presentation(args.pptx_path)
    problems = 0

    for idx, slide in enumerate(prs.slides, start=1):
        if args.slide and idx != args.slide:
            continue
        tables = [s for s in slide.shapes if getattr(s, "has_table", False)]
        frames = [s for s in slide.shapes
                  if getattr(s, "has_text_frame", False)
                  and (getattr(s, "name", "") or "").startswith("textMainBullets")]
        if not tables and not args.slide:
            continue
        print("=" * 74)
        print(f"  SLIDE {idx}   {len(tables)} table(s), {len(frames)} commentary frame(s)")
        print("=" * 74)

        for t in tables:
            print(f"  [table {t.name!r}] {Emu(t.top).pt:7.1f} -> "
                  f"{Emu(t.top).pt + Emu(t.height).pt:7.1f} pt   left={Emu(t.left).inches:.2f}in")
        print()

        for shape in frames:
            text = shape.text_frame.text or ""
            if not text.strip():
                continue
            m = measurers[_is_chinese(text)]
            box = text_box_from_shape(shape)
            line_h = m.line_height_pt()
            hang_w = max(10.0, box.width_pt - BULLET_HANG_PT)
            raw_h = Emu(shape.height).pt
            top_inset = max(0.0, (raw_h - box.height_pt) / 2.0)
            base = Emu(shape.top).pt + top_inset

            print(f"  {shape.name}  top={Emu(shape.top).pt:.1f}pt  usable={box.height_pt:.1f}pt  "
                  f"top_inset={top_inset:.1f}pt  content starts {base:.1f}pt")

            # tables in the same horizontal half as this frame
            mine = [t for t in tables
                    if abs(Emu(t.left).inches - Emu(shape.left).inches) < 1.0]

            y = 0.0
            for i, p in enumerate(shape.text_frame.paragraphs):
                ptext = p.text or ""
                sizes = [r.font.size.pt for r in p.runs if r.font.size is not None]
                pitch = (max(sizes) if sizes else FONT_PT) * 1.2
                gap = p.space_after.pt if p.space_after is not None else 0.0
                if not ptext.strip():
                    a, b = base + y, base + y + pitch
                    print(f"    p{i:<3} BLANK {a:7.1f} ->{b:7.1f}   font={max(sizes) if sizes else FONT_PT:.2f}")
                    y += pitch + gap
                    continue
                n = max(1, len(m.wrap(
                    ptext, hang_w,
                    first_line_width_pt=box.width_pt if ptext.lstrip().startswith("■") else None,
                )))
                a, b = base + y, base + y + n * line_h
                hits = [t for t in mine
                        if Emu(t.top).pt < b - 0.5
                        and Emu(t.top).pt + Emu(t.height).pt > a + 0.5]
                flag = ""
                if hits:
                    problems += 1
                    t = hits[0]
                    tt, tb = Emu(t.top).pt, Emu(t.top).pt + Emu(t.height).pt
                    ov = min(b, tb) - max(a, tt)
                    flag = f"   <-- COLLIDES with {t.name!r} by {ov:.1f}pt"
                print(f"    p{i:<3} TEXT  {a:7.1f} ->{b:7.1f}   {n}L gap={gap:.0f}  "
                      f"{ptext[:30]!r}{flag}")
                y += n * line_h + gap
            print(f"    content ends {base + y:.1f}pt  (box bottom {base + box.height_pt:.1f}pt)")
            print()

    print("=" * 74)
    print(f"{problems} real collision(s) found."
          if problems else "No collision: every table sits in blank space only.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
