"""PPTX export inspection tool — run this against a FILE ALREADY EXPORTED by
the pipeline (Streamlit UI or CLI test script) to catch the layout classes of
bug that used to require opening PowerPoint by eye:

  1. L/R column collision — a BS/IS content page rendered as one full-width
     commentary box instead of two side-by-side halves (template only had a
     single textMainBullets box for a page that needed two logical slots).
  2. Table/commentary overlap — the embedded financial statement table drawn
     on top of (instead of beside) the commentary text box.
  3. Overflow risk — commentary text that, at production font size/spacing,
     wraps to more lines than the box has vertical room for (same Pillow
     glyph-width measurement the packer itself uses, so a "no overflow" here
     means the packer's own capacity estimate agrees with this independent
     re-check).
  4. Fill ratio (utilisation) — how full each slot ended up, matching the
     packer's own target_fill_min_ratio concept; flags slots that are
     suspiciously empty EXCEPT the last slot of a statement (a lighter tail
     page is normal, not a bug).

Everything here is DETERMINISTIC, geometry + font-metrics only — no AI calls,
no PowerPoint installation required. It does NOT replace opening the real
PPTX in PowerPoint at least once per template/font change (subtle rendering
quirks — kerning, hinting, OS-level font substitution — are still only fully
verifiable there), but it catches the structural classes of bug automatically
so that step becomes a spot-check instead of a full manual read-through.

Usage:
    python inspect_pptx.py path/to/exported.pptx
    python inspect_pptx.py path/to/exported.pptx --config fdd_utils/config.yml
    python inspect_pptx.py path/to/a/folder/          # loops every .pptx in it, prints a summary
"""
from __future__ import annotations

import argparse
import re
import sys
from dataclasses import dataclass
from pathlib import Path
from typing import Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.util import Emu

from fdd_utils.financial_common import load_yaml_file
from fdd_utils.text_metrics import get_measurer, text_box_from_shape

DEFAULT_CONFIG_CANDIDATES = ["fdd_utils/config.yml", "fdd_utils/config.example.yml"]

# Mirrors _fill_text_main_bullets_with_category_and_key -- the function that
# ACTUALLY sets a textMainBullets run's formatting -- not get_font_size_for_
# text/get_line_spacing_for_text/get_space_after_for_text/get_space_before_
# for_text, which belong to a separate, legacy code path (_fill_content_shape,
# only reached from the unused markdown generate() flow) and were never the
# values really applied to a live commentary bullet. Caught via a real
# Windows client-metrics export + inspect_single_slot.py: assuming a 9pt
# inter-paragraph gap (this file's old PARA_GAP_CHI, itself copied from a
# since-fixed pptx.py bug) against a real hardcoded 3pt gap inflated "capacity
# used" by roughly 30% -- a box the user could still type 5-7 more lines into
# was already being reported as 94% full.
FONT_SIZE_ENG = 9.0
FONT_SIZE_CHI = 9.0
LINE_SPACING_ENG = 1.0
LINE_SPACING_CHI = 1.0
# 2.2, not 3.0 (2026-08-04) -- kept in lockstep with fdd_utils/pptx.py's
# _real_para_gap_pt, which moved to 2.2 the same day (back-solved from real
# empirical spare-capacity measurements on 2 independent real boxes). This
# file already carries a scar from the OPPOSITE class of drift (see the
# comment above: an old 9pt copy here overstated fill by ~30% against a
# real 3pt render) -- if pptx.py's value ever moves again, update this one
# in the SAME commit, not a follow-up, or this tool starts silently
# reporting fill%/overflow against a formula generation no longer uses.
PARA_GAP_ENG = 3.0
PARA_GAP_CHI = 3.0
MIN_FILL_RATIO_WARN = 0.40  # below this on a non-last slot -> utilisation flag
# Kept in lockstep with fdd_utils/pptx.py's _TAIL_OVERFLOW_TOLERANCE_UNITS;
# overwritten from config at inspect time. See the `overflow=` assignment.
_OVERFLOW_TOLERANCE_LINES = 2.0


def _is_chinese_text(text: str, threshold: float = 0.3) -> bool:
    # Predominantly-Chinese (mirrors fdd_utils.financial_common's
    # contains_predominantly_chinese_text / pptx.py's account-level
    # is_chinese flag), not "contains any CJK character" -- an English
    # commentary box that merely names a Chinese counterparty/person
    # still wraps as Latin-script prose in the real render, and measuring
    # it with CJK metrics here would misreport its true fill.
    if not text:
        return False
    chinese_chars = sum(1 for ch in text if "一" <= ch <= "鿿")
    return (chinese_chars / len(text)) > threshold


def _slot_of(shape_name: str) -> str:
    name = (shape_name or "").lower()
    if name.endswith("_l"):
        return "L"
    if name.endswith("_r"):
        return "R"
    return "single"


def _load_config(path: Optional[str]) -> dict:
    candidates = [path] if path else DEFAULT_CONFIG_CANDIDATES
    for candidate in candidates:
        if not candidate:
            continue
        try:
            cfg = load_yaml_file(candidate)
            if cfg:
                return cfg
        except (FileNotFoundError, OSError):
            continue
    return {}


@dataclass
class ShapeInfo:
    name: str
    slot: str
    left_in: float
    top_in: float
    width_in: float
    height_in: float
    text: str
    n_chars: int
    capacity_lines: float
    wrapped_lines: int
    content_units: float
    fill_ratio: float
    overflow: bool
    font_sizes_pt: tuple


def _actual_font_sizes_pt(shape) -> tuple:
    """Every DISTINCT font size actually set on a run in this shape's text,
    in points, sorted. Reads the real saved XML rather than assuming the
    nominal size a caller intended -- the only way a report like "the text
    size looks wrong" is verifiable from pasted text output instead of a
    screenshot. `None` (a run with no explicit size, inheriting from the
    placeholder/theme) is reported as the string 'inherited' rather than
    silently dropped, since an unexpectedly-inherited size is itself
    sometimes the actual bug."""
    sizes = set()
    try:
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                sizes.add(run.font.size.pt if run.font.size is not None else "inherited")
    except Exception:
        pass
    return tuple(sorted(sizes, key=lambda v: (isinstance(v, str), v)))


#: Template boilerplate that must never survive into a deliverable. A real
#: export leaked "Placeholder – placeholder" into every table column: the
#: renderer had stopped clearing the slot's frame, so the text both showed
#: in the deck AND occupied a full line that the table floated over it was
#: positioned without -- clipping the real lead-in by 5.8pt. The template is
#: gitignored and per-machine, so a local template without this text cannot
#: reproduce it; only a check on the exported file can.
_PLACEHOLDER_PATTERNS = (
    "placeholder – placeholder",
    "placeholder - placeholder",
    "click to edit",
    "lorem ipsum",
)


def _leaked_placeholder(text: str) -> Optional[str]:
    low = (text or "").lower()
    for pat in _PLACEHOLDER_PATTERNS:
        if pat in low:
            return pat
    return None


#: Legal-form suffixes the reference deck never writes in running prose --
#: it uses the short form and reserves the full legal name for the first
#: mention of a specific contract counterparty. A repeated full name costs
#: most of a line each time, which is why this is worth flagging rather
#: than leaving to the eye.
_FULL_LEGAL_NAME_MARKERS = (
    "有限公司", "有限责任公司", "股份有限公司", "会计师事务所",
    "Co., Ltd", "Co.,Ltd", "Company Limited",
)

#: NOTE: a "line would start with punctuation" check was written here and
#: REMOVED. It could never fire: text_metrics' own wrapper already avoids
#: putting closing punctuation at a line start, so the check was a
#: permanently-green light while the real deck showed the defect.
#:
#: That difference is itself worth recording -- OUR line counts assume
#: 禁则处理 is applied, and PowerPoint's render (before eaLnBrk/hangingPunct/
#: <a:ea> were set) did not apply it. So this class of defect is NOT
#: observable from the exported file at all; it can only be seen in a real
#: render, and any future attempt must be verified there.


def _full_legal_names(text: str) -> List[str]:
    """Distinct full-legal-name occurrences, with a little context so the
    report names the actual party rather than just the suffix."""
    found = []
    for marker in _FULL_LEGAL_NAME_MARKERS:
        start = 0
        while True:
            i = (text or "").find(marker, start)
            if i < 0:
                break
            snippet = (text[max(0, i - 12):i + len(marker)]).strip()
            if snippet not in found:
                found.append(snippet)
            start = i + 1
    return found


def _bbox_overlap(a, b) -> bool:
    """True if two (left, top, width, height) EMU boxes overlap by area."""
    ax1, ay1, ax2, ay2 = a[0], a[1], a[0] + a[2], a[1] + a[3]
    bx1, by1, bx2, by2 = b[0], b[1], b[0] + b[2], b[1] + b[3]
    return ax1 < bx2 and bx1 < ax2 and ay1 < by2 and by1 < ay2


def inspect_pptx(pptx_path: str, config: dict, *, quiet: bool = False, dump_text: bool = False) -> dict:
    """Runs the full layout inspection and returns a structured summary
    (used both by this file's own CLI and by inspect_databook.py's combined
    export+inspect flow). Pass quiet=True to suppress the per-slide print
    lines and keep only the final summary print (caller still gets full
    detail back in the returned dict either way)."""
    _print = (lambda *a, **k: None) if quiet else print
    packing_cfg = ((config.get("pptx") or {}).get("commentary_packing") or {})
    global _OVERFLOW_TOLERANCE_LINES
    _OVERFLOW_TOLERANCE_LINES = float(packing_cfg.get("tail_overflow_tolerance_lines", 2.0) or 0.0)
    metrics_eng = packing_cfg.get("font_metrics_path_eng") or "fdd_utils/font_metrics/arial_eng.json"
    metrics_chi = packing_cfg.get("font_metrics_path_chi") or "fdd_utils/font_metrics/msyh_chi.json"
    family_eng = packing_cfg.get("font_family_eng") or "Arial"
    family_chi = packing_cfg.get("font_family_chi") or "Microsoft YaHei"

    eng_measurer = get_measurer(family_eng, FONT_SIZE_ENG, is_cjk=False,
                                 line_spacing=LINE_SPACING_ENG, metrics_path=metrics_eng)
    chi_measurer = get_measurer(family_chi, FONT_SIZE_CHI, is_cjk=True,
                                 line_spacing=LINE_SPACING_CHI, metrics_path=metrics_chi)
    _print(f"Measurement source: ENG={eng_measurer.source}  CHI={chi_measurer.source}")
    if eng_measurer.source == "system-font" or chi_measurer.source == "system-font":
        _print("⚠️  Falling back to a system-installed font for at least one language — "
               "this machine's font may not match the client's PowerPoint font. Check "
               f"that {metrics_eng!r} / {metrics_chi!r} exist and are readable.")

    prs = Presentation(pptx_path)
    _print(f"\nTotal slides: {len(prs.slides)}\n")

    total_warnings = 0
    warning_details: List[str] = []
    slide_reports: List[dict] = []

    for slide_idx, slide in enumerate(prs.slides):
        commentary_shapes = [
            s for s in slide.shapes
            if "textmainbullets" in (getattr(s, "name", "") or "").lower() and s.has_text_frame
        ]
        table_shapes = [s for s in slide.shapes if getattr(s, "has_table", False)]
        summary_shapes = [
            s for s in slide.shapes
            if "summary" in (getattr(s, "name", "") or "").lower() and s.has_text_frame
        ]
        has_summary = bool(summary_shapes)

        if not commentary_shapes:
            continue

        # A presentation-table account (_render_table_accounts_stack in
        # pptx.py) writes its lead-in and post-table explanation into fresh,
        # auto-named textboxes (python-pptx assigns "TextBox N") rather than
        # the well-known textMainBullets*/coSummaryShape names this file
        # otherwise matches by -- so without this, those two boxes (exactly
        # the ones sized against real AI text, i.e. exactly the ones that
        # can actually overflow) were invisible to this check entirely; only
        # this instance's OWN sizing math was ever "verifying" them. Plain
        # ordinary accounts flowed into a table's own leftover space (e.g.
        # 投资收益/营业外支出 after 财务费用's table -- see pptx.py's
        # _append_table_accounts_to_distribution trailing_items) are the
        # same blind spot one level further down the column.
        #
        # Matched by geometry, not name: same column (left) as a table AND
        # vertically touching it -- either ending right where the table
        # starts (the lead-in, single check against the table's own top) or
        # part of an unbroken chain starting at the table's own bottom
        # (source line -> explanation -> any number of trailing plain
        # accounts, each one's top touching the PREVIOUS shape's bottom).
        # The first hop of that chain uses a looser tolerance (the
        # source-line box's own height sits between the table and the
        # explanation); every hop after that uses the tight fixed gap
        # _TABLE_GAP_BELOW_PT actually renders with. A same-column-only
        # check (no vertical constraint) first tried here also matched the
        # slide title and the small fixed "Commentary" header label purely
        # because they happen to share the table's left -- both sit far
        # above the table, nowhere near touching it, so the vertical check
        # excludes them.
        # Dedup/identity is keyed on shape_id (the stable id python-pptx
        # stores in the XML itself), NOT Python's id() -- shape objects are
        # freshly re-wrapped on every separate `slide.shapes` iteration, so
        # two accesses of "the same" shape are different Python objects
        # with different id()s (confirmed empirically); id()-based tracking
        # silently failed to prevent a shape from being matched twice when
        # it happened to sit within tolerance of two different tables (e.g.
        # one account's lead-in sitting just below a PRECEDING account's
        # table in a shared slot, close enough to also look like that
        # table's own trailing explanation).
        LEFT_TOLERANCE_EMU = int(0.05 * 914400)
        ABOVE_TOLERANCE_EMU = int(2 * 12700)      # lead-in's bottom == table's top, ~exactly
        FIRST_BELOW_TOLERANCE_EMU = int(20 * 12700)   # table's bottom -> source line -> explanation
        CHAIN_TOLERANCE_EMU = int(6 * 12700)          # explanation/account bottom -> next trailing account
        table_stack_shapes = []
        _seen_ids = {s.shape_id for s in commentary_shapes}

        def _find_touching(reference_left, floor_bottom, tolerance):
            for s in slide.shapes:
                if s.shape_id in _seen_ids:
                    continue
                if not getattr(s, "has_text_frame", False) or getattr(s, "has_table", False):
                    continue
                if s.left is None or s.top is None or s.height is None:
                    continue
                if abs(s.left - reference_left) > LEFT_TOLERANCE_EMU:
                    continue
                if not (0 <= (s.top - floor_bottom) <= tolerance):
                    continue
                _text = (s.text_frame.text or "").strip()
                if not _text or _text.startswith(("资料来源", "Source:")):
                    continue
                return s
            return None

        for _t in table_shapes:
            if _t.left is None or _t.top is None or _t.height is None:
                continue
            t_bottom = _t.top + _t.height

            for s in slide.shapes:
                if s.shape_id in _seen_ids or s.shape_id == _t.shape_id:
                    continue
                if not getattr(s, "has_text_frame", False) or getattr(s, "has_table", False):
                    continue
                if s.left is None or s.top is None or s.height is None:
                    continue
                if abs(s.left - _t.left) > LEFT_TOLERANCE_EMU:
                    continue
                s_bottom = s.top + s.height
                if 0 <= (_t.top - s_bottom) <= ABOVE_TOLERANCE_EMU:
                    _text = (s.text_frame.text or "").strip()
                    if _text:
                        _seen_ids.add(s.shape_id)
                        table_stack_shapes.append(s)
                    break

            frontier_bottom = t_bottom
            tolerance = FIRST_BELOW_TOLERANCE_EMU
            while True:
                found = _find_touching(_t.left, frontier_bottom, tolerance)
                if found is None:
                    break
                _seen_ids.add(found.shape_id)
                table_stack_shapes.append(found)
                frontier_bottom = found.top + found.height
                tolerance = CHAIN_TOLERANCE_EMU

        _print(f"=== Slide {slide_idx + 1} ===  (table={bool(table_shapes)}  coSummaryShape={has_summary})")
        for s_shape in summary_shapes:
            s_text = s_shape.text_frame.text.strip()
            if s_text:
                if dump_text:
                    # Full text, untruncated -- a 120-char preview cuts an
                    # executive summary off mid-sentence, which is useless
                    # for actually reading it for tone/wording.
                    _print(f"  [{s_shape.name}] executive summary ({len(s_text)} chars): {s_text!r}")
                else:
                    preview = s_text[:120] + ("..." if len(s_text) > 120 else "")
                    _print(f"  [{s_shape.name}] executive summary ({len(s_text)} chars): {preview!r}")
            else:
                _print(f"  ⚠️  [{s_shape.name}] executive summary shape is EMPTY (0 chars)")
                total_warnings += 1
                warning_details.append(f"Slide {slide_idx + 1}: executive summary ({s_shape.name}) is empty")

        # Presentation-table (subtable) width + wrap-risk check. Reports
        # each native table's own width against the widest sibling
        # commentary column on the same slide (so "is this full-width or
        # content-width" is a printable number, not something you have to
        # open PowerPoint to eyeball), and cross-checks every cell's REAL
        # measured text width (same glyph-metrics measurer everything else
        # in this file uses) against its assigned column's available width.
        # A cell that fails this WILL wrap when PowerPoint actually opens
        # the file and auto-grow its row past the nominal height set in the
        # XML -- this predicts that deterministically from the saved file
        # alone, without needing PowerPoint itself to render it.
        for _t in table_shapes:
            if _t.left is None or _t.width is None:
                continue
            tbl = _t.table
            n_rows, n_cols = len(tbl.rows), len(tbl.columns)
            col_widths_in = [round(Emu(c.width).inches, 3) for c in tbl.columns]
            table_width_in = round(Emu(_t.width).inches, 3)

            # Excluded tokens mirror pptx.py's _is_commentary_text_shape --
            # without this, the slide TITLE (which often shares the same
            # left edge as a table sitting under it) gets picked up as the
            # "sibling slot", reporting a nonsense full-slide-width ratio.
            _excluded_sibling_tokens = ("title", "summary", "table", "subtitle")
            sibling_widths_in = [
                round(Emu(s.width).inches, 3) for s in slide.shapes
                if s.shape_id != _t.shape_id and getattr(s, "has_text_frame", False)
                and not getattr(s, "has_table", False)
                and not any(tok in (s.name or "").lower() for tok in _excluded_sibling_tokens)
                and abs((s.left or 0) - _t.left) <= LEFT_TOLERANCE_EMU
                and s.width
            ]
            ratio_note = ""
            if sibling_widths_in:
                slot_w = max(sibling_widths_in)
                ratio_note = f"  (slot width ~{slot_w}in, table/slot={table_width_in / slot_w:.0%})"

            _print(f"  [table {_t.name!r}] left={Emu(_t.left).inches:.2f}in width={table_width_in}in "
                   f"{n_rows}x{n_cols} columns(in)={col_widths_in}{ratio_note}")

            all_text = " ".join(tbl.cell(r, 0).text for r in range(n_rows))
            table_is_chi = _is_chinese_text(all_text)
            fam = family_chi if table_is_chi else family_eng
            mpath = metrics_chi if table_is_chi else metrics_eng
            m_title = get_measurer(fam, 8.0, is_cjk=table_is_chi, metrics_path=mpath)
            m_header = get_measurer(fam, 7.5, is_cjk=table_is_chi, metrics_path=mpath)
            m_data = get_measurer(fam, 7.0, is_cjk=table_is_chi, metrics_path=mpath)
            CELL_PADDING_PT = 5.76  # matches _TABLE_CELL_PADDING_PT in pptx.py

            # Font sizes actually used are measured AT (not asserted to
            # equal) -- the overview grid table (_fill_table_placeholder)
            # and the presentation/subtable (_render_presentation_table)
            # legitimately use DIFFERENT font-size schemes (the overview
            # table density-tiers by row count, 7-8pt; the subtable is a
            # fixed 8.0/7.5/7.0 by row role) and this check can't reliably
            # tell which kind of table it's looking at from the saved file
            # alone. So rather than assert a specific expected value (which
            # produced a false positive against the overview table's own
            # legitimate 8pt header when first tried), measure with THIS
            # table's own row-0/row-1/data sizes and flag only INTERNAL
            # inconsistency -- e.g. some data-row cells at 7.0pt and others
            # at 6.5pt -- which is a real signal (partial shrink, a stray
            # style) regardless of what the "correct" absolute value is for
            # this particular table's own convention.
            title_avail_pt = sum(col_widths_in) * 72 - CELL_PADDING_PT  # title spans ALL columns (merged)
            wrap_risks = []
            sizes_by_role = {"title": set(), "header": set(), "data": set()}
            for r in range(n_rows):
                role = "title" if r == 0 else ("header" if r == 1 else "data")
                measurer_r = m_title if r == 0 else (m_header if r == 1 else m_data)
                for c in range(n_cols):
                    if r == 0 and c > 0:
                        continue  # title row is merged across all columns
                    cell = tbl.cell(r, c)
                    text = cell.text
                    if not text:
                        continue
                    needed_pt = measurer_r.text_width_pt(text)
                    avail_pt = title_avail_pt if r == 0 else (col_widths_in[c] * 72 - CELL_PADDING_PT)
                    if needed_pt > avail_pt + 0.5:
                        wrap_risks.append((r, c, text[:30], round(needed_pt, 1), round(avail_pt, 1)))
                    sizes_by_role[role].update(_actual_font_sizes_pt(cell))

            if wrap_risks:
                _print(f"  ⚠️  {len(wrap_risks)} cell(s) will likely WRAP (real text width exceeds its "
                       f"column's width) -- PowerPoint will auto-grow that row past its nominal height:")
                for r, c, text, needed_pt, avail_pt in wrap_risks[:10]:
                    _print(f"      row={r} col={c} text={text!r} needs={needed_pt}pt has={avail_pt}pt")
                total_warnings += 1
                warning_details.append(
                    f"Slide {slide_idx + 1}: table {_t.name!r} has {len(wrap_risks)} wrap-risk cell(s)"
                )
            else:
                _print(f"  ✅ table {_t.name!r}: no cell wrap risk detected.")

            _print(f"  [table {_t.name!r}] font sizes actually used -- "
                   f"title={sorted(sizes_by_role['title'], key=str)} "
                   f"header={sorted(sizes_by_role['header'], key=str)} "
                   f"data={sorted(sizes_by_role['data'], key=str)}")
            inconsistent_roles = [role for role, sizes in sizes_by_role.items() if len(sizes) > 1]
            if inconsistent_roles:
                _print(f"  ⚠️  font size is INCONSISTENT within the same row role: {inconsistent_roles} "
                       f"-- some cells that should match are rendering at different sizes.")
                total_warnings += 1
                warning_details.append(
                    f"Slide {slide_idx + 1}: table {_t.name!r} has inconsistent font size within {inconsistent_roles}"
                )

        all_checked_shapes = list(commentary_shapes) + table_stack_shapes
        infos: List[ShapeInfo] = []
        text_bands_by_shape: Dict[int, List[Tuple[float, float]]] = {}
        for shape in all_checked_shapes:
            if shape in table_stack_shapes:
                # No name to key off of -- fall back to which half of the
                # slide it's positioned in, purely for a readable label.
                slot = "tblL" if shape.left < prs.slide_width / 2 else "tblR"
            else:
                slot = _slot_of(shape.name)
            text = shape.text_frame.text
            box = text_box_from_shape(shape)
            is_chi = _is_chinese_text(text)
            measurer = chi_measurer if is_chi else eng_measurer
            line_h = measurer.line_height_pt()
            para_gap = PARA_GAP_CHI if is_chi else PARA_GAP_ENG
            std_lh = line_h + para_gap
            # Float, not int()-floored -- matches fdd_utils/pptx.py's
            # _calculate_max_lines_for_textbox (fixed in 5bbec43) and
            # inspect_single_slot.py (fixed in 3ceb2a4). This file's own
            # copy of the capacity formula was missed by both of those
            # fixes, so it kept discarding up to a full std_lh unit of
            # real box height and understating fill_ratio against what the
            # live packer actually computes.
            capacity = (box.height_pt / std_lh) if std_lh > 0 else 0.0

            # Bullet paragraphs (the "■ key - ..." line, rendered as p_key)
            # hang-indent: left_indent=0.15" / first_line_indent=-0.15", so
            # LINE 1 spans the box's FULL width and only WRAPPED continuation
            # lines (2+) are 10.8pt narrower. Continuation paragraphs (a
            # second '\n'-split commentary line, rendered as p_text with
            # first_line_indent=0) are narrow on EVERY line, no exception.
            # Mirror fdd_utils/pptx.py's _BULLET_HANGING_INDENT_PT /
            # first_line_width_pt so this independent re-check uses the same
            # effective width as the packer AND the real render.
            hang_w = max(10.0, box.width_pt - 10.8)

            # Literal wrapped-line count, for display only (chars=/wraps_to=).
            wrapped = measurer.wrap(text, hang_w) if text.strip() else []
            n_lines = len(wrapped)

            # Content cost in the SAME std_lh units as capacity: one para_gap
            # PER PARAGRAPH (not per wrapped physical line), mirroring
            # fdd_utils/pptx.py's _calculate_content_lines. Comparing a
            # literal physical-line count against a std_lh-unit capacity
            # is apples-to-oranges (std_lh bundles a full para_gap into every
            # "line", so it under-counts how many literal lines actually fit)
            # and produced false OVERFLOW RISK flags on ordinary multi-line
            # paragraphs before this was unit-matched.
            paras = [p for p in text.split("\n") if p.strip()] if text.strip() else []
            # A CATEGORY HEADER ("Expenses" / "流动资产") renders with
            # space_after = Pt(0), not Pt(3) -- see
            # _fill_text_main_bullets_with_category_and_key. fdd_utils/
            # pptx.py's own cost function already accounts for that by
            # handling the category separately; this file charged a gap
            # for every paragraph uniformly, so it over-counted every
            # lead-in box carrying a category line by a full 3pt.
            # A category line is one with no bullet/arrow marker.
            _MARKERS = ("■", "➢", "-", "•")

            # Walk the REAL paragraph objects, not text.split("\n"), because
            # a table's vertical space is now reserved as BLANK paragraphs
            # inside this same frame with the table floated over them (see
            # fdd_utils/pptx.py's _render_table_accounts_stack). Splitting
            # the frame's flat text dropped every one of those spacers, so a
            # column holding a big table measured as nearly empty -- a real
            # export reported 24% full for a column that is physically full.
            # The last spacer of a band is deliberately sized to a FRACTION
            # of a line, so each paragraph's own font size has to be read
            # rather than assumed.
            para_bands: List[Tuple[bool, float, float]] = []  # (is_blank, top_pt, bottom_pt)
            content_pt = 0.0
            try:
                real_paras = list(shape.text_frame.paragraphs)
            except Exception:
                real_paras = []
            if real_paras:
                _y = 0.0
                for p_obj in real_paras:
                    p_text = p_obj.text or ""
                    sizes = [r.font.size.pt for r in p_obj.runs if r.font.size is not None]
                    p_pitch = (max(sizes) if sizes else 9.0) * 1.2
                    p_gap = p_obj.space_after.pt if p_obj.space_after is not None else 0.0
                    if not p_text.strip():
                        para_bands.append((True, _y, _y + p_pitch))
                        content_pt += p_pitch + p_gap
                        _y += p_pitch + p_gap
                        continue
                    n = max(1, len(measurer.wrap(
                        p_text, hang_w,
                        first_line_width_pt=box.width_pt if p_text.lstrip().startswith("■") else None,
                    )))
                    h = n * line_h
                    para_bands.append((False, _y, _y + h))
                    content_pt += h + p_gap
                    _y += h + p_gap
            else:
                content_pt = sum(
                    len(measurer.wrap(
                        p, hang_w,
                        first_line_width_pt=box.width_pt if p.lstrip().startswith("■") else None,
                    )) * line_h + (para_gap if p.lstrip().startswith(_MARKERS) else 0.0)
                    for p in paras
                )
            # Same correction as fdd_utils/pptx.py's _calculate_content_lines:
            # the final paragraph's trailing space_after is invisible padding
            # at the bottom of the frame, not occupied height. Counting it
            # produced false OVERFLOW warnings -- a real 27-line box that
            # PowerPoint's own BoundHeight puts at 93.5% full was reported
            # here as 102%.
            # Only if the LAST paragraph actually received a gap above --
            # a box ending on a category line never had one to remove.
            if paras and paras[-1].lstrip().startswith(_MARKERS):
                content_pt -= para_gap
            content_units = (content_pt / std_lh) if std_lh > 0 else 0.0
            # Absolute vertical extent of every real-text paragraph, for the
            # table-overlap check further down. A table sitting over a BLANK
            # band is the intended layout, not a collision.
            _top_inset = max(0.0, (Emu(shape.height).pt - box.height_pt) / 2.0) if shape.height is not None else 3.6
            _base = (Emu(shape.top).pt if shape.top is not None else 0.0) + _top_inset
            text_bands_by_shape[id(shape._element)] = [
                (_base + a, _base + b) for is_blank, a, b in para_bands if not is_blank
            ]

            fill_ratio = (content_units / capacity) if capacity > 0 else 0.0
            infos.append(ShapeInfo(
                name=shape.name, slot=slot,
                left_in=Emu(shape.left).inches if shape.left is not None else -1,
                top_in=Emu(shape.top).inches if shape.top is not None else -1,
                width_in=Emu(shape.width).inches if shape.width is not None else -1,
                height_in=Emu(shape.height).inches if shape.height is not None else -1,
                text=text, n_chars=len(text), capacity_lines=capacity,
                wrapped_lines=n_lines, content_units=content_units, fill_ratio=fill_ratio,
                # The packer now DELIBERATELY lets a box protrude by up to
                # commentary_packing.tail_overflow_tolerance_lines rather than
                # take an ugly split (the project team accepts 1-2 lines
                # sticking out). Flagging at capacity + 0.0 would therefore
                # warn on every well-packed deck; only a protrusion BEYOND
                # that allowance is a real problem.
                overflow=content_units > capacity + _OVERFLOW_TOLERANCE_LINES,
                font_sizes_pt=_actual_font_sizes_pt(shape),
            ))

        for i, info in enumerate(infos):
            is_last_slot_on_slide = (i == len(infos) - 1)
            flags = []
            if info.overflow:
                flags.append("⚠️ OVERFLOW RISK")
            _leak = _leaked_placeholder(info.text)
            if _leak:
                flags.append(f"❌ TEMPLATE PLACEHOLDER LEAKED ({_leak!r})")
            _legal = _full_legal_names(info.text)
            if _legal:
                flags.append(f"📛 FULL LEGAL NAME x{len(_legal)}")
            if info.n_chars > 0 and info.fill_ratio < MIN_FILL_RATIO_WARN and not is_last_slot_on_slide:
                flags.append(f"📉 under-filled ({info.fill_ratio:.0%})")
            flag_str = ("  " + "  ".join(flags)) if flags else ""
            if flags:
                total_warnings += 1
                warning_details.append(f"Slide {slide_idx + 1} [{info.slot}] {info.name}: {', '.join(flags)}")
            _print(f"  [{info.slot:6s}] {info.name:24s} left={info.left_in:5.2f}in width={info.width_in:5.2f}in "
                   f"chars={info.n_chars:4d} capacity={info.capacity_lines:5.1f}L used={info.content_units:5.1f}L "
                   f"fill={info.fill_ratio:.0%} font_pt={list(info.font_sizes_pt)} "
                   f"(raw wraps_to={info.wrapped_lines:3d}L, NOT comparable to capacity){flag_str}")

        # 1. L/R collision: a page with no table/summary (i.e. NOT the
        # designed single-column table slide) but only a single unsplit slot.
        slots_seen = {info.slot for info in infos}
        if slots_seen == {"single"} and not table_shapes and not has_summary:
            _print("  ❌ L/R COLLISION SUSPECTED — only a 'single' (full-width) commentary "
                   "slot on a page with no table/coSummaryShape, i.e. this looks like an L/R "
                   "content page that collapsed into one box instead of two side-by-side halves.")
            total_warnings += 1
            warning_details.append(f"Slide {slide_idx + 1}: L/R collision suspected")

        # 2. Table/commentary bounding-box overlap.
        for table_shape in table_shapes:
            t_box = (table_shape.left, table_shape.top, table_shape.width, table_shape.height)
            for info, shape in zip(infos, all_checked_shapes):
                if shape is table_shape:
                    continue
                if info.n_chars == 0:
                    continue
                c_box = (shape.left, shape.top, shape.width, shape.height)
                if None in t_box or None in c_box:
                    continue
                if not _bbox_overlap(t_box, c_box):
                    continue
                # A bounding-box hit is EXPECTED now: a table's vertical space
                # is reserved as blank paragraphs inside the commentary frame
                # and the table is floated over them, so it is inside that
                # frame's box by construction. Only a table landing on a
                # paragraph that has REAL TEXT in it is a genuine collision.
                t_top = Emu(table_shape.top).pt
                t_bot = t_top + Emu(table_shape.height).pt
                bands = text_bands_by_shape.get(id(shape._element))
                if bands is not None:
                    hits = [(a, b) for a, b in bands if a < t_bot - 0.5 and b > t_top + 0.5]
                    if not hits:
                        continue  # sits in its own reserved blank band
                    _print(f"  ❌ TABLE OVERLAPS REAL TEXT — '{table_shape.name}' covers "
                           f"{len(hits)} paragraph(s) of '{info.name}' "
                           f"(table {t_top:.0f}-{t_bot:.0f}pt).")
                else:
                    _print(f"  ❌ TABLE/COMMENTARY OVERLAP — '{table_shape.name}' overlaps '{info.name}'.")
                total_warnings += 1
                warning_details.append(f"Slide {slide_idx + 1}: table overlaps '{info.name}'")

        slide_reports.append({
            "slide": slide_idx + 1, "table": bool(table_shapes), "coSummaryShape": has_summary,
            "shapes": [i.__dict__ for i in infos],
        })
        _print()

    _print("=" * 78)
    if not slide_reports and len(prs.slides) > 0:
        # Every slide got skipped at the "if not commentary_shapes: continue"
        # gate above -- meaning NOT ONE shape anywhere in this deck had a name
        # containing "textmainbullets". That's this project's own template
        # convention (fdd_utils/template.pptx); a PPTX that never matches it
        # anywhere is very likely NOT generated by this pipeline at all (e.g.
        # a client's own hand-authored report), not an empty/broken deck --
        # so "0 warnings" here means "nothing was actually inspected", not
        # "everything is clean". Surfaced clearly instead of silently
        # reporting a misleading ✅.
        _print("⚠️  0 of this deck's slide(s) matched this project's own template shape-naming")
        _print("   convention ('textMainBullets'/'coSummaryShape') -- NOTHING was actually")
        _print("   inspected below. This usually means the PPTX wasn't generated by this")
        _print("   pipeline (e.g. a client's own report). Use --dump-text for a generic")
        _print("   per-shape text dump instead of this geometry check.")
    elif total_warnings == 0:
        _print("✅ No layout warnings found across all slides.")
    else:
        _print(f"⚠️  {total_warnings} warning(s) found — see ❌/⚠️/📉 markers above.")
    _print("Reminder: this is a geometry + font-metrics check, not a substitute for opening")
    _print("the file in real PowerPoint at least once per template/font change.")

    return {
        "total_slides": len(prs.slides),
        "content_slides": len(slide_reports),
        "total_warnings": total_warnings,
        "warning_details": warning_details,
        "measurement_source_eng": eng_measurer.source,
        "measurement_source_chi": chi_measurer.source,
        "slide_reports": slide_reports,
    }


# A duplicate bullet whose only numeric content is "no balance in any period"
# language, with no comma-grouped/CNY figure anywhere, is usually NOT a
# copy-paste bug -- in a multi-entity deck, several different entities can
# genuinely carry a zero balance in the same account, and the AI writes the
# same near-boilerplate sentence for each. Flag these separately (informational)
# instead of alongside real duplicates (same account, same entity, real figures
# repeated verbatim -- which IS always a bug).
_ZERO_BALANCE_PHRASES = ("无余额", "未形成余额", "未发生", "均无", "无变动", "为零", "未产生")


def _is_likely_zero_balance_boilerplate(text: str) -> bool:
    if not any(phrase in text for phrase in _ZERO_BALANCE_PHRASES):
        return False
    if _ENG_CNY_INT_RE.search(text):
        return False
    if re.search(r"\d,\d{3}", text):
        return False
    return True


def _dump_all_shapes_generic(pptx_path: str) -> None:
    """--dump-text fallback for a PPTX that doesn't match this project's own
    template shape-naming convention anywhere (inspect_pptx() found 0
    content_slides) -- e.g. a client's own hand-authored report used as a
    wording/tone reference, not something generated by this pipeline. Every
    other function in this file keys off specific shape names
    ('textMainBullets'/'coSummaryShape'), so none of them can read a foreign
    deck's content at all; this dumps every text-bearing shape's raw text,
    by slide, with no naming assumption. No layout/duplicate/number-
    formatting checks apply here -- those all assume our own template."""
    print(
        "\nℹ️  Falling back to a generic per-shape text dump (no shape-name assumption) "
        "since this deck doesn't match this project's own template convention -- see the "
        "⚠️ note above. No layout/duplicate/number-formatting checks apply to this dump.\n"
    )
    print("=" * 78)
    print("  GENERIC TEXT DUMP (any shape with text, any name, per slide)")
    print("=" * 78)
    prs = Presentation(pptx_path)
    any_text = False
    for slide_idx, slide in enumerate(prs.slides):
        shapes_with_text = [
            s for s in slide.shapes if getattr(s, "has_text_frame", False) and s.text_frame.text.strip()
        ]
        if not shapes_with_text:
            continue
        print(f"\n--- Slide {slide_idx + 1} ---")
        for shape in shapes_with_text:
            any_text = True
            name = getattr(shape, "name", "") or "(unnamed)"
            print(f"  [{name}]")
            print(f"  {shape.text_frame.text}")
    if not any_text:
        print(
            "\n⚠️  No text found in ANY shape on ANY slide. This deck may be image-based "
            "(text baked into pictures rather than real PowerPoint text boxes), which no "
            "text-extraction tool (including this one) can read -- open it in PowerPoint "
            "to confirm before assuming this tool is broken."
        )


def _dump_text_and_check_duplicates(result: dict) -> int:
    """Prints every commentary shape's full text (not just the char-count
    summary line inspect_pptx() already prints), and scans every '■ '-led
    bullet across the WHOLE file for duplicates -- i.e. the same account's
    commentary appearing on more than one slide/slot. Returns the number of
    duplicate bullets found (0 = clean), NOT counting likely-benign
    zero-balance boilerplate (see _is_likely_zero_balance_boilerplate)."""
    print("\n" + "=" * 78)
    print("  FULL TEXT DUMP (per slide/shape)")
    print("=" * 78)
    bullet_locations: dict[str, list[str]] = {}
    for slide_report in result["slide_reports"]:
        slide_no = slide_report["slide"]
        for shape in slide_report["shapes"]:
            text = shape.get("text") or ""
            if not text.strip():
                continue
            print(f"\n--- Slide {slide_no} [{shape['name']}] ---")
            print(text)
            for line in text.split("\n"):
                stripped = line.strip()
                if stripped.startswith("■"):
                    # Key on the FULL line, not a truncated prefix -- a
                    # truncated key false-flags two DIFFERENT bullets (same
                    # account, different entity, same generic opening template
                    # but different trailing figures) as duplicates just
                    # because they share their first N characters.
                    bullet_locations.setdefault(stripped, []).append(f"slide {slide_no} [{shape['name']}]")

    print("\n" + "=" * 78)
    print("  DUPLICATE-BULLET SCAN (same bullet text appearing on >1 slide/slot)")
    print("=" * 78)
    all_duplicates = {k: v for k, v in bullet_locations.items() if len(v) > 1}
    real_duplicates = {k: v for k, v in all_duplicates.items() if not _is_likely_zero_balance_boilerplate(k)}
    benign_duplicates = {k: v for k, v in all_duplicates.items() if _is_likely_zero_balance_boilerplate(k)}
    if not all_duplicates:
        print("✅ No duplicate bullets found -- every account's commentary appears exactly once.")
    else:
        for key, locations in real_duplicates.items():
            print(f"  ❌ DUPLICATE: {key[:90]!r}{'...' if len(key) > 90 else ''}")
            for loc in locations:
                print(f"      - {loc}")
        for key, locations in benign_duplicates.items():
            print(f"  ℹ️  LIKELY BENIGN, zero-balance boilerplate repeated across different entities: "
                  f"{key[:90]!r}{'...' if len(key) > 90 else ''}")
            for loc in locations:
                print(f"      - {loc}")
    return len(real_duplicates)


# English "CNY<comma-grouped integer>" -- requires at least one comma group,
# which only ever appears on the exact-integer form ("CNY238,366"), never on
# the "CNY<X> million/thousand" decimal form ("CNY7.9 million" has no comma
# at all, so it can never match here regardless of its own decimal digit).
_ENG_CNY_INT_RE = re.compile(r"CNY\s?(\d{1,3}(?:,\d{3})+)\b")

# Currency-amount tokens, captured as a plain (possibly decimal) number so the
# VALUE can be checked for zero in Python rather than pattern-matching "0" as
# text -- matching text like "0" would also match the trailing ".0" inside an
# ordinary non-zero number such as "457.0万元" or "CNY570.0 million".
# Digit-group patterns accept comma-grouping ("3,000") as ONE number, and the
# negative lookbehind excludes a preceding comma too -- without this, a comma-
# grouped amount like "人民币3,000万元" gets mis-split: the digits after the
# comma ("000") pass the lookbehind on their own and get misread as a
# standalone zero-value match ("000万元" -> 0), even though the real number is
# 3000, not 0.
_GROUPED_DIGITS = r"(\d{1,3}(?:,\d{3})*(?:\.\d+)?)"
_NUMBER_WITH_UNIT_RES = [
    re.compile(rf"CNY\s?{_GROUPED_DIGITS}\s?(?:million|thousand)?"),
    re.compile(rf"(?<![\d.,]){_GROUPED_DIGITS}万元"),
    re.compile(rf"(?<![\d.,]){_GROUPED_DIGITS}亿元"),
    # Excludes only a following letter/digit (avoids matching "元" as part of a
    # longer alphanumeric token) -- NOT `(?!\S)`, which also blocked a match
    # right before Chinese punctuation like "。"/"，", i.e. the overwhelmingly
    # common case for how a real bullet sentence actually ends ("...为人民币0元。").
    re.compile(rf"人民币\s?{_GROUPED_DIGITS}元(?![A-Za-z0-9])"),
]


# Chinese sentence boundary, for the repeated-人民币 check below -- the rule
# (confirmed against redd_patterns.xlsx, the project's real-report reference
# corpus) is per-SENTENCE, not per-bullet: a multi-period list should state
# "人民币" once ("分别为人民币2,216万元、2,421万元及968万元"), not before every
# number in the list ("...人民币2,216万元、人民币2,421万元及人民币968万元").
_CHI_SENTENCE_SPLIT_RE = re.compile(r"[。！？；]")


def _find_zero_currency_mentions(text: str) -> List[str]:
    hits: List[str] = []
    for pattern in _NUMBER_WITH_UNIT_RES:
        for m in pattern.finditer(text):
            try:
                if float(m.group(1).replace(",", "")) == 0:
                    hits.append(m.group(0))
            except ValueError:
                continue
    return hits


# Phrases that cite the SOURCE of a fact rather than stating the fact.
# Found 5x in one real exported deck ("根据备注信息，该科目在2023年已全部转入
# 固定资产"). No prompt rule currently mentions them, unlike the repeated-人民币
# and zero-wording conventions above. They leak the pipeline's own data
# structure -- the databook's remarks column -- into a client deliverable,
# where an FDD consultant would simply assert the fact, or attribute it to
# management ("管理层表示...") which IS an established convention here.
_SOURCE_META_RE = re.compile(
    r"(根据|依据|按照)\s*(补充|管理层)?(备注|说明|备注信息)(信息|内容|说明)?|"
    r"(备注|说明)(显示|表明|提到|中提及|仅说明|未说明|未进一步说明)|"
    r"(as|per)\s+(stated\s+in|noted\s+in|the)\s+(the\s+)?(supplementary\s+)?remarks?\b|"
    r"according\s+to\s+the\s+(supplementary\s+)?(remarks?|notes?)\b",
    re.IGNORECASE,
)


def _check_number_formatting_and_zero_wording(result: dict) -> int:
    """Scans every '■ '-led bullet for three classes of issue flagged from
    real reports: (1) English sub-million CNY amounts more precise than the
    intended nearest-thousand rounding (e.g. 'CNY238,366' instead of
    'CNY238,000') -- these read as excessive, inconsistent-with-Chinese
    detail; (2) a literal zero-value currency mention that should have been
    reworded as 'nil'/'未发生' instead (e.g. 'CNY0', '人民币0.0万元'); (3) '人民币'
    repeated before every number in a Chinese multi-period list within the
    SAME sentence (e.g. '分别为人民币2,216万元、人民币2,421万元及人民币968万元')
    instead of stated once for the whole list -- confirmed against
    redd_patterns.xlsx, the project's real-report reference corpus, as a
    Chinese-specific convention (English correctly repeats 'CNY' per item).
    Returns the total number of flagged bullets (0 = clean)."""
    print("\n" + "=" * 78)
    print("  NUMBER-FORMATTING / ZERO-WORDING SCAN")
    print("=" * 78)
    print(
        "Flags three things per bullet: (a) English sub-million CNY amounts that\n"
        "aren't rounded to the nearest thousand (over-precise vs the Chinese\n"
        "report's own 万-unit rounding for the same figure), (b) a literal zero\n"
        "currency mention that should read as 'nil'/'未发生' instead, (c) '人民币'\n"
        "repeated before every number in one Chinese multi-period list sentence\n"
        "instead of stated once for the whole list. None of these is necessarily\n"
        "wrong on its own -- a genuinely sub-CNY10,000 amount is correctly exact,\n"
        "a materiality-threshold '0%' isn't a currency mention, and a single\n"
        "'人民币' per DIFFERENT sentence is fine -- so treat this as a worklist\n"
        "to skim, not an automatic fail.\n"
    )
    flagged = 0
    for slide_report in result["slide_reports"]:
        slide_no = slide_report["slide"]
        for shape in slide_report["shapes"]:
            text = shape.get("text") or ""
            for line in text.split("\n"):
                stripped = line.strip()
                if not stripped.startswith("■"):
                    continue
                label = stripped[:60]
                issues: List[str] = []

                for m in _ENG_CNY_INT_RE.finditer(stripped):
                    value = int(m.group(1).replace(",", ""))
                    if 10_000 <= value < 1_000_000 and value % 1000 != 0:
                        issues.append(f"over-precise amount {m.group(0)!r} (not rounded to nearest thousand)")

                for hit in _find_zero_currency_mentions(stripped):
                    issues.append(f"literal zero mention {hit!r} (should read as 'nil'/'未发生')")

                if "人民币" in stripped:
                    for sentence in _CHI_SENTENCE_SPLIT_RE.split(stripped):
                        count = sentence.count("人民币")
                        if count >= 2:
                            issues.append(
                                f"'人民币' repeated {count}x in one sentence (state it once for the "
                                f"whole list): {sentence.strip()[:150]!r}"
                            )

                for m in _SOURCE_META_RE.finditer(stripped):
                    issues.append(
                        f"source meta-reference {m.group(0)!r} -- state the fact directly; "
                        f"a deliverable shouldn't cite the databook's own remarks column"
                    )

                if issues:
                    flagged += 1
                    print(f"  Slide {slide_no} [{shape['name']}] {label!r}...")
                    for issue in issues:
                        print(f"      - {issue}")

    if not flagged:
        print("✅ No over-precise amounts, literal zero mentions, or repeated-人民币 sentences found.")
    return flagged


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("pptx_path", help="Path to an already-exported .pptx file, or a directory to scan "
                                       "for every .pptx in it (e.g. a folder of batch-exported decks)")
    ap.add_argument("--config", default=None, help="Path to config.yml (default: tries fdd_utils/config.yml then config.example.yml)")
    ap.add_argument("--dump-text", action="store_true",
                     help="Also print every commentary shape's full text and scan for (1) duplicate "
                          "bullets across the whole file (same account's commentary appearing on "
                          "more than one slide/slot), (2) over-precise English sub-million CNY amounts "
                          "(not rounded to the nearest thousand), (3) literal zero-value currency "
                          "mentions that should read as 'nil'/'未发生' instead, (4) '人民币' repeated "
                          "before every number in one Chinese multi-period list sentence.")
    args = ap.parse_args()

    config = _load_config(args.config)

    input_path = Path(args.pptx_path)
    if input_path.is_dir():
        pptx_files = sorted(input_path.glob("*.pptx"))
        if not pptx_files:
            print(f"No .pptx files found in {input_path}")
            return 1
    else:
        pptx_files = [input_path]

    total_warnings = 0
    total_duplicates = 0
    total_wording_flags = 0
    per_file_summary: List[tuple] = []
    for pptx_file in pptx_files:
        if len(pptx_files) > 1:
            print(f"\n{'=' * 90}\n{pptx_file.name}\n{'=' * 90}")
        result = inspect_pptx(str(pptx_file), config, dump_text=args.dump_text)
        duplicate_count = 0
        wording_flag_count = 0
        if args.dump_text:
            if result["content_slides"] == 0 and result["total_slides"] > 0:
                _dump_all_shapes_generic(str(pptx_file))
            else:
                duplicate_count = _dump_text_and_check_duplicates(result)
                wording_flag_count = _check_number_formatting_and_zero_wording(result)
        total_warnings += result["total_warnings"]
        total_duplicates += duplicate_count
        total_wording_flags += wording_flag_count
        per_file_summary.append((pptx_file.name, result["total_warnings"], duplicate_count, wording_flag_count))

    if len(pptx_files) > 1:
        print(f"\n{'=' * 90}\nSUMMARY ({len(pptx_files)} file(s))\n{'=' * 90}")
        for name, warnings, duplicates, wording_flags in per_file_summary:
            flags = []
            if warnings:
                flags.append(f"{warnings} layout warning(s)")
            if duplicates:
                flags.append(f"{duplicates} duplicate bullet(s)")
            if wording_flags:
                flags.append(f"{wording_flags} wording flag(s)")
            status = "⚠️ " + ", ".join(flags) if flags else "✅ clean"
            print(f"  {name}: {status}")

    return 1 if (total_warnings or total_duplicates or total_wording_flags) else 0


if __name__ == "__main__":
    sys.exit(main())
