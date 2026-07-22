#!/usr/bin/env python3
"""Combines the raw-data bridge extractor (extract_bridge_from_raw.py) with
the waterfall chart builder (bridge_chart_prototype.py): computes the
price/occupancy/days factor decomposition per phase per year-transition
directly from a raw AB-* tab (no pre-built '<entity>-量价桥图' helper tab
needed), then renders a real PPTX waterfall chart -- across EVERY AB- tab
in the workbook in one pass, covering the entities that don't have a
manually-built bridge tab (16 of the 17 real ones, per this session's
17-tab scan -- only 成都/AB-CD has one).

Decomposition formula -- confirmed bit-exact against AB-CD's OWN real
bridge-tab formulas (2024->2025 transition, all 3 phases) this session,
sequential substitution in price -> area -> days order:
    price_effect = (price_B - price_A) * area_A * days_A / 1000
    area_effect  = price_B * (area_B - area_A) * days_A / 1000
    days_effect  = price_B * area_B * (days_B - days_A) / 1000
Each phase contributes 3 factor bars per year-transition; the bridge starts
at year_A's total revenue (summed across all phases) and ends at year_B's.

The LATEST year in a real workbook is usually a partial/LTM period (e.g.
"2026" only has data through June). Comparing it directly against a full
prior calendar year would make the 'days effect' swing by roughly negative
half a year -- a period-length artifact, not a real change. The REAL
'成都-量价桥图' bridge tab's own second transition avoids this WITHOUT any
annualization/extrapolation assumption: it compares the last full calendar
year against a trailing-12-month (LTM) window ending at the latest
available month (e.g. "2025年收入" -> "2025年7月至2026年6月收入") -- both
sides then have the same ~365 days, using only actual monthly data already
on hand, no assumption about how the rest of the year will look. This
script reproduces that exact convention for the tail transition.

Usage:
    python generate_bridge_waterfall_batch.py "databooks/xx.xlsx"
        # every AB- tab: full-year-to-full-year transitions for all
        # complete years, plus one LTM transition at the tail if the
        # latest year is partial
    python generate_bridge_waterfall_batch.py "databooks/xx.xlsx" --tab AB-CD --validate
        # cross-checks BOTH real bridge blocks (2024->2025 and the LTM
        # 2025->Jul25-Jun26 transition) against AB-CD's own real bridge-tab
        # factor values (the "answer key")
"""
import argparse
import sys
from typing import Dict, List, Optional

from openpyxl import load_workbook

from extract_bridge_from_raw import (
    find_phase_blocks, extract_annual_series, extract_ltm_series,
    find_month_row, format_ltm_label, PhaseBlock,
)
from inspect_ab_tabs_structure import find_labeled_rows
from bridge_chart_prototype import BridgeItem, BridgeBlock, build_waterfall_chart

from pptx import Presentation
from pptx.util import Inches

_FACTOR_SUFFIX = {"price": "单价变动", "area": "出租率/面积变动", "days": "运营天数变动"}
_PARTIAL_YEAR_DAYS_THRESHOLD = 350


def _is_partial_year(days: float) -> bool:
    return 0 < days < _PARTIAL_YEAR_DAYS_THRESHOLD


def decompose_transition(phase_label: str, series_a: Dict[str, float], series_b: Dict[str, float]) -> List[BridgeItem]:
    price_a, price_b = series_a["unit_rent"], series_b["unit_rent"]
    area_a, area_b = series_a["area"], series_b["area"]
    days_a, days_b = series_a["days"], series_b["days"]

    price_effect = (price_b - price_a) * area_a * days_a / 1000
    area_effect = price_b * (area_b - area_a) * days_a / 1000
    days_effect = price_b * area_b * (days_b - days_a) / 1000

    return [
        BridgeItem(label=f"{phase_label}{_FACTOR_SUFFIX['price']}", kind="delta", value=price_effect),
        BridgeItem(label=f"{phase_label}{_FACTOR_SUFFIX['area']}", kind="delta", value=area_effect),
        BridgeItem(label=f"{phase_label}{_FACTOR_SUFFIX['days']}", kind="delta", value=days_effect),
    ]


def find_year_days_rows(ws_values) -> Dict[str, Optional[int]]:
    labeled = find_labeled_rows(ws_values)
    year_row = days_row = None
    for r in sorted(labeled):
        cats = {cat for _, _, cat in labeled[r]}
        if year_row is None and "period_year" in cats:
            year_row = r
        if days_row is None and "period_days" in cats:
            days_row = r
    return {"year_row": year_row, "days_row": days_row}


def _assemble_bridge(blocks: List[PhaseBlock], series_a: List[Dict[str, float]], series_b: List[Dict[str, float]],
                      start_label: str, end_label: str) -> BridgeBlock:
    total_a = sum(s["revenue_k"] for s in series_a)
    total_b = sum(s["revenue_k"] for s in series_b)

    items: List[BridgeItem] = [BridgeItem(label=start_label, kind="total", value=total_a)]
    for block, sa, sb in zip(blocks, series_a, series_b):
        items.extend(decompose_transition(block.label, sa, sb))
    items.append(BridgeItem(label=end_label, kind="total", value=total_b))

    reconstructed = total_a + sum(it.value for it in items[1:-1])
    # A real client-data residual is expected here (confirmed in AB-CD's own
    # notes: revenue is "all-in" but leased area excludes some pallet-priced
    # space, so price*area*days never reconstructs revenue PERFECTLY) --
    # this is informational, not a reason to withhold the chart.
    check_ok = abs(reconstructed - total_b) < max(5.0, abs(total_b) * 0.02)
    return BridgeBlock(header_row=-1, label_col=-1, base_col=-1, change_col=-1, items=items, check_ok=check_ok)


def build_bridge_for_transition(ws_values, blocks: List[PhaseBlock], year_row: int, days_row: int,
                                 year_a: int, year_b: int, start_label: str, end_label: str) -> Optional[BridgeBlock]:
    all_series = [extract_annual_series(ws_values, b, year_row, days_row) for b in blocks]
    if any(year_a not in s or year_b not in s for s in all_series):
        return None
    series_a = [s[year_a] for s in all_series]
    series_b = [s[year_b] for s in all_series]
    return _assemble_bridge(blocks, series_a, series_b, start_label, end_label)


def build_ltm_bridge(ws_values, blocks: List[PhaseBlock], year_row: int, month_row: int, days_row: int,
                      last_full_year: int, end_year: int, end_month: int) -> Optional[BridgeBlock]:
    all_series_a = [extract_annual_series(ws_values, b, year_row, days_row) for b in blocks]
    if any(last_full_year not in s for s in all_series_a):
        return None
    series_a = [s[last_full_year] for s in all_series_a]
    series_b = [extract_ltm_series(ws_values, b, year_row, month_row, days_row, end_year, end_month) for b in blocks]
    if any(s is None for s in series_b):
        return None
    start_label = f"{last_full_year}年收入"
    end_label = format_ltm_label(end_year, end_month)
    return _assemble_bridge(blocks, series_a, series_b, start_label, end_label)


# Hardcoded from the REAL '成都-量价桥图' bridge tab's own values (already
# confirmed this session against the real file) -- the "answer key" for
# --validate, covering BOTH real bridge blocks.
_AB_CD_EXPECTED_FACTORS_2024_2025 = {
    "干仓": {"price": 2059.545458900859, "area": 1463.5320655100995, "days": 3298.960105589041},
    "综合楼": {"price": 66.9393853357857, "area": -55.45101747277201, "days": 29.103442136986306},
    "冷库": {"price": 0.0, "area": 0.0, "days": 3061.83285},
}
_AB_CD_EXPECTED_FACTORS_LTM = {
    "干仓": {"price": 1185.257478889836, "area": 1082.3849911101638, "days": 0.0},
    "综合楼": {"price": -0.2670081709406935, "area": 19.893338170940705, "days": 0.0},
    "冷库": {"price": -614.8229452141594, "area": 818.0862552141591, "days": 0.0},
}


def _fmt_series(label: str, s: Dict[str, float]) -> str:
    return (f"{label}: revenue={s['revenue_k']:,.2f}k area={s['area']:,.2f} "
            f"days={s['days']:.0f} unit_rent={s['unit_rent']:.4f}")


def dump_transition_detail(blocks: List[PhaseBlock], series_a: List[Dict[str, float]],
                            series_b: List[Dict[str, float]], label_a: str, label_b: str):
    for block, sa, sb in zip(blocks, series_a, series_b):
        print(f"    [{block.label}]")
        print(f"      {_fmt_series(label_a, sa)}")
        print(f"      {_fmt_series(label_b, sb)}")
        items = decompose_transition(block.label, sa, sb)
        for it in items:
            print(f"      -> {it.label}: {it.value:,.2f}k")


def _validate_against(bridge: BridgeBlock, expected: Dict[str, Dict[str, float]], label: str) -> bool:
    all_ok = True
    for phase_name, exp in expected.items():
        got_items = [it for it in bridge.items if it.label.startswith(phase_name)]
        if len(got_items) < 3:
            print(f"      ⚠️  {phase_name}: expected 3 factor items, found {len(got_items)}")
            all_ok = False
            continue
        got = {"price": got_items[0].value, "area": got_items[1].value, "days": got_items[2].value}
        for factor in ("price", "area", "days"):
            tol = max(0.5, abs(exp[factor]) * 0.005)
            ok = abs(got[factor] - exp[factor]) <= tol
            all_ok = all_ok and ok
            status = "✅" if ok else "❌"
            print(f"      {status} [{label}] {phase_name} {factor}: got={got[factor]:.4f} expected={exp[factor]:.4f}")
    return all_ok


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("path", help="path to the databook .xlsx")
    ap.add_argument("--tab", default=None, help="only process this one AB- tab (default: every AB- tab)")
    ap.add_argument("--validate", action="store_true",
                     help="cross-check --tab AB-CD's 2024->2025 AND the tail LTM transition against "
                          "its own real bridge-tab factor values")
    ap.add_argument("--dump-detail", action="store_true",
                     help="print each phase's raw (revenue/area/days/unit_rent) for both sides of "
                          "every transition, plus the resulting 3 factor values -- use this to trace "
                          "a chart number that looks wrong back to its exact source inputs")
    ap.add_argument("--skip-partial", action="store_true",
                     help="skip the tail transition entirely instead of building an LTM window for it "
                          "(e.g. if monthly data isn't available to construct one)")
    ap.add_argument("--out", default="bridge_waterfall_batch_output.pptx", help="output pptx path")
    args = ap.parse_args()

    print(f"Loading {args.path!r}...")
    wb_values = load_workbook(args.path, data_only=True)
    all_sheets = wb_values.sheetnames
    ab_tabs = [s for s in all_sheets if s.startswith("AB-") or s.startswith("AB")]
    if args.tab:
        ab_tabs = [s for s in ab_tabs if s == args.tab]
        if not ab_tabs:
            print(f"❌ tab {args.tab!r} not found among AB- sheets.")
            return 1

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    slides_added = 0
    all_ok = True
    for tab in ab_tabs:
        ws = wb_values[tab]
        yd = find_year_days_rows(ws)
        if not yd["year_row"] or not yd["days_row"]:
            print(f"--- {tab}: ⚠️ no Year/Days row found, skipping ---")
            continue
        blocks = find_phase_blocks(ws)
        if not blocks:
            print(f"--- {tab}: ⚠️ no phase blocks found, skipping ---")
            continue

        all_series = {b.label: extract_annual_series(ws, b, yd["year_row"], yd["days_row"]) for b in blocks}
        years = sorted({y for series in all_series.values() for y in series.keys()})
        print(f"--- {tab}: {len(blocks)} phase(s) [{', '.join(b.label for b in blocks)}], years {years} ---")
        if not years:
            continue

        max_year = years[-1]
        max_year_days = max((all_series[b.label].get(max_year, {}).get("days", 0) for b in blocks), default=0)
        tail_is_partial = _is_partial_year(max_year_days)
        # full-year-to-full-year transitions cover everything EXCEPT a
        # partial tail year, which gets the LTM treatment below instead.
        full_years = years[:-1] if tail_is_partial else years
        transitions = list(zip(full_years, full_years[1:]))

        for year_a, year_b in transitions:
            bridge = build_bridge_for_transition(
                ws, blocks, yd["year_row"], yd["days_row"], year_a, year_b,
                start_label=f"{year_a}年收入", end_label=f"{year_b}年收入",
            )
            if bridge is None:
                continue
            check_msg = "✅" if bridge.check_ok else "⚠️ residual > tolerance"
            print(f"    {year_a}->{year_b}: start={bridge.items[0].value:,.1f}k "
                  f"end={bridge.items[-1].value:,.1f}k  {check_msg}")

            if args.dump_detail:
                series_a = [all_series[b.label][year_a] for b in blocks]
                series_b = [all_series[b.label][year_b] for b in blocks]
                dump_transition_detail(blocks, series_a, series_b, f"{year_a} (raw)", f"{year_b} (raw)")

            if args.validate and tab == "AB-CD" and (year_a, year_b) == (2024, 2025):
                all_ok = _validate_against(bridge, _AB_CD_EXPECTED_FACTORS_2024_2025, "2024->2025") and all_ok

            slide = prs.slides.add_slide(blank_layout)
            title = f"{tab}: {year_a}年→{year_b}年 收入量价桥图"
            build_waterfall_chart(slide, bridge, title,
                                   left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(6.3))
            slides_added += 1

        if tail_is_partial:
            last_full_year = full_years[-1] if full_years else None
            month_row = find_month_row(ws)
            if not args.skip_partial and last_full_year is not None and month_row is not None:
                # anchor the LTM window at the latest (year, month) actually present
                latest_month = None
                for m in range(12, 0, -1):
                    days = 0
                    for b in blocks:
                        s = extract_ltm_series(ws, b, yd["year_row"], month_row, yd["days_row"], max_year, m, window=1)
                        if s:
                            days = max(days, s["days"])
                    if days > 0:
                        latest_month = m
                        break
                if latest_month is not None:
                    bridge = build_ltm_bridge(ws, blocks, yd["year_row"], month_row, yd["days_row"],
                                               last_full_year, max_year, latest_month)
                    if bridge is not None:
                        print(f"    {last_full_year}->LTM({format_ltm_label(max_year, latest_month)}): "
                              f"start={bridge.items[0].value:,.1f}k end={bridge.items[-1].value:,.1f}k "
                              f"{'✅' if bridge.check_ok else '⚠️ residual > tolerance'}")

                        if args.dump_detail:
                            series_a = [all_series[b.label][last_full_year] for b in blocks]
                            series_b = [extract_ltm_series(ws, b, yd["year_row"], month_row, yd["days_row"],
                                                            max_year, latest_month) for b in blocks]
                            dump_transition_detail(blocks, series_a, series_b, f"{last_full_year} (raw)",
                                                    f"LTM ending {max_year}-{latest_month:02d}")

                        if args.validate and tab == "AB-CD" and last_full_year == 2025 and max_year == 2026:
                            all_ok = _validate_against(bridge, _AB_CD_EXPECTED_FACTORS_LTM, "LTM") and all_ok

                        slide = prs.slides.add_slide(blank_layout)
                        title = f"{tab}: {last_full_year}年→LTM 收入量价桥图"
                        build_waterfall_chart(slide, bridge, title,
                                               left=Inches(0.5), top=Inches(0.5), width=Inches(12.3), height=Inches(6.3))
                        slides_added += 1
                    else:
                        print(f"    ⚠️ tail year {max_year} is partial ({max_year_days:.0f}d) but LTM window "
                              f"couldn't be built (incomplete monthly data) -- skipping tail transition")
                else:
                    print(f"    ⚠️ tail year {max_year} is partial but no month with data found -- skipping")
            elif not args.skip_partial:
                print(f"    ⚠️ tail year {max_year} is partial ({max_year_days:.0f}d) but no Month row / "
                      f"prior full year found -- skipping tail transition (pass --skip-partial to silence)")

    prs.save(args.out)
    print(f"\nSaved {slides_added} chart(s) to {args.out!r}.")
    if args.validate:
        print("✅ Factor decomposition matches AB-CD's own real values." if all_ok
              else "❌ MISMATCH -- do not trust this decomposition yet.")
        return 0 if all_ok else 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
