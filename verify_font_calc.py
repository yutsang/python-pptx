"""Verify the font-based text-fit calculation against the template — and against
PowerPoint ground truth.

What it does (no AI, no databook needed):
  For every commentary shape (textMainBullets*) in the template, prints the
  predicted line CAPACITY and the predicted WRAPPED LINE COUNT of a sample text
  under three calculation modes side by side:
      1. legacy CPL        — chars-per-line constant (the old heuristic)
      2. system-font       — real glyph widths from THIS machine's font
      3. client-metrics    — real glyph widths from metrics_*.json (if configured)
  If 2/3 differ from 1, the font-aware path is ACTIVE and changing the numbers.

Ground-truth protocol (the real acceptance test):
  1. Run this script, note the predicted wrapped-line count for a shape.
  2. Open the template in PowerPoint, paste the SAME sample text into that shape
     at 9pt (EN) / 9pt (CN), count the actual rendered lines.
  3. PASS if |predicted - actual| <= 1 line for client-metrics/system-font.
     The legacy CPL is typically off by more — that gap is the fix.

Usage:
    python verify_font_calc.py [template.pptx] [--eng metrics_eng.json] [--chi metrics_chi.json]
    (defaults: fdd_utils/template.pptx; metrics paths read from fdd_utils/config.yml)
"""
from __future__ import annotations

import argparse
import math
import os
import sys

SAMPLE_ENG = (
    "Cash and bank balances increased by RMB 12.3 million from RMB 45.6 "
    "million as at 31 Dec 2024 to RMB 57.9 million as at 30 Jun 2025, "
    "primarily driven by a one-off receipt of RMB 8.2 million from the "
    "disposal of a logistics property and improved collection from key "
    "customers in Q2 2025. Trade receivables turnover days improved from "
    "62 days to 48 days over the same period."
)
SAMPLE_CHI = (
    "现金及银行存款由2024年12月31日的人民币4,560万元增加人民币1,230万元至"
    "2025年6月30日的人民币5,790万元,主要由于出售一处物流物业的一次性收款"
    "人民币820万元以及2025年第二季度对主要客户应收款项回收情况改善所致。"
    "应收账款周转天数同期由62天改善至48天。"
)

# Legacy CPL constants (mirror config.yml pptx.commentary_packing.chars_per_line)
LEGACY_CPL = {
    "single": {"eng": 100, "chi": 50},
    "L": {"eng": 56, "chi": 30},
    "R": {"eng": 56, "chi": 30},
}


def slot_of(shape_name: str) -> str:
    n = (shape_name or "").lower()
    if n.endswith("_l"):
        return "L"
    if n.endswith("_r"):
        return "R"
    return "single"


def load_config_paths() -> dict:
    try:
        import yaml
        cfg_path = os.path.join("fdd_utils", "config.yml")
        with open(cfg_path, encoding="utf-8") as fh:
            cfg = yaml.safe_load(fh) or {}
        packing = ((cfg.get("pptx") or {}).get("commentary_packing") or {})
        return {
            "eng": packing.get("font_metrics_path_eng") or "",
            "chi": packing.get("font_metrics_path_chi") or "",
        }
    except Exception:
        return {"eng": "", "chi": ""}


def describe(measurer, box, text: str) -> str:
    line_h = measurer.line_height_pt()
    cap = int(box.height_pt // line_h) if line_h > 0 else 0
    lines = measurer.wrap(text, box.width_pt)
    return f"line_h={line_h:5.2f}pt  capacity={cap:>3}  sample_wraps_to={len(lines):>2} lines"


def main() -> int:
    ap = argparse.ArgumentParser()
    ap.add_argument("template", nargs="?", default=os.path.join("fdd_utils", "template.pptx"))
    ap.add_argument("--eng", default=None, help="metrics json for EN (default: from config.yml)")
    ap.add_argument("--chi", default=None, help="metrics json for CN (default: from config.yml)")
    args = ap.parse_args()

    from pptx import Presentation
    from fdd_utils.text_metrics import get_measurer, text_box_from_shape

    cfg = load_config_paths()
    eng_path = args.eng if args.eng is not None else cfg["eng"]
    chi_path = args.chi if args.chi is not None else cfg["chi"]

    print(f"template     : {args.template}")
    print(f"metrics (EN) : {eng_path or '(none — will use system font only)'}"
          f"{'' if not eng_path or os.path.exists(eng_path) else '   ⚠️ FILE NOT FOUND'}")
    print(f"metrics (CN) : {chi_path or '(none — will use system font only)'}"
          f"{'' if not chi_path or os.path.exists(chi_path) else '   ⚠️ FILE NOT FOUND'}")

    prs = Presentation(args.template)
    seen = set()
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name = (getattr(shape, "name", "") or "")
            if "textmainbullets" not in name.lower() or name in seen:
                continue
            seen.add(name)
            slot = slot_of(name)
            box = text_box_from_shape(shape)
            print(f"\n== slide {slide_idx + 1} | {name} (slot={slot}) | "
                  f"box {box.width_pt:.0f} x {box.height_pt:.0f} pt ==")

            for lang, is_cjk, size, spacing, sample, mpath in (
                ("ENG", False, 9, 1.0, SAMPLE_ENG, eng_path),
                ("CHI", True, 9, 0.95, SAMPLE_CHI, chi_path),
            ):
                cpl = LEGACY_CPL[slot]["chi" if is_cjk else "eng"]
                legacy_lines = math.ceil(len(sample) / cpl)
                print(f"  [{lang}] legacy CPL={cpl:<3}                  "
                      f"sample_wraps_to={legacy_lines:>2} lines  (old heuristic)")
                sysm = get_measurer("Microsoft YaHei" if is_cjk else "Arial", size,
                                    is_cjk=is_cjk, line_spacing=spacing)
                print(f"  [{lang}] system-font   : {describe(sysm, box, sample)}")
                if mpath and os.path.exists(mpath):
                    climt = get_measurer("ignored", size, is_cjk=is_cjk,
                                         line_spacing=spacing, metrics_path=mpath)
                    tag = "client-metrics" if climt.source == "client-metrics" else "FALLBACK(sys)"
                    print(f"  [{lang}] {tag:<14}: {describe(climt, box, sample)}")

    print("\n—— Ground-truth check (do this once per template) ——")
    print("1. Open the template in PowerPoint; paste the SAME sample text into a")
    print("   textMainBullets box at 9pt.")
    print("2. Count the actual rendered lines; compare with 'sample_wraps_to' above.")
    print("3. PASS: client-metrics/system-font within ±1 line. The legacy CPL row")
    print("   shows how far the old heuristic was off — that gap is the improvement.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
