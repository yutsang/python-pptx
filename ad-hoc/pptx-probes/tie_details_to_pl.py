#!/usr/bin/env python3
"""Does each detail table's 合计 equal the P&L line it sits beside?

A 营业成本 breakdown printed next to an income statement that says something
else is the kind of thing a reader checks first and the deck cannot survive.
Hand-checking found 28 such breaks in one portfolio and none in the other two,
so it is worth having the run say it rather than someone adding up columns.

    python ad-hoc/pptx-probes/tie_details_to_pl.py DECK.pptx
    python ad-hoc/pptx-probes/tie_details_to_pl.py PREVIEW_DIR

Each detail table is matched to the nearest 利润表 at or before its own slide,
by table title against P&L row label. Periods align on the four-digit year, so
a detail column headed 2026年1-6月 lines up with a P&L column headed
2026年6月30日. Signs are ignored: the statement prints costs negative and the
breakdown prints them positive.

Tolerance is ±2, because everything is rounded to 千元 before it is printed and
a total of seven rounded rows can legitimately sit a couple of units off its
own rounded total.

Reads only. Writes nothing.
"""
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from pptx import Presentation

TOLERANCE = 2.0
TOTAL_LABELS = ("合计", "总计", "Total")
PL_TITLE = "利润表"
_YEAR = re.compile(r"(\d{4})")


def _num(text):
    t = (text or "").strip().replace(",", "").replace("，", "")
    if not t or t in ("-", "—", "–"):
        return None
    neg = t.startswith("(") and t.endswith(")")
    t = t.strip("()")
    try:
        v = float(t)
    except ValueError:
        return None
    return -v if neg else v


def _years(table):
    """Period year per column index, from the header row. None where absent."""
    return [(_YEAR.search(table.cell(1, c).text or "") or [None, None])[1]
            if _YEAR.search(table.cell(1, c).text or "") else None
            for c in range(len(table.columns))]


def _rows(table):
    return {(table.cell(r, 0).text or "").strip(): r for r in range(2, len(table.rows))}


def main() -> int:
    if len(sys.argv) != 2:
        print(__doc__)
        return 2
    target = Path(sys.argv[1])
    decks = (sorted(p for p in target.glob("*.pptx") if not p.name.startswith("~$"))
             if target.is_dir() else [target])

    grand_checked = grand_bad = 0
    for deck in decks:
        prs = Presentation(str(deck))
        pl = None            # (slide_no, title, table, years)
        checked = bad = 0
        lines = []
        for s_i, slide in enumerate(prs.slides, 1):
            tables = [sh for sh in slide.shapes if getattr(sh, "has_table", False)]
            for sh in tables:
                t = sh.table
                title = (t.cell(0, 0).text or "").strip()
                if PL_TITLE in title:
                    pl = (s_i, title, t, _years(t))
            for sh in tables:
                t = sh.table
                title = (t.cell(0, 0).text or "").strip()
                if PL_TITLE in title or pl is None:
                    continue
                rows = _rows(t)
                total_r = next((rows[k] for k in rows if k in TOTAL_LABELS), None)
                if total_r is None:
                    continue
                # the P&L row this table breaks down, matched on the table's title
                pl_slide, _pl_title, pl_t, pl_years = pl
                pl_rows = _rows(pl_t)
                pl_r = pl_rows.get(title)
                if pl_r is None:
                    # a title like "营业成本 - 示意性调整后" still names its line
                    pl_r = next((r for lbl, r in pl_rows.items()
                                 if lbl and (lbl in title or title.startswith(lbl))), None)
                if pl_r is None:
                    continue
                det_years = _years(t)
                for c, yr in enumerate(det_years):
                    if not yr or c == 0:
                        continue
                    pc = next((i for i, y in enumerate(pl_years) if y == yr and i), None)
                    if pc is None:
                        continue
                    d = _num(t.cell(total_r, c).text)
                    p = _num(pl_t.cell(pl_r, pc).text)
                    if d is None or p is None:
                        continue
                    checked += 1
                    if abs(abs(d) - abs(p)) > TOLERANCE:
                        bad += 1
                        lines.append(f"    slide {s_i:>3}  {title[:14]:<14} {yr}: "
                                     f"P&L (slide {pl_slide}) {abs(p):>11,.0f}   "
                                     f"合计 {abs(d):>11,.0f}   diff {abs(abs(d) - abs(p)):>10,.0f}")
        grand_checked += checked
        grand_bad += bad
        mark = "✅" if bad == 0 else "⚠️"
        print(f"\n{mark} {deck.name}: {bad} of {checked} comparison(s) do NOT tie")
        for ln in lines:
            print(ln)

    if len(decks) > 1:
        print(f"\n{'=' * 70}\nTOTAL: {grand_bad} of {grand_checked} comparisons do not tie")
    return 1 if grand_bad else 0


if __name__ == "__main__":
    raise SystemExit(main())
