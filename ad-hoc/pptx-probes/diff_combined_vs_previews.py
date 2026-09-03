#!/usr/bin/env python3
"""Does the combined deck actually hold less than the previews it was built from?

Reported symptom: every per-entity .preview.pptx has commentary text, and the
Portfolio_*.pptx merged from them shows tables only. combine_presentations was
tested locally against a real export and against a text-box-plus-native-table
fixture, and text survived both, so this reads the two sides of the user's own
run instead of guessing which one is wrong.

    python ad-hoc/pptx-probes/diff_combined_vs_previews.py PREVIEW_DIR
    python ad-hoc/pptx-probes/diff_combined_vs_previews.py COMBINED.pptx PREVIEW_DIR
    python ad-hoc/pptx-probes/diff_combined_vs_previews.py COMBINED.pptx a.pptx b.pptx

Given one directory it does the whole folder: every Portfolio_<label>_<stamp>
deck in it is paired with the .preview decks carrying the same portfolio label,
newest combined file per label, and each pair is reported in turn.

Previews are concatenated in filename order, which is the order
inspect_databook.py merges them in. If your run merged a different order the
per-slide pairing shifts; the TOTALS line is order-independent, so read that
first.

Three outcomes, and they point at different code:
  totals match          -> the text is IN the file. It is a rendering problem
                           (autofit shrink, a shape off-slide, something drawn
                           over it), not a merge problem. Open the deck.
  combined has less     -> the merge is dropping content. The per-slide table
                           below names the first slide where it happens.
  previews have less    -> the previews are the problem and the merge is
                           faithful; the export ran before the text existed.

Reads only. Writes nothing.
"""
import re
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from pptx import Presentation


def slide_stats(slide):
    texts, chars, sizes = 0, 0, []
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False):
            body = sh.text_frame.text or ""
            if body.strip():
                texts += 1
                chars += len(body)
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size is not None:
                        sizes.append(round(run.font.size.pt, 2))
    tables = sum(1 for sh in slide.shapes if getattr(sh, "has_table", False))
    return {"shapes": len(slide.shapes), "tables": tables,
            "textshapes": texts, "chars": chars,
            "minfont": min(sizes) if sizes else None}


def deck_stats(path):
    return [slide_stats(s) for s in Presentation(str(path)).slides]


#: Portfolio_I_20260901_214703.pptx -> "I";  Portfolio_I-II-III_....pptx -> "I-II-III"
_COMBINED = re.compile(r"^(?:Portfolio_(.+?)|Combined)_\d{8}_\d{6}\.pptx$", re.I)
#: Project Mint.Portfolio I.南通通海.preview.pptx -> "I"
_PREVIEW = re.compile(r"Portfolio[ _]([^.\s_]+)", re.I)


def _pairs_in_folder(folder: Path):
    """(label, combined, [previews]) for every Portfolio deck in one folder."""
    combined_by_label, previews_by_label = {}, {}
    for f in sorted(folder.glob("*.pptx")):
        if f.name.startswith("~$"):
            continue
        m = _COMBINED.match(f.name)
        if m:
            label = (m.group(1) or "").upper()
            # newest wins: the stamp sorts lexically, and sorted() got here first
            combined_by_label[label] = f
            continue
        m = _PREVIEW.search(f.name)
        if m:
            previews_by_label.setdefault(m.group(1).upper(), []).append(f)
    out = []
    for label in sorted(combined_by_label):
        # A combined file whose label is a join of several (I-II-III) collects
        # every one of them, so a whole-run merge still finds its inputs.
        parts = [p for p in re.split(r"[-,]", label) if p] or [label]
        previews = [p for part in parts for p in previews_by_label.get(part, [])]
        out.append((label, combined_by_label[label], sorted(set(previews))))
    return out


def main() -> int:
    if len(sys.argv) == 2 and Path(sys.argv[1]).is_dir():
        folder = Path(sys.argv[1])
        pairs = _pairs_in_folder(folder)
        if not pairs:
            print(f"No Portfolio_<label>_<stamp>.pptx found in {folder}. Name the combined "
                  f"deck explicitly instead:\n    ... COMBINED.pptx {folder}")
            return 1
        rc = 0
        for i, (label, combined, previews) in enumerate(pairs):
            print(f"\n{'=' * 78}\nPortfolio {label or '(unlabelled)'}\n{'=' * 78}")
            if not previews:
                print(f"  {combined.name}: no .preview deck carries this label -- nothing to "
                      f"compare against.")
                rc = 1
                continue
            rc = _report(combined, previews) or rc
        return rc

    if len(sys.argv) < 3:
        print(__doc__)
        return 2
    combined = Path(sys.argv[1])
    rest = [Path(a) for a in sys.argv[2:]]
    previews = []
    for r in rest:
        if r.is_dir():
            previews += sorted(p for p in r.glob("*.pptx")
                               if not p.name.startswith("~$") and p.resolve() != combined.resolve())
        else:
            previews.append(r)
    if not previews:
        print("No preview .pptx found.")
        return 1
    return _report(combined, previews)


def _report(combined: Path, previews) -> int:
    c_stats = deck_stats(combined)
    p_stats, p_origin = [], []
    for p in previews:
        s = deck_stats(p)
        p_stats += s
        p_origin += [p.name] * len(s)

    print(f"COMBINED  {combined.name}: {len(c_stats)} slide(s)")
    print(f"PREVIEWS  {len(previews)} file(s), {len(p_stats)} slide(s) total")
    for p in previews:
        print(f"            {p.name}")

    ct, pt = sum(s["chars"] for s in c_stats), sum(s["chars"] for s in p_stats)
    ctab = sum(s["tables"] for s in c_stats)
    ptab = sum(s["tables"] for s in p_stats)
    print(f"\nTOTALS    combined {ct:>7} chars, {ctab:>3} tables")
    print(f"          previews {pt:>7} chars, {ptab:>3} tables")
    if len(c_stats) != len(p_stats):
        print(f"\n  !! slide COUNT differs ({len(c_stats)} vs {len(p_stats)}) -- the merge did not "
              f"take every slide, or these previews are not the ones it merged.")
    if ct == pt:
        print("\n  Character totals MATCH. Every character is in the combined file, so nothing")
        print("  was dropped in the merge. If the page looks empty in PowerPoint the text is")
        print("  being rendered invisibly or covered -- open it and check that page.")
    elif ct < pt:
        print(f"\n  !! combined holds {pt - ct} FEWER characters. The merge is losing text.")
    else:
        print(f"\n  !! combined holds {ct - pt} MORE characters than the previews.")

    print("\nper slide (combined | preview), rows that differ marked !!")
    print(f"{'#':>4}  {'chars':>16}  {'tables':>9}  {'textshapes':>11}  {'min font pt':>13}  source")
    for i in range(max(len(c_stats), len(p_stats))):
        c = c_stats[i] if i < len(c_stats) else None
        p = p_stats[i] if i < len(p_stats) else None
        src = p_origin[i] if i < len(p_origin) else "-"
        f = lambda a, b, k: f"{(a[k] if a else '-')!s:>7}|{(b[k] if b else '-')!s:<7}"
        differs = (c or {}).get("chars") != (p or {}).get("chars")
        print(f"{i + 1:>4}  {f(c, p, 'chars'):>16}  {f(c, p, 'tables'):>9}  "
              f"{f(c, p, 'textshapes'):>11}  {f(c, p, 'minfont'):>13}  {src}"
              + ("   !!" if differs else ""))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
