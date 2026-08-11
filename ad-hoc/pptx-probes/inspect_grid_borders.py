#!/usr/bin/env python3
"""Which rows of the BS/IS overview grids actually carry a border, and in what
colour and weight.

Answers "the balance sheet has its blue rules but the income statement does
not" without guessing. Both grids are drawn by the same function
(_fill_table_placeholder), which puts a thin top rule on any row matching a
total keyword and adds a heavy bottom rule on the statement-level grand
totals -- so a difference between the two statements is a difference in their
ROW LABELS, not in the drawing code.

Prints, per grid row: the label, and whether lnT/lnB are present, with the
colour and weight actually written. "(no border)" means the row matched no
total keyword. "=none" means an explicit <a:noFill/> (the header band's own
seams, which are cleared on purpose).

Read-only. No AI, no export -- runs against an already-exported .pptx.

Usage:
    python ad-hoc/pptx-probes/inspect_grid_borders.py
    python ad-hoc/pptx-probes/inspect_grid_borders.py <file.pptx>
"""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation

A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
DEFAULT = "for_test/pptx_previews/Crescent-databook.preview.pptx"
# The overview grids span the full commentary width (~4.78in); the per-account
# subtables are ~2.94in. Width is what tells them apart on a mixed slide.
GRID_MIN_WIDTH_IN = 4.0


def _border(tcPr, tag: str) -> str | None:
    if tcPr is None:
        return None
    el = tcPr.find(f"{A}{tag}")
    if el is None:
        return None
    if el.find(f"{A}noFill") is not None:
        return "none"
    clr = el.find(f".//{A}srgbClr")
    if clr is None:
        return "inherit"
    try:
        width_pt = int(el.get("w", 0)) / 12700
    except (TypeError, ValueError):
        width_pt = 0.0
    return f"#{clr.get('val')} @{width_pt:.2f}pt"


def main() -> int:
    path = sys.argv[1] if len(sys.argv) > 1 else DEFAULT
    if not Path(path).exists():
        print(f"Not found: {path}")
        return 1

    prs = Presentation(path)
    found = 0
    for slide_no, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not getattr(shape, "has_table", False):
                continue
            width_in = shape.width / 914400
            if width_in < GRID_MIN_WIDTH_IN:
                continue
            found += 1
            table = shape.table
            print(f"\n=== slide {slide_no}  {shape.name}  "
                  f"{len(table.rows)}x{len(table.columns)}  width={width_in:.2f}in ===")
            ruled = 0
            for row_no in range(len(table.rows)):
                cell = table.cell(row_no, 0)
                tcPr = cell._tc.find(f"{A}tcPr")
                marks = []
                for tag, label in (("lnT", "TOP"), ("lnB", "BOT")):
                    value = _border(tcPr, tag)
                    if value:
                        marks.append(f"{label}={value}")
                if any(m.split("=", 1)[1] not in ("none", "inherit") for m in marks):
                    ruled += 1
                text = (cell.text or "").strip().replace("\n", " ")
                print(f"  r{row_no:<3} {text[:26]:<28} {'  '.join(marks) or '(no border)'}")
            print(f"  -> {ruled} row(s) carry a real rule")

    if not found:
        print(f"No overview grid (>= {GRID_MIN_WIDTH_IN}in wide) found in {path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
