"""What would a different packing objective have done? Answered off an
already-exported deck, with no AI run and nothing written.

The measurement layer is settled -- 1,636 of 1,638 paragraphs match real
PowerPoint -- so the remaining problem is not how tall the text is, it is which
slot it lands in. Across a full Portfolio I+II+III run: mean fill 70.7%, every
deck carrying an entirely empty commentary column, and three of four carrying
slots that overflow. Overflowing and empty AT THE SAME TIME is the pathology.

That is the DP objective behaving exactly as written, not an arithmetic error
(gen_packing.py, `slot_cost` / the lexicographic state below it):

    state = (number of non-empty slots, underfill penalty)   -- compared as a tuple

so fewer slots ALWAYS wins first, and the penalty exempts the last non-empty
slot entirely. Cram the earlier slots past 95% and whatever is left can go in
the final slot at any fill at all, for free.

This script does not change that. It re-runs the same partition problem over
the same measured block heights under several candidate objectives and prints
what each would have produced, so the trade-off is a table rather than an
argument. `docs/failed-attempts.md` records two reverted attempts in this area;
the point of a dry run is to not become the third.

    python ad-hoc/pptx-probes/dryrun_packing_objective.py <deck.pptx | folder>

READ THE `A vs shipped` LINE FIRST. The script re-solves the reconstructed
input under objective A, the one production uses, and reports how far that
lands from what actually shipped. Expect a gap: the shipped layout is the DP
*plus* every rebalance pass that runs after it (_maximize_forward_fill,
_rebalance_overflowing_boundaries, ...), and this models the DP alone. A small
gap means the reconstruction is sound and the objective columns are a fair A/B
on the same input. A large one means the post-DP passes are doing most of the
work, and changing the objective would matter less than the table suggests --
which is a finding about where to intervene, not a reason to ignore the run.

Known limits, all of which weaken the reconstruction rather than the comparison
between objectives on the same input:
  * blocks are read POST-split, so an account production cut in half arrives
    here as two blocks that can no longer be recombined;
  * a table's reserved band is measured as the blank paragraphs it renders as,
    which is its height but not its indivisibility;
  * statement boundaries are inferred from coSummaryShape, which is where
    _apply_content_to_slides starts a statement.
"""

from __future__ import annotations

import os
import sys
from dataclasses import dataclass, field
from typing import Dict, List, Optional, Sequence, Tuple

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__)))))

from pptx import Presentation

from fdd_utils.text_metrics import POWERPOINT_LINE_PITCH_FACTOR, get_measurer, text_box_from_shape
from fdd_utils.pptx.helpers import (
    _measurer_family, _real_font_size_pt, _real_line_spacing, _real_para_gap_pt,
    _resolve_font_metrics_path,
)
from fdd_utils.pptx.payloads import _load_pptx_settings
from inspect_pptx import _is_chinese_text

BULLET = "■"
HANG_PT = 10.8
TARGET = 0.95          # gen_packing.py's target_fill_min_ratio default
# gen_packing.py runs the DP at STRICT capacity first and only relaxes when
# that has no feasible answer -- the 1.05 floor and the tiers above it. Giving
# the solver the tail tolerance up front instead let it cram slot 0 to 106% on
# the first pass, which production would never do, and the sanity check
# correctly rejected the result.
RELAX_LADDER = (1.0, 1.05, 1.35, 1.6, 10.0)
#: A slot holding only the template's leading empty paragraph reads ~3%, not 0.
EMPTY_FILL = 0.05


@dataclass
class Block:
    """One account as the packer sees it: a contiguous run of paragraphs that
    must stay together, plus the height it really occupies."""
    label: str
    height_pt: float
    category: str
    slot_index: int        # where it actually shipped, for the sanity check


@dataclass
class Statement:
    name: str
    blocks: List[Block] = field(default_factory=list)
    caps: List[float] = field(default_factory=list)        # usable height per slot
    shipped: List[float] = field(default_factory=list)     # pt actually in each slot
    # Height in a slot that no Block ends up owning. This started as an
    # accounting curiosity and turned out to be the reason every "A vs shipped:
    # 7pp" row existed: `used` charged the category-header paragraph, the block
    # model did not, so the solver was handed a problem ~4.5pp per slot lighter
    # than reality and every "this fits" it printed was optimistic by a line.
    # Keep this reported. A silent leftover is a biased solver.
    unmodelled: List[float] = field(default_factory=list)
    unmodelled_text: List[str] = field(default_factory=list)


# ---------------------------------------------------------------------------
# Reconstruct the packing input from a rendered deck
# ---------------------------------------------------------------------------

def _read_statements(deck: str) -> List[Statement]:
    prs = Presentation(deck)
    packing = (_load_pptx_settings() or {})
    packing = packing.get("packing", packing)
    measurers = {
        is_chi: get_measurer(
            _measurer_family(is_chi, packing), _real_font_size_pt(is_chi), is_cjk=is_chi,
            line_spacing=_real_line_spacing(is_chi),
            metrics_path=_resolve_font_metrics_path(is_chi, packing),
        )
        for is_chi in (False, True)
    }

    statements: List[Statement] = []
    current: Optional[Statement] = None

    for s_idx, slide in enumerate(prs.slides, start=1):
        names = {str(getattr(sh, "name", "") or "") for sh in slide.shapes}
        # A statement starts on the slide carrying the executive-summary band --
        # that is the statement-first page _apply_content_to_slides writes.
        if "coSummaryShape" in names or current is None:
            current = Statement(name=f"{os.path.basename(deck)} #{len(statements) + 1}")
            statements.append(current)

        # L before R, matching the reading order the packer fills in.
        slots = sorted(
            (sh for sh in slide.shapes
             if str(getattr(sh, "name", "") or "").startswith("textMainBullets")),
            key=lambda sh: (int(sh.left or 0)),
        )
        for shape in slots:
            box = text_box_from_shape(shape)
            slot_i = len(current.caps)
            current.caps.append(box.height_pt)
            tf = shape.text_frame
            text = tf.text or ""
            is_chi = _is_chinese_text(text)
            measurer = measurers[is_chi]
            hang_w = max(10.0, box.width_pt - HANG_PT)

            used = 0.0
            open_block: Optional[Block] = None
            pending_category = ""
            # A category header belongs to the accounts under it -- the packer
            # cannot leave the header at the foot of one column and the first
            # account at the head of the next -- so its height rides forward
            # onto the next block. Anything still pending when the slot ends is
            # height nothing owns, and gets reported rather than dropped.
            pending_pt = 0.0
            pending_texts: List[str] = []
            for para in tf.paragraphs:
                p_text = para.text or ""
                sa = para.space_after.pt if para.space_after is not None else 0.0
                sb = para.space_before.pt if para.space_before is not None else 0.0
                sizes = [r.font.size.pt for r in para.runs if r.font.size is not None]
                font_pt = max(sizes) if sizes else _real_font_size_pt(is_chi)
                if not p_text.strip():
                    n = 1
                else:
                    first_w = box.width_pt if p_text.lstrip().startswith((BULLET, "➢")) else None
                    n = max(1, len(measurer.wrap(p_text, hang_w, first_line_width_pt=first_w)))
                height = n * font_pt * POWERPOINT_LINE_PITCH_FACTOR + sa + sb
                used += height

                stripped = p_text.strip()
                if stripped.startswith(BULLET):
                    open_block = Block(label=stripped[:28],
                                       height_pt=height + pending_pt,
                                       category=pending_category, slot_index=slot_i)
                    current.blocks.append(open_block)
                    pending_category = ""
                    pending_pt = 0.0
                    pending_texts.clear()
                elif stripped and sa == 0.0 and not stripped.startswith("➢"):
                    pending_category = stripped[:16]     # category header
                    pending_pt += height
                    pending_texts.append(stripped[:16])
                elif open_block is not None:
                    open_block.height_pt += height       # continuation / explain / spacer
                else:
                    pending_pt += height                 # spacer ahead of any block
                    pending_texts.append(repr(p_text)[:14])
            current.shipped.append(used)
            current.unmodelled.append(pending_pt)
            current.unmodelled_text.append(", ".join(pending_texts))
    return [st for st in statements if st.blocks]


# ---------------------------------------------------------------------------
# The partition problem, once, under a pluggable objective
# ---------------------------------------------------------------------------

def _penalty_last_slot_exempt(fills: Sequence[float]) -> float:
    """Objective A -- what production does today. Every non-empty slot except
    the LAST non-empty one is charged for falling short of TARGET."""
    nz = [i for i, f in enumerate(fills) if f > EMPTY_FILL]
    if not nz:
        return 0.0
    return sum(max(0.0, TARGET - fills[i]) for i in nz[:-1])


def _penalty_no_exemption(fills: Sequence[float]) -> float:
    """Objective B -- the last slot is charged like every other."""
    return sum(max(0.0, TARGET - f) for f in fills if f > EMPTY_FILL)


def _penalty_last_slot_floor(fills: Sequence[float], floor: float = 0.50) -> float:
    """Objective C -- the last slot keeps its exemption but only down to a
    floor. A final column at 20% is charged; one at 60% is not."""
    nz = [i for i, f in enumerate(fills) if f > EMPTY_FILL]
    if not nz:
        return 0.0
    p = sum(max(0.0, TARGET - fills[i]) for i in nz[:-1])
    return p + max(0.0, floor - fills[nz[-1]])


def _penalty_worst_slot(fills: Sequence[float]) -> float:
    """Objective D -- minimise the emptiest slot. This is what the code
    comment says was REPLACED for giving 45%/72%; included so the dry run can
    show whether that judgement still holds rather than restating it."""
    nz = [f for f in fills if f > EMPTY_FILL]
    return (1.0 - min(nz)) if nz else 0.0


#: (label, penalty, slot_count_first)
#: slot_count_first mirrors gen_packing.py's lexicographic tuple, where the
#: number of non-empty slots is compared BEFORE the penalty.
#:
#: An earlier note here said A-D all return the identical layout on real decks.
#: That was measured on decks that predated the current packer and it is wrong.
#: On the 2026-08-27 exports (15 entities, 30 statements) A and D differ on 11
#: of the 30, and where they differ D is far flatter: mean spread between the
#: fullest and emptiest slot is 18pp for D against 59pp for A and 62pp as
#: shipped -- at the same slot count. The slot-count term does dominate, but it
#: leaves more room underneath it than that note claimed. B and E remain
#: identical to each other by construction (same penalty, and demoting the slot
#: count changes nothing once the slot count is already forced).
OBJECTIVES = (
    ("A last-slot exempt (today)", _penalty_last_slot_exempt, True),
    ("B no exemption", _penalty_no_exemption, True),
    ("C last-slot floor 50%", _penalty_last_slot_floor, True),
    ("D minimise worst slot", _penalty_worst_slot, True),
    ("E balance first, slots 2nd", _penalty_no_exemption, False),
)


def _solve(st: Statement, penalty_fn,
           slot_count_first: bool = True) -> Optional[Tuple[List[int], float]]:
    """Contiguous partition of blocks into slots, minimising
    (non-empty slots, penalty) exactly as gen_packing.py compares them.

    Brute-forced over cut positions rather than reusing the production DP:
    the point is to vary the objective, and a small independent solver that
    can be read in one screen is worth more here than sharing code with the
    thing under test.
    """
    n, S = len(st.blocks), len(st.caps)
    if n == 0 or S == 0:
        return None
    best: Optional[Tuple[Tuple[float, float], List[int]]] = None
    cap_limit = [c for c in st.caps]

    def walk(slot: int, start: int, assign: List[int]) -> None:
        nonlocal best
        if slot == S:
            if start != n:
                return
            fills = []
            for s in range(S):
                h = sum(st.blocks[i].height_pt for i in range(n) if assign[i] == s)
                fills.append(h / st.caps[s] if st.caps[s] else 0.0)
            n_slots = float(sum(1 for f in fills if f > EMPTY_FILL))
            pen = round(penalty_fn(fills), 6)
            state = (n_slots, pen) if slot_count_first else (pen, n_slots)
            if best is None or state < best[0]:
                best = (state, list(assign))
            return
        remaining_slots = S - slot
        for end in range(start, n + 1):
            if n - end > 0 and remaining_slots == 1:
                continue
            # Slots fill in reading order, so a slot may only be empty once
            # every slot after it is empty too. Without this the solver
            # discovers that skipping slot 0 costs nothing -- fewer non-empty
            # slots wins lexicographically -- and returns layouts like
            # "0% 0% 97% 91% 55%" that production would never produce. The DP
            # guards the same case: see the j == -1 bypass in gen_packing.py,
            # allowed only when the previous slot's whole row is infeasible.
            if end == start and start < n:
                continue
            h = sum(st.blocks[i].height_pt for i in range(start, end))
            if h > cap_limit[slot]:
                break
            for i in range(start, end):
                assign[i] = slot
            walk(slot + 1, end, assign)
        return

    # Same progressive relax the DP uses: strict first, widen only on failure.
    for mult in RELAX_LADDER:
        cap_limit[:] = [c * mult for c in st.caps]
        best = None
        walk(0, 0, [0] * n)
        if best is not None:
            return best[1], mult
    return None


def _fills_from(st: Statement, assign: Sequence[int]) -> List[float]:
    out = []
    for s in range(len(st.caps)):
        h = sum(st.blocks[i].height_pt for i in range(len(st.blocks)) if assign[i] == s)
        out.append(h / st.caps[s] if st.caps[s] else 0.0)
    return out


def _fmt(fills: Sequence[float]) -> str:
    return " ".join(f"{f * 100:5.0f}%" for f in fills)


def _score(fills: Sequence[float]) -> Tuple[int, int, float]:
    nz = [f for f in fills if f > EMPTY_FILL]
    return (sum(1 for f in fills if f <= EMPTY_FILL),
            sum(1 for f in fills if f > 1.0),
            (sum(nz) / len(nz)) if nz else 0.0)


def run(deck: str) -> None:
    try:
        statements = _read_statements(deck)
    except Exception as exc:
        print(f"  {os.path.basename(deck)}: could not read -- {type(exc).__name__}: {exc}")
        return
    if not statements:
        print(f"  {os.path.basename(deck)}: no commentary blocks found")
        return

    print(f"\n{'=' * 92}\n{os.path.basename(deck)}\n{'=' * 92}")
    for st in statements:
        shipped_fill = [h / c if c else 0.0 for h, c in zip(st.shipped, st.caps)]
        print(f"\n  {len(st.blocks)} account block(s) over {len(st.caps)} slot(s)")
        print(f"    {'as shipped':<28}{_fmt(shipped_fill)}")

        leftover = sum(st.unmodelled)
        if leftover > 1.0:
            where = ", ".join(f"slot {i + 1}: {t}" for i, t in
                              enumerate(st.unmodelled_text) if st.unmodelled[i] > 1.0)
            print(f"    leftover {leftover:.1f}pt no block owns -- {where}")

        solved: Dict[str, List[float]] = {}
        degenerate: Dict[str, float] = {}
        for name, fn, first in OBJECTIVES:
            got = _solve(st, fn, first)
            if got is None:
                print(f"    {name:<28}(no feasible partition)")
                continue
            assign, mult = got
            solved[name] = _fills_from(st, assign)
            degenerate[name] = mult

        a_name = OBJECTIVES[0][0]
        a_fills = solved.get(a_name)
        gap = (max(abs(x - y) for x, y in zip(a_fills, shipped_fill))
               if a_fills else None)
        for name, _fn, _first in OBJECTIVES:
            if name not in solved:
                continue
            empty, over, mean = _score(solved[name])
            mark = f"   <- A vs shipped: worst slot differs by {gap * 100:.0f}pp" \
                if (name == a_name and gap is not None) else ""
            if degenerate.get(name, 1.0) >= 10.0:
                mark = "   <- DEGENERATE, see below"
            print(f"    {name:<28}{_fmt(solved[name])}"
                  f"   empty={empty} over={over} mean={mean * 100:.0f}%{mark}")

        if any(m >= 10.0 for m in degenerate.values()):
            # Not a dry-run artefact: gen_packing.py:2118 really does end its
            # ladder at x10, and at that rung "fewest non-empty slots" is won
            # outright by putting everything in slot 0. So the DP contributes
            # nothing here and the shipped layout is the rebalance passes'
            # work end to end. Do not read the row above as a proposal.
            print("    !! Solved only at relax x10, where the objective degenerates to")
            print("       'one slot holds everything'. Production reaches the same rung")
            print("       (gen_packing.py:2118) -- so for this statement the DP is not")
            print("       choosing the layout at all. The content simply does not fit.")
        elif gap is not None and gap > 0.15:
            print(f"    !! A lands {gap * 100:.0f}pp from what shipped on its worst slot, so most")
            print("       of this layout came from the rebalance passes AFTER the DP, not from")
            print("       the objective. Changing the objective would move it less than the")
            print("       rows above imply -- read this as pointing at those passes instead.")


def main() -> int:
    if len(sys.argv) < 2:
        print(__doc__)
        return 2
    target = sys.argv[1]
    decks = ([os.path.join(target, f) for f in sorted(os.listdir(target))
              if f.lower().endswith(".pptx") and not f.startswith("~$")]
             if os.path.isdir(target) else [target])
    print("Dry run only. Nothing is written and production is unchanged.")
    for deck in decks:
        run(deck)
    print("\nRead the SANITY marker on each statement before believing any row.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
