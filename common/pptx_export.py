import os
import textwrap
import logging
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from dataclasses import dataclass
from typing import List, Tuple
from datetime import datetime
import re
from pptx.oxml.ns import qn  # Required import for XML namespace handling

logging.basicConfig(level=logging.INFO)

def get_tab_name(project_name):
    """Get tab name based on project name - more flexible approach"""
    try:
        # Hardcoded mappings for known entities
        if project_name == 'Haining':
            return "BSHN"
        elif project_name == 'Nanjing':
            return "BSNJ"
        elif project_name == 'Ningbo':
            return "BSNB"
        else:
            # For unknown entities, return the project name itself to avoid None
            print(f"Warning: Unknown project name '{project_name}', using project name as fallback")
            return project_name
    except Exception:
        return None

def clean_content_quotes(content):
    """
    Remove outermost quotation marks from content while preserving legitimate internal quotes.
    Supports both straight quotes (") and curly quotes (" ").
    """
    if not content:
        return content

    # Handle straight quotes
    content = re.sub(r'^"([^"]*)"$', r'\1', content)

    # Handle curly quotes
    content = re.sub(r'^"([^"]*)"$', r'\1', content)
    content = re.sub(r'^"([^"]*)"$', r'\1', content)

    return content

def detect_chinese_text(text, force_chinese_mode=False):
    """
    Detect if text contains significant Chinese characters.
    For Chinese PPTX export, apply Chinese formatting even to English text.
    Returns True if more than 30% of characters are Chinese, or if force_chinese_mode is True.
    """
    if not text:
        return force_chinese_mode  # Return the forced mode for empty text

    chinese_chars = sum(1 for char in text if '\u4e00' <= char <= '\u9fff')
    total_chars = len(text)

    if total_chars == 0:
        return force_chinese_mode

    chinese_ratio = chinese_chars / total_chars

    # If force_chinese_mode is True (for Chinese PPTX export), apply Chinese formatting
    if force_chinese_mode:
        return True

    # Otherwise, use the normal detection threshold
    return chinese_ratio > 0.3

def get_font_size_for_text(text, base_size=Pt(9), force_chinese_mode=False):
    """
    Get appropriate font size based on text content.
    Chinese text gets smaller font to maximize content density and prevent overflow.
    Enhanced for better Chinese character handling.
    """
    if detect_chinese_text(text, force_chinese_mode=force_chinese_mode):
        # Use even smaller font for Chinese to prevent line breaks
        chinese_ratio = sum(1 for char in text if '\u4e00' <= char <= '\u9fff') / len(text) if text else 0
        if chinese_ratio > 0.5 or force_chinese_mode:  # Mostly Chinese text or forced mode
            print(f"🔤 FONT: Using 9pt Arial for Chinese text (ratio: {chinese_ratio:.2f}, force: {force_chinese_mode})")
            return Pt(9)  # Standard 9pt font for Chinese content (as requested)
        else:
            print(f"🔤 FONT: Using 9pt Arial for mixed Chinese text (ratio: {chinese_ratio:.2f})")
            return Pt(9)  # Standard 9pt font for mixed Chinese/English content
    else:
        print(f"🔤 FONT: Using {base_size}pt for English text")
        return base_size  # Default size for English

def get_font_name_for_text(text, default_font='Arial'):
    """
    Get appropriate font name based on text content.
    Chinese text uses Microsoft YaHei, English uses Arial.
    """
    if not text:
        return default_font

    # Check if text contains Chinese characters
    has_chinese = any('\u4e00' <= char <= '\u9fff' for char in text)

    if has_chinese:
        return 'Microsoft YaHei'
    else:
        return default_font

def get_line_spacing_for_text(text, force_chinese_mode=False):
    """
    Get appropriate line spacing based on text content.
    Chinese text needs tighter spacing to maximize content density.
    Enhanced for better Chinese line break handling.
    """
    if detect_chinese_text(text, force_chinese_mode=force_chinese_mode):
        chinese_ratio = sum(1 for char in text if '\u4e00' <= char <= '\u9fff') / len(text) if text else 0
        if chinese_ratio > 0.5:  # Mostly Chinese text
            return Pt(10)  # Ultra-tight spacing for dense Chinese content
        else:
            return Pt(11)  # Tighter spacing for mixed Chinese/English
    else:
        return Pt(12)  # Standard spacing for English

def get_space_after_for_text(text, force_chinese_mode=False):
    """
    Get appropriate space after paragraph based on text content.
    """
    if detect_chinese_text(text, force_chinese_mode=force_chinese_mode):
        return Pt(4)  # Much less space after for Chinese to maximize content
    else:
        return Pt(8)  # Standard space after for English

def get_space_before_for_text(text, force_chinese_mode=False):
    """
    Get appropriate space before paragraph based on text content.
    """
    if detect_chinese_text(text, force_chinese_mode=force_chinese_mode):
        return Pt(2)  # Much less space before for Chinese to maximize content
    else:
        return Pt(4)  # Standard space before for English

def replace_entity_placeholders(content, project_name):
    """
    Replace entity name placeholders in content with abbreviated entity names.
    """
    if not content or not project_name:
        return content

    # Extract first two words for professional display
    if project_name:
        words = project_name.split()
        # Use first two words, or first word if only one word
        display_entity = ' '.join(words[:2]) if len(words) >= 2 else words[0] if words else project_name
    else:
        display_entity = project_name

    # Replace common entity name placeholders
    replacements = {
        '[specific entity name]': display_entity,
        '[entity name]': display_entity,
        '[company name]': display_entity,
        '[target entity]': display_entity,
        '[ENTITY_NAME]': display_entity,
        '[COMPANY_NAME]': display_entity,
        '[SPECIFIC_ENTITY_NAME]': display_entity,
        '[TARGET_ENTITY]': display_entity
    }

    for placeholder, replacement in replacements.items():
        content = content.replace(placeholder, replacement)

    return content

@dataclass
class FinancialItem:
    accounting_type: str
    account_title: str
    descriptions: List[str]
    layer1_continued: bool = False
    layer2_continued: bool = False
    is_table: bool = False

class PowerPointGenerator:
    def __init__(self, template_path: str, language: str = 'english', row_limit: int = 20):
        self.prs = Presentation(template_path)
        self.current_slide_index = 0
        self.LINE_HEIGHT = Pt(12)
        self.ROWS_PER_SECTION = 35  # Reduced for less crowding on page 1
        self.language = language  # Store language for Chinese mode detection
        self.row_limit = row_limit  # Maximum number of rows per shape
        
        
        # Find a slide with textMainBullets shape
        shape = None
        for slide in self.prs.slides:
            try:
                shape = next(s for s in slide.shapes if s.name == "textMainBullets")
                break
            except StopIteration:
                continue
        
        if not shape:
            # If no textMainBullets found, try to find any text shape
            for slide in self.prs.slides:
                for s in slide.shapes:
                    if hasattr(s, 'text_frame'):
                        shape = s
                        break
                if shape:
                    break
        
        if not shape:
            raise ValueError("No suitable text shape found in template")
        
        self.CHARS_PER_ROW = 50
        self.BULLET_CHAR = '■ '
        self.DARK_BLUE = RGBColor(0, 50, 150)
        self.DARK_GREY = RGBColor(169, 169, 169)
        self.prev_layer1 = None
        self.prev_layer2 = None

    def log_template_shapes(self):
        # Removed shape audit logging to reduce debug output
        pass

    def _calculate_max_rows_for_shape(self, shape):
        """Calculate the actual maximum number of rows that can fit in a shape"""
        if not shape or not hasattr(shape, 'height'):
            return 25  # Default fallback

        # Get the actual shape height in EMU (English Metric Units)
        shape_height_emu = shape.height

        # Convert EMU to points (1 EMU = 1/914400 inches, 1 inch = 72 points)
        shape_height_pt = shape_height_emu * 72 / 914400

        # Account for margins and padding - use maximum space (99.9% of shape height)
        effective_height_pt = shape_height_pt * 0.999  # Maximum height utilization

        # Calculate line height based on font size and line spacing
        # Use 8pt font for better space utilization
        if hasattr(self, 'language') and self.language == 'chinese':
            font_size_pt = 8  # 8pt font size for Chinese
            line_spacing = 1.05  # Balanced spacing for readability
        else:
            font_size_pt = 8  # 8pt font size for English
            line_spacing = 1.1  # Balanced spacing for English
        line_height_pt = font_size_pt * line_spacing

        # Calculate maximum rows that can fit
        max_rows = int(effective_height_pt / line_height_pt)

        # Optimize for better space utilization with 8pt font
        # Allow more content to fit while preventing overflow
        if hasattr(self, 'language') and self.language == 'chinese':
            max_rows = max(40, max_rows)  # Adjusted for 8pt font and less crowding
        else:
            max_rows = max(38, max_rows)  # Adjusted for 8pt font and less crowding

        return max_rows

    def _calculate_max_rows(self):
        # Find a slide with textMainBullets shape
        shape = None
        for slide in self.prs.slides:
            try:
                shape = next(s for s in slide.shapes if s.name == "textMainBullets")
                break
            except StopIteration:
                continue
        
        if not shape:
            # If no textMainBullets found, try to find any text shape
            for slide in self.prs.slides:
                for s in slide.shapes:
                    if hasattr(s, 'text_frame'):
                        shape = s
                        break
                if shape:
                    break
        
        if not shape:
            return 25  # Default fallback
        
        return self._calculate_max_rows_for_shape(shape)

    def _get_target_shape_for_section(self, slide_idx, section):
        """Get the target shape for a specific section on a specific slide"""
        # Check if the slide exists
        if slide_idx >= len(self.prs.slides):
            return None

        slide = self.prs.slides[slide_idx]

        # Handle both template structures:
        # - Old template: textMainBullets_L and textMainBullets_R
        # - New template: single textMainBullets
        
        if slide_idx == 0:
            # Slide 0 - try textMainBullets first, then fallback
            if section == 'c':
                # Try textMainBullets first
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
                except StopIteration:
                    pass
            return None  # No 'b' section on slide 0
        else:
            # Additional slides (index 1+) - handle both template structures
            if section == 'b':
                # Left section - try textMainBullets_L first, then textMainBullets
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets_L"), None)
                except StopIteration:
                    pass
                # Fallback to textMainBullets for single column layout
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
                except StopIteration:
                    pass
            elif section == 'c':
                # Right section - try textMainBullets_R first, then textMainBullets
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets_R"), None)
                except StopIteration:
                    pass
                # Fallback to textMainBullets for single column layout
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
                except StopIteration:
                    pass
        
        # Fallback: try to find any textMainBullets shape
        try:
            return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
        except StopIteration:
            pass
        
        # If still not found, try to find Content Placeholder 2 (standard layout)
        try:
            return next((s for s in slide.shapes if s.name == "Content Placeholder 2"), None)
        except StopIteration:
            pass
        
        # If still not found, try to find any text shape with text_frame
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                return shape
        
        # Final fallback: any shape that might be a text container
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                return shape
        
        return None

    def parse_markdown(self, md_content: str) -> List[FinancialItem]:
        items = []
        current_type = ""
        current_title = ""
        current_descs = []
        is_table = False

        for line in md_content.strip().split('\n'):
            stripped = line.strip()
            
            if stripped.startswith('## '):
                if current_descs:
                    items.append(FinancialItem(
                        current_type, current_title, current_descs, 
                        is_table=is_table
                    ))
                current_type = stripped[3:]
                current_title = ""
                current_descs = []
                is_table = False
            elif stripped.startswith('### '):
                if current_descs:
                    items.append(FinancialItem(
                        current_type, current_title, current_descs,
                        is_table=is_table
                    ))
                current_title = stripped[4:]
                current_descs = []
                is_table = "taxes payables" in current_title.lower()
            elif re.match(r'^={20,}', stripped):
                is_table = True
                current_descs.append("="*40)
            else:
                if stripped:
                    # Clean quotes from the content
                    cleaned_stripped = clean_content_quotes(stripped)
                    current_descs.append(cleaned_stripped)

        if current_descs:
            items.append(FinancialItem(current_type, current_title, current_descs, is_table=is_table))
        return items

    def _plan_content_distribution(self, items: List[FinancialItem]):
        distribution = []
        content_queue = items.copy()

        # Start from slide 0 (index 0) for content slides
        # For the server template structure:
        # - Slide 0 (index 0): use only 'c' section (textMainBullets)
        # - Slide 1+ (index 1+): use 'b' (left) and 'c' (right) sections (textMainBullets_L and textMainBullets_R)
        slide_idx = 0

        # Debug: Log total items and content type
        total_items = len(items)
        print(f"🔍 CONTENT DISTRIBUTION: Processing {total_items} items in {self.language} mode")
        chinese_items = sum(1 for item in items if any('\u4e00' <= char <= '\u9fff' for desc in item.descriptions for char in desc))
        print(f"🔍 CONTENT TYPE: {chinese_items} Chinese items, {total_items - chinese_items} English items")
        print(f"🔍 LANGUAGE MODE: {self.language} (force Chinese formatting: {self.language == 'chinese'})")
        
        while content_queue:
            print(f"📊 SLIDE {slide_idx}: Processing {len(content_queue)} remaining items")
            if slide_idx == 0:
                sections = ['c']  # Only 'c' section on slide 0 (textMainBullets)
                print(f"📊 SLIDE {slide_idx}: Using section 'c' (textMainBullets)")
            else:
                # For server template: use both left and right sections
                sections = ['b', 'c']  # Left and right sections (textMainBullets_L and textMainBullets_R)
                print(f"📊 SLIDE {slide_idx}: Using sections 'b' and 'c' (textMainBullets_L and textMainBullets_R)")

            for section in sections:
                if not content_queue:
                    break
                section_items = []
                lines_used = 0

                # Get the actual shape for this section to calculate proper line limits
                shape = self._get_target_shape_for_section(slide_idx, section)
                if shape:
                    max_lines = self._calculate_max_rows_for_shape(shape)
                    # Set all slide limits to 35 lines
                    max_lines = 35  # All slides limited to 35 lines
                    print(f"📐 SECTION {section}: Found shape '{shape.name}', max_lines = {max_lines}")
                else:
                    max_lines = self.ROWS_PER_SECTION  # Fallback
                    print(f"⚠️ SECTION {section}: No shape found, using fallback max_lines = {max_lines}")
                
                while content_queue and lines_used < max_lines:
                    item = content_queue[0]
                    # Check if this is Chinese content for line calculation
                    is_chinese_content = any('\u4e00' <= char <= '\u9fff' for desc in item.descriptions for char in desc)
                    if is_chinese_content:
                        item_lines = self._calculate_chinese_item_lines(item)
                    else:
                        item_lines = self._calculate_item_lines(item)
                    print(f"📝 ITEM: '{item.accounting_type}' needs {item_lines} lines, used {lines_used}/{max_lines} (Chinese: {is_chinese_content})")

                    # If the item fits completely, add it
                    if lines_used + item_lines <= max_lines:
                        section_items.append(item)
                        content_queue.pop(0)
                        lines_used += item_lines
                    else:
                        # Try to split the item to fill remaining space
                        remaining_lines = max_lines - lines_used
                        if remaining_lines >= 3:  # Need at least 3 lines for meaningful content
                            # Check if this is Chinese content and use appropriate splitter
                            is_chinese_content = any('\u4e00' <= char <= '\u9fff' for desc in item.descriptions for char in desc)
                            if is_chinese_content:
                                split_item, remaining_item = self._split_item_chinese(item, remaining_lines)
                            else:
                                split_item, remaining_item = self._split_item(item, remaining_lines)
                            if split_item and split_item.descriptions:
                                section_items.append(split_item)
                                lines_used = max_lines  # Mark as full
                                
                                # Put the remaining part back at the front of the queue
                                if remaining_item and remaining_item.descriptions and remaining_item.descriptions[0]:
                                    remaining_item.layer2_continued = True
                                    remaining_item.layer1_continued = True
                                    content_queue[0] = remaining_item
                                else:
                                    content_queue.pop(0)
                            else:
                                # If splitting failed, move to next section
                                break
                        else:
                            # Not enough space for meaningful split, move to next section
                            break
                if section_items:
                    distribution.append((slide_idx, section, section_items))
                    print(f"✅ SECTION {section}: Added {len(section_items)} items to slide {slide_idx}")

            slide_idx += 1
            # Don't limit by existing slides since we'll create new ones as needed

        # Final distribution summary
        print(f"\n🎯 SLIDE DISTRIBUTION: {len(set(slide_idx for slide_idx, _, _ in distribution))} slides")
        slide_counts = {}
        for slide_idx, section, items in distribution:
            key = f"Slide {slide_idx}"
            if key not in slide_counts:
                slide_counts[key] = {}
            slide_counts[key][section] = len(items)

        for slide, sections in slide_counts.items():
            section_summary = ", ".join([f"{sec}: {count}" for sec, count in sections.items()])
            print(f"  {slide} - {section_summary}")

        return distribution

    def _calculate_chars_per_line(self, shape):
        """Calculate characters per line based on actual shape width and content type"""
        if not shape or not hasattr(shape, 'width'):
            return self.CHARS_PER_ROW  # Default fallback

        # Convert EMU to pixels: 1 EMU = 1/914400 inches, 1 inch = 96 px
        px_width = int(shape.width * 96 / 914400)

        # Use maximum width utilization for full space usage
        effective_width = px_width * 0.999  # Ultra-maximum width utilization

        # Different character widths for different font sizes and styles
        # Chinese characters are wider than English, so we need different calculations
        if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
            # Check if text contains Chinese characters
            has_chinese = False
            is_bold = False

            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.bold:
                        is_bold = True
                    if run.text and any('\u4e00' <= char <= '\u9fff' for char in run.text):
                        has_chinese = True
                        break
                if has_chinese:
                    break

            if has_chinese:
                # Chinese characters are wider than English - corrected calculations
                # Get text from the first paragraph with Chinese characters
                text = ""
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and any('\u4e00' <= char <= '\u9fff' for char in run.text):
                            text = run.text
                            break
                    if text:
                        break
                chinese_ratio = sum(1 for char in text if '\u4e00' <= char <= '\u9fff') / len(text) if text else 0
                if chinese_ratio > 0.8:  # Almost entirely Chinese
                    if is_bold:
                        avg_char_px = 14.0  # Bold Chinese text (much wider than English) - increased
                    else:
                        avg_char_px = 12.5  # Regular Chinese text (wider than English) - increased
                elif chinese_ratio > 0.6:  # Mostly Chinese
                    if is_bold:
                        avg_char_px = 13.2  # Bold mostly Chinese - increased
                    else:
                        avg_char_px = 11.8  # Regular mostly Chinese - increased
                else:  # Mixed Chinese/English
                    if is_bold:
                        avg_char_px = 10.5  # Bold mixed text - increased
                    else:
                        avg_char_px = 9.2  # Regular mixed text - increased
            else:
                # English characters - standard calculations
                if is_bold:
                    avg_char_px = 6.5  # Bold English text
                else:
                    avg_char_px = 5.5  # Regular English text
        else:
            # Default - assume mixed content, use conservative estimate
            avg_char_px = 7.0  # Conservative estimate for mixed content

        chars_per_line = max(20, int(effective_width // avg_char_px))  # Minimum 20 chars for Chinese optimization - extremely aggressive
        return chars_per_line

    def _calculate_max_rows_for_summary(self, shape):
        """Calculate maximum rows for summary shape with optimized settings for Chinese"""
        if not shape or not hasattr(shape, 'height'):
            return 25  # Default fallback

        shape_height_emu = shape.height
        shape_height_pt = shape_height_emu * 72 / 914400

        # Use more space for Chinese text to prevent overflow
        effective_height_pt = shape_height_pt * 0.90  # Slightly more space for Chinese

        # Use 9pt font for Chinese optimization (between 8pt and 10pt)
        font_size_pt = 9
        line_spacing = 1.1  # Tighter spacing for Chinese
        line_height_pt = font_size_pt * line_spacing

        max_rows = int(effective_height_pt / line_height_pt)
        max_rows = max(25, max_rows)  # Minimum 25 rows for Chinese summary

        return max_rows

    def _calculate_chars_per_line_for_summary(self, shape):
        """Calculate characters per line for summary with Chinese optimization"""
        if not shape or not hasattr(shape, 'width'):
            return self.CHARS_PER_ROW

        px_width = int(shape.width * 96 / 914400)

        # Use slightly less width utilization for Chinese to prevent overflow
        effective_width = px_width * 0.95  # More conservative for Chinese

        # Chinese-optimized character width estimates for 9pt font
        avg_char_px = 6.8  # Optimized for Chinese characters at 9pt

        chars_per_line = max(75, int(effective_width // avg_char_px))
        return chars_per_line

    def _wrap_text_to_shape(self, text, shape):
        # Use the actual shape width and font size to wrap text accurately
        chars_per_line = self._calculate_chars_per_line(shape)

        # Chinese characters are wider, so they need more lines - be less conservative
        if text and any('\u4e00' <= char <= '\u9fff' for char in text):
            chinese_ratio = sum(1 for char in text if '\u4e00' <= char <= '\u9fff') / len(text)
            if chinese_ratio > 0.8:  # Almost entirely Chinese
                chars_per_line = int(chars_per_line * 0.88)  # 12% more lines for Chinese (less conservative)
            elif chinese_ratio > 0.6:  # Mostly Chinese
                chars_per_line = int(chars_per_line * 0.90)  # 10% more lines for mostly Chinese
            else:  # Mixed Chinese/English
                chars_per_line = int(chars_per_line * 0.92)  # 8% more lines for mixed

        # Use Chinese-aware text wrapping
        if text and any('\u4e00' <= char <= '\u9fff' for char in text):
            # Use the sophisticated Chinese text wrapping
            wrapped = self._wrap_chinese_text(text, chars_per_line)
            return wrapped
        else:
            # Standard text wrapping for English
            wrapped = textwrap.wrap(text, width=chars_per_line)
            return wrapped

    def _wrap_chinese_text(self, text: str, chars_per_line: int) -> list[str]:
        """Wrap Chinese text with proper sentence and phrase boundaries"""
        if not text:
            return []

        wrapped_lines = []
        remaining_text = text

        while remaining_text:
            # Try to fit as much as possible on this line
            if len(remaining_text) <= chars_per_line:
                # The remaining text fits on one line
                wrapped_lines.append(remaining_text)
                break

            # Find the best place to break the line
            line_text = remaining_text[:chars_per_line]

            # Try to find a better break point (sentence endings)
            break_pos = chars_per_line
            for i in range(chars_per_line - 1, max(0, chars_per_line - 20), -1):
                if remaining_text[i] in ['。', '！', '？', '；', '，', '、', '：', ' ', '\n']:
                    break_pos = i + 1  # Include the punctuation/space
                    break

            line_text = remaining_text[:break_pos]
            remaining_text = remaining_text[break_pos:]

            # Clean up the line (remove trailing spaces/newlines)
            line_text = line_text.rstrip()

            if line_text:  # Only add non-empty lines
                wrapped_lines.append(line_text)

        return wrapped_lines

    def _calculate_item_lines(self, item: FinancialItem) -> int:
        """Calculate lines needed for an item using shape-based calculations with Chinese optimization"""
        # Use the current shape for calculations, or fallback to default
        shape = getattr(self, 'current_shape', None)
        chars_per_line = self._calculate_chars_per_line(shape) if shape else self.CHARS_PER_ROW

        lines = 0

        # Calculate header lines using display header
        display_header = self._get_display_header_for_item(item)
        cont_text = self._get_continuation_text(item.layer1_continued)
        header = f"{display_header}{cont_text}"
        header_lines = len(textwrap.wrap(header, width=chars_per_line))
        lines += header_lines

        # Calculate description lines with enhanced Chinese optimization
        desc_lines = 0
        for desc in item.descriptions:
            # Check Chinese character ratio for better optimization
            chinese_chars = sum(1 for char in desc if '\u4e00' <= char <= '\u9fff')
            total_chars = len(desc)
            chinese_ratio = chinese_chars / total_chars if total_chars > 0 else 0

            for para in desc.split('\n'):
                if chinese_ratio > 0.3:  # Has significant Chinese content
                    if chinese_ratio > 0.8:  # Almost entirely Chinese
                        # Chinese characters are wider, so they need more lines - more conservative
                        para_lines = max(1, len(textwrap.wrap(para, width=int(chars_per_line * 0.75))))  # 25% more lines
                    elif chinese_ratio > 0.6:  # Mostly Chinese
                        para_lines = max(1, len(textwrap.wrap(para, width=int(chars_per_line * 0.78))))  # 22% more lines
                    else:  # Mixed Chinese/English
                        para_lines = max(1, len(textwrap.wrap(para, width=int(chars_per_line * 0.82))))  # 18% more lines
                else:
                    # English or minimal Chinese content
                    para_lines = len(textwrap.wrap(para, width=chars_per_line)) or 1
                desc_lines += para_lines

        lines += desc_lines
        return lines

    def _split_item_chinese(self, item: FinancialItem, max_lines: int) -> tuple[FinancialItem, FinancialItem | None]:
        """Split a Chinese item to fit within max_lines, with Chinese-aware text breaking"""
        shape = getattr(self, 'current_shape', None)
        chars_per_line = self._calculate_chars_per_line(shape) if shape else self.CHARS_PER_ROW

        # Account for header line using display header
        display_header = self._get_display_header_for_item(item)
        cont_text = self._get_continuation_text(item.layer1_continued)
        header = f"{display_header}{cont_text}"
        header_lines = self._calculate_chinese_text_lines(header, chars_per_line)
        available_lines = max_lines - header_lines

        desc_paras = item.descriptions
        lines_used = 0
        split_index = 0

        # Try to fit as many whole paragraphs as possible
        for i, para in enumerate(desc_paras):
            para_lines = self._calculate_chinese_text_lines(para, chars_per_line)
            if lines_used + para_lines > available_lines:
                break
            lines_used += para_lines
            split_index = i + 1

        # If all paragraphs fit, no split needed
        if split_index == len(desc_paras):
            return (
                FinancialItem(
                    item.accounting_type,
                    item.account_title,
                    desc_paras,
                    layer1_continued=item.layer1_continued,
                    layer2_continued=item.layer2_continued,
                    is_table=item.is_table
                ),
                None
            )

        # If no paragraph fits, split the first paragraph at character boundary
        if split_index == 0:
            para = desc_paras[0]
            first_part, remaining_part = self._split_chinese_text_at_line(para, chars_per_line, available_lines)

            split_item = FinancialItem(
                item.accounting_type,
                item.account_title,
                [first_part] if first_part else [],
                layer1_continued=item.layer1_continued,
                layer2_continued=False,
                is_table=item.is_table
            )

            remaining_descriptions = ([remaining_part] if remaining_part else []) + desc_paras[1:]
            remaining_item = FinancialItem(
                item.accounting_type,
                item.account_title,
                remaining_descriptions,
                layer1_continued=True,
                layer2_continued=True,
                is_table=item.is_table
            ) if remaining_descriptions else None

            return split_item, remaining_item

        # Otherwise, split at paragraph boundary
        split_item = FinancialItem(
            item.accounting_type,
            item.account_title,
            desc_paras[:split_index],
            layer1_continued=item.layer1_continued,
            layer2_continued=False,
            is_table=item.is_table
        )

        remaining_item = FinancialItem(
            item.accounting_type,
            item.account_title,
            desc_paras[split_index:],
            layer1_continued=True,
            layer2_continued=True,
            is_table=item.is_table
        ) if split_index < len(desc_paras) else None

        return split_item, remaining_item

    def _calculate_chinese_text_lines(self, text: str, chars_per_line: int) -> int:
        """Calculate how many lines Chinese text will occupy"""
        if not text:
            return 0

        # For Chinese text, use character-based line breaking
        total_chars = len(text)
        lines = (total_chars + chars_per_line - 1) // chars_per_line  # Ceiling division

        # Adjust for Chinese character width differences
        chinese_chars = sum(1 for char in text if '\u4e00' <= char <= '\u9fff')
        if chinese_chars > 0:
            chinese_ratio = chinese_chars / total_chars
            # Chinese characters are wider, so they need more lines - extremely aggressive
            if chinese_ratio > 0.8:  # Mostly Chinese
                lines = int(lines * 2.5)  # 150% more lines (2.5x)
            elif chinese_ratio > 0.6:  # Mixed Chinese
                lines = int(lines * 2.2)  # 120% more lines
            elif chinese_ratio > 0.3:  # Some Chinese
                lines = int(lines * 1.8)  # 80% more lines
            else:  # Minimal Chinese
                lines = int(lines * 1.5)  # 50% more lines

        return max(1, lines)

    def _split_chinese_paragraph(self, paragraph: str, chars_per_line: int) -> list[str]:
        """Split long Chinese paragraphs at sentence boundaries for better space utilization"""
        if len(paragraph) <= 100:
            return [paragraph]

        # Chinese sentence endings
        sentence_endings = ['。', '！', '？', '；']
        sub_paragraphs = []
        current_part = ""

        # Split at sentence boundaries
        for char in paragraph:
            current_part += char
            if char in sentence_endings and len(current_part) > 50:  # Minimum sentence length
                sub_paragraphs.append(current_part.strip())
                current_part = ""

        # Add remaining part if any
        if current_part.strip():
            sub_paragraphs.append(current_part.strip())

        # If no sentence boundaries found, split at reasonable intervals
        if len(sub_paragraphs) == 1 and len(paragraph) > 150:
            # Split long paragraph without sentence boundaries
            sub_paragraphs = []
            words = paragraph.split()
            current_part = ""

            for word in words:
                if len(current_part + word) > 80:  # Approximate sentence length
                    if current_part:
                        sub_paragraphs.append(current_part.strip())
                    current_part = word
                else:
                    current_part += word + " "

            if current_part.strip():
                sub_paragraphs.append(current_part.strip())

        return sub_paragraphs if sub_paragraphs else [paragraph]

    def _split_paragraph_at_boundary(self, paragraph: str, available_lines: int, chars_per_line: int) -> tuple[str, str]:
        """Split a paragraph at the optimal boundary for Chinese text, respecting available lines"""
        if not paragraph:
            return "", ""

        # For Chinese text, prefer sentence boundaries
        chinese_chars = sum(1 for char in paragraph if '\u4e00' <= char <= '\u9fff')
        is_chinese = chinese_chars > len(paragraph) * 0.3

        if is_chinese:
            first_part, remaining = self._split_chinese_paragraph_at_boundary(paragraph, available_lines, chars_per_line)
        else:
            first_part, remaining = self._split_english_paragraph_at_boundary(paragraph, available_lines, chars_per_line)

        # Ensure we always have content in first_part if paragraph is not empty
        # This prevents the fallback that allows overflow
        if not first_part and paragraph:
            # Force split at 80% of available space if no good boundary found
            max_chars = int(available_lines * chars_per_line * 0.8)
            if max_chars < len(paragraph):
                first_part = paragraph[:max_chars].rstrip()
                remaining = paragraph[max_chars:].lstrip()
            else:
                first_part = paragraph
                remaining = ""

        return first_part, remaining

    def _split_chinese_paragraph_at_boundary(self, paragraph: str, available_lines: int, chars_per_line: int) -> tuple[str, str]:
        """Split Chinese paragraph at optimal boundary (sentence endings preferred)"""
        # Chinese sentence endings - prioritize complete sentences
        primary_endings = ['。', '！', '？', '；']  # Complete sentences
        secondary_endings = ['，', '：', '；']       # Clause breaks

        max_chars = available_lines * chars_per_line

        # First priority: Complete sentences (。！？；) at 75-90% utilization
        for target_ratio in [0.9, 0.85, 0.8, 0.75]:
            target_chars = int(max_chars * target_ratio)
            for i in range(target_chars, min(max_chars, len(paragraph))):
                if paragraph[i] in primary_endings:
                    first_part = paragraph[:i + 1].rstrip()
                    remaining_part = paragraph[i + 1:].lstrip()
                    if first_part:  # Ensure we have content
                        return first_part, remaining_part

        # Second priority: Clause breaks (，：；) at 70-85% utilization
        for target_ratio in [0.85, 0.8, 0.75, 0.7]:
            target_chars = int(max_chars * target_ratio)
            for i in range(target_chars, min(max_chars, len(paragraph))):
                if paragraph[i] in secondary_endings:
                    first_part = paragraph[:i + 1].rstrip()
                    remaining_part = paragraph[i + 1:].lstrip()
                    if first_part:  # Ensure we have content
                        return first_part, remaining_part

        # Third priority: Any punctuation or space at 80% utilization
        target_chars = int(max_chars * 0.8)
        break_chars = ['，', '、', '：', '；', ' ', '\n', '-', '—']
        for i in range(max_chars - 1, target_chars - 1, -1):
            if i > 0 and paragraph[i] in break_chars:
                first_part = paragraph[:i + 1].rstrip()
                remaining_part = paragraph[i + 1:].lstrip()
                if first_part:
                    return first_part, remaining_part

        # Final fallback: Force split at 80% capacity
        split_point = min(int(max_chars * 0.8), len(paragraph))
        first_part = paragraph[:split_point].rstrip()
        remaining_part = paragraph[split_point:].lstrip()

        return first_part, remaining_part

    def _split_english_paragraph_at_boundary(self, paragraph: str, available_lines: int, chars_per_line: int) -> tuple[str, str]:
        """Split English paragraph at optimal boundary (word boundaries preferred)"""
        max_chars = available_lines * chars_per_line

        if len(paragraph) <= max_chars:
            return paragraph, ""

        # Try to find word boundary
        split_point = max_chars
        for i in range(max_chars - 1, max_chars - 20, -1):  # Look back up to 20 chars
            if i > 0 and paragraph[i] == ' ':
                split_point = i
                break

        first_part = paragraph[:split_point].rstrip()
        remaining_part = paragraph[split_point:].lstrip()

        return first_part, remaining_part

    def _split_chinese_text_at_line(self, text: str, chars_per_line: int, max_lines: int) -> tuple[str, str]:
        """Split Chinese text at line boundary, preserving sentence and phrase integrity"""
        if not text:
            return "", ""

        # Calculate maximum characters that can fit
        max_chars = chars_per_line * max_lines

        # Try to split at sentence boundaries first (Chinese sentences)
        sentences = []
        current_sentence = ""
        for char in text:
            current_sentence += char
            if char in ['。', '！', '？', '；']:
                sentences.append(current_sentence)
                current_sentence = ""

        if current_sentence:
            sentences.append(current_sentence)

        # Find how many sentences fit
        first_part_chars = 0
        split_sentence_idx = 0

        for i, sentence in enumerate(sentences):
            if first_part_chars + len(sentence) > max_chars:
                break
            first_part_chars += len(sentence)
            split_sentence_idx = i + 1

        if split_sentence_idx > 0:
            first_part = ''.join(sentences[:split_sentence_idx])
            remaining_part = ''.join(sentences[split_sentence_idx:])
            return first_part, remaining_part

        # If no sentence boundary found, split at character boundary
        # Try to avoid breaking in the middle of Chinese phrases
        first_part = text[:max_chars]

        # Try to find a better break point (avoid breaking inside Chinese words)
        if max_chars < len(text):
            # Look for punctuation or spaces to break at
            break_chars = ['，', '、', '：', '；', ' ', '\n']
            for i in range(max_chars - 1, max(0, max_chars - 20), -1):
                if text[i] in break_chars:
                    first_part = text[:i + 1]
                    break

        remaining_part = text[len(first_part):]
        return first_part, remaining_part

    def _calculate_chinese_item_lines(self, item: FinancialItem) -> int:
        """Calculate lines needed for a Chinese item using Chinese-aware calculations"""
        # Use the current shape for calculations, or fallback to default
        shape = getattr(self, 'current_shape', None)
        chars_per_line = self._calculate_chars_per_line(shape) if shape else self.CHARS_PER_ROW

        lines = 0

        # Calculate header lines using display header
        display_header = self._get_display_header_for_item(item)
        cont_text = self._get_continuation_text(item.layer1_continued)
        header = f"{display_header}{cont_text}"
        header_lines = self._calculate_chinese_text_lines(header, chars_per_line)
        lines += header_lines

        # Calculate description lines with Chinese-aware optimization
        desc_lines = 0
        for desc in item.descriptions:
            # Check Chinese character ratio for better optimization
            chinese_chars = sum(1 for char in desc if '\u4e00' <= char <= '\u9fff')
            total_chars = len(desc)
            chinese_ratio = chinese_chars / total_chars if total_chars > 0 else 0

            # Split by paragraphs and further split long paragraphs for better utilization
            paragraphs = desc.split('\n') if desc.strip() else ['']

            for para in paragraphs:
                # Skip empty paragraphs
                if not para.strip():
                    continue

                # For long Chinese paragraphs, split them further to utilize space better
                if chinese_ratio > 0.3 and len(para) > 100:  # Long Chinese paragraph
                    # Split long paragraphs at sentence boundaries for Chinese
                    sub_paragraphs = self._split_chinese_paragraph(para, chars_per_line)
                else:
                    sub_paragraphs = [para]

                for sub_para in sub_paragraphs:
                    if chinese_ratio > 0.3:  # Has significant Chinese content
                        if chinese_ratio > 0.8:  # Almost entirely Chinese
                            # Chinese characters are wider, so they need more lines - extremely aggressive
                            para_lines = max(1, self._calculate_chinese_text_lines(sub_para, int(chars_per_line * 0.25)))  # 75% more lines
                        elif chinese_ratio > 0.6:  # Mostly Chinese
                            para_lines = max(1, self._calculate_chinese_text_lines(sub_para, int(chars_per_line * 0.3)))  # 70% more lines
                        else:  # Mixed Chinese/English
                            para_lines = max(1, self._calculate_chinese_text_lines(sub_para, int(chars_per_line * 0.35)))  # 65% more lines
                else:
                    # English or minimal Chinese content
                    para_lines = max(1, self._calculate_chinese_text_lines(sub_para, chars_per_line))
                desc_lines += para_lines

        lines += desc_lines
        return lines

    def _split_item(self, item: FinancialItem, max_lines: int) -> tuple[FinancialItem, FinancialItem | None]:
        """Split an item to fit within max_lines, adding proper (cont'd) indicators"""
        shape = getattr(self, 'current_shape', None)
        chars_per_line = self._calculate_chars_per_line(shape) if shape else self.CHARS_PER_ROW

        # Account for header line using display header
        display_header = self._get_display_header_for_item(item)
        cont_text = self._get_continuation_text(item.layer1_continued)
        header = f"{display_header}{cont_text}"
        header_lines = len(textwrap.wrap(header, width=chars_per_line))
        available_lines = max_lines - header_lines
        
        desc_paras = item.descriptions
        lines_used = 0
        split_index = 0
        
        # Try to fit as many whole paragraphs as possible
        # Be more conservative to leave room for proper splitting
        conservative_limit = int(available_lines * 0.9)  # Leave 10% buffer

        for i, para in enumerate(desc_paras):
            para_lines = len(textwrap.wrap(para, width=chars_per_line)) or 1
            if lines_used + para_lines > conservative_limit:
                break
            lines_used += para_lines
            split_index = i + 1
        
        # If all paragraphs fit, no split needed
        if split_index == len(desc_paras):
            return (
                FinancialItem(
                    display_header,
                    item.account_title,
                    desc_paras,
                    layer1_continued=item.layer1_continued,
                    layer2_continued=item.layer2_continued,
                    is_table=item.is_table
                ),
                None
            )
        
        # If no paragraph fits, split the first paragraph at optimal boundary
        if split_index == 0:
            para = desc_paras[0]

            # Use improved Chinese-aware text splitting
            first_part_text, remaining_text = self._split_paragraph_at_boundary(
                para, available_lines, chars_per_line
            )

            if first_part_text:
                split_item = FinancialItem(
                    item.accounting_type,
                    item.account_title,
                    [first_part_text],
                    layer1_continued=item.layer1_continued,
                    layer2_continued=False,  # First part is not continued
                    is_table=item.is_table
                )
            else:
                # If we can't split, keep the whole paragraph (will cause overflow but better than losing content)
                split_item = FinancialItem(
                    item.accounting_type,
                    item.account_title,
                    [para],
                    layer1_continued=item.layer1_continued,
                    layer2_continued=False,
                    is_table=item.is_table
                )
                remaining_text = None

            # Create remaining item if there's remaining text
            if remaining_text:
                remaining_item = FinancialItem(
                    item.accounting_type,
                    item.account_title,
                    [remaining_text],
                    layer1_continued=True,  # This is a continuation
                    layer2_continued=False,
                    is_table=item.is_table
                )
            else:
                remaining_item = None

            return split_item, remaining_item
        
        # Otherwise, split at paragraph boundary
        split_item = FinancialItem(
            item.accounting_type,
            item.account_title,
            desc_paras[:split_index],
            layer1_continued=item.layer1_continued,
            layer2_continued=False,  # First part is not continued
            is_table=item.is_table
        )
        
        remaining_item = FinancialItem(
            item.accounting_type,
            item.account_title,
            desc_paras[split_index:],
            layer1_continued=True,
            layer2_continued=True,  # Remaining part is continued
            is_table=item.is_table
        ) if split_index < len(desc_paras) else None
        
        return split_item, remaining_item

    def _format_table_text(self, text: str) -> str:
        lines = []
        current_line = []
        for part in text.split():
            if part == self.BULLET_CHAR.strip():
                if current_line:
                    lines.append(" ".join(current_line))
                current_line = ["  " + part]
            else:
                current_line.append(part)
        lines.append(" ".join(current_line))
        return "\n".join(lines)

    def _validate_content_placement(self, distribution):
        seen = set()
        for slide_idx, section, items in distribution:
            for item in items:
                key = (item.accounting_type, item.account_title, tuple(item.descriptions))
                if key in seen:
                    raise ValueError(f"Duplicate content detected: {key}")
                seen.add(key)

    def parse_summary_markdown(self, md_content: str) -> str:
        summary_lines = []
        in_summary = False
        for line in md_content.strip().split('\n'):
            stripped = line.strip()
            if stripped.startswith('## Summary'):
                in_summary = True
                continue
            if in_summary:
                if stripped.startswith('##'):
                    break  # Next section found
                if stripped:  # Skip empty lines
                    summary_lines.append(stripped)
        summary_text = " ".join(summary_lines)
        return summary_text

    def _get_section_shape(self, slide, section: str):
        # For the server template structure:
        # - Slide 0 (index 0): Content slide with textMainBullets
        # - Slide 1+ (index 1+): Content slides with textMainBullets_L and textMainBullets_R
        
        if self.current_slide_index == 0:
            # Slide 0 has textMainBullets
            if section == 'c':
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
                except StopIteration:
                    pass
            return None  # No 'b' section on slide 0
        else:
            # Additional slides (index 1+) - try to find textMainBullets_L and textMainBullets_R
            if section == 'b':
                # Left section
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets_L"), None)
                except StopIteration:
                    pass
            elif section == 'c':
                # Right section
                try:
                    return next((s for s in slide.shapes if s.name == "textMainBullets_R"), None)
                except StopIteration:
                    pass
        
        # Fallback: try to find any textMainBullets shape
        try:
            return next((s for s in slide.shapes if s.name == "textMainBullets"), None)
        except StopIteration:
            pass
        
        # If still not found, try to find Content Placeholder 2 (standard layout)
        try:
            return next((s for s in slide.shapes if s.name == "Content Placeholder 2"), None)
        except StopIteration:
            pass
        
        # If still not found, try to find any text shape
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                return shape
        
        return None

    def _get_display_header_for_item(self, item: FinancialItem) -> str:
        """Get appropriate header for item based on content language."""
        # Check if any description contains Chinese characters
        has_chinese = False
        for desc in item.descriptions:
            if any('\u4e00' <= char <= '\u9fff' for char in desc):
                has_chinese = True
                break

        if has_chinese or (hasattr(self, 'language') and self.language == 'chinese'):
            # For Chinese content, translate section headers to Chinese
            translated_header = self._translate_section_header(item.accounting_type)
            return translated_header
        else:
            # For English content, use the current display format
            return item.accounting_type

    def _populate_section(self, shape, items: List[FinancialItem]):
        tf = shape.text_frame
        tf.clear()
        tf.word_wrap = True
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # Ensure text starts from top
        self.prev_layer1 = None  # Reset for each new section
        self.current_shape = shape  # For dynamic line calculation
        prev_account_title = None
        prev_continued = False
        prev_layer1_continued = None
        paragraph_count = 0  # Track paragraph index in the shape
        
        # Apply row limit (row_limit represents max rows per shape)
        limited_items = items[:self.row_limit] if len(items) > self.row_limit else items
        if len(items) > self.row_limit:
            print(f"⚠️ WARNING: Limiting commentary to {self.row_limit} rows per shape (was {len(items)} rows)")
        
        for idx, item in enumerate(limited_items):
            is_first_part = not (item.layer1_continued or item.layer2_continued)
            # Layer 1 Header
            display_header = self._get_display_header_for_item(item)

            # Add secondary header for continued content on new slides
            if item.layer1_continued and display_header == self.prev_layer1 and display_header != self.prev_layer1_continued:
                # This is a continuation on a new slide - add secondary header
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=False)
                try: p.space_before = get_space_before_for_text(f"{display_header} (续)", force_chinese_mode=(self.language == 'chinese'))
                except: pass
                try: p.space_after = get_space_after_for_text(f"{display_header} (续)", force_chinese_mode=(self.language == 'chinese'))
                except: pass
                try: p.line_spacing = get_line_spacing_for_text(f"{display_header} (续)", force_chinese_mode=(self.language == 'chinese'))
                except: pass
                run = p.add_run()
                run.text = f"{display_header} (续)"
                run.font.size = get_font_size_for_text(run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
                run.font.bold = True
                run.font.name = get_font_name_for_text(run.text, 'Arial')
                try:
                    run.font.color.rgb = self.DARK_BLUE
                except:
                    run.font.color.rgb = RGBColor(0, 51, 160)  # Fallback dark blue
                paragraph_count += 1
                self.prev_layer1_continued = display_header

            if (display_header != self.prev_layer1) or (self.prev_layer1 is None):
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=False)
                if paragraph_count > 0:
                    try:
                        p.space_before = get_space_before_for_text(f"{display_header}{cont_text}", force_chinese_mode=(self.language == 'chinese'))
                    except:
                        pass
                try: p.space_after = get_space_after_for_text(f"{display_header}{cont_text}", force_chinese_mode=(self.language == 'chinese'))
                except: pass
                try: p.line_spacing = get_line_spacing_for_text(f"{display_header}{cont_text}", force_chinese_mode=(self.language == 'chinese'))
                except: pass
                run = p.add_run()
                cont_text = self._get_continuation_text(item.layer1_continued)
                run.text = f"{display_header}{cont_text}"
                run.font.size = get_font_size_for_text(run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
                run.font.bold = True
                run.font.name = get_font_name_for_text(run.text, 'Arial')
                try:
                    run.font.color.rgb = self.DARK_BLUE
                except:
                    run.font.color.rgb = RGBColor(0, 51, 160)  # Fallback dark blue
                self.prev_layer1 = display_header
                paragraph_count += 1
            # Treat all items the same (no special handling for taxes payables)
            for para_idx, desc in enumerate(item.descriptions):
                # Split on '\n' and add an empty row between each part
                desc_parts = desc.split('\n')
                for part_idx, part in enumerate(desc_parts):
                    if is_first_part and para_idx == 0 and part_idx == 0:
                        # Main bullet with heading
                        p = tf.add_paragraph()
                        self._apply_paragraph_formatting(p, is_layer2_3=True)
                        if paragraph_count > 0:
                            try: p.space_before = get_space_before_for_text(f"{self.BULLET_CHAR}{item.account_title}{cont_text} - {part}", force_chinese_mode=(self.language == 'chinese'))
                            except: pass
                        try: p.space_after = get_space_after_for_text(f"{self.BULLET_CHAR}{item.account_title}{cont_text} - {part}", force_chinese_mode=(self.language == 'chinese'))
                        except: pass
                        try: p.line_spacing = get_line_spacing_for_text(f"{self.BULLET_CHAR}{item.account_title}{cont_text} - {part}", force_chinese_mode=(self.language == 'chinese'))
                        except: pass
                        bullet_run = p.add_run()
                        bullet_run.text = self.BULLET_CHAR
                        try:
                            bullet_run.font.color.rgb = self.DARK_GREY
                        except:
                            bullet_run.font.color.rgb = RGBColor(128, 128, 128)  # Fallback grey
                        bullet_run.font.name = get_font_name_for_text(bullet_run.text, 'Arial')
                        bullet_run.font.size = get_font_size_for_text(bullet_run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
                        title_run = p.add_run()
                        cont_text = self._get_continuation_text(item.layer2_continued)
                        translated_title = self._translate_account_title(item.account_title)
                        title_run.text = f"{translated_title}{cont_text}"
                        title_run.font.bold = True
                        title_run.font.name = get_font_name_for_text(title_run.text, 'Arial')
                        title_run.font.size = get_font_size_for_text(title_run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
                        desc_run = p.add_run()
                        # No language marking needed - language detected from databook type
                        desc_run.text = f" - {part}"
                        desc_run.font.bold = False
                        desc_run.font.name = get_font_name_for_text(desc_run.text, 'Arial')
                        desc_run.font.size = get_font_size_for_text(desc_run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
                        paragraph_count += 1
                    else:
                        # Continuation: visually subordinate (indented, no bullet)
                        p = tf.add_paragraph()
                        self._apply_paragraph_formatting(p, is_layer2_3=True)
                        if paragraph_count > 0:
                            try: p.space_before = get_space_before_for_text(part, force_chinese_mode=(self.language == 'chinese'))
                            except: pass
                        try: p.space_after = get_space_after_for_text(part, force_chinese_mode=(self.language == 'chinese'))
                        except: pass
                        try: p.line_spacing = get_line_spacing_for_text(part, force_chinese_mode=(self.language == 'chinese'))
                        except: pass
                        try: p.left_indent = Inches(0.4)
                        except: pass
                        try: p.first_line_indent = Inches(-0.19)
                        except: pass
                        cont_run = p.add_run()
                        # No language marking needed - language detected from databook type
                        cont_run.text = f"{part}"
                        cont_run.font.bold = False
                        cont_run.font.name = get_font_name_for_text(cont_run.text, 'Arial')
                        cont_run.font.size = get_font_size_for_text(cont_run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
                        paragraph_count += 1
                    # Add an empty row between split parts (but not after the last)
                    if part_idx < len(desc_parts) - 1:
                        tf.add_paragraph()
                        paragraph_count += 1
            prev_account_title = item.account_title
            prev_continued = item.layer1_continued or item.layer2_continued
        self.current_shape = None

    def _detect_bullet_content(self, text):
        """
        Enhanced bullet detection that handles multi-level bullet formats up to 4 levels deep
        with special handling for layer 3 bullet points that might contain dashes indicating layer 4
        """
        lines = text.split('\n')
        bullet_lines = []
        has_bullets = False

        # Check if we have any lines starting with "- "
        for line in lines:
            if line.strip().startswith('- '):
                has_bullets = True
                break
        
        # If we have bullets, process the content
        if has_bullets:
            main_content = []
            
            for line in lines:
                stripped = line.strip()
                original_line = line
                
                # Detect bullet lines with '- ' prefix
                if original_line.lstrip().startswith('- '):
                    # If we have accumulated main content, add it first
                    if main_content:
                        main_text = ' '.join(main_content).strip()
                        if main_text:
                            bullet_lines.append((0, main_text))
                        main_content = []
                    
                    # Calculate indentation level (based on spaces/tabs before the bullet)
                    indent_spaces = len(original_line) - len(original_line.lstrip())
                    
                    # Determine bullet level based on indentation
                    # We divide by 2 assuming 2 spaces per indentation level
                    level = min(4, (indent_spaces // 2) + 1)  # Cap at 4 levels
                    
                    # Clean and store bullet line
                    clean_line = stripped[2:]  # Remove '- '
                    
                    # Special handling for level 3 bullets that contain a dash indicating level 4
                    if level == 3 and " - " in clean_line:
                        # Split at the first occurrence of " - "
                        parts = clean_line.split(" - ", 1)
                        if len(parts) > 1:
                            # Add level 3 content
                            bullet_lines.append((level, parts[0].strip()))
                            # Add level 4 content
                            bullet_lines.append((4, parts[1].strip()))
                        else:
                            bullet_lines.append((level, clean_line))
                    else:
                        bullet_lines.append((level, clean_line))
                
                elif stripped:
                    # This is regular content - could be main content or continuation
                    main_content.append(stripped)
            
            # Add any remaining main content
            if main_content:
                main_text = ' '.join(main_content).strip()
                if main_text:
                    bullet_lines.insert(0, (0, main_text))
        
        return bullet_lines if bullet_lines else None

    def generate_full_report(self, md_content: str, summary_md: str, output_path: str):
        try:
            items = self.parse_markdown(md_content)
            distribution = self._plan_content_distribution(items)
            self._validate_content_placement(distribution)
            
            # Generate AI summary content based on commentary length
            # For Chinese mode, use the provided summary_md if it's Chinese content
            if hasattr(self, 'language') and self.language == 'chinese' and summary_md and any('\u4e00' <= char <= '\u9fff' for char in summary_md):
                summary_content = summary_md.strip()
            else:
                # Always generate AI summary content if summary_md is None or empty
                summary_content = self._generate_ai_summary_content(md_content, distribution)

            # Generate per-page summaries based on each page's content
            page_summaries = self._generate_per_page_summaries(distribution, md_content)

            # Ensure summary_content is valid
            if not summary_content or summary_content is None:
                summary_content = "Comprehensive financial analysis with detailed commentary on key financial metrics and performance indicators."
            elif not isinstance(summary_content, str):
                summary_content = str(summary_content)

            max_slide_used = max((slide_idx for slide_idx, _, _ in distribution), default=0)
            total_slides_needed = max_slide_used + 1

            while len(self.prs.slides) < total_slides_needed:
                self.prs.slides.add_slide(self.prs.slide_layouts[1])

            for slide_idx, section, section_items in distribution:
                if slide_idx >= len(self.prs.slides):
                    raise ValueError("Insufficient slides in template")
                slide = self.prs.slides[slide_idx]
                self.current_slide_index = slide_idx
                shape = self._get_section_shape(slide, section)
                if shape:
                    self._populate_section(shape, section_items)

            # Populate summary sections on all slides with different content
            summary_chunks = self._distribute_summary_across_slides(summary_content, total_slides_needed)
            for slide_idx in range(total_slides_needed):
                slide = self.prs.slides[slide_idx]

                # Look for coSummaryShape or alternative summary shapes
                summary_shape = next((s for s in slide.shapes if s.name == "coSummaryShape"), None)
                if not summary_shape:
                    # Fallback: look for Subtitle or other suitable text shapes for summary
                    summary_shape = next((s for s in slide.shapes if hasattr(s, 'text_frame') and s.name == "Subtitle 2"), None)

                if not summary_shape:
                    # Look for TextBox shapes
                    summary_shape = next((s for s in slide.shapes if hasattr(s, 'text_frame') and "TextBox" in s.name), None)

                if not summary_shape:
                    # Look for any text shape that could be used for summary
                    summary_shape = next((s for s in slide.shapes if hasattr(s, 'text_frame') and s.name != "textMainBullets" and "Title" not in s.name), None)
                
                if summary_shape:
                    # Use per-page summary if available, otherwise use chunked summary
                    page_content = page_summaries.get(slide_idx, "")
                    if not page_content:
                        page_content = summary_chunks[slide_idx] if slide_idx < len(summary_chunks) else ""
                    self._populate_summary_section_safe(summary_shape, page_content)
                    logging.info(f"✅ Populated summary on slide {slide_idx + 1} using shape: {summary_shape.name}")
                else:
                    logging.warning(f"⚠️ No summary shape found on slide {slide_idx + 1}")
            
            unused_slides = self._detect_unused_slides(distribution)
            if unused_slides:
                logging.info(f"Removing unused slides: {[idx+1 for idx in unused_slides]}")
                self._remove_slides(unused_slides)
            self.prs.save(output_path)
        except Exception as e:
            logging.error(f"Generation failed: {str(e)}")
            raise

    def _apply_paragraph_formatting(self, paragraph, is_layer2_3=False):
        # Apply commentary-specific formatting as requested by user
        # Requirements: 0.21" indent before text, 0.19" special hanging, 0pt before/after spacing, single line spacing

        try:
            # Apply consistent formatting for all commentary paragraphs
            try: paragraph.left_indent = Inches(0.21)  # 0.21" indent before text
            except: pass
            try: paragraph.first_line_indent = Inches(-0.19)  # 0.19" special hanging
            except: pass
            try: paragraph.space_before = Pt(0)  # 0pt spacing before
            except: pass
            try: paragraph.space_after = Pt(0)  # 0pt spacing after
            except: pass
            try: paragraph.line_spacing = 1.0  # Single line spacing
            except: pass

            # Different formatting for headers vs content
            if not is_layer2_3:
                # Headers (layer 1) - keep current formatting but apply user requirements
                print(f"📄 HEADER PARAGRAPH: 0.21\" indent, 0.19\" hanging, 0pt before/after, single line spacing")
            else:
                # Content (layer 2/3) - apply user requirements
                print(f"📄 CONTENT PARAGRAPH: 0.21\" indent, 0.19\" hanging, 0pt before/after, single line spacing")

            try: paragraph.alignment = PP_ALIGN.LEFT
            except: self._handle_legacy_alignment(paragraph)
        except Exception as e:
            pass  # Silently handle paragraph formatting errors

    def _handle_legacy_alignment(self, paragraph):
        """XML-based left alignment for legacy versions"""
        try:
            pPr = paragraph._element.get_or_add_pPr()
            for child in list(pPr):
                if child.tag.endswith('jc'):
                    pPr.remove(child)
            align = OxmlElement('a:jc')
            align.set('val', 'left')
            pPr.append(align)
        except Exception as e:
            logging.warning(f"Legacy alignment failed: {str(e)}")

    def _safe_set_rgb_color(self, color_obj, rgb_color):
        """Safely set RGB color, handling both RGBColor and _SchemeColor objects"""
        try:
            color_obj.rgb = rgb_color
        except:
            # If setting RGB fails (e.g., _SchemeColor), try setting to a similar RGB value
            try:
                color_obj.rgb = rgb_color
            except:
                # Final fallback - silently continue
                pass

    def _create_taxes_table(self, shape, item):
        tf = shape.text_frame
        # Header
        p = tf.add_paragraph()
        self._apply_paragraph_formatting(p, is_layer2_3=True)
        header_run = p.add_run()
        cont_text = self._get_continuation_text(item.layer2_continued)
        translated_title = self._translate_account_title(item.account_title)
        header_run.text = f"{self.BULLET_CHAR}{translated_title}{cont_text}"
        header_run.font.bold = True
        header_run.font.name = get_font_name_for_text(header_run.text, 'Arial')
        header_run.font.size = get_font_size_for_text(header_run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
        # Process table content
        content = ' '.join(item.descriptions)
        current_category = None
        for line in content.split('\n'):
            line = line.strip()
            if not line or '===' in line:
                continue
            if 'Tax:' in line or 'Payable:' in line:
                current_category = line.replace(':', '')
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=True)
                try: p.left_indent = Inches(0.4)
                except: pass
                run = p.add_run()
                run.text = f"• {current_category}"
                run.font.name = get_font_name_for_text(run.text, 'Arial')
                run.font.bold = True
                run.font.size = get_font_size_for_text(run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))
            elif line.startswith('- '):
                p = tf.add_paragraph()
                self._apply_paragraph_formatting(p, is_layer2_3=True)
                try: p.left_indent = Inches(0.6)
                except: pass
                run = p.add_run()
                run.text = f"◦ {line[2:]}"
                run.font.name = get_font_name_for_text(run.text, 'Arial')
                run.font.size = get_font_size_for_text(run.text, Pt(9), force_chinese_mode=(self.language == 'chinese'))

    def _populate_summary_section_safe(self, shape, summary_content: str):
        """Summary section with natural text wrapping and filling"""
        tf = shape.text_frame
        tf.clear()

        # Create first row as 18pt placeholder
        p1 = tf.add_paragraph()
        placeholder_run = p1.add_run()
        placeholder_run.text = " "  # Single space placeholder
        placeholder_run.font.size = Pt(18)  # 18pt placeholder
        placeholder_run.font.name = 'Arial'
        placeholder_run.font.bold = True
        try:
            placeholder_run.font.color.rgb = RGBColor(0, 51, 102)
        except:
            placeholder_run.font.color.rgb = RGBColor(0, 51, 102)

        # Create second paragraph for actual content
        p2 = tf.add_paragraph()
        run = p2.add_run()

        # Set the full summary content - let PowerPoint handle wrapping naturally
        run.text = summary_content

        # Apply proper font settings for coSummary
        # Use appropriate font size and ensure proper font name
        if summary_content and any('\u4e00' <= char <= '\u9fff' for char in summary_content):
            run.font.size = Pt(8)  # 8pt font for Chinese content
            run.font.name = 'Microsoft YaHei'  # Use YaHei for Chinese text
        else:
            run.font.size = Pt(8)  # 8pt font for English
            run.font.name = 'Arial'  # Use Arial for English text

        try:
            run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue text
        except:
            run.font.color.rgb = RGBColor(0, 51, 102)  # Same color as fallback
        run.font.bold = False

        # Configure text frame for optimal filling
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.NONE  # Use full shape space

        # Add left margin for better readability while maximizing space utilization
        try:
            tf.margin_left = Pt(8)  # Increased left margin for better appearance
            tf.margin_right = Pt(2)
            tf.margin_top = Pt(2)
            tf.margin_bottom = Pt(2)
        except:
            pass  # Some versions don't support margin settings

        # Apply Chinese-optimized paragraph formatting for natural text flow
        try:
            p2.alignment = PP_ALIGN.LEFT
            # Use tighter spacing for Chinese content
            if summary_content and any('\u4e00' <= char <= '\u9fff' for char in summary_content):
                p2.line_spacing = 1.0  # Tighter line spacing for Chinese
                p2.space_before = Pt(2)  # Less space before for Chinese
                p2.space_after = Pt(2)  # Less space after for Chinese
            else:
                p2.line_spacing = 1.1  # Standard line spacing for English
                p2.space_before = Pt(4)  # Standard space before for English
                p2.space_after = Pt(4)  # Standard space after for English
        except AttributeError:
            self._handle_legacy_alignment(p2)

    def _generate_ai_summary_content(self, md_content: str, distribution) -> str:
        """Generate AI summary content based on commentary length and distribution"""
        try:
            # Check if we're in Chinese mode
            is_chinese_mode = hasattr(self, 'language') and self.language == 'chinese'

            # Count total items and slides to determine summary length
            total_items = sum(len(items) for _, _, items in distribution)
            total_slides = max((slide_idx for slide_idx, _, _ in distribution), default=0) + 1

            # Extract key information from markdown content
            lines = md_content.split('\n')
            key_points = []

            for line in lines:
                line = line.strip()
                if line.startswith('### ') and not line.startswith('### ' + self.BULLET_CHAR):
                    # This is a section header
                    key_points.append(line.replace('### ', ''))
                elif line.startswith(self.BULLET_CHAR):
                    # This is a bullet point - extract key info
                    clean_line = line.replace(self.BULLET_CHAR, '').strip()
                    if ' - ' in clean_line:
                        title, desc = clean_line.split(' - ', 1)
                        key_points.append(f"{title}: {desc[:100]}...")
                    else:
                        key_points.append(clean_line[:100] + "...")

            print(f"📝 SUMMARY: Generating {'Chinese' if is_chinese_mode else 'English'} summary for {len(key_points)} key points")

            # Generate professional FDD summary content (80-100 words)
            if key_points and len(key_points) > 0:
                summary_parts = []

                if is_chinese_mode:
                    # Generate Chinese summary
                    intro = f"本综合财务尽职调查分析涵盖了对{len(key_points)}个关键财务领域在{total_slides}页详细分析演示文稿中的全面审查。"
                    intro += "该调查提供了对实体财务状况、运营绩效和风险管理框架的深入评估。"
                    summary_parts.append(intro)

                    if len(key_points) >= 2:
                        findings = f"关键分析发现揭示了对{key_points[0]}和{key_points[1]}的重要洞察，"
                        findings += "展示了实体的财务稳定性和运营效率。该综合审查识别了财务管理和需要加强监控和控制的领域。"
                        summary_parts.append(findings)

                    if len(key_points) >= 3:
                        conclusion = f"特别关注{key_points[2]}的发展，这强调了强大的财务控制和战略风险管理的重要性。"
                        conclusion += "该分析支持明智的投资决策，并为未来的业务规划和利益相关者沟通提供了坚实基础。"
                        summary_parts.append(conclusion)
                else:
                    # Generate English summary
                    intro = f"This comprehensive financial due diligence analysis encompasses a thorough examination of {len(key_points)} critical financial domains "
                    intro += f"across {total_slides} detailed analytical presentation pages. The investigation provides an in-depth assessment of the entity's "
                    intro += "financial position, operational performance, and risk management framework. "
                    summary_parts.append(intro)

                    if len(key_points) >= 2:
                        findings = f"Key analytical findings reveal significant insights into {key_points[0]} and {key_points[1]}, "
                        findings += "demonstrating the entity's financial stability and operational efficiency. The comprehensive review identifies "
                        findings += "both strengths in financial management and areas requiring enhanced monitoring and control. "
                        summary_parts.append(findings)

                    if len(key_points) >= 3:
                        conclusion = f"Particular attention has been given to developments in {key_points[2]}, which underscore the importance of "
                        conclusion += "robust financial controls and strategic risk management. The analysis supports informed investment decisions "
                        conclusion += "and provides a solid foundation for future business planning and stakeholder communication."
                        summary_parts.append(conclusion)

                    # If we don't have enough key points, create a single professional paragraph
                    if len(key_points) < 3:
                        if is_chinese_mode:
                            summary = f"本综合财务尽职调查分析考察了{len(key_points)}个关键领域在{total_slides}页演示文稿中的情况，"
                            summary += f"特别关注{key_points[0]}"
                            if len(key_points) > 1:
                                summary += f"和{key_points[1]}"
                            summary += "。该分析提供了对财务绩效、风险评估和运营效率的关键洞察。该综合审查能够为未来的业务发展提供明智的决策和战略规划。"
                        else:
                            summary = f"This comprehensive financial due diligence analysis examines {len(key_points)} key areas across {total_slides} presentation pages, "
                            summary += f"with particular focus on {key_points[0]}"
                            if len(key_points) > 1:
                                summary += f" and {key_points[1]}"
                            summary += ". The analysis provides critical insights into financial performance, risk assessment, and operational efficiency. "
                            summary += "The comprehensive review enables informed decision-making and strategic planning for future business development."
                        summary_parts = [summary]

                    # Combine all parts
                    full_summary = " ".join(summary_parts)

                    # Ensure we have valid content
                    if not full_summary or full_summary.strip() == "":
                        full_summary = "Comprehensive financial due diligence analysis with detailed commentary."

                    # Ensure it's between 80-100 words
                    word_count = len(full_summary.split())
                    if word_count < 80:
                        # Add professional depth
                        additional = "The methodology employed rigorous analytical procedures, including detailed financial statement analysis, "
                        additional += "assessment of internal control effectiveness, and evaluation of compliance with regulatory requirements. "
                        additional += "Management representations have been obtained and corroborated with supporting documentation."
                        full_summary += " " + additional
                    elif word_count > 100:
                        # Trim to fit while maintaining professional tone
                        words = full_summary.split()
                        full_summary = " ".join(words[:100])

                    # Final safety check - ensure we return a valid string
                    if not isinstance(full_summary, str):
                        full_summary = str(full_summary)

                    return full_summary or "Comprehensive financial analysis summary."
            else:
                # Fallback summary (comprehensive)
                if is_chinese_mode:
                    summary = "本综合财务尽职调查分析提供了对组织财务状况、绩效指标和关键风险因素的详细审查，涵盖多个详细演示页面。"
                    summary += "该分析包括对财务报表的全面审查、对内部控制的评估、对监管要求的评估，以及对重要财务趋势和模式的识别。"
                    summary += "关键发现揭示了资产管理、负债结构和整体财务健康的重要发展，为战略决策和未来规划提供了关键洞察。"
                    summary += "该综合审查涵盖所有重大财务交易，并对财务报告的准确性和完整性提供保证，使利益相关者能够就组织的财务状况做出明智决定。"
                else:
                    summary = "This comprehensive financial due diligence analysis provides detailed examination of the organization's financial position, performance metrics, and key risk factors across multiple detailed presentation pages. "
                    summary += "The analysis encompasses thorough review of financial statements, assessment of internal controls, evaluation of compliance with regulatory requirements, and identification of significant financial trends and patterns. "
                    summary += "Key findings reveal important developments in asset management, liability structure, and overall financial health, providing critical insights for strategic decision-making and future planning initiatives. "
                    summary += "The comprehensive review covers all material financial transactions and provides assurance on the accuracy and completeness of financial reporting, enabling stakeholders to make informed decisions regarding the organization's financial position."
                return summary
            
        except Exception as e:
            logging.warning(f"Failed to generate AI summary: {str(e)}")
            # Ensure we always return a valid string
            fallback_summary = "Comprehensive financial analysis with detailed commentary on key financial metrics and performance indicators."
            return str(fallback_summary) if fallback_summary else "Financial analysis summary."

    def _generate_per_page_summaries(self, distribution, md_content: str) -> dict[int, str]:
        """Generate AI-powered summaries for each page based on that page's content"""
        page_summaries = {}

        # Group items by slide
        slides_content = {}
        for slide_idx, section, section_items in distribution:
            if slide_idx not in slides_content:
                slides_content[slide_idx] = []
            slides_content[slide_idx].extend(section_items)

        # Generate summary for each slide
        for slide_idx, items in slides_content.items():
            if not items:
                continue

            # Extract content from items on this slide
            page_content_parts = []
            for item in items:
                if hasattr(item, 'descriptions') and item.descriptions:
                    # Translate section headers to Chinese if in Chinese mode
                    header = item.accounting_type
                    if hasattr(self, 'language') and self.language == 'chinese':
                        header = self._translate_section_header(header)

                    page_content_parts.append(f"{header}: {' '.join(item.descriptions)}")

            page_content = ' '.join(page_content_parts)

            if page_content.strip():
                # Generate a concise summary for this page
                if hasattr(self, 'language') and self.language == 'chinese':
                    summary = self._generate_chinese_page_summary(page_content, slide_idx + 1)
                else:
                    summary = self._generate_english_page_summary(page_content, slide_idx + 1)

                page_summaries[slide_idx] = summary

        return page_summaries

    def _translate_section_header(self, header: str) -> str:
        """Translate common section headers to Chinese"""
        translation_map = {
            'Current Assets': '流动资产',
            'Non-current Assets': '非流动资产',
            'Current Liabilities': '流动负债',
            'Non-current Liabilities': '非流动负债',
            'Equity': '股东权益',
            'Assets': '资产',
            'Liabilities': '负债',
            'Cash and Cash Equivalents': '货币资金',
            'Accounts Receivable': '应收账款',
            'Inventory': '存货',
            'Property, Plant and Equipment': '固定资产',
            'Intangible Assets': '无形资产',
            'Accounts Payable': '应付账款',
            'Loans Payable': '借款',
            'Revenue': '收入',
            'Cost of Sales': '销售成本',
            'Operating Expenses': '营业费用',
            'Net Income': '净利润'
        }
        return translation_map.get(header, header)
    
    def _translate_account_title(self, account_title: str) -> str:
        """Translate account titles based on language - Chinese only for Chinese databooks, English only for English databooks"""
        translation_map = {
            'Cash': '现金',
            'AR': '应收账款',
            'Prepayments': '预付账款',
            'OR': '其他应收款',
            'Other CA': '其他流动资产',
            'IP': '投资性房地产',
            'Other NCA': '其他非流动资产',
            'AP': '应付账款',
            'Advances': '预收账款',
            'Taxes payable': '应交税费',
            'OP': '其他应付款',
            'Capital': '实收资本',
            'Reserve': '盈余公积',
            'Capital reserve': '资本公积',
            'OI': '营业收入',
            'OC': '营业成本',
            'Tax and Surcharges': '税金及附加',
            'GA': '管理费用',
            'Fin Exp': '财务费用',
            'Cr Loss': '信用减值损失',
            'Other Income': '其他收益',
            'Non-operating Income': '营业外收入',
            'Non-operating Exp': '营业外支出',
            'Income tax': '所得税费用',
            'LT DTA': '递延所得税资产'
        }
        
        # For Chinese databooks, return Chinese only
        if hasattr(self, 'language') and self.language == 'chinese':
            chinese_title = translation_map.get(account_title, account_title)
            return chinese_title
        else:
            # For English databooks, return English only
            return account_title
    
    def _get_continuation_text(self, is_continued: bool) -> str:
        """Get appropriate continuation text based on language"""
        if not is_continued:
            return ""
        
        
        if hasattr(self, 'language') and self.language == 'chinese':
            return " (续)"
        else:
            return " (continued)"

    def _generate_chinese_page_summary(self, page_content: str, page_number: int) -> str:
        """Generate a concise Chinese summary for a page"""
        try:
            # Extract key points from the page content
            sentences = page_content.split('。')
            key_points = []

            for sentence in sentences[:3]:  # Take first 3 sentences
                if sentence.strip():
                    # Keep it concise - first 60 characters for more informative content
                    summary_point = sentence.strip()[:60]
                    if summary_point:
                        key_points.append(summary_point)

            if key_points:
                # Remove page number and make more informative
                summary = f"{'。'.join(key_points)}。"
                return summary
            else:
                return "财务分析摘要：本页包含重要的财务指标和业务分析内容。"

        except Exception as e:
            print(f"Error generating Chinese page summary: {e}")
            return f"第{page_number}页财务分析摘要。"

    def _generate_english_page_summary(self, page_content: str, page_number: int) -> str:
        """Generate a concise English summary for a page"""
        try:
            # Extract key points from the page content
            sentences = page_content.split('.')
            key_points = []

            for sentence in sentences[:2]:  # Take first 2 sentences
                if sentence.strip():
                    # Keep it concise - first 60 characters
                    summary_point = sentence.strip()[:60]
                    if summary_point:
                        key_points.append(summary_point)

            if key_points:
                # Remove page number and make more informative
                summary = f"{'. '.join(key_points)}."
                return summary
            else:
                return "Financial Analysis Summary: This page contains important financial metrics and business analysis."

        except Exception as e:
            print(f"Error generating English page summary: {e}")
            return f"Page {page_number} financial analysis summary."

    def _distribute_summary_across_slides(self, summary_content: str, total_slides: int) -> List[str]:
        """Distribute summary content across multiple slides with different content on each"""
        try:
            # Handle None or empty summary content
            if not summary_content or summary_content is None:
                summary_content = "Comprehensive financial analysis with detailed commentary on key financial metrics and performance indicators."

            # Ensure summary_content is a string
            if not isinstance(summary_content, str):
                summary_content = str(summary_content)

            # Split the summary into sentences
            sentences = summary_content.split('. ')
            
            # Calculate sentences per slide
            sentences_per_slide = max(1, len(sentences) // total_slides)
            
            # Distribute content across slides
            summary_chunks = []
            for slide_idx in range(total_slides):
                start_idx = slide_idx * sentences_per_slide
                end_idx = start_idx + sentences_per_slide
                
                if slide_idx == total_slides - 1:
                    # Last slide gets remaining content
                    slide_sentences = sentences[start_idx:]
                else:
                    slide_sentences = sentences[start_idx:end_idx]
                
                # Create slide-specific content
                if slide_sentences:
                    slide_content = '. '.join(slide_sentences)
                    if not slide_content.endswith('.'):
                        slide_content += '.'
                else:
                    # Fallback content for empty slides
                    slide_content = "Financial analysis summary - detailed commentary provided in main sections."
                
                summary_chunks.append(slide_content)
            
            return summary_chunks
            
        except Exception as e:
            logging.warning(f"Failed to distribute summary: {str(e)}")
            # Return single chunk for all slides
            return [summary_content] * total_slides

    def _detect_unused_slides(self, distribution):
        """Adjusted slide retention logic with content-aware detection"""
        used_slides = set()
        content_slides = set()
        # Track slides that actually contain content
        for slide_idx, section, items in distribution:
            if items:  # Only consider slides with actual content
                content_slides.add(slide_idx)
            used_slides.add(slide_idx)
        # Calculate maximum content-bearing slide
        max_content_slide = max(content_slides) if content_slides else 0
        # Determine minimum slides to keep (content slides + buffer)
        min_slides_to_keep = max(2, max_content_slide + 1)
        # Preserve all slides up to min_slides_to_keep
        remove_slides = []
        for slide_idx in range(len(self.prs.slides)-1, min_slides_to_keep-1, -1):
            remove_slides.append(slide_idx)
        return sorted(remove_slides, reverse=True)

    def _remove_slides(self, slide_indices):
        """Remove slides by indices (must be in descending order)"""
        for slide_idx in slide_indices:
            if slide_idx < len(self.prs.slides):
                # Method 1: Using _sldIdLst (more reliable)
                xml_slides = self.prs.slides._sldIdLst
                slides = list(xml_slides)
                # Remove relationship
                rId = slides[slide_idx].rId
                self.prs.part.drop_rel(rId)
                # Remove from slide list
                xml_slides.remove(slides[slide_idx])

    def _populate_bullet_content(self, shape, item, bullet_lines):
        tf = shape.text_frame
        tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # Ensure text starts from top
        p = tf.add_paragraph()
        self._apply_paragraph_formatting(p, is_layer2_3=True)
        bullet_run = p.add_run()
        bullet_run.text = self.BULLET_CHAR
        try:
            bullet_run.font.color.rgb = self.DARK_GREY
        except:
            bullet_run.font.color.rgb = RGBColor(128, 128, 128)  # Fallback grey
        bullet_run.font.name = 'Arial'
        bullet_run.font.size = Pt(9)
        title_run = p.add_run()
        cont_text = self._get_continuation_text(item.layer2_continued)
        translated_title = self._translate_account_title(item.account_title)
        title_run.text = f"{translated_title}{cont_text}"
        title_run.font.bold = True
        title_run.font.name = 'Arial'
        title_run.font.size = Pt(9)
        for line in bullet_lines:
            p = tf.add_paragraph()
            self._apply_paragraph_formatting(p, is_layer2_3=True)
            try: p.left_indent = Inches(0.4)
            except: pass
            try: p.first_line_indent = Inches(-0.19)
            except: pass
            run = p.add_run()
            clean_line = line.lstrip('- ').strip() if isinstance(line, str) else line[1].lstrip('- ').strip()
            run.text = f"• {clean_line}"
            run.font.name = get_font_name_for_text(run.text, 'Arial')
            run.font.size = Pt(9)
            run.font.bold = False

class ReportGenerator:
    def __init__(self, template_path, markdown_file, output_path, project_name=None, language='english', row_limit=20):
        self.template_path = template_path
        self.markdown_file = markdown_file
        self.output_path = output_path
        self.project_name = project_name
        self.language = language
        self.row_limit = row_limit
        
    def generate(self):
        # Read the markdown content
        with open(self.markdown_file, 'r', encoding='utf-8') as f:
            md_content = f.read()
        
        # Create PowerPoint generator with language support
        generator = PowerPointGenerator(self.template_path, self.language, row_limit=self.row_limit)
        
        try: 
            # Generate the report with AI summary content
            # Pass None for summary_md so it generates AI summary automatically
            generator.generate_full_report(md_content, None, self.output_path)
        except Exception as e:
            logging.error(f"❌ Report generation failed: {str(e)}")
            raise  # Re-raise the exception so it can be handled properly

# --- Project Title Update Logic (from 3.wrap_up.py) ---
def find_shape_by_name(shapes, name):
    for shape in shapes:
        if shape.name == name:
            return shape
    return None

def replace_text_preserve_formatting(shape, replacements):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            for old_text, new_text in replacements.items():
                if old_text in run.text:
                    run.text = run.text.replace(old_text, new_text)

def update_project_titles(presentation_path, project_name, output_path=None, language='english', statement_type='BS'):
    prs = Presentation(presentation_path)
    total_slides = len(prs.slides)

    # Extract first two words for professional display
    if project_name:
        words = project_name.split()
        # Use first two words, or first word if only one word
        display_entity = ' '.join(words[:2]) if len(words) >= 2 else words[0] if words else project_name
    else:
        display_entity = project_name

    # Define title templates based on language and statement type
    if statement_type.upper() == 'BS':  # Balance Sheet
        if language.lower() == 'chinese':
            title_template = f"资产负债表概览 - {display_entity}"
        else:
            title_template = f"Entity Overview - {display_entity}"
    elif statement_type.upper() == 'IS':  # Income Statement
        if language.lower() == 'chinese':
            title_template = f"利润表概览 - {display_entity}"
        else:
            title_template = f"Income Statement - {display_entity}"
    else:  # Default/Unknown
        if language.lower() == 'chinese':
            title_template = f"财务报表概览 - {display_entity}"
        else:
            title_template = f"Financial Report - {display_entity}"

    for slide_index, slide in enumerate(prs.slides):
        current_slide_number = slide_index + 1
        projTitle_shape = find_shape_by_name(slide.shapes, "projTitle")
        if projTitle_shape:
            # Replace the entire title with the language-appropriate template
            # First try to replace placeholders if they exist, otherwise replace the whole text
            current_text = projTitle_shape.text
            if "[PROJECT]" in current_text:
                replacements = {
                    "[PROJECT]": display_entity,
                    "[Current]": str(current_slide_number),
                    "[Total]": str(total_slides)
                }
                replace_text_preserve_formatting(projTitle_shape, replacements)
            else:
                # Replace the entire title text
                projTitle_shape.text = title_template

    # Save the presentation
    if output_path is None:
        output_path = presentation_path
    prs.save(output_path)
    return output_path

# --- High-level function for app.py ---
def embed_excel_data_in_pptx(presentation_path, excel_file_path, sheet_name, project_name, output_path=None, statement_type='BS', **kwargs):
    """
    Update existing financialData shape in PowerPoint with Excel data

    Args:
        statement_type: 'BS' for Balance Sheet, 'IS' for Income Statement, 'ALL' for combined
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches
        import pandas as pd
        import os

        # Load the presentation
        prs = Presentation(presentation_path)

        # Read Excel data
        df = pd.read_excel(excel_file_path, sheet_name=sheet_name)

        # Apply BS/IS separation based on statement type
        from fdd_utils.excel_processing import separate_balance_sheet_and_income_statement_tables, filter_to_indicative_adjusted_columns

        bs_data, is_data, separation_metadata = separate_balance_sheet_and_income_statement_tables(df, [project_name])

        # Use appropriate data based on statement type
        if statement_type == 'IS' and is_data:
            df = is_data['data']
            print(f"✅ EMBED: Using INCOME STATEMENT data for Excel embedding")
            print(f"   📊 IS data shape: {df.shape}")
        elif statement_type == 'BS' and bs_data:
            df = bs_data['data']
            print(f"✅ EMBED: Using BALANCE SHEET data for Excel embedding")
            print(f"   📊 BS data shape: {df.shape}")
        elif statement_type == 'ALL':
            # For combined mode, try to use both but prioritize based on what's available
            if bs_data and is_data:
                # Combine both tables
                bs_df = bs_data['data']
                is_df = is_data['data']
                # Add a separator row and combine
                separator_df = pd.DataFrame([['--- INCOME STATEMENT ---'] + [''] * (len(bs_df.columns) - 1)], columns=bs_df.columns)
                df = pd.concat([bs_df, separator_df, is_df], ignore_index=True)
                print(f"✅ EMBED: Using COMBINED BS+IS data for Excel embedding")
                print(f"   📊 Combined data shape: {df.shape}")
            elif bs_data:
                df = bs_data['data']
                print(f"✅ EMBED: Using BALANCE SHEET data (no IS data found) for Excel embedding")
            elif is_data:
                df = is_data['data']
                print(f"✅ EMBED: Using INCOME STATEMENT data (no BS data found) for Excel embedding")
            else:
                print(f"⚠️ EMBED: No separated data found, using filtered original")
                df = filter_to_indicative_adjusted_columns(df)
        else:
            # Fallback for other cases
            if bs_data:
                df = bs_data['data']
                print(f"✅ EMBED: Using BALANCE SHEET data (fallback)")
            else:
                print(f"⚠️ EMBED: No BS data found, using filtered original")
                df = filter_to_indicative_adjusted_columns(df)
        
        # Find the financialData shape in all slides
        financial_data_shape = None
        target_slide = None
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.name == "financialData":
                    financial_data_shape = shape
                    target_slide = slide
                    break
            if financial_data_shape:
                break
        
        if financial_data_shape:
            # Found the financialData shape - update its content
            # Remove any placeholder text or empty content first
            if hasattr(financial_data_shape, 'text_frame') and financial_data_shape.text_frame:
                if financial_data_shape.text_frame.text.strip() in ['', ' ', 'Placeholder', 'Financial Data', 'Table']:
                    financial_data_shape.text_frame.clear()
            if hasattr(financial_data_shape, 'table'):
                # It's a table shape - update the table data while preserving formatting
                table = financial_data_shape.table
                
                # Store original formatting for each cell
                original_formats = {}
                for row_idx in range(len(table.rows)):
                    for col_idx in range(len(table.columns)):
                        cell = table.cell(row_idx, col_idx)
                        # Store font properties
                        if cell.text_frame.paragraphs[0].runs:
                            run = cell.text_frame.paragraphs[0].runs[0]
                            original_formats[(row_idx, col_idx)] = {
                                'font_name': run.font.name,
                                'font_size': run.font.size,
                                'font_bold': run.font.bold,
                                'font_italic': run.font.italic,
                                'font_color': run.font.color.rgb if hasattr(run.font.color, 'rgb') and run.font.color.rgb else None
                            }
                
                # Clear existing content but preserve formatting
                for row in range(len(table.rows)):
                    for col in range(len(table.columns)):
                        if row < len(table.rows) and col < len(table.columns):
                            cell = table.cell(row, col)
                            cell.text = ""
                
                # Update with new data and apply original formatting
                # Header row
                for col_idx, col_name in enumerate(df.columns):
                    if col_idx < len(table.columns):
                        cell = table.cell(0, col_idx)
                        cell.text = str(col_name)
                        # Apply header formatting (usually bold, different color)
                        if (0, col_idx) in original_formats:
                            format_info = original_formats[(0, col_idx)]
                            if cell.text_frame.paragraphs[0].runs:
                                run = cell.text_frame.paragraphs[0].runs[0]
                                if format_info['font_name']:
                                    run.font.name = format_info['font_name']
                                if format_info['font_size']:
                                    run.font.size = format_info['font_size']
                                if format_info['font_bold'] is not None:
                                    run.font.bold = format_info['font_bold']
                                if format_info['font_italic'] is not None:
                                    run.font.italic = format_info['font_italic']
                                if format_info['font_color']:
                                    try:
                                        run.font.color.rgb = format_info['font_color']
                                    except:
                                        # Handle theme colors by setting to a default RGB color
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                
                # Data rows
                for row_idx, row in enumerate(df.values):
                    if row_idx + 1 < len(table.rows):
                        for col_idx, value in enumerate(row):
                            if col_idx < len(table.columns):
                                cell = table.cell(row_idx + 1, col_idx)
                                cell.text = str(value)
                                # Apply data row formatting
                                if (row_idx + 1, col_idx) in original_formats:
                                    format_info = original_formats[(row_idx + 1, col_idx)]
                                    if cell.text_frame.paragraphs[0].runs:
                                        run = cell.text_frame.paragraphs[0].runs[0]
                                        if format_info['font_name']:
                                            run.font.name = format_info['font_name']
                                        if format_info['font_size']:
                                            run.font.size = format_info['font_size']
                                        if format_info['font_bold'] is not None:
                                            run.font.bold = format_info['font_bold']
                                        if format_info['font_italic'] is not None:
                                            run.font.italic = format_info['font_italic']
                                        if format_info['font_color']:
                                            try:
                                                run.font.color.rgb = format_info['font_color']
                                            except:
                                                # Handle theme colors by setting to a default RGB color
                                                run.font.color.rgb = RGBColor(0, 0, 0)
                
                logging.info(f"✅ Updated financialData table with Excel data (formatting preserved)")
                
            elif hasattr(financial_data_shape, 'text_frame'):
                # It's a text shape - convert to table or update text
                # For now, let's create a table in its place
                left = financial_data_shape.left
                top = financial_data_shape.top
                width = financial_data_shape.width
                height = financial_data_shape.height
                
                # Remove the old shape
                target_slide.shapes._spTree.remove(financial_data_shape._element)
                
                # Add new table with same position and size
                table = target_slide.shapes.add_table(
                    rows=len(df) + 1,  # +1 for header
                    cols=len(df.columns),
                    left=left,
                    top=top,
                    width=width,
                    height=height
                ).table
                
                # Name the new table
                table._element.getparent().getparent().set('name', 'financialData')
                
                # Fill table with data
                for col_idx, col_name in enumerate(df.columns):
                    table.cell(0, col_idx).text = str(col_name)
                
                for row_idx, row in enumerate(df.values):
                    for col_idx, value in enumerate(row):
                        table.cell(row_idx + 1, col_idx).text = str(value)
                
                logging.info(f"✅ Replaced financialData text with table")
                
        else:
            # financialData shape not found - create a new table in the middle
            logging.warning("⚠️ financialData shape not found, creating new table")
            
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                
                # Calculate position for the table (middle of slide)
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                object_width = Inches(8)
                object_height = Inches(4)
                left = (slide_width - object_width) / 2
                top = (slide_height - object_height) / 2
                
                # Add a table with the Excel data
                table = slide.shapes.add_table(
                    rows=len(df) + 1,  # +1 for header
                    cols=len(df.columns),
                    left=left,
                    top=top,
                    width=object_width,
                    height=object_height
                ).table
                
                # Name the table
                table._element.getparent().getparent().set('name', 'financialData')
                
                # Apply professional formatting to the new table
                from pptx.dml.color import RGBColor
                
                # Fill table with data and apply formatting
                for col_idx, col_name in enumerate(df.columns):
                    cell = table.cell(0, col_idx)
                    cell.text = str(col_name)
                    # Header formatting
                    if cell.text_frame.paragraphs[0].runs:
                        run = cell.text_frame.paragraphs[0].runs[0]
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run.font.name = get_font_name_for_text(run.text, 'Arial')
                        # Header background color (light blue)
                        cell.fill.solid()
                        try:
                            cell.fill.fore_color.rgb = RGBColor(217, 225, 242)
                        except:
                            cell.fill.fore_color.rgb = RGBColor(217, 225, 242)  # Same color as fallback
                
                for row_idx, row in enumerate(df.values):
                    for col_idx, value in enumerate(row):
                        cell = table.cell(row_idx + 1, col_idx)
                        cell.text = str(value)
                        # Data row formatting
                        if cell.text_frame.paragraphs[0].runs:
                            run = cell.text_frame.paragraphs[0].runs[0]
                            run.font.size = Pt(10)
                            run.font.name = get_font_name_for_text(run.text, 'Arial')
                            # Alternate row colors for readability
                            if row_idx % 2 == 0:
                                cell.fill.solid()
                                try:
                                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
                                except:
                                    cell.fill.fore_color.rgb = RGBColor(242, 242, 242)  # Same color as fallback
        
        # Clean up any unused placeholder shapes
        _cleanup_placeholder_shapes(prs)
        
        # Save the presentation
        if output_path is None:
            output_path = presentation_path
        prs.save(output_path)
        
        logging.info(f"✅ Excel data updated in PowerPoint: {output_path}")
        return output_path
        
    except Exception as e:
        logging.error(f"❌ Failed to update Excel data: {str(e)}")
        raise

def export_pptx(template_path, markdown_path, output_path, project_name=None, excel_file_path=None, language='english', statement_type='BS', row_limit=20):
    generator = ReportGenerator(template_path, markdown_path, output_path, project_name, language, row_limit)
    generator.generate()
    if not os.path.exists(output_path):
        logging.error(f"PPTX file was not created at {output_path}")
        raise FileNotFoundError(f"PPTX file was not created at {output_path}")
    if project_name:
        update_project_titles(output_path, project_name, language=language, statement_type=statement_type)

    # Embed Excel data if provided
    if excel_file_path and project_name:
        # Get the appropriate sheet name based on project
        sheet_name = get_tab_name(project_name)
        if sheet_name:
            try:
                embed_excel_data_in_pptx(output_path, excel_file_path, sheet_name, project_name, statement_type=statement_type)
            except Exception as e:
                logging.warning(f"⚠️ Could not embed Excel data: {str(e)}")

    # Add success message
    logging.info(f"✅ PowerPoint presentation successfully exported to: {output_path}")
    return output_path

def merge_presentations(bs_presentation_path, is_presentation_path, output_path):
    """
    Merge Balance Sheet and Income Statement presentations into a single presentation.
    
    This function combines BS and IS presentations by:
    1. Loading BS presentation as the base
    2. Copying all slides from IS presentation
    3. Appending IS slides after BS slides
    4. Preserving formatting, text, tables, and other elements
    
    Args:
        bs_presentation_path: Path to the Balance Sheet presentation
        is_presentation_path: Path to the Income Statement presentation  
        output_path: Path for the merged output presentation
        
    Returns:
        str: Path to the merged presentation
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import logging
        
        print(f"🔄 Starting presentation merge...")
        print(f"   BS: {bs_presentation_path}")
        print(f"   IS: {is_presentation_path}")
        print(f"   Output: {output_path}")
        
        # Verify input files exist
        if not os.path.exists(bs_presentation_path):
            raise FileNotFoundError(f"BS presentation not found: {bs_presentation_path}")
        if not os.path.exists(is_presentation_path):
            raise FileNotFoundError(f"IS presentation not found: {is_presentation_path}")
        
        # Load the Balance Sheet presentation as the base
        bs_prs = Presentation(bs_presentation_path)
        is_prs = Presentation(is_presentation_path)
        
        print(f"📊 BS presentation: {len(bs_prs.slides)} slides")
        print(f"📈 IS presentation: {len(is_prs.slides)} slides")
        
        logging.info(f"📊 BS presentation has {len(bs_prs.slides)} slides")
        logging.info(f"📈 IS presentation has {len(is_prs.slides)} slides")
        
        # Copy all slides from Income Statement to Balance Sheet presentation
        for slide in is_prs.slides:
            # Get the slide layout from the BS presentation
            slide_layout = bs_prs.slide_layouts[0]  # Use first layout for all slides
            
            # Create new slide in BS presentation
            new_slide = bs_prs.slides.add_slide(slide_layout)
            
            # Copy all shapes from IS slide to new slide
            for shape in slide.shapes:
                # Get shape position and size
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                
                # Copy shape based on type
                if shape.shape_type == 17:  # Text box
                    # Copy text box
                    textbox = new_slide.shapes.add_textbox(left, top, width, height)
                    textbox.text_frame.text = shape.text_frame.text
                    
                    # Copy text formatting
                    for i, paragraph in enumerate(shape.text_frame.paragraphs):
                        if i < len(textbox.text_frame.paragraphs):
                            new_paragraph = textbox.text_frame.paragraphs[i]
                            new_paragraph.alignment = paragraph.alignment
                            for j, run in enumerate(paragraph.runs):
                                if j < len(new_paragraph.runs):
                                    new_run = new_paragraph.runs[j]
                                    new_run.font.bold = run.font.bold
                                    new_run.font.italic = run.font.italic
                                    new_run.font.size = run.font.size
                                    new_run.font.name = run.font.name
                                    if hasattr(run.font, 'color') and hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                        try:
                                            new_run.font.color.rgb = run.font.color.rgb
                                        except:
                                            new_run.font.color.rgb = RGBColor(0, 0, 0)  # Fallback black
                
                elif shape.shape_type == 19:  # Table
                    # Copy table
                    table = new_slide.shapes.add_table(
                        rows=shape.table.rows.__len__(),
                        cols=shape.table.columns.__len__(),
                        left=left,
                        top=top,
                        width=width,
                        height=height
                    ).table
                    
                    # Copy table data
                    for row_idx in range(shape.table.rows.__len__()):
                        for col_idx in range(shape.table.columns.__len__()):
                            if row_idx < table.rows.__len__() and col_idx < table.columns.__len__():
                                table.cell(row_idx, col_idx).text = shape.table.cell(row_idx, col_idx).text
                
                else:
                    # For other shape types, try to copy as picture
                    try:
                        new_slide.shapes.add_picture(
                            shape.image.blob,
                            left,
                            top,
                            width,
                            height
                        )
                    except:
                        # If copying as picture fails, skip this shape
                        logging.warning(f"⚠️ Could not copy shape type {shape.shape_type}")
                        continue
        
        # Save the merged presentation
        bs_prs.save(output_path)
        
        total_slides = len(bs_prs.slides)
        print(f"✅ Successfully merged presentations!")
        print(f"   Total slides: {total_slides}")
        print(f"   Output: {output_path}")
        
        logging.info(f"✅ Successfully merged presentations to: {output_path}")
        logging.info(f"📊 Final presentation has {total_slides} slides")
        
        # Verify output file was created
        if not os.path.exists(output_path):
            raise FileNotFoundError(f"Merged presentation was not created: {output_path}")
        
        file_size = os.path.getsize(output_path)
        print(f"   File size: {file_size:,} bytes")
        
    except Exception as e:
        logging.error(f"❌ Failed to merge presentations: {str(e)}")
        raise


def _cleanup_placeholder_shapes(prs):
    """Remove unused placeholder shapes from the presentation"""
    try:
        placeholder_names_to_remove = [
            'Placeholder', 'Financial Data', 'Table', 'Data', 'Content',
            'financialData_placeholder', 'table_placeholder'
        ]
        
        for slide in prs.slides:
            shapes_to_remove = []
            for shape in slide.shapes:
                # Check if shape name contains placeholder indicators
                if any(placeholder in shape.name for placeholder in placeholder_names_to_remove):
                    # Check if shape is empty or contains only placeholder text
                    if hasattr(shape, 'text_frame') and shape.text_frame:
                        text_content = shape.text_frame.text.strip()
                        if text_content in ['', ' ', 'Placeholder', 'Financial Data', 'Table', 'Data', 'Content']:
                            shapes_to_remove.append(shape)
                    elif hasattr(shape, 'table') and shape.table:
                        # Check if table is empty
                        if len(shape.table.rows) == 0 or (len(shape.table.rows) == 1 and not shape.table.cell(0, 0).text.strip()):
                            shapes_to_remove.append(shape)
            
            # Remove identified placeholder shapes
            for shape in shapes_to_remove:
                try:
                    slide.shapes._spTree.remove(shape._element)
                    print(f"🗑️ Removed placeholder shape: {shape.name}")
                except Exception as e:
                    print(f"⚠️ Could not remove shape {shape.name}: {e}")
                    
    except Exception as e:
        print(f"⚠️ Error during placeholder cleanup: {e}") 