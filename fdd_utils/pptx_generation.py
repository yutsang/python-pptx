#!/usr/bin/env python3
"""
PowerPoint Generation Module for Financial Reports
Based on the backup methods but implemented fresh for the new system
"""

import os
import re
import logging
from typing import Dict, List, Optional, Tuple
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def get_tab_name(project_name: str) -> Optional[str]:
    """Get tab name from project name for Excel embedding"""
    if not project_name:
        return None

    # Extract key identifier from project name
    # e.g., "东莞联洋利润表" -> "东莞联洋"
    words = project_name.split()
    if words:
        return words[0]
    return None


def clean_content_quotes(content: str) -> str:
    """Clean and format content quotes"""
    if not content:
        return ""

    # Remove excessive quotes and clean formatting
    content = re.sub(r'^"*|"*$', '', content.strip())
    content = re.sub(r'""+', '"', content)

    return content


def detect_chinese_text(text: str, force_chinese_mode: bool = False) -> bool:
    """Detect if text contains significant Chinese characters"""
    if force_chinese_mode:
        return True

    if not text:
        return False

    chinese_chars = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
    total_chars = len(text)

    if total_chars == 0:
        return False

    # Consider it Chinese if > 30% Chinese characters
    return (chinese_chars / total_chars) > 0.3


def get_font_size_for_text(text: str, base_size: int = 9, force_chinese_mode: bool = False) -> Pt:
    """Get appropriate font size for text content"""
    is_chinese = detect_chinese_text(text, force_chinese_mode)

    if is_chinese:
        # Slightly larger for Chinese readability
        return Pt(base_size + 1)
    else:
        return Pt(base_size)


def get_font_name_for_text(text: str, default_font: str = 'Arial') -> str:
    """Get appropriate font name for text"""
    is_chinese = detect_chinese_text(text)

    if is_chinese:
        return 'Microsoft YaHei'  # or 'SimSun' as fallback
    else:
        return default_font


def get_line_spacing_for_text(text: str, force_chinese_mode: bool = False) -> float:
    """Get line spacing for text"""
    is_chinese = detect_chinese_text(text, force_chinese_mode)
    return 0.9 if is_chinese else 1.0


def get_space_after_for_text(text: str, force_chinese_mode: bool = False) -> Pt:
    """Get space after for text paragraphs"""
    is_chinese = detect_chinese_text(text, force_chinese_mode)
    return Pt(6) if is_chinese else Pt(4)


def get_space_before_for_text(text: str, force_chinese_mode: bool = False) -> Pt:
    """Get space before for text paragraphs"""
    is_chinese = detect_chinese_text(text, force_chinese_mode)
    return Pt(3) if is_chinese else Pt(2)


def replace_entity_placeholders(content: str, project_name: str) -> str:
    """Replace entity placeholders in content"""
    if not content or not project_name:
        return content

    # Replace common placeholders
    replacements = {
        '[PROJECT]': project_name,
        '[Entity]': project_name,
        '[Company]': project_name,
    }

    for placeholder, replacement in replacements.items():
        content = content.replace(placeholder, replacement)

    return content


class PowerPointGenerator:
    """Main PowerPoint generator class"""

    def __init__(self, template_path: str, language: str = 'english', row_limit: int = 20):
        self.template_path = template_path
        self.language = language.lower()
        self.row_limit = row_limit
        self.presentation = None

    def load_template(self):
        """Load the PowerPoint template"""
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template not found: {self.template_path}")

        self.presentation = Presentation(self.template_path)
        logger.info(f"Loaded template: {self.template_path}")

    def find_shape_by_name(self, shapes, name: str):
        """Find shape by name in slide (case-insensitive), recursive"""
        name_lower = name.lower()
        for shape in shapes:
            if hasattr(shape, 'name') and (shape.name == name or shape.name.lower() == name_lower):
                return shape
            
            # Check for group
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                found = self.find_shape_by_name(shape.shapes, name)
                if found:
                    return found
        return None
    
    def find_content_shape(self, shapes):
        """Find content shape by trying multiple possible names"""
        # Try different possible names for content shapes
        possible_names = [
            'Content',
            'Text-commentary',
            'textMainBullets',
            'Text',
            'Commentary',
            'MainContent',
            'Body'
        ]
        
        for name in possible_names:
            shape = self.find_shape_by_name(shapes, name)
            if shape and shape.has_text_frame:
                return shape
        
        # If no named shape found, try to find any text frame shape that's not a title
        for shape in shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                shape_name = getattr(shape, 'name', '')
                # Skip title shapes and other non-content shapes
                if shape_name and 'title' not in shape_name.lower() and 'proj' not in shape_name.lower():
                    return shape
        
        return None

    def replace_text_preserve_formatting(self, shape, replacements: Dict[str, str]):
        """Replace text while preserving formatting"""
        if not shape.has_text_frame:
            return

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                for old_text, new_text in replacements.items():
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)

    def update_project_titles(self, project_name: str, statement_type: str = 'BS'):
        """Update project titles in presentation"""
        if not self.presentation:
            return

        # Extract first two words for professional display
        if project_name:
            words = project_name.split()
            display_entity = ' '.join(words[:2]) if len(words) >= 2 else words[0] if words else project_name
        else:
            display_entity = project_name

        # Define title templates based on language and statement type
        if statement_type.upper() == 'BS':
            if self.language == 'chinese':
                title_template = f"资产负债表概览 - {display_entity}"
            else:
                title_template = f"Entity Overview - {display_entity}"
        elif statement_type.upper() == 'IS':
            if self.language == 'chinese':
                title_template = f"利润表概览 - {display_entity}"
            else:
                title_template = f"Income Statement - {display_entity}"
        else:
            if self.language == 'chinese':
                title_template = f"财务报表概览 - {display_entity}"
            else:
                title_template = f"Financial Report - {display_entity}"

        # Update titles in all slides
        for slide_index, slide in enumerate(self.presentation.slides):
            current_slide_number = slide_index + 1
            proj_title_shape = self.find_shape_by_name(slide.shapes, "projTitle")

            if proj_title_shape:
                current_text = proj_title_shape.text
                if "[PROJECT]" in current_text:
                    replacements = {
                        "[PROJECT]": display_entity,
                        "[Current]": str(current_slide_number),
                        "[Total]": str(len(self.presentation.slides))
                    }
                    self.replace_text_preserve_formatting(proj_title_shape, replacements)
                else:
                    # Replace the entire title
                    if proj_title_shape.has_text_frame:
                        proj_title_shape.text_frame.text = title_template

    def generate_full_report(self, markdown_content: str, summary_md: Optional[str] = None,
                           output_path: str = None):
        """Generate full PowerPoint report from markdown content"""
        if not self.presentation:
            self.load_template()

        # Process markdown content
        processed_content = self._process_markdown_content(markdown_content)

        # Apply content to presentation
        self._apply_content_to_presentation(processed_content)

        # Save if output path provided
        if output_path:
            self.save(output_path)

    def _process_markdown_content(self, content: str) -> Dict:
        """Process markdown content into structured data"""
        if not content:
            logger.warning("Empty content provided to _process_markdown_content")
            return {}

        logger.info(f"Processing markdown content, length: {len(content)}")
        logger.debug(f"Content preview (first 500 chars): {content[:500]}")

        # Split by headers (## Account Name)
        sections = re.split(r'^##\s+(.+)$', content, flags=re.MULTILINE)

        logger.info(f"Found {len(sections)} sections after splitting")

        processed_sections = {}

        # Process each section
        for i in range(1, len(sections), 2):
            if i + 1 < len(sections):
                account_name = sections[i].strip()
                account_content = sections[i + 1].strip()

                logger.info(f"Processing section: {account_name}, content length: {len(account_content)}")

                processed_sections[account_name] = {
                    'content': account_content,
                    'is_chinese': detect_chinese_text(account_content)
                }

        logger.info(f"Processed {len(processed_sections)} sections")
        return processed_sections

    def _apply_content_to_presentation(self, sections: Dict):
        """Apply processed content to presentation slides"""
        if not self.presentation:
            logger.warning("No presentation loaded")
            return

        logger.info(f"Applying {len(sections)} sections to presentation with {len(self.presentation.slides)} slides")

        # Find content placeholders and fill them
        slide_idx = 0
        for slide in self.presentation.slides:
            if slide_idx >= len(sections):
                logger.warning(f"More slides ({len(self.presentation.slides)}) than sections ({len(sections)})")
                break

            account_name = list(sections.keys())[slide_idx]
            section_data = sections[account_name]

            logger.info(f"Processing slide {slide_idx + 1} for account: {account_name}")

            # Find content shape using flexible name matching
            content_shape = self.find_content_shape(slide.shapes)
            if content_shape:
                logger.info(f"Found content shape '{content_shape.name}' on slide {slide_idx + 1}")
                if content_shape.has_text_frame:
                    # Apply content to shape
                    self._fill_content_shape(content_shape, section_data)
                    logger.info(f"Applied content to slide {slide_idx + 1}")
                else:
                    logger.warning(f"Content shape found but has no text_frame on slide {slide_idx + 1}")
            else:
                logger.warning(f"No content shape found on slide {slide_idx + 1}, available shapes: {[s.name if hasattr(s, 'name') else 'unnamed' for s in slide.shapes]}")
                # Try to use the first available text frame as fallback
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        shape_name = getattr(shape, 'name', 'unnamed')
                        if 'title' not in shape_name.lower() and 'proj' not in shape_name.lower():
                            logger.info(f"Using fallback shape '{shape_name}' on slide {slide_idx + 1}")
                            self._fill_content_shape(shape, section_data)
                            break

            slide_idx += 1

    def _fill_content_shape(self, shape, section_data: Dict):
        """Fill content shape with processed data"""
        if not shape.has_text_frame:
            logger.warning("Shape does not have text_frame")
            return

        content = section_data.get('content', '')
        is_chinese = section_data.get('is_chinese', False)

        logger.info(f"Filling shape with content length: {len(content)}")

        # Clear existing content
        shape.text_frame.clear()
        
        if not content or not content.strip():
            logger.warning("No content to fill")
            return
        
        # Split content into paragraphs if it contains newlines
        content_lines = content.split('\n')
        
        # Add content with proper formatting
        for idx, line in enumerate(content_lines):
            line = line.strip()
            if not line and idx > 0:
                # Skip empty lines except add a paragraph break
                continue
            
            if idx == 0:
                # Use first paragraph or create one
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
            else:
                p = shape.text_frame.add_paragraph()
            
            p.text = line
            
            # Apply formatting to runs
            for run in p.runs:
                run.font.size = get_font_size_for_text(line, force_chinese_mode=is_chinese)
                run.font.name = get_font_name_for_text(line)

            # Set paragraph formatting
            p.space_after = get_space_after_for_text(line, force_chinese_mode=is_chinese)
            p.space_before = get_space_before_for_text(line, force_chinese_mode=is_chinese)
            p.line_spacing = get_line_spacing_for_text(line, force_chinese_mode=is_chinese)
        
        logger.info(f"Successfully filled shape with {len([l for l in content_lines if l.strip()])} paragraphs")

    def _calculate_max_lines_for_textbox(self, shape):
        """Calculate maximum lines that can fit in textbox"""
        if not shape or not hasattr(shape, 'height'):
            return 35  # Default fallback
        
        # Get shape height in EMU (English Metric Units)
        shape_height_emu = shape.height
        
        # Convert EMU to points (1 EMU = 1/914400 inches, 1 inch = 72 points)
        shape_height_pt = shape_height_emu * 72 / 914400
        
        # Account for margins and padding - use safer space (85% utilization)
        effective_height_pt = shape_height_pt * 0.85
        
        # Calculate line height based on font size (9pt) and line spacing (1.0)
        font_size_pt = 9
        line_spacing = 1.0
        line_height_pt = font_size_pt * line_spacing
        
        # Calculate maximum rows that can fit
        max_rows = int(effective_height_pt / line_height_pt)
        
        return max(20, max_rows)  # At least 20 lines (reduced from 25 to prevent overflow)
    
    def _calculate_content_lines(self, category: str, mapping_key: str, commentary: str) -> int:
        """Calculate how many lines a piece of content will take"""
        lines = 0
        
        # Category line (if exists)
        if category:
            lines += 1
        
        # Key line (grey char + key name + dash + first line of commentary)
        lines += 1
        
        # Commentary lines (estimate based on text length)
        # Average 50 chars per line for 9pt Arial
        commentary_lines = commentary.split('\n')
        for line in commentary_lines:
            line = line.strip()
            if line:
                # Estimate lines needed for this text (50 chars per line)
                chars_per_line = 50
                line_count = max(1, (len(line) + chars_per_line - 1) // chars_per_line)
                lines += line_count
        
        return lines
    
    def _distribute_content_across_slides(self, structured_data: List[Dict], max_slides: int = 4):
        """Distribute content across slides based on textbox capacity with page break logic"""
        if not structured_data:
            return []
        
        # Find a textbox shape to calculate capacity
        sample_shape = None
        for slide in self.presentation.slides:
            shape = self.find_shape_by_name(slide.shapes, "textMainBullets")
            if shape:
                sample_shape = shape
                break
        
        if not sample_shape:
            # Fallback: try alternative names
            for slide in self.presentation.slides:
                for alt_name in ["textMainBullets_L", "textMainBullets_R", "Content Placeholder 2"]:
                    shape = self.find_shape_by_name(slide.shapes, alt_name)
                    if shape:
                        sample_shape = shape
                        break
                if sample_shape:
                    break
        
        max_lines_per_slide = self._calculate_max_lines_for_textbox(sample_shape) if sample_shape else 35
        
        # Distribute content with page break support
        distribution = []  # List of (slide_idx, [account_data_list], needs_continuation)
        current_slide_idx = 0
        current_slide_content = []
        current_slide_lines = 0
        previous_category = None
        
        for account_data in structured_data:
            category = account_data.get('category', '')
            mapping_key = account_data.get('mapping_key', account_data.get('account_name', ''))
            commentary = account_data.get('commentary', '')
            
            # Category header line (only if category changes)
            category_lines = 1 if (category and category != previous_category) else 0
            previous_category = category
            
            # Calculate content lines (key + commentary, category is handled separately)
            content_lines = self._calculate_content_lines('', mapping_key, commentary)  # Don't count category here
            
            # Check if this fits on current slide
            if current_slide_lines + category_lines + content_lines <= max_lines_per_slide and current_slide_idx < max_slides:
                # Add to current slide
                current_slide_content.append(account_data)
                current_slide_lines += category_lines + content_lines
            else:
                # Save current slide if it has content
                if current_slide_content:
                    # Mark if we're at max and have more content
                    needs_continuation = (current_slide_idx < max_slides - 1)
                    distribution.append((current_slide_idx, current_slide_content, needs_continuation))
                
                # Start new slide
                if current_slide_idx < max_slides - 1:
                    current_slide_idx += 1
                    current_slide_content = [account_data]
                    current_slide_lines = category_lines + content_lines
                else:
                    # Max slides reached, add continuation marker to last slide
                    if current_slide_content:
                        distribution.append((current_slide_idx, current_slide_content, True))
                    break
        
        # Add last slide if it has content
        if current_slide_content:
            distribution.append((current_slide_idx, current_slide_content, False))
        
        return distribution
    
    def apply_structured_data_to_slides(self, structured_data: List[Dict], start_slide: int, 
                                       project_name: str, statement_type: str, is_chinese_databook: bool = False):
        """Apply structured data directly to slides (slides 1-4 for BS, 5-8 for IS)"""
        if not self.presentation:
            self.load_template()
        
        logger.info(f"Applying {len(structured_data)} accounts to slides starting at {start_slide}")
        
        # Distribute content across slides based on textbox capacity (max 4 slides)
        distribution = self._distribute_content_across_slides(structured_data, max_slides=4)
        
        # Ensure we have enough slides
        max_slide_idx = max((slide_idx for slide_idx, _, _ in distribution), default=0)
        needed_slides = start_slide + max_slide_idx
        current_slide_count = len(self.presentation.slides)
        
        if needed_slides > current_slide_count:
            # Add slides if needed (use the same layout as existing slides)
            if current_slide_count > 0:
                slide_layout = self.presentation.slides[0].slide_layout
                for _ in range(needed_slides - current_slide_count):
                    self.presentation.slides.add_slide(slide_layout)
        
        # Track which slides are used
        used_slide_indices = set()
        
        # Apply content to slides
        for slide_idx, account_data_list, needs_continuation in distribution:
            actual_slide_idx = start_slide - 1 + slide_idx  # Convert to 0-based
            if actual_slide_idx >= len(self.presentation.slides):
                logger.warning(f"Slide index {actual_slide_idx + 1} exceeds available slides")
                continue
            
            used_slide_indices.add(actual_slide_idx)
            slide = self.presentation.slides[actual_slide_idx]
            
            # Note: Financial tables are filled by embed_financial_tables() which is called after
            # applying all data. This ensures the full BS/IS tables are embedded, not just individual account data.
            # So we skip filling tables here to avoid conflicts.
            
            # Update projTitle
            proj_title_shape = self.find_shape_by_name(slide.shapes, "projTitle")
            if proj_title_shape and proj_title_shape.has_text_frame:
                if project_name:
                    proj_title_shape.text_frame.text = project_name
            
            # Fill textMainBullets with all accounts for this slide
            bullets_shape = self.find_shape_by_name(slide.shapes, "textMainBullets")
            if not bullets_shape:
                for alt_name in ["textMainBullets_L", "textMainBullets_R", "Content Placeholder 2"]:
                    bullets_shape = self.find_shape_by_name(slide.shapes, alt_name)
                    if bullets_shape:
                        break
                if not bullets_shape:
                    for shape in slide.shapes:
                        if hasattr(shape, 'text_frame'):
                            bullets_shape = shape
                            break
            
            if bullets_shape and bullets_shape.has_text_frame:
                tf = bullets_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                from pptx.enum.text import MSO_VERTICAL_ANCHOR
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                
                # Fill with all accounts for this slide, grouped by category
                # Show category header only once per category group
                current_category = None
                for account_idx, account_data in enumerate(account_data_list):
                    category = account_data.get('category', '')
                    mapping_key = account_data.get('mapping_key', account_data.get('account_name', ''))
                    display_name = account_data.get('display_name', mapping_key)  # Use proper name from financial statement
                    commentary = account_data.get('commentary', '')
                    is_chinese = account_data.get('is_chinese', False)
                    
                    # Show category header only when category changes
                    if category and category != current_category:
                        # Add category header - use Chinese if databook is Chinese
                        p_category = tf.add_paragraph()
                        p_category.level = 0
                        try:
                            p_category.left_indent = Inches(0.21)
                            p_category.first_line_indent = Inches(-0.19)
                            p_category.space_before = Pt(6) if current_category else Pt(0)  # Space before if not first
                            p_category.space_after = Pt(0)
                            p_category.line_spacing = 1.0
                        except:
                            pass
                        
                        run_category = p_category.add_run()
                        # Use Chinese category name if databook is Chinese
                        category_text = category
                        if is_chinese_databook or is_chinese:
                            # Translate category to Chinese - comprehensive list
                            category_translations = {
                                # Balance Sheet - Assets
                                'Current assets': '流动资产',
                                'Current Assets': '流动资产',
                                'Non-current assets': '非流动资产',
                                'Non-Current Assets': '非流动资产',
                                'Non current assets': '非流动资产',
                                'Assets': '资产',
                                # Balance Sheet - Liabilities
                                'Current liabilities': '流动负债',
                                'Current Liabilities': '流动负债',
                                'Non-current liabilities': '非流动负债',
                                'Non-Current Liabilities': '非流动负债',
                                'Non current liabilities': '非流动负债',
                                'Liabilities': '负债',
                                # Balance Sheet - Equity
                                'Equity': '所有者权益',
                                "Owner's equity": '所有者权益',
                                "Owners' equity": '所有者权益',
                                'Shareholders equity': '股东权益',
                                "Shareholders' equity": '股东权益',
                                # Income Statement - Revenue
                                'Revenue': '营业收入',
                                'Sales': '销售收入',
                                'Income': '收入',
                                'Operating revenue': '营业收入',
                                'Operating Revenue': '营业收入',
                                # Income Statement - Costs
                                'Cost of sales': '营业成本',
                                'Cost of Sales': '营业成本',
                                'Cost of goods sold': '销售成本',
                                'COGS': '销售成本',
                                # Income Statement - Expenses
                                'Operating expenses': '营业费用',
                                'Operating Expenses': '营业费用',
                                'Selling expenses': '销售费用',
                                'Administrative expenses': '管理费用',
                                'General and administrative': '管理费用',
                                'G&A': '管理费用',
                                # Income Statement - Other
                                'Other income': '其他收入',
                                'Other Income': '其他收入',
                                'Other expenses': '其他费用',
                                'Other Expenses': '其他费用',
                                'Finance costs': '财务费用',
                                'Finance Costs': '财务费用',
                                'Financial expenses': '财务费用',
                                'Interest expense': '利息费用',
                                'Tax': '税费',
                                'Income tax': '所得税',
                                'Taxes': '税费',
                                'Tax expense': '所得税费用',
                                # Profit items
                                'Gross profit': '毛利',
                                'Operating profit': '营业利润',
                                'Net profit': '净利润',
                                'Profit before tax': '利润总额',
                            }
                            # Try direct match first, then case-insensitive match
                            category_text = category_translations.get(category)
                            if category_text is None:
                                # Try case-insensitive match
                                category_lower = category.lower()
                                for eng_cat, chi_cat in category_translations.items():
                                    if eng_cat.lower() == category_lower:
                                        category_text = chi_cat
                                        break
                                else:
                                    category_text = category  # Keep original if no match
                        
                        # Add "(continued)" or "(续)" if this slide needs continuation
                        # BUT NOT on the first category of the first slide (account_idx == 0 and current_category is None)
                        if needs_continuation and account_idx == 0 and current_category is not None:  # Not first category
                            cont_text = " (续)" if (is_chinese or is_chinese_databook) else " (continued)"
                            category_text += cont_text
                        
                        run_category.text = category_text
                        run_category.font.size = Pt(9)
                        run_category.font.name = 'Arial'
                        run_category.font.bold = False
                        try:
                            from pptx.dml.color import RGBColor
                            run_category.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                        except:
                            pass
                        
                        current_category = category
                    
                    # Fill commentary with key formatting (no category, already shown)
                    # Use display_name (from financial statement) instead of mapping_key
                    # Check if this is the last account and we need continuation
                    is_last_account = (account_idx == len(account_data_list) - 1)
                    needs_cont = needs_continuation and is_last_account
                    
                    self._fill_text_main_bullets_with_category_and_key(
                        tf, None, display_name, commentary, is_chinese, 
                        is_chinese_databook=is_chinese_databook, needs_continuation=needs_cont
                    )
                
                # Generate AI summary for this slide from all commentary
                all_commentary = []
                for account_data in account_data_list:
                    commentary = account_data.get('commentary', '')
                    if commentary:
                        all_commentary.append(commentary)
                
                # Fill coSummaryShape with AI-generated summary
                summary_shape = self.find_shape_by_name(slide.shapes, "coSummaryShape")
                if summary_shape and summary_shape.has_text_frame:
                    summary_shape.text_frame.clear()
                    if all_commentary:
                        # Combine all commentary for this page
                        page_commentary = '\n\n'.join(all_commentary)
                        # Generate summary using AI
                        ai_summary = self._generate_ai_summary(page_commentary, is_chinese_databook)
                        if ai_summary:
                            p = summary_shape.text_frame.paragraphs[0] if summary_shape.text_frame.paragraphs else summary_shape.text_frame.add_paragraph()
                            p.text = ai_summary
                        else:
                            # Fallback to simple summary using _generate_page_summary which is smarter
                            # Combine all summaries first
                            raw_combined = ' '.join([acc.get('summary', '') for acc in account_data_list if acc.get('summary')])
                            # Use the smarter generation function instead of hard truncation
                            combined_summary = self._generate_page_summary(raw_combined, is_chinese_databook)
                            
                            if combined_summary:
                                p = summary_shape.text_frame.paragraphs[0] if summary_shape.text_frame.paragraphs else summary_shape.text_frame.add_paragraph()
                                p.text = combined_summary
                        is_chinese = account_data_list[0].get('is_chinese', False) if account_data_list else False
                        for run in p.runs:
                            run.font.size = get_font_size_for_text(combined_summary, force_chinese_mode=is_chinese)
                            run.font.name = get_font_name_for_text(combined_summary)
                
                logger.info(f"Filled slide {actual_slide_idx + 1} with {len(account_data_list)} accounts")
        
        # Note: Unused slides will be removed at the end, after all content and tables are embedded
        # Store unused slides for later removal
        statement_slide_range = list(range(start_slide - 1, min(start_slide + 3, len(self.presentation.slides))))
        unused_slides = [idx for idx in statement_slide_range if idx not in used_slide_indices]
        if unused_slides:
            # Store for later removal - don't remove now
            if not hasattr(self, '_unused_slides_to_remove'):
                self._unused_slides_to_remove = []
            self._unused_slides_to_remove.extend(unused_slides)
            logger.info(f"Marked {len(unused_slides)} unused slides for {statement_type} for later removal: {[idx + 1 for idx in unused_slides]}")
    
    def _remove_slides(self, slide_indices):
        """Remove slides by indices (from backup method)"""
        # Sort in reverse order to maintain indices while removing
        for slide_idx in sorted(slide_indices, reverse=True):
            if slide_idx < len(self.presentation.slides):
                try:
                    # Use XML-based removal (from backup method)
                    xml_slides = self.presentation.slides._sldIdLst
                    slides = list(xml_slides)
                    
                    if slide_idx < len(slides):
                        # Get the slide element to remove
                        slide_element = slides[slide_idx]
                        # Remove relationship
                        rId = slide_element.rId
                        self.presentation.part.drop_rel(rId)
                        # Remove from XML
                        xml_slides.remove(slide_element)
                        logger.info(f"Removed slide {slide_idx + 1}")
                    else:
                        logger.warning(f"Slide index {slide_idx} out of range (only {len(slides)} slides)")
                except Exception as e:
                    logger.warning(f"Could not remove slide {slide_idx + 1}: {e}")
                    import traceback
                    logger.debug(traceback.format_exc())
    
    def _set_cell_border(self, cell, border_position='top', color_rgb=None, width=Pt(1)):
        """Set cell border"""
        from pptx.oxml.xmlchemy import OxmlElement
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        # Map position to tag name
        tag_map = {'top': 'lnT', 'bottom': 'lnB', 'left': 'lnL', 'right': 'lnR'}
        tag_name = tag_map.get(border_position)
        if not tag_name:
            return
            
        # Check if line element exists
        ln = tcPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag_name}")
        if ln is None:
            ln = OxmlElement(f"a:{tag_name}")
            tcPr.append(ln)
            
        # Set properties
        ln.set('w', str(int(width)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')
        
        # Set color
        if color_rgb:
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            # Convert RGBColor or tuple to hex string
            hex_color = "000000"
            if isinstance(color_rgb, str):
                hex_color = color_rgb.replace('#', '')
            elif isinstance(color_rgb, tuple) and len(color_rgb) == 3:
                hex_color = f"{color_rgb[0]:02x}{color_rgb[1]:02x}{color_rgb[2]:02x}"
            # If it's an RGBColor object, user should pass str or tuple for this low-level func
                
            srgbClr.set('val', hex_color)
            solidFill.append(srgbClr)
            ln.append(solidFill)
            
            prstDash = OxmlElement('a:prstDash')
            prstDash.set('val', 'solid')
            ln.append(prstDash)
            
            round_ = OxmlElement('a:round')
            ln.append(round_)
            
            headEnd = OxmlElement('a:headEnd')
            headEnd.set('type', 'none')
            headEnd.set('w', 'med')
            headEnd.set('len', 'med')
            ln.append(headEnd)
            
            tailEnd = OxmlElement('a:tailEnd')
            tailEnd.set('type', 'none')
            tailEnd.set('w', 'med')
            tailEnd.set('len', 'med')
            ln.append(tailEnd)

    def _fill_table_placeholder(self, shape, df, table_name: str = None, currency_unit: str = None):
        """Fill table placeholder with DataFrame data, preserving original formatting
        Args:
            shape: Table shape or placeholder
            df: DataFrame with data
            table_name: Name of the table (e.g., "示意性调整后资产负债表 - xxxx")
            currency_unit: Currency unit (e.g., "人民币千元" or "CNY'000") to replace "Description"
        """
        try:
            # Debug: Log DataFrame content
            logger.info(f"Filling table with DF shape: {df.shape}")
            if not df.empty:
                logger.info(f"First row data: {df.iloc[0].to_dict()}")
                # Check if any data is non-zero
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    non_zero_count = (df[numeric_cols] != 0).sum().sum()
                    logger.info(f"Non-zero values in DF: {non_zero_count}")
            
            # Find parent slide
            slide = None
            for s in self.presentation.slides:
                for shp in s.shapes:
                    if shp == shape:
                        slide = s
                        break
                if slide:
                    break
            
            # Adjust position and width
            try:
                target_width = Inches(4.78)
                shape.width = target_width
                
                # Align top with Text-commentary if possible
                if slide:
                    ref_shape = self.find_shape_by_name(slide.shapes, "Text-commentary")
                    if not ref_shape:
                        ref_shape = self.find_shape_by_name(slide.shapes, "textMainBullets")
                    
                    if ref_shape:
                        shape.top = ref_shape.top
                        logger.info(f"Aligned table top to {ref_shape.name} at {shape.top}")
            except Exception as e:
                logger.warning(f"Could not adjust table position/width: {e}")

            # Check if shape is a TablePlaceholder (textbox placeholder)
            from pptx.shapes.placeholder import TablePlaceholder
            
            table = None
            # Check if it's a TablePlaceholder (textbox placeholder named "Table Placeholder")
            is_table_placeholder = False
            try:
                is_table_placeholder = isinstance(shape, TablePlaceholder)
            except:
                # Check by name
                if hasattr(shape, 'name') and 'Table' in shape.name and 'Placeholder' in shape.name:
                    is_table_placeholder = True
            
            if is_table_placeholder:
                # It's a table placeholder - insert a table into it
                logger.info(f"Found TablePlaceholder ({shape.name if hasattr(shape, 'name') else 'unnamed'}), inserting table with {len(df)} rows and {len(df.columns)} columns")
                try:
                    # Get placeholder dimensions - Override with requested width
                    left = shape.left
                    top = shape.top
                    width = Inches(4.78) # Fixed width
                    height = shape.height
                    
                    # Find the slide containing this shape (already found above)
                    if slide:
                        # Remove the placeholder shape
                        sp = shape._element
                        slide.shapes._spTree.remove(sp)
                        
                        # Add new table at the same position
                        # Need: 1 row for title (if table_name), 1 for header, N for data
                        total_rows = len(df) + 2 if table_name else len(df) + 1
                        table_shape = slide.shapes.add_table(
                            rows=total_rows,
                            cols=len(df.columns),
                            left=left,
                            top=top,
                            width=width,
                            height=height
                        )
                        table = table_shape.table
                        logger.info(f"✅ Inserted new table: {len(table.rows)} rows, {len(table.columns)} columns")
                except Exception as e:
                    logger.error(f"Could not insert table into placeholder: {e}")
                    import traceback
                    logger.debug(traceback.format_exc())
            elif hasattr(shape, 'table'):
                # Try to access existing table
                try:
                    table = shape.table
                    logger.info(f"Found existing table with {len(table.rows)} rows and {len(table.columns)} columns")
                except ValueError:
                    # Shape doesn't contain a table
                    logger.warning(f"Shape has table attribute but doesn't contain a table")
                    table = None
            
            if table:
                # Colors
                DARK_BLUE = RGBColor(0, 51, 102)
                TIFFANY_BLUE = RGBColor(10, 186, 181)
                GREY = RGBColor(217, 217, 217)
                WHITE = RGBColor(255, 255, 255)
                BLACK = RGBColor(0, 0, 0)
                
                # Adjust column widths
                # Make first column (description) wider, distribute rest
                if len(table.columns) > 0:
                    try:
                        # Make first column (description) wider by 80% (1.8x)
                        table.columns[0].width = int(table.columns[0].width * 1.8)
                    except:
                        pass
                
                # Add table name as first row if provided
                if table_name:
                    # Insert a new row at the top for table name
                    try:
                        # Ensure table has at least one row
                        if len(table.rows) == 0:
                            table.rows.add_row()
                            
                        name_row = table.rows[0]  # Use first row for name
                        # Merge all cells in first row for table name
                        if len(table.columns) > 1:
                            name_row.cells[0].merge(name_row.cells[len(table.columns) - 1])
                        name_cell = name_row.cells[0]
                        name_cell.text = table_name
                        # Format table name: Arial 9, bold, centered, Dark Blue bg, White font
                        if name_cell.text_frame.paragraphs:
                            p = name_cell.text_frame.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER  # Center alignment
                            if p.runs:
                                run = p.runs[0]
                            else:
                                run = p.add_run()
                            run.font.name = 'Arial'
                            run.font.size = Pt(9)
                            run.font.bold = True
                            run.font.color.rgb = WHITE
                            
                            name_cell.fill.solid()
                            name_cell.fill.fore_color.rgb = DARK_BLUE
                            
                        # Shift data down - we'll use rows starting from index 1
                        data_start_row = 1
                    except:
                        data_start_row = 0
                else:
                    data_start_row = 0
                
                # Fill header row with formatting
                max_cols = min(len(df.columns), len(table.columns))
                header_row_idx = data_start_row
                
                # Ensure header row exists
                if len(table.rows) <= header_row_idx:
                    table.rows.add_row()
                    
                for col_idx, col_name in enumerate(df.columns[:max_cols]):
                    if col_idx < len(table.columns):
                        cell = table.cell(header_row_idx, col_idx)
                        # Replace "Description" with currency unit if found
                        if currency_unit and (col_name.lower() == 'description' or '描述' in str(col_name) or '项目' in str(col_name)):
                            cell.text = currency_unit
                        else:
                            cell.text = str(col_name)
                        # Apply header formatting: Arial 9, bold, Tiffany Blue bg, White font
                        if cell.text_frame.paragraphs:
                            p = cell.text_frame.paragraphs[0]
                            p.alignment = PP_ALIGN.CENTER
                            
                            if p.runs:
                                run = p.runs[0]
                            else:
                                run = p.add_run()
                            
                            run.font.name = 'Arial'
                            run.font.size = Pt(9)
                            run.font.bold = True
                            run.font.color.rgb = WHITE # White font for header
                            
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = TIFFANY_BLUE
                        
                        logger.debug(f"Filled header cell {col_idx}: {cell.text}")
                
                # Fill data rows with formatting - show ALL rows (no limit)
                # Check if table has enough rows, if not, limit to available rows
                max_rows = len(df)  # Show all rows
                rows_needed = max_rows + data_start_row + 1  # +1 for header row
                available_rows = len(table.rows)
                
                if available_rows < rows_needed:
                    logger.warning(f"Table has {available_rows} rows but needs {rows_needed}. Will only fill {available_rows - data_start_row - 1} data rows.")
                    max_rows = min(max_rows, available_rows - data_start_row - 1)
                    if max_rows < 0:
                        max_rows = 0
                
                logger.info(f"Table has {available_rows} rows available, will fill {max_rows} data rows")
                
                # Now fill all rows with Arial 9 font
                # Check for title, date, total, and subtotal rows to highlight
                logger.info(f"Filling {max_rows} data rows, starting at row index {header_row_idx + 1}, table has {len(table.rows)} rows")
                for row_idx in range(max_rows):
                    if row_idx >= len(df):
                        break
                    df_row = df.iloc[row_idx]
                    first_col_value = str(df_row.iloc[0]) if len(df_row) > 0 else ""
                    
                    # Check if this is a title, date, total, or subtotal row
                    is_special_row = False
                    is_total_row = False
                    first_col_lower = first_col_value.lower()
                    total_keywords = ['total', '合计', '总计', '小计', 'subtotal', 'sub-total', 'sub total']
                    special_keywords = total_keywords + ['title', '标题', 'date', '日期', '年', '月']
                    
                    if any(keyword in first_col_lower for keyword in special_keywords):
                        is_special_row = True
                    
                    if any(keyword in first_col_lower for keyword in total_keywords):
                        is_total_row = True
                    
                    # Data row index = header_row_idx + 1 + row_idx
                    data_row_idx = header_row_idx + 1 + row_idx
                    if data_row_idx >= len(table.rows):
                        logger.warning(f"Data row index {data_row_idx} exceeds table rows {len(table.rows)}, skipping")
                        break
                    
                    # Log first row processing
                    if row_idx == 0:
                        logger.info(f"Processing first data row: {df_row.values[:3]}")

                    # Determine if we need to divide by 1000
                    # NOTE: We already disabled multiplication in extraction, so values should be in thousands
                    # But if for some reason we still need to handle it (e.g. different source), we check here
                    # User asked to remove *1000 logic, so we assume values are correct as-is (in thousands)
                    divide_by_1000 = False 
                    # if currency_unit and ("千" in currency_unit or "000" in currency_unit):
                    #     divide_by_1000 = True
                        
                    for col_idx, col_name in enumerate(df.columns[:max_cols]):
                        if col_idx >= len(table.columns):
                            break
                        cell = table.cell(data_row_idx, col_idx)
                        
                        # Get value from DataFrame safely
                        value = df_row[col_name] if col_name in df_row.index else ""
                        
                        # Special handling for Description column (index 0)
                        if col_idx == 0:
                            # Shorten description to prevent row height expansion
                            str_val = str(value)
                            if len(str_val) > 35:
                                value = str_val[:32] + "..."
                        
                        # Format value to string
                        text_val = ""
                        try:
                            if pd.isna(value):
                                text_val = ""
                            elif isinstance(value, (int, float)):
                                if divide_by_1000 and col_idx > 0: # Don't divide description
                                    value = value / 1000
                                text_val = f"{value:,.0f}" if value != 0 else "-" # Round to 0 d.p., use dash for zero
                            else:
                                text_val = str(value).strip()
                                # Handle potential string representations of numbers
                                if col_idx > 0 and text_val.replace('.','',1).isdigit(): # Only for data columns
                                    # It's a number string, try to format it
                                    try:
                                        float_val = float(text_val)
                                        if divide_by_1000:
                                            float_val = float_val / 1000
                                        text_val = f"{float_val:,.0f}" if float_val != 0 else "-"
                                    except:
                                        pass
                        except Exception as e:
                            text_val = str(value)
                        
                        # Set text
                        cell.text = text_val
                        
                        # Log first cell value of first row
                        if row_idx == 0 and col_idx < 2:
                            logger.info(f"Setting cell ({data_row_idx}, {col_idx}) to: '{text_val}'")
                        
                        # Apply formatting: Arial 7pt (reduced from 9pt) for all cells
                        # Note: Always access paragraphs[0] AFTER setting text
                        if not cell.text_frame.paragraphs:
                            cell.text_frame.add_paragraph()
                            
                        p = cell.text_frame.paragraphs[0]
                        if not p.runs:
                            p.add_run()
                            
                        # Apply formatting to ALL runs (setting cell.text might create one run, but best to be safe)
                        for run in p.runs:
                            run.text = text_val # Ensure text is set in the run
                            run.font.name = 'Arial'
                            run.font.size = Pt(7) # Reduced to 7pt
                            
                            # Force Black color for data rows
                            try:
                                run.font.color.rgb = BLACK
                            except:
                                pass
                            
                            # Bold for special rows
                            run.font.bold = is_special_row
                        
                        # Highlight special rows
                        if is_special_row:
                            try:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = GREY
                            except:
                                pass
                        else:
                            # White background for normal rows
                            try:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = WHITE
                            except:
                                pass
                                
                        # Add bold top horizontal line for total/subtotal rows
                        if is_total_row:
                            try:
                                # Top border in Dark Blue, 2pt width. Pass hex string "003366"
                                self._set_cell_border(cell, 'top', color_rgb="003366", width=Pt(2))
                            except:
                                pass
                    
                    logger.debug(f"Filled table row {row_idx + 1} (data_row_idx: {data_row_idx}, special: {is_special_row})")
                
                logger.info(f"✅ Updated table with Excel data (formatting preserved)")
            else:
                # If no table, this is an error - table placeholder should be a table shape
                logger.error("Table Placeholder is not a table shape! Cannot embed financial table.")
                logger.error(f"Shape type: {type(shape)}, has_table: {hasattr(shape, 'table')}")
                logger.error(f"Shape name: {shape.name if hasattr(shape, 'name') else 'unnamed'}")
                # Check if shape has table attribute but it's None
                if hasattr(shape, 'table'):
                    logger.error(f"shape.table is: {shape.table}")
                # Try to create a table representation in text frame as last resort
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    # Convert DataFrame to formatted text table - show ALL rows
                    try:
                        # Show all rows, no limit
                        text_table = df.to_string(index=False)
                    except:
                        text_table = str(df)
                    
                    p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                    p.text = text_table
                    logger.warning(f"Added text table representation with all {len(df)} rows ({len(text_table)} chars) - NOT IDEAL, should be table format")
        except Exception as e:
            logger.error(f"Could not fill table placeholder: {e}")
            import traceback
            logger.error(traceback.format_exc())
            # Fallback: add text representation - show ALL rows
            if shape.has_text_frame:
                shape.text_frame.clear()
                # Show all rows, not just first 10
                text_repr = df.to_string(index=False)
                p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                p.text = text_repr
    
    def _detect_bullet_levels(self, text: str) -> List[Tuple[int, str]]:
        """
        Detect bullet levels (1-3) from commentary text
        Returns list of (level, text) tuples where level 0 = no bullet, 1-3 = bullet levels
        """
        lines = text.split('\n')
        bullet_lines = []
        
        for line in lines:
            stripped = line.strip()
            original_line = line
            
            # Detect bullet lines with '- ' prefix
            if original_line.lstrip().startswith('- '):
                # Calculate indentation level (based on spaces/tabs before the bullet)
                indent_spaces = len(original_line) - len(original_line.lstrip())
                
                # Determine bullet level based on indentation (2 spaces per level)
                level = min(3, (indent_spaces // 2) + 1)  # Cap at level 3
                
                # Clean and store bullet line
                clean_line = stripped[2:]  # Remove '- '
                
                # Special handling for level 3 bullets that contain a dash indicating sub-level
                if level == 3 and " - " in clean_line:
                    # Split at the first occurrence of " - "
                    parts = clean_line.split(" - ", 1)
                    if len(parts) > 1:
                        # Add level 3 content
                        bullet_lines.append((level, parts[0].strip()))
                        # Add continuation as level 3 (indented)
                        bullet_lines.append((level, parts[1].strip()))
                    else:
                        bullet_lines.append((level, clean_line))
                else:
                    bullet_lines.append((level, clean_line))
            elif stripped:
                # Regular content (no bullet) - level 0
                bullet_lines.append((0, stripped))
        
        return bullet_lines
    
    def _fill_text_main_bullets_with_category_and_key(self, text_frame, category: str, display_name: str, 
                                                      commentary: str, is_chinese: bool, is_chinese_databook: bool = False,
                                                      needs_continuation: bool = False):
        """
        Fill textMainBullets shape with commentary formatted as:
        - Category as first level (dark blue Arial 9) - only if category is provided
        - Key name with filled round bullet + space + key name (black bold Arial 9) + "-" (not bold) + plain text
        - Indentation 0.15" with special hanging 0.15", spacing after 6pt
        """
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Add category as first level (if category exists and is not None)
        # Note: category is now handled at slide level, so this is only for individual calls
        if category:
            p_category = text_frame.add_paragraph()
            p_category.level = 0
            try:
                p_category.left_indent = Inches(0.21)
                p_category.first_line_indent = Inches(-0.19)
                p_category.space_before = Pt(0)
                p_category.space_after = Pt(0)
                p_category.line_spacing = 1.0
            except:
                pass
            
            run_category = p_category.add_run()
            run_category.text = category
            run_category.font.size = Pt(9)  # Arial 9
            run_category.font.name = 'Arial'
            run_category.font.bold = False
            try:
                run_category.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            except:
                pass
        
        # Add key name with grey char + space + key name (black bold) + "-" (not bold) + plain text
        p_key = text_frame.add_paragraph()
        p_key.level = 0  # No bullet level, we'll use grey character
        try:
            # Set formatting
            p_key.left_indent = Inches(0.15)  # 0.15" indent
            p_key.first_line_indent = Inches(-0.15)  # 0.15" special hanging
            p_key.space_before = Pt(0)
            p_key.space_after = Pt(6)  # 6pt spacing after
            p_key.line_spacing = 1.0
        except Exception as e:
            logger.warning(f"Could not set paragraph formatting: {e}")
            pass
        
        # Grey char (U+25A0) + space
        run_bullet = p_key.add_run()
        run_bullet.text = '\u25A0 '  # U+25A0 (black square) + space
        run_bullet.font.size = Pt(9)
        run_bullet.font.name = 'Arial'
        run_bullet.font.bold = False
        try:
            run_bullet.font.color.rgb = RGBColor(128, 128, 128)  # Grey
        except:
            pass
        
        # Key name (black bold Arial 9)
        run_key = p_key.add_run()
        run_key.text = display_name
        run_key.font.size = Pt(9)
        run_key.font.name = 'Arial'
        run_key.font.bold = True
        try:
            run_key.font.color.rgb = RGBColor(0, 0, 0)  # Black
        except:
            pass
        
        # "-" (not bold)
        run_dash = p_key.add_run()
        run_dash.text = " - "
        run_dash.font.size = Pt(9)
        run_dash.font.name = 'Arial'
        run_dash.font.bold = False
        try:
            run_dash.font.color.rgb = RGBColor(0, 0, 0)  # Black
        except:
            pass
        
        # Plain text (commentary content)
        # Split commentary into paragraphs and add each as continuation
        commentary_lines = commentary.split('\n')
        first_line_added = False
        for line_idx, line in enumerate(commentary_lines):
            line = line.strip()
            if not line:
                continue
            
            if not first_line_added:
                # First line continues on same paragraph after the dash
                run_text = p_key.add_run()
                run_text.text = line
                first_line_added = True
            else:
                # Subsequent lines as new paragraphs (indented continuation)
                p_text = text_frame.add_paragraph()
                p_text.level = 0  # No bullet for continuation
                try:
                    p_text.left_indent = Inches(0.15)  # 0.15" indent (same as key text)
                    p_text.first_line_indent = Inches(0)  # No hanging for continuation lines
                    p_text.space_before = Pt(0)
                    p_text.space_after = Pt(6)  # 6pt spacing after
                    p_text.line_spacing = 1.0
                except:
                    pass
                run_text = p_text.add_run()
                run_text.text = line
            
            # Apply formatting to the run
            run_text.font.size = Pt(9)
            run_text.font.name = 'Arial'
            run_text.font.bold = False
            try:
                run_text.font.color.rgb = RGBColor(0, 0, 0)  # Black
            except:
                pass
        
        # Note: "(continued)" is now added to category header, not here
    
    def _fill_text_main_bullets_with_levels(self, text_frame, commentary: str, is_chinese: bool):
        """
        Fill textMainBullets shape with commentary using detailed line break logic
        and level 1-3 text handling with page breaks (legacy method, kept for compatibility)
        """
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Detect bullet levels
        bullet_lines = self._detect_bullet_levels(commentary)
        
        # Calculate max lines that can fit in the shape
        # Estimate based on shape height (conservative estimate)
        max_lines = 35  # Default conservative estimate
        
        lines_added = 0
        
        for level, text in bullet_lines:
            if not text.strip():
                continue
            
            # Check if we need a page break (if shape is getting full)
            # Note: Actual page breaks would require creating new slides, which is handled
            # at a higher level. Here we just ensure content fits.
            if lines_added >= max_lines:
                # Add continuation indicator
                p = text_frame.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = "... (continued on next page)" if not is_chinese else "... (续下页)"
                run.font.size = get_font_size_for_text(run.text, force_chinese_mode=is_chinese)
                run.font.name = get_font_name_for_text(run.text)
                run.font.italic = True
                break
            
            # Create paragraph with appropriate level
            p = text_frame.add_paragraph()
            p.level = level  # Set bullet level (0-3)
            
            # Apply paragraph formatting based on level
            try:
                # Level 0 (no bullet) or Level 1 (main bullet)
                if level == 0 or level == 1:
                    p.left_indent = Inches(0.21)  # 0.21" indent before text
                    p.first_line_indent = Inches(-0.19)  # 0.19" special hanging
                    p.space_before = Pt(0)  # 0pt spacing before
                    p.space_after = Pt(0)  # 0pt spacing after
                    p.line_spacing = 1.0  # Single line spacing
                elif level == 2:
                    # Level 2 - more indented
                    p.left_indent = Inches(0.4)
                    p.first_line_indent = Inches(-0.19)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    p.line_spacing = 1.0
                elif level == 3:
                    # Level 3 - most indented
                    p.left_indent = Inches(0.6)
                    p.first_line_indent = Inches(-0.19)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    p.line_spacing = 1.0
            except:
                pass  # Silently handle formatting errors
            
            # Add text with proper formatting
            run = p.add_run()
            run.text = text
            run.font.size = get_font_size_for_text(text, force_chinese_mode=is_chinese)
            run.font.name = get_font_name_for_text(text)
            
            # Apply level-specific formatting
            if level == 1:
                run.font.bold = True
                try:
                    run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue for level 1
                except:
                    pass
            elif level == 0:
                # Regular text - no special formatting
                pass
            
            lines_added += 1
    
    def _generate_ai_summary(self, commentary: str, is_chinese: bool) -> str:
        """Generate AI summary from page commentary - AI controls length (~200 words target)"""
        try:
            from fdd_utils.ai_helper import AIHelper
            
            # Create summary prompt - ask AI to write ~200 words summary
            # Don't hard limit - let AI produce a coherent summary
            if is_chinese:
                prompt = f"""请为以下财务评论内容生成一个概要摘要。

目标长度：约200字（可根据内容适当增减，确保完整表达要点）

要求：
1. 涵盖所有关键财务数据和变动
2. 保持专业、简洁的语言风格
3. 不要截断句子中途
4. 确保摘要是完整、可读的段落

原始内容：
{commentary[:3000]}"""
            else:
                prompt = f"""Please generate a summary for the following financial commentary.

Target length: approximately 200 words (adjust as needed to cover key points completely)

Requirements:
1. Cover all key financial data and changes
2. Maintain professional, concise language
3. Do not truncate sentences mid-way
4. Ensure the summary is a complete, readable paragraph

Original content:
{commentary[:3000]}"""
            
            # Call AI with enough tokens to allow proper summary generation
            ai_helper = AIHelper()
            response = ai_helper.call_ai(prompt, max_tokens=500)
            
            if response and response.strip():
                summary = response.strip()
                return summary
        except Exception as e:
            logger.warning(f"Could not generate AI summary: {e}")
            import traceback
            logger.debug(traceback.format_exc())
        
        return None
    
    def _generate_page_summary(self, commentary: str, is_chinese: bool) -> str:
        """
        Generate a per-page summary from commentary text
        This is a fallback when AI summary is not available
        Extracts first few complete sentences without hard truncation
        """
        if not commentary or not commentary.strip():
            return ""
        
        # Split commentary into sentences/paragraphs
        # For Chinese, split by sentence endings
        if is_chinese or detect_chinese_text(commentary):
            # Chinese sentence endings
            sentences = []
            current_sentence = ""
            for char in commentary:
                current_sentence += char
                if char in ['。', '！', '？', '；']:
                    sentences.append(current_sentence.strip())
                    current_sentence = ""
            if current_sentence.strip():
                sentences.append(current_sentence.strip())
            
            # Take first 5-8 complete sentences as summary (~300 words equivalent)
            # Don't hard limit by characters - take complete sentences
            summary_parts = []
            word_count = 0
            for sentence in sentences[:8]:  # Up to 8 sentences
                summary_parts.append(sentence)
                word_count += len(sentence)
                if word_count >= 250:  # Target ~250-300 Chinese characters
                    break
            
            summary = ''.join(summary_parts)
        else:
            # English - split by periods
            sentences = commentary.split('.')
            summary_parts = []
            word_count = 0
            for sentence in sentences[:8]:  # Up to 8 sentences
                sentence = sentence.strip()
                if sentence:
                    summary_parts.append(sentence + '.')
                    word_count += len(sentence.split())
                    if word_count >= 100:  # Target ~100-150 words for English
                        break
            
            summary = ' '.join(summary_parts)
        
        return summary.strip()

    def embed_financial_tables(self, excel_path: str, sheet_name: str, project_name: str, language: str):
        """Embed financial tables: BS to page 1, IS to page 5"""
        try:
            import pandas as pd
            from fdd_utils.financial_extraction import extract_balance_sheet_and_income_statement
            
            logger.info(f"Embedding financial tables from {excel_path}, sheet: {sheet_name}")
            
            # Validate inputs
            if not excel_path or not sheet_name:
                logger.warning(f"Missing excel_path ({excel_path}) or sheet_name ({sheet_name}), skipping table embedding")
                return
            
            # Use the existing extraction function with DEBUG enabled
            # multiply_values=False to keep original units (thousands) for PPTX table
            bs_is_results = extract_balance_sheet_and_income_statement(excel_path, sheet_name, debug=True, multiply_values=False)
            if not bs_is_results:
                logger.warning("No BS/IS data extracted")
                return
            
            # Extract BS and IS DataFrames from results
            # The structure should have 'balance_sheet' and 'income_statement' keys with DataFrames
            bs_df = bs_is_results.get('balance_sheet')
            is_df = bs_is_results.get('income_statement')
            
            # Extract table names and currency units from databook
            # Try to find table name and currency from the sheet
            bs_table_name = None
            is_table_name = None
            currency_unit = None
            
            try:
                # Read the Excel file to find table names and currency
                excel_df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None)
                # Look for table name patterns (e.g., "示意性调整后资产负债表")
                for idx, row in excel_df.iterrows():
                    row_text = ' '.join([str(cell) for cell in row if pd.notna(cell)])
                    if '资产负债表' in row_text or 'Balance Sheet' in row_text:
                        # Extract table name - remove any entity name suffix if already present
                        bs_table_name = row_text.strip()
                        # If it already contains " - " followed by entity name, keep it as is
                        # Otherwise, it's just the table name without entity
                        if project_name and project_name not in bs_table_name:
                            # Only add entity name if it's not already there
                            bs_table_name = f"{bs_table_name} - {project_name}"
                    if '利润表' in row_text or 'Income Statement' in row_text:
                        is_table_name = row_text.strip()
                        # If it already contains " - " followed by entity name, keep it as is
                        if project_name and project_name not in is_table_name:
                            # Only add entity name if it's not already there
                            is_table_name = f"{is_table_name} - {project_name}"
                    # Look for currency unit
                    if '人民币千元' in row_text or "CNY'000" in row_text or "CNY 000" in row_text:
                        if '人民币千元' in row_text:
                            currency_unit = '人民币千元'
                        elif "CNY'000" in row_text or "CNY 000" in row_text:
                            currency_unit = "CNY'000"
            except:
                pass
            
            logger.info(f"Extracted BS: {bs_df.shape if bs_df is not None else 'None'}, IS: {is_df.shape if is_df is not None else 'None'}")
            logger.info(f"Table names - BS: {bs_table_name}, IS: {is_table_name}, Currency: {currency_unit}")
            
            # Embed BS table to slide 0 (page 1)
            if bs_df is not None and not bs_df.empty and len(self.presentation.slides) > 0:
                slide_0 = self.presentation.slides[0]
                logger.info(f"Looking for table shape on slide 1, available shapes: {[s.name if hasattr(s, 'name') else type(s).__name__ for s in slide_0.shapes]}")
                
                table_shape = self.find_shape_by_name(slide_0.shapes, "Table Placeholder")
                if not table_shape:
                    # Try alternative names
                    for name in ["Table Placeholder 2", "Table", "table", "TABLE", "Content Placeholder 2"]:
                        table_shape = self.find_shape_by_name(slide_0.shapes, name)
                        if table_shape:
                            logger.info(f"Found table shape with name: {name}")
                            break
                
                # If still not found, try to find any shape with "table" in name or has table attribute
                if not table_shape:
                    for shape in slide_0.shapes:
                        shape_name = getattr(shape, 'name', '').lower()
                        if 'table' in shape_name:
                            table_shape = shape
                            logger.info(f"Found table shape by partial name match: {shape.name}")
                            break
                        # Check if shape has table
                        try:
                            if hasattr(shape, 'table') and shape.table is not None:
                                table_shape = shape
                                logger.info(f"Found shape with table attribute: {getattr(shape, 'name', 'unnamed')}")
                                break
                            # Also check via BaseShape has_table property
                            if hasattr(shape, 'has_table') and shape.has_table:
                                table_shape = shape
                                logger.info(f"Found shape with has_table property: {getattr(shape, 'name', 'unnamed')}")
                                break
                        except:
                            pass
                
                if table_shape:
                    logger.info(f"✅ Found table shape on slide 1: {getattr(table_shape, 'name', 'unnamed')}")
                    logger.info(f"BS DataFrame shape: {bs_df.shape}, columns: {list(bs_df.columns)}")
                    # Ensure we have at least a header row in the shape
                    if hasattr(table_shape, 'table') and table_shape.table and len(table_shape.table.rows) == 0:
                         table_shape.table.rows.add_row()
                    self._fill_table_placeholder(table_shape, bs_df, table_name=bs_table_name, currency_unit=currency_unit)
                else:
                    logger.warning(f"❌ Table Placeholder not found on slide 1 for BS table")
                    logger.warning(f"Available shapes: {[s.name if hasattr(s, 'name') else type(s).__name__ for s in slide_0.shapes]}")
                    # LAST RESORT: Create a new table if none found
                    logger.info("Attempting to create new table on slide 1 as last resort")
                    try:
                        # Default position for BS table
                        left = Inches(0.5)
                        top = Inches(1.5)
                        width = Inches(12.33)
                        height = Inches(4.0)
                        # Need: 1 for title (if table_name), 1 for header, N for data
                        total_rows = len(bs_df) + 2 if bs_table_name else len(bs_df) + 1
                        table_shape = slide_0.shapes.add_table(total_rows, len(bs_df.columns), left, top, width, height)
                        self._fill_table_placeholder(table_shape, bs_df, table_name=bs_table_name, currency_unit=currency_unit)
                        logger.info("✅ Created new BS table on slide 1")
                    except Exception as e:
                        logger.error(f"Failed to create new BS table: {e}")
            else:
                logger.warning(f"BS DataFrame is None or empty, skipping BS table embedding")
            
            # Embed IS table to slide 5 (1-indexed) = slide index 4 (0-indexed)
            # IS slides are 5-8 (1-indexed) = indices 4-7 (0-indexed)
            # Just paste it to slide index 4 (page 5) - unused slides will be removed later
            if is_df is not None and not is_df.empty:
                is_slide_idx = 4  # First IS slide (page 5, 1-indexed)
                
                logger.info(f"Looking for IS table on slide {is_slide_idx + 1} (index {is_slide_idx}), total slides: {len(self.presentation.slides)}")
                
                if len(self.presentation.slides) > is_slide_idx:
                    is_slide = self.presentation.slides[is_slide_idx]
                    logger.info(f"Slide {is_slide_idx + 1} exists, searching for Table Placeholder...")
                    
                    # Try to find table shape - check all shapes for table capability
                    table_shape = None
                    # First try by name
                    for name in ["Table Placeholder", "Table Placeholder 2", "Table", "table", "TABLE"]:
                        table_shape = self.find_shape_by_name(is_slide.shapes, name)
                        if table_shape:
                            logger.info(f"Found table shape by name: {name}")
                            break
                    
                    # If not found by name, try to find any shape that has a table
                    if not table_shape:
                        for shape in is_slide.shapes:
                            if hasattr(shape, 'table') and shape.table is not None:
                                table_shape = shape
                                logger.info(f"Found table shape by table attribute: {shape.name if hasattr(shape, 'name') else 'unnamed'}")
                                break
                            # Also check if it's a table shape type
                            from pptx.shapes.base import BaseShape
                            if hasattr(shape, 'has_table') and shape.has_table:
                                table_shape = shape
                                logger.info(f"Found table shape by has_table: {shape.name if hasattr(shape, 'name') else 'unnamed'}")
                                break
                    
                    if table_shape:
                        logger.info(f"✅ Found table shape on slide {is_slide_idx + 1}, embedding IS table ({is_df.shape})")
                        if hasattr(table_shape, 'table') and table_shape.table and len(table_shape.table.rows) == 0:
                             table_shape.table.rows.add_row()
                        self._fill_table_placeholder(table_shape, is_df, table_name=is_table_name, currency_unit=currency_unit)
                    else:
                        logger.error(f"❌ Table Placeholder not found on slide {is_slide_idx + 1} for IS table")
                        # LAST RESORT: Create a new table if none found
                        logger.info(f"Attempting to create new table on slide {is_slide_idx + 1} as last resort")
                        try:
                            # Default position for IS table
                            left = Inches(0.5)
                            top = Inches(1.5)
                            width = Inches(12.33)
                            height = Inches(4.0)
                            # Need: 1 for title (if table_name), 1 for header, N for data
                            total_rows = len(is_df) + 2 if is_table_name else len(is_df) + 1
                            table_shape = is_slide.shapes.add_table(total_rows, len(is_df.columns), left, top, width, height)
                            self._fill_table_placeholder(table_shape, is_df, table_name=is_table_name, currency_unit=currency_unit)
                            logger.info(f"✅ Created new IS table on slide {is_slide_idx + 1}")
                        except Exception as e:
                            logger.error(f"Failed to create new IS table: {e}")
                else:
                    logger.error(f"❌ Slide {is_slide_idx + 1} does not exist for IS table (only {len(self.presentation.slides)} slides)")
                    logger.error(f"IS data should be on slide 5, but presentation only has {len(self.presentation.slides)} slides")
                    
        except Exception as e:
            logger.error(f"Error embedding financial tables: {e}")
            import traceback
            logger.error(traceback.format_exc())

    def save(self, output_path: str):
        """Save the presentation"""
        if not self.presentation:
            raise ValueError("No presentation loaded")

        # Ensure output directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)

        self.presentation.save(output_path)
        logger.info(f"Presentation saved to: {output_path}")


class ReportGenerator:
    """Report generator that orchestrates the PPTX creation"""

    def __init__(self, template_path: str, markdown_file: str, output_path: str,
                 project_name: Optional[str] = None, language: str = 'english', row_limit: int = 20):
        self.template_path = template_path
        self.markdown_file = markdown_file
        self.output_path = output_path
        self.project_name = project_name
        self.language = language
        self.row_limit = row_limit

    def generate(self):
        """Generate the report"""
        logger.info(f"Starting PPTX generation...")
        logger.info(f"Template: {self.template_path}")
        logger.info(f"Markdown: {self.markdown_file}")
        logger.info(f"Output: {self.output_path}")
        logger.info(f"Language: {self.language}")
        logger.info(f"Project: {self.project_name}")

        # Read markdown content
        with open(self.markdown_file, 'r', encoding='utf-8') as f:
            md_content = f.read()

        logger.info(f"Content length: {len(md_content)} characters")

        # Create PowerPoint generator
        generator = PowerPointGenerator(self.template_path, self.language, self.row_limit)

        try:
            # Generate the report
            generator.generate_full_report(md_content, None, self.output_path)

            # Update project titles if project name provided
            if self.project_name:
                generator.update_project_titles(self.project_name, 'BS')  # Default to BS

        except Exception as e:
            logger.error(f"Report generation failed: {e}")
            raise

        logger.info(f"✅ PPTX generation completed: {self.output_path}")


def export_pptx(template_path: str, markdown_path: str, output_path: str,
                project_name: Optional[str] = None, excel_file_path: Optional[str] = None,
                language: str = 'english', statement_type: str = 'BS', row_limit: int = 20):
    """
    Export PowerPoint presentation from markdown content

    Args:
        template_path: Path to PPTX template
        markdown_path: Path to markdown content file
        output_path: Output PPTX file path
        project_name: Project/entity name for titles
        excel_file_path: Optional Excel file for embedding data
        language: Language ('english' or 'chinese')
        statement_type: Statement type ('BS' or 'IS')
        row_limit: Maximum rows per slide
    """
    generator = ReportGenerator(template_path, markdown_path, output_path,
                              project_name, language, row_limit)
    generator.generate()

    if not os.path.exists(output_path):
        raise FileNotFoundError(f"PPTX file was not created at {output_path}")

    # Update project titles with correct statement type
    if project_name:
        temp_presentation = Presentation(output_path)
        pptx_gen = PowerPointGenerator(template_path, language, row_limit)
        pptx_gen.presentation = temp_presentation
        pptx_gen.update_project_titles(project_name, statement_type)
        temp_presentation.save(output_path)
        # Ensure presentation is properly closed
        del temp_presentation
        del pptx_gen

    # Note: Excel embedding functionality would need to be implemented
    # if excel_file_path and project_name:
    #     embed_excel_data_in_pptx(output_path, excel_file_path, sheet_name, project_name, statement_type=statement_type)

    logger.info(f"✅ PowerPoint presentation successfully exported to: {output_path}")
    return output_path


def export_pptx_from_structured_data_combined(template_path: str, bs_data: List[Dict], is_data: List[Dict], 
                                              output_path: str, project_name: Optional[str] = None, 
                                              language: str = 'english', temp_path: Optional[str] = None,
                                              selected_sheet: Optional[str] = None, is_chinese_databook: bool = False):
    """
    Export ONE combined PowerPoint presentation with both BS and IS
    
    Args:
        template_path: Path to PPTX template
        bs_data: List of BS account data dicts
        is_data: List of IS account data dicts
        output_path: Output PPTX file path
        project_name: Project/entity name for titles
        language: Language ('english' or 'chinese')
        temp_path: Path to Excel file for table embedding
        selected_sheet: Sheet name for table embedding
    """
    try:
        logger.info(f"Starting COMBINED PPTX generation...")
        logger.info(f"Template: {template_path}")
        logger.info(f"Output: {output_path}")
        logger.info(f"Language: {language}")
        logger.info(f"BS accounts: {len(bs_data)}, IS accounts: {len(is_data)}")

        # Load template
        generator = PowerPointGenerator(template_path, language, row_limit=20)
        generator.load_template()

        # Apply BS data to slides 1-4
        if bs_data:
            generator.apply_structured_data_to_slides(bs_data, 1, project_name, 'BS', is_chinese_databook=is_chinese_databook)
        
        # Apply IS data to slides 5-8
        if is_data:
            generator.apply_structured_data_to_slides(is_data, 5, project_name, 'IS', is_chinese_databook=is_chinese_databook)
        
        # Embed financial tables: BS to page 1, IS to page 5
        # IMPORTANT: Do this BEFORE removing unused slides
        # This ensures tables are embedded on the correct slides
        if temp_path and selected_sheet:
            generator.embed_financial_tables(temp_path, selected_sheet, project_name, language)
        
        # NOW remove unused slides at the very end, after all content and tables are embedded
        if hasattr(generator, '_unused_slides_to_remove') and generator._unused_slides_to_remove:
            # Sort in reverse order to maintain indices while removing
            unused_slides_sorted = sorted(set(generator._unused_slides_to_remove), reverse=True)
            logger.info(f"Removing {len(unused_slides_sorted)} unused slides at the end: {[idx + 1 for idx in unused_slides_sorted]}")
            generator._remove_slides(unused_slides_sorted)
            logger.info(f"✅ Removed {len(unused_slides_sorted)} unused slides")
        
        # Save presentation
        generator.save(output_path)
        
        logger.info(f"✅ Combined PPTX generation completed: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"PPTX generation failed: {e}")
        import traceback
        logger.error(traceback.format_exc())
        raise


def export_pptx_from_structured_data(template_path: str, structured_data: List[Dict], output_path: str,
                                     project_name: Optional[str] = None, language: str = 'english',
                                     statement_type: str = 'BS', start_slide: int = 1):
    """
    Export PowerPoint presentation from structured data (not markdown)
    
    Args:
        template_path: Path to PPTX template
        structured_data: List of account data dicts with keys: account_name, financial_data, commentary, summary
        output_path: Output PPTX file path
        project_name: Project/entity name for titles
        language: Language ('english' or 'chinese')
        statement_type: Statement type ('BS' or 'IS')
        start_slide: Starting slide index (1-4 for BS, 5-8 for IS)
    """
    try:
        logger.info(f"Starting PPTX generation from structured data...")
        logger.info(f"Template: {template_path}")
        logger.info(f"Output: {output_path}")
        logger.info(f"Language: {language}")
        logger.info(f"Statement type: {statement_type}, Start slide: {start_slide}")
        logger.info(f"Accounts to process: {len(structured_data)}")

        # Load template
        generator = PowerPointGenerator(template_path, language, row_limit=20)
        generator.load_template()

        # Apply structured data to slides
        generator.apply_structured_data_to_slides(structured_data, start_slide, project_name, statement_type)

        # Save presentation
        generator.save(output_path)
        
        logger.info(f"✅ PPTX generation completed: {output_path}")
        return output_path
        
    except Exception as e:
        logger.error(f"PPTX generation failed: {e}")
        raise


def merge_presentations(bs_presentation_path: str, is_presentation_path: str, output_path: str):
    """
    Merge Balance Sheet and Income Statement presentations into a single presentation.

    Args:
        bs_presentation_path: Path to Balance Sheet presentation
        is_presentation_path: Path to Income Statement presentation
        output_path: Path for merged output presentation
    """
    try:
        logger.info("🔄 Starting presentation merge...")
        logger.info(f"   BS: {bs_presentation_path}")
        logger.info(f"   IS: {is_presentation_path}")

        # Load BS presentation as base
        merged_prs = Presentation(bs_presentation_path)

        # Load IS presentation
        is_prs = Presentation(is_presentation_path)

        # Copy all slides from IS to BS
        # Use XML-level copying for reliable slide duplication
        import xml.etree.ElementTree as ET
        from copy import deepcopy
        
        for slide_idx, slide in enumerate(is_prs.slides):
            try:
                # Get the slide layout
                slide_layout = slide.slide_layout
                
                # Create new slide with same layout
                new_slide = merged_prs.slides.add_slide(slide_layout)
                
                # Get XML elements
                source_slide_xml = slide._element
                target_slide_xml = new_slide._element
                
                # Remove placeholder shapes from new slide (from layout)
                # We'll replace them with actual content
                shapes_to_remove = list(new_slide.shapes)
                for shape in shapes_to_remove:
                    try:
                        sp_tree = target_slide_xml.get_or_add_spTree()
                        sp_tree.remove(shape._element)
                    except:
                        pass
                
                # Copy all shapes from source slide
                source_sp_tree = source_slide_xml.get_or_add_spTree()
                target_sp_tree = target_slide_xml.get_or_add_spTree()
                
                for shape_element in source_sp_tree:
                    # Deep copy the shape element
                    new_shape_element = deepcopy(shape_element)
                    # Add to target slide
                    target_sp_tree.append(new_shape_element)
                    
            except Exception as e:
                logger.error(f"Error copying slide {slide_idx}, using fallback method: {e}")
                # Fallback: simple text copy
                slide_layout = slide.slide_layout
                new_slide = merged_prs.slides.add_slide(slide_layout)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for new_shape in new_slide.shapes:
                            if (hasattr(new_shape, 'name') and hasattr(shape, 'name') and
                                new_shape.name == shape.name and new_shape.has_text_frame):
                                new_shape.text_frame.text = shape.text_frame.text
                                break

        # Save merged presentation
        merged_prs.save(output_path)
        
        # Ensure presentation objects are properly closed
        del merged_prs
        del is_prs
        
        # Force garbage collection to ensure file handles are released
        import gc
        gc.collect()

        logger.info("✅ Presentation merge completed successfully")
    except Exception as e:
        logger.error(f"❌ Presentation merge failed: {e}")
        raise
