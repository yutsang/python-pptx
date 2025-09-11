"""
PowerPoint export functionality moved from common/
"""

import os
import textwrap
import logging
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from dataclasses import dataclass
from typing import List, Tuple
from datetime import datetime
import re
from pptx.oxml.ns import qn

logging.basicConfig(level=logging.INFO)

def get_tab_name(project_name):
    """Get tab name based on project name - more flexible approach"""
    try:
        # Hardcoded mappings for known entities
        if project_name == 'Haining':
            return "BSHN"
        elif project_name == 'Nanjing':
            return "BSNJ"
        elif project_name == 'Ningbo':
            return "BSNB"
        else:
            # For unknown entities, return the project name itself to avoid None
            print(f"Warning: Unknown project name '{project_name}', using project name as fallback")
            return project_name
    except Exception:
        return None

def clean_content_quotes(content):
    """Remove outermost quotation marks from content while preserving legitimate internal quotes."""
    if not content:
        return content

    # Handle straight quotes
    content = re.sub(r'^"([^"]*)"$', r'\1', content)
    # Handle curly quotes
    content = re.sub(r'^"([^"]*)"$', r'\1', content)
    content = re.sub(r'^"([^"]*)"$', r'\1', content)

    return content

def detect_chinese_text(text):
    """Detect if text contains Chinese characters"""
    if not text:
        return False
    return any('\u4e00' <= char <= '\u9fff' for char in text)

def get_font_name_for_text(text, language='english'):
    """Get appropriate font name based on text content and language"""
    if language == 'chinese' or detect_chinese_text(text):
        return "Microsoft YaHei"  # Good Chinese font
    else:
        return "Calibri"  # Standard English font

def get_font_size_for_text(text, base_size=14):
    """Get appropriate font size based on text content"""
    from pptx.util import Pt
    
    if detect_chinese_text(text):
        # Chinese text might need slightly larger size for readability
        return Pt(base_size + 1)
    else:
        return Pt(base_size)

def export_pptx(template_path, markdown_path, output_path, project_name, excel_file_path=None, language='english'):
    """Export PowerPoint presentation from markdown content"""
    try:
        # Load template
        prs = Presentation(template_path)
        
        # Read markdown content
        with open(markdown_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        # Parse markdown sections
        sections = parse_markdown_sections(content)
        
        # Create slides for each section
        for section_title, section_content in sections.items():
            slide_layout = prs.slide_layouts[1]  # Title and Content layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Set title
            title = slide.shapes.title
            title.text = section_title
            
            # Set content
            content_placeholder = slide.placeholders[1]
            content_placeholder.text = section_content
            
            # Apply language-specific formatting
            if language == 'chinese':
                apply_chinese_formatting(slide)
        
        # Save presentation
        prs.save(output_path)
        print(f"✅ PowerPoint exported successfully: {output_path}")
        
    except Exception as e:
        print(f"❌ PowerPoint export failed: {e}")
        raise

def parse_markdown_sections(content):
    """Parse markdown content into sections"""
    sections = {}
    current_section = None
    current_content = []
    
    for line in content.split('\n'):
        if line.startswith('## '):
            # Save previous section
            if current_section:
                sections[current_section] = '\n'.join(current_content).strip()
            
            # Start new section
            current_section = line[3:].strip()
            current_content = []
        elif line.startswith('### '):
            # Subsection
            current_content.append(line[4:].strip() + ':')
        else:
            # Regular content
            if line.strip():
                current_content.append(line.strip())
    
    # Save last section
    if current_section:
        sections[current_section] = '\n'.join(current_content).strip()
    
    return sections

def apply_chinese_formatting(slide):
    """Apply Chinese-specific formatting to slide"""
    try:
        # Apply to title
        if slide.shapes.title and slide.shapes.title.text:
            title_frame = slide.shapes.title.text_frame
            for paragraph in title_frame.paragraphs:
                for run in paragraph.runs:
                    if detect_chinese_text(run.text):
                        run.font.name = "Microsoft YaHei"
                        run.font.size = Pt(24)
        
        # Apply to content
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if detect_chinese_text(run.text):
                            run.font.name = "Microsoft YaHei"
                            run.font.size = Pt(14)
    except Exception as e:
        print(f"Warning: Could not apply Chinese formatting: {e}")

def merge_presentations(bs_path, is_path, output_path):
    """Merge Balance Sheet and Income Statement presentations"""
    try:
        # Load BS presentation as base
        prs = Presentation(bs_path)
        
        # Load IS presentation
        is_prs = Presentation(is_path)
        
        # Add IS slides to BS presentation (skip title slide)
        for i, slide in enumerate(is_prs.slides):
            if i == 0:  # Skip title slide
                continue
            
            # Copy slide layout and content
            slide_layout = prs.slide_layouts[1]
            new_slide = prs.slides.add_slide(slide_layout)
            
            # Copy title
            if slide.shapes.title and new_slide.shapes.title:
                new_slide.shapes.title.text = slide.shapes.title.text
            
            # Copy content (simplified)
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text and shape != slide.shapes.title:
                    if len(new_slide.placeholders) > 1:
                        new_slide.placeholders[1].text = shape.text
                        break
        
        # Save merged presentation
        prs.save(output_path)
        print(f"✅ Presentations merged successfully: {output_path}")
        
    except Exception as e:
        print(f"❌ Presentation merge failed: {e}")
        raise