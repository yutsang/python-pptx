"""
fdd_utils/text_metrics.py — Pillow-based text fitting for PPTX shapes.

Standalone module: predicts how text will wrap inside a PPTX text frame by
reading the actual font's glyph advance widths (and GPOS kerning, via
FreeType under the hood) instead of guessing with a chars-per-line constant.

Why this exists
---------------
The legacy heuristic in pptx.py multiplies a hardcoded chars-per-line value
(e.g. 100 for single-slot ENG @ 9pt) by a width scale factor. That model
treats every glyph as the same width — fine for monospace, wrong for any
real font. "WWWW" and "iiii" land 2-3x apart in actual rendered width, and
mixed CJK + Latin commentary makes the error worse. This module replaces
the guess with a real measurement: for each candidate font size, it walks
the text word-by-word, asks Pillow for the rendered width, and decides
where to break lines.

Public API
----------
    text_box_from_shape(shape, *, bullet_indent_emu=0) -> TextBox
        Effective writable text box of a python-pptx shape, in points.

    fit_text_to_box(text, box, *, family, is_cjk, ...) -> FitResult
        Pick the largest candidate font size whose wrapped lines fit.

    wrap_text(text, font, max_width_pt) -> list[str]
        Greedy line wrap using real glyph widths.

    add_norm_autofit(text_frame, font_scale_pct=100, ...) -> None
        Belt-and-braces XML hint so PowerPoint can shrink further on open.

Run this file directly to see predictions against the project template:

    python -m fdd_utils.text_metrics
"""

from __future__ import annotations

import os
import platform
import shutil
import subprocess
import sys
from dataclasses import dataclass
from functools import lru_cache
from typing import List, Optional, Sequence, Tuple

from PIL import ImageFont


# ---------------------------------------------------------------------------
# Unit conversions and PPTX defaults
# ---------------------------------------------------------------------------

EMU_PER_INCH = 914400
EMU_PER_POINT = 12700

# OOXML bodyPr default insets (used when the attribute is absent).
# Confirmed from the ECMA-376 spec, NOT 0.05" all-around as some docs say.
DEFAULT_LIns_EMU = 91440   # 0.10"
DEFAULT_RIns_EMU = 91440   # 0.10"
DEFAULT_TIns_EMU = 45720   # 0.05"
DEFAULT_BIns_EMU = 45720   # 0.05"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"


def emu_to_pt(emu: float) -> float:
    return emu / EMU_PER_POINT


def pt_to_emu(pt: float) -> int:
    return int(round(pt * EMU_PER_POINT))


# ---------------------------------------------------------------------------
# Font path resolution
# ---------------------------------------------------------------------------
#
# We try a curated, OS-specific fallback chain first, then fall back to
# fontconfig (fc-match) if the family is not in the chain. This avoids
# silently substituting a Latin font for a CJK request, which would give
# wildly wrong widths.

_LATIN_FALLBACKS = {
    "Darwin": [
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
    ],
    "Windows": [
        r"C:\Windows\Fonts\arial.ttf",
        r"C:\Windows\Fonts\Arial.ttf",
    ],
    "Linux": [
        "/usr/share/fonts/truetype/msttcorefonts/Arial.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ],
}

_CJK_FALLBACKS = {
    "Darwin": [
        "/Library/Fonts/Microsoft/Microsoft YaHei.ttf",
        "/Library/Fonts/Microsoft YaHei.ttf",
        "/System/Library/Fonts/PingFang.ttc",
        "/System/Library/Fonts/STHeiti Light.ttc",
        "/System/Library/Fonts/Hiragino Sans GB.ttc",
    ],
    "Windows": [
        r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\msyh.ttf",
        r"C:\Windows\Fonts\msyhbd.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
    ],
    "Linux": [
        "/usr/share/fonts/truetype/wqy/wqy-microhei.ttc",
        "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
        "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    ],
}


def _fc_match(family: str) -> Optional[str]:
    """Use fontconfig fc-match to find a font file path. Returns None if
    fc-match is not installed or returns nothing useful."""
    if not shutil.which("fc-match"):
        return None
    try:
        out = subprocess.check_output(
            ["fc-match", "-f", "%{file}", family],
            text=True,
            stderr=subprocess.DEVNULL,
            timeout=2,
        ).strip()
        return out or None
    except (subprocess.SubprocessError, OSError):
        return None


@lru_cache(maxsize=64)
def resolve_font_path(family: str, *, is_cjk: bool) -> Optional[str]:
    """Find a TTF/OTF/TTC file for the requested family.

    Returns None if nothing usable is found — the caller should warn or
    hard-fail rather than silently using a wrong font.
    """
    if os.path.isabs(family) and os.path.exists(family):
        return family

    system = platform.system()
    candidates = list(
        _CJK_FALLBACKS.get(system, []) if is_cjk else _LATIN_FALLBACKS.get(system, [])
    )

    for path in candidates:
        if os.path.exists(path):
            return path

    fc = _fc_match(family)
    if fc and os.path.exists(fc):
        return fc

    return None


# ---------------------------------------------------------------------------
# Font cache
# ---------------------------------------------------------------------------
#
# Loading a TTF takes ~5ms; reuse the loaded face across all measurements.
# Pillow's `truetype(path, size=N)` interprets N in the same unit it returns
# from getlength(). If we pass `size=9` (a point value), getlength() returns
# advance widths in points. Convenient: no DPI math needed.

_FONT_CACHE: dict = {}


def get_font(
    family: str,
    size_pt: float,
    *,
    is_cjk: bool,
) -> ImageFont.FreeTypeFont:
    key = (family, round(size_pt, 3), is_cjk)
    cached = _FONT_CACHE.get(key)
    if cached is not None:
        return cached

    path = resolve_font_path(family, is_cjk=is_cjk)
    if path is None:
        raise FileNotFoundError(
            f"Could not resolve a font file for family={family!r} "
            f"is_cjk={is_cjk}. Install the font, pass an absolute path "
            f"as the family, or extend text_metrics._LATIN_FALLBACKS / "
            f"_CJK_FALLBACKS for this OS."
        )

    font = ImageFont.truetype(path, size=size_pt)
    _FONT_CACHE[key] = font
    return font


# ---------------------------------------------------------------------------
# Shape geometry
# ---------------------------------------------------------------------------

@dataclass(frozen=True)
class TextBox:
    """Effective writable region of a PPTX text frame, in points."""

    width_pt: float
    height_pt: float
    bullet_indent_pt: float


def text_box_from_shape(
    shape,
    *,
    bullet_indent_emu: int = 0,
) -> TextBox:
    """Compute the effective text box of a python-pptx shape.

    Reads bodyPr lIns/rIns/tIns/bIns from the shape's text frame and
    subtracts them from the shape size. Optionally subtracts an extra
    bullet/indent margin (pass marL from the paragraph properties).
    """
    body_pr = None
    try:
        body_pr = shape.text_frame._txBody.bodyPr
    except Exception:
        pass

    def _attr(name: str, default: int) -> int:
        if body_pr is None:
            return default
        v = body_pr.get(name)
        if v is None:
            return default
        try:
            return int(v)
        except (TypeError, ValueError):
            return default

    l_ins = _attr("lIns", DEFAULT_LIns_EMU)
    r_ins = _attr("rIns", DEFAULT_RIns_EMU)
    t_ins = _attr("tIns", DEFAULT_TIns_EMU)
    b_ins = _attr("bIns", DEFAULT_BIns_EMU)

    width_emu = max(0, int(shape.width) - l_ins - r_ins - int(bullet_indent_emu))
    height_emu = max(0, int(shape.height) - t_ins - b_ins)

    return TextBox(
        width_pt=emu_to_pt(width_emu),
        height_pt=emu_to_pt(height_emu),
        bullet_indent_pt=emu_to_pt(bullet_indent_emu),
    )


# ---------------------------------------------------------------------------
# Word wrap simulator
# ---------------------------------------------------------------------------
#
# Mixed CJK + Latin handling:
#   - Each CJK character is its own break point (no spaces between).
#   - Latin runs are split on whitespace; words wrap as units.
#   - Joining a Latin word to another Latin word inserts a single space;
#     joining to/from a CJK char inserts nothing.
#
# We always measure the full prospective line, never deltas, because Pillow's
# getlength is non-additive: getlength("A") + getlength("B") != getlength("AB")
# whenever the font has GPOS kerning between A and B.

# Unicode ranges that should wrap per-character.
_CJK_RANGES = (
    (0x3000, 0x303F),   # CJK Symbols and Punctuation
    (0x3040, 0x309F),   # Hiragana
    (0x30A0, 0x30FF),   # Katakana
    (0x3400, 0x4DBF),   # CJK Unified Ideographs Extension A
    (0x4E00, 0x9FFF),   # CJK Unified Ideographs
    (0xAC00, 0xD7AF),   # Hangul Syllables
    (0xF900, 0xFAFF),   # CJK Compatibility Ideographs
    (0xFF00, 0xFFEF),   # Halfwidth and Fullwidth Forms
)


def _is_cjk_char(ch: str) -> bool:
    if not ch:
        return False
    cp = ord(ch[0])
    for lo, hi in _CJK_RANGES:
        if lo <= cp <= hi:
            return True
    return False


def _measure(font: ImageFont.FreeTypeFont, text: str) -> float:
    if not text:
        return 0.0
    return float(font.getlength(text))


def _atomize(text: str) -> List[Tuple[str, str]]:
    """Split a paragraph into (separator, atom) pairs.

    An atom is either a single CJK character or a Latin "word" (a maximal
    run of non-CJK, non-whitespace characters). The separator is " " when
    joining two Latin atoms, "" otherwise. Whitespace runs in the source
    are collapsed to a single separator and never become atoms themselves.
    """
    atoms: List[Tuple[str, str]] = []
    i = 0
    n = len(text)
    while i < n:
        ch = text[i]
        if ch.isspace():
            i += 1
            continue
        if _is_cjk_char(ch):
            atoms.append(("", ch))
            i += 1
            continue
        j = i
        while j < n and not text[j].isspace() and not _is_cjk_char(text[j]):
            j += 1
        word = text[i:j]
        if atoms and not _is_cjk_char(atoms[-1][1][-1]):
            sep = " "
        else:
            sep = ""
        atoms.append((sep, word))
        i = j
    return atoms


def wrap_paragraph(
    text: str,
    font: ImageFont.FreeTypeFont,
    max_width_pt: float,
) -> List[str]:
    """Greedy word wrap of a single paragraph using real glyph widths."""
    if not text or not text.strip():
        return [""]
    if max_width_pt <= 0:
        return [text]

    atoms = _atomize(text)
    if not atoms:
        return [""]

    lines: List[str] = []
    current = ""

    for sep, word in atoms:
        candidate = word if not current else current + sep + word
        if _measure(font, candidate) <= max_width_pt:
            current = candidate
            continue

        # Doesn't fit. Flush the current line, then try the word alone.
        if current:
            lines.append(current)

        if _measure(font, word) > max_width_pt:
            # Word itself is wider than the line — break per character.
            sub = ""
            for ch in word:
                if _measure(font, sub + ch) <= max_width_pt:
                    sub += ch
                else:
                    if sub:
                        lines.append(sub)
                    sub = ch
            current = sub
        else:
            current = word

    if current:
        lines.append(current)

    return lines or [""]


def wrap_text(
    text: str,
    font: ImageFont.FreeTypeFont,
    max_width_pt: float,
) -> List[str]:
    """Wrap multi-paragraph text. Empty source paragraphs are preserved as
    blank lines so the vertical capacity check counts them."""
    out: List[str] = []
    for para in str(text or "").split("\n"):
        if para.strip() == "":
            out.append("")
        else:
            out.extend(wrap_paragraph(para, font, max_width_pt))
    return out


# ---------------------------------------------------------------------------
# Vertical capacity & font-size search
# ---------------------------------------------------------------------------

def line_height_pt(
    font: ImageFont.FreeTypeFont,
    *,
    line_spacing: float = 1.0,
    extra_leading_pt: float = 0.0,
) -> float:
    """Compute baseline-to-baseline distance in points.

    Uses the font's ascent + descent (from the OS/2 / hhea tables) rather
    than a fixed multiplier of the font size. This is closer to PowerPoint's
    actual behavior. Calibrate the residual error with `extra_leading_pt`.
    """
    ascent, descent = font.getmetrics()
    return (ascent + descent) * line_spacing + extra_leading_pt


def lines_that_fit(height_pt: float, line_h_pt: float) -> int:
    if line_h_pt <= 0:
        return 0
    return int(height_pt // line_h_pt)


@dataclass
class FitResult:
    font_size_pt: float
    family: str
    is_cjk: bool
    lines: List[str]
    line_count: int
    line_height_pt: float
    capacity_lines: int
    fits: bool
    text_box: TextBox


def fit_text_to_box(
    text: str,
    box: TextBox,
    *,
    family: str,
    is_cjk: bool,
    candidate_sizes: Sequence[float] = (9.0, 8.5, 8.0, 7.5, 7.0),
    line_spacing: float = 1.0,
    extra_leading_pt: float = 0.0,
) -> FitResult:
    """Try each candidate size from largest to smallest. Return the first
    one whose wrapped line count is within the box's vertical capacity. If
    none fit, return the smallest size; the caller can then truncate, accept
    overflow, or rely on a normAutofit safety net.
    """
    last: Optional[FitResult] = None
    for size in candidate_sizes:
        font = get_font(family, size, is_cjk=is_cjk)
        lines = wrap_text(text, font, box.width_pt)
        line_h = line_height_pt(
            font,
            line_spacing=line_spacing,
            extra_leading_pt=extra_leading_pt,
        )
        capacity = lines_that_fit(box.height_pt, line_h)
        result = FitResult(
            font_size_pt=size,
            family=family,
            is_cjk=is_cjk,
            lines=lines,
            line_count=len(lines),
            line_height_pt=line_h,
            capacity_lines=capacity,
            fits=len(lines) <= capacity,
            text_box=box,
        )
        if result.fits:
            return result
        last = result

    assert last is not None
    return last


# ---------------------------------------------------------------------------
# normAutofit XML safety net
# ---------------------------------------------------------------------------

def add_norm_autofit(
    text_frame,
    *,
    font_scale_pct: float = 100.0,
    line_space_reduction_pct: float = 0.0,
) -> None:
    """Inject `<a:normAutofit fontScale=... lnSpcReduction=.../>` into the
    text frame's bodyPr.

    PowerPoint reads this on open and shrinks text accordingly. Pass 100.0
    / 0.0 for a no-op marker. LibreOffice and Keynote ignore this hint, so
    treat it as a safety net on top of an already-correct Pillow estimate,
    not the primary fitting mechanism.
    """
    from lxml import etree  # local import: only needed when this is called

    body_pr = text_frame._txBody.bodyPr
    for tag in ("normAutofit", "spAutoFit", "noAutofit"):
        for el in body_pr.findall(f"{{{A_NS}}}{tag}"):
            body_pr.remove(el)

    autofit = etree.SubElement(body_pr, f"{{{A_NS}}}normAutofit")
    autofit.set("fontScale", str(int(round(font_scale_pct * 1000))))
    if line_space_reduction_pct > 0:
        autofit.set(
            "lnSpcReduction",
            str(int(round(line_space_reduction_pct * 1000))),
        )


# ---------------------------------------------------------------------------
# Demo / sanity check
# ---------------------------------------------------------------------------

# Legacy CPL constants from config.yml, for side-by-side comparison.
_LEGACY_CPL = {
    ("single", False): 100, ("single", True): 50,
    ("L", False): 56, ("L", True): 30,
    ("R", False): 56, ("R", True): 30,
}


def _legacy_predict_lines(text: str, slot_key: str, is_cjk: bool) -> int:
    cpl = _LEGACY_CPL.get((slot_key, is_cjk), 60)
    n = len(text)
    return max(1, (n + cpl - 1) // cpl)


def _slot_key_for_shape_name(name: str) -> str:
    n = (name or "").lower()
    if n.endswith("_l"):
        return "L"
    if n.endswith("_r"):
        return "R"
    return "single"


def _demo() -> int:
    from pptx import Presentation

    template_path = os.path.join(os.path.dirname(__file__), "template.pptx")
    if not os.path.exists(template_path):
        print(f"Template not found at {template_path}", file=sys.stderr)
        return 1

    prs = Presentation(template_path)

    sample_eng = (
        "Cash and bank balances increased by RMB 12.3 million from RMB 45.6 "
        "million as at 31 Dec 2024 to RMB 57.9 million as at 30 Jun 2025, "
        "primarily driven by a one-off receipt of RMB 8.2 million from the "
        "disposal of the Foshan property and improved collection from key "
        "customers in Q2 2025. Trade receivables turnover days improved from "
        "62 days to 48 days over the same period."
    )
    sample_chi = (
        "现金及银行存款由2024年12月31日的人民币4,560万元增加人民币1,230万元至"
        "2025年6月30日的人民币5,790万元,主要由于佛山物业处置的一次性收款人民"
        "币820万元以及2025年第二季度对主要客户应收款项回收情况改善所致。应收"
        "账款周转天数同期由62天改善至48天。"
    )

    target_names = {"textmainbullets", "textmainbullets_l", "textmainbullets_r"}

    found = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            name = (shape.name or "").lower()
            if name in target_names and shape.has_text_frame:
                found.append((slide_idx, shape))

    if not found:
        print("No textMainBullets shapes found in template.", file=sys.stderr)
        return 1

    print("=" * 78)
    print("Pillow text-fitting demo — comparing real measurement vs. legacy CPL")
    print("=" * 78)

    seen_names = set()
    shown = 0
    for slide_idx, shape in found:
        if shape.name in seen_names:
            continue
        seen_names.add(shape.name)
        shown += 1
        if shown > 3:
            break

        slot_key = _slot_key_for_shape_name(shape.name)
        box = text_box_from_shape(shape)

        print()
        print(f"Slide {slide_idx + 1}, shape={shape.name!r}, slot={slot_key}")
        print(
            f"  raw size:        {int(shape.width) / EMU_PER_INCH:.2f}\" x "
            f"{int(shape.height) / EMU_PER_INCH:.2f}\""
        )
        print(
            f"  effective box:   {box.width_pt:.1f} pt x {box.height_pt:.1f} pt "
            f"(insets stripped)"
        )

        for label, text, is_cjk, family, line_spacing in (
            ("ENG", sample_eng, False, "Arial", 1.0),
            ("CHI", sample_chi, True, "Microsoft YaHei", 0.95),
        ):
            print()
            print(f"  --- {label} ({len(text)} chars) ---")
            try:
                fit = fit_text_to_box(
                    text,
                    box,
                    family=family,
                    is_cjk=is_cjk,
                    line_spacing=line_spacing,
                )
            except FileNotFoundError as e:
                print(f"  ERROR: {e}")
                continue

            legacy_lines = _legacy_predict_lines(text, slot_key, is_cjk)

            resolved = resolve_font_path(family, is_cjk=is_cjk)
            print(f"  font file:       {resolved}")
            print(
                f"  Pillow:          size={fit.font_size_pt} pt   "
                f"line_h={fit.line_height_pt:.2f} pt   "
                f"capacity={fit.capacity_lines}   "
                f"used={fit.line_count}   fits={fit.fits}"
            )
            cpl = _LEGACY_CPL.get((slot_key, is_cjk), 60)
            print(
                f"  Legacy CPL:      cpl={cpl}   "
                f"predicted_lines={legacy_lines}   "
                f"delta_vs_pillow={legacy_lines - fit.line_count:+d}"
            )
            for i, line in enumerate(fit.lines, 1):
                marker = " " if i <= fit.capacity_lines else "!"
                print(f"  {marker}L{i:>2}: {line}")

    print()
    print("Legend: '!' marks lines that overflow the box at the chosen size.")
    return 0


if __name__ == "__main__":
    sys.exit(_demo())
