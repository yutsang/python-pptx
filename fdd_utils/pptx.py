from __future__ import annotations

# --- begin pptx/text.py ---
from concurrent.futures import ThreadPoolExecutor, as_completed
import re
from typing import Optional

from pptx.util import Pt


def get_tab_name(project_name: str) -> Optional[str]:
    if not project_name:
        return None
    words = project_name.split()
    if words:
        return words[0]
    return None


def clean_content_quotes(content: str) -> str:
    if not content:
        return ""
    content = re.sub(r'^"*|"*$', "", content.strip())
    content = re.sub(r'""+', '"', content)
    return content


def detect_chinese_text(text: str, force_chinese_mode: bool = False) -> bool:
    if force_chinese_mode:
        return True
    return contains_predominantly_chinese_text(text)


def get_font_size_for_text(text: str, base_size: int = 9, force_chinese_mode: bool = False) -> Pt:
    # Deck-wide typography: every commentary run, every slide, every
    # language renders at a single fixed size. We intentionally ignore the
    # text, base_size, and force_chinese_mode arguments — any caller that
    # asked for something else would reintroduce the size-variation bug.
    return Pt(9)


def get_font_name_for_text(text: str, default_font: str = "Arial") -> str:
    # Same philosophy: one font for the whole deck. Arial has CJK fallback
    # glyphs via the system's default font substitution, so Chinese content
    # still renders correctly without switching to Microsoft YaHei (which
    # would change glyph width / baseline on some slides).
    return "Arial"


def get_line_spacing_for_text(text: str, force_chinese_mode: bool = False) -> float:
    return 0.9 if detect_chinese_text(text, force_chinese_mode) else 1.0


def get_space_after_for_text(text: str, force_chinese_mode: bool = False) -> Pt:
    return Pt(6) if detect_chinese_text(text, force_chinese_mode) else Pt(4)


def get_space_before_for_text(text: str, force_chinese_mode: bool = False) -> Pt:
    return Pt(3) if detect_chinese_text(text, force_chinese_mode) else Pt(2)


def replace_entity_placeholders(content: str, project_name: str) -> str:
    if not content or not project_name:
        return content
    replacements = {
        "[PROJECT]": project_name,
        "[Entity]": project_name,
        "[Company]": project_name,
    }
    for placeholder, replacement in replacements.items():
        content = content.replace(placeholder, replacement)
    return content
# --- end pptx/text.py ---

# --- begin pptx/payloads.py ---
from typing import Any, Dict, Iterable, List, Optional

import pandas as pd

from .financial_common import (
    contains_chinese_text,
    contains_predominantly_chinese_text,
    get_pipeline_result_text,
    load_yaml_file,
    package_file_path,
)
from .keyword_registry import (
    STATEMENT_ORDER_SKIP_KEYWORDS,
    SUMMARY_ACCOUNT_SKIP_KEYWORDS,
    translate_category_to_chinese,
    translate_statement_line_to_chinese,
)
from .workbook import find_mapping_key


PPTX_DEFAULT_SETTINGS: Dict[str, Any] = {
    "max_commentary_slides_per_statement": 4,
    "executive_summary": {
        "target_words_eng": 110,
        "target_chars_chi": 144,
        "max_sentences_eng": 4,
        "max_sentences_chi": 4,
        "max_tokens": 240,
        "validation_max_tokens": 180,
        "max_input_chars": 1400,
        "max_numeric_sentences": 2,
        "max_workers": 2,
        "enable_validation": True,
        "generation_temperature": 0.2,
        "validation_temperature": 0.1,
    },
    "commentary_packing": {
        "use_pillow_text_fitting": True,
        # Repeatedly bumped up (1.08 -> 1.15 -> 1.25, and the BS override
        # further to 1.13 -> 1.47) in response to page fill plateauing too
        # low -- but the real cause turned out to be _real_para_gap_pt/
        # _real_line_spacing assuming a 6-9pt inter-paragraph gap and 0.9
        # Chinese line spacing that _fill_text_main_bullets_with_category_
        # and_key never actually applies (it hardcodes a flat 3pt gap and
        # 1.0 spacing for every paragraph, any language) plus a capacity
        # formula that floored away up to a full extra line on every box.
        # Once those are fixed at the source (real capacity ~40-50% higher
        # than before), this relax tier goes back to being a small genuine
        # second-chance buffer instead of a compensating hack for an
        # undersized first tier.
        "shape_height_utilization": 1.00,
        "minimum_slot_lines": 22,
        "split_min_remaining_lines": 3,
        "split_min_content_lines": 5,
        # Lowered: pull a whole bullet forward into a slot even when the slot
        # is already 50% full (was 74%). This stops the first IS slot from
        # sitting at 35% fill while later slots are full.
        "move_whole_min_fill_ratio": 0.50,
        "target_fill_min_ratio": 0.95,
        "target_fill_max_ratio": 1.00,
        "ppt_length_ratio": 0.84,
        "ppt_min_chars_eng": 190,
        "ppt_min_chars_chi": 110,
        "ppt_max_sentences_eng": 6,
        "ppt_max_sentences_chi": 5,
        "ppt_max_numeric_sentences": 2,
        "category_line_cost": 0.95,
        "key_line_cost": 1.0,
        "continuation_spacing_penalty": 0.15,
        "line_height_padding_pt": 1.6,
        "split_slot_height_penalty": 1.02,
        "width_scale_min": 0.9,
        "width_scale_max": 1.22,
        "chars_per_line": {
            "single": {"eng": 100, "chi": 50},
            "L": {"eng": 56, "chi": 30},
            "R": {"eng": 56, "chi": 30},
            "default": {"eng": 66, "chi": 36},
        },
        "statement_overrides": {
            "BS": {
                "shape_height_utilization": 1.00,
                "line_height_padding_pt": 1.3,
                "chars_per_line": {
                    "single": {"eng": 106},
                    "L": {"eng": 59},
                    "R": {"eng": 59},
                    "default": {"eng": 69},
                },
            },
        },
    },
}


def _merge_nested_dict(base: Dict[str, Any], overrides: Dict[str, Any]) -> Dict[str, Any]:
    merged = dict(base or {})
    for key, value in (overrides or {}).items():
        if isinstance(value, dict) and isinstance(merged.get(key), dict):
            merged[key] = _merge_nested_dict(merged[key], value)
        else:
            merged[key] = value
    return merged


def _load_pptx_settings(config_path: Optional[str] = None) -> Dict[str, Any]:
    config = load_yaml_file(config_path or package_file_path("config.yml"))
    return _merge_nested_dict(PPTX_DEFAULT_SETTINGS, (config or {}).get("pptx") or {})


def _split_text_sentences(text: str, is_chinese: bool) -> List[str]:
    normalized = str(text or "").strip()
    if not normalized:
        return []
    if is_chinese:
        parts = re.split(r"(?<=[。！？；])", normalized)
    else:
        parts = re.split(r"(?<=[.!?;])\s+", normalized)
    return [part.strip() for part in parts if part and part.strip()]


def _join_text_sentences(sentences: List[str], is_chinese: bool) -> str:
    cleaned = [str(sentence or "").strip() for sentence in sentences if str(sentence or "").strip()]
    if not cleaned:
        return ""
    return "".join(cleaned) if is_chinese else " ".join(cleaned)


def _sentence_is_numeric_heavy(sentence: str) -> bool:
    text = str(sentence or "")
    numeric_tokens = re.findall(r"\d[\d,.\-]*%?|USD|HKD|RMB|CNY|EUR|JPY|\$", text, flags=re.IGNORECASE)
    return len(numeric_tokens) >= 2


def _build_compact_summary_text(
    text: str,
    *,
    is_chinese: bool,
    max_sentences: int,
    max_chars: int,
    max_numeric_sentences: int,
) -> str:
    sentences = _split_text_sentences(text, is_chinese)
    if not sentences:
        return str(text or "").strip()

    selected: List[str] = []
    numeric_sentences = 0
    for sentence in sentences:
        heavy = _sentence_is_numeric_heavy(sentence)
        if heavy and numeric_sentences >= max_numeric_sentences:
            continue
        candidate = _join_text_sentences(selected + [sentence], is_chinese)
        if selected and len(candidate) > max_chars:
            break
        selected.append(sentence)
        if heavy:
            numeric_sentences += 1
        if len(selected) >= max_sentences:
            break

    if not selected:
        selected = sentences[:1]

    summary = _join_text_sentences(selected, is_chinese).strip()
    if len(summary) > max_chars:
        summary = summary[:max_chars].rstrip(" ,;:/-") + "..."
    return summary.strip()


#: Legal-form tails and place-of-registration brackets carry NO identifying
#: information -- they say what kind of company it is and where it was
#: registered, not which company. The reference deck strips both throughout
#: (浙江卓圣, not 浙江卓圣物业管理有限公司) and only writes a full legal name
#: when first naming a specific contract counterparty.
_LEGAL_FORM_TAILS = (
    "股份有限责任公司", "股份有限公司", "有限责任公司", "私人有限公司",
    "有限合伙企业", "会计师事务所", "律师事务所", "有限公司",
)
#: Removed only when the bracket holds a short place name. A long bracket is
#: usually a real qualifier and gets left alone.
_REGISTRATION_BRACKET = re.compile(r"[（(][一-鿿]{2,4}[）)](?=[一-鿿]*(?:%s))"
                                   % "|".join(_LEGAL_FORM_TAILS))


def shorten_company_names(text: str) -> str:
    """Strip legal-form tails and registration brackets from company names.

    Done deterministically, NOT by prompt. Two separate prompt attempts
    failed outright: the databook's own row labels ARE the full legal names
    (e.g. a row literally reading *某某咨询管理有限公司), and the prompts
    carry a much older, much stronger instruction to reproduce entity names
    from the data exactly. The model correctly followed the stronger rule.
    Shortening a name is a mechanical string operation, so it belongs here
    rather than in an instruction that has to win an argument.

    Deliberately conservative: only the legal form and the registration
    bracket are removed, because neither identifies the party. Trailing
    business descriptors (物业管理 / 企业管理) are NOT stripped -- that would
    reach the reference deck's brevity but risks making two counterparties
    read the same.
    """
    body = str(text or "")
    if not body:
        return body
    body = _REGISTRATION_BRACKET.sub("", body)
    for tail in _LEGAL_FORM_TAILS:
        body = body.replace(tail, "")
    return body


def _normalize_slide_commentary_text(text: str) -> str:
    normalized = clean_content_quotes(str(text or ""))
    if not normalized:
        return ""
    normalized = normalized.replace("\r\n", "\n")
    normalized = re.sub(r"[ \t]+", " ", normalized)
    normalized = re.sub(r"\n{3,}", "\n\n", normalized)
    normalized = shorten_company_names(normalized)
    return normalized.strip()


def _extract_summary(content):
    text = str(content or "").strip()
    if not text:
        return ""
    if _looks_like_blocked_ai_content(text):
        return ""
    return text


def _looks_like_blocked_ai_content(text: str) -> bool:
    normalized = str(text or "").strip()
    if not normalized:
        return False
    lowered = normalized.lower()
    blocked_markers = (
        "<!doctype html",
        "<html",
        "ac_block_page",
        "sp.eagleyun.cn",
        "form.submit()",
        "api.deepseek.com",
        "request_uri",
        "request_user_agent",
    )
    return any(marker in lowered for marker in blocked_markers)


def _extract_final_content(result_dict):
    # Defence in depth: strip any Qwen3 <think> block that slipped through (e.g.
    # via run_generator_reprompt, which skips the _ensure pass) before it can
    # render into a no-autofit text box and overflow / leak reasoning.
    from fdd_utils.ai import strip_thinking
    return strip_thinking(get_pipeline_result_text(result_dict))


def _build_statement_order(
    financial_statement_df: Optional[pd.DataFrame],
    mappings: Dict[str, Any],
) -> tuple[Dict[str, int], Dict[str, str]]:
    financial_statement_order: Dict[str, int] = {}
    statement_display_names: Dict[str, str] = {}
    if financial_statement_df is None or financial_statement_df.empty or len(financial_statement_df.columns) == 0:
        return financial_statement_order, statement_display_names

    first_col = financial_statement_df.iloc[:, 0]
    skip_keywords = STATEMENT_ORDER_SKIP_KEYWORDS
    for idx, account_name_in_statement in enumerate(first_col):
        if pd.isna(account_name_in_statement):
            continue

        account_name_str = str(account_name_in_statement).strip()
        if not account_name_str or any(skip in account_name_str.lower() for skip in skip_keywords):
            continue

        mapping_key = find_mapping_key(account_name_str, mappings)
        if mapping_key:
            financial_statement_order[mapping_key] = idx
            statement_display_names[mapping_key] = account_name_str

        financial_statement_order[account_name_str] = idx

    return financial_statement_order, statement_display_names


# Common Traditional -> Simplified character pairs that show up in
# mappings.yml's own Chinese aliases (e.g. "貨幣資金" vs "货币资金"). Not a
# general S/T converter -- just enough of this narrow FDD-account-name
# vocabulary to (a) normalize a Traditional alias to its Simplified spelling
# for detecting "these two aliases are the same concept, just different
# script" and (b) prefer Simplified for the final label, matching the
# convention CATEGORY_TRANSLATIONS_ZH already uses for section headers
# (fdd_utils/keyword_registry.py).
_TRADITIONAL_TO_SIMPLIFIED_PAIRS = {
    "貨": "货", "應": "应", "產": "产", "負": "负", "稅": "税", "資": "资",
    "積": "积", "準": "准", "幣": "币", "讓": "让", "長": "长", "遞": "递",
    "認": "认", "賬": "账", "債": "债", "現": "现", "後": "后", "裡": "里",
    "歸": "归", "屬": "属", "歷": "历", "業": "业", "當": "当", "項": "项",
    "餘": "余", "繳": "缴", "會": "会", "計": "计", "師": "师", "貴": "贵",
    "賣": "卖", "買": "买", "須": "须", "廠": "厂", "聯": "联", "繫": "系",
    "務": "务", "單": "单", "帳": "帐", "報": "报",
    # Audited (extracted directly from mappings.yml's own alias vocabulary,
    # see: "python3 -c 'grep aliases + CJK char diff'" in this commit's history)
    "內": "内", "動": "动", "發": "发", "實": "实", "損": "损", "據": "据",
    "攤": "摊", "權": "权", "減": "减", "潤": "润", "無": "无", "營": "营",
    "綜": "综", "職": "职", "譽": "誉", "財": "财", "費": "费", "賃": "赁",
    "預": "预",
}
_TRADITIONAL_TO_SIMPLIFIED = str.maketrans(_TRADITIONAL_TO_SIMPLIFIED_PAIRS)


def _find_chinese_display_name(mapping_key: str, fallback: str, mappings: Dict[str, Any]) -> str:
    """display_name (built in _build_statement_order) is whatever literal
    text sits in that account's row in the Financials-summary sheet -- for
    an English-labelled source databook (e.g. Kunshan's "Cash at bank and
    on hand"), that stays English even when the REPORT is being generated
    in Chinese, since nothing translates it. mappings.yml already carries a
    Chinese alias for essentially every account (used for matching Chinese
    source sheets) -- reuse one of those as the Chinese-report label instead
    of leaving the English source text in an otherwise fully-translated
    Chinese bullet.

    A single mapping_key's aliases list often mixes a precise term with
    broader catch-all synonyms for MATCHING purposes (e.g. Capital's aliases
    include both "实收资本"/paid-in capital AND "股东权益"/shareholders'
    equity -- correct for fuzzy-matching a client's sheet name, but "股东权益"
    would be a semantically WRONG display label for a paid-in-capital
    account). Picking "any CJK alias" isn't safe.

    Approach: group the CJK aliases by their Simplified-normalized spelling
    (so a Traditional/Simplified pair like "實收資本"/"实收资本" is treated as
    ONE concept, not two competing candidates), rank the resulting concepts
    by how close any of their members sit to where `fallback` itself
    appears in the alias list (aliases are typically authored in loosely
    paired EN/CN blocks, e.g. "...实收资本", "Paid-in capital" sit adjacent),
    and return the Simplified spelling of the winning concept."""
    config = mappings.get(mapping_key) if isinstance(mappings, dict) else None
    aliases = config.get("aliases") if isinstance(config, dict) else None
    if not isinstance(aliases, list) or not aliases:
        return fallback
    aliases = [str(a).strip() for a in aliases]
    chinese_indices = [i for i, a in enumerate(aliases) if contains_chinese_text(a)]
    if not chinese_indices:
        return fallback

    normalized_fallback = str(fallback or "").strip().lower()
    fallback_idx = next(
        (i for i, a in enumerate(aliases) if a.lower() == normalized_fallback),
        len(aliases) // 2,  # unknown position -- fall back to the list midpoint
    )

    concepts: Dict[str, Dict[str, Any]] = {}
    for i in chinese_indices:
        simplified = aliases[i].translate(_TRADITIONAL_TO_SIMPLIFIED)
        entry = concepts.setdefault(simplified, {"best_distance": None, "simplified_form": None})
        distance = abs(i - fallback_idx)
        if entry["best_distance"] is None or distance < entry["best_distance"]:
            entry["best_distance"] = distance
        if aliases[i] == simplified:
            entry["simplified_form"] = aliases[i]

    best_concept = min(concepts.items(), key=lambda kv: kv[1]["best_distance"])
    simplified_key, entry = best_concept
    return entry["simplified_form"] or simplified_key


def _translate_statement_row_label(label: str, mappings: Optional[Dict[str, Any]]) -> str:
    """Translates ONE row label from the embedded BS/IS summary table
    (embed_financial_tables) to Chinese, for a Chinese-language report.

    Unlike commentary bullets (which carry their own resolved mapping_key
    via build_pptx_structured_payloads), these rows come straight from
    extract_balance_sheet_and_income_statement's own parse of the raw
    Financials-sheet text -- there was previously NO translation path for
    them at all, so a Chinese report's embedded table stayed 100% English
    even though the title/commentary around it were fully translated.

    Two label classes need two different lookups: (1) individual account
    rows (e.g. "Cash at bank and on hand") resolve via the same
    mappings.yml alias machinery _find_chinese_display_name already uses
    for commentary, by first recovering the mapping_key with
    find_mapping_key; (2) statement-structure total/subtotal rows (e.g.
    "Total current assets") aren't mapping_key accounts at all, so they
    fall back to the small fixed STATEMENT_TOTAL_LINE_TRANSLATIONS_ZH
    table. Returns `label` unchanged if neither resolves (e.g. an
    already-Chinese source label, or a genuinely unmapped line) --
    partial coverage beats a blank or crashed cell.
    """
    label = str(label or "")
    if not label.strip() or contains_chinese_text(label):
        return label
    if mappings:
        mapping_key = find_mapping_key(label, mappings)
        if mapping_key:
            return _find_chinese_display_name(mapping_key, label, mappings)
    return translate_statement_line_to_chinese(label) or label


def _has_significant_balance(financial_data: Optional[pd.DataFrame]) -> bool:
    """Does this account carry a balance worth commenting on?

    The dataframe handed in is the PROJECTION frame -- one stage, ONE date
    column -- so scanning it answers "is the latest period non-zero", not the
    question actually being asked. An account that ran to nil by the reporting
    date but had real activity in earlier periods (a real Crescent 投资收益
    reads 0 at 2026-03-31 with prior-year movement) was silently dropped from
    the deck after its commentary had already been generated and paid for.

    workbook.py already computes exactly the right signal while parsing --
    any_period_nonzero_by_description, |v| >= 0.01 over ALL periods of the
    stage -- and filter_zero_value_rows already uses it to decide which ROWS
    survive. Use it here for the same decision at ACCOUNT level; fall back to
    the single-column scan only when the attribute is absent.
    """
    if financial_data is None or financial_data.empty:
        return True

    any_period = financial_data.attrs.get("any_period_nonzero_by_description")
    if isinstance(any_period, dict) and any_period:
        return any(bool(v) for v in any_period.values())

    numeric_cols = financial_data.select_dtypes(include=[float, int]).columns
    if len(numeric_cols) == 0:
        return True

    for col in numeric_cols:
        if (financial_data[col].abs() >= 0.01).any():
            return True
    return False


def build_pptx_structured_payloads(
    ai_results,
    mappings,
    bs_is_results=None,
    dfs=None,
):
    if not ai_results:
        return {"BS": [], "IS": []}

    balance_sheet_df = bs_is_results.get("balance_sheet") if bs_is_results else None
    income_statement_df = bs_is_results.get("income_statement") if bs_is_results else None
    bs_order, bs_display_names = _build_statement_order(balance_sheet_df, mappings)
    is_order, is_display_names = _build_statement_order(income_statement_df, mappings)

    payloads = {"BS": [], "IS": []}
    sortable_items = {"BS": [], "IS": []}

    for account_key, result in ai_results.items():
        mapping_key = find_mapping_key(account_key, mappings)
        if not mapping_key:
            continue

        account_type = mappings[mapping_key].get("type")
        if account_type not in {"BS", "IS"}:
            continue

        financial_data = dfs.get(account_key) if dfs and account_key in dfs else None
        if not _has_significant_balance(financial_data):
            # Dropping an account here throws away commentary the AI pipeline
            # has already produced and been paid for, and it happens silently
            # -- the only trace was a smaller "IS items: N". Say which account
            # and on what evidence, so the next run answers "why is 投资收益
            # missing" instead of another round of guessing.
            _attrs = getattr(financial_data, "attrs", {}) or {}
            _any = _attrs.get("any_period_nonzero_by_description")
            logger.warning(
                "Dropping %s from the deck: no period carries a balance. "
                "any_period_nonzero_by_description=%s, stage=%s, columns=%s",
                mapping_key,
                dict(_any) if isinstance(_any, dict) else _any,
                _attrs.get("prompt_analysis_stage"),
                list(getattr(financial_data, "columns", [])) or None,
            )
            continue

        final_content = _extract_final_content(result)
        commentary_text = (
            str(final_content).strip()
            if final_content and str(final_content).strip()
            else f"[No content generated for {account_key}]"
        )

        clause_reviews: List[Dict[str, Any]] = []
        if isinstance(result, dict):
            validator_metadata = result.get("agent_4_validation") or {}
            if isinstance(validator_metadata, dict):
                raw_reviews = validator_metadata.get("clause_reviews") or []
                if isinstance(raw_reviews, list):
                    clause_reviews = [r for r in raw_reviews if isinstance(r, dict)]

        statement_order = bs_order if account_type == "BS" else is_order
        statement_display_names = bs_display_names if account_type == "BS" else is_display_names
        order = statement_order.get(mapping_key, statement_order.get(account_key, 9999))
        display_name = statement_display_names.get(mapping_key, account_key)

        sortable_items[account_type].append(
            (
                order,
                mappings[mapping_key].get("category", ""),
                mapping_key,
                {
                    "account_name": account_key,
                    "mapping_key": mapping_key,
                    "display_name": display_name,
                    "display_name_zh": _find_chinese_display_name(mapping_key, display_name, mappings),
                    "category": mappings[mapping_key].get("category", ""),
                    "financial_data": financial_data,
                    "commentary": commentary_text,
                    "clause_reviews": clause_reviews,
                    "summary": _extract_summary(final_content) if final_content else "",
                    # Predominantly-Chinese (>30%), not "contains any CJK
                    # character" -- this flag drives CJK-vs-Latin text
                    # WRAPPING/measurement throughout the packing pipeline
                    # (_calculate_content_lines, _calculate_max_lines_for_
                    # textbox's whole-statement is_chinese_any). An English
                    # commentary that merely names a Chinese counterparty/
                    # person (e.g. "...payable to the related party 维彧")
                    # is still fundamentally Latin-script prose; measuring
                    # it (and, via is_chinese_any, EVERY slot's capacity in
                    # the whole statement) with CJK line-height/spacing/
                    # wrap rules instead of Arial's produced a systematic
                    # believed-vs-actually-rendered fill gap.
                    "is_chinese": contains_predominantly_chinese_text(commentary_text),
                },
            )
        )

    for statement_type in ["BS", "IS"]:
        payloads[statement_type] = [
            item
            for _order, _category, _mapping_key, item in sorted(
                sortable_items[statement_type],
                key=lambda row: (row[0], row[1], row[2]),
            )
        ]

    return payloads
# --- end pptx/payloads.py ---

# --- begin pptx/exporters.py ---
import copy
import logging
import os
import posixpath
import time
import traceback
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.oxml.ns import qn

logger = logging.getLogger(__name__)


class ReportGenerator:
    """Report generator that orchestrates PPTX creation from markdown."""

    def __init__(
        self,
        template_path: str,
        markdown_file: str,
        output_path: str,
        project_name: Optional[str] = None,
        language: str = "english",
        row_limit: int = 20,
    ):
        self.template_path = template_path
        self.markdown_file = markdown_file
        self.output_path = output_path
        self.project_name = project_name
        self.language = language
        self.row_limit = row_limit

    def generate(self):
        logger.info("Starting PPTX generation...")
        logger.info("Template: %s", self.template_path)
        logger.info("Markdown: %s", self.markdown_file)
        logger.info("Output: %s", self.output_path)
        logger.info("Language: %s", self.language)
        logger.info("Project: %s", self.project_name)

        with open(self.markdown_file, "r", encoding="utf-8") as handle:
            md_content = handle.read()

        logger.info("Content length: %s characters", len(md_content))
        generator = PowerPointGenerator(self.template_path, self.language, self.row_limit)

        try:
            generator.generate_full_report(md_content, None, self.output_path)
            if self.project_name:
                generator.update_project_titles(self.project_name, "BS")
        except Exception as exc:
            logger.error("Report generation failed: %s", exc)
            raise

        logger.info("PPTX generation completed: %s", self.output_path)


def export_pptx(
    template_path: str,
    markdown_path: str,
    output_path: str,
    project_name: Optional[str] = None,
    _excel_file_path: Optional[str] = None,
    language: str = "english",
    statement_type: str = "BS",
    row_limit: int = 20,
    model_type: Optional[str] = None,
):
    generator = ReportGenerator(template_path, markdown_path, output_path, project_name, language, row_limit)
    generator.generate()

    if not os.path.exists(output_path):
        raise FileNotFoundError(f"PPTX file was not created at {output_path}")

    if project_name:
        temp_presentation = Presentation(output_path)
        pptx_gen = PowerPointGenerator(template_path, language, row_limit, model_type=model_type)
        pptx_gen.presentation = temp_presentation
        pptx_gen.update_project_titles(project_name, statement_type)
        temp_presentation.save(output_path)

    logger.info("PowerPoint presentation successfully exported to: %s", output_path)
    return output_path


def export_pptx_from_structured_data_combined(
    template_path: str,
    bs_data: List[Dict],
    is_data: List[Dict],
    output_path: str,
    project_name: Optional[str] = None,
    language: str = "english",
    temp_path: Optional[str] = None,
    selected_sheet: Optional[str] = None,
    is_chinese_databook: bool = False,
    bs_is_results: Optional[Dict[str, Any]] = None,
    model_type: Optional[str] = None,
    model_name: Optional[str] = None,
    skip_summary_ai: bool = False,  # AI summary needed for coSummaryShape; parallelized at max_workers=4
    pre_generated_summaries: Optional[Dict[str, str]] = None,  # {"BS": str, "IS": str} — bypass AI in PPTX export
    mappings: Optional[Dict[str, Any]] = None,  # for translating the embedded BS/IS table's row labels when Chinese
):
    try:
        export_started_at = time.perf_counter()
        def _stage_log(msg: str) -> None:
            logger.info(msg)

        _stage_log(f"Starting export | BS={len(bs_data)} IS={len(is_data)} skip_summary_ai={skip_summary_ai}")

        generator = PowerPointGenerator(template_path, language, row_limit=20, model_type=model_type, model_name=model_name)
        if skip_summary_ai:
            generator.pptx_settings.setdefault("executive_summary", {})["enable_ai"] = False
        stage_started_at = time.perf_counter()
        generator.load_template()
        _stage_log(f"load_template: {time.perf_counter() - stage_started_at:.2f}s")

        pre_summaries = pre_generated_summaries or {}
        if bs_data:
            stage_started_at = time.perf_counter()
            generator.apply_structured_data_to_slides(
                bs_data, 1, project_name, "BS",
                is_chinese_databook=is_chinese_databook,
                pre_generated_summary=pre_summaries.get("BS"),
            )
            _stage_log(f"apply_bs_slides: {time.perf_counter() - stage_started_at:.2f}s")
        if is_data:
            stage_started_at = time.perf_counter()
            generator.apply_structured_data_to_slides(
                is_data, 5, project_name, "IS",
                is_chinese_databook=is_chinese_databook,
                pre_generated_summary=pre_summaries.get("IS"),
            )
            _stage_log(f"apply_is_slides: {time.perf_counter() - stage_started_at:.2f}s")
        # bs_is_results being already-computed is sufficient on its own --
        # requiring selected_sheet too silently skipped the embedded table
        # whenever the caller had no sheet name to give (roll-up-sourced
        # financials with a blank own-file sheet, or a synthesized BS/IS
        # built purely from schedule tabs with no Financials sheet at all)
        # even though there was real BS/IS data ready to embed.
        if temp_path and (selected_sheet or bs_is_results):
            stage_started_at = time.perf_counter()
            generator.embed_financial_tables(
                temp_path,
                selected_sheet,
                project_name,
                language,
                bs_is_results=bs_is_results,
                mappings=mappings,
            )
            _stage_log(f"embed_financial_tables: {time.perf_counter() - stage_started_at:.2f}s")
        if hasattr(generator, "_unused_slides_to_remove") and generator._unused_slides_to_remove:
            stage_started_at = time.perf_counter()
            unused_slides_sorted = sorted(set(generator._unused_slides_to_remove), reverse=True)
            generator._remove_slides(unused_slides_sorted)
            _stage_log(f"remove_unused_slides ({len(unused_slides_sorted)}): {time.perf_counter() - stage_started_at:.2f}s")
        if project_name:
            stage_started_at = time.perf_counter()
            generator.refresh_project_placeholders(project_name)
            _stage_log(f"refresh_project_placeholders: {time.perf_counter() - stage_started_at:.2f}s")

        stage_started_at = time.perf_counter()
        generator.save(output_path)
        _stage_log(f"save_presentation: {time.perf_counter() - stage_started_at:.2f}s")
        _stage_log(f"TOTAL export: {time.perf_counter() - export_started_at:.2f}s")
        logger.info("Combined PPTX generation completed: %s", output_path)
        return output_path
    except Exception as exc:
        logger.error("PPTX generation failed: %s", exc)
        logger.error(traceback.format_exc())
        raise


def export_pptx_from_structured_data(
    template_path: str,
    structured_data: List[Dict],
    output_path: str,
    project_name: Optional[str] = None,
    language: str = "english",
    statement_type: str = "BS",
    start_slide: int = 1,
    model_type: Optional[str] = None,
):
    try:
        logger.info("Starting PPTX generation from structured data...")
        logger.info("Template: %s", template_path)
        logger.info("Output: %s", output_path)
        logger.info("Language: %s", language)
        logger.info("Statement type: %s, Start slide: %s", statement_type, start_slide)
        logger.info("Accounts to process: %s", len(structured_data))

        generator = PowerPointGenerator(template_path, language, row_limit=20, model_type=model_type)
        generator.load_template()
        generator.apply_structured_data_to_slides(structured_data, start_slide, project_name, statement_type)
        generator.save(output_path)

        logger.info("PPTX generation completed: %s", output_path)
        return output_path
    except Exception as exc:
        logger.error("PPTX generation failed: %s", exc)
        raise


def merge_presentations(bs_presentation_path: str, is_presentation_path: str, output_path: str):
    try:
        logger.info("🔄 Starting presentation merge...")
        logger.info("   BS: %s", bs_presentation_path)
        logger.info("   IS: %s", is_presentation_path)

        merged_prs = Presentation(bs_presentation_path)
        is_prs = Presentation(is_presentation_path)

        from copy import deepcopy

        for slide_idx, slide in enumerate(is_prs.slides):
            try:
                slide_layout = slide.slide_layout
                new_slide = merged_prs.slides.add_slide(slide_layout)

                source_slide_xml = slide._element
                target_slide_xml = new_slide._element

                shapes_to_remove = list(new_slide.shapes)
                for shape in shapes_to_remove:
                    try:
                        sp_tree = target_slide_xml.get_or_add_spTree()
                        sp_tree.remove(shape._element)
                    except Exception:
                        pass

                source_sp_tree = source_slide_xml.get_or_add_spTree()
                target_sp_tree = target_slide_xml.get_or_add_spTree()
                for shape_element in source_sp_tree:
                    target_sp_tree.append(deepcopy(shape_element))

            except Exception as exc:
                logger.error("Error copying slide %s, using fallback method: %s", slide_idx, exc)
                slide_layout = slide.slide_layout
                new_slide = merged_prs.slides.add_slide(slide_layout)
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for new_shape in new_slide.shapes:
                            if (
                                hasattr(new_shape, "name")
                                and hasattr(shape, "name")
                                and new_shape.name == shape.name
                                and new_shape.has_text_frame
                            ):
                                new_shape.text_frame.text = shape.text_frame.text
                                break

        merged_prs.save(output_path)
        del merged_prs
        del is_prs

        import gc

        gc.collect()
        logger.info("✅ Presentation merge completed successfully")
    except Exception as exc:
        logger.error("Presentation merge failed: %s", exc)
        raise


def _dedupe_part_name(dest_prs: "Presentation", target_part, renamed_part_ids: set) -> None:
    """Rename `target_part` in-place if its partname collides with a part
    already present in dest_prs's package.

    python-pptx's Package.save() writes every part reachable from the
    package's own relationship graph using each Part object's OWN
    `.partname` -- it never re-derives a name. When _copy_slide_into()
    relates a destination slide directly to a Part object still owned by a
    DIFFERENT source Presentation (e.g. a picture's blipFill target), that
    part keeps the partname it was assigned in ITS OWN package (e.g.
    "/ppt/media/image3.png"). Since every batch entity's deck is built by
    the same export code, two different source decks landing on the same
    numbered partname is common, not a corner case -- and when that
    happens, the combined package ends up with two different parts both
    claiming "/ppt/media/image3.png", which produces a zip with a
    duplicate member name: invalid OPC, which is exactly what makes
    PowerPoint prompt "repair this presentation" (the media is
    unrecoverable/misattributed, not merely cosmetically wrong).
    Renaming the incoming part to a partname that's actually free in the
    destination package's namespace (via next_partname, the same
    mechanism python-pptx itself uses when adding new parts) avoids the
    collision. Only checked once per distinct source Part object
    (tracked by id() in `renamed_part_ids`, shared across an entire
    combine_presentations() call) -- once resolved, a part's identity/
    partname pairing is stable for the rest of the run.
    """
    if id(target_part) in renamed_part_ids:
        return
    renamed_part_ids.add(id(target_part))
    existing_partnames = {p.partname for p in dest_prs.part.package.iter_parts()}
    if target_part.partname not in existing_partnames:
        return
    partname = target_part.partname
    name_part = re.sub(r"\d+$", "", posixpath.splitext(partname.filename)[0]) or "part"
    tmpl = posixpath.join(partname.baseURI, f"{name_part}%d.{partname.ext}") if partname.ext else posixpath.join(partname.baseURI, f"{name_part}%d")
    target_part.partname = dest_prs.part.package.next_partname(tmpl)


def _copy_slide_into(dest_prs: "Presentation", source_slide, renamed_part_ids: Optional[set] = None) -> None:
    """Deep-copy one slide from a DIFFERENT Presentation (built from the
    same template.pptx) onto the end of dest_prs, preserving every shape
    including images and native tables.

    python-pptx has no built-in "append an existing slide" API, so this
    clones the slide's shape-tree XML directly -- the same technique
    merge_presentations() above uses. The one thing that technique is
    missing (and why it's not reused as-is here): the copied XML still
    references relationship IDs (r:embed / r:id / r:link, used by
    pictures and hyperlinks) that only exist in the SOURCE file's part.
    Left unmapped, those would point at nothing in the destination part --
    copied images would come through as silently broken/missing rather
    than raising an error. Every non-slideLayout relationship the source
    slide owns is re-created on the destination slide's own part first,
    and every r:embed/r:id/r:link attribute in the copied XML is
    rewritten to the new relationship id.

    Embedded/linked OLE objects (MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT /
    LINKED_OLE_OBJECT -- e.g. a "TCLayout.ActiveDocument.1" marker some
    add-ins like ThinkCell/UpSlide leave on every slide) are deliberately
    SKIPPED entirely, not copied or relationship-remapped: a real batch
    combine produced blank/whited-out pages specifically where these
    existed, and this codebase has no template with such an object to
    debug the exact OLE relationship mechanics against locally. These
    markers are consistently 0.001in x 0.001in (invisible, carry no
    reader-facing content) in every template seen so far, so dropping them
    trades an add-in bookkeeping artifact for guaranteed-correct visible
    content -- the safer side of that tradeoff.
    """
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    renamed_part_ids = renamed_part_ids if renamed_part_ids is not None else set()

    layout_name = source_slide.slide_layout.name
    dest_layout = next(
        (layout for layout in dest_prs.slide_layouts if layout.name == layout_name),
        dest_prs.slide_layouts[0],
    )
    dest_slide = dest_prs.slides.add_slide(dest_layout)

    # The layout auto-populates placeholder shapes -- clear them, the
    # source slide's own shape tree (copied below) already carries
    # everything that should be on the page.
    for shape in list(dest_slide.shapes):
        shape._element.getparent().remove(shape._element)

    r_ns = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    ole_shape_elements = set()
    ole_rel_ids = set()
    for shape in source_slide.shapes:
        try:
            is_ole = shape.shape_type in (MSO_SHAPE_TYPE.EMBEDDED_OLE_OBJECT, MSO_SHAPE_TYPE.LINKED_OLE_OBJECT)
        except (ValueError, NotImplementedError):
            is_ole = False
        if is_ole:
            ole_shape_elements.add(shape._element)
            for el in shape._element.iter():
                for attr_name in ("embed", "link", "id"):
                    rid = el.get(f"{r_ns}{attr_name}")
                    if rid:
                        ole_rel_ids.add(rid)

    rel_id_map: Dict[str, str] = {}
    for rel_id, rel in source_slide.part.rels.items():
        if rel.reltype.endswith("/slideLayout") or rel_id in ole_rel_ids:
            continue  # layout relationship isn't copied; OLE ones are deliberately dropped
        if rel.is_external:
            new_rel_id = dest_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            _dedupe_part_name(dest_prs, rel.target_part, renamed_part_ids)
            new_rel_id = dest_slide.part.relate_to(rel.target_part, rel.reltype)
        rel_id_map[rel_id] = new_rel_id

    for shape_elm in list(source_slide.shapes._spTree):
        if shape_elm.tag in (qn("p:nvGrpSpPr"), qn("p:grpSpPr")):
            continue  # spTree's two fixed non-shape children, not content
        if shape_elm in ole_shape_elements:
            continue
        new_elm = copy.deepcopy(shape_elm)
        for el in new_elm.iter():
            for attr_name in ("embed", "link", "id"):
                old_rid = el.get(f"{r_ns}{attr_name}")
                if old_rid and old_rid in rel_id_map:
                    el.set(f"{r_ns}{attr_name}", rel_id_map[old_rid])
        dest_slide.shapes._spTree.append(new_elm)


def combine_presentations(pptx_sources: List, output_path) -> "str | None":
    """Combine several already-exported .pptx decks (e.g. one per batch
    entity, all built from the same template.pptx) into a single deck --
    every slide from every source, in order, via _copy_slide_into().

    pptx_sources: file paths (str) and/or file-like objects (e.g.
    io.BytesIO of already-in-memory PPTX bytes -- python-pptx's own
    Presentation() constructor accepts either, so no temp files are needed
    when combining straight from a batch run's cached pptx_download_data).
    output_path: a path (str) to save to, OR a file-like object (e.g.
    io.BytesIO) to write into instead of touching disk -- returns the path
    string in the former case, None in the latter (caller already holds
    the buffer it passed in).

    Deliberately NOT a general-purpose "merge any two PPTX files" utility:
    it assumes every input shares the same template (true for every batch
    entity, since they all come from export_pptx_from_structured_data_combined
    with the same template_path), which is what makes layout-name matching
    a safe way to pick the destination layout for each copied slide.
    """
    if not pptx_sources:
        raise ValueError("combine_presentations requires at least one input source")

    combined_prs = Presentation(pptx_sources[0])
    renamed_part_ids: set = set()
    for source in pptx_sources[1:]:
        source_prs = Presentation(source)
        for source_slide in source_prs.slides:
            _copy_slide_into(combined_prs, source_slide, renamed_part_ids)

    combined_prs.save(output_path)
    logger.info("Combined %s presentation(s)", len(pptx_sources))
    if isinstance(output_path, str):
        return output_path
    return None
# --- end pptx/exporters.py ---

# --- begin pptx/generation.py ---
"""
PowerPoint Generation Module for Financial Reports
Based on the backup methods but implemented fresh for the new system
"""

import os
import re
import logging
import threading
import time
from typing import Any, Dict, List, Optional, Tuple
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
logger.setLevel(logging.WARNING)


class PowerPointGenerator:
    """Main PowerPoint generator class"""

    def __init__(
        self,
        template_path: str,
        language: str = 'english',
        row_limit: int = 20,
        model_type: Optional[str] = None,
        model_name: Optional[str] = None,
    ):
        self.template_path = template_path
        self.language = language.lower()
        self.row_limit = row_limit
        self.model_type = str(model_type or "").strip() or None
        self.model_name = str(model_name or "").strip() or None
        self.presentation = None
        self.pptx_settings = _load_pptx_settings()

    def load_template(self):
        """Load the PowerPoint template"""
        if not os.path.exists(self.template_path):
            raise FileNotFoundError(f"Template not found: {self.template_path}")

        self.presentation = Presentation(self.template_path)
        logger.info("Loaded template: %s", self.template_path)

    def find_shape_by_name(self, shapes, name: str):
        """Find shape by name in slide (case-insensitive), recursive"""
        name_lower = name.lower()
        for shape in shapes:
            if hasattr(shape, 'name') and (shape.name == name or shape.name.lower() == name_lower):
                return shape
            
            # Check for group
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                found = self.find_shape_by_name(shape.shapes, name)
                if found:
                    return found
        return None

    @staticmethod
    def _is_commentary_text_shape(shape) -> bool:
        if not getattr(shape, "has_text_frame", False):
            return False
        shape_name = str(getattr(shape, "name", "") or "").lower()
        excluded_tokens = (
            "title",
            "projtitle",
            "summary",
            "cosummaryshape",
            "table",
            "subtitle",
        )
        return not any(token in shape_name for token in excluded_tokens)

    def _split_single_into_lr(self, slide, source_shape):
        """Clone a full-width commentary box into two half-width boxes side by
        side (named textMainBullets_L / textMainBullets_R), replacing the
        original. Used when a page was assigned two logical slots (L, R) but
        the underlying template slide only has ONE commentary text box —
        without this, both slots resolve to the SAME shape (see
        _resolve_commentary_slot_shape), so L-content and R-content silently
        collide into one full-width box instead of sitting side by side like
        the BS pages. Idempotent per slide (cached by slide element id)."""
        cache = getattr(self, "_split_lr_cache", None)
        if cache is None:
            cache = self._split_lr_cache = {}
        slide_key = id(slide._element)
        cached = cache.get(slide_key)
        if cached:
            return cached

        from copy import deepcopy
        orig_left = int(source_shape.left)
        orig_width = int(source_shape.width)
        gutter = max(0, int(orig_width * 0.03))
        half_width = max(1, (orig_width - gutter) // 2)

        # Left half = the original shape, resized in place.
        source_shape.left = orig_left
        source_shape.width = half_width
        try:
            source_shape.name = "textMainBullets_L"
        except Exception:
            pass

        # Right half = a deep-copied XML clone, repositioned and renamed. The
        # clone starts with a copy of the original's text — clear it so the
        # packer fills it fresh rather than duplicating whatever was there.
        new_element = deepcopy(source_shape._element)
        source_shape._element.addnext(new_element)
        right_shape = None
        for shape in slide.shapes:
            if shape._element is new_element:
                right_shape = shape
                break
        if right_shape is None:
            # Clone insertion failed for some reason — fall back to treating
            # the (now half-width) original as both slots rather than crashing.
            cache[slide_key] = (source_shape, source_shape)
            return cache[slide_key]

        right_shape.left = orig_left + half_width + gutter
        right_shape.width = half_width
        try:
            right_shape.name = "textMainBullets_R"
        except Exception:
            pass
        if getattr(right_shape, "has_text_frame", False):
            right_shape.text_frame.clear()

        cache[slide_key] = (source_shape, right_shape)
        logger.info(
            "Split single full-width commentary box into L/R halves on a slide "
            "(page was assigned two content slots but the template only had one box)."
        )
        return cache[slide_key]

    def _resolve_commentary_slot_shape(self, slide, slot_name: str, used_shape_ids=None):
        """Resolve the best text box for a commentary slot on a slide."""
        used_shape_ids = used_shape_ids or set()
        preferred_names = {
            "single": [
                "textMainBullets",
                "textMainBullets_L",
                "textMainBullets_R",
                "Text-commentary",
                "Content",
                "MainContent",
                "Body",
            ],
            "L": [
                "textMainBullets_L",
                "Text-commentary",
                "Content",
                "MainContent",
                "Body",
            ],
            "R": [
                "textMainBullets_R",
                "Text-commentary",
                "Content",
                "MainContent",
                "Body",
            ],
        }

        for name in preferred_names.get(slot_name, []):
            shape = self.find_shape_by_name(slide.shapes, name)
            if shape and getattr(shape, "has_text_frame", False) and id(shape) not in used_shape_ids:
                return shape

        # No dedicated _L/_R box. Only fall back to a single full-width box for
        # an ACTUAL single-column slot; for "L"/"R" that would make both slots
        # resolve to the same physical shape (content collision -> renders as
        # one full-width box instead of two side-by-side columns). Split it
        # into two half-width boxes instead, mirroring the BS page layout.
        if slot_name in ("L", "R"):
            single_shape = self.find_shape_by_name(slide.shapes, "textMainBullets")
            if single_shape and getattr(single_shape, "has_text_frame", False) and id(single_shape) not in used_shape_ids:
                left_shape, right_shape = self._split_single_into_lr(slide, single_shape)
                return left_shape if slot_name == "L" else right_shape

        generic_candidates = [
            shape for shape in slide.shapes
            if self._is_commentary_text_shape(shape) and id(shape) not in used_shape_ids
        ]
        if not generic_candidates:
            return None

        if slot_name == "L":
            return min(generic_candidates, key=lambda shape: (getattr(shape, "left", 0), -getattr(shape, "width", 0)))
        if slot_name == "R":
            return max(generic_candidates, key=lambda shape: (getattr(shape, "left", 0), getattr(shape, "width", 0)))
        return max(generic_candidates, key=lambda shape: (getattr(shape, "width", 0), -getattr(shape, "left", 0)))

    def _add_commentary_slot_shape(self, slide, slot_name: str):
        top = Inches(2.22)
        width = Inches(4.78)
        height = Inches(4.13)
        if slot_name == "L":
            left = Inches(0.13)
        elif slot_name == "R":
            left = Inches(5.09)
        else:
            # Page 1 template uses a single commentary box on the right beside the table.
            left = Inches(5.09)
        return slide.shapes.add_textbox(left, top, width, height)

    def _summary_settings(self) -> Dict[str, Any]:
        return dict(self.pptx_settings.get("executive_summary") or {})

    def _summary_length_targets(self, is_chinese: bool) -> Tuple[int, int]:
        """(target_chars_chi, target_words_eng) sized to the REAL summary box
        rather than to a static config number.

        The static values had to be re-tuned by hand for every template and
        were consistently short -- a real deck filled about 2 of the 4 lines
        the box can hold. Measuring the actual `coSummaryShape` width and
        multiplying by `target_lines` self-corrects, and keeps the AI's own
        length instruction honest about the space it is writing into.

        Falls back to the configured values whenever the shape can't be
        resolved or measured."""
        settings = self._summary_settings()
        static_chi = int(settings.get("target_chars_chi", 144))
        static_eng = int(settings.get("target_words_eng", 110))
        target_lines = max(1, int(settings.get("target_lines", 4) or 4))
        try:
            shape = None
            for slide in self.presentation.slides:
                shape = self.find_shape_by_name(slide.shapes, "coSummaryShape")
                if shape is not None:
                    break
            if shape is None:
                return static_chi, static_eng
            from fdd_utils.text_metrics import get_measurer, text_box_from_shape
            box = text_box_from_shape(shape)
            packing = self._packing_settings()
            font_pt = self._real_font_size_pt(is_chinese)
            measurer = get_measurer(
                self._measurer_family(is_chinese, packing), font_pt, is_cjk=is_chinese,
                line_spacing=self._real_line_spacing(is_chinese),
                metrics_path=self._resolve_font_metrics_path(is_chinese, packing),
            )
            probe = ("财务尽职调查评述示例文字内容测算每行可容纳字数上限" * 12) if is_chinese else (
                "financial due diligence commentary sample text measuring how many words fit " * 6)
            wrapped = measurer.wrap(probe, box.width_pt)
            if not wrapped:
                return static_chi, static_eng
            per_line = max(1.0, len(probe) / len(wrapped))
            # How many lines the box HOLDS, measured -- not the configured
            # guess. The width here has always been measured while the line
            # count was a static 4, which is wrong in both directions on a
            # per-machine template: this Mac's band holds 2.99 lines (so 4
            # over-fills it), while the band that prompted "exe sum 非常短"
            # is evidently taller. target_lines survives only as the fallback
            # for when the shape can't be measured, which is exactly the role
            # the static width fallback already plays.
            from fdd_utils.text_metrics import POWERPOINT_LINE_PITCH_FACTOR
            pitch = font_pt * POWERPOINT_LINE_PITCH_FACTOR * self._real_line_spacing(is_chinese)
            fitting_lines = (box.height_pt / pitch) if pitch > 0 else 0
            lines = fitting_lines if fitting_lines >= 1 else target_lines
            # 0.92: aim just under a full box so a slightly long sentence
            # doesn't push a whole extra line past the table beneath it.
            chars = int(per_line * lines * 0.92)
            logger.info(
                "Executive summary target: %d chars (band %.0fx%.0fpt holds %.2f lines "
                "at %.1fpt, ~%.0f chars/line)",
                chars, box.width_pt, box.height_pt, fitting_lines, font_pt, per_line,
            )
            # No max(static, ...) floor once BOTH dimensions are measured.
            # The floor existed because the line count was a guess; raising a
            # measured 127-char band to a static 144 just overflows it, which
            # is the same mistake as shrinking a box to absorb 0.3 lines. If
            # the resulting summary is too short for the reader's taste, the
            # band itself is too small -- the log line above says by how much,
            # so that is now a decision that can be made on a real number.
            if fitting_lines >= 1:
                return (chars, static_eng) if is_chinese else (static_chi, int(chars / 6))
            if is_chinese:
                return max(static_chi, chars), static_eng
            # English targets are in WORDS; ~6 chars per word including space.
            return static_chi, max(static_eng, int(chars / 6))
        except Exception as exc:
            logger.debug("Could not size summary targets from the real shape: %s", exc)
            return static_chi, static_eng

    def _packing_settings(self, statement_type: Optional[str] = None) -> Dict[str, Any]:
        packing = dict(self.pptx_settings.get("commentary_packing") or {})
        if not statement_type:
            return packing
        overrides = ((packing.get("statement_overrides") or {}).get(statement_type) or {})
        if not overrides:
            return packing
        return _merge_nested_dict(packing, overrides)

    # How far an account may protrude BELOW its box, in std_lh line-units,
    # when the only alternative is splitting it across slots. The project
    # team explicitly accepts 1-2 lines sticking out; they do not accept a
    # split landing mid-name. Applied only at the split decision, never as a
    # general capacity increase -- every slot is still packed to its real
    # capacity, and a slot that simply has more content than fits still
    # splits as before.
    _TAIL_OVERFLOW_TOLERANCE_UNITS = 2.0

    def _tail_overflow_tolerance_units(self, statement_type: Optional[str] = None) -> float:
        raw = self._packing_settings(statement_type).get(
            "tail_overflow_tolerance_lines", self._TAIL_OVERFLOW_TOLERANCE_UNITS
        )
        try:
            return max(0.0, float(raw))
        except (TypeError, ValueError):
            return self._TAIL_OVERFLOW_TOLERANCE_UNITS

    def _resolve_summary_model_type(self, is_chinese: bool) -> str:
        cached = getattr(self, "_summary_model_type_cache", None)
        if cached:
            return cached

        if self.model_type:
            resolved_model_type = str(self.model_type).strip()
            self._summary_model_type_cache = resolved_model_type
            return resolved_model_type

        try:
            from fdd_utils.ai import FDDConfig

            config = load_yaml_file(package_file_path("config.yml")) or {}
            requested_model_type = (
                str(self.model_type).strip()
                if self.model_type
                else str(((config.get("default") or {}).get("ai_provider")) or "deepseek")
            )
            config_manager = FDDConfig(
                language="Chi" if is_chinese else "Eng",
                model_type=requested_model_type,
            )
            resolved_model_type = str(config_manager.model_type or requested_model_type)
        except Exception as exc:
            logger.warning("Could not resolve PPTX summary model type, defaulting to deepseek: %s", exc)
            resolved_model_type = "deepseek"

        self._summary_model_type_cache = resolved_model_type
        return resolved_model_type

    def _summary_max_workers(self, summary_jobs: List[Dict[str, Any]]) -> int:
        if not summary_jobs:
            return 1

        summary_settings = self._summary_settings()
        configured_workers = int(summary_settings.get("max_workers", 4) or 4)
        model_type = self._resolve_summary_model_type(bool(summary_jobs[0].get("is_chinese")))
        if model_type == "local":
            configured_workers = int(summary_settings.get("local_max_workers", 1) or 1)
        return max(1, min(configured_workers, len(summary_jobs)))

    def _call_with_timeout(self, func, timeout_seconds: float, timeout_label: str):
        if timeout_seconds <= 0:
            return func()

        result_container = {"value": None, "error": None, "completed": False}

        def _run():
            try:
                result_container["value"] = func()
            except Exception as exc:
                result_container["error"] = exc
            finally:
                result_container["completed"] = True

        worker = threading.Thread(target=_run, daemon=True)
        worker.start()
        worker.join(timeout=timeout_seconds)

        if not result_container["completed"]:
            raise TimeoutError(f"{timeout_label} timed out after {timeout_seconds:.1f} seconds")
        if result_container["error"] is not None:
            raise result_container["error"]
        return result_container["value"]

    def _call_with_timeout_retry(
        self,
        func,
        timeout_seconds: float,
        max_retries: int,
        timeout_label: str,
    ):
        """Call ``func`` with a per-attempt timeout and retry on TimeoutError
        or other transient failures. Raises the last exception if all retries
        fail. Use ``max_retries >= 1`` (1 means "no retry, just run once")."""
        attempts = max(1, int(max_retries or 1))
        last_error: Optional[BaseException] = None
        for attempt in range(1, attempts + 1):
            label = (
                timeout_label
                if attempts == 1
                else f"{timeout_label} (attempt {attempt}/{attempts})"
            )
            try:
                return self._call_with_timeout(func, timeout_seconds, label)
            except TimeoutError as te:
                last_error = te
                logger.warning(
                    "%s timed out after %.1fs; %s",
                    label,
                    timeout_seconds,
                    "retrying" if attempt < attempts else "giving up",
                )
            except Exception as exc:
                last_error = exc
                logger.warning(
                    "%s errored (%s); %s",
                    label,
                    exc,
                    "retrying" if attempt < attempts else "giving up",
                )
        assert last_error is not None
        raise last_error

    def _generate_slide_summaries(self, summary_jobs: List[Dict[str, Any]]) -> Dict[int, str]:
        if not summary_jobs:
            return {}

        max_workers = self._summary_max_workers(summary_jobs)
        model_type = self._resolve_summary_model_type(bool(summary_jobs[0].get("is_chinese")))
        jobs_by_slide = {job["slide_idx"]: job for job in summary_jobs}
        results: Dict[int, str] = {}

        logger.info(
            "Generating %s PPTX slide summaries with model_type=%s, max_workers=%s",
            len(summary_jobs),
            model_type,
            max_workers,
        )

        def _generate_summary(job: Dict[str, Any]) -> str:
            slide_number = int(job["slide_idx"]) + 1
            summary_started_at = time.perf_counter()
            ai_summary = self._generate_ai_summary(
                job["page_commentary"] or job["page_summary_source"],
                job["page_summary_source"],
                job["is_chinese"],
            )
            if ai_summary:
                logger.info(
                    "PPTX summary slide %s completed via AI in %.2fs",
                    slide_number,
                    time.perf_counter() - summary_started_at,
                )
                return ai_summary
            fallback_summary = self._generate_page_summary(job["page_summary_source"], job["is_chinese"])
            logger.info(
                "PPTX summary slide %s completed via fallback in %.2fs",
                slide_number,
                time.perf_counter() - summary_started_at,
            )
            return fallback_summary

        if max_workers == 1:
            for slide_idx, job in jobs_by_slide.items():
                results[slide_idx] = _generate_summary(job)
            return results

        with ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_slide = {
                executor.submit(_generate_summary, job): slide_idx
                for slide_idx, job in jobs_by_slide.items()
            }
            for future in as_completed(future_to_slide):
                slide_idx = future_to_slide[future]
                job = jobs_by_slide[slide_idx]
                try:
                    results[slide_idx] = future.result()
                except Exception as exc:
                    logger.warning(
                        "Slide %s summary generation failed, using fallback summary: %s",
                        slide_idx + 1,
                        exc,
                    )
                    results[slide_idx] = self._generate_page_summary(
                        job["page_summary_source"],
                        job["is_chinese"],
                    )

        return results

    def _prepare_commentary_blocks(self, commentary: str) -> List[str]:
        normalized = str(commentary or "").replace("\r\n", "\n").strip()
        if not normalized:
            return []

        blocks: List[str] = []
        for raw_block in re.split(r"\n\s*\n", normalized):
            lines = [line.strip() for line in raw_block.split("\n") if line.strip()]
            if not lines:
                continue
            if len(lines) == 1:
                blocks.append(lines[0])
                continue

            rebuilt: List[str] = []
            current = ""
            for line in lines:
                is_bullet_like = bool(re.match(r"^([-*•]|\d+[.)])\s+", line))
                if is_bullet_like:
                    if current:
                        rebuilt.append(current.strip())
                        current = ""
                    rebuilt.append(line)
                    continue
                current = f"{current} {line}".strip() if current else line
            if current:
                rebuilt.append(current.strip())
            blocks.extend(rebuilt)
        return blocks

    def _compact_commentary_for_ppt(self, commentary: str, is_chinese: bool) -> str:
        normalized = _normalize_slide_commentary_text(commentary)
        if not normalized:
            return ""

        packing = self._packing_settings()
        min_chars = int(
            packing.get("ppt_min_chars_chi" if is_chinese else "ppt_min_chars_eng", 110 if is_chinese else 190)
        )
        if len(normalized) <= min_chars:
            return normalized

        target_ratio = float(packing.get("ppt_length_ratio", 0.72) or 0.72)
        target_chars = max(min_chars, int(len(normalized) * target_ratio))
        max_sentences = int(
            packing.get("ppt_max_sentences_chi" if is_chinese else "ppt_max_sentences_eng", 3)
        )
        max_numeric_sentences = int(packing.get("ppt_max_numeric_sentences", 2) or 2)

        compact = _build_compact_summary_text(
            normalized,
            is_chinese=is_chinese,
            max_sentences=max_sentences,
            max_chars=target_chars,
            max_numeric_sentences=max_numeric_sentences,
        )
        compact = _normalize_slide_commentary_text(compact)
        if not compact:
            return normalized

        minimum_retained_chars = max(90 if is_chinese else 140, int(len(normalized) * 0.35))
        if len(compact) < minimum_retained_chars:
            return normalized
        return compact if len(compact) < len(normalized) else normalized

    def _prepare_structured_data_for_slides(self, structured_data: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        prepared: List[Dict[str, Any]] = []
        for account_data in structured_data or []:
            item = dict(account_data or {})
            commentary = _normalize_slide_commentary_text(item.get("commentary", ""))
            if commentary:
                item["original_commentary"] = commentary
            item["commentary"] = commentary  # Keep full length; fill optimizer handles fit
            prepared.append(item)
        return prepared

    # Average rendered character width (pt) for the fonts we use.
    # English: Arial 9pt mixed text ≈ 5.0 pt/char (incl. spaces & punctuation).
    # Chinese: YaHei 10pt CJK characters are square — 1 em ≈ 10 pt/char.
    # A small word-wrap slack (≈8 %) is subtracted because lines always break
    # at a word/character boundary, not at the exact pixel edge.
    _AVG_CHAR_WIDTH_ENG = 5.0
    _AVG_CHAR_WIDTH_CHI = 10.0
    _WORD_WRAP_SLACK    = 0.92   # use 92 % of the theoretical line width

    def _estimate_chars_per_line(
        self,
        slot_name: str,
        is_chinese: bool,
        shape=None,
        *,
        statement_type: Optional[str] = None,
    ) -> int:
        """Return the number of characters that fit on one line.

        When the actual shape is available we measure directly from its width
        and the text-frame insets, then divide by the known average character
        width for the font in use.  This removes all dependency on the
        ``chars_per_line`` config block for shapes that exist in the template.

        Falls back to the config-based estimate only when no shape is supplied.
        """
        if shape is not None and hasattr(shape, "width"):
            width_pt = shape.width * 72 / 914400
            # Read actual text-frame left/right insets; default is 0.1 in = 7.2 pt.
            left_pt = right_pt = 7.2
            try:
                tf = shape.text_frame
                if tf.margin_left is not None:
                    left_pt  = tf.margin_left  * 72 / 914400
                if tf.margin_right is not None:
                    right_pt = tf.margin_right * 72 / 914400
            except Exception:
                pass
            effective_pt = max(10.0, width_pt - left_pt - right_pt)
            avg_char = self._AVG_CHAR_WIDTH_CHI if is_chinese else self._AVG_CHAR_WIDTH_ENG
            return max(16, int(effective_pt * self._WORD_WRAP_SLACK / avg_char))

        # No shape — fall back to config values.
        packing = self._packing_settings(statement_type)
        chars_per_line = packing.get("chars_per_line") or {}
        slot_key = slot_name if slot_name in {"single", "L", "R"} else "default"
        language_key = "chi" if is_chinese else "eng"
        base_value = (
            ((chars_per_line.get(slot_key) or {}).get(language_key))
            or ((chars_per_line.get("default") or {}).get(language_key))
            or (32 if is_chinese else 60)
        )
        return int(base_value)

    def _build_page_summary_source(self, slide_accounts: List[Dict]) -> Tuple[str, str]:
        """Build the exact slide commentary set used for summary generation.

        The summary source is each account's LEAD-IN only, never its table
        detail. strip_table_detail_for_summary was written for exactly this
        and its docstring records the symptom it prevents -- the naive
        first-sentence splitter in _generate_page_summary doesn't know
        "明细如下：" isn't a sentence end, so a table account's raw commentary
        yields a "first sentence" that runs through the handoff phrase and
        swallows the whole first "➢" bullet, splicing table detail into the
        executive summary and cutting off mid-sentence. Both ui.py call
        sites already strip; the export path never did, so the deterministic
        in-export summary reintroduced the same defect the moment it started
        filling the band rather than leaving it blank.

        page_commentary stays UNSTRIPPED -- it is the AI path's own input and
        should still see everything.
        """
        commentary_parts = []
        summary_source_parts = []
        summary_parts = []

        for account_data in slide_accounts or []:
            commentary = str(account_data.get("commentary", "") or "").strip()
            summary = str(account_data.get("summary", "") or "").strip()
            if commentary:
                commentary_parts.append(commentary)
                lead_in = self.strip_table_detail_for_summary(
                    commentary, self._account_is_chinese(account_data),
                )
                summary_source_parts.append(lead_in or commentary)
            if summary:
                summary_parts.append(summary)

        page_commentary = "\n\n".join(commentary_parts).strip()
        page_summary_source = (
            "\n\n".join(summary_source_parts).strip() or " ".join(summary_parts).strip()
        )
        return page_commentary, page_summary_source

    @staticmethod
    def _shape_name(shape) -> str:
        return str(getattr(shape, "name", "") or "")

    @staticmethod
    def _shape_has_table(shape) -> bool:
        try:
            if getattr(shape, "has_table", False):
                return True
        except Exception:
            pass

        try:
            table = getattr(shape, "table", None)
            return table is not None
        except Exception:
            return False

    def _resolve_table_target_shape(self, slide, statement_type: str):
        """Resolve the best existing target for a BS/IS table on a slide."""
        statement_type = (statement_type or "").upper()
        preferred_names = [
            "Table Placeholder",
            "Table Placeholder 2",
            "Content Placeholder 2",
        ]
        if statement_type == "IS":
            preferred_names.extend(["Table 3", "Table 2"])
        preferred_names.extend(["Table", "table", "TABLE"])

        for name in preferred_names:
            shape = self.find_shape_by_name(slide.shapes, name)
            if shape:
                return shape

        named_table_candidates = []
        table_candidates = []
        text_placeholder_candidates = []
        for shape in slide.shapes:
            shape_name = self._shape_name(shape)
            shape_name_lower = shape_name.lower()
            if "table" in shape_name_lower and "placeholder" in shape_name_lower:
                text_placeholder_candidates.append(shape)
                continue
            if self._shape_has_table(shape):
                table_candidates.append(shape)
                continue
            if "table" in shape_name_lower:
                named_table_candidates.append(shape)

        if text_placeholder_candidates:
            return text_placeholder_candidates[0]
        if table_candidates:
            return table_candidates[0]
        if named_table_candidates:
            return named_table_candidates[0]
        return None

    def _calculate_table_bounds(self, slide, target_shape=None, statement_type: str = "BS") -> Dict[str, int]:
        """Use target geometry when available, otherwise derive bounds from slide layout.

        The table top is aligned with the textMainBullets commentary body so the
        financial table and the Commentary blue box sit on the same horizontal
        baseline. The blue "Commentary"/"Table" label boxes (TextBox 10/11 in
        the template) act as headers above this baseline and are not covered
        by the table.
        """
        if target_shape is not None:
            return {
                "left": target_shape.left,
                "top": target_shape.top,
                "width": target_shape.width,
                "height": target_shape.height,
            }

        slide_width = getattr(self.presentation, "slide_width", Inches(10))
        slide_height = getattr(self.presentation, "slide_height", Inches(7.5))

        title_like_shapes = []
        body_like_shapes = []
        subtitle_shapes = []
        label_shapes = []  # Blue "Commentary"/"Table" label boxes — used as baseline anchor.
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            name = self._shape_name(shape).lower()
            try:
                label_text = (shape.text_frame.text or "").strip().lower()
            except Exception:
                label_text = ""
            if "subtitle" in name:
                subtitle_shapes.append(shape)
                body_like_shapes.append(shape)
            elif any(token in name for token in ["title", "projtitle"]):
                title_like_shapes.append(shape)
            elif any(token in name for token in ["text-commentary", "textmainbullets", "content"]):
                body_like_shapes.append(shape)
            elif label_text in ("commentary", "table"):
                label_shapes.append(shape)

        left = Inches(0.5)
        width = max(Inches(5.5), slide_width - Inches(1.0))
        top = Inches(1.45 if statement_type.upper() == "BS" else 1.6)
        height = slide_height - top - Inches(0.45)

        if title_like_shapes:
            bottom = max(shape.top + shape.height for shape in title_like_shapes)
            top = max(top, bottom + Inches(0.15))

        generic_is_layout = statement_type.upper() == "IS" and subtitle_shapes and target_shape is None
        if generic_is_layout:
            earliest_subtitle_top = min(shape.top for shape in subtitle_shapes)
            left = Inches(0.5)
            top = Inches(1.5)
            width = min(slide_width - left - Inches(0.35), int((slide_width - Inches(1.0)) * 0.5))
            height = max(Inches(2.0), earliest_subtitle_top - top - Inches(0.12))

        # Horizontal extent: widen left to include any body/label shape that
        # starts further left than the default, but STOP before any such
        # shape that sits to the right (a commentary column) — never extend
        # the table's right edge OUT to it. body_like_shapes/label_shapes are
        # always commentary text or its label, so there's no case where the
        # table should legitimately render under/over one of them. The old
        # `right_edge = max(anchor right edges)` did exactly that on a layout
        # with a single right-positioned commentary box and no left-side
        # counterpart (e.g. the first BS/IS slide, which has no separate
        # "Table" label once the table's own title row replaces it) —
        # right_edge became the commentary box's own right edge, stretching
        # the table full-width so it rendered on top of the commentary text
        # instead of stopping at the gutter before it.
        horizontal_anchors = list(body_like_shapes) + list(label_shapes)
        if horizontal_anchors:
            left = min(left, min(shape.left for shape in horizontal_anchors))
            anchors_to_the_right = [shape for shape in horizontal_anchors if shape.left > left]
            if anchors_to_the_right and not generic_is_layout:
                gutter = Inches(0.15)
                right_edge = min(shape.left for shape in anchors_to_the_right) - gutter
                width = min(width, max(Inches(2.0), right_edge - left))

        # Vertical alignment: anchor the table TOP to the "Commentary" /
        # "Table" blue label box (TextBox 10/11 in the template). This puts
        # the navy title row of the table at the exact same visible level as
        # the Commentary label on the right, replacing the need for a separate
        # "Table" label above the table.
        if label_shapes:
            label_top = min(shape.top for shape in label_shapes)
            top = max(top, label_top)
            if not generic_is_layout:
                height = max(Inches(2.5), slide_height - top - Inches(0.35))
        elif body_like_shapes:
            earliest_body_top = min(shape.top for shape in body_like_shapes)
            if not generic_is_layout:
                # No label shapes (dynamically added slide): fall back to the
                # commentary body top as the anchor.
                top = max(top, earliest_body_top)
                height = max(Inches(2.5), slide_height - top - Inches(0.35))
            else:
                height = min(height, max(Inches(2.0), earliest_body_top - top - Inches(0.12)))

        width = min(width, slide_width - left - Inches(0.25))
        height = max(Inches(2.5), min(height, slide_height - top - Inches(0.2)))
        return {
            "left": int(left),
            "top": int(top),
            "width": int(width),
            "height": int(height),
        }

    @staticmethod
    def _read_table_style_id(tbl_element) -> Optional[str]:
        """Read <a:tableStyleId> (the style GUID) from a table's XML, or None."""
        try:
            from pptx.oxml.ns import qn
            tblPr = tbl_element.find(qn("a:tblPr"))
            if tblPr is None:
                return None
            el = tblPr.find(qn("a:tableStyleId"))
            return el.text.strip() if (el is not None and el.text) else None
        except Exception:
            return None

    @staticmethod
    def _set_table_style_id(tbl_element, style_id: str) -> None:
        """Set the table's style GUID so PowerPoint renders it with that (e.g.
        UpSlide) table style instead of the python-pptx default."""
        from pptx.oxml.ns import qn
        tblPr = tbl_element.find(qn("a:tblPr"))
        if tblPr is None:
            tblPr = tbl_element.makeelement(qn("a:tblPr"), {})
            tbl_element.insert(0, tblPr)  # tblPr must be the first child of <a:tbl>
        for el in tblPr.findall(qn("a:tableStyleId")):
            tblPr.remove(el)
        style_el = tblPr.makeelement(qn("a:tableStyleId"), {})
        style_el.text = style_id
        tblPr.append(style_el)

    def _resolve_table_style_id(self) -> Optional[str]:
        """The table style GUID to apply to BS/IS tables.

        Priority: explicit config (pptx.table_style_id) > the style GUID of any
        table already present in the template (i.e. the firm's UpSlide style).
        Cached after first resolution. Returns None if none found (keep default).
        """
        if hasattr(self, "_cached_table_style_id"):
            return self._cached_table_style_id
        style_id = None
        try:
            style_id = (str(self.pptx_settings.get("table_style_id") or "").strip()) or None
        except Exception:
            style_id = None
        if not style_id and self.presentation is not None:
            for slide in self.presentation.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_table", False):
                        sid = self._read_table_style_id(shape.table._tbl)
                        if sid:
                            style_id = sid
                            break
                if style_id:
                    break
        if style_id:
            logger.info("Applying table style GUID to BS/IS tables: %s", style_id)
        self._cached_table_style_id = style_id
        return style_id

    def _add_table_to_slide(self, slide, df, bounds: Dict[str, int], table_name: str = None):
        total_rows = len(df) + 2 if table_name else len(df) + 1
        graphic_frame = slide.shapes.add_table(
            total_rows,
            len(df.columns),
            bounds["left"],
            bounds["top"],
            bounds["width"],
            bounds["height"],
        )
        # Auto-apply the firm's UpSlide table style (detected from the
        # template, or set via config) so new tables match existing ones.
        style_id = self._resolve_table_style_id()
        if style_id:
            try:
                self._set_table_style_id(graphic_frame.table._tbl, style_id)
            except Exception as exc:
                logger.debug("Could not apply table style %s: %s", style_id, exc)
        return graphic_frame

    def _fit_table_columns(self, table, df):
        """Allocate width by role and content length for better readability."""
        if len(table.columns) == 0:
            return

        try:
            total_width = sum(col.width for col in table.columns)
        except Exception:
            total_width = 0
        if total_width <= 0:
            return

        # A CJK character renders roughly 2x as wide as a Latin
        # character/digit at the same point size, but max_len here is a
        # raw character COUNT -- an 11-character Chinese row label like
        # "一年内到期的非流动负债" measures the same "length" as an
        # 11-character Latin one that's actually half as wide on the page.
        # Using the same /10 divisor for both meant max_len/10 almost never
        # exceeded the 2.0 floor for realistic Chinese labels (would need
        # 20+ characters), so column 0's weight was effectively a FIXED
        # 2.0 regardless of actual label length -- combined with the 0.12in
        # left-indent every leaf row gets, longer labels routinely wrapped
        # to 2 lines, each one rendering roughly 2x its neighbors' height
        # (PowerPoint auto-grows a row to fit wrapped text; the nominal
        # row.height set elsewhere is a floor, not a cap).
        col0_series = df.iloc[:, 0].astype(str) if len(df.columns) else pd.Series(dtype=str)
        is_cjk_labels = any(
            any('一' <= ch <= '鿿' for ch in str(v)) for v in col0_series.head(25).tolist()
        )

        weights = []
        for col_idx, col_name in enumerate(df.columns[: len(table.columns)]):
            col_series = df.iloc[:, col_idx].astype(str) if col_idx < len(df.columns) else pd.Series(dtype=str)
            max_len = max([len(str(col_name))] + [len(val) for val in col_series.head(25).tolist()]) if len(col_series) else len(str(col_name))
            col_name_str = str(col_name).lower()
            if col_idx == 0:
                weight = (
                    max(2.6, min(4.2, max_len / 5)) if is_cjk_labels
                    else max(2.0, min(3.2, max_len / 10))
                )
            elif any(token in col_name_str for token in ["20", "19", "date", "年", "月"]):
                weight = max(1.4, min(2.0, max_len / 10))
            else:
                weight = max(1.2, min(1.9, max_len / 9))
            weights.append(weight)

        total_weight = sum(weights) or 1
        assigned = 0
        for col_idx, weight in enumerate(weights):
            if col_idx == len(weights) - 1:
                width = total_width - assigned
            else:
                width = int(total_width * weight / total_weight)
                assigned += width
            table.columns[col_idx].width = max(int(Inches(0.7)), width)

    @staticmethod
    def _format_table_value(value, is_numeric_column: bool) -> str:
        def _fmt_number(n: float) -> str:
            if n == 0:
                return "-"
            # Accounting convention: negatives in parentheses, not with a minus sign.
            return f"({abs(n):,.0f})" if n < 0 else f"{n:,.0f}"

        if pd.isna(value):
            return ""
        if isinstance(value, (int, float)) and is_numeric_column:
            return _fmt_number(float(value))

        text_val = str(value).strip()
        if is_numeric_column:
            numeric_candidate = text_val.replace(",", "")
            if re.fullmatch(r"-?\d+(\.\d+)?", numeric_candidate):
                try:
                    return _fmt_number(float(numeric_candidate))
                except Exception:
                    return text_val
        return text_val

    @staticmethod
    def _set_paragraph_left_indent(paragraph, left_indent_emu: int) -> None:
        """Set a table-cell paragraph's left indent (marL) directly on its
        <a:pPr> XML, with indent (first-line offset) pinned to 0.

        _Paragraph has NO left_indent property in this python-pptx version
        (only alignment/level/line_spacing/font are exposed) -- `paragraph.
        left_indent = Inches(...)` doesn't raise, but that's because plain
        Python objects accept arbitrary ad-hoc attribute assignment; it
        silently creates a throwaway instance attribute with ZERO effect on
        the underlying XML, discarded the moment the object is garbage
        collected. Confirmed by round-tripping through a real save+reload:
        the "set" value reads back fine within the SAME Python session (the
        fake attribute is still sitting right there), but a freshly loaded
        Presentation() from that same saved file raises AttributeError on
        the same read -- proof nothing was ever written. marL/indent are
        real OOXML attributes on <a:pPr> (ECMA-376 CT_TextParagraphProperties)
        that python-pptx just doesn't wrap with a friendly property; setting
        them via the raw element (same get_or_add_pPr() pattern python-pptx's
        own oxml layer uses internally) is the only way that actually
        persists.
        """
        pPr = paragraph._p.get_or_add_pPr()
        pPr.set('marL', str(int(left_indent_emu)))
        pPr.set('indent', '0')

    # -- Presentation-detail tables (per-account report-ready breakdowns) --
    # workbook.py's extract_presentation_detail_table finds a small,
    # human-readable breakdown table some account sheets carry alongside
    # their GL-coded main schedule (e.g. 管理费用's 会计服务费/审计费/...). This
    # renders it as a native PowerPoint table next to that account's own
    # commentary bullet. Config-gated (pptx_settings.presentation_tables.
    # enabled), off by default until verified against a real export.
    #
    # Reuses the generic table primitives above (add_table, _fit_table_
    # columns, _format_table_value, _set_paragraph_left_indent,
    # _resolve_table_style_id) rather than _embed_statement_table/
    # _fill_table_placeholder -- those are hard-coupled to the full BS/IS
    # overview grid (category-header rows, row-role tiering by keyword) that
    # a flat ~3-15 row account breakdown doesn't have and doesn't need.
    _TABLE_TITLE_ROW_PT = 16.0
    _TABLE_HEADER_ROW_PT = 14.0
    _TABLE_DATA_ROW_PT = 12.0
    _TABLE_CHILD_ROW_PT = 11.0
    _TABLE_TOTAL_ROW_PT = 14.0
    _TABLE_SOURCE_LINE_PT = 12.0
    # Pure whitespace, not text-safety margin -- trimmed 2026-07-31 (6->4,
    # and the source-line box's own "+4" padding at its two call sites ->
    # "+2") to buy back a few points of bottom-of-slide margin on the
    # tallest real stack (营业成本: 16-row table + long lead-in + long
    # explanation measured at only 0.313in/22.6pt clear of the slide's
    # bottom edge on a real Crescent export). Safe to shrink because
    # neither is sized against variable AI text -- the source line is a
    # single fixed short string, and this gap is pure spacing between two
    # already-independently-sized shapes.
    # 4.0 -> 2.0 (2026-08-04). Real PowerPoint BoundHeight measurements
    # decomposed the visible gap under a lead-in exactly: 2.2pt of
    # paragraph spacing we count but BoundHeight doesn't show, 3.5pt from
    # the safety factor, and this 4pt spacer -- 9.7pt total, matching the
    # measured 9.7pt to the decimal. 2pt still reads as a deliberate
    # separation between the text and the table below it.
    _TABLE_GAP_ABOVE_PT = 2.0
    _TABLE_GAP_BELOW_PT = 4.0
    # Shared margin for both text blocks sized against a table (lead-in
    # above it, explanatory bullets below it): inspect_pptx.py re-derives
    # capacity/usage independently from the exported XML rather than
    # sharing this instance's live measurement state, so its verdict can
    # disagree with what was used to size the shape. 1.3x still left a
    # real generated lead-in at 104% (⚠️ OVERFLOW RISK) on a real Crescent
    # export -- raised rather than chasing exact parity between two
    # independently-implemented measurers, since the safe direction here
    # is a touch more headroom, not exact agreement.
    #
    # NOT the lever for a stack running close to the SLIDE's own bottom
    # edge (as opposed to text overflowing its own box) -- raising this
    # makes every box it sizes (lead-in, explanation) TALLER, which pushes
    # whatever comes after it further DOWN the slide. For a single tall
    # stack like 营业成本 that's already tight against the slide bottom,
    # raising this shrinks that margin further, not the other way round.
    _TEXT_HEIGHT_SAFETY_FACTOR = 1.6
    # Render-time counterpart for the two ACTUAL shape heights this feature
    # sets directly (table lead-in, table explanation). The full 1.6x
    # applied to an ACTUAL rendered box instead (as it was until
    # 2026-08-03) creates VISIBLE blank space below the real text and
    # before the table -- confirmed against a real Crescent export where a
    # ~3-line lead-in got sized for ~4.8 "line units", roughly 2 blank
    # lines, matching a real user report ("表格頂部有一些空白...大約2-3行").
    # 1.6x's own history (raised from 1.3x after THAT still read 104%
    # overflow-risk on inspect_pptx.py's independent re-check) means this
    # can't just be zeroed out -- some margin is real. Iteratively tuned
    # (1.35 -> 1.15 -> 1.25) against TWO real texts from an actual Crescent
    # export (营业成本's and 财务费用's own lead-in/explanation) via
    # inspect_pptx.py's independent re-check each time -- 1.15 looked fine
    # on the first (98%/94% fill) but the SECOND, shorter-but-differently-
    # wrapping explanation hit 101% (⚠️ OVERFLOW RISK), proving a single
    # real sample isn't enough evidence to land on a value. 1.25 clears
    # both with real margin (85-91% fill, not razor-thin) -- if a future
    # real case still overflows at 1.25, that's a real, reproducible data
    # point to retune from, not a guess.
    #
    # ALSO used (2026-08-03, same day) by the PLANNING-time estimates
    # below (_estimate_lead_in_pt / _estimate_table_account_block_height_
    # pt) -- these used to stay at the full 1.6x on the theory that a more
    # conservative planning estimate is the safe direction (better to
    # under-pack than overflow). That theory stopped holding once render
    # itself moved to 1.25x: planning at 1.6x now systematically believes
    # LESS room remains after a table block than render will actually
    # need, so a trailing account that would genuinely have fit was never
    # even tried -- confirmed against two real Crescent table pages
    # showing 1.5-2.2in of real, measured, unexplained blank column below
    # a finished table+explanation(+trailing account) stack. Planning and
    # render now use the SAME factor so the estimate that drives "does
    # this fit" decisions actually predicts what render will produce,
    # instead of two different numbers pulling in opposite directions.
    # 1.25 -> 1.10 (2026-08-04). At 1.25 EVERY box rendered at exactly
    # 1/1.25 = 78-80% fill -- the user saw that constant 20-25% tail as
    # "無論當中的字是多少行 下面都會有一行的空間". An earlier attempt at 1.10
    # DID overflow (103-154% on real client metrics) and had to be
    # reverted, but that was before the missing-inset bug was found: each
    # box was silently 7.2pt short of its own content, and this factor was
    # the only thing covering it. With insets now added properly, 1.10 was
    # re-validated by sweeping every real lead-in AND explanation text
    # from a real export, each scored with inspect_pptx.py's own
    # independent formula rather than the renderer's: worst case 87%
    # (lead) / 91% (explanation), i.e. real margin. The same sweep shows
    # 1.05 reaching 95% and 1.00 hitting 100% -- so 1.10 is the last value
    # with genuine headroom, not merely the first that happens to pass.
    _TABLE_RENDER_HEIGHT_SAFETY_FACTOR = 1.10

    def _presentation_tables_enabled(self) -> bool:
        try:
            return bool((self.pptx_settings.get("presentation_tables") or {}).get("enabled", False))
        except Exception:
            return False

    def _presentation_table_style(self) -> str:
        """"table" (default, native PPTX table + dedicated stack shapes,
        see _render_table_accounts_stack) or "sublist" (config-gated
        fallback: the table dict is converted to plain indented text lines
        appended to the account's own commentary, and the account is NEVER
        pulled out of the normal packing pool at all -- it inherits every
        existing text module -- packing, splitting, cross-column
        continuation, overflow handling -- for free, at the cost of losing
        the native-table look and per-period-per-component detail). See
        _sublist_text_for_table."""
        try:
            value = (self.pptx_settings.get("presentation_tables") or {}).get("style", "table")
        except Exception:
            value = "table"
        value = str(value or "table").strip().lower()
        return value if value in ("table", "sublist") else "table"

    def _sublist_text_for_table(
        self, table: Dict[str, Any], is_chinese_databook: bool, source_multiplier: float = 1,
        max_items: int = 5,
    ) -> str:
        """Converts a presentation_detail_table dict (extract_presentation_
        detail_table's return shape) into plain text lines for "sublist"
        style. Component lines show ONLY the LATEST period's figure -- a
        full table's worth of per-period, per-component detail written out
        as prose would be worse than the empty space this whole feature
        exists to fill. The total line keeps every period inline, matching
        how OI/OC accounts already state multi-year figures in this
        project's own reference style. Top-level rows only: nested children
        are already rolled into their own parent's total (e.g. 物业管理费's
        第三方/上海熙麦 sub-vendors), and a plain-text bullet list is not the
        place for two levels of indentation -- this already keeps
        same-nature items merged under their shared parent.

        When there are more than max_items top-level rows, only the
        max_items-1 largest (by the latest period's absolute value) are
        shown individually, ranked descending; everything past that is
        rolled into one final "其他"/"Other" line (summed, not dropped --
        the account's own real total still fully accounts for it even
        though this specific line doesn't itemise it). Keeps a long
        component list (e.g. 管理费用's 8 rows) from producing an equally
        long, table-like bullet list -- exactly what "sublist" style trades
        the native table's own per-component precision for.

        Values are in the same raw-yuan internal scale every account's df
        uses (see _render_presentation_table's own docstring) -- divided
        back down by source_multiplier here, at text-building time only,
        same as the native-table path does at render time (cadbce8)."""
        divisor = source_multiplier if source_multiplier and source_multiplier != 0 else 1

        def _scaled(v):
            return v / divisor if isinstance(v, (int, float)) else v

        periods = table.get("periods") or []
        period_labels = table.get("period_labels") or {}
        rows = table.get("rows") or []
        total_row = table.get("total_row") or {}
        if not periods or not rows:
            return ""

        latest_period = periods[-1]
        marker = "- "
        items: List[Tuple[str, float]] = []
        for row in rows:
            label = row.get("label", "")
            value = _scaled((row.get("values") or {}).get(latest_period))
            if value is None or not label:
                continue
            items.append((label, value))

        lines: List[str] = []
        if len(items) > max(1, max_items):
            ranked = sorted(items, key=lambda item: abs(item[1]), reverse=True)
            shown, rest = ranked[: max(1, max_items - 1)], ranked[max(1, max_items - 1):]
            for label, value in shown:
                lines.append(f"{marker}{label}：{self._format_table_value(value, is_numeric_column=True)}")
            if rest:
                other_label = "其他" if is_chinese_databook else "Other"
                other_value = sum(v for _l, v in rest)
                lines.append(f"{marker}{other_label}：{self._format_table_value(other_value, is_numeric_column=True)}")
        else:
            for label, value in items:
                lines.append(f"{marker}{label}：{self._format_table_value(value, is_numeric_column=True)}")

        if total_row:
            total_label = total_row.get("label") or ("合计" if is_chinese_databook else "Total")
            total_values = total_row.get("values") or {}
            parts = []
            for period in periods:
                v = _scaled(total_values.get(period))
                if v is None:
                    continue
                label = period_labels.get(period, period)
                text_val = self._format_table_value(v, is_numeric_column=True)
                parts.append(f"{label}{text_val}" if is_chinese_databook else f"{text_val} in {label}")
            if parts:
                joiner = "，" if is_chinese_databook else ", "
                sep = "：" if is_chinese_databook else ": "
                lines.append(f"{marker}{total_label}{sep}{joiner.join(parts)}")

        return "\n".join(lines)

    @staticmethod
    def _presentation_table_for_account(account_data: Dict[str, Any]) -> Optional[Dict[str, Any]]:
        financial_data = (account_data or {}).get("financial_data")
        if not isinstance(financial_data, pd.DataFrame):
            return None
        try:
            table = (financial_data.attrs or {}).get("presentation_detail_table")
        except Exception:
            return None
        if not table or not table.get("rows"):
            return None
        return table

    @classmethod
    def _presentation_table_height_pt(cls, table: Dict[str, Any]) -> float:
        """Total table height in points, from the SAME row-height constants
        _render_presentation_table sets as real row heights -- so the space
        reserved during packing and the space actually drawn can't drift
        apart from each other."""
        rows = table.get("rows") or []
        pts = cls._TABLE_TITLE_ROW_PT + cls._TABLE_HEADER_ROW_PT
        for row in rows:
            pts += cls._TABLE_DATA_ROW_PT
            pts += len(row.get("children") or []) * cls._TABLE_CHILD_ROW_PT
        if table.get("total_row"):
            pts += cls._TABLE_TOTAL_ROW_PT
        return pts

    # The literal handoff phrase ai.py's _detail_table_guidance asks the
    # model to end its short lead-in with -- the one point in the real
    # convention's two-part structure (lead-in, then optional "-"/"➢"
    # explanatory bullets) that's reliably the SAME string every time,
    # so it's what _split_table_commentary splits on.
    _TABLE_HANDOFF_CHI = "明细如下："
    _TABLE_HANDOFF_ENG = "the breakdown is set out below"

    @staticmethod
    def _truncate_text_at_boundary(text: str, limit: int, is_chinese: bool) -> str:
        """Cuts `text` to at most `limit` chars at a sentence boundary where
        possible. Shared by the lead-in and the post-table explanation --
        same safety-net shape, different caps (see _split_table_commentary)."""
        text = (text or "").strip()
        if len(text) <= limit:
            return text
        boundary_chars = "。；;.!?！？"
        cut = text[:limit]
        # A "." between two digits is a DECIMAL POINT, not a sentence end.
        # Taking it as one truncated a real deck's 营业成本 lead-in at
        # "...较2025年度下降74." -- the last "boundary" in the string was the
        # point inside 74.9%, so the figure was cut in half and the rest of
        # the sentence, including the "明细如下：" handoff, was thrown away.
        # Same defect class as the mid-number split _snap_split_before_number
        # already guards in the packing path; this truncation path never had
        # the guard.
        best = -1
        for pos in range(len(cut) - 1, -1, -1):
            ch = cut[pos]
            if ch not in boundary_chars:
                continue
            if (
                ch == "."
                and pos > 0
                and cut[pos - 1].isdigit()
                and pos + 1 < len(text)
                and text[pos + 1].isdigit()
            ):
                continue
            best = pos
            break
        if best >= int(limit * 0.4):
            return cut[: best + 1]
        return cut.rstrip() + ("…" if is_chinese else "...")

    def _split_table_commentary(self, commentary: str, is_chinese: bool) -> Tuple[str, str]:
        """Splits an account's raw commentary into (lead_in, post_table_text)
        at the handoff phrase ai.py's _detail_table_guidance asks the model
        to end its short lead-in with. Everything up to and including that
        phrase is the lead-in (rendered above the table); anything after it
        is the optional "-"/"➢" explanatory bullets a real deliverable puts
        BELOW the table (provider, charging basis, contract terms) -- a real
        photo comparison against the project team's own Crescent deck showed
        this second part entirely missing before: the whole commentary was
        being treated as lead-in and hard-capped at ~220/340 chars, silently
        dropping any such bullets the model had written past that point with
        nowhere for them to ever render.

        Each side is truncated separately and more generously than the old
        single whole-commentary cap, since each now has its own dedicated
        vertical space (lead-in above the table, explanation below it) --
        see _presentation_table_extra_text_height_pt, which sizes both.

        Falls back to (whole commentary, '') if the model didn't include the
        handoff phrase -- exactly the old behaviour, so a non-compliant
        generation degrades to "no post-table text" rather than misplacing
        content."""
        handoff = self._TABLE_HANDOFF_CHI if is_chinese else self._TABLE_HANDOFF_ENG
        text = (commentary or "").strip()
        idx = text.lower().find(handoff.lower())
        if idx < 0:
            lead_limit = 220 if is_chinese else 340
            return self._truncate_text_at_boundary(text, lead_limit, is_chinese), ""
        split_at = idx + len(handoff)
        lead_in = text[:split_at].strip()
        post_table = text[split_at:].strip()
        lead_limit = 140 if is_chinese else 220
        # The real, confirmed-fitting Crescent example (营业成本, the TALLEST
        # of the 4 known tables -- 15 rows) carried ~250 chars of "-"/"➢"
        # explanation. This cap allows meaningfully more room than that
        # observed-good case without being open-ended -- a table already
        # near the tallest slot capacity plus an unbounded explanation is
        # the one real overflow risk this feature has, since it's genuinely
        # new vertical space that wasn't being claimed before.
        post_limit = 450 if is_chinese else 700
        return (
            self._truncate_text_at_boundary(lead_in, lead_limit, is_chinese),
            self._truncate_text_at_boundary(post_table, post_limit, is_chinese),
        )

    def _slot_names_for_actual_slide(self, actual_slide_idx: int, start_slide: int) -> List[str]:
        """Which slot names this slide's OWN template shapes actually
        support. Mirrors _distribute_content_across_slots's own identically
        -named local closure (slot_names_for_actual_slide) exactly,
        including its "first slide of statement -> single" fallback rule --
        this MUST agree with that closure's notion of what a slide offers,
        or a table account could be "placed" onto a slot name the packer's
        own accounting disagrees with, silently overlapping real content.
        `actual_slide_idx` is an absolute index into self.presentation.
        slides, NOT the slide_idx slot_distribution tuples carry (those are
        0-based from start_slide -- see _append_table_accounts_to_distribution).
        """
        if self.presentation is not None and 0 <= actual_slide_idx < len(self.presentation.slides):
            slide = self.presentation.slides[actual_slide_idx]
            has_left = self.find_shape_by_name(slide.shapes, "textMainBullets_L") is not None
            has_right = self.find_shape_by_name(slide.shapes, "textMainBullets_R") is not None
            has_single = self.find_shape_by_name(slide.shapes, "textMainBullets") is not None
            if has_left and has_right:
                return ["L", "R"]
            if has_single:
                return ["single"]
        return ["single"] if actual_slide_idx == start_slide - 1 else ["L", "R"]

    # Planning-time-only assumption (no real shape exists yet for a slot a
    # table account hasn't been placed into) of how much vertical room a
    # slot has, used purely to decide whether a SECOND table account can
    # join one that already has a first -- render time positions every
    # account's own block at a running offset regardless of this estimate
    # (see _render_table_accounts_stack), so getting this exactly right
    # isn't load-bearing for correctness, only for how good the packing
    # decision is. ~331pt (single) to ~360pt (L/R) measured directly from
    # the real template (scratchpad/measure_slot_capacity.py); table
    # accounts land in L/R far more often than single (single is usually
    # claimed by the statement's own first/executive-summary slide), so
    # this leans toward that number rather than splitting the difference.
    _TABLE_SLOT_CAPACITY_PT = 355.0
    # Real accounts were only ADDED to a slot's already-claimed footprint at
    # 90% of the assumed capacity, not 100% -- deliberate margin for the
    # planning estimate (heuristic, no real font metrics) to disagree with
    # the real-metrics render-time measurement without the combined result
    # actually overflowing the slot. Not tighter than this: each account's
    # OWN estimate already carries the full _TEXT_HEIGHT_SAFETY_FACTOR
    # margin individually, so stacking a second, separate margin on top of
    # that here would mostly just suppress legitimate packing rather than
    # add real safety -- confirmed against a real Crescent export that two
    # of the four real accounts (税金及附加 ~167pt, 财务费用 ~222pt once its
    # own real, fairly long variance-flagging explanation is counted) don't
    # fit together at ANY reasonable threshold (389pt combined vs a
    # 331-360pt slot) -- that's a genuine size fact, not a threshold this
    # constant should be tuned to paper over.
    #
    # 0.90 -> 0.92 (2026-08-04), a deliberately small raise. Two errors
    # this threshold had been quietly absorbing were fixed the same day:
    # the planning estimates were ~14% low (font x 1.0 instead of
    # PowerPoint's real 1.2 line pitch -- see _planning_std_lh_pt) and it
    # multiplied one shared capacity constant instead of each slot's real
    # height (see _slot_capacity_pt). Measured after both fixes, the
    # estimate still runs ~5% under what render produces, so the SAFE
    # ceiling is about 0.95/1.05 ~= 0.90; 0.92 takes the small genuine
    # gain and no more. A first attempt at 0.96 was rejected on this
    # arithmetic alone -- it works out to ~101% of a column's real
    # capacity, i.e. overflow by construction.
    # 0.92 -> 0.98 (2026-08-04). Every reason this sat below 1.0 has now
    # been removed one at a time: the shared-capacity constant (now each
    # slot's real height, _slot_capacity_pt), the ~14% planning line-pitch
    # error (_planning_std_lh_pt), and finally the trailing-paragraph-gap
    # over-count -- the cost model now reproduces PowerPoint's own
    # BoundHeight exactly, so there is no measurement bias left for this
    # threshold to absorb. 0.98 keeps a 2% cushion for the render-time
    # rounding this estimate can't see, while letting a table's own block
    # actually use the column it was measured against; the user's
    # standing preference after seeing a slot land at exactly 0pt spare
    # is "用盡空間 比較好".
    _TABLE_SLOT_PACK_THRESHOLD = 0.98

    def _estimate_lead_in_pt(self, item: Dict[str, Any], is_chinese_databook: bool = False) -> float:
        """Heuristic (no real shape) estimate of one account's lead-in block
        alone -- category header + "key - commentary" bullet -- in points.
        Shared by _estimate_table_account_block_height_pt (as its lead-in
        arm) and directly by trailing plain-text accounts appended to a
        table slot's leftover space (see _append_table_accounts_to_
        distribution's trailing_items), which have no table/source/
        explanation arms at all."""
        is_chinese = bool(item.get("is_chinese"))
        lead_in_units = self._calculate_content_lines(
            # The rendered label, NOT the raw mapping_key -- see
            # _rendered_bullet_label for why that difference is a whole
            # wrapped line on real Chinese content.
            item.get("category") or "", self._rendered_bullet_label(item, is_chinese_databook),
            item.get("commentary", ""), slot_name="single",
            shape=self._measurement_slot_shape(), is_chinese=is_chinese,
            # whole_box=False: this block sits INSIDE the column's shared
            # frame with more content after it, so it really does carry its
            # own trailing paragraph gap. Only the frame's last block gets
            # that gap refunded, and that refund belongs at frame level
            # (slot_cost), not here -- charging it per account is the exact
            # bug that took 15 rounds to find in the capacity investigation.
            whole_box=False,
        )
        lead_std_lh_pt = self._planning_std_lh_pt(is_chinese)
        # No safety factor and no insets. Both existed because a lead-in used
        # to render into its OWN textbox, which carries insets and needed
        # margin against a height this file computed but PowerPoint laid out.
        # Since _render_table_accounts_stack writes every lead-in as ordinary
        # paragraphs in the column's single shared frame, a lead-in costs
        # exactly what any other bullet costs -- its measured text height.
        # (Do NOT add a category-header line here: _calculate_content_lines
        # already charges one line pitch for it when `category` is non-empty,
        # verified at 10.80pt, matching the header paragraph's space_after=0.
        # Adding it again over-charged every account by a full line.)
        return lead_in_units * lead_std_lh_pt

    def _estimate_table_account_block_height_pt(
        self, item: Dict[str, Any], table: Dict[str, Any], is_chinese_databook: bool,
    ) -> float:
        """Heuristic (no real shape) estimate of one table-bearing account's
        whole block -- lead-in + table + source line + explanation -- for
        the bin-packing decision in _append_table_accounts_to_distribution.
        Mirrors the real render-time formula (_add_presentation_detail_
        table_below_text / _render_presentation_table) closely enough to be
        directionally right, without needing a real shape to measure
        against (none exists yet at planning time)."""
        return sum(self._estimate_table_account_parts_pt(item, table, is_chinese_databook))

    def _table_block_reserved_pt(self, table: Dict[str, Any]) -> float:
        """Vertical space one table block occupies inside the column's shared
        frame: the gap above it, the table itself, and the source caption
        _render_presentation_table draws underneath (whose box is
        _TABLE_SOURCE_LINE_PT + 2 tall).

        Single definition on purpose -- the renderer reserves this and the
        planner charges it, and the two silently drifting is exactly the
        class of bug that produced a table overhanging its column.

        No _TABLE_GAP_BELOW_PT: separation from the paragraph that follows
        is already provided by that paragraph's own space_before/after in
        the shared frame, and adding it here was pushing the reservation
        over a line boundary -- the extra blank line reported under every
        subtable."""
        return (
            self._TABLE_GAP_ABOVE_PT
            + self._presentation_table_height_pt(table)
            + self._TABLE_SOURCE_LINE_PT + 2.0
        )

    def _estimate_table_account_parts_pt(
        self, item: Dict[str, Any], table: Dict[str, Any], is_chinese_databook: bool,
    ) -> Tuple[float, float, float]:
        """The same estimate broken into its three independently-placeable
        parts -- (lead-in, table+source, explanation) -- so flow() can try
        splitting an account at either boundary rather than only
        accepting or rejecting it whole."""
        lead_in_pt = self._estimate_lead_in_pt(item, is_chinese_databook)

        # Exactly what the renderer reserves -- single definition, so the two
        # can't drift (see _table_block_reserved_pt).
        table_pt = self._table_block_reserved_pt(table)

        post_table_text = item.get("_post_table_text", "")
        explain_pt = 0.0
        if post_table_text:
            explain_units = self._calculate_content_lines(
                # Measure the marker-prefixed text the renderer actually
                # writes (_append_explanation_to_frame), not the bare text --
                # "➢ " is two more characters and on a full line that is
                # enough to force one extra wrap, measured at 10.8pt (a whole
                # line) of under-estimate on a real-shaped account.
                "", "", self._explanation_render_text(post_table_text, is_chinese_databook),
                slot_name="single", shape=self._measurement_slot_shape(),
                is_chinese=is_chinese_databook, whole_box=False,
            )
            # Like _estimate_lead_in_pt: no safety factor, no insets. The
            # explanation is ordinary flowing text in the column's shared
            # frame now, not a floating textbox of its own.
            explain_pt = explain_units * self._planning_std_lh_pt(is_chinese_databook)

        return lead_in_pt, table_pt, explain_pt

    def _append_table_accounts_to_distribution(
        self, table_items: List[Dict[str, Any]],
        slot_distribution: List[Tuple[int, str, List[Dict[str, Any]]]],
        *, max_slides: int, start_slide: int, is_chinese_databook: bool = False,
        trailing_items: Optional[List[Dict[str, Any]]] = None,
    ) -> List[Tuple[int, str, List[Dict[str, Any]]]]:
        """Assigns each table-bearing account to a (slide_idx, slot_name),
        preferring to join an EARLIER table account's own slot when there's
        room (see _TABLE_SLOT_CAPACITY_PT/_TABLE_SLOT_PACK_THRESHOLD) rather
        than always claiming a fresh one -- a small table (e.g. 税金及附加,
        3 components) claiming a WHOLE column regardless of the normal
        packer's own output left most of that column empty; multiple small
        table accounts can share a column exactly the way multiple ordinary
        accounts already share one, just via a running EMU offset instead
        of shared-textframe paragraphs (see _render_table_accounts_stack).

        Never shares a slot with a NON-table account from the normal
        packer's own pool, and never touches that pool's already-finished
        distribution -- both of those are what caused the original
        two-tables-on-top-of-each-other bug this design replaced.

        trailing_items (optional) are ordinary accounts that come AFTER the
        LAST table account in the statement's own reading order -- e.g. 投资
        收益/营业外支出 following 财务费用 in a real Crescent IS deck. The
        caller (apply_structured_data_to_slides) excludes these from the
        normal packer's input pool entirely before it ever runs, precisely
        so this function can try them here instead: each is placed AFTER all
        table_items have claimed their slots, using ONLY its own lead-in
        height (no table/source/explanation arms -- it's plain commentary),
        against the exact same slot_fill_pt pool table_items just filled --
        which is exactly the empty space below a table that prompted this
        whole feature. A trailing item that fits into no existing slot
        claims a fresh one via the identical fallback table_items already
        use below, so it always renders somewhere; this function still never
        reopens or second-guesses the normal pool's own finished output.

        slide_idx here is 0-based FROM start_slide, matching exactly what
        _distribute_content_across_slots itself returns (its caller,
        apply_structured_data_to_slides, converts to an absolute slide via
        `start_slide - 1 + slide_idx` when it actually renders) -- using an
        absolute index here instead double-applies that offset downstream.
        """
        used = {(slide_idx, slot_name) for slide_idx, slot_name, _ in slot_distribution}
        result = list(slot_distribution)
        slot_fill_pt: Dict[Tuple[int, str], float] = {}
        def _cap_for(key: Tuple[int, str]) -> float:
            # Per-slot real capacity, not one shared constant -- a "single"
            # first-page slot and an L/R continuation slot differ by ~8%.
            return (self._slot_capacity_pt(key[0], key[1], start_slide)
                    * self._TABLE_SLOT_PACK_THRESHOLD)


        def _append_to_slot(key: Tuple[int, str], item: Dict[str, Any], block_pt: float) -> None:
            slide_idx, slot_name = key
            for i, (s, n, accounts) in enumerate(result):
                if s == slide_idx and n == slot_name:
                    result[i] = (s, n, accounts + [item])
                    break
            slot_fill_pt[key] = slot_fill_pt.get(key, 0.0) + block_pt

        def _open_new_slot(item: Dict[str, Any], block_pt: float) -> Optional[Tuple[int, str]]:
            for slide_idx in range(max_slides):
                actual_slide_idx = start_slide - 1 + slide_idx
                for slot_name in self._slot_names_for_actual_slide(actual_slide_idx, start_slide):
                    if (slide_idx, slot_name) in used:
                        continue
                    used.add((slide_idx, slot_name))
                    result.append((slide_idx, slot_name, [item]))
                    slot_fill_pt[(slide_idx, slot_name)] = block_pt
                    return (slide_idx, slot_name)
            logger.warning(
                "No free slot for account %s within %s slides -- not rendered.",
                item.get("mapping_key"), max_slides,
            )
            return None

        # Content flows in READING ORDER through one column at a time, the
        # way a document does -- fill the current column, then move to the
        # next. This replaces a first-fit search over every open slot,
        # which could place a LATER account into an EARLIER column that
        # still had room: a real export had 投资收益 (last in the income
        # statement) rendered on slide 4 between 税金及附加 and 管理费用,
        # because it happened to fit there. Sequential flow makes that
        # impossible by construction rather than by an extra ordering
        # check, and is also what makes the lead/table split below safe --
        # nothing can ever land between a lead-in and its own table.
        cursor: Optional[Tuple[int, str]] = None

        # Start the flow in the LAST slot the normal packer used, rather than
        # on a brand-new one. That slot is very often barely filled -- a real
        # deck left the IS statement page's commentary column at 33% while
        # the next page's right column sat COMPLETELY EMPTY, roughly two
        # blank columns between them.
        #
        # This slot used to be off-limits: every packer slot went into `used`
        # so a table could never share one, because two renderers writing the
        # same slot once produced overlapping tables. That objection is gone
        # -- since the shared-frame rewrite, _render_table_accounts_stack
        # writes its lead-ins as PARAGRAPHS in the slot's own text frame and
        # reserves table space inside it, and the render dispatch already
        # routes a slot holding any table account through that one path. A
        # mixed slot therefore has a single writer.
        #
        # Only the LAST slot is eligible, which is what keeps reading order:
        # the table accounts come after the plain ones in the statement, so
        # appending to the final packer slot continues the document; appending
        # to an earlier one would interleave them.
        if slot_distribution:
            for slide_idx, slot_name, accounts in reversed(slot_distribution):
                if not accounts:
                    continue
                key = (slide_idx, slot_name)
                # Same currency flow() measures in, so the fit test below is
                # comparing like with like.
                slot_fill_pt[key] = sum(
                    self._estimate_lead_in_pt(a, is_chinese_databook) for a in accounts
                )
                cursor = key
                logger.debug(
                    "Table flow continues in the packer's last slot %s (%.0f of %.0fpt used)",
                    key, slot_fill_pt[key], _cap_for(key),
                )
                break

        def _fits(key: Optional[Tuple[int, str]], block_pt: float) -> bool:
            return key is not None and slot_fill_pt.get(key, 0.0) + block_pt <= _cap_for(key)

        def flow(item: Dict[str, Any], block_pt: float,
                 parts_pt: Optional[Tuple[float, float, float]] = None) -> None:
            """Places one account at the flow cursor, splitting it at a
            natural boundary when the whole block won't fit in the current
            column but a leading portion of it still will.

            ONE split point: lead-in + table stay here, and only the
            trailing explanation continues in the next column under a
            "（续）" heading. That is the real reference deck's own
            convention for its 营业成本 bullets.

            A lead-in never separates from its own table -- see the removed
            branch below for why. The table itself is never divided (no
            repeated header) per explicit user instruction. parts_pt=None
            marks an account with nothing to split (a plain trailing
            account), which flows whole."""
            nonlocal cursor
            if _fits(cursor, block_pt):
                _append_to_slot(cursor, item, block_pt)
                return

            if parts_pt is not None and self._TABLE_ALLOW_LEAD_TABLE_SPLIT:
                lead_pt, table_pt, explain_pt = parts_pt
                heading = self._TABLE_CONTINUATION_HEADING_PT
                # 1. Keep the lead-in AND table here; only the explanation
                #    moves. Worth it only if there IS an explanation to move.
                if explain_pt > 0 and _fits(cursor, lead_pt + table_pt):
                    _append_to_slot(cursor, dict(item, _render_parts=("lead", "table")),
                                    lead_pt + table_pt)
                    cursor = _open_new_slot(dict(item, _render_parts=("expl",)),
                                            explain_pt + heading)
                    return
                # 2. Keep only the lead-in here; table + explanation move.
                #    REMOVED after reviewing a real deck. This was built to
                #    the letter of "表格前如果這個point有文字 而這一邊放不下表格
                #    那表格可以放下一個", but rendered it is the worst of the
                #    three outcomes and it fired on EVERY table account:
                #      - the column ends on "...明细如下：" with no 明细 under
                #        it, which simply reads as broken;
                #      - the table then needs a "（续）" heading in the next
                #        column -- the repeated header the same instruction
                #        asked us to avoid;
                #      - and the stranded lead-in leaves the rest of its own
                #        column empty (one real column sat ~50% blank).
                #    Falling through to _open_new_slot moves the whole
                #    lead+table+explanation block together instead. On the
                #    reviewed deck that is exactly one account per column
                #    with no continuation headings at all.

            cursor = _open_new_slot(item, block_pt)

        for item in table_items:
            table = item.get("_presentation_table") or {}
            parts_pt = self._estimate_table_account_parts_pt(item, table, is_chinese_databook)
            flow(item, sum(parts_pt), parts_pt=parts_pt)

        for item in (trailing_items or []):
            flow(item, self._estimate_lead_in_pt(item, is_chinese_databook))

        # Represent every OTHER slot on a slide this pass wrote to, as an
        # explicit empty entry. The render loop fills and clears purely by
        # walking the distribution, so a slot that simply never appears in
        # it is visited by neither branch and keeps whatever raw text the
        # template shipped with -- a real export leaked a literal
        # "Placeholder - placeholder" into slide 5's empty right column
        # this way. _distribute_content_across_slots already does exactly
        # this for the ordinary packing pool; the table pass needed its
        # own copy because it can leave a slot untouched that the ordinary
        # pass never saw either.
        touched_slides = {slide_idx for slide_idx, _slot, accounts in result if accounts}
        present = {(slide_idx, slot) for slide_idx, slot, _accounts in result}
        for slide_idx in sorted(touched_slides):
            actual_slide_idx = start_slide - 1 + slide_idx
            for slot_name in self._slot_names_for_actual_slide(actual_slide_idx, start_slide):
                if (slide_idx, slot_name) not in present:
                    result.append((slide_idx, slot_name, []))

        return result

    def _render_continuation_heading(
        self, slide, left: int, top: int, width: int, display_name: str,
        is_chinese_databook: bool,
    ) -> int:
        """Draws the compact "科目名（续）" / "name (cont'd)" line that stands
        in for an account's intro sentence when only its TABLE renders in
        this slot (the intro finished the previous column -- see
        place_table_item). Returns the bottom EMU position.

        Deliberately its own small renderer rather than reusing
        _fill_text_main_bullets_with_category_and_key with empty
        commentary: that path always writes a " - " separator after the
        account name, which would render as a dangling dash here."""
        suffix = "（续）" if is_chinese_databook else " (cont'd)"
        box = slide.shapes.add_textbox(left, top, width, Pt(self._TABLE_CONTINUATION_HEADING_PT))
        tf = box.text_frame
        tf.word_wrap = True
        # Zero the default 3.6pt top/bottom insets -- a one-line heading is
        # tight enough that they alone pushed it over its own box height.
        try:
            tf.margin_top = 0
            tf.margin_bottom = 0
        except Exception:
            pass
        p = tf.paragraphs[0]
        try:
            p.left_indent = Inches(0.15)
            p.first_line_indent = Inches(-0.15)
            p.space_before = Pt(0)
            p.space_after = Pt(0)
            p.line_spacing = 1.0
            self._apply_east_asian_line_breaking(p)
        except Exception:
            pass
        run_bullet = p.add_run()
        run_bullet.text = '■ '
        run_bullet.font.size = Pt(9)
        run_bullet.font.name = 'Arial'
        self._set_east_asian_typeface(run_bullet)
        self._declare_run_language(run_bullet)
        try:
            run_bullet.font.color.rgb = RGBColor(128, 128, 128)
        except Exception:
            pass
        run_name = p.add_run()
        run_name.text = f"{display_name}{suffix}"
        run_name.font.size = Pt(9)
        run_name.font.name = 'Arial'
        self._set_east_asian_typeface(run_name)
        self._declare_run_language(run_name)
        run_name.font.bold = True
        try:
            run_name.font.color.rgb = RGBColor(0, 0, 0)
        except Exception:
            pass
        self._force_no_autofit(tf)
        return top + int(Pt(self._TABLE_CONTINUATION_HEADING_PT))

    def _render_table_accounts_stack(
        self, slide, bullets_shape, account_data_list: List[Dict[str, Any]],
        is_chinese_databook: bool, statement_type: Optional[str] = None,
    ) -> None:
        """Renders one or more accounts stacked within a single slot -- one
        or more, since _append_table_accounts_to_distribution now packs
        multiple small table accounts into the same column when there's
        room, rather than always giving every table account a whole column
        regardless of how little of it real content needs (a real Crescent
        export showed 税金及附加 -- 3 components, one short explanation --
        leaving most of its column empty), AND now also flows plain
        (non-table) trailing accounts into a table's own leftover space --
        e.g. 投资收益/营业外支出 after 财务费用, matching how the real deck's
        own last IS page reads (table, then more prose, same column).

        Each account gets its own dedicated lead-in textbox (mimicking the
        "■ key - text" bullet style via _fill_text_main_bullets_with_
        category_and_key, but writing into a FRESH text frame rather than
        the slot's shared one), sized to its own measured content. A
        table-bearing account then gets its table and any post-table
        explanation immediately below (_render_presentation_table); a plain
        trailing account (no `_presentation_table`) stops after its own
        lead-in -- there is nothing else to render for it. Either way the
        next account in account_data_list starts at a running vertical
        offset below whatever the previous one drew. This is deliberately
        NOT the shared-text-frame/paragraph-flow mechanism ordinary
        accounts use below -- that assumes a slot's whole content lives in
        one text frame with nothing else interleaved, which doesn't hold
        once a table (a separate shape) sits between one account's lead-in
        and the next account's own lead-in: writing a second account's
        lead-in as another paragraph in the same shared frame would render
        it flowing directly under the FIRST account's lead-in, ignoring
        that account's table sitting (as a different shape) in between.

        bullets_shape's own text frame is left empty (matches how an
        intentionally-unused slot is already handled elsewhere) -- every
        real line of content here lives in a shape this method creates.
        Pixel-perfect positioning isn't achievable from python-pptx alone
        -- PowerPoint itself does the actual text layout, not this library
        -- so each account's own height is a reasoned estimate with the
        same margin the rest of this feature uses, not a ground-truth
        read-back.
        """
        left = int(bullets_shape.left)
        width = int(bullets_shape.width)
        current_category = None

        # ---- One shared text frame for the whole column ------------------
        # A table is a separate shape and cannot flow inside a text frame, so
        # its vertical space is RESERVED with blank paragraphs and the table
        # is floated over them. Everything else -- each account's lead-in,
        # its post-table explanation, and the NEXT account's lead-in -- stays
        # in one continuous flow inside bullets_shape's own frame.
        #
        # The previous design gave every account its own textbox stacked down
        # the column. That could never reuse the space a table left over, and
        # when a table did not fit it stranded the account's "...明细如下："
        # with nothing under it. Reserving the space inside the shared frame
        # removes both: text simply resumes below the table, in the same box.
        tf = bullets_shape.text_frame
        try:
            # MUST come first. The template ships these slots carrying
            # "Placeholder – placeholder"; the previous per-account-textbox
            # renderer cleared it and the shared-frame rewrite dropped that
            # call. Left in, it is a real rendered LINE that every table's
            # position was then computed without: a real deck put the table
            # 10.8pt (exactly one line) above where the text actually
            # started, so it clipped the lead-in by 5.8pt on every table
            # column. It also leaks the literal placeholder text into the
            # deliverable.
            tf.clear()
            tf.word_wrap = True
            self._force_no_autofit(tf)
            from pptx.enum.text import MSO_VERTICAL_ANCHOR
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
        except Exception:
            pass
        template_empty_p = tf.paragraphs[0]._p

        shape_name = self._shape_name(bullets_shape) or ""
        slot_name = ("L" if shape_name.endswith("_L")
                     else "R" if shape_name.endswith("_R") else "single")
        std_lh_pt = self._planning_std_lh_pt(is_chinese_databook)
        line_pitch_pt = max(1.0, std_lh_pt - self._real_para_gap_pt(is_chinese_databook))
        _usable_pt, _inset_pt = self._textbox_usable_and_inset_pt(bullets_shape)
        top_inset_pt = _inset_pt / 2.0

        cursor_pt = 0.0
        deferred_tables: List[Tuple[Dict[str, Any], Dict[str, Any], float]] = []

        def _measured_content_pt() -> float:
            """Height of everything CURRENTLY in the frame, measured from the
            frame's own paragraphs.

            Deliberately re-measures the whole frame instead of accumulating
            per-block estimates. A block estimate and the frame's real
            rendered height only have to disagree by one wrapped line for the
            table floated over the band below to land on top of real text --
            and the two are measured by different code paths, on different
            machines, with different font metrics (this repo's own checker
            falls back to system fonts where production uses client metrics),
            so they WILL drift. Measuring what is actually in the frame is
            the only version that can't. Same walk inspect_pptx.py and
            inspect_table_bands.py perform, so all three agree by
            construction."""
            total = 0.0
            for para in tf.paragraphs:
                ptext = para.text or ""
                sizes = [r.font.size.pt for r in para.runs if r.font.size is not None]
                gap = para.space_after.pt if para.space_after is not None else 0.0
                if not ptext.strip():
                    total += (max(sizes) if sizes else 9.0) * POWERPOINT_LINE_PITCH_FACTOR + gap
                    continue
                total += self._measure_paragraph_pt(ptext, bullets_shape, is_chinese_databook) + gap
            return total

        from fdd_utils.text_metrics import POWERPOINT_LINE_PITCH_FACTOR

        def _reserve_blank_lines(count: int, font_pt: float = 9.0) -> None:
            """`count` blank paragraphs of exactly font_pt * 1.2 each -- no
            paragraph gap and no text, so the reserved height is simply
            count * that pitch and the table dropped on top of it lands
            where the arithmetic says it will."""
            for _ in range(max(0, int(count))):
                blank_p = tf.add_paragraph()
                try:
                    blank_p.space_before = Pt(0)
                    blank_p.space_after = Pt(0)
                    blank_p.line_spacing = 1.0
                    self._apply_east_asian_line_breaking(blank_p)
                except Exception:
                    pass
                blank_run = blank_p.add_run()
                blank_run.text = " "
                blank_run.font.size = Pt(font_pt)
                blank_run.font.name = 'Arial'
                self._set_east_asian_typeface(blank_run)
                self._declare_run_language(blank_run)

        def _reserve_exact(needed_pt: float) -> float:
            """Reserve `needed_pt` almost exactly, and return what was
            actually reserved.

            Rounding every reservation UP to a whole 9pt line left up to a
            full blank line visible under a table -- reported and manually
            confirmed on a real deck. A line's height is just its font size
            x 1.2, so the leftover fraction gets its own blank paragraph at
            whatever font size makes it come out right (floored at
            PowerPoint's own 1pt minimum). Over-reservation drops from up to
            10.8pt to at most ~1.2pt."""
            full = int(needed_pt // line_pitch_pt)
            _reserve_blank_lines(full)
            reserved = full * line_pitch_pt
            remainder = needed_pt - reserved
            if remainder > 0.05:
                font_pt = max(1.0, remainder / POWERPOINT_LINE_PITCH_FACTOR)
                _reserve_blank_lines(1, font_pt)
                reserved += font_pt * POWERPOINT_LINE_PITCH_FACTOR
            return reserved

        for account_data in account_data_list:
            table = account_data.get("_presentation_table")
            category = account_data.get('category', '')
            mapping_key = account_data.get('mapping_key', account_data.get('account_name', ''))
            display_name = (
                account_data.get('display_name_zh') or account_data.get('display_name', mapping_key)
                if is_chinese_databook
                else account_data.get('display_name', mapping_key)
            )
            commentary = account_data.get('commentary', '')
            clause_reviews = account_data.get('clause_reviews', [])
            is_chinese = account_data.get('is_chinese', False)

            # Which parts of this account render HERE. Absent (the normal
            # case) means the whole account. Set by
            # _append_table_accounts_to_distribution's flow().
            render_parts = account_data.get("_render_parts")
            wants_lead = render_parts is None or "lead" in render_parts
            wants_table = render_parts is None or "table" in render_parts
            wants_expl = render_parts is None or "expl" in render_parts
            post_table_text = account_data.get("_post_table_text", "")

            category_to_write = category if (category and category != current_category) else None
            if category:
                current_category = category

            if wants_lead:
                try:
                    self._fill_text_main_bullets_with_category_and_key(
                        tf, category_to_write, display_name, commentary, is_chinese,
                        is_chinese_databook=is_chinese_databook, needs_continuation=False,
                        font_size_pt=9, clause_reviews=clause_reviews,
                    )
                    # The template ships one empty paragraph; drop it as soon
                    # as real content exists so every measurement below sees
                    # only what will render.
                    if (template_empty_p is not None
                            and template_empty_p.getparent() is not None
                            and not (template_empty_p.text or "").strip()):
                        template_empty_p.getparent().remove(template_empty_p)
                        template_empty_p = None
                    cursor_pt = _measured_content_pt()
                except Exception as exc:
                    logger.warning("Could not render lead-in for account %s: %s", mapping_key, exc)
            else:
                # Continued fragment: a compact "科目名（续）" line stands in
                # for the intro sentence that ended the previous column.
                self._append_continuation_line_to_frame(
                    tf, display_name, is_chinese_databook, std_lh_pt,
                )
                cursor_pt = _measured_content_pt()

            if table and wants_table:
                deferred_tables.append(
                    (table, account_data, cursor_pt + self._TABLE_GAP_ABOVE_PT)
                )
                cursor_pt += _reserve_exact(
                    self._table_block_reserved_pt(table)
                )

            if post_table_text and wants_expl:
                self._append_explanation_to_frame(
                    tf, post_table_text, is_chinese_databook, std_lh_pt,
                    shape=bullets_shape, slot_name=slot_name, statement_type=statement_type,
                )
                cursor_pt = _measured_content_pt()

        # The frame's own template paragraph was never counted in cursor_pt,
        # so it has to go before the tables are positioned against it -- one
        # stray empty line at the top would shift every reserved band down.
        try:
            if (template_empty_p is not None
                    and template_empty_p.getparent() is not None
                    and not (template_empty_p.text or "").strip()):
                template_empty_p.getparent().remove(template_empty_p)
        except Exception:
            pass

        base_top = int(bullets_shape.top) + int(Pt(top_inset_pt))
        for table, account_data, y_pt in deferred_tables:
            try:
                self._render_presentation_table(
                    slide, left, base_top + int(Pt(y_pt)), width, table, is_chinese_databook,
                    source_multiplier=self._table_source_multiplier(account_data),
                    # The explanation is now real flowing text in the shared
                    # frame, not a floating box under the table.
                    post_table_text="",
                )
            except Exception as exc:
                logger.warning(
                    "Could not render presentation table for %s: %s",
                    account_data.get("mapping_key", ""), exc,
                )

    # Cell padding (0.04in left + 0.04in right, matching _set_cell's
    # margin_left/margin_right) and the child-row left indent (0.12in,
    # matching _set_cell's indent_emu for kind=="child") -- both added on
    # top of the raw measured text width below.
    _TABLE_CELL_PADDING_PT = 5.76
    _TABLE_CHILD_INDENT_PT = 8.64
    _TABLE_MIN_COLUMN_PT = 25.2  # 0.35in floor, guards a pathologically short column

    # A textbox's own top+bottom insets (OOXML default 0.05in each = 7.2pt
    # total). Height formulas that size a box from its TEXT's measured
    # height must add this back, or the box ends up exactly one inset
    # short of holding its own content -- the direct cause of lead-in and
    # explanation boxes sitting at 88-101% fill (and a real 101% OVERFLOW
    # RISK on a real export) instead of the intended ~80%.
    _TEXTBOX_INSET_PT = 7.2

    # Lets a table account's intro sentence finish one column while its
    # table starts the next (see place_table_item). Flip to False to fall
    # straight back to the previous whole-block-only behaviour without
    # touching anything else.
    _TABLE_ALLOW_LEAD_TABLE_SPLIT = True
    # Budget for the "科目名（续）" heading drawn above a continued fragment.
    # Must hold one full 9pt std_lh line (font x 1.2 pitch + paragraph gap
    # ~= 13pt) with the render safety margin -- a first attempt at 16pt
    # measured 148% full on real client metrics, because the textbox's own
    # default 3.6pt top/bottom insets eat into it. The heading renders with
    # those insets zeroed (see _render_continuation_heading) so this budget
    # is available to the text itself.
    _TABLE_CONTINUATION_HEADING_PT = 18.0

    def _planning_std_lh_pt(self, is_chinese: bool) -> float:
        """One std_lh unit as RENDER actually produces it, for the
        shape-less planning estimates.

        The estimates used to compute font_size x line_spacing + para_gap
        (9 x 1.0 + 2.2 = 11.2pt), but PowerPoint's real line pitch is
        1.2 x the point size (POWERPOINT_LINE_PITCH_FACTOR -- researched
        and confirmed separately, see project memory), so render uses
        10.8 + 2.2 = 13.0pt. Planning was therefore under-estimating every
        block by ~14%, which _TABLE_SLOT_PACK_THRESHOLD was quietly
        absorbing -- two compensating errors that together left real
        columns filled to only ~60% of their true capacity."""
        from fdd_utils.text_metrics import POWERPOINT_LINE_PITCH_FACTOR
        return (self._real_font_size_pt(is_chinese) * POWERPOINT_LINE_PITCH_FACTOR
                * self._real_line_spacing(is_chinese) + self._real_para_gap_pt(is_chinese))

    def _slot_capacity_pt(self, slide_idx: int, slot_name: str, start_slide: int) -> float:
        """The REAL usable height of one commentary slot, read from the
        template shape it will actually render into. Falls back to the
        shared constant when the shape can't be resolved.

        A single shared constant can't be right for both slot types -- a
        real template's first-page "single" slot is ~330pt while its
        continuation L/R slots are ~359pt, so one number is either 8%
        too small for one or too large for the other."""
        try:
            slide = self.presentation.slides[start_slide - 1 + slide_idx]
            shape = self._resolve_commentary_slot_shape(slide, slot_name)
            if shape is not None:
                usable_pt, _inset_pt = self._textbox_usable_and_inset_pt(shape)
                if usable_pt > 1.0:
                    return usable_pt
        except Exception as exc:
            logger.debug("Could not resolve real capacity for slot %s/%s: %s",
                         slide_idx, slot_name, exc)
        return self._TABLE_SLOT_CAPACITY_PT

    def _measurement_slot_shape(self):
        """A representative commentary slot shape, used ONLY to give text
        measurement a real width when a block's height must be estimated
        before its column has been chosen.

        Passing shape=None instead makes the measurer fall back to a default
        width, and on real content that is worth a whole wrapped line: a
        measured lead-in came out 24.1pt against a real 35.4pt, which is
        exactly the amount a table then overshot the column by. Resolving a
        specific slot would be circular (the estimate decides the column),
        but it doesn't need to be: every commentary slot in the template is
        the same 4.78in wide -- only their HEIGHTS differ -- and wrapping
        depends on width alone."""
        cached = getattr(self, "_measurement_slot_shape_cache", "unset")
        if cached != "unset":
            return cached
        shape = None
        try:
            for slide in self.presentation.slides:
                for candidate in slide.shapes:
                    if (self._shape_name(candidate) or "").startswith("textMainBullets"):
                        shape = candidate
                        break
                if shape is not None:
                    break
        except Exception as exc:
            logger.debug("Could not resolve a measurement slot shape: %s", exc)
            shape = None
        self._measurement_slot_shape_cache = shape
        return shape

    @staticmethod
    def _rendered_bullet_label(account_data: Dict[str, Any], is_chinese_databook: bool) -> str:
        """The label a bullet ACTUALLY renders with ("■ <label> - ...").

        Cost estimates must measure this, not the raw mapping_key: in a
        Chinese deck the mapping_key is the English short code
        ("Tax and Surcharges"), which is far wider than the Chinese name
        that really renders ("税金及附加") -- 352pt vs 315pt against a
        329.8pt box for one real lead-in, i.e. the estimate believed the
        line wrapped when it doesn't. Every such lead-in box came out one
        whole line too tall, which is exactly the "height 似乎是固定的...
        表格不是緊貼comments" the user reported."""
        mapping_key = account_data.get("mapping_key", account_data.get("account_name", ""))
        if is_chinese_databook:
            return account_data.get("display_name_zh") or account_data.get("display_name", mapping_key)
        return account_data.get("display_name", mapping_key)

    @staticmethod
    def _textbox_usable_and_inset_pt(shape) -> Tuple[float, float]:
        """(usable text height, total top+bottom inset) in points for a
        shape, read from its real bodyPr insets. Falls back to the OOXML
        default when they aren't declared."""
        raw_pt = int(shape.height) / 12700
        try:
            from fdd_utils.text_metrics import text_box_from_shape
            usable_pt = text_box_from_shape(shape).height_pt
        except Exception:
            usable_pt = max(1.0, raw_pt - PowerPointGenerator._TEXTBOX_INSET_PT)
        return usable_pt, max(0.0, raw_pt - usable_pt)

    @staticmethod
    def _table_unit_label(is_chinese_databook: bool) -> str:
        return "人民币千元" if is_chinese_databook else "CNY'000"

    @staticmethod
    def _table_source_multiplier(account_data: Dict[str, Any]) -> float:
        """The account's own raw-yuan -> display-unit divisor. Single
        definition, since the renderer, the width precompute and the
        AI-prompt side all need the identical value (an earlier 1000x
        display bug came from exactly this being derived twice)."""
        financial_data = (account_data or {}).get("financial_data")
        if hasattr(financial_data, "attrs"):
            return financial_data.attrs.get("source_multiplier") or 1
        return 1

    @classmethod
    def _build_presentation_table_plan(
        cls, table: Dict[str, Any], is_chinese_databook: bool, source_multiplier: float,
    ) -> List[Dict[str, Any]]:
        """Flattens a presentation table's rows -> children (indented) ->
        total into the single ordered render plan both the renderer and
        the uniform-width precompute measure against. Values are divided
        back down to display units here (see _render_presentation_table's
        docstring for why that division belongs at display time)."""
        divisor = source_multiplier if source_multiplier and source_multiplier != 0 else 1

        def _scaled(values: Dict[str, float]) -> Dict[str, float]:
            return {period: (v / divisor if isinstance(v, (int, float)) else v)
                    for period, v in (values or {}).items()}

        plan: List[Dict[str, Any]] = []
        for row in (table.get("rows") or []):
            plan.append({"label": row.get("label", ""), "values": _scaled(row.get("values")), "kind": "data"})
            for child in (row.get("children") or []):
                plan.append({"label": child.get("label", ""), "values": _scaled(child.get("values")), "kind": "child"})
        total_row = table.get("total_row")
        if total_row:
            plan.append({"label": total_row.get("label", "合计" if is_chinese_databook else "Total"),
                         "values": _scaled(total_row.get("values")), "kind": "total"})
        return plan

    def _clamp_column_widths_to_available(
        self, widths_pt: List[float], available_pt: Optional[float],
    ) -> List[float]:
        """Floors each column at the legibility minimum, then scales the
        whole set down proportionally if it would overflow the slot."""
        widths_pt = [max(self._TABLE_MIN_COLUMN_PT, w) for w in widths_pt]
        total_pt = sum(widths_pt)
        if available_pt is not None and total_pt > available_pt and total_pt > 0:
            scale = available_pt / total_pt
            widths_pt = [w * scale for w in widths_pt]
        return widths_pt

    def _precompute_uniform_table_column_widths(
        self, table_items: List[Dict[str, Any]], is_chinese_databook: bool,
    ) -> None:
        """Computes one shared set of column widths per column-count, as
        the element-wise MAX of every table's own measured need, and
        caches it on the instance for _render_presentation_table.

        Sizing each subtable independently (80e7a70) correctly guaranteed
        no cell ever wraps, but made two subtables on the SAME page render
        at visibly different widths whenever their longest label differed
        -- flagged by the user as "表格要固定一種format 不要一個大一個細".
        Taking the max (rather than, say, a fixed width) keeps the
        no-wrap guarantee intact by construction: every table gets at
        least the width its own content needed.

        Deliberately grouped by column count -- tables with different
        period counts aren't visually comparable side by side anyway, and
        forcing them to a shared width would mean padding or squeezing
        columns that have no counterpart."""
        by_cols: Dict[int, List[float]] = {}
        for item in (table_items or []):
            table = item.get("_presentation_table")
            if not table:
                continue
            periods = table.get("periods") or []
            try:
                plan = self._build_presentation_table_plan(
                    table, is_chinese_databook, self._table_source_multiplier(item),
                )
                widths = self._measure_presentation_table_column_widths_pt(
                    plan, periods, table.get("period_labels") or {},
                    self._table_unit_label(is_chinese_databook), is_chinese_databook,
                    available_pt=None,   # raw need; the per-slot clamp happens at render
                )
            except Exception as exc:
                logger.debug("Could not measure uniform width for %s: %s", item.get("mapping_key"), exc)
                continue
            n_cols = len(widths)
            existing = by_cols.get(n_cols)
            by_cols[n_cols] = widths if existing is None else [max(a, b) for a, b in zip(existing, widths)]
        self._uniform_table_col_widths_pt = by_cols

    def _measure_presentation_table_column_widths_pt(
        self, plan: List[Dict[str, Any]], periods: List[str], period_labels: Dict[str, str],
        unit_label: str, is_chinese_databook: bool, available_pt: Optional[float] = None,
    ) -> List[float]:
        """Real-glyph-metrics column widths for a presentation (subtable)
        breakdown -- guarantees every cell's content fits on ONE line at its
        real rendered font size, rather than 2c6e2b1's crude per-character-
        count estimate. That estimate is what actually broke: it under-
        guessed real glyph width, a column ended up narrower than its
        content needed, PowerPoint wrapped the cell to a second line, and
        the row auto-grew past the fixed _TABLE_*_ROW_PT height that
        assumed one line -- text looked shrunk/cramped and the table's
        actual rendered height stopped matching _presentation_table_height_
        pt's estimate (the two symptoms that got 2c6e2b1 reverted).
        Reuses the SAME text_metrics measurer commentary sizing already
        trusts instead of a second, separate heuristic, so "guaranteed
        no-wrap" is a real guarantee, not another guess.

        Returns one width in points per column (label column first),
        summing to at most `available_pt` -- when content genuinely needs
        less than that (the common case, and the whole point: the real
        deck's own subtables don't fill their column either), the returned
        widths sum to LESS, and the caller renders a narrower, left-aligned
        table rather than stretching to fill the slot."""
        from fdd_utils.text_metrics import get_measurer

        packing = self._packing_settings()
        family = self._measurer_family(is_chinese_databook, packing)
        metrics_path = self._resolve_font_metrics_path(is_chinese_databook, packing)
        measurer_header = get_measurer(family, 7.5, is_cjk=is_chinese_databook, metrics_path=metrics_path)
        measurer_data = get_measurer(family, 7.0, is_cjk=is_chinese_databook, metrics_path=metrics_path)

        # Column 0 (label): widest of the unit-label header and every row's
        # own label text (plus child indent when that row is indented).
        label_candidates_pt = [measurer_header.text_width_pt(unit_label)]
        for entry in plan:
            indent_pt = self._TABLE_CHILD_INDENT_PT if entry.get("kind") == "child" else 0.0
            label_candidates_pt.append(measurer_data.text_width_pt(entry.get("label", "")) + indent_pt)
        widths_pt = [max(label_candidates_pt) + self._TABLE_CELL_PADDING_PT]

        # One column per period: widest of its header label and every row's
        # own formatted value in that column.
        for period in periods:
            candidates_pt = [measurer_header.text_width_pt(period_labels.get(period, period))]
            for entry in plan:
                value = entry.get("values", {}).get(period)
                text_val = self._format_table_value(value, is_numeric_column=True) if value is not None else ""
                candidates_pt.append(measurer_data.text_width_pt(text_val))
            widths_pt.append(max(candidates_pt) + self._TABLE_CELL_PADDING_PT)

        # available_pt=None returns the RAW measured need with no slot
        # clamp -- used by _precompute_uniform_table_column_widths, which
        # needs comparable per-table numbers before any per-slot scaling.
        return self._clamp_column_widths_to_available(widths_pt, available_pt)

    def _measure_paragraph_pt(self, text: str, shape, is_chinese: bool) -> float:
        """Rendered height of ONE paragraph's own lines (excluding its
        space_after), measured with the same wrap rules the renderer uses:
        a "■ ..." bullet hangs, so its FIRST line spans the full box width
        and only continuation lines are narrower."""
        try:
            from fdd_utils.text_metrics import get_measurer, text_box_from_shape
            packing = self._packing_settings()
            measurer = get_measurer(
                self._measurer_family(is_chinese, packing),
                self._real_font_size_pt(is_chinese), is_cjk=is_chinese,
                line_spacing=self._real_line_spacing(is_chinese),
                metrics_path=self._resolve_font_metrics_path(is_chinese, packing),
            )
            box = text_box_from_shape(shape)
            hang_w = max(10.0, box.width_pt - self._BULLET_HANGING_INDENT_PT)
            n = max(1, len(measurer.wrap(
                text, hang_w,
                first_line_width_pt=box.width_pt if text.lstrip().startswith("■") else None,
            )))
            return n * measurer.line_height_pt()
        except Exception as exc:
            logger.debug("Could not measure paragraph: %s", exc)
            return self._planning_std_lh_pt(is_chinese)

    def _append_continuation_line_to_frame(
        self, text_frame, display_name: str, is_chinese_databook: bool, std_lh_pt: float,
    ) -> float:
        """Writes "科目名（续）" as a paragraph in the shared column frame and
        returns the height it consumes. The standalone
        _render_continuation_heading textbox is only used by callers that
        still draw their own boxes; inside the shared frame a heading has to
        be a real paragraph or the text after it would not flow past it."""
        label = f"{display_name}（续）" if is_chinese_databook else f"{display_name} (cont'd)"
        para = text_frame.add_paragraph()
        try:
            para.space_before = Pt(0)
            para.space_after = Pt(3)
            para.line_spacing = 1.0
            self._apply_east_asian_line_breaking(para)
        except Exception:
            pass
        run = para.add_run()
        run.text = label
        run.font.size = Pt(9)
        run.font.name = 'Arial'
        self._set_east_asian_typeface(run)
        self._declare_run_language(run)
        run.font.bold = True
        try:
            run.font.color.rgb = RGBColor(0, 0, 0)
        except Exception:
            pass
        return std_lh_pt

    @staticmethod
    def _explanation_render_text(post_table_text: str, is_chinese_databook: bool) -> str:
        """The post-table explanation exactly as it RENDERS -- one marker-
        prefixed line per source line. Single definition so the planner
        measures the same string the renderer writes; the two-character
        "➢ " prefix is worth a whole wrapped line on a full-width line."""
        marker = "➢ " if is_chinese_databook else "- "
        raw_lines = [ln.strip() for ln in (post_table_text or "").split("\n") if ln.strip()]
        if not raw_lines:
            raw_lines = [(post_table_text or "").strip()]
        return "\n".join(
            ln if ln.startswith(("➢", "-", "•")) else f"{marker}{ln}" for ln in raw_lines
        )

    def _append_explanation_to_frame(
        self, text_frame, post_table_text: str, is_chinese_databook: bool, std_lh_pt: float,
        *, shape=None, slot_name: str = "single", statement_type: Optional[str] = None,
    ) -> float:
        """Writes the "➢"/"-" explanatory bullets that follow a presentation
        table as real paragraphs in the shared column frame, and returns the
        height they consume. Same content _render_post_table_explanation
        draws into its own textbox -- but as flowing text, so the next
        account's lead-in continues underneath instead of needing a whole
        new column."""
        lines = self._explanation_render_text(post_table_text, is_chinese_databook).split("\n")

        for line_text in lines:
            para = text_frame.add_paragraph()
            try:
                para.space_before = Pt(0)
                para.space_after = Pt(3)
                para.line_spacing = 1.0
                self._apply_east_asian_line_breaking(para)
            except Exception:
                pass
            run = para.add_run()
            run.text = line_text
            run.font.size = Pt(9)
            run.font.name = 'Arial'
            self._set_east_asian_typeface(run)
            self._declare_run_language(run)
            try:
                run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception:
                pass

        return self._calculate_content_lines(
            "", "", "\n".join(lines), slot_name=slot_name, shape=shape,
            is_chinese=is_chinese_databook, statement_type=statement_type,
        ) * std_lh_pt

    def _render_presentation_table(
        self, slide, left: int, top: int, width: int, table: Dict[str, Any],
        is_chinese_databook: bool, source_multiplier: float = 1, post_table_text: str = "",
    ) -> int:
        """Renders `table` (extract_presentation_detail_table's return
        shape) as a native table at the given EMU position, a source-line
        caption beneath it, and -- when the model wrote any (see ai.py's
        _detail_table_guidance and _split_table_commentary above) -- the
        "-"/"➢" explanatory bullets a real deliverable places below that,
        naming each named component's provider/charging basis/contract
        terms. Returns the bottom EMU position.

        table["rows"]/["total_row"] hold values in the SAME raw-yuan
        internal scale every account's df uses (normalize_financial_schedule
        multiplies by source_multiplier -- 1000 whenever the sheet's own
        header says CNY'000/千元 -- so cross-account math and the block's
        own tie-out against the account total both work in one consistent
        unit). The table itself is rendered under its own "人民币千元"/
        "CNY'000" header for a human to read directly, so those raw-yuan
        figures are divided back down here, at render time only -- dividing
        upstream in extract_presentation_detail_table would break its own
        tie-out against the (also raw-yuan-scale) account total.
        """
        periods = table.get("periods") or []
        period_labels = table.get("period_labels") or {}
        plan = self._build_presentation_table_plan(table, is_chinese_databook, source_multiplier)

        n_cols = 1 + len(periods)
        n_rows = 2 + len(plan)  # title + header + plan rows
        height = int(self._presentation_table_height_pt(table) * 12700)

        unit_label = self._table_unit_label(is_chinese_databook)
        available_pt = width / 12700.0
        # Prefer the deck-wide uniform widths when they've been precomputed
        # (see _precompute_uniform_table_column_widths): sizing each table
        # to its OWN content made otherwise-identical subtables render at
        # visibly different widths on the same page (2.94in vs 2.586in on
        # a real export), which the user flagged as "一個大一個細". The
        # uniform set is the element-wise MAX across every table sharing
        # this column count, so no table's own content can wrap under it.
        uniform = (getattr(self, "_uniform_table_col_widths_pt", None) or {}).get(n_cols)
        if uniform:
            column_widths_pt = self._clamp_column_widths_to_available(list(uniform), available_pt)
        else:
            column_widths_pt = self._measure_presentation_table_column_widths_pt(
                plan, periods, period_labels, unit_label, is_chinese_databook, available_pt,
            )
        table_width = max(int(Inches(0.7)), int(sum(column_widths_pt) * 12700))

        graphic_frame = slide.shapes.add_table(n_rows, n_cols, left, top, table_width, height)
        table_shape = graphic_frame.table
        for col_idx, col_width_pt in enumerate(column_widths_pt):
            table_shape.columns[col_idx].width = Pt(col_width_pt)
        style_id = self._resolve_table_style_id()
        if style_id:
            try:
                self._set_table_style_id(table_shape._tbl, style_id)
            except Exception as exc:
                logger.debug("Could not apply table style %s: %s", style_id, exc)

        DARK_BLUE = RGBColor(0x00, 0x33, 0x8D)
        WHITE = RGBColor(255, 255, 255)
        BLACK = RGBColor(0, 0, 0)
        GREY_TOTAL_FILL = RGBColor(0xD9, 0xD9, 0xD9)
        CHILD_BLUE = RGBColor(0x1F, 0x4E, 0x96)
        CHILD_ROW_FILL = RGBColor(0xF2, 0xF2, 0xF2)
        # The reference report's two banner rows are DIFFERENT blues sitting
        # directly on top of each other: a dark navy statement title, then a
        # lighter blue period header. Rendering both in the same navy (tried
        # first) merges them into one slab and loses the distinction the
        # reference draws; a separating rule loses the adjacency instead.
        # Different fills, no rule between them.
        HEADER_BAND_BLUE = RGBColor(0x1F, 0x4E, 0x96)

        def _set_cell(cell, text, *, bold=False, color=BLACK, fill=None, size_pt=7.0,
                      align=PP_ALIGN.LEFT, indent_emu=0):
            cell.text = text
            if not cell.text_frame.paragraphs:
                cell.text_frame.add_paragraph()
            p = cell.text_frame.paragraphs[0]
            if not p.runs:
                p.add_run()
            for run in p.runs:
                run.font.name = 'Arial'
                self._set_east_asian_typeface(run)
                self._declare_run_language(run)
                run.font.size = Pt(size_pt)
                run.font.bold = bold
                try:
                    run.font.color.rgb = color
                except Exception:
                    pass
            try:
                p.line_spacing = 1.0
                self._apply_east_asian_line_breaking(p)
                p.alignment = align
                if indent_emu:
                    self._set_paragraph_left_indent(p, indent_emu)
            except Exception:
                pass
            try:
                cell.margin_left = Inches(0.04)
                cell.margin_right = Inches(0.04)
                cell.margin_top = Inches(0.01)
                cell.margin_bottom = Inches(0.01)
                if fill is not None:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = fill
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = WHITE
            except Exception:
                pass

        # Row 0: title band, merged across every column.
        table_shape.rows[0].height = Pt(self._TABLE_TITLE_ROW_PT)
        if n_cols > 1:
            table_shape.cell(0, 0).merge(table_shape.cell(0, n_cols - 1))
        title_text = table.get("title") or ""
        _set_cell(table_shape.cell(0, 0), title_text, bold=True, color=WHITE, fill=DARK_BLUE, size_pt=8.0)

        # Row 1: period header -- navy fill, white bold text, matching the
        # reference deck (its "人民币千元 | 2023年 | ..." row is filled navy,
        # not white).
        #
        # This was reverted on 2026-08-04 as a PRECAUTION while isolating a
        # blank-page bug, not because it was implicated: those blank pages
        # were the BS/IS overview pages, which this function never draws on
        # -- they come from _fill_table_placeholder, whose own revert note
        # names the same two slides and which stays reverted. Slides 4-5,
        # the only ones this function touches, never went blank. Re-applied
        # here alone so the two can be told apart if it ever recurs.
        table_shape.rows[1].height = Pt(self._TABLE_HEADER_ROW_PT)
        _set_cell(table_shape.cell(1, 0), unit_label, bold=True,
                  color=WHITE, fill=HEADER_BAND_BLUE, size_pt=7.5)
        for j, period in enumerate(periods, start=1):
            _set_cell(table_shape.cell(1, j), period_labels.get(period, period),
                      bold=True, color=WHITE, fill=HEADER_BAND_BLUE, size_pt=7.5,
                      align=PP_ALIGN.CENTER)

        # Data / child / total rows.
        for row_idx, entry in enumerate(plan, start=2):
            kind = entry["kind"]
            is_total = kind == "total"
            is_child = kind == "child"
            row_h = (self._TABLE_TOTAL_ROW_PT if is_total
                     else self._TABLE_CHILD_ROW_PT if is_child
                     else self._TABLE_DATA_ROW_PT)
            table_shape.rows[row_idx].height = Pt(row_h)
            label_color = CHILD_BLUE if is_child else BLACK
            # Reference deck (IMG_0229, pixel-sampled): indented child
            # rows carry a light NEUTRAL-grey band (~5-7% darker than
            # the white parent rows, no blue cast) on top of their blue
            # text; parent/plain rows stay white.
            label_fill = (GREY_TOTAL_FILL if is_total
                          else CHILD_ROW_FILL if is_child
                          else None)
            _set_cell(table_shape.cell(row_idx, 0), entry["label"], bold=is_total,
                      color=label_color, fill=label_fill, size_pt=7.0,
                      indent_emu=int(Inches(0.12)) if is_child else 0)
            for j, period in enumerate(periods, start=1):
                value = entry["values"].get(period)
                text_val = self._format_table_value(value, is_numeric_column=True) if value is not None else ""
                _set_cell(table_shape.cell(row_idx, j), text_val, bold=is_total,
                          color=label_color, fill=label_fill, size_pt=7.0, align=PP_ALIGN.RIGHT)

        # Thin vertical borders between columns, a rule under the header
        # row, and top+bottom rules bracketing the total row -- matches the
        # BS/IS grid table's own convention (no rule on every data row).
        try:
            for r in range(n_rows):
                for c in range(n_cols):
                    cell = table_shape.cell(r, c)
                    if c > 0:
                        self._set_cell_border(cell, 'left', color_rgb=RGBColor(0xBF, 0xBF, 0xBF), width=Pt(0.5))
            for c in range(n_cols):
                self._set_cell_border(table_shape.cell(1, c), 'bottom', color_rgb=BLACK, width=Pt(1))
            total_row_idx = next((i for i, e in enumerate(plan, start=2) if e["kind"] == "total"), None)
            if total_row_idx is not None:
                for c in range(n_cols):
                    self._set_cell_border(table_shape.cell(total_row_idx, c), 'top', color_rgb=BLACK, width=Pt(1))
                    self._set_cell_border(table_shape.cell(total_row_idx, c), 'bottom', color_rgb=BLACK, width=Pt(1.25))
        except Exception as exc:
            logger.debug("Could not apply presentation-table borders: %s", exc)

        bottom = top + height
        source_box = slide.shapes.add_textbox(left, bottom, width, Pt(self._TABLE_SOURCE_LINE_PT + 2))
        source_tf = source_box.text_frame
        source_tf.word_wrap = True
        # Same spAutoFit problem as the explanation box below: left on, this
        # one-line 7pt caption grows itself and pushes everything under it down.
        self._force_no_autofit(source_tf)
        source_p = source_tf.paragraphs[0]
        source_run = source_p.add_run()
        source_run.text = "资料来源：管理层信息；毕马威分析" if is_chinese_databook else "Source: Management information; KPMG analysis"
        source_run.font.name = 'Arial'
        self._set_east_asian_typeface(source_run)
        self._declare_run_language(source_run)
        source_run.font.size = Pt(7.0)
        source_run.font.italic = True
        try:
            source_run.font.color.rgb = RGBColor(0x59, 0x59, 0x59)
        except Exception:
            pass
        after_source = bottom + int(Pt(self._TABLE_SOURCE_LINE_PT + 2))

        if not post_table_text:
            return after_source

        return self._render_post_table_explanation(
            slide, left, after_source, width, post_table_text, is_chinese_databook,
        )

    def _render_post_table_explanation(
        self, slide, left: int, top: int, width: int, post_table_text: str,
        is_chinese_databook: bool,
    ) -> int:
        """Draws the "➢"/"-" explanatory bullets that follow a presentation
        table. Extracted from _render_presentation_table so the same
        rendering can also be used on its own, when the explanation is
        carried into the NEXT column while its table stays behind (see
        _append_table_accounts_to_distribution's flow()). Returns the
        bottom EMU position."""
        marker = "➢ " if is_chinese_databook else "- "
        raw_lines = [ln.strip() for ln in post_table_text.split("\n") if ln.strip()]
        if not raw_lines:
            raw_lines = [post_table_text.strip()]
        lines = [ln if ln.startswith(("➢", "-", "•")) else f"{marker}{ln}" for ln in raw_lines]

        # Placeholder height, resized below via the same real-metrics
        # measurement _render_table_accounts_stack uses for the lead-in --
        # needs a real shape to measure against first.
        explain_box = slide.shapes.add_textbox(left, top, width, Pt(200))
        explain_tf = explain_box.text_frame
        explain_tf.word_wrap = True
        # python-pptx's add_textbox ships <a:spAutoFit/> ("resize shape to fit
        # text"), which makes PowerPoint DISCARD the height computed below and
        # substitute its own -- consistently taller, because it charges the
        # last paragraph's space_after plus both insets. That is the phantom
        # blank line under this box, and it is also why the box cannot be
        # dragged smaller by hand: PowerPoint snaps it straight back.
        self._force_no_autofit(explain_tf)
        for i, line_text in enumerate(lines):
            p = explain_tf.paragraphs[0] if i == 0 else explain_tf.add_paragraph()
            run = p.add_run()
            run.text = line_text
            run.font.name = 'Arial'
            self._set_east_asian_typeface(run)
            self._declare_run_language(run)
            run.font.size = Pt(9.0)
            try:
                run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception:
                pass
            p.line_spacing = 1.0
            self._apply_east_asian_line_breaking(p)
            p.space_after = Pt(2)

        try:
            combined = "\n".join(lines)
            used_units = self._calculate_content_lines(
                "", "", combined, slot_name="single", shape=explain_box,
                is_chinese=is_chinese_databook, whole_box=True,
            )
            capacity_units = self._calculate_max_lines_for_textbox(
                explain_box, is_chinese=is_chinese_databook, slot_name="single",
            )
            # See _render_table_accounts_stack's lead-in sizing: derive
            # std_lh from the USABLE height capacity was measured against,
            # then add the insets back on top of the text's own height.
            usable_pt, inset_pt = self._textbox_usable_and_inset_pt(explain_box)
            std_lh_pt = (
                (usable_pt / capacity_units) if capacity_units > 0 else
                self._real_font_size_pt(is_chinese_databook) * self._real_line_spacing(is_chinese_databook)
                + self._real_para_gap_pt(is_chinese_databook)
            )
            explain_height_pt = (
                max(used_units, 1.0) * std_lh_pt * self._TABLE_RENDER_HEIGHT_SAFETY_FACTOR + inset_pt
            )
            explain_box.height = int(explain_height_pt * 12700)
        except Exception as exc:
            logger.debug("Could not size presentation-table explanatory text: %s", exc)

        return top + int(explain_box.height)

    @staticmethod
    def _insert_category_header_rows(df, mappings: Optional[Dict[str, Any]], is_chinese_mode: bool):
        """Insert a blank-figures header row ("流动资产" / "Current assets"
        / etc.) into `df` whenever a leaf line item's mapped category
        (mappings.yml -- the SAME per-account "category" field the
        commentary bullets already group by) changes from the previous
        one. Total/subtotal rows (same keyword detection the later styling
        pass uses) never update the running category tracker and never
        trigger an insertion themselves -- a subtotal belongs to whatever
        category the items above it were in, not a category of its own.

        A real Financials-sheet check this session (inspect_financials_
        structure.py against the Kunshan databook) confirmed the RAW
        extracted sheet has no such header rows at all -- straight from a
        leaf item to "Total current assets" -- so this is what actually
        produces the reference format's ("IMG_0035") header rows, since
        nothing upstream of this table provides them on its own.

        Returns `df` unchanged if there's no mappings to categorise
        against (never silently drops rows in that case).
        """
        if not mappings or df is None or df.empty:
            return df

        total_keywords = list(
            {'total', '合计', '总计', '小计', 'subtotal', 'sub-total', 'sub total'}
            | set(SUMMARY_ACCOUNT_SKIP_KEYWORDS)
        )

        new_rows = []
        current_category = None
        for _, row in df.iterrows():
            label = str(row.iloc[0]).strip()
            label_lower = label.lower()
            is_total = any(kw in label_lower for kw in total_keywords)

            if not is_total and label:
                mapping_key = find_mapping_key(label, mappings)
                category = str((mappings.get(mapping_key) or {}).get('category', '') or '') if mapping_key else ''
                if category and category != current_category:
                    header_label = translate_category_to_chinese(category) if is_chinese_mode else category
                    header_row = {col: (header_label if i == 0 else pd.NA) for i, col in enumerate(df.columns)}
                    new_rows.append(header_row)
                    current_category = category

            new_rows.append(row.to_dict())

        return pd.DataFrame(new_rows, columns=df.columns)

    def _embed_statement_table(
        self, slide, df, statement_type: str, table_name: str = None, currency_unit: str = None,
        mappings: Optional[Dict[str, Any]] = None, is_chinese_mode: bool = False,
    ):
        # Insert category-header rows BEFORE anything below sizes the table
        # off len(df) -- the table's own row count (_add_table_to_slide /
        # the TablePlaceholder branch) is derived directly from len(df), so
        # doing this here means the extra rows are already accounted for by
        # the time the table shape itself gets created, not squeezed in
        # after the fact.
        df = self._insert_category_header_rows(df, mappings, is_chinese_mode)

        target_shape = self._resolve_table_target_shape(slide, statement_type)
        bounds = self._calculate_table_bounds(slide, target_shape=target_shape, statement_type=statement_type)
        target_name = self._shape_name(target_shape) if target_shape is not None else "(new table)"
        logger.info(
            f"Resolved {statement_type} table target on slide using {target_name} "
            f"at left={bounds['left']} top={bounds['top']} width={bounds['width']} height={bounds['height']}"
        )

        # Remove the redundant "Table" label box (TextBox 11 in the template).
        # The table's navy title row now lives at the same vertical position
        # as that label, so keeping the label would double-print the header.
        # Leave the right-side "Commentary" label intact — there's still a
        # commentary box on this slide that needs its header.
        for shape in list(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                label_text = (shape.text_frame.text or "").strip().lower()
            except Exception:
                continue
            if label_text != "table":
                continue
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception as e:
                logger.debug("Could not remove 'Table' label shape: %s", e)

        if target_shape is None:
            table_shape = self._add_table_to_slide(slide, df, bounds, table_name=table_name)
            self._fill_table_placeholder(
                table_shape,
                df,
                table_name=table_name,
                currency_unit=currency_unit,
                bounds=bounds,
                mappings=mappings,
                is_chinese_mode=is_chinese_mode,
            )
            return

        self._fill_table_placeholder(
            target_shape,
            df,
            table_name=table_name,
            currency_unit=currency_unit,
            bounds=bounds,
            mappings=mappings,
            is_chinese_mode=is_chinese_mode,
        )
    
    def find_content_shape(self, shapes):
        """Find content shape by trying multiple possible names"""
        # Try different possible names for content shapes
        possible_names = [
            'Content',
            'Text-commentary',
            'textMainBullets',
            'Text',
            'Commentary',
            'MainContent',
            'Body'
        ]
        
        for name in possible_names:
            shape = self.find_shape_by_name(shapes, name)
            if shape and shape.has_text_frame:
                return shape
        
        # If no named shape found, try to find any text frame shape that's not a title
        for shape in shapes:
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                shape_name = getattr(shape, 'name', '')
                # Skip title shapes and other non-content shapes
                if shape_name and 'title' not in shape_name.lower() and 'proj' not in shape_name.lower():
                    return shape
        
        return None

    def replace_text_preserve_formatting(self, shape, replacements: Dict[str, str]) -> bool:
        """Replace text while preserving formatting"""
        if not shape.has_text_frame:
            return False

        replaced = False

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                original_text = run.text
                for old_text, new_text in replacements.items():
                    if old_text in run.text:
                        run.text = run.text.replace(old_text, new_text)
                if run.text != original_text:
                    replaced = True

        if not replaced:
            current_text = shape.text_frame.text
            updated_text = current_text
            for old_text, new_text in replacements.items():
                updated_text = updated_text.replace(old_text, new_text)
            if updated_text != current_text:
                shape.text_frame.text = updated_text
                replaced = True

        return replaced

    def refresh_project_placeholders(self, project_name: str):
        """Refresh placeholder tokens such as [PROJECT], [Current], and [Total]."""
        if not self.presentation or not project_name:
            return

        display_entity = str(project_name).strip()
        total_slides = len(self.presentation.slides)
        if not display_entity or total_slides <= 0:
            return

        base_replacements = {
            "[PROJECT]": display_entity,
            "[Entity]": display_entity,
            "[Company]": display_entity,
        }

        for slide_index, slide in enumerate(self.presentation.slides):
            proj_title_shape = self.find_shape_by_name(slide.shapes, "projTitle")
            if not proj_title_shape or not proj_title_shape.has_text_frame:
                continue

            replacements = dict(base_replacements)
            replacements["[Current]"] = str(slide_index + 1)
            replacements["[Total]"] = str(total_slides)
            current_text = proj_title_shape.text

            if any(token in current_text for token in replacements):
                self.replace_text_preserve_formatting(proj_title_shape, replacements)

            # Chinese reports: the template's own title text is an English
            # scaffold ("Entity overview - Project [PROJECT] (1/4)") with
            # only the [PROJECT] token substituted above -- for a Chinese
            # entity name that reads as English label + Chinese name mixed
            # mid-title. Strip the English lead-in so the title is just the
            # entity name (plus any pagination suffix the template kept),
            # e.g. "昆明经开 (1/4)" instead of "Entity overview - Project
            # 昆明经开 (1/4)".
            if self.language == 'chinese':
                stripped = proj_title_shape.text
                stripped = re.sub(
                    r'(?i)^\s*entity\s+overview\s*[-–:]\s*(project\s*)?',
                    '', stripped,
                )
                if stripped != proj_title_shape.text:
                    self.replace_text_preserve_formatting(
                        proj_title_shape, {proj_title_shape.text: stripped},
                    )

    def update_project_titles(self, project_name: str, statement_type: str = 'BS'):
        """Update project titles in presentation"""
        if not self.presentation:
            return

        display_entity = str(project_name or "").strip()
        self.refresh_project_placeholders(display_entity)

        # Define title templates based on language and statement type
        if statement_type.upper() == 'BS':
            if self.language == 'chinese':
                title_template = f"资产负债表概览 - {display_entity}"
            else:
                title_template = f"Entity Overview - {display_entity}"
        elif statement_type.upper() == 'IS':
            if self.language == 'chinese':
                title_template = f"利润表概览 - {display_entity}"
            else:
                title_template = f"Income Statement - {display_entity}"
        else:
            if self.language == 'chinese':
                title_template = f"财务报表概览 - {display_entity}"
            else:
                title_template = f"Financial Report - {display_entity}"

        # Update titles in all slides
        for slide_index, slide in enumerate(self.presentation.slides):
            current_slide_number = slide_index + 1
            proj_title_shape = self.find_shape_by_name(slide.shapes, "projTitle")

            if proj_title_shape:
                current_text = proj_title_shape.text
                if "[PROJECT]" in current_text:
                    replacements = {
                        "[PROJECT]": display_entity,
                        "[Current]": str(current_slide_number),
                        "[Total]": str(len(self.presentation.slides))
                    }
                    self.replace_text_preserve_formatting(proj_title_shape, replacements)
                else:
                    # Replace the entire title
                    if proj_title_shape.has_text_frame:
                        proj_title_shape.text_frame.text = title_template

    def generate_full_report(self, markdown_content: str, summary_md: Optional[str] = None,
                           output_path: str = None):
        """Generate full PowerPoint report from markdown content"""
        if not self.presentation:
            self.load_template()

        # Process markdown content
        processed_content = self._process_markdown_content(markdown_content)

        # Apply content to presentation
        self._apply_content_to_presentation(processed_content)

        # Save if output path provided
        if output_path:
            self.save(output_path)

    def _process_markdown_content(self, content: str) -> Dict:
        """Process markdown content into structured data"""
        if not content:
            logger.warning("Empty content provided to _process_markdown_content")
            return {}

        logger.info("Processing markdown content, length: %s", len(content))
        logger.debug("Content preview (first 500 chars): %s", content[:500])

        # Split by headers (## Account Name)
        sections = re.split(r'^##\s+(.+)$', content, flags=re.MULTILINE)

        logger.info("Found %s sections after splitting", len(sections))

        processed_sections = {}

        # Process each section
        for i in range(1, len(sections), 2):
            if i + 1 < len(sections):
                account_name = sections[i].strip()
                account_content = sections[i + 1].strip()

                logger.info("Processing section: %s, content length: %s", account_name, len(account_content))

                processed_sections[account_name] = {
                    'content': account_content,
                    'is_chinese': detect_chinese_text(account_content)
                }

        logger.info("Processed %s sections", len(processed_sections))
        return processed_sections

    def _apply_content_to_presentation(self, sections: Dict):
        """Apply processed content to presentation slides"""
        if not self.presentation:
            logger.warning("No presentation loaded")
            return

        logger.info("Applying %s sections to presentation with %s slides", len(sections), len(self.presentation.slides))

        # Find content placeholders and fill them
        slide_idx = 0
        for slide in self.presentation.slides:
            if slide_idx >= len(sections):
                logger.warning("More slides (%s) than sections (%s)", len(self.presentation.slides), len(sections))
                break

            account_name = list(sections.keys())[slide_idx]
            section_data = sections[account_name]

            logger.info("Processing slide %s for account: %s", slide_idx + 1, account_name)

            # Find content shape using flexible name matching
            content_shape = self.find_content_shape(slide.shapes)
            if content_shape:
                logger.info("Found content shape '%s' on slide %s", content_shape.name, slide_idx + 1)
                if content_shape.has_text_frame:
                    # Apply content to shape
                    self._fill_content_shape(content_shape, section_data)
                    logger.info("Applied content to slide %s", slide_idx + 1)
                else:
                    logger.warning("Content shape found but has no text_frame on slide %s", slide_idx + 1)
            else:
                logger.warning("No content shape found on slide %s, available shapes: %s", slide_idx + 1, [s.name if hasattr(s, 'name') else 'unnamed' for s in slide.shapes])
                # Try to use the first available text frame as fallback
                for shape in slide.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        shape_name = getattr(shape, 'name', 'unnamed')
                        if 'title' not in shape_name.lower() and 'proj' not in shape_name.lower():
                            logger.info("Using fallback shape '%s' on slide %s", shape_name, slide_idx + 1)
                            self._fill_content_shape(shape, section_data)
                            break

            slide_idx += 1

    def _fill_content_shape(self, shape, section_data: Dict):
        """Fill content shape with processed data"""
        if not shape.has_text_frame:
            logger.warning("Shape does not have text_frame")
            return

        content = section_data.get('content', '')
        is_chinese = section_data.get('is_chinese', False)

        logger.info("Filling shape with content length: %s", len(content))

        # Clear existing content
        shape.text_frame.clear()
        
        if not content or not content.strip():
            logger.warning("No content to fill")
            return
        
        # Split content into paragraphs if it contains newlines
        content_lines = content.split('\n')
        
        # Add content with proper formatting
        for idx, line in enumerate(content_lines):
            line = line.strip()
            if not line and idx > 0:
                # Skip empty lines except add a paragraph break
                continue
            
            if idx == 0:
                # Use first paragraph or create one
                if shape.text_frame.paragraphs:
                    p = shape.text_frame.paragraphs[0]
                else:
                    p = shape.text_frame.add_paragraph()
            else:
                p = shape.text_frame.add_paragraph()
            
            p.text = line
            
            # Apply formatting to runs
            for run in p.runs:
                run.font.size = get_font_size_for_text(line, force_chinese_mode=is_chinese)
                run.font.name = get_font_name_for_text(line)

            # Set paragraph formatting
            p.space_after = get_space_after_for_text(line, force_chinese_mode=is_chinese)
            p.space_before = get_space_before_for_text(line, force_chinese_mode=is_chinese)
            p.line_spacing = get_line_spacing_for_text(line, force_chinese_mode=is_chinese)
        
        logger.info("Successfully filled shape with %s paragraphs", len([l for l in content_lines if l.strip()]))

    def _pillow_fitting_enabled(self, packing: Dict[str, Any]) -> bool:
        if os.environ.get("FDD_USE_PILLOW_FITTING") == "1":
            return True
        if os.environ.get("FDD_USE_PILLOW_FITTING") == "0":
            return False
        return bool(packing.get("use_pillow_text_fitting", False))

    def _resolve_font_metrics_path(self, is_chinese: bool, packing: Dict[str, Any]) -> Optional[str]:
        """Path to the client-font metrics.json (dumped via dump_font_metrics.py),
        so line-fitting measures with the font the client's PowerPoint renders.
        Language-specific key wins; falls back to a single shared path. Relative
        paths resolve against the repo root."""
        key = "font_metrics_path_chi" if is_chinese else "font_metrics_path_eng"
        path = packing.get(key) or packing.get("font_metrics_path")
        if not path:
            return None
        p = str(path)
        if not os.path.isabs(p):
            p = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), p)
        return p if os.path.exists(p) else None

    def _measurer_family(self, is_chinese: bool, packing: Dict[str, Any]) -> str:
        """System-font family for the Pillow fallback (overridable in config)."""
        key = "font_family_chi" if is_chinese else "font_family_eng"
        return str(packing.get(key) or ("Microsoft YaHei" if is_chinese else "Arial"))

    def _log_measurer_source_once(self, measurer, metrics_path: Optional[str], is_chinese: bool) -> None:
        """INFO-log the text-measurement source once per language per export, so a
        server log shows unambiguously whether client-font metrics are active."""
        key = "CHI" if is_chinese else "ENG"
        logged = getattr(self, "_measurer_sources_logged", None)
        if logged is None:
            logged = self._measurer_sources_logged = set()
        if key in logged:
            return
        logged.add(key)
        detail = f" ({metrics_path})" if measurer.source == "client-metrics" else ""
        logger.info("Text measurement [%s]: %s%s", key, measurer.source, detail)

    def _pillow_measure(
        self,
        shape,
        text: str,
        *,
        is_chinese: bool,
    ) -> Optional[Tuple[int, int]]:
        """Returns (used_lines, capacity_lines) using real font metrics, or
        None on any failure (caller falls back to legacy CPL heuristic)."""
        if not shape or not hasattr(shape, "height") or not hasattr(shape, "width"):
            return None
        try:
            from fdd_utils.text_metrics import (
                get_measurer,
                lines_that_fit,
                text_box_from_shape,
            )
        except Exception:
            return None
        try:
            packing = self._packing_settings(None)
            font_size_pt = self._real_font_size_pt(is_chinese)
            line_spacing = self._real_line_spacing(is_chinese)
            _mpath = self._resolve_font_metrics_path(is_chinese, packing)
            measurer = get_measurer(
                self._measurer_family(is_chinese, packing), font_size_pt,
                is_cjk=is_chinese, line_spacing=line_spacing,
                metrics_path=_mpath,
            )
            self._log_measurer_source_once(measurer, _mpath, is_chinese)
            box = text_box_from_shape(shape)
            capacity = lines_that_fit(box.height_pt, measurer.line_height_pt())
            if not text:
                return (0, capacity)
            lines = measurer.wrap(text, box.width_pt)
            return (len(lines), capacity)
        except Exception:
            return None

    # Space-after (pt) applied to every paragraph in _fill_text_main_bullets.
    # Matches _fill_text_main_bullets_with_category_and_key's own hardcoded
    # p_key.space_after = Pt(3) (see the "Matches _PARA_SPACE_AFTER (cost
    # estimator)" comment right there in that function) -- the ACTUAL,
    # currently-live commentary-bullet renderer, unlike get_space_after_for_
    # text/get_space_before_for_text below, which belong to a different,
    # legacy code path (_fill_content_shape, reached only from the unused
    # markdown generate() flow) and were never actually applied to a
    # textMainBullets run.
    _PARA_SPACE_AFTER = 3.0

    @staticmethod
    def _real_font_size_pt(is_chinese: bool) -> float:
        """Font size actually applied to the run (get_font_size_for_text) —
        a single deck-wide 9pt regardless of language, NOT the 10pt some
        capacity/content code used to assume for Chinese."""
        return get_font_size_for_text("", force_chinese_mode=is_chinese).pt

    @staticmethod
    def _real_line_spacing(is_chinese: bool) -> float:
        """Line spacing actually applied to a commentary bullet run.

        _fill_text_main_bullets_with_category_and_key hardcodes
        line_spacing = 1.0 on every paragraph it creates (category header,
        key line, and continuation lines alike) -- unconditionally, not
        gated on is_chinese at all. get_line_spacing_for_text's 0.9-for-
        Chinese value belongs to the separate, legacy _fill_content_shape
        path (markdown generate() flow) and was never actually the value
        applied to a live textMainBullets paragraph. A user-supplied real-
        client-metrics capacity check (inspect_single_slot.py against a
        real Windows export) directly caught this: assuming 0.9 line
        spacing + a 6-9pt inter-paragraph gap that never actually renders
        made the computed capacity roughly 30% smaller than the box's true
        capacity -- "the tool says 94% full" against a box the user could
        still visibly type 5-7 more lines into.
        """
        return 1.0

    @staticmethod
    def _real_para_gap_pt(is_chinese: bool) -> float:
        """Total vertical gap PowerPoint actually renders between two
        consecutive bullet paragraphs.

        _fill_text_main_bullets_with_category_and_key hardcodes
        space_before = Pt(0) and space_after = Pt(3) on every paragraph
        (category header, key line, continuation line) -- 3pt total,
        REGARDLESS of language. It never calls get_space_after_for_text /
        get_space_before_for_text at all; those getters' 4-9pt values
        belong to the separate legacy _fill_content_shape path. See
        _real_line_spacing's docstring for how this was actually caught.

        2.2, not 3.0 (2026-08-04): the requested space_after XML value is
        still literally Pt(3) at render time, unchanged -- this is not a
        claim that PowerPoint renders less than what's asked for. It's a
        correction to how much of that requested space this codebase's
        OWN capacity/content-cost formula should count against a box's
        available room, back-solved from real, empirical spare-capacity
        measurements the user made in real PowerPoint on two independent,
        differently-sized, differently-shaped boxes (a single-column table
        page's textMainBullets and a plain L-column continuation page) --
        both independently implied a real std_lh of ~13.0pt against this
        formula's previous 13.8pt (line_h 10.8, PROVEN correct separately,
        see POWERPOINT_LINE_PITCH_FACTOR's own history -- so the gap
        isolates to para_gap specifically: implied ~2.2pt, not 3.0pt).
        Deliberately not landing exactly on 2.2 without a second real
        cross-check on the render side too -- see the commit message this
        shipped in for the full reasoning and what still needs re-verifying.
        """
        return 3.0

    def _calculate_max_lines_for_textbox(
        self,
        shape,
        *,
        is_chinese: bool = False,
        slot_name: str = "single",
        statement_type: Optional[str] = None,
    ):
        """Return the number of 'line units' that fit in this text box.

        Measures the effective height directly from the shape and its
        text-frame insets (top/bottom margins), then divides by the standard
        line height used by ``_calculate_content_lines``:

            std_line_h = font_size × line_spacing + PARA_SPACE_AFTER

        Both capacity and content are expressed in the same unit so the fill
        ratios are accurate without any fudge factors.
        """
        packing = self._packing_settings(statement_type)
        if not shape or not hasattr(shape, "height"):
            return int(packing.get("minimum_slot_lines", 20) or 20)

        font_size_pt = self._real_font_size_pt(is_chinese)
        line_spacing = self._real_line_spacing(is_chinese)
        family       = self._measurer_family(is_chinese, packing)

        # ── Real font metrics via text_metrics ───────────────────────────────────
        # Prefer the client's font (metrics.json) so line height matches what the
        # client's PowerPoint renders; else the resolved system font. text_box_from_shape
        # reads bodyPr tIns/bIns directly from shape XML.
        try:
            from fdd_utils.text_metrics import get_measurer, text_box_from_shape
            _mpath   = self._resolve_font_metrics_path(is_chinese, packing)
            measurer = get_measurer(
                family, font_size_pt, is_cjk=is_chinese, line_spacing=line_spacing,
                metrics_path=_mpath,
            )
            self._log_measurer_source_once(measurer, _mpath, is_chinese)
            box      = text_box_from_shape(shape)
            std_lh   = measurer.line_height_pt() + self._real_para_gap_pt(is_chinese)
            # Float, not int(...) floored -- content cost (_calculate_content_
            # lines / _compute_slot_used_lines) is ALSO measured in these same
            # std_lh units and never rounds, so "N lines fit" and "the content
            # that's been packed in totals N.xx units" are already compared as
            # floats everywhere except here. Flooring only this side throws
            # away up to a full std_lh unit (line + its trailing gap) of real,
            # physically-there box height on every single slot -- verified via
            # a real Windows client-metrics export: a box the DP stopped
            # filling at "believed" 95.2% (15.23 / floor(16.77)=16) was
            # actually only 90.8% full relative to its own true height
            # (15.23 / 16.77) once the discarded 0.77-line margin is counted.
            # This value is exactly self-consistent with the box's real pt
            # height by construction (capacity_float * std_lh == box.height_pt),
            # so packing right up to it can never overflow the box -- the
            # floor was never a safety margin, just an unnecessary display-
            # style rounding that leaked into the fit-decision math.
            return max(1.0, box.height_pt / std_lh)
        except Exception:
            pass   # font file missing — fall through to heuristic

        # ── Heuristic fallback ───────────────────────────────────────────────────
        height_pt    = shape.height * 72 / 914400
        top_pt = bottom_pt = 3.6          # OOXML default tIns/bIns = 0.05" = 3.6 pt
        try:
            tf = shape.text_frame
            if tf.margin_top    is not None: top_pt    = tf.margin_top    * 72 / 914400
            if tf.margin_bottom is not None: bottom_pt = tf.margin_bottom * 72 / 914400
        except Exception:
            pass
        effective_pt = max(1.0, height_pt - top_pt - bottom_pt)
        std_lh       = font_size_pt * line_spacing + self._PARA_SPACE_AFTER
        max_rows     = int(effective_pt / std_lh)
        return max(int(packing.get("minimum_slot_lines", 20) or 20), max_rows)

    def _calculate_content_lines(
        self,
        category: str,
        mapping_key: str,
        commentary: str,
        *,
        slot_name: str = "single",
        shape=None,
        is_chinese: Optional[bool] = None,
        statement_type: Optional[str] = None,
        whole_box: bool = False,
    ) -> float:
        """Return the physical height of this content expressed in std_lh units.

        whole_box=True means this call measures ALL the text in one shape,
        so the final paragraph's trailing space_after can be dropped -- it
        renders as invisible padding at the bottom of the frame and pushes
        nothing. Default False, because the packer calls this ONCE PER
        ACCOUNT for slots that hold several: there, every account's
        trailing gap except the very last one is real spacing that
        separates it from the next account. Dropping it per-account
        under-counted a 7-account slot by 6 gaps (~1.3 lines), which is
        exactly the residual that had the DP reporting 100% on a slot
        PowerPoint renders at 103%.

        Returns a *float* (no ceil) so that the DP and greedy fill can track
        actual physical space consumed.  Using ceil was inflating every
        multi-line account to the next integer boundary, causing the DP to
        report 100 % fill when the box was only ~75 % physically used.

        One "unit" = std_lh = line_h + PARA_SPACE_AFTER (17 pt for English).
        Capacity from _calculate_max_lines_for_textbox is int(box_height/std_lh),
        so comparing float content against int capacity gives physically accurate
        fill ratios.
        """
        is_chinese = any('\u4e00' <= c <= '\u9fff' for c in commentary) if is_chinese is None else is_chinese

        # Memoize per-instance — Pillow font measurement runs ~80 calls per
        # paragraph and the same (account, slot, shape) tuple is asked many
        # times across the greedy distribute, the DP optimizer, and any
        # post-processing logging. Cache lookup is keyed on shape width
        # (the only shape attribute that affects line wrapping).
        if not hasattr(self, "_content_lines_cache"):
            self._content_lines_cache = {}
        shape_w = int(getattr(shape, "width", 0) or 0) if shape is not None else 0
        cache_key = (
            bool(category), mapping_key, commentary, slot_name, shape_w,
            is_chinese, str(statement_type or ""), whole_box,
        )
        cached = self._content_lines_cache.get(cache_key)
        if cached is not None:
            return cached

        font_size_pt = self._real_font_size_pt(is_chinese)
        line_spacing = self._real_line_spacing(is_chinese)
        para_gap     = self._real_para_gap_pt(is_chinese)
        packing = self._packing_settings(statement_type)
        family = self._measurer_family(is_chinese, packing)

        # -- Real glyph metrics via text_metrics --
        # Uses get_measurer (client metrics.json when configured, else the
        # resolved system font) -- the SAME measurer _calculate_max_lines_for_textbox
        # uses for capacity. Previously this called get_font()/wrap_paragraph()
        # directly with a hardcoded "Arial"/"Microsoft YaHei" system font,
        # ignoring font_metrics_path_eng/chi entirely -- capacity and content
        # were measured with two different rulers whenever a client metrics
        # file was configured, which is exactly the kind of quiet mismatch
        # that produces a "few rows off" fill-ratio gap despite already having
        # real font metrics available.
        if shape is not None:
            try:
                from fdd_utils.text_metrics import get_measurer, text_box_from_shape
                _mpath   = self._resolve_font_metrics_path(is_chinese, packing)
                measurer = get_measurer(
                    family, font_size_pt, is_cjk=is_chinese, line_spacing=line_spacing,
                    metrics_path=_mpath,
                )
                self._log_measurer_source_once(measurer, _mpath, is_chinese)
                box      = text_box_from_shape(shape)
                line_h   = measurer.line_height_pt()
                std_lh   = line_h + para_gap

                total_pt = 0.0
                if category:
                    total_pt += line_h          # category header: no space_after

                paras = [p for p in commentary.split('\n') if p.strip()] if commentary else []
                key_prefix = f"\u25a0 {mapping_key} - "
                # Hanging indent: wrapped CONTINUATION lines render 0.15"
                # narrower than the box (see _BULLET_HANGING_INDENT_PT), but
                # line 1 of the account's own first paragraph (p_key) spans
                # the FULL box width -- first_line_indent=-0.15" cancels
                # left_indent=0.15" for that one line only. Only paras[0]
                # (rendered as p_key) gets this exception; paras[1:] render
                # as p_text with first_line_indent=0, i.e. narrow throughout.
                wrap_w = max(10.0, box.width_pt - self._BULLET_HANGING_INDENT_PT)
                if paras:
                    first_wrapped = measurer.wrap(
                        key_prefix + paras[0], wrap_w, first_line_width_pt=box.width_pt,
                    )
                    total_pt += len(first_wrapped) * line_h + para_gap
                    for para in paras[1:]:
                        wrapped = measurer.wrap(para, wrap_w)
                        total_pt += len(wrapped) * line_h + para_gap
                else:
                    total_pt += line_h + para_gap

                # Drop the LAST paragraph's trailing gap: space_after on the
                # final paragraph renders as invisible padding at the
                # bottom of the frame and pushes nothing, so counting it
                # inflates every block by one para_gap. Confirmed against
                # PowerPoint's own BoundHeight on a real export -- a
                # 27-line, 8-paragraph box measured 309.6pt where this
                # function claimed ~327pt, and inspect_pptx.py's matching
                # copy of the same over-count was reporting a false
                # "102% OVERFLOW RISK" on a box that is really 93.5% full.
                if paras and whole_box:
                    total_pt -= para_gap

                # Return float -- no ceil so actual physical proportion is preserved.
                result = total_pt / std_lh
                self._content_lines_cache[cache_key] = result
                return result
            except Exception:
                pass    # font file missing -- fall through to heuristic

        # ── Heuristic fallback (no shape or font unavailable) ───────────────────
        space_after  = self._PARA_SPACE_AFTER
        std_lh       = font_size_pt * line_spacing + space_after
        cpl          = self._estimate_chars_per_line(slot_name, is_chinese, shape=shape,
                                                     statement_type=statement_type)
        total_pt     = 0.0
        if category:
            total_pt += font_size_pt * line_spacing
        paras = [p for p in commentary.split('\n') if p.strip()] if commentary else []
        key_pfx_len  = len(str(mapping_key)) + 5
        if paras:
            first_len    = key_pfx_len + len(paras[0])
            first_w      = max(1, (first_len + cpl - 1) // cpl)
            total_pt    += first_w * font_size_pt * line_spacing + space_after
            for para in paras[1:]:
                w         = max(1, (len(para) + cpl - 1) // cpl)
                total_pt += w * font_size_pt * line_spacing + space_after
        else:
            total_pt += font_size_pt * line_spacing + space_after
        result = total_pt / std_lh
        self._content_lines_cache[cache_key] = result
        return result

    def _distribute_content_across_slots(
        self,
        structured_data: List[Dict],
        max_slides: int = 4,
        start_slide: int = 1,
        statement_type: Optional[str] = None,
    ):
        """
        Distribute content across textbox slots based on capacity.
        Slot structure is derived from the actual template slides when they exist.
        Auto-added slides follow the template convention: page 1 keeps a single right-side
        commentary box beside the table, while later slides use left/right commentary slots.
        
        Returns: List of (slide_idx, slot_idx, [account_data], is_partial, continuation_of)
        """
        if not structured_data:
            return []

        # Find a textbox shape to calculate capacity
        sample_shape = None
        for slide in self.presentation.slides:
            for alt_name in ["textMainBullets", "textMainBullets_L", "textMainBullets_R"]:
                shape = self.find_shape_by_name(slide.shapes, alt_name)
                if shape:
                    sample_shape = shape
                    break
            if sample_shape:
                break
        max_lines_per_textbox = (
            self._calculate_max_lines_for_textbox(sample_shape, statement_type=statement_type)
            if sample_shape
            else 40
        )
        
        logger.info("\n%s", '='*80)
        logger.info("CONTENT DISTRIBUTION STARTING")
        logger.info("%s", '='*80)
        logger.info("Total accounts: %s", len(structured_data))
        logger.info("Max lines per textbox: %s", max_lines_per_textbox)
        if sample_shape:
            logger.info("Sample shape height: %.2f inches", sample_shape.height / 914400)
            logger.info("Estimated capacity: %s lines", max_lines_per_textbox)
        logger.info("%s\n", '='*80)
        
        # _slot_names_for_actual_slide is the single source of truth for
        # this (also used by _append_table_accounts_to_distribution) --
        # see its docstring for why the two must never diverge.
        def slot_names_for_actual_slide(actual_slide_idx: int) -> List[str]:
            return self._slot_names_for_actual_slide(actual_slide_idx, start_slide)

        # Define slot structure: (slide_idx, slot_name)
        slots: List[Tuple[int, str]] = []
        for slide_idx in range(max_slides):
            actual_slide_idx = start_slide - 1 + slide_idx
            for slot_name in slot_names_for_actual_slide(actual_slide_idx):
                slots.append((slide_idx, slot_name))

        slot_shapes: Dict[int, Any] = {}
        for slot_idx, (slide_idx, slot_name) in enumerate(slots):
            actual_slide_idx = start_slide - 1 + slide_idx
            slot_shape = None
            if 0 <= actual_slide_idx < len(self.presentation.slides):
                slide = self.presentation.slides[actual_slide_idx]
                slot_shape = self._resolve_commentary_slot_shape(slide, slot_name)
            slot_shapes[slot_idx] = slot_shape or sample_shape

        logger.info("Total slots available: %s", len(slots))
        
        # Distribution result: [(slide_idx, slot_name, [account_data])]
        distribution = []
        
        current_slot_idx = 0
        current_slot_content = []
        current_slot_lines = 0
        previous_category = None

        def slot_capacity_for(slot_idx: int, *, is_chinese: bool, slot_name_override: Optional[str] = None) -> int:
            _slide_idx, derived_slot_name = slots[slot_idx]
            slot_name_local = slot_name_override or derived_slot_name
            slot_shape_local = slot_shapes.get(slot_idx)
            capacity = self._calculate_max_lines_for_textbox(
                slot_shape_local,
                is_chinese=is_chinese,
                slot_name=slot_name_local,
                statement_type=statement_type,
            )
            if slot_name_local == 'L':
                capacity = int(capacity * 0.98)
            return capacity
        
        for account_idx, account_data in enumerate(structured_data):
            mapping_key_debug = account_data.get('mapping_key', account_data.get('account_name', ''))
            logger.info("\nAccount %s/%s: %s", account_idx + 1, len(structured_data), mapping_key_debug)
            if current_slot_idx >= len(slots):
                dropped_accounts = len(structured_data) - account_idx
                logger.warning(
                    "Ran out of commentary slots; dropping %s remaining account(s) starting from '%s'",
                    dropped_accounts,
                    mapping_key_debug,
                )
                break
            
            category = account_data.get('category', '')
            mapping_key = account_data.get('mapping_key', account_data.get('account_name', ''))
            commentary = account_data.get('commentary', '')

            slide_idx_check, slot_name_check = slots[current_slot_idx]
            is_chinese_content = any('\u4e00' <= c <= '\u9fff' for c in commentary)
            chars_setting = 35 if is_chinese_content else 70
            category_lines = 1 if (category and category != previous_category) else 0
            content_lines = self._calculate_content_lines(
                '',
                mapping_key,
                commentary,
                slot_name=slot_name_check,
                shape=slot_shapes.get(current_slot_idx),
                is_chinese=is_chinese_content,
                statement_type=statement_type,
            )
            total_lines = category_lines + content_lines
            logger.info("  Category: '%s', Lines: cat=%s, content=%s, total=%s", category, category_lines, content_lines, total_lines)
            logger.info("  Commentary length: %s chars, Language: %s, Chars/line: %s", len(commentary), 'Chinese' if is_chinese_content else 'English', chars_setting)

            adjusted_capacity = slot_capacity_for(current_slot_idx, is_chinese=is_chinese_content, slot_name_override=slot_name_check)
            logger.info("  Current slot %s (%s): %s/%s lines used", current_slot_idx, slot_name_check, current_slot_lines, adjusted_capacity)

            if current_slot_lines + total_lines <= adjusted_capacity:
                current_slot_content.append(account_data)
                current_slot_lines += total_lines
                previous_category = category
                logger.info("  Slot %s (%s): Added '%s' (%s lines), now %s/%s lines used", current_slot_idx, slot_name_check, mapping_key, total_lines, current_slot_lines, adjusted_capacity)
            else:
                remaining_lines = adjusted_capacity - current_slot_lines
                logger.info("  Doesn't fit. Remaining: %s lines, Content: %s lines", remaining_lines, content_lines)

                # NOTE: the tail-overflow tolerance is deliberately NOT applied
                # here. Measured on a 14-account workload, keeping an account
                # whole at THIS point made the final deck worse (4 split
                # fragments -> 5): this greedy pass only produces the starting
                # point for _optimize_slot_fill, which re-packs everything
                # anyway, so suppressing a split here just perturbs the DP's
                # input and it re-splits elsewhere. The tolerance belongs at
                # the pass that actually decides to cut an account in half --
                # see _rebalance_overflowing_boundaries.
                next_slot_idx = current_slot_idx + 1

                split_remaining_min = float(self._packing_settings(statement_type).get("split_min_remaining_lines", 3))
                split_content_min = int(self._packing_settings(statement_type).get("split_min_content_lines", 5))
                if remaining_lines > split_remaining_min and content_lines > split_content_min:
                    logger.info("  Attempting to split content...")
                    paragraphs = commentary.split('\n\n')
                    if len(paragraphs) == 1:
                        paragraphs = commentary.split('\n')

                    chars_per_line = self._estimate_chars_per_line(
                        slot_name_check,
                        is_chinese_content,
                        shape=slot_shapes.get(current_slot_idx),
                        statement_type=statement_type,
                    )
                    available_for_commentary = remaining_lines - category_lines - 1

                    # Convert float line-units to visual display lines for the
                    # paragraph-fitting loop below.  available_for_commentary is
                    # in "std_lh units" (one unit = line_h + space_after pt),
                    # but para_lines is computed via chars_per_line and counts
                    # visual display lines.  Multiply by (std_lh / line_h) so
                    # both are in the same unit.
                    #   Both languages: std_lh=9+3=12pt, line_h=9pt → factor ≈ 1.333
                    #   (matches _real_font_size_pt/_real_line_spacing/_PARA_SPACE_AFTER
                    #   -- _fill_text_main_bullets_with_category_and_key hardcodes
                    #   Pt(3) space_after / line_spacing=1.0 for every paragraph,
                    #   any language; there's no separate Chinese 10pt/0.95 value.)
                    _lh_est = 9.0
                    _std_lh_est = _lh_est + self._PARA_SPACE_AFTER
                    available_visual = available_for_commentary * (_std_lh_est / _lh_est)

                    if available_for_commentary > 0:
                        part1_commentary = None
                        part2_commentary = None
                        part1_paragraphs = []
                        part1_lines_used = 0
                        split_index = 0

                        for i, para in enumerate(paragraphs):
                            para_lines = max(1, (len(para) + chars_per_line - 1) // chars_per_line)
                            if part1_lines_used + para_lines <= available_visual:
                                part1_paragraphs.append(para)
                                part1_lines_used += para_lines
                                split_index = i + 1
                            else:
                                break

                        if split_index == len(paragraphs) and part1_paragraphs:
                            # All paragraphs fit in remaining space per heuristic.
                            # Pillow slightly over-counted (< ~1-2 lines) — tolerate
                            # it and force-add to current slot rather than leave a
                            # half-empty slot with nothing below.
                            current_slot_content.append(account_data)
                            current_slot_lines += total_lines
                            previous_category = category
                            logger.info(
                                "  Heuristic fit: forced '%s' into slot %s "
                                "(Pillow over-count tolerated, total_lines=%.1f)",
                                mapping_key, current_slot_idx, total_lines,
                            )
                            continue  # skip move-to-next-slot fallthrough
                        elif part1_paragraphs and split_index < len(paragraphs):
                            # Clean paragraph-boundary split — always safe.
                            part1_commentary = '\n\n'.join(part1_paragraphs).strip()
                            part2_commentary = '\n\n'.join(paragraphs[split_index:]).strip()
                        elif not part1_paragraphs and len(paragraphs) > 0:
                            para = paragraphs[0]
                            chars_available = int(max(1, available_visual * chars_per_line))

                            if len(para) > chars_available:
                                # Only split at SENTENCE boundaries (period,
                                # Chinese full-stop, "!", "?"). No commas,
                                # word-breaks, or hard char cuts — those
                                # produce ugly mid-row fragments.
                                #
                                # Strategy: collect EVERY sentence ending up
                                # to 5 % past chars_available, then pick the
                                # one closest to chars_available without
                                # going over it. This packs the current slot
                                # tight — like a human would — rather than
                                # grabbing the first boundary found and
                                # leaving 5–6 rows of empty space below.
                                hard_cap = min(len(para), int(chars_available * 1.05))
                                _split_min_ratio = float(
                                    self._packing_settings(statement_type).get("split_min_fill_ratio", 0.6)
                                )
                                min_fill = max(1, int(chars_available * _split_min_ratio))
                                end_positions: List[int] = []
                                for end_char in ['. ', '。', '! ', '！', '? ', '？']:
                                    start = 0
                                    while True:
                                        pos = para.find(end_char, start, hard_cap)
                                        if pos < 0:
                                            break
                                        end_positions.append(pos + len(end_char))
                                        start = pos + 1
                                # Keep only splits that still leave the slot
                                # at least 15 % filled — avoids cutting after
                                # the first tiny opening sentence.
                                candidates = [p for p in end_positions if p >= min_fill]
                                best_split = max(candidates) if candidates else None

                                if best_split is None:
                                    # No sentence boundary fits — fall back to
                                    # word boundary to use all available lines
                                    # rather than leave the slot empty.
                                    word_end = para.rfind(' ', 0, hard_cap)
                                    if word_end > 0:
                                        best_split = word_end + 1
                                    else:
                                        # Chinese has no spaces, so the word
                                        # fallback above NEVER fires and this
                                        # dropped straight to a raw character
                                        # cut. That is the real defect behind
                                        # the reported "...某某系统" | "工程第四建设
                                        # 有限公司" -- a cut
                                        # through the middle of a company
                                        # name. _snap_split_before_number
                                        # can't rescue it either: "系统" is a
                                        # legitimate jieba token boundary, so
                                        # nothing downstream sees a problem.
                                        # A clause boundary is always a better
                                        # cut than an arbitrary character, and
                                        # is only used when no sentence
                                        # boundary was available at all.
                                        clause_end = max(
                                            (para.rfind(c, min_fill, hard_cap)
                                             for c in ('，', '；', '、', '：', ',', ';')),
                                            default=-1,
                                        )
                                        if clause_end > 0:
                                            best_split = clause_end + 1
                                        elif chars_available < len(para):
                                            best_split = chars_available  # last-resort hard cut

                                # Never split inside a number -- this is a
                                # SEPARATE, independent split implementation
                                # from _split_commentary_at_boundary (this is
                                # the FIRST-ever split an account gets, during
                                # the initial greedy distribution, before any
                                # rebalance pass runs) and _snap_split_before_
                                # number's fix there never touched this one --
                                # confirmed as the real source of a production
                                # case reading "...室外工程人民币2," at the
                                # bottom of one slide and "818.7万元..."
                                # continuing the next, AFTER that other fix
                                # already shipped.
                                if best_split:
                                    best_split = self._snap_split_before_number(para, best_split)

                                # Slice ONCE on the finalised best_split. Previously the
                                # word-boundary/hard-cut fallback recomputed best_split but
                                # never sliced, leaving part1/part2 unset → UnboundLocalError
                                # or stale text from a prior account bleeding onto the slide.
                                if best_split:
                                    part1_commentary = para[:best_split].strip()
                                    remaining_para = para[best_split:].strip()
                                    if len(paragraphs) > 1:
                                        part2_commentary = remaining_para + '\n\n' + '\n\n'.join(paragraphs[1:])
                                    else:
                                        part2_commentary = remaining_para
                                else:
                                    # No boundary and no cut possible — keep the whole
                                    # paragraph in part1 rather than corrupting the slide.
                                    part1_commentary = para
                                    part2_commentary = '\n\n'.join(paragraphs[1:]) if len(paragraphs) > 1 else None
                            else:
                                part1_commentary = para
                                part2_commentary = '\n\n'.join(paragraphs[1:]) if len(paragraphs) > 1 else ""
                        else:
                            part1_commentary = commentary
                            part2_commentary = None

                        if part1_commentary and part2_commentary:
                            account_part1 = account_data.copy()
                            account_part1['commentary'] = part1_commentary
                            account_part1['is_partial'] = True
                            account_part1['part_num'] = 1
                            current_slot_content.append(account_part1)
                            
                            # Save current slot
                            slide_idx, slot_name = slots[current_slot_idx]
                            distribution.append((slide_idx, slot_name, current_slot_content))
                            logger.info("Split '%s': Part 1 (%s chars) to slot %s, Part 2 (%s chars) to next slot", mapping_key, len(part1_commentary), current_slot_idx, len(part2_commentary))

                            if current_slot_idx + 1 >= len(slots):
                                logger.warning(
                                    "Ran out of commentary slots after splitting '%s'; dropping the remaining continuation",
                                    mapping_key,
                                )
                                break
                            current_slot_idx += 1

                            account_part2 = account_data.copy()
                            account_part2['commentary'] = part2_commentary
                            account_part2['is_continuation'] = True
                            account_part2['part_num'] = 2
                            account_part2['original_key'] = mapping_key

                            next_slot_name = slots[current_slot_idx][1]
                            part2_lines = self._calculate_content_lines(
                                '',
                                mapping_key,
                                part2_commentary,
                                slot_name=next_slot_name,
                                shape=slot_shapes.get(current_slot_idx),
                                is_chinese=is_chinese_content,
                            )
                            current_slot_content = [account_part2]
                            current_slot_lines = part2_lines
                            previous_category = None
                            continue
                else:
                    logger.info("  Not splitting: remaining_lines=%s, content_lines=%s", remaining_lines, content_lines)

                if current_slot_content:
                    slide_idx, slot_name = slots[current_slot_idx]
                    distribution.append((slide_idx, slot_name, current_slot_content))
                    logger.info("  Slot %s (%s): FULL with %s accounts, %s lines used", current_slot_idx, slot_name, len(current_slot_content), current_slot_lines)

                current_slot_idx += 1
                if current_slot_idx >= len(slots):
                    logger.warning(
                        "Ran out of commentary slots while placing '%s'; dropping that account from the remaining slides",
                        mapping_key,
                    )
                    break

                slide_idx_new, slot_name_new = slots[current_slot_idx]
                moved_account = account_data.copy()
                moved_account["commentary"] = commentary
                moved_category_lines = 1 if category else 0
                moved_lines = self._calculate_content_lines(
                    '',
                    mapping_key,
                    moved_account["commentary"],
                    slot_name=slot_name_new,
                    shape=slot_shapes.get(current_slot_idx),
                    is_chinese=is_chinese_content,
                    statement_type=statement_type,
                )
                current_slot_content = [moved_account]
                current_slot_lines = moved_category_lines + moved_lines
                previous_category = category
                logger.info("  Moving '%s' to next slot %s (%s), %s lines", mapping_key, current_slot_idx, slot_name_new, current_slot_lines)
        
        # Save last slot if it has content
        if current_slot_content and current_slot_idx < len(slots):
            slide_idx, slot_name = slots[current_slot_idx]
            distribution.append((slide_idx, slot_name, current_slot_content))
            logger.info("  Slot %s (%s): FINAL with %s accounts, %s lines", current_slot_idx, slot_name, len(current_slot_content), current_slot_lines)
        
        slot_position_map = {slot: idx for idx, slot in enumerate(slots)}

        logger.info("\nDistribution complete: %s slots filled", len(distribution))
        # Diagnostic fill-ratio logging used to recompute Pillow measurements
        # for every (slot, account) pair after the packer was already done.
        # That added 1-3s per export with no functional value. Skip the
        # recompute and just log slot composition.
        if logger.isEnabledFor(logging.DEBUG):
            for distribution_idx, (slide_idx, slot_name, accounts) in enumerate(distribution):
                slot_idx = slot_position_map.get((slide_idx, slot_name), distribution_idx)
                logger.debug(
                    "  Slot %s (Slide %s, %s): %s accounts",
                    slot_idx, slide_idx, slot_name, len(accounts),
                )
        
        # --- Fill optimization pass: pull accounts forward into under-filled slots ---
        distribution = self._optimize_slot_fill(
            distribution,
            slot_shapes=slot_shapes,
            slot_meta=slots,
            statement_type=statement_type,
        )

        # A slide that ends up with real content in one slot (e.g. L) but
        # never had ANY entry at all for one of its OTHER slots (e.g. R) --
        # not drained to empty by a rebalance pass, just never touched by
        # the greedy first pass above, which only appends a distribution
        # entry "if current_slot_content" (line ~2776) and so silently
        # skips a slot that ran out of content before reaching it -- needs
        # that other slot represented too, even with zero accounts.
        # _optimize_slot_fill can't fix this on its own: it only ever
        # builds its internal slot list FROM the `distribution` it's given
        # (see its "for slide_idx, slot_name, _accounts in distribution"
        # loop), so a slot that was never in `distribution` to begin with
        # never reaches it as input, let alone its output -- unlike the
        # "consolidated to empty" case _optimize_slot_fill's own output
        # filter already handles (5da5d38), this one is invisible to it
        # entirely. Confirmed via a real export: the resulting shape is
        # left holding whatever raw template placeholder text it shipped
        # with (a literal "Placeholder – placeholder"), because the
        # render loop's per-slot fill AND its "clear if empty" fallback
        # both key off `slot_contents.items()` -- a slot name that's
        # simply not a key in that dict is never visited by either branch.
        slides_with_content = {slide_idx for slide_idx, _slot_name, accounts in distribution if accounts}
        present_pairs = {(slide_idx, slot_name) for slide_idx, slot_name, _accounts in distribution}
        for slide_idx, slot_name in slots:
            if slide_idx in slides_with_content and (slide_idx, slot_name) not in present_pairs:
                distribution.append((slide_idx, slot_name, []))

        return distribution

    # Rendered bullet paragraphs hang-indent their WRAPPED lines by 0.15"
    # (p_key.left_indent = Inches(0.15), first_line_indent = Inches(-0.15)):
    # line 1 starts at the box margin and spans the full width, every wrapped
    # line starts 10.8pt further right. Cost estimates wrap at
    # (width - 10.8pt) -- exact for wrapped lines, ~1 char conservative for
    # line 1. Measuring every line at full width under-counted long
    # paragraphs by up to a line each; that error was masked while the old
    # line-height model over-charged the vertical axis, and became real
    # (uncaught-by-autofit) overflow once the 1.2x-size pitch fix landed.
    _BULLET_HANGING_INDENT_PT = 10.8  # 0.15 inch

    def _account_is_chinese(self, account: Dict) -> bool:
        """Language flag for MEASUREMENT (which glyph-width table to wrap
        with). Uses the account's own is_chinese when present (set by the
        payload builder via contains_predominantly_chinese_text); otherwise
        detects from the commentary instead of silently defaulting to
        English -- measuring CJK text with Arial's advance table (which has
        no CJK glyphs) under-counted lines badly enough that genuinely
        overflowing slots passed the render-time autofit gate as 'fits'."""
        v = (account or {}).get("is_chinese")
        if v is not None:
            return bool(v)
        return contains_predominantly_chinese_text(str((account or {}).get("commentary", "")))

    def _account_cost_key(self, account: Dict) -> str:
        """The key text whose rendered width the cost model should charge:
        mapping_key plus the continuation marker the renderer appends
        (' (续)' / \" (cont'd)\") -- previously never charged, so every
        continuation's first paragraph was measured ~4-9 chars short."""
        key = str(account.get("mapping_key", account.get("account_name", "")) or "")
        if account.get("is_continuation"):
            key += " (续)" if self._account_is_chinese(account) else " (cont'd)"
        return key

    def _compute_slot_used_lines(
        self,
        accounts: List[Dict],
        slot_name: str,
        slot_shape=None,
        statement_type: Optional[str] = None,
    ) -> float:
        """Return used line-units for *accounts* in this slot (float).

        Uses the same accounting as ``slot_cost`` in the DP: each category
        header costs 1 line unit, and each account's commentary costs the
        float value returned by ``_calculate_content_lines`` (actual pt /
        std_lh, no ceil).  Comparing against int capacity from
        ``_calculate_max_lines_for_textbox`` gives accurate fill ratios.
        """
        used: float = 0.0
        prev_cat = None
        for account in accounts:
            cat = str(account.get("category", "") or "")
            # A continuation account never gets its own category-header
            # paragraph rendered (the fill loop explicitly skips it via
            # "not is_continuation" -- a "(cont'd)" fragment belongs to
            # whatever category its first part already introduced, not a
            # second one) -- so it must not be charged the 1.0-line cost
            # for one here either. A continuation is very often the FIRST
            # account placed in a slot (that's how it got split off), so
            # skipping this made every such slot's "used" belief exactly
            # 1.0 line higher than what actually gets rendered -- a
            # believed-vs-actual gap that compounds with the DP's own
            # front-loading and was a real contributor to slots reading as
            # under-filled despite the packing math saying they were full.
            # prev_cat only advances on an actually-rendered header (mirrors
            # the render loop's own current_category, which likewise never
            # moves off a skipped continuation) -- otherwise a later account
            # sharing the continuation's category would wrongly believe its
            # own header was already shown further up.
            if cat and cat != prev_cat and not account.get("is_continuation"):
                if prev_cat is None:
                    # The slot's OWN first category header gets NO space_before
                    # (matches the render's "p_category.space_before = Pt(3)
                    # if current_category else Pt(0)" -- current_category is
                    # still unset the very first time). Charging the full
                    # 1.0-unit (line_h+gap) here overstated every slot's
                    # opening header by one gap's worth -- confirmed by
                    # replicating the render's own paragraph-by-paragraph pt
                    # math for a real 7-account/2-category slot and comparing
                    # totals directly (0.23-unit gap, ~3pt, matched exactly).
                    _approx_line_h = self._real_font_size_pt(False) * self._real_line_spacing(False)
                    _approx_std_lh = _approx_line_h + self._real_para_gap_pt(False)
                    used += (_approx_line_h / _approx_std_lh) if _approx_std_lh > 0 else 1.0
                else:
                    used += 1.0   # category header (same as slot_cost)
                prev_cat = cat
            used += self._calculate_content_lines(
                "",
                self._account_cost_key(account),
                account.get("commentary", ""),
                slot_name=slot_name,
                shape=slot_shape,
                is_chinese=self._account_is_chinese(account),
                statement_type=statement_type,
            )
            # A presentation-table account's `commentary` is only its lead-in;
            # the table itself and the explanation below it are reserved as
            # blank lines by _render_table_accounts_stack and were not counted
            # here at all. On a real deck that read a slot the renderer fills
            # to 94% as 10% used -- and once inspect_databook's fill
            # diagnostic started modelling the real slot assignment, that
            # under-count immediately produced false "GENUINE GAP" findings
            # (an account "fits in the 16.7 remaining lines" whose table alone
            # needs 24). Safe to charge here because table accounts never
            # reach the packer: they are pulled out before
            # _distribute_content_across_slots and appended afterwards, so
            # the only callers that ever see one are the render-time autofit
            # gate and the diagnostic -- both of which want the real height.
            # ...charged per PART, matching _render_parts, which is how a
            # split table account is actually rendered: the first fragment
            # draws ("lead", "table") and the continuation draws ("expl",)
            # only. Charging the whole block to both put 管理费用's table on
            # two slots and read a 67%-full column as 117%. Absent
            # _render_parts means the account renders whole.
            table = account.get("_presentation_table")
            if table:
                _is_chi = self._account_is_chinese(account)
                _parts = account.get("_render_parts")
                _lead_pt, _table_pt, _explain_pt = self._estimate_table_account_parts_pt(
                    account, table, _is_chi,
                )
                _std_lh = self._planning_std_lh_pt(_is_chi)
                if _std_lh > 0:
                    _extra = 0.0
                    if _parts is None or "table" in _parts:
                        _extra += _table_pt
                    if _parts is None or "expl" in _parts:
                        _extra += _explain_pt
                    used += _extra / _std_lh
        return max(0.0, used)

    def _rebalance_lopsided_lr_pairs(
        self,
        assignment: List[List[Dict[str, Any]]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> None:
        """Mutates `assignment` in place. Fixes same-page L/R pairs where the
        DP left one column completely empty while the other is significantly
        full.

        Root cause: _optimize_slot_fill's lexicographic objective is
        (num_nonempty_slots, underfill_penalty). An EMPTY slot contributes
        NEITHER — it costs the DP nothing — while a non-empty slot under
        target_fill_min_ratio incurs a real penalty. So whenever a page's
        total content doesn't reach ~2x a single slot's capacity, "one
        slot full + one slot empty" scores strictly better than "two slots
        each moderately full" (e.g. one at 90% + one empty beats two at
        45% each: 5% penalty vs 50%+50%). That is a reasonable objective
        for deciding how many PAGES to use, but applied to the two columns
        of a single page it produces a visually broken half-blank layout a
        reader would never expect.

        This is a post-pass, not a DP objective change, to keep the fix
        narrowly scoped: find same-slide (L, R) pairs with exactly one
        side empty, and if the full side has 2+ accounts, look for the
        split point whose two halves are most evenly balanced (by the same
        real line-cost function the DP itself uses) without overflowing
        either box. If no such split exists, leave the pair as the DP
        produced it — this never risks introducing an overflow that wasn't
        there before.
        """
        by_slide: Dict[int, Dict[str, int]] = {}
        for s_i, slot in enumerate(slots):
            if slot.get("slot_name") in ("L", "R"):
                by_slide.setdefault(slot["slide_idx"], {})[slot["slot_name"]] = s_i

        for slide_idx, pair in by_slide.items():
            if "L" not in pair or "R" not in pair:
                continue
            l_i, r_i = pair["L"], pair["R"]
            l_accts, r_accts = assignment[l_i], assignment[r_i]
            if bool(l_accts) == bool(r_accts):
                continue  # both empty, or both already non-empty -- nothing to rebalance
            full_i, empty_i = (l_i, r_i) if l_accts else (r_i, l_i)
            full_accts = assignment[full_i]
            if len(full_accts) < 2:
                continue  # can't split a single account here -- that's the greedy first pass's job

            full_name = slots[full_i]["slot_name"]
            empty_name = slots[empty_i]["slot_name"]
            full_cap = slots[full_i]["capacity"]
            empty_cap = slots[empty_i]["capacity"]
            full_shape = slots[full_i]["shape"]
            empty_shape = slots[empty_i]["shape"]

            best_k, best_diff = None, None
            for k in range(1, len(full_accts)):
                first, rest = full_accts[:k], full_accts[k:]
                first_lines = self._compute_slot_used_lines(
                    first, empty_name, slot_shape=empty_shape, statement_type=statement_type,
                )
                rest_lines = self._compute_slot_used_lines(
                    rest, full_name, slot_shape=full_shape, statement_type=statement_type,
                )
                if first_lines > empty_cap or rest_lines > full_cap:
                    continue
                diff = abs(first_lines - rest_lines)
                if best_diff is None or diff < best_diff:
                    best_diff, best_k = diff, k

            if best_k is None:
                continue  # no split keeps both sides within capacity -- leave as-is

            assignment[empty_i] = full_accts[:best_k]
            assignment[full_i] = full_accts[best_k:]
            logger.info(
                "  Rebalanced lopsided L/R pair on slide %s: moved %s account(s) from slot %s "
                "into previously-empty slot %s",
                slide_idx, best_k, full_i, empty_i,
            )

    def _consolidate_tiny_stub_lr_pairs(
        self,
        assignment: List[List[Dict[str, Any]]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> None:
        """Mutates `assignment` in place. Front-loading -- filling L to
        near-capacity before spilling anything into R -- is the intended
        packing philosophy here (minimises total page count), not a bug to
        fix. But confirmed against real production decks: it regularly
        leaves R holding a single orphaned fragment at ~5-20% fill next to
        an L column in the 90s% -- not what a human preparing the same
        page would produce (they'd either use both columns for real or
        fold a leftover sliver into the fuller one, never leave a
        near-blank second column standing next to a full one).

        This only ever fires on a same-slide (L, R) pair where R is
        NON-empty (bool(r_accts) is True, so _rebalance_lopsided_lr_pairs'
        "one side fully empty" case never overlaps with this one) but its
        fill ratio reads as an orphaned stub rather than real content, AND
        R's entire content fits into L's own remaining capacity whole --
        folds it in and leaves R properly, cleanly empty (a genuinely
        blank box now correctly renders that way -- see the
        "if not account_data_list" shape-clearing fix alongside this).
        Never attempts a partial move here: trading one awkward stub for a
        smaller one doesn't serve the "look like a human made this" goal.
        """
        STUB_FILL_THRESHOLD = 0.20  # below this, a slot reads as leftover, not real content

        by_slide: Dict[int, Dict[str, int]] = {}
        for s_i, slot in enumerate(slots):
            if slot.get("slot_name") in ("L", "R"):
                by_slide.setdefault(slot["slide_idx"], {})[slot["slot_name"]] = s_i

        for slide_idx, pair in by_slide.items():
            if "L" not in pair or "R" not in pair:
                continue
            l_i, r_i = pair["L"], pair["R"]
            l_accts, r_accts = assignment[l_i], assignment[r_i]
            if not l_accts or not r_accts:
                continue  # nothing to fold, or _rebalance_lopsided_lr_pairs's job

            l_slot, r_slot = slots[l_i], slots[r_i]
            r_used = self._compute_slot_used_lines(
                r_accts, r_slot["slot_name"], slot_shape=r_slot["shape"], statement_type=statement_type,
            )
            r_cap = r_slot["capacity"]
            if r_cap <= 0 or (r_used / r_cap) >= STUB_FILL_THRESHOLD:
                continue  # R holds real content, not just a leftover stub -- leave it alone

            combined_used = self._compute_slot_used_lines(
                l_accts + r_accts, l_slot["slot_name"], slot_shape=l_slot["shape"], statement_type=statement_type,
            )
            if combined_used > l_slot["capacity"]:
                continue  # doesn't fit whole -- a partial move would just trade one stub for another

            assignment[l_i] = l_accts + r_accts
            assignment[r_i] = []
            logger.info(
                "  Consolidated tiny stub on slide %s: folded R (%.0f%% fill) into L -- "
                "avoids an orphaned near-empty trailing column",
                slide_idx, r_used / r_cap * 100,
            )

    def _split_commentary_at_boundary(
        self,
        commentary: str,
        available_std_lh_units: float,
        *,
        slot_name: str,
        is_chinese: bool,
        shape=None,
        statement_type: Optional[str] = None,
        min_fill_ratio: float = 0.5,
        key_prefix: str = "",
        min_available_visual: float = 1.0,
        overflow_allowance_units: Optional[float] = None,
    ) -> Optional[Tuple[str, str]]:
        """Find a clean split point (paragraph boundary, else sentence
        boundary, else word boundary) so the head of `commentary` fits
        within `available_std_lh_units` std_lh-units. Pure/side-effect
        free version of the paragraph-fit + sentence-boundary logic the
        initial forward-fill pass uses for its one-time split, so a
        boundary rebalance can invoke it repeatedly -- the same account's
        remainder may need splitting again at a later pass, which is how
        a single account ends up spread across more than two slots.

        Returns None if nothing usable can be carved off (already fits
        whole, or every candidate split leaves a sliver below
        `min_fill_ratio` of the available space).

        `min_available_visual` is the floor below which this refuses to
        even try -- default 1.0 (roughly one real line) for
        balance-oriented callers. _maximize_forward_fill, whose whole
        purpose is squeezing out the LAST bit of a nearly-full slot,
        passes a lower value; every other caller keeps the safer default.
        """
        commentary = str(commentary or "")
        if not commentary.strip():
            return None

        chars_per_line = self._estimate_chars_per_line(
            slot_name, is_chinese, shape=shape, statement_type=statement_type,
        )
        # Both languages: matches _real_font_size_pt/_real_line_spacing (see
        # the sibling occurrence in _distribute_content_across_slots for the
        # full explanation) -- no separate Chinese 10pt/0.95 value actually
        # gets applied to a rendered paragraph.
        _lh_est = 9.0
        _std_lh_est = _lh_est + self._PARA_SPACE_AFTER
        available_visual = max(0.0, available_std_lh_units) * (_std_lh_est / _lh_est)
        if available_visual < min_available_visual:
            return None

        paragraphs = commentary.split('\n\n')
        if len(paragraphs) == 1:
            paragraphs = commentary.split('\n')

        part1_paragraphs: List[str] = []
        part1_lines_used = 0
        split_index = 0
        for para in paragraphs:
            para_lines = max(1, (len(para) + chars_per_line - 1) // chars_per_line)
            if part1_lines_used + para_lines <= available_visual:
                part1_paragraphs.append(para)
                part1_lines_used += para_lines
                split_index += 1
            else:
                break

        if part1_paragraphs and split_index == len(paragraphs):
            return None  # whole thing already fits -- caller shouldn't be splitting

        if part1_paragraphs:
            part1 = '\n\n'.join(part1_paragraphs).strip()
            part2 = '\n\n'.join(paragraphs[split_index:]).strip()
            if part1 and part2 and part1_lines_used >= available_visual * min_fill_ratio:
                return part1, part2
            return None

        # No whole paragraph fits by the coarse ceil-per-line estimate above --
        # try a sentence boundary within paragraph 0 using a finer character
        # budget instead. The two estimates can disagree (ceil-per-line rounds
        # up to a whole line even for a small fractional overage) -- that's
        # expected, not an error, so don't bail out just because this looser
        # measure says the paragraph would fit; let the boundary search below
        # decide, and downstream capacity re-validation is the real safety net.
        para = paragraphs[0]
        chars_available = int(max(1, available_visual * chars_per_line))
        hard_cap = min(len(para), int(chars_available * 1.05))
        min_fill = max(1, int(chars_available * min_fill_ratio))
        end_positions: List[int] = []
        for end_char in ['. ', '。', '! ', '！', '? ', '？']:
            start = 0
            while True:
                pos = para.find(end_char, start, hard_cap)
                if pos < 0:
                    break
                end_positions.append(pos + len(end_char))
                start = pos + 1
        # Exclude a candidate sitting at (or past) the paragraph's own end --
        # when the finer char-budget is generous enough to reach the last
        # sentence in the paragraph, picking it would leave nothing for the
        # tail, i.e. no real split at all.
        candidates = [p for p in end_positions if min_fill <= p < len(para)]
        best_split = max(candidates) if candidates else None

        if best_split is None:
            word_end = para.rfind(' ', 0, min(hard_cap, len(para) - 1))
            if word_end > 0 and word_end >= min_fill:
                best_split = word_end + 1

        if best_split is None and is_chinese:
            # Chinese has no spaces between words, so the Latin word-boundary
            # fallback above always misses on CJK text -- a long single
            # sentence with no "." within budget (e.g. one continuous
            # amount-heavy clause well past 72 chars before its first "。")
            # had NO fallback at all here and always returned None, silently
            # giving up on pulling any of it forward even with real budget
            # to spare. Comma is the natural next-best CJK clause boundary.
            #
            # The line below used to end "...if even that isn't present in
            # range, a hard character cut is fine for CJK (unlike Latin text,
            # cutting between two Chinese characters doesn't break a word)".
            # A real deck disproved that: it cut
            #   "...主要为应付某某系统" | "工程第四建设有限公司..."
            # straight through a company name. jieba can't rescue it either --
            # 系统|工程 IS a legitimate token boundary, so _snap_split_before_
            # number sees nothing wrong. CJK absolutely does have words; it
            # just doesn't mark them with spaces.
            _cut = min(hard_cap, len(para) - 1)
            comma_end = max(
                para.rfind('，', 0, _cut), para.rfind(',', 0, _cut),
                para.rfind('；', 0, _cut), para.rfind(';', 0, _cut),
                para.rfind('、', 0, _cut),
            )
            if comma_end > 0:
                # Accept a clause boundary even BELOW min_fill: it still
                # beats a cut through the middle of a name.
                #
                # An earlier version returned None here instead, on the
                # theory that "no split, let the box protrude" was the
                # project team's stated preference. Measured, that made
                # things WORSE (per-slot fill 96.9% -> 93.6%): returning
                # None doesn't make the account protrude, it makes the
                # caller decline to pull it in at all, so the slot is left
                # emptier than before -- the reported "放少了1-2行". The
                # tolerance only helps where content STAYS PUT, which is
                # _rebalance_overflowing_boundaries, not here.
                best_split = comma_end + 1
            else:
                hard_end = _cut
                if hard_end >= min_fill:
                    best_split = hard_end

        if not best_split:
            return None

        # Validate against the REAL measurer, not the crude chars_per_line
        # estimate used to find best_split above -- the two can disagree,
        # and because wrapped-line cost is quantized (a candidate a few
        # characters shorter often still wraps to the SAME line count),
        # _maximize_forward_fill's retry loop would see this function
        # propose the exact same "still too big" candidate across several
        # attempts (a small budget cut rarely crosses the crude estimate's
        # own boundary even when the real candidate measurably overshoots),
        # then give up once its overage-scaled reduction blew past the
        # retry floor -- confirmed via a real production case: a 0.44-unit
        # (well under half a line) overage was enough to end the whole
        # retry after a single attempt, discarding a boundary that
        # genuinely had 1.3 real units of room. Back off through the same
        # candidate boundaries already found above (largest to smallest),
        # accepting the first whose ACCURATE wrapped cost actually fits.
        try:
            from fdd_utils.text_metrics import get_measurer, text_box_from_shape
            if shape is not None:
                packing = self._packing_settings(statement_type)
                family = self._measurer_family(is_chinese, packing)
                font_size_pt = self._real_font_size_pt(is_chinese)
                line_spacing = self._real_line_spacing(is_chinese)
                _mpath = self._resolve_font_metrics_path(is_chinese, packing)
                measurer = get_measurer(
                    family, font_size_pt, is_cjk=is_chinese, line_spacing=line_spacing,
                    metrics_path=_mpath,
                )
                box = text_box_from_shape(shape)
                line_h = measurer.line_height_pt()
                para_gap = self._real_para_gap_pt(is_chinese)
                std_lh = line_h + para_gap
                budget_pt = max(0.0, available_std_lh_units) * std_lh
                min_fill_pt = budget_pt * min_fill_ratio

                # Every candidate boundary found above, largest first --
                # sentence ends, then the word/comma/hard-cut fallback
                # position, each de-duplicated and capped at best_split
                # (never consider a candidate BIGGER than the crude
                # estimate already chose).
                backoff_candidates = sorted(
                    {p for p in candidates if p <= best_split} | {best_split}, reverse=True,
                )
                accepted = None
                for cand_pos in backoff_candidates:
                    cand_head = para[:cand_pos].strip()
                    if not cand_head:
                        continue
                    # key_prefix mirrors _calculate_content_lines' own
                    # "■ {mapping_key} - " prepended to an account's first
                    # paragraph -- omitting it here (the original bug) made
                    # every candidate measure smaller than what the caller's
                    # actual _compute_slot_used_lines check would find, so
                    # a "fixed" head that trimmed real characters off could
                    # still fail the caller's check by the exact same margin.
                    wrapped = measurer.wrap(
                        key_prefix + cand_head,
                        max(10.0, box.width_pt - self._BULLET_HANGING_INDENT_PT),
                        first_line_width_pt=box.width_pt,
                    )
                    cand_pt = len(wrapped) * line_h + para_gap
                    if cand_pt <= budget_pt and cand_pt >= min_fill_pt:
                        accepted = cand_pos
                        break
                if accepted is not None and accepted != best_split:
                    best_split = accepted
                elif accepted is None:
                    # Nothing in the existing candidate list fits even
                    # accurately-measured -- fall back to trimming the
                    # crude choice down word-by-word (Latin) or char-by-
                    # char (CJK) until it does, rather than giving up
                    # outright on a boundary that may still have room. This
                    # is the last-resort branch (no sentence/comma boundary
                    # produced anything accurately-fitting either) -- accept
                    # the first fit regardless of min_fill_ratio here, since
                    # by construction it's the LARGEST candidate that fits
                    # (we're trimming down from best_split), so nothing
                    # smaller would meet min_fill_ratio either; a small-but-
                    # real split still beats the caller getting nothing back
                    # and abandoning a boundary that had real room to spare.
                    trimmed = best_split
                    while trimmed > 0:
                        cand_head = para[:trimmed].strip()
                        if cand_head:
                            wrapped = measurer.wrap(
                                key_prefix + cand_head,
                                max(10.0, box.width_pt - self._BULLET_HANGING_INDENT_PT),
                                first_line_width_pt=box.width_pt,
                            )
                            cand_pt = len(wrapped) * line_h + para_gap
                            if cand_pt <= budget_pt:
                                best_split = trimmed
                                break
                        trimmed -= max(1, len(cand_head) // 10) if cand_head else 5
                    else:
                        return None  # nothing, down to zero, ever fit

                # Extend FORWARD past whatever best_split the crude search
                # above settled on, up to the TRUE available budget. The
                # crude per-character estimate that located candidates/
                # hard_cap above (_AVG_CHAR_WIDTH_CHI/_WORD_WRAP_SLACK)
                # systematically under-counts real Chinese glyph density
                # (~10.9pt/char effective vs the real ~9pt at this font
                # size), and that same 8% "word-wrap slack" margin is
                # justified for Latin text (room for a whole word that
                # might not fit) but doesn't apply to CJK at all (every
                # character is its own valid break point) -- so even the
                # "largest candidate that fits" from the search above
                # routinely stops well short of the box's true capacity.
                # This is exactly the "沒有用盡那一行才cut" a user report
                # confirmed by hand-counting real rendered lines against
                # this file's own prediction: a whole further line of real,
                # usable room was never even considered as a candidate.
                # Binary-search directly against the accurate measurer (not
                # the crude heuristic) for the true maximum that fits, then
                # prefer snapping to the nearest natural sentence/comma
                # boundary AT OR BEFORE that true maximum so the cut still
                # reads naturally wherever that costs nothing.
                lo, hi = best_split, len(para)
                true_max = best_split
                while lo <= hi:
                    mid = (lo + hi) // 2
                    cand_head = para[:mid].strip()
                    if not cand_head:
                        lo = mid + 1
                        continue
                    wrapped = measurer.wrap(
                        key_prefix + cand_head,
                        max(10.0, box.width_pt - self._BULLET_HANGING_INDENT_PT),
                        first_line_width_pt=box.width_pt,
                    )
                    cand_pt = len(wrapped) * line_h + para_gap
                    if cand_pt <= budget_pt:
                        true_max = mid
                        lo = mid + 1
                    else:
                        hi = mid - 1

                # Choose the split point that best FILLS the budget, out of
                # every real sentence/clause boundary within reach.
                #
                # The old rule -- "take any boundary that fits" -- silently
                # wasted space: measured on the reviewed deck's own 应付账款
                # text, a 1.5-line budget stopped at a 1.00-line boundary and
                # left half a line empty, because the next boundary needed
                # 1.78 lines, i.e. 0.28 over. That is nowhere near the ~2
                # lines of protrusion the project team accepts, and it is the
                # reported "能夠再放1-2行才滿". Callers re-validate with the
                # same tolerance (see _try_partial_split_into_gap), so an
                # overshoot chosen here survives rather than being trimmed.
                # Callers that already widened their own budget by the
                # protrusion allowance pass 0 here; adding it a second time
                # let the search run to capacity + 2x the allowance, which
                # overshot every real boundary and fell back to a raw
                # character cut (measured: one appeared the moment
                # _try_partial_split_into_gap started pre-inflating its gap).
                _allow = (self._tail_overflow_tolerance_units(statement_type)
                          if overflow_allowance_units is None
                          else float(overflow_allowance_units))
                tol_pt = _allow * (line_h + para_gap)

                def _head_pt(pos: int) -> Optional[float]:
                    head_txt = para[:pos].strip()
                    if not head_txt:
                        return None
                    return len(measurer.wrap(
                        key_prefix + head_txt,
                        max(10.0, box.width_pt - self._BULLET_HANGING_INDENT_PT),
                        first_line_width_pt=box.width_pt,
                    )) * line_h + para_gap

                # Largest position still within budget+tolerance, so the
                # boundary scan below never has to look further than useful.
                lo2, hi2, tol_max = best_split, len(para) - 1, best_split
                while lo2 <= hi2:
                    mid2 = (lo2 + hi2) // 2
                    cand_pt = _head_pt(mid2)
                    if cand_pt is None:
                        lo2 = mid2 + 1
                        continue
                    if cand_pt <= budget_pt + tol_pt:
                        tol_max = mid2
                        lo2 = mid2 + 1
                    else:
                        hi2 = mid2 - 1

                positions = set()
                for end_char in ('. ', '。', '! ', '！', '? ', '？',
                                 '，', ',', '；', ';', '、'):
                    scan = min_fill
                    while True:
                        found = para.find(end_char, scan, tol_max + 1)
                        if found < 0:
                            break
                        cand = found + len(end_char)
                        if min_fill <= cand < len(para):
                            positions.add(cand)
                        scan = found + 1

                best_under = best_over = None
                under_pt = over_pt = 0.0
                for cand in sorted(positions):
                    cand_pt = _head_pt(cand)
                    if cand_pt is None:
                        continue
                    if cand_pt <= budget_pt:
                        if best_under is None or cand > best_under:
                            best_under, under_pt = cand, cand_pt
                    elif cand_pt <= budget_pt + tol_pt:
                        if best_over is None or cand < best_over:
                            best_over, over_pt = cand, cand_pt
                # Only spend the protrusion when staying inside the budget
                # would actually waste something -- half a line is the point
                # at which the gap is visible in a rendered deck.
                _WASTE_PT = 0.5 * (line_h + para_gap)
                if best_under is not None and (budget_pt - under_pt) < _WASTE_PT:
                    best_split = best_under
                elif best_over is not None:
                    best_split = best_over
                elif best_under is not None:
                    best_split = best_under
                elif true_max > best_split:
                    best_split = true_max
        except Exception:
            pass  # measurer unavailable -- fall through with the crude best_split

        # Never split inside a number -- a real production case did exactly
        # this ("...室外工程人民币2," at the end of one slide, "818.7万元..."
        # continuing the next): the hard-cut/character-trim fallbacks above
        # have no concept of a numeric literal, so a raw character offset
        # can land between the "2" and the "8" of "2,818.7". Backing up
        # only ever REMOVES characters from head (never adds), so whatever
        # capacity check already accepted this position stays satisfied.
        best_split = self._snap_split_before_number(para, best_split)

        head = para[:best_split].strip()
        tail_rest = para[best_split:].strip()
        tail = (tail_rest + '\n\n' + '\n\n'.join(paragraphs[1:])).strip() if len(paragraphs) > 1 else tail_rest
        if not head or not tail:
            return None
        return head, tail

    # Currency markers that read as visually broken when stranded at the end
    # of a slide with their amount on the next one (real production case:
    # "...bad debt losses of CNY" / "484,000 in FY25..." -- the number itself
    # was intact, only the currency prefix got separated from it).
    _CURRENCY_MARKERS = ('CNY', 'RMB', 'USD', 'HKD', '人民币', '¥', '$', '£', '€')

    # Common Chinese words/compounds this domain's commentary is full of --
    # protected the same way as _CURRENCY_MARKERS (never split strictly
    # inside one). The mid-word/mid-marker fix for "人民币" alone wasn't
    # enough: the SAME hard-cut/character-trim fallback that used to land
    # inside "人民币" lands just as easily inside any other bare CJK
    # compound, since Chinese has no spaces to signal a word boundary in
    # the first place. Confirmed in real production output on the SAME
    # page: "...人民币784万" | "元变为..." (万元 split), "...2026年06月30日
    # 分" | "别为人民币16万元..." (分别 split), "...2025年" | "度为人民币
    # 17万元..." (年度 split) -- three different compounds, not a single
    # fixable special case. Not exhaustive (full Chinese word segmentation
    # is a different-sized problem) -- this is the domain-specific set
    # observed breaking in real financial-commentary text plus the most
    # obviously analogous high-frequency terms; the mechanism only ever
    # backs a split up (never forces a worse one), so a broader list here
    # is safe to extend further if new cases turn up.
    _PROTECTED_CJK_COMPOUNDS = (
        '万元', '亿元', '年度', '分别', '期间', '变为', '转为', '增至', '降至',
        '合计', '其中', '构成', '包括', '管理层', '余额', '截至', '备注',
        '数据', '期末', '期初', '账款', '借款', '贷款', '利息', '费用',
        '资产', '负债', '权益', '收入', '成本', '折旧', '摊销', '核对',
        '差异', '发生', '形成', '增加', '维持', '原值', '净值', '账面',
        '进一步', '未发生', '无余额', '未形成',
    )

    # Organisation-name tails. A split anywhere inside a company name reads
    # as a mistake even though every generic protection above passes it:
    # jieba tokenizes 某某系统|工程|第四|建设|有限公司 into legitimate words,
    # so a cut between any two of them looks like a clean word boundary.
    # Detecting the name's START is unreliable; detecting that pos is
    # heading INTO one of these tails is not.
    _ORG_NAME_TAILS = (
        '有限公司', '有限责任公司', '股份有限公司', '公司', '银行', '集团',
        '事务所', '研究院', '设计院', '工程局', '合伙企业', '分行', '支行',
    )
    _SENTENCE_PUNCT = '。，；、：！？,;.:!?'

    @classmethod
    def _snap_before_org_name(cls, text: str, pos: int) -> int:
        """Back `pos` up to before an organisation name it would cut into.

        Fires only when the text between `pos` and an org-name tail ahead of
        it contains no punctuation -- i.e. pos really is inside one name, not
        merely before a sentence that happens to mention a company. Backs up
        to just after the previous punctuation (where the name can only have
        started), and declines when that would throw away most of the head:
        losing the whole slot is a worse outcome than the cut it avoids.
        Only ever moves `pos` earlier, so it can't invalidate a capacity
        check that already passed.
        """
        if pos <= 0 or pos >= len(text):
            return pos
        window = text[pos:pos + 24]
        hits = [window.find(s) for s in cls._ORG_NAME_TAILS]
        hits = [h for h in hits if h >= 0]
        if not hits:
            return pos
        # Punctuation before the tail means pos isn't inside that name.
        if any(ch in cls._SENTENCE_PUNCT for ch in window[:min(hits)]):
            return pos
        prev = max((text.rfind(ch, 0, pos) for ch in cls._SENTENCE_PUNCT), default=-1)
        if prev < 0:
            return pos
        snapped = prev + 1
        return snapped if snapped >= pos * 0.3 else pos

    @classmethod
    def _jieba_word_boundary_snap(cls, text: str, pos: int) -> Optional[int]:
        """If jieba is installed, segment `text` and return the start index
        of whichever word strictly contains `pos`, or None if pos already
        sits on a word boundary (or jieba is unavailable/errors).

        This is the GENERAL version of the curated _PROTECTED_CJK_COMPOUNDS
        list below -- it was found to have a real gap: "结清" (settle) split
        as "...或交割前结" / "清安排..." in real production output, the
        SECOND compound found broken after the first round of fixes
        (人民币/万元/分别/年度) -- confirming a fixed list will always be
        one case behind whatever the AI writes next, since Chinese has no
        spaces to mark word boundaries structurally. jieba is a real
        Chinese-word-segmentation library (context-aware -- correctly
        keeps "784"/"万元" as separate tokens but "分别"/"年度"/"结清" as
        single ones); used here ONLY for its segmentation, no other
        behaviour change. Optional dependency, imported lazily so a
        machine without it still runs (falls back to the curated list
        below, unchanged) rather than failing PPTX generation outright.
        """
        try:
            import jieba  # type: ignore
        except ImportError:
            return None
        try:
            offset = 0
            for word in jieba.cut(text):
                word_len = len(word)
                if offset < pos < offset + word_len:
                    return offset
                offset += word_len
                if offset >= pos:
                    break
            return None
        except Exception:
            return None

    @classmethod
    def _snap_split_before_number(cls, text: str, pos: int) -> int:
        """If `pos` sits inside a numeric literal (digits, thousands commas,
        decimal point), back up to the start of that literal so the whole
        number stays together on one side of the split. Also backs up past
        an immediately-preceding currency marker (CNY/RMB/人民币/$/¥/...) so
        a bare currency symbol doesn't get stranded with its amount on the
        far side of the split, past any word jieba's segmentation says
        `pos` would split in half (see _jieba_word_boundary_snap), and past
        any _PROTECTED_CJK_COMPOUNDS entry as a fallback/supplement when
        jieba isn't installed or doesn't treat it as one token. Only ever
        moves pos earlier, so it can't undo a capacity check that already
        accepted this position."""
        if pos <= 0 or pos >= len(text):
            return pos
        # Organisation names first: this is the one case jieba actively
        # MASKS rather than merely misses, because every internal boundary of
        # a company name is a real word boundary to it. Applied before the
        # other snaps (and not at the end) because this function has several
        # early returns; each snap only ever moves pos earlier, so running
        # this first is safe and covers every exit path.
        pos = cls._snap_before_org_name(text, pos)
        if pos <= 0:
            return pos
        numeric_chars = set('0123456789,.')
        jieba_snap = cls._jieba_word_boundary_snap(text, pos)
        if jieba_snap is not None:
            pos = jieba_snap
            if pos <= 0:
                return pos
        # Units glued directly to the number immediately before them --
        # magnitude units (784万元) and date suffixes (2024年/12月/31日)
        # are the same shape of bug: jieba correctly tokenizes the number
        # and the unit as SEPARATE tokens (confirmed: jieba.cut gives
        # ['784', '万元'] and ['31', '日']), so once jieba has already
        # snapped `pos` to that clean token boundary, the straddle-scan
        # below (which only fires when `pos` sits INSIDE the unit itself)
        # never gets a chance to also pull the number back -- jieba
        # correctly stops a mid-unit split but, on its own, still leaves
        # the number stranded from its own unit one token earlier. Two
        # real production cases hit this: "784万" | "元..." before jieba
        # existed, and "...2024年12月31" | "日及2025年12月31日..." AFTER
        # jieba was added (jieba's OWN snap masked the marker loop's
        # number-pullback special case below, since by the time that loop
        # runs, `pos` no longer straddles "万元" -- it sits cleanly right
        # before it). Checking here, unconditionally, right after the
        # jieba snap and before the straddle-scan, closes both cases.
        for _suffix in ('万元', '亿元', '年', '月', '日'):
            if text[pos:pos + len(_suffix)] == _suffix and text[pos - 1] in numeric_chars:
                num_start = pos
                while num_start > 0 and text[num_start - 1] in numeric_chars:
                    num_start -= 1
                if num_start < pos:
                    return num_start
        # `pos` landing strictly INSIDE a multi-character marker/compound
        # itself (e.g. between "人民" and "币", or "分" and "别") is a
        # different failure mode from the number/marker-boundary case
        # below -- neither text[pos-1] nor text[pos] is a digit, so the
        # numeric-literal check never fires and this position was never
        # being checked at all.
        for marker in cls._CURRENCY_MARKERS + cls._PROTECTED_CJK_COMPOUNDS:
            if len(marker) < 2:
                continue  # single-char symbols ($/¥/£/€) can't be split mid-marker
            # Check every position `start` such that the marker, if present
            # there, would straddle pos (start < pos < start+len(marker)).
            # str.find's end argument is EXCLUSIVE, so searching up to `pos`
            # can never match an occurrence whose last character IS at pos
            # -- exactly the case being guarded against -- hence the direct
            # index scan instead.
            for start in range(max(0, pos - len(marker) + 1), pos):
                if text[start:start + len(marker)] == marker:
                    # 万元/亿元 are magnitude units directly attached to the
                    # number before them ("784万元") -- stranding the bare
                    # number from its own unit reads almost as broken as
                    # splitting the unit itself would. Back up further past
                    # any digits/comma/decimal-point run immediately before
                    # the unit, same treatment as a currency marker glued
                    # to its amount below.
                    if marker in ('万元', '亿元'):
                        num_start = start
                        while num_start > 0 and text[num_start - 1] in numeric_chars:
                            num_start -= 1
                        if num_start < start:
                            return num_start
                    return start
        if text[pos - 1] in numeric_chars and text[pos] in numeric_chars:
            start = pos
            while start > 0 and text[start - 1] in numeric_chars:
                start -= 1
            pos = start
        # pos is now either unchanged (wasn't mid-number) or at the start of
        # a numeric literal -- either way, check for a currency marker (and
        # an optional single space) immediately before it, and back up past
        # that too if found.
        if pos < len(text) and text[pos] in set('0123456789'):
            for marker in cls._CURRENCY_MARKERS:
                candidate_start = pos - len(marker)
                if candidate_start >= 0 and text[candidate_start:pos] == marker:
                    return candidate_start
                candidate_start_sp = pos - len(marker) - 1
                if candidate_start_sp >= 0 and text[candidate_start_sp:pos] == marker + ' ':
                    return candidate_start_sp
        return pos

    def _try_partial_split_into_gap(
        self,
        cur_accts: List[Dict[str, Any]],
        nxt_accts: List[Dict[str, Any]],
        cur_used: float,
        cur_cap: int,
        cur_name: str,
        cur_shape,
        nxt_cap: int,
        nxt_name: str,
        nxt_shape,
        is_last_cur: bool,
        is_last_nxt: bool,
        target_fill: float,
        statement_type: Optional[str],
    ) -> bool:
        """Mutates `cur_accts`/`nxt_accts` in place if it commits. Attempts
        to move the FRONT part of nxt_accts[0]'s commentary back into
        cur_accts' remaining capacity, leaving the rest as a continuation
        at the front of nxt_accts. Generalizes the initial forward-fill
        pass's one-time split into a repeatable operation any boundary can
        trigger -- so a single account can end up split across more than
        two slots when that's what it takes to close a gap, instead of
        being force-moved whole (leaving the gap unfilled) whenever it
        doesn't fit. Only commits when it strictly lowers the pair's
        summed underfill penalty and both fragments clear a minimum-fill
        safeguard -- never overflows a box, never drops text.
        """
        # This floor is what caps the achievable fill: refusing to act below
        # it leaves exactly that much of every slot empty. At the old
        # hardcoded 1.5 a 23.9-line slot could never exceed
        # (23.9-1.5)/23.9 = 94%, which is precisely the fill reported on a
        # real deck. 1.0 is the meaningful floor -- below one rendered line
        # there is nothing worth moving, and _split_commentary_at_boundary's
        # own min_available_visual already refuses smaller budgets.
        # Budget against capacity PLUS the protrusion allowance, not capacity
        # alone. The acceptance test below already lets a candidate exceed
        # cur_cap by that much, so gating entry on the strict gap made the
        # two disagree: a real deck left a slot at 25.6/26.0 -- a 0.4-line
        # gap, under the 1.0 floor -- so nothing was even attempted, while a
        # ~2-line sentence from the next column would have been accepted had
        # it been offered. "小小溢出其實問題也不大" is exactly this case.
        tail_tol = self._tail_overflow_tolerance_units(statement_type)
        gap = (cur_cap + tail_tol) - cur_used
        if gap < float(self._packing_settings(statement_type).get("min_gap_to_split_lines", 1.0) or 1.0):
            return False

        head_acct = nxt_accts[0]
        cur_last_cat = str(cur_accts[-1].get("category", "") or "") if cur_accts else ""
        head_cat = str(head_acct.get("category", "") or "")
        category_gap_cost = 1.0 if (head_cat and head_cat != cur_last_cat) else 0.0
        text_budget = gap - category_gap_cost
        if text_budget < 1.0:
            return False

        is_chinese = self._account_is_chinese(head_acct)
        # The split-point estimate (char-count based) and the real measurement
        # (_compute_slot_used_lines, font-metric based) don't perfectly agree --
        # shrink the requested budget and retry a few times if the chosen head
        # measures slightly over cur_cap, rather than giving up on the first
        # miss. Each retry backs off by the exact measured overage (plus a
        # small margin), so this converges in 1-2 extra tries in practice.
        part1 = None
        trial_cur_used = None
        remaining_budget = text_budget
        for _attempt in range(4):
            split_result = self._split_commentary_at_boundary(
                str(head_acct.get("commentary", "") or ""),
                remaining_budget,
                slot_name=cur_name,
                is_chinese=is_chinese,
                shape=cur_shape,
                statement_type=statement_type,
                key_prefix=f"■ {self._account_cost_key(head_acct)} - ",
                # `gap` above is already capacity + allowance.
                overflow_allowance_units=0.0,
            )
            if not split_result:
                return False
            head_text, tail_text = split_result

            candidate = head_acct.copy()
            candidate["commentary"] = head_text
            candidate["is_partial"] = True
            candidate["part_num"] = int(head_acct.get("part_num") or 1)
            candidate["original_key"] = head_acct.get("original_key", head_acct.get("mapping_key"))

            candidate_used = self._compute_slot_used_lines(
                cur_accts + [candidate], cur_name, slot_shape=cur_shape, statement_type=statement_type,
            )
            # The tail tolerance belongs HERE, on the acceptance test, not
            # inside the splitter. This pass PULLS content into an underfilled
            # slot, so letting the result protrude a line or two fills the box
            # -- which is what "能夠再放1-2行才滿" asks for. (An earlier attempt
            # put the tolerance inside _split_commentary_at_boundary as a
            # "refuse to split" branch; that made the caller decline the pull
            # entirely and left the slot EMPTIER. Accepting more content and
            # declining to split are opposite things.)
            if candidate_used <= cur_cap + self._tail_overflow_tolerance_units(statement_type):
                part1, trial_cur_used = candidate, candidate_used
                break
            overage = candidate_used - cur_cap
            remaining_budget -= overage + 0.25
            if remaining_budget < 1.0:
                return False

        if part1 is None:
            return False  # estimate/actual mismatch never converged -- bail out safely

        part2 = head_acct.copy()
        part2["commentary"] = tail_text
        part2["is_continuation"] = True
        part2["part_num"] = int(head_acct.get("part_num") or 1) + 1
        part2["original_key"] = head_acct.get("original_key", head_acct.get("mapping_key"))

        trial_nxt_accts = [part2] + nxt_accts[1:]
        trial_nxt_used = self._compute_slot_used_lines(
            trial_nxt_accts, nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
        )
        if trial_nxt_used > nxt_cap:
            return False

        orig_nxt_used = self._compute_slot_used_lines(
            nxt_accts, nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
        )
        cur_penalty = 0.0 if is_last_cur else max(0.0, target_fill - (cur_used / cur_cap if cur_cap else 0.0))
        nxt_penalty = 0.0 if is_last_nxt else max(0.0, target_fill - (orig_nxt_used / nxt_cap if nxt_cap else 0.0))
        current_total = cur_penalty + nxt_penalty

        alt_cur_penalty = 0.0 if is_last_cur else max(0.0, target_fill - (trial_cur_used / cur_cap if cur_cap else 0.0))
        alt_nxt_penalty = 0.0 if is_last_nxt else max(0.0, target_fill - (trial_nxt_used / nxt_cap if nxt_cap else 0.0))
        alt_total = alt_cur_penalty + alt_nxt_penalty

        if alt_total >= current_total - 1e-9:
            return False

        cur_accts.append(part1)
        nxt_accts[0] = part2
        logger.info(
            "  Split '%s' across boundary: moved a fragment into slot %s, continuation stays in slot %s "
            "(pair penalty %.3f -> %.3f)",
            head_acct.get("mapping_key", "?"), cur_name, nxt_name, current_total, alt_total,
        )
        return True

    def _rebalance_underfilled_boundaries(
        self,
        assignment: List[List[Dict[str, Any]]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> None:
        """Mutates `assignment` in place. Generalizes
        _rebalance_lopsided_lr_pairs beyond the "one slot totally empty"
        case: for every adjacent slot boundary (in packing order,
        including same-slide L/R pairs), if the next slot's first account
        would both fit in the current slot's remaining capacity AND
        strictly lower the DP's own lexicographic underfill penalty summed
        across the pair, move it back. If the whole account doesn't fit,
        falls back to _try_partial_split_into_gap to move just enough of
        its commentary to close the gap, leaving the rest as a
        continuation -- so a single account can be split across more than
        two slots when that's what filling every boundary takes.

        This is the exact same boundary check inspect_databook.py's
        analyze_population_fill uses to verify (not just guess) that a
        boundary is a genuine gap rather than already optimal -- turned
        from a passive diagnostic into an active fix. Runs to a fixed
        point (bounded by len(assignment) passes) since a move can only
        ever enable ANOTHER move at the same or a later boundary, never
        undo one, so a short chain of small imbalances in a row all get
        resolved, not just the first.
        """
        packing = self._packing_settings(statement_type)
        target_fill = float(packing.get("target_fill_min_ratio", 0.95) or 0.95)
        n = len(assignment)

        for _pass in range(n):
            changed = False
            for i in range(n - 1):
                nxt_accts = assignment[i + 1]
                if not nxt_accts:
                    continue
                cur_accts = assignment[i]
                cur_slot, nxt_slot = slots[i], slots[i + 1]
                cur_cap, nxt_cap = cur_slot["capacity"], nxt_slot["capacity"]
                cur_shape, nxt_shape = cur_slot["shape"], nxt_slot["shape"]
                cur_name, nxt_name = cur_slot["slot_name"], nxt_slot["slot_name"]

                cur_used = self._compute_slot_used_lines(
                    cur_accts, cur_name, slot_shape=cur_shape, statement_type=statement_type,
                )
                alt_cur_used = self._compute_slot_used_lines(
                    cur_accts + [nxt_accts[0]], cur_name, slot_shape=cur_shape, statement_type=statement_type,
                )
                is_last_cur = (i == n - 1)
                is_last_nxt = (i + 1 == n - 1)

                if alt_cur_used > cur_cap:
                    if self._try_partial_split_into_gap(
                        cur_accts, nxt_accts, cur_used, cur_cap, cur_name, cur_shape,
                        nxt_cap, nxt_name, nxt_shape, is_last_cur, is_last_nxt,
                        target_fill, statement_type,
                    ):
                        changed = True
                    continue  # whole move doesn't fit -- partial split (if any) already handled

                # A same-slide L/R pair with nxt down to its LAST account is exactly
                # what _rebalance_lopsided_lr_pairs exists to prevent (a visibly empty
                # column next to a full one, exposing the template's raw placeholder
                # text) -- but when nxt is also the statement's last slot overall, its
                # penalty is unconditionally exempt (0.0 below regardless of fill), so
                # this whole-move would look like a free win and drain nxt right back
                # to empty, silently undoing that earlier fix. Refuse this specific
                # move; every other boundary (including L/R pairs with 2+ accounts
                # left in nxt) is untouched.
                same_slide_lr_pair = (
                    cur_slot["slide_idx"] == nxt_slot["slide_idx"]
                    and {cur_name, nxt_name} == {"L", "R"}
                )
                if same_slide_lr_pair and len(nxt_accts) == 1:
                    continue

                nxt_used = self._compute_slot_used_lines(
                    nxt_accts, nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
                )
                alt_nxt_used = self._compute_slot_used_lines(
                    nxt_accts[1:], nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
                )

                cur_penalty = 0.0 if is_last_cur else max(0.0, target_fill - (cur_used / cur_cap if cur_cap else 0.0))
                nxt_penalty = 0.0 if is_last_nxt else max(0.0, target_fill - (nxt_used / nxt_cap if nxt_cap else 0.0))
                current_total = cur_penalty + nxt_penalty

                alt_cur_penalty = 0.0 if is_last_cur else max(0.0, target_fill - (alt_cur_used / cur_cap if cur_cap else 0.0))
                alt_nxt_penalty = 0.0 if is_last_nxt else max(0.0, target_fill - (alt_nxt_used / nxt_cap if nxt_cap else 0.0))
                alt_total = alt_cur_penalty + alt_nxt_penalty

                if alt_total < current_total - 1e-9:
                    moved = nxt_accts.pop(0)
                    cur_accts.append(moved)
                    changed = True
                    logger.info(
                        "  Rebalanced underfilled boundary: moved '%s' from slot %s into slot %s "
                        "(pair penalty %.3f -> %.3f)",
                        moved.get("mapping_key", "?"), i + 1, i, current_total, alt_total,
                    )
            if not changed:
                break

    def _try_partial_split_overflow_forward(
        self,
        cur_accts: List[Dict[str, Any]],
        nxt_accts: List[Dict[str, Any]],
        cur_cap: int,
        cur_name: str,
        cur_shape,
        nxt_cap: int,
        nxt_name: str,
        nxt_shape,
        statement_type: Optional[str],
    ) -> bool:
        """Mutates `cur_accts`/`nxt_accts` in place if it commits. Mirrors
        _try_partial_split_into_gap but for the opposite direction: cur is
        already overflowing, so shrink its LAST account down to whatever
        genuinely fits within cur's own capacity and push the trimmed-off
        tail forward as a continuation at the FRONT of nxt. Only commits
        when the split fits both sides — never leaves cur still over
        capacity, never pushes nxt over its own.
        """
        tail_acct = cur_accts[-1]
        other_used = self._compute_slot_used_lines(
            cur_accts[:-1], cur_name, slot_shape=cur_shape, statement_type=statement_type,
        )
        available_for_tail = cur_cap - other_used
        if available_for_tail < 1.0:
            return False

        is_chinese = self._account_is_chinese(tail_acct)
        part1 = None
        head_text = tail_text = ""
        remaining_budget = available_for_tail
        for _attempt in range(4):
            split_result = self._split_commentary_at_boundary(
                str(tail_acct.get("commentary", "") or ""),
                remaining_budget,
                slot_name=cur_name,
                is_chinese=is_chinese,
                shape=cur_shape,
                statement_type=statement_type,
                min_fill_ratio=0.3,
                key_prefix=f"■ {self._account_cost_key(tail_acct)} - ",
            )
            if not split_result:
                return False
            head_text, tail_text = split_result

            candidate = tail_acct.copy()
            candidate["commentary"] = head_text
            candidate["is_partial"] = True
            candidate["part_num"] = int(tail_acct.get("part_num") or 1)
            candidate["original_key"] = tail_acct.get("original_key", tail_acct.get("mapping_key"))

            candidate_used = self._compute_slot_used_lines(
                cur_accts[:-1] + [candidate], cur_name, slot_shape=cur_shape, statement_type=statement_type,
            )
            if candidate_used <= cur_cap:
                part1 = candidate
                break
            overage = candidate_used - cur_cap
            remaining_budget -= overage + 0.25
            if remaining_budget < 1.0:
                return False

        if part1 is None:
            return False  # estimate/actual mismatch never converged -- bail out safely

        part2 = tail_acct.copy()
        part2["commentary"] = tail_text
        part2["is_continuation"] = True
        part2["part_num"] = int(tail_acct.get("part_num") or 1) + 1
        part2["original_key"] = tail_acct.get("original_key", tail_acct.get("mapping_key"))

        trial_nxt_accts = [part2] + nxt_accts
        trial_nxt_used = self._compute_slot_used_lines(
            trial_nxt_accts, nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
        )
        if trial_nxt_used > nxt_cap:
            return False

        cur_accts[-1] = part1
        nxt_accts.insert(0, part2)
        logger.info(
            "  Split overflowing '%s' forward across boundary: head stays in current slot, "
            "continuation moves into next slot",
            tail_acct.get("mapping_key", "?"),
        )
        return True

    def _rebalance_overflowing_boundaries(
        self,
        assignment: List[List[Dict[str, Any]]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> None:
        """Mutates `assignment` in place. Complements
        _rebalance_underfilled_boundaries, which only ever pulls a
        FOLLOWING slot's content BACKWARD into a preceding slot's positive
        gap — it never pushes content FORWARD out of an already-
        overflowing slot into a following slot with spare room, because
        the DP's own lexicographic penalty formula caps an overflowing
        slot's penalty at 0.0 (identical to a perfectly-filled slot), so
        overflow never registers as something worth fixing to either of
        the existing passes (confirmed via a real GPT-5.5/workbench IS
        export: one column measured 103%/OVERFLOW RISK while its same-
        slide sibling sat at 34% underfilled, and the population
        diagnostic's own genuine-gap check — which reuses that same
        penalty formula — reported no fixable gap at all).

        This pass looks at raw used-vs-capacity instead of the penalty
        formula specifically to catch that blind spot: for every adjacent
        boundary where the current slot's real content exceeds its
        capacity and the next slot has spare room, tries moving the
        current slot's LAST account forward whole; if that doesn't fit
        (or cur only has one account to begin with), falls back to
        _try_partial_split_overflow_forward. Never commits a move that
        would newly overflow the next slot.
        """
        n = len(assignment)
        for _pass in range(n):
            changed = False
            for i in range(n - 1):
                cur_accts = assignment[i]
                if not cur_accts:
                    continue
                nxt_accts = assignment[i + 1]
                cur_slot, nxt_slot = slots[i], slots[i + 1]
                cur_cap, nxt_cap = cur_slot["capacity"], nxt_slot["capacity"]
                cur_shape, nxt_shape = cur_slot["shape"], nxt_slot["shape"]
                cur_name, nxt_name = cur_slot["slot_name"], nxt_slot["slot_name"]

                cur_used = self._compute_slot_used_lines(
                    cur_accts, cur_name, slot_shape=cur_shape, statement_type=statement_type,
                )
                if cur_used <= cur_cap:
                    continue  # not overflowing -- nothing for this pass to do at this boundary

                nxt_used = self._compute_slot_used_lines(
                    nxt_accts, nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
                ) if nxt_accts else 0.0
                if nxt_used >= nxt_cap:
                    continue  # next slot has no spare room either -- can't help here

                # Moving just the SINGLE last account used to be the only
                # whole-account move attempted -- if the accounts BEFORE it
                # already exceed cur_cap on their own (several trailing
                # accounts collectively overflow, not just the tail one),
                # `rest_used <= cur_cap` was never true, this whole branch
                # fell through, and _try_partial_split_overflow_forward's own
                # `other_used` (the same all-but-last cost) was equally over
                # cur_cap, so IT bailed immediately too -- nothing ever moved
                # at all. Confirmed via a real production case: a slot at
                # 109% sat next to a COMPLETELY EMPTY next slot (0% -- not
                # merely underfilled) because neither path could ever fire.
                # Grow the batch of trailing whole accounts to move one at a
                # time until what remains in cur actually fits, or until
                # moving the next one would overflow nxt -- whichever comes
                # first. Single-account overflow (the original case) is
                # k=1 here, unchanged in behavior.
                move_count = 0
                for k in range(1, len(cur_accts)):
                    rest_used = self._compute_slot_used_lines(
                        cur_accts[:-k], cur_name, slot_shape=cur_shape, statement_type=statement_type,
                    )
                    trial_nxt_accts = cur_accts[-k:] + nxt_accts
                    trial_nxt_used = self._compute_slot_used_lines(
                        trial_nxt_accts, nxt_name, slot_shape=nxt_shape, statement_type=statement_type,
                    )
                    if trial_nxt_used > nxt_cap:
                        break
                    move_count = k
                    if rest_used <= cur_cap:
                        break
                if move_count:
                    moved = cur_accts[-move_count:]
                    del cur_accts[-move_count:]
                    nxt_accts[0:0] = moved
                    changed = True
                    logger.info(
                        "  Rebalanced overflowing boundary: moved %d whole account(s) forward from slot %s "
                        "(was %.1f/%s) into slot %s",
                        move_count, i, cur_used, cur_cap, i + 1,
                    )
                    continue

                # No whole account could move. The only remaining option is
                # cutting one in half -- and a small overflow is not worth
                # that. The project team accepts 1-2 lines protruding below
                # the box; what they rejected was a split landing mid-name
                # ("某某系统" | "工程第四建设有限公司"). Note this gate is
                # deliberately AFTER the whole-account move above: moving an
                # account forward is free and still happens for any overflow.
                if cur_used <= cur_cap + self._tail_overflow_tolerance_units(statement_type):
                    logger.info(
                        "  Slot %s left overflowing %.2f line(s) over %s on purpose — "
                        "within tail tolerance, so not split",
                        i, cur_used - cur_cap, cur_cap,
                    )
                    continue

                if self._try_partial_split_overflow_forward(
                    cur_accts, nxt_accts, cur_cap, cur_name, cur_shape,
                    nxt_cap, nxt_name, nxt_shape, statement_type,
                ):
                    changed = True
            if not changed:
                break

    def _extend_continuation_in_place(
        self,
        cur_accts: List[Dict[str, Any]],
        nxt_accts: List[Dict[str, Any]],
        cur_name: str,
        cur_shape,
        cur_cap: float,
        statement_type: Optional[str],
    ) -> bool:
        """Mutates cur_accts/nxt_accts in place. Pulls more of nxt_accts[0]
        (a direct continuation of cur_accts[-1], same original account)
        into cur_accts[-1]'s own paragraph, instead of adding it as a new
        one. Returns True if anything changed.

        Why this exists: a NEW paragraph costs at minimum line_h+para_gap
        (~1.0 std_lh unit) even for a single character, since the wrap+
        cost formula charges a full line height for any non-empty text
        plus the fixed per-paragraph gap. _maximize_forward_fill's normal
        "add next account as a new paragraph" path can therefore never
        fill a gap smaller than that floor, however much real room is
        left. Extending an EXISTING paragraph's text only adds wrapped
        LINES (no extra para_gap), so it can use up a fractional-line gap
        the new-paragraph path mathematically cannot.
        """
        head_acct = cur_accts[-1]
        tail_acct = nxt_accts[0]
        head_commentary = str(head_acct.get("commentary", "") or "")
        tail_commentary = str(tail_acct.get("commentary", "") or "")
        if not tail_commentary.strip():
            return False
        combined = (head_commentary + tail_commentary).strip()

        other_used = self._compute_slot_used_lines(
            cur_accts[:-1], cur_name, slot_shape=cur_shape, statement_type=statement_type,
        )
        budget = cur_cap - other_used
        if budget <= 0:
            return False

        # First check whether the WHOLE combined text now fits -- if so,
        # absorb tail_acct entirely rather than leaving an artificially
        # truncated remainder.
        candidate_whole = dict(head_acct)
        candidate_whole["commentary"] = combined
        whole_used = self._compute_slot_used_lines(
            cur_accts[:-1] + [candidate_whole], cur_name, slot_shape=cur_shape, statement_type=statement_type,
        )
        if whole_used <= cur_cap:
            cur_accts[-1] = candidate_whole
            nxt_accts.pop(0)
            return True

        # Deliberately NOT using _split_commentary_at_boundary here: its
        # rough char-count pre-check can decide "the whole paragraph fits"
        # (a coarser, disagreeing estimate from the accurate whole_used
        # check just above) and return None on that basis alone, before
        # ever reaching its own accurate-measurer validation -- a real
        # false negative observed here (whole_used said no, but the
        # function still gave up immediately). A direct binary search
        # against the same accurate _compute_slot_used_lines this whole
        # pass already trusts sidesteps that disagreement entirely.
        lo, hi = len(head_commentary), len(combined) - 1
        best_len = len(head_commentary)
        while lo <= hi:
            mid = (lo + hi) // 2
            candidate = dict(head_acct)
            candidate["commentary"] = combined[:mid].strip()
            used = self._compute_slot_used_lines(
                cur_accts[:-1] + [candidate], cur_name, slot_shape=cur_shape, statement_type=statement_type,
            )
            if used <= cur_cap:
                best_len = mid
                lo = mid + 1
            else:
                hi = mid - 1

        if best_len <= len(head_commentary):
            return False

        # The binary search above optimises for FILL only, so best_len is a
        # raw character offset -- and _snap_split_before_number below only
        # rescues numbers and jieba tokens, not names (系统|工程 is a valid
        # token boundary, which is how a real deck ended up cutting
        # 某某系统 | 工程第四建设有限公司). Snap back to the last natural
        # sentence/clause boundary at or before best_len -- it necessarily
        # still fits, and it must still be a real extension of the head.
        _marks = ('。', '！', '？', '，', '；', '、', '.', ',', ';')
        _back = max((combined.rfind(c, 0, best_len + 1) for c in _marks), default=-1)
        if _back >= 0 and _back + 1 > len(head_commentary):
            best_len = _back + 1
        else:
            # No clause boundary anywhere inside the text this pass wanted to
            # pull forward. Extending anyway would keep the binary search's
            # raw character offset -- a mid-word cut, which is what showed up
            # the moment the gap gate started admitting smaller gaps. There
            # is nothing here worth a bad cut, so decline the extension.
            return False

        best_len = self._snap_split_before_number(combined, best_len)
        if best_len <= len(head_commentary):
            return False
        new_head = combined[:best_len].strip()
        new_tail = combined[best_len:].strip()
        if not new_head or not new_tail:
            return False

        candidate = dict(head_acct)
        candidate["commentary"] = new_head
        candidate_used = self._compute_slot_used_lines(
            cur_accts[:-1] + [candidate], cur_name, slot_shape=cur_shape, statement_type=statement_type,
        )
        if candidate_used > cur_cap:
            return False

        cur_accts[-1] = candidate
        tail_acct["commentary"] = new_tail
        return True

    def _maximize_forward_fill(
        self,
        assignment: List[List[Dict[str, Any]]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> None:
        """Mutates `assignment` in place. User-confirmed requirement:
        given a FIXED amount of content, pack it the way a person would --
        fill the first slot as close to its own capacity as possible,
        THEN move on to the next slot with whatever's left, and so on.
        This is a plain greedy maximal-fill-in-reading-order pass,
        deliberately NOT balance-seeking (an empty or lightly-filled
        TRAILING slot is fine and expected once content genuinely runs
        out -- that's a separate, not-yet-requested question of whether
        to use fewer total slides). It directly supersedes _optimize_
        slot_fill's own DP objective ("minimise the MAXIMUM fill ratio
        across slots"), which spreads content for evenness rather than
        maximising each slot in turn -- runs LAST, after every other
        rebalance pass, since it's the final word on "is slot i as full
        as slot i's own capacity allows" and nothing after it should
        undo that.

        For each boundary in order: repeatedly pull the next slot's
        first account into the current slot, whole if it fits; if it
        doesn't fit whole, split it (same paragraph/sentence/word
        boundary search _try_partial_split_into_gap uses) to use up
        whatever room remains, leaving the rest as a continuation at the
        front of the next slot. Keeps going until the current slot has
        no meaningful room left (<0.5 lines) or the next slot is empty,
        then moves on to the next boundary.
        """
        n = len(assignment)
        for i in range(n - 1):
            cur_slot, nxt_slot = slots[i], slots[i + 1]
            cur_name = cur_slot["slot_name"]
            cur_shape = cur_slot["shape"]
            cur_cap = cur_slot["capacity"]

            while assignment[i + 1]:
                cur_accts = assignment[i]
                nxt_accts = assignment[i + 1]
                cur_used = self._compute_slot_used_lines(
                    cur_accts, cur_name, slot_shape=cur_shape, statement_type=statement_type,
                )
                gap = cur_cap - cur_used
                # Lowered from 0.5 -- with the accurate-measurer backoff/trim
                # now in _split_commentary_at_boundary (5fd871e) and the
                # number/currency-safe split points (3149651/f88de44), a
                # small gap can safely be attempted rather than written off;
                # the 1.0-unit floors below this used to compound into a
                # worst-case ~2-line leftover (this gate PLUS a category-
                # transition cost eating into the split attempt), which is
                # exactly the "some pages 2 lines short" the user reported
                # even after those other fixes landed.
                if gap < 0.2:
                    break

                head_acct = nxt_accts[0]

                # A brand-new paragraph costs AT LEAST line_h+para_gap
                # (~1.0 std_lh unit) no matter how short its text is -- the
                # wrap+cost formula charges a full line height for even a
                # single character, plus the fixed per-paragraph gap. So
                # whenever gap < ~1.0 and the next slot's head is a DIRECT
                # continuation of the account already at the end of THIS
                # slot, extending that existing paragraph's own text (which
                # only adds wrapped LINES, no extra para_gap) can use up a
                # gap the "add as a new paragraph" logic below can never
                # fill by construction -- confirmed via a real production
                # case: gap=0.641 units (8.85pt) sat permanently unfilled
                # because even a 1-character new split costs >=13.8pt
                # (1 line + the paragraph gap), a mathematical floor, not a
                # tuning knob. Try this FIRST, before the new-paragraph path.
                if (
                    cur_accts
                    and head_acct.get("is_continuation")
                    and cur_accts[-1].get("mapping_key") == head_acct.get("mapping_key")
                    and cur_accts[-1].get("original_key", cur_accts[-1].get("mapping_key"))
                        == head_acct.get("original_key", head_acct.get("mapping_key"))
                ):
                    if self._extend_continuation_in_place(
                        cur_accts, nxt_accts, cur_name, cur_shape, cur_cap, statement_type,
                    ):
                        continue

                whole_used = self._compute_slot_used_lines(
                    cur_accts + [head_acct], cur_name, slot_shape=cur_shape, statement_type=statement_type,
                )
                # Prefer keeping the account WHOLE over a tidy split, up to
                # the tail tolerance. Strict capacity here was the last place
                # still ignoring it: the team accepts 1-2 lines protruding and
                # has said so repeatedly, but this line split an account the
                # moment it exceeded capacity by any amount at all. On a real
                # deck that cut 其他非流动资产 across the page-1/page-2
                # boundary while page 1 rendered at 94% -- roughly 1.5 unused
                # lines under the split. The tolerance is the same value
                # inspect_pptx.py already declines to flag as overflow, so
                # this cannot produce a warning that was not already accepted.
                if whole_used <= cur_cap + self._tail_overflow_tolerance_units(statement_type):
                    cur_accts.append(nxt_accts.pop(0))
                    continue

                # Doesn't fit whole -- split it to use up exactly what room remains.
                cur_last_cat = str(cur_accts[-1].get("category", "") or "") if cur_accts else ""
                head_cat = str(head_acct.get("category", "") or "")
                category_gap_cost = 1.0 if (head_cat and head_cat != cur_last_cat) else 0.0
                text_budget = gap - category_gap_cost
                # Lowered from 1.0 in step with _split_commentary_at_boundary's
                # own internal floor (also lowered below) -- see the gap<0.2
                # comment above for why 1.0 here was leaving up to ~2 real
                # lines on the table in the worst case.
                if text_budget < 0.5:
                    break

                is_chinese = self._account_is_chinese(head_acct)
                part1 = None
                head_text = tail_text = ""
                remaining_budget = text_budget
                # 8 attempts, and each failed candidate shrinks the budget by
                # 1.5x its overage (not 1x) -- _split_commentary_at_boundary's
                # own split-point search is a coarse chars_per_line estimate,
                # not the accurate Pillow/client-metrics measurer this loop's
                # own candidate_used check uses, so the two can disagree by a
                # small, fairly consistent margin. A 1x reduction often lands
                # on the exact same sentence/comma boundary again (nothing
                # between the old and new budget crossed a punctuation mark),
                # burning attempts without changing the answer; overshooting
                # the reduction converges in fewer tries.
                for _attempt in range(8):
                    split_result = self._split_commentary_at_boundary(
                        str(head_acct.get("commentary", "") or ""),
                        remaining_budget,
                        slot_name=cur_name,
                        is_chinese=is_chinese,
                        shape=cur_shape,
                        statement_type=statement_type,
                        # Lower than the other split passes' 0.3 -- this
                        # pass's whole point is "use up as much of the
                        # remaining gap as possible," so a smaller carved-
                        # off fragment is still worth taking here even
                        # where it wouldn't be for a balance-oriented move.
                        min_fill_ratio=0.15,
                        key_prefix=f"■ {self._account_cost_key(head_acct)} - ",
                        # Also lower than the shared 1.0 default, in step
                        # with this pass's own gap<0.2/text_budget<0.5 gates
                        # above -- same reasoning: this is the last-word,
                        # maximize-fill pass, not a balance-oriented one.
                        min_available_visual=0.5,
                    )
                    if not split_result:
                        break
                    head_text, tail_text = split_result
                    candidate = head_acct.copy()
                    candidate["commentary"] = head_text
                    candidate["is_partial"] = True
                    candidate["part_num"] = int(head_acct.get("part_num") or 1)
                    candidate["original_key"] = head_acct.get("original_key", head_acct.get("mapping_key"))
                    candidate_used = self._compute_slot_used_lines(
                        cur_accts + [candidate], cur_name, slot_shape=cur_shape, statement_type=statement_type,
                    )
                    if candidate_used <= cur_cap:
                        part1 = candidate
                        break
                    overage = candidate_used - cur_cap
                    remaining_budget -= (overage * 1.5) + 0.25
                    if remaining_budget < 0.5:
                        break

                if part1 is None:
                    break  # can't split to fit either -- nothing more fits in this slot

                part2 = head_acct.copy()
                part2["commentary"] = tail_text
                part2["is_continuation"] = True
                part2["part_num"] = int(head_acct.get("part_num") or 1) + 1
                part2["original_key"] = head_acct.get("original_key", head_acct.get("mapping_key"))

                cur_accts.append(part1)
                nxt_accts[0] = part2
                logger.info(
                    "  Maximized forward fill: split '%s' to top up slot %s (gap was %.1f)",
                    head_acct.get("mapping_key", "?"), i, gap,
                )
                continue

    def _consolidate_trailing_near_empty_slot(
        self,
        assignment: List[List[Dict[str, Any]]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> None:
        """User-reported case: a whole trailing slide held only a single
        leftover sentence spilled from the slide before it (e.g. one
        continuation of an already-mostly-placed account), while that
        prior slide was already packed to ~100% -- _maximize_forward_fill
        correctly refuses to grow a slot past its own capacity, so this
        sliver never had anywhere to go and sat alone on an otherwise
        blank page. Eliminating a whole near-empty trailing page is worth
        more than staying strictly under 100% on the page before it, so
        this runs LAST and explicitly accepts a small bounded overflow
        (up to 15% over nominal capacity -- well within what
        _apply_bounded_autofit's shrink, floored at _BOUNDED_AUTOFIT_MIN_
        SCALE=0.70, can visually absorb) specifically to empty this slot
        out entirely. An emptied slot's slide becomes eligible for the
        existing unused-slide removal, collapsing the whole page.

        Only the true tail of the WHOLE statement -- not every small
        slot -- since folding a small-but-legitimate MIDDLE slot forward
        would just recreate the same problem one slot earlier.
        """
        n = len(assignment)
        if n < 2:
            return
        last_i = n - 1
        while last_i > 0 and not assignment[last_i]:
            last_i -= 1
        if last_i <= 0:
            return  # nothing trails, or the whole statement is one slot

        last_slot = slots[last_i]
        last_accts = assignment[last_i]
        last_used = self._compute_slot_used_lines(
            last_accts, last_slot["slot_name"], slot_shape=last_slot["shape"], statement_type=statement_type,
        )
        # Only a genuinely tiny leftover -- not a legitimately-substantial
        # trailing slot that just happens to be under-filled because its
        # own statement's content ran out (that's expected and fine, per
        # _maximize_forward_fill's own docstring).
        if last_used > 3.0 or last_used > 0.25 * last_slot["capacity"]:
            return

        prev_i = last_i - 1
        prev_slot = slots[prev_i]
        prev_accts = assignment[prev_i]
        if not prev_accts:
            return
        combined_used = self._compute_slot_used_lines(
            prev_accts + last_accts, prev_slot["slot_name"], slot_shape=prev_slot["shape"], statement_type=statement_type,
        )
        if combined_used <= prev_slot["capacity"] * 1.15:
            assignment[prev_i] = prev_accts + last_accts
            assignment[last_i] = []
            logger.info(
                "  Consolidated trailing near-empty slot %s (%.1f units) into slot %s (now %.1f/%.1f)",
                last_i, last_used, prev_i, combined_used, prev_slot["capacity"],
            )

    def _optimize_slot_fill(
        self,
        distribution: List[tuple],
        slot_shapes: Optional[Dict[int, Any]] = None,
        slot_meta: Optional[List[Tuple[int, str]]] = None,
        statement_type: Optional[str] = None,
    ) -> List[tuple]:
        """Dynamic-programming balanced partition.

        Flattens all accounts into reading order, then partitions them into
        contiguous groups (one per slot) so that the maximum slot fill ratio
        is minimised. Line counts come from _compute_slot_used_lines measured
        against each slot's actual shape, so when Pillow fitting is enabled
        this uses real font metrics. Preserves reading order; drops trailing
        empty slots.

        DP: dp[s][i] = min achievable "max fill ratio" when placing
        accounts[0..i] into slots[0..s]. Transition: slot s takes a suffix
        accounts[j+1..i]; combine with dp[s-1][j]. O(S * N^2) states, but
        N ≤ ~20 accounts and S ≤ ~8 slots in practice, so this is trivial.
        """
        if not distribution:
            return distribution

        slot_lookup: Dict[Tuple[int, str], Any] = {}
        if slot_meta and slot_shapes:
            for slot_idx, (s_idx, s_name) in enumerate(slot_meta):
                slot_lookup[(s_idx, s_name)] = slot_shapes.get(slot_idx)

        def _resolve_shape(slide_idx: int, slot_name: str):
            shape = slot_lookup.get((slide_idx, slot_name))
            if shape is not None:
                return shape
            try:
                slide = self.presentation.slides[slide_idx]
            except Exception:
                return None
            return self._resolve_commentary_slot_shape(slide, slot_name)

        flat_accounts: List[Dict[str, Any]] = []
        for _slide_idx, _slot_name, accounts in distribution:
            flat_accounts.extend(accounts)

        if not flat_accounts:
            return distribution

        slots: List[Dict[str, Any]] = []
        is_chinese_any = any(self._account_is_chinese(a) for a in flat_accounts)
        for slide_idx, slot_name, _accounts in distribution:
            shape = _resolve_shape(slide_idx, slot_name)
            capacity = self._calculate_max_lines_for_textbox(
                shape,
                is_chinese=is_chinese_any,
                slot_name=slot_name,
                statement_type=statement_type,
            )
            slots.append({
                "slide_idx": slide_idx,
                "slot_name": slot_name,
                "shape": shape,
                # Float, not int(...) -- see _calculate_max_lines_for_textbox;
                # re-flooring its already-precise return value here threw the
                # same up-to-one-line margin away a second time.
                "capacity": max(1.0, float(capacity or 1)),
            })

        N = len(flat_accounts)
        S = len(slots)

        # ── Pre-compute per-account content lines for each unique slot type ──
        # Key: (slot_name, shape_width_emu).  Two slots that share the same
        # name and width get the same measurements, so we only call
        # _calculate_content_lines (and Pillow when enabled) once per
        # (account, slot_type) pair — O(N × slot_types) total instead of the
        # O(S × N²) calls that the old range-based approach produced.
        _acct_cost: Dict[Tuple[int, str, int], float] = {}
        seen_slot_types: set = set()
        for slot in slots:
            shape = slot["shape"]
            w_key = int(shape.width) if shape and hasattr(shape, "width") else 0
            type_key = (slot["slot_name"], w_key)
            if type_key in seen_slot_types:
                continue
            seen_slot_types.add(type_key)
            for a_i, account in enumerate(flat_accounts):
                _acct_cost[(a_i, slot["slot_name"], w_key)] = self._calculate_content_lines(
                    "",
                    self._account_cost_key(account),
                    account.get("commentary", ""),
                    slot_name=slot["slot_name"],
                    shape=shape,
                    is_chinese=self._account_is_chinese(account),
                    statement_type=statement_type,
                )

        # ── Tight packing: use minimum slots, expand only if infeasible ─────────
        # The DP objective (min max-fill-ratio) spreads content across ALL
        # available slots at ~80% fill.  We want ~90-95%.  Fix: compute the
        # minimum number of slots that can hold all content, try that first;
        # if infeasible (split accounts can push content above S_min capacity),
        # expand by one and retry until feasible or S_orig is reached.
        _slots_all = list(slots)
        _S_orig = S

        import math as _math
        _est_sname = slots[0]["slot_name"] if slots else "single"
        _est_wkey = (
            int(slots[0]["shape"].width)
            if slots and slots[0].get("shape") and hasattr(slots[0]["shape"], "width")
            else 0
        )
        _total_est: float = 0.0
        _prev_cat_e: Optional[str] = None
        for _a_i, _acct_e in enumerate(flat_accounts):
            _cat_e = str(_acct_e.get("category", "") or "")
            if _cat_e and _cat_e != _prev_cat_e:
                _total_est += 1.0
            _prev_cat_e = _cat_e
            _total_est += _acct_cost.get((_a_i, _est_sname, _est_wkey), 2.0)
        _min_cap = min(slot["capacity"] for slot in slots) if slots else 1
        S_min = max(1, _math.ceil(_total_est / _min_cap))

        # cost_cache and slot_cost are defined before the retry loop.
        # slot_cost captures `slots` by reference — updating slots = _slots_all[:S_try]
        # inside the loop automatically changes what slot_cost sees.
        cost_cache: Dict[Tuple[int, int, int], float] = {}

        # Progressive relax factors for DP feasibility. Start at 1.0 (strict
        # capacity) and widen until the DP finds a partition. The final
        # factor (very large) guarantees feasibility, so the DP ALWAYS
        # returns a balanced result — we never fall through to greedy, which
        # would force-place oversized accounts and break the slide layout.
        # shape_height_utilization is the "natural" first relaxation because
        # PPT auto-fit can absorb that much overflow at render time.
        _packing_relax = self._packing_settings(statement_type)
        # Default 1.15 -> 1.00 (2026-08-04). This relax lets the packer
        # treat a slot as N% taller than it is, to avoid dropping content
        # or adding a page. It was safe only because the cost model used
        # to OVER-count height (a trailing paragraph gap per block), so a
        # "105% full" slot really rendered at ~100%. Now that the model
        # reproduces PowerPoint's own BoundHeight exactly, that hidden
        # buffer is gone and the same relax produced a real, measured
        # 6.6pt overflow (slide 2 L: 366.6pt of text in a 358.9pt box).
        # 1.00 means "fill the slot completely, never past it", which is
        # exactly the behaviour the user asked for after seeing slide 1
        # land at precisely 0pt spare: "可以用盡空間 比較好".
        _shape_util = float(_packing_relax.get("shape_height_utilization", 1.00) or 1.00)
        # NOTE: the second tier is a hardcoded 1.05 FLOOR, so configuring
        # shape_height_utilization BELOW 1.05 has no effect. Removing the
        # floor was tried (2026-08-05) and made things materially worse:
        # with it gone the next tier is 1.35, and a real 6-account export
        # front-loaded slide 1 to 109% while slide 2 sat at 63%/0%. The
        # 1.05 tier is what keeps a near-fit from falling all the way to
        # 1.35. Residual effect: a slot can still render ~1.05x full (a
        # real export measured -3.6pt, about a third of a line). Closing
        # that properly means making the DP add a slot instead of
        # relaxing -- i.e. the S_min expansion logic, not this ladder.
        _relax_factors: List[float] = [1.0, max(1.05, _shape_util), 1.35, 1.6, 10.0]

        # Front-loading target: slots before the LAST used one should be packed
        # to at least this fill ratio. target_fill_min_ratio existed in config
        # already (default 0.95) but was never read anywhere — this is the
        # first real consumer of it.
        _target_min_fill = float(_packing_relax.get("target_fill_min_ratio", 0.95) or 0.95)

        def slot_cost(s: int, j: int, i: int) -> float:
            """Return float line-units for placing flat_accounts[j..i] in slot s.
            Category headers cost 1.0; account content costs the float from
            _calculate_content_lines (actual_pt / std_lh, no ceil)."""
            if j > i:
                return 0.0
            key = (s, j, i)
            if key in cost_cache:
                return cost_cache[key]
            slot = slots[s]
            shape = slot["shape"]
            w_key = int(shape.width) if shape and hasattr(shape, "width") else 0
            sname = slot["slot_name"]
            used: float = 0.0
            prev_cat = None
            for a_i in range(j, i + 1):
                account = flat_accounts[a_i]
                cat = str(account.get("category", "") or "")
                if cat and cat != prev_cat:
                    used += 1.0  # category header line
                prev_cat = cat
                used += _acct_cost.get((a_i, sname, w_key), 0.0)
            # Each account was costed with its own trailing paragraph gap
            # included (whole_box=False), which is right for every account
            # except the LAST one in the slot -- that final gap renders as
            # invisible padding at the bottom of the frame. Refund exactly
            # one, once, per slot.
            if i >= j:
                used -= self._real_para_gap_pt(True) / (
                    self._planning_std_lh_pt(True) or 1.0
                )
            cost_cache[key] = used
            return used

        INF = float("inf")
        # Lexicographic DP state: (num_nonempty_slots, underfill_penalty).
        #
        # underfill_penalty is a TEXT-JUSTIFICATION-style cost (like the classic
        # "optimal paragraph layout" problem), NOT a load-balancing cost. Every
        # slot except the LAST one in this attempt's range is penalised for
        # falling short of _target_min_fill; the last slot is exempt (a lighter
        # final page is normal and expected — a lighter FIRST/MIDDLE page is the
        # bug this replaces). Compared as Python tuples: fewer non-empty slots
        # wins first, then lower total penalty. This front-loads content instead
        # of spreading it evenly — the previous "minimise the single worst slot"
        # objective was solved by keeping EVERY slot moderately empty (see
        # commit history: 45%/72% instead of 83%/30%), which is the opposite of
        # what a reader expects from a paginated document.
        INF_ST = (INF, INF)
        dp: List[List[Tuple[float, float]]] = []
        # split[s][i] = j such that slot s holds flat_accounts[j+1..i]; j == i
        # means slot s is empty (carries i through from previous slot).
        split: List[List[int]] = []

        # Progressive relax loop. For each factor, run the DP at full S_orig
        # slots — because the DP's own tight-packing objective will already
        # leave trailing slots empty when the content fits in fewer. The
        # final 10× factor guarantees feasibility, so we never need to fall
        # through to greedy.
        _dp_solved = False
        _solved_factor = 1.0
        S = _S_orig
        slots = _slots_all[:S]
        for _cap_mult in _relax_factors:
            cost_cache.clear()

            dp = [[INF_ST] * N for _ in range(S)]
            split = [[-1] * N for _ in range(S)]

            # Slot 0 is exempt from the underfill penalty only if it's also the
            # LAST slot in this attempt (S == 1) — a single-slot statement is
            # allowed to be light. Otherwise it must justify to the target.
            _cap0 = slots[0]["capacity"] * _cap_mult
            for i in range(N):
                lines = slot_cost(0, 0, i)
                if lines <= _cap0:
                    ratio0 = lines / slots[0]["capacity"]
                    penalty0 = 0.0 if S == 1 else max(0.0, _target_min_fill - ratio0)
                    dp[0][i] = (1.0, penalty0)
                split[0][i] = -1

            for s in range(1, S):
                cap_true = slots[s]["capacity"]
                cap_check = cap_true * _cap_mult
                is_last_slot = (s == S - 1)
                # j == -1 below means "slot s starts fresh, slots 0..s-1
                # contribute nothing" — it always compares as strictly better
                # than any real dp[s-1][j] (fewer non-empty slots wins
                # lexicographically first), so whenever it's available it's
                # ALWAYS chosen, even when an earlier slot could genuinely
                # hold content. That skipped a real coSummaryShape+table
                # slide (slot 0, smaller capacity than a plain L/R breakdown
                # slide) whenever its first content chunk was too big for
                # slot 0 alone but fit a bigger slot 1 — the DP then routed
                # everything through slot 1+, leaving slot 0 with ZERO
                # accounts. That slide was then dropped as "unused" and the
                # table-embedding code fell back to the next used slide (an
                # ordinary breakdown slide with no room for a table),
                # producing the table-overlaps-commentary bug. Only allow the
                # free (0,0) bypass when slot s-1's entire row is infeasible
                # at this relax factor — i.e. skipping ahead is the only way
                # to make progress, not a free lexicographic win over a slot
                # that could actually hold something. Real dp[s-1][j]
                # transitions (j >= 0 below) already cover every case where
                # an earlier slot legitimately holds some/all prefix content.
                prev_row_feasible = any(dp[s - 1][x][0] < INF for x in range(N))
                for i in range(N):
                    # Case A: slot s non-empty, holds accounts[j+1..i]
                    for j in range(-1, i):
                        if j < 0:
                            if prev_row_feasible:
                                continue
                            prev_state: Tuple[float, float] = (0.0, 0.0)
                        else:
                            prev_state = dp[s - 1][j]
                            if prev_state[0] >= INF:
                                continue
                        lines = slot_cost(s, j + 1, i)
                        if lines > cap_check:
                            continue
                        ratio = lines / cap_true
                        penalty = 0.0 if is_last_slot else max(0.0, _target_min_fill - ratio)
                        curr_state = (
                            prev_state[0] + 1.0,
                            prev_state[1] + penalty,
                        )
                        if curr_state <= dp[s][i]:
                            dp[s][i] = curr_state
                            split[s][i] = j
                    # Case B: slot s empty — carry dp[s-1][i] forward unchanged
                    if dp[s - 1][i] < dp[s][i]:
                        dp[s][i] = dp[s - 1][i]
                        split[s][i] = i  # marker: slot s is empty

            if dp[S - 1][N - 1][0] < INF:
                _dp_solved = True
                _solved_factor = _cap_mult
                break

            logger.info(
                "  DP infeasible at relax × %.2f; widening tolerance",
                _cap_mult,
            )

        _used_slots = int(dp[S - 1][N - 1][0]) if _dp_solved else S
        _final_penalty = dp[S - 1][N - 1][1] if _dp_solved else INF
        # logger.warning, not info: inspect_databook.py's export-log
        # analysis only surfaces WARNING/ERROR, so at INFO this line was
        # captured but never shown -- which is exactly the number needed
        # to tell "the DP thinks it packed tight but render disagrees"
        # apart from "the DP had to relax". Cheap to emit (once per
        # statement) and directly actionable, so it earns the level.
        # Also written straight to a file next to the export. Two
        # successive attempts to surface this through logging (INFO, then
        # WARNING) produced nothing on the user's machine -- something in
        # their logging setup filters it -- and this one number is what
        # decides between two opposite fixes for the residual overflow.
        # A file write can't be filtered by a handler or level.
        try:
            _msg = ("tight-pack at 1.0" if _solved_factor <= 1.0
                    else f"RELAXED to x{_solved_factor:.2f}")
            with open("dp_packing_report.txt", "a", encoding="utf-8") as _fh:
                _fh.write(
                    f"{statement_type}: {_msg} | slots used {_used_slots} of {_S_orig} "
                    f"(min {S_min}) | underfill penalty {_final_penalty:.2f} "
                    f"| target_min {_target_min_fill*100:.0f}%\n"
                )
        except Exception:
            pass

        # DEBUG, not WARNING: these were escalated during the capacity-gap
        # investigation purely so they'd survive a filtered logging setup.
        # That investigation is closed, and this console is shared, so a
        # routine packing decision must not read as a warning. The same
        # numbers still go to dp_packing_report.txt above.
        if _solved_factor > 1.0:
            logger.debug(
                "  DP feasible with relax × %.2f; using %s of %s slots, underfill penalty %.2f (target_min=%.0f%%)",
                _solved_factor, _used_slots, _S_orig, _final_penalty, _target_min_fill * 100,
            )
        else:
            logger.debug(
                "  DP tight-pack: using %s of %s slots (min=%s), underfill penalty %.2f (target_min=%.0f%%)",
                _used_slots, _S_orig, S_min, _final_penalty, _target_min_fill * 100,
            )

        # Reconstruct the assignment.
        assignment: List[List[Dict[str, Any]]] = [[] for _ in range(S)]
        i = N - 1
        for s in range(S - 1, -1, -1):
            j = split[s][i]
            if j == i:
                assignment[s] = []
                continue
            assignment[s] = list(flat_accounts[j + 1 : i + 1])
            i = j
            if i < 0:
                break

        # The DP's own result is only the starting point -- six passes
        # below mutate it, so the "tight-pack at 1.0" figure reported above
        # describes the DP, NOT what finally renders. A real export showed
        # the DP at 99% and the rendered slot at 110%; the difference is
        # necessarily introduced here. Record each pass's effect on every
        # slot's cost so the responsible one is identifiable from a single
        # run instead of by elimination.
        def _slot_fill_snapshot() -> List[Tuple[float, float]]:
            out: List[Tuple[float, float]] = []
            for _s_i, _slot in enumerate(slots):
                _used = 0.0
                _prev = None
                for _a in assignment[_s_i]:
                    _c = str(_a.get("category", "") or "")
                    if _c and _c != _prev:
                        _used += 1.0
                    _prev = _c
                    _is_chi = bool(_a.get("is_chinese"))
                    _used += self._calculate_content_lines(
                        "", self._rendered_bullet_label(_a, _is_chi),
                        _a.get("commentary", ""), slot_name=_slot["slot_name"],
                        shape=_slot["shape"], is_chinese=_is_chi,
                    )
                if assignment[_s_i]:
                    _used -= self._real_para_gap_pt(True) / (self._planning_std_lh_pt(True) or 1.0)
                out.append((_used, float(_slot["capacity"])))
            return out

        _pass_trace: List[Tuple[str, List[Tuple[float, float]]]] = [
            ("after DP", _slot_fill_snapshot())
        ]
        for _pass_name, _pass_fn in (
            ("rebalance_lopsided_lr_pairs", self._rebalance_lopsided_lr_pairs),
            ("consolidate_tiny_stub_lr_pairs", self._consolidate_tiny_stub_lr_pairs),
            ("rebalance_underfilled_boundaries", self._rebalance_underfilled_boundaries),
            ("rebalance_overflowing_boundaries", self._rebalance_overflowing_boundaries),
            ("maximize_forward_fill", self._maximize_forward_fill),
            ("consolidate_trailing_near_empty_slot", self._consolidate_trailing_near_empty_slot),
        ):
            _pass_fn(assignment, slots, statement_type)
            _pass_trace.append((_pass_name, _slot_fill_snapshot()))

        try:
            with open("dp_packing_report.txt", "a", encoding="utf-8") as _fh:
                _fh.write(f"\n--- {statement_type}: per-pass slot fill (used/capacity) ---\n")
                _prev_snap = None
                for _name, _snap in _pass_trace:
                    _cells = []
                    for _k, (_u, _c) in enumerate(_snap):
                        _pct = (_u / _c * 100) if _c else 0.0
                        _mark = "!" if _u > _c else " "
                        _chg = ""
                        if _prev_snap and abs(_prev_snap[_k][0] - _u) > 0.01:
                            _chg = f"({_prev_snap[_k][0]:.1f}->{_u:.1f})"
                        _cells.append(f"s{_k}:{_pct:5.1f}%{_mark}{_chg}")
                    _fh.write(f"  {_name:<38} " + "  ".join(_cells) + "\n")
                    _prev_snap = _snap
                _fh.write("  ('!' = over capacity. The first pass where a '!' appears "
                          "is the one that caused the overflow.)\n")
        except Exception as _exc:
            logger.debug("Could not write per-pass fill trace: %s", _exc)

        for s_i, slot in enumerate(slots):
            lines = slot_cost(s_i, 0, -1) if not assignment[s_i] else self._compute_slot_used_lines(
                assignment[s_i],
                slot["slot_name"],
                slot_shape=slot["shape"],
                statement_type=statement_type,
            )
            logger.info(
                "  Balanced DP slot %s (%s): %s/%s lines, accts=%s",
                s_i, slot["slot_name"], lines, slot["capacity"],
                [a.get("mapping_key", "?") for a in assignment[s_i]],
            )

        # A slot with no accounts is normally dropped entirely -- e.g. an
        # unused TRAILING slide the DP decided wasn't needed at all. But an
        # L/R pair where ONE side is intentionally empty (e.g.
        # _consolidate_tiny_stub_lr_pairs folding a tiny stub into its
        # sibling) is NOT that case: the slide itself is genuinely in use
        # (its other slot has real content), so the empty slot must still
        # be included here -- otherwise slides_content never gets an entry
        # for it at all, the per-slot render loop's "if not
        # account_data_list: clear the shape" code never even runs (there's
        # nothing in the dict to iterate over), and the shape is left
        # showing whatever raw template placeholder text it shipped with
        # (confirmed via a real export: a literal "Placeholder –
        # placeholder" on the untouched R box). Dropping is still correct
        # for a slide where EVERY slot is empty (truly unused).
        slide_has_content = {
            slot["slide_idx"] for s_i, slot in enumerate(slots) if assignment[s_i]
        }
        rebuilt = [
            (slot["slide_idx"], slot["slot_name"], self._merge_contd_pairs(assignment[s_i]))
            for s_i, slot in enumerate(slots)
            if assignment[s_i] or slot["slide_idx"] in slide_has_content
        ]
        return rebuilt

    @staticmethod
    def _merge_contd_pairs(accounts: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
        """Merge any consecutive run of (part1, cont'd-part2, cont'd-part3, ...)
        fragments that landed in the same slot.  This happens when the DP
        re-balances: a split was created because an earlier slot was almost
        full, but the resulting pieces all fit together in the slot the DP
        actually chose.  Merging removes the spurious (cont'd) label(s) and
        restores the original single account.

        Only ever merged a single PAIR until a real screenshot (IMG_0076)
        showed an orphaned "(续)" bullet sitting right after its own
        already-rendered head -- a 3-way split (a middle fragment that is
        BOTH is_partial [it got re-split by a later rebalance pass] AND
        is_continuation [it continues the fragment before it]) only had its
        first two pieces merged; the third was never considered because the
        old loop looked at exactly one `nxt`, not the whole chain. Confirmed
        via direct reproduction + tracing: a 3-part split landed as
        [is_partial, is_partial+is_continuation, is_continuation] in one
        slot, and the old pairwise merge produced "merged(1+2)" followed by
        untouched "3" as its own still-(续)-labelled bullet."""
        result: List[Dict[str, Any]] = []
        i = 0
        n = len(accounts)
        while i < n:
            acct = accounts[i]
            if not acct.get("is_partial"):
                result.append(acct)
                i += 1
                continue
            base_key = acct.get("mapping_key")
            run = [acct]
            j = i + 1
            while (
                j < n
                and accounts[j].get("is_continuation")
                and accounts[j].get("original_key", accounts[j].get("mapping_key")) == base_key
            ):
                run.append(accounts[j])
                j += 1
            if len(run) > 1:
                combined = run[0].copy()
                combined["commentary"] = " ".join(
                    str(a.get("commentary", "") or "").strip() for a in run
                ).strip()
                combined.pop("is_partial", None)
                combined.pop("part_num", None)
                # A middle fragment re-split by a later rebalance pass can
                # itself be is_continuation=True (it continues a head that
                # sits in an EARLIER slot) as well as is_partial=True (it
                # got split again, with its own tail in THIS run) -- if
                # run[0] is one of those, keep its is_continuation/
                # original_key on the merged result so the "(续)" label
                # renders correctly against the real earlier-slot head;
                # only drop them when run[0] is a genuine, non-continuation
                # first part (the common case).
                if not run[0].get("is_continuation"):
                    combined.pop("is_continuation", None)
                    combined.pop("original_key", None)
                result.append(combined)
            else:
                result.append(acct)
            i = j
        return result

    def _greedy_forward_fill(
        self,
        flat_accounts: List[Dict[str, Any]],
        slots: List[Dict[str, Any]],
        statement_type: Optional[str],
    ) -> List[tuple]:
        """Fallback: fill each slot to capacity greedily. Used only if DP
        can't find a feasible partition (e.g. a single account overflows a
        slot). Always places every account — if an account alone exceeds a
        slot's capacity it is force-placed rather than dropped."""
        def measure(accts, slot):
            return self._compute_slot_used_lines(
                accts, slot["slot_name"], slot_shape=slot["shape"],
                statement_type=statement_type,
            )

        idx = 0
        assignment: List[List[Dict[str, Any]]] = [[] for _ in slots]
        for s_i, slot in enumerate(slots):
            while idx < len(flat_accounts):
                trial = assignment[s_i] + [flat_accounts[idx]]
                if measure(trial, slot) > slot["capacity"] and assignment[s_i]:
                    # Slot already has content and adding this account overflows — move on
                    break
                # Place the account: either the slot is empty (force-place to avoid
                # dropping) or it still fits within capacity.
                assignment[s_i] = trial
                idx += 1

        # If any accounts are still unplaced (more accounts than slots can absorb),
        # append them to the last slot rather than silently dropping them.
        if idx < len(flat_accounts) and slots:
            for remaining in flat_accounts[idx:]:
                assignment[-1].append(remaining)

        return [
            (slot["slide_idx"], slot["slot_name"], self._merge_contd_pairs(assignment[s_i]))
            for s_i, slot in enumerate(slots)
            if assignment[s_i]
        ]

    def _expand_commentary_to_cover_summary(self, slide) -> bool:
        """Remove coSummaryShape from a continuation slide and expand the
        commentary box(es) upward to fill the freed area.

        Returns True if the operation modified the slide. Called only on
        continuation slides (i.e., not the first slide of a BS/IS statement)
        so the AI executive summary stays on the first slide only.
        """
        summary_shape = self.find_shape_by_name(slide.shapes, "coSummaryShape")
        if summary_shape is None:
            return False
        try:
            co_top = int(summary_shape.top)
            co_height = int(summary_shape.height)
        except Exception:
            return False
        co_bottom = co_top + co_height

        for slot_name in ("textMainBullets", "textMainBullets_L", "textMainBullets_R"):
            box = self.find_shape_by_name(slide.shapes, slot_name)
            if box is None:
                continue
            try:
                box_top = int(box.top)
                box_height = int(box.height)
            except Exception:
                continue
            # Only expand boxes located below the summary shape — avoid
            # accidentally covering tables / titles that sit above it.
            if box_top >= co_bottom:
                extension = box_top - co_top
                box.top = co_top
                box.height = box_height + extension

        try:
            sp = summary_shape._element
            sp.getparent().remove(sp)
        except Exception as exc:
            logger.warning("Could not remove coSummaryShape on continuation slide: %s", exc)
            return False
        return True

    def _plan_slot_distribution(
        self,
        structured_data: List[Dict],
        *,
        max_slides: int,
        start_slide: int,
        statement_type: str,
        is_chinese_databook: bool,
    ) -> List[Tuple[int, str, List[Dict]]]:
        """Decide which accounts land in which slot -- the whole plan,
        including presentation-table accounts.

        Extracted so the export and inspect_databook.py's fill diagnostic
        run the SAME planner. They did not: the diagnostic called
        _distribute_content_across_slots on every account, while the export
        pulls table accounts out of that pool first and appends them
        afterwards. On a real deck that made the diagnostic report an IS
        page as 100% full with 4 accounts when the shipped page held one
        account at 33%, and conclude "no genuine packing gaps" about a
        layout it was not describing. A diagnostic that models a different
        algorithm than the one shipping is worse than no diagnostic.

        Callers are expected to have run _prepare_structured_data_for_slides
        already, since the export needs the prepared items for other things
        too.
        """
        tables_enabled = self._presentation_tables_enabled()
        table_style = self._presentation_table_style() if tables_enabled else "table"
        table_items: List[Dict[str, Any]] = []
        normal_items: List[Dict[str, Any]] = []
        last_table_pos: Optional[int] = None
        tagged_normal: List[Tuple[int, Dict[str, Any]]] = []
        for pos, item in enumerate(structured_data):
            table = self._presentation_table_for_account(item) if tables_enabled else None
            if table and table_style == "sublist":
                # Fallback style: the account is NEVER pulled out of the
                # normal packing pool -- the table dict becomes plain text
                # appended to its own commentary, so it inherits the whole
                # existing text pipeline (packing, splitting, cross-column
                # continuation, overflow handling) instead of the dedicated
                # native-table rendering path below.
                item = dict(item)
                is_chinese = bool(item.get("is_chinese"))
                lead_in, post_table_text = self._split_table_commentary(
                    item.get("commentary", ""), is_chinese,
                )
                source_multiplier = 1
                financial_data = item.get("financial_data")
                if hasattr(financial_data, "attrs"):
                    source_multiplier = financial_data.attrs.get("source_multiplier") or 1
                sublist_text = self._sublist_text_for_table(table, is_chinese, source_multiplier)
                parts = [p for p in (lead_in, sublist_text, post_table_text) if p]
                item["commentary"] = "\n".join(parts)
                tagged_normal.append((pos, item))
            elif table:
                item = dict(item)
                lead_in, post_table_text = self._split_table_commentary(
                    item.get("commentary", ""), bool(item.get("is_chinese")),
                )
                item["commentary"] = lead_in
                item["_presentation_table"] = table
                item["_post_table_text"] = post_table_text
                table_items.append(item)
                last_table_pos = pos
            else:
                tagged_normal.append((pos, item))

        # Accounts positioned AFTER the LAST table account in the
        # statement's own reading order (e.g. 投资收益/营业外支出 following
        # 财务费用 on a real Crescent IS) are withheld from the normal
        # packer's input pool and instead flowed into the table slots'
        # own leftover space below -- see _append_table_accounts_to_
        # distribution's trailing_items. Excluding them BEFORE the packer
        # ever runs (rather than trying to reopen its finished output
        # afterwards) mirrors exactly how table_items themselves have
        # always been excluded from this same pool. Accounts BEFORE or
        # BETWEEN table accounts are untouched -- only the true reading-
        # order tail is eligible, so this can't reorder anything the
        # packer would otherwise have placed earlier.
        trailing_items: List[Dict[str, Any]] = []
        if last_table_pos is not None:
            for pos, item in tagged_normal:
                if pos > last_table_pos:
                    trailing_items.append(item)
                else:
                    normal_items.append(item)
        else:
            normal_items = [item for _pos, item in tagged_normal]

        # Distribute content across textbox slots based on capacity
        slot_distribution = self._distribute_content_across_slots(
            normal_items,
            max_slides=max_slides,
            start_slide=start_slide,
            statement_type=statement_type,
        )

        # One shared column-width set for every subtable in this statement,
        # computed before any of them render (see the method's docstring).
        self._precompute_uniform_table_column_widths(table_items, is_chinese_databook)

        if table_items or trailing_items:
            slot_distribution = self._append_table_accounts_to_distribution(
                table_items, slot_distribution, max_slides=max_slides, start_slide=start_slide,
                is_chinese_databook=is_chinese_databook, trailing_items=trailing_items,
            )

        return slot_distribution

    def apply_structured_data_to_slides(self, structured_data: List[Dict], start_slide: int,
                                       project_name: str, statement_type: str, is_chinese_databook: bool = False,
                                       pre_generated_summary: Optional[str] = None):
        """Apply structured data directly to slides (slides 1-4 for BS, 5-8 for IS).

        If ``pre_generated_summary`` is provided, it's used directly for the
        first slide's coSummaryShape — no AI call from inside PPTX export.
        """
        if not self.presentation:
            self.load_template()

        # Remembered so save() can decide, once for the whole deck, whether
        # to strip the template's static English "Commentary" label bands
        # (see _localize_label_bands) -- both statement passes
        # (BS then IS) report the same databook language, so whichever
        # runs last simply re-confirms the same value.
        self._is_chinese_mode = is_chinese_databook

        stage_started_at = time.perf_counter()
        logger.info("Applying %s accounts to slides starting at %s", len(structured_data), start_slide)

        # Normalize commentary and store originals for fill optimization
        structured_data = self._prepare_structured_data_for_slides(structured_data)

        # Continuation slides (every slide of this statement after the first)
        # lose their coSummaryShape and gain that area as extra commentary
        # space. The executive summary stays only on the first slide of
        # each statement, which cuts AI summary calls from up to 8 to 2.
        max_slides = int(self.pptx_settings.get("max_commentary_slides_per_statement", 4) or 4)
        first_slide_idx = start_slide - 1
        for offset in range(1, max_slides):
            cont_idx = first_slide_idx + offset
            if cont_idx >= len(self.presentation.slides):
                break
            self._expand_commentary_to_cover_summary(self.presentation.slides[cont_idx])

        # Presentation-table accounts are pulled out of the packing pool
        # entirely rather than fed to the DP/greedy packer with an inflated
        # cost. A first attempt padded commentary cost to make a table
        # account "claim" its whole slot -- verified against a real render
        # that this does NOT work: the packer still placed a second account
        # in the same slot, then SPLIT it across two slides when the
        # combined cost overflowed (the overflow-split logic works purely
        # off the commentary text, with no concept of "this text is padding
        # for a table, do not cut it"). Two tables ended up drawn at the
        # exact same on-slide position. Removing them before packing runs
        # sidesteps that entirely -- the packer never knows they exist, so
        # its own sharing/splitting behaviour can't touch them.
        slot_distribution = self._plan_slot_distribution(
            structured_data, max_slides=max_slides, start_slide=start_slide,
            statement_type=statement_type, is_chinese_databook=is_chinese_databook,
        )

        # Group slot distribution by slide for easier processing
        slides_content = {}  # {slide_idx: {'single': [...], 'L': [...], 'R': [...]}}
        for slot_slide_idx, slot_name, account_list in slot_distribution:
            if slot_slide_idx not in slides_content:
                slides_content[slot_slide_idx] = {}
            slides_content[slot_slide_idx][slot_name] = account_list
        
        # Ensure we have enough slides
        if slides_content:
            max_slide_idx = max(slides_content.keys())
            needed_slides = start_slide + max_slide_idx
            current_slide_count = len(self.presentation.slides)
            
            if needed_slides > current_slide_count:
                # Add slides if needed
                if current_slide_count > 0:
                    slide_layout = self.presentation.slides[0].slide_layout
                    for _ in range(needed_slides - current_slide_count):
                        self.presentation.slides.add_slide(slide_layout)
        
        # Track which slides are used
        used_slide_indices = set()
        summary_jobs: List[Dict[str, Any]] = []
        
        # Apply content to slides
        for slide_idx in sorted(slides_content.keys()):
            actual_slide_idx = start_slide - 1 + slide_idx  # Convert to 0-based
            if actual_slide_idx >= len(self.presentation.slides):
                logger.warning("Slide index %s exceeds available slides", actual_slide_idx + 1)
                continue
            
            slide = self.presentation.slides[actual_slide_idx]
            slot_contents = slides_content[slide_idx]  # {'single': [...], 'L': [...], 'R': [...]}
            # A slide only counts as "used" if at least one of its slots has
            # real content -- a slide whose slots are ALL empty (e.g. after
            # _consolidate_trailing_near_empty_slot folds its one leftover
            # sliver into the previous slide) should be eligible for the
            # unused-slide removal below, not kept around as a blank page.
            if any(slot_contents.values()):
                used_slide_indices.add(actual_slide_idx)

            # Note: Financial tables are filled by embed_financial_tables()

            # Collect all accounts on this slide for summary generation
            all_slide_accounts = []
            for slot_name, accounts in slot_contents.items():
                all_slide_accounts.extend(accounts)

            used_slot_shape_ids = set()
            
            # Fill each slot (single, L, or R) on this slide
            for slot_name, account_data_list in slot_contents.items():
                if not account_data_list:
                    # An INTENTIONALLY empty slot (e.g.
                    # _consolidate_tiny_stub_lr_pairs folded its one tiny
                    # fragment into the earlier slot rather than leave an
                    # orphaned near-empty column) must still have its
                    # shape's text actively cleared -- otherwise the shape
                    # keeps whatever raw template placeholder sample text
                    # it shipped with (e.g. a visible "Placeholder –
                    # placeholder"), which reads as broken output, not as
                    # the clean blank box an intentionally-unused slot
                    # should be.
                    empty_shape = self._resolve_commentary_slot_shape(
                        slide, slot_name, used_shape_ids=used_slot_shape_ids,
                    )
                    if empty_shape is not None and empty_shape.has_text_frame:
                        empty_shape.text_frame.clear()
                        used_slot_shape_ids.add(id(empty_shape))
                    continue

                # Find the appropriate shape based on slot_name
                bullets_shape = self._resolve_commentary_slot_shape(
                    slide,
                    slot_name,
                    used_shape_ids=used_slot_shape_ids,
                )
                if not bullets_shape and slot_name == "single":
                    bullets_shape = self._resolve_commentary_slot_shape(
                        slide,
                        "L",
                        used_shape_ids=used_slot_shape_ids,
                    ) or self._resolve_commentary_slot_shape(
                        slide,
                        "R",
                        used_shape_ids=used_slot_shape_ids,
                    )
                if not bullets_shape:
                    bullets_shape = self._add_commentary_slot_shape(slide, slot_name)
                
                if not bullets_shape.has_text_frame:
                    logger.warning("Slide %s: Shape for slot '%s' has no text frame", actual_slide_idx + 1, slot_name)
                    continue
                used_slot_shape_ids.add(id(bullets_shape))

                # A slot containing at least one table account (possibly
                # followed by plain trailing accounts flowed into its
                # leftover space -- see _append_table_accounts_to_
                # distribution's trailing_items) renders via its own
                # dedicated-shape-per-account path instead of the
                # shared-text-frame one below, since that one assumes all of
                # a slot's content lives in one text frame with nothing else
                # interleaved -- not true once a table (a separate shape)
                # sits between one account's lead-in and the next account's
                # own lead-in. A slot with ONLY plain accounts (no table at
                # all) still goes through the normal shared-text-frame path
                # unchanged -- that's the common case and needs none of this.
                if account_data_list and any(a.get("_presentation_table") for a in account_data_list):
                    self._render_table_accounts_stack(
                        slide, bullets_shape, account_data_list, is_chinese_databook,
                        statement_type=statement_type,
                    )
                    continue

                # Fill this slot
                tf = bullets_shape.text_frame
                tf.clear()
                tf.word_wrap = True
                # TextFrame.clear() always leaves exactly one (now-empty)
                # paragraph behind; every real line below is added via
                # add_paragraph() rather than reusing it, so it survives as a
                # permanent blank leading line PowerPoint actually renders
                # (~1 line-height of capacity silently lost on every slot in
                # the deck) even though it contributes zero to this file's
                # own content-line accounting (empty paragraphs are filtered
                # out by both the DP/packing math and inspect_pptx.py's
                # measurement) -- a real vs. believed capacity mismatch that
                # was part of why slots read as "stuck" under their true
                # capacity. Removed once real content exists, right before
                # moving to the next slot.
                _leading_empty_p = tf.paragraphs[0]._p

                # Most slots fit within strict 1.0x capacity and keep the
                # exact 9pt/10pt size (noAutofit) -- no drift, no surprise
                # shrink. The DP occasionally packs a slot beyond that
                # (shape_height_utilization / the looser relax tiers) on
                # the assumption PowerPoint's autofit absorbs the excess;
                # noAutofit alone would silently clip that content instead,
                # so give exactly those slots a bounded normAutofit shrink
                # (never below _BOUNDED_AUTOFIT_MIN_SCALE) so the overflow
                # the DP already decided was "worth it" is actually visible.
                #
                # But ONLY past the tail tolerance. Writing a normAutofit does
                # not shrink text by the fontScale computed here: PowerPoint
                # re-runs its own shrink with its own font metrics and steps
                # down in coarse increments, so it lands wherever it lands.
                # Measured against a real deck (VBA BoundHeight, the ground
                # truth): a slot overflowing by 0.3 lines -- fontScale ~99%
                # written here -- rendered at ~87%, 27 lines at 10.31pt
                # instead of 10.8pt, leaving 52.8pt of visible blank under a
                # box the packer had filled to 101%. Every other box in that
                # deck measured >= 10.8pt/line; this was the only shrunk one
                # and the only one over capacity. Below the tolerance, letting
                # the tail protrude is both what the project team asked for
                # ("1-2 lines sticking out is fine") and strictly better
                # looking than a whole column shrunk to leave a gap.
                slot_is_chinese = any(self._account_is_chinese(a) for a in account_data_list)
                slot_capacity = self._calculate_max_lines_for_textbox(
                    bullets_shape, is_chinese=slot_is_chinese, slot_name=slot_name,
                    statement_type=statement_type,
                )
                slot_used = self._compute_slot_used_lines(
                    account_data_list, slot_name, slot_shape=bullets_shape,
                    statement_type=statement_type,
                )
                _tail_tol = self._tail_overflow_tolerance_units(statement_type)
                if slot_capacity > 0 and (slot_used - slot_capacity) > _tail_tol:
                    self._apply_bounded_autofit(tf, slot_capacity / slot_used)
                else:
                    self._force_no_autofit(tf)  # keep text at 9pt/10pt, never shrink
                from pptx.enum.text import MSO_VERTICAL_ANCHOR
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
                
                # Fixed deck-wide font size — 9pt Arial, no exceptions,
                # no per-slot / per-language variation. This matches the
                # hardcoded return from get_font_size_for_text() and
                # guarantees identical typography on every slide.
                slot_font_size = 9
                logger.info(
                    "Slide %s, slot '%s': Filling with %s accounts at %spt",
                    actual_slide_idx + 1, slot_name, len(account_data_list), slot_font_size,
                )

                # Fill with accounts, grouped by category
                # Show category header only once per category group
                current_category = None
                for account_idx, account_data in enumerate(account_data_list):
                    category = account_data.get('category', '')
                    mapping_key = account_data.get('mapping_key', account_data.get('account_name', ''))
                    # display_name is whatever text sat in the source Financials
                    # sheet's row (often English even for a Chinese-language
                    # report, e.g. Kunshan's own sheets are English-labelled) --
                    # use the pre-resolved Chinese alias for a Chinese report
                    # instead of leaving that English label in an otherwise
                    # fully-translated bullet.
                    display_name = (
                        account_data.get('display_name_zh') or account_data.get('display_name', mapping_key)
                        if is_chinese_databook
                        else account_data.get('display_name', mapping_key)
                    )
                    commentary = account_data.get('commentary', '')
                    clause_reviews = account_data.get('clause_reviews', [])
                    is_chinese = account_data.get('is_chinese', False)
                    is_continuation = account_data.get('is_continuation', False)
                    
                    # Skip category header if this is a continuation of a split account
                    # Show category header only when category changes
                    if category and category != current_category and not is_continuation:
                        # Add category header - use Chinese if databook is Chinese
                        p_category = tf.add_paragraph()
                        p_category.level = 0
                        try:
                            p_category.left_indent = Inches(0.21)
                            p_category.first_line_indent = Inches(-0.19)
                            p_category.space_before = Pt(3) if current_category else Pt(0)
                            p_category.space_after = Pt(0)
                            p_category.line_spacing = 1.0
                            self._apply_east_asian_line_breaking(p_category)
                        except:
                            pass
                        
                        run_category = p_category.add_run()
                        # Use Chinese category name if databook is Chinese
                        category_text = translate_category_to_chinese(category) if is_chinese_databook else category
                        
                        run_category.text = category_text
                        run_category.font.size = Pt(slot_font_size)
                        run_category.font.name = 'Arial'
                        self._set_east_asian_typeface(run_category)
                        self._declare_run_language(run_category)
                        run_category.font.bold = False
                        try:
                            from pptx.dml.color import RGBColor
                            run_category.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
                        except:
                            pass
                        
                        current_category = category
                    
                    # Fill commentary with key formatting
                    # For continuation accounts, show "(cont'd)" or "(续)" after key name
                    if is_continuation:
                        # More prominent continuation marker
                        if is_chinese_databook:
                            display_name_with_cont = f"{display_name} (续)"
                        else:
                            display_name_with_cont = f"{display_name} (cont'd)"
                        
                        # Log continuation for debugging
                        logger.info("Displaying continuation: %s", display_name_with_cont)
                    else:
                        display_name_with_cont = display_name
                    
                    self._fill_text_main_bullets_with_category_and_key(
                        tf, None, display_name_with_cont, commentary, is_chinese,
                        is_chinese_databook=is_chinese_databook, needs_continuation=False,
                        font_size_pt=slot_font_size,
                        clause_reviews=clause_reviews,
                    )
                    # Table-bearing accounts never reach this loop -- they're
                    # intercepted by the table-only-slot dispatch above and
                    # rendered via _render_table_accounts_stack instead.

                if _leading_empty_p.getparent() is not None and not (_leading_empty_p.text or "").strip():
                    _leading_empty_p.getparent().remove(_leading_empty_p)

            page_commentary, page_summary_source = self._build_page_summary_source(all_slide_accounts)

            # Collect coSummaryShape jobs and fill after summaries are generated.
            summary_shape = self.find_shape_by_name(slide.shapes, "coSummaryShape")
            if summary_shape and summary_shape.has_text_frame:
                summary_shape.text_frame.clear()
                self._force_no_autofit(summary_shape.text_frame)
                if page_summary_source:
                    summary_jobs.append({
                        "slide_idx": actual_slide_idx,
                        "summary_shape": summary_shape,
                        "page_commentary": page_commentary,
                        "page_summary_source": page_summary_source,
                        "is_chinese": is_chinese_databook,
                        "font_is_chinese": all_slide_accounts[0].get('is_chinese', False) if all_slide_accounts else False,
                    })
                else:
                    # Every way this band ends up blank currently looks
                    # identical from the outside: shape missing, no source
                    # text, generator returned nothing, a blank pre-generated
                    # summary. Each one silently does nothing, so "exe sum
                    # 一直未有東西" cost several rounds of guessing. Name the
                    # cause at WARNING level -- inspect_databook.py's export
                    # log analysis surfaces warnings, so the next real run
                    # says which one it was instead of just showing a blank.
                    logger.warning(
                        "Slide %s: coSummaryShape found but NO summary source text "
                        "(%s accounts on this slide, %s with commentary) — band left blank",
                        actual_slide_idx + 1, len(all_slide_accounts),
                        sum(1 for a in all_slide_accounts if str(a.get("commentary", "") or "").strip()),
                    )
            elif actual_slide_idx == start_slide - 1:
                logger.warning(
                    "Slide %s: coSummaryShape NOT FOUND on the first slide of %s — band "
                    "left blank. The lookup is an exact name match; check the template.",
                    actual_slide_idx + 1, statement_type,
                )
            else:
                # Continuation slides have their band removed on purpose by
                # _expand_commentary_to_cover_summary, so its absence here is
                # the designed layout, not a fault. Warning on it put three
                # false lines into the export log analysis of a run whose
                # summaries had in fact all worked.
                logger.info(
                    "Slide %s: continuation slide, no summary band by design",
                    actual_slide_idx + 1,
                )

            logger.info("Filled slide %s with %s accounts across %s slots", actual_slide_idx + 1, len(all_slide_accounts), len(slot_contents))

        # If a pre-generated summary was supplied (computed during the AI
        # commentary phase), use it directly and skip the in-PPTX AI call.
        # The pre-generated summary applies to the FIRST slide of this
        # statement only — continuation slides have coSummaryShape removed.
        pre_summary_text = str(pre_generated_summary or "").strip() if pre_generated_summary else ""
        if pre_summary_text and summary_jobs:
            summary_jobs.sort(key=lambda j: j["slide_idx"])
            first_job = summary_jobs[0]
            summary_results = {first_job["slide_idx"]: pre_summary_text}
            jobs_to_fill = [first_job]
        elif summary_jobs:
            # No pre-generated summary supplied. Calling LLM during PPTX export
            # is slow when the API is flaky (3 retries × 30s × N slides) and
            # the user reported 10+ min hangs. Skip the in-export AI call —
            # leave coSummaryShape blank rather than wait. The user can re-run
            # the AI generation step to refresh summaries when the API is
            # responsive.
            # Skipping the LLM call is right; leaving the box EMPTY was not.
            # The front page of a deliverable shipped with a blank summary
            # band, and because the template's own placeholder in that shape
            # is set to 2.17pt, what showed instead was an illegible sliver
            # of "Entity overview – Project [PROJECT]" -- reported as "exe
            # sum 一直未有東西".
            #
            # _generate_page_summary already exists for exactly this case:
            # it takes the opening sentence of each account paragraph, so it
            # spans the whole page, costs no tokens and no wall time. Use it
            # rather than shipping nothing.
            logger.info(
                "No pre-generated summary supplied; using the deterministic "
                "page summary instead of an in-export LLM call (which would "
                "add 1-3 min per slide)."
            )
            summary_jobs.sort(key=lambda j: j["slide_idx"])
            summary_results = {}
            jobs_to_fill = []
            for job in summary_jobs:
                _src = job.get("page_summary_source") or job.get("page_commentary") or ""
                text = self._generate_page_summary(_src, bool(job.get("is_chinese")))
                if not text:
                    logger.warning(
                        "Slide %s: _generate_page_summary returned nothing from %s chars "
                        "of source — summary band left blank",
                        job["slide_idx"] + 1, len(_src),
                    )
                if text:
                    summary_results[job["slide_idx"]] = text
                    jobs_to_fill.append(job)
        else:
            # No job was ever collected. Either no slide in this statement has
            # a coSummaryShape, or none had source text -- both already warned
            # per-slide above, so this only records that the fill stage ran
            # with nothing to do rather than silently falling through.
            logger.warning(
                "%s: no coSummaryShape jobs collected — every summary band for "
                "this statement stays blank", statement_type,
            )
            summary_results = {}
            jobs_to_fill = []
        for job in jobs_to_fill:
            final_summary = str(summary_results.get(job["slide_idx"]) or "").strip()
            if not final_summary:
                logger.warning("Slide %s: summary resolved to empty text — band left blank",
                               job["slide_idx"] + 1)
                continue
            logger.info("Slide %s: wrote %s-char executive summary into coSummaryShape",
                        job["slide_idx"] + 1, len(final_summary))
            summary_shape = job["summary_shape"]
            p = summary_shape.text_frame.paragraphs[0] if summary_shape.text_frame.paragraphs else summary_shape.text_frame.add_paragraph()
            p.text = final_summary
            self._apply_east_asian_line_breaking(p)
            for run in p.runs:
                run.font.size = get_font_size_for_text(final_summary, force_chinese_mode=job["font_is_chinese"])
                run.font.name = get_font_name_for_text(final_summary)
                # This fill site sets font.name from a helper rather than the
                # literal 'Arial' the other 17 sites use, so it was missed
                # when the CJK typeface and language were added everywhere
                # else. The summary band is the first Chinese text on the
                # page; it needs them at least as much.
                self._set_east_asian_typeface(run)
                self._declare_run_language(run)
        
        # Record the FIRST used commentary slide of this statement as a slide
        # OBJECT (not an index). embed_financial_tables targets this so the BS/IS
        # table lands on the correct page even after slides are added/removed.
        if used_slide_indices:
            first_used_idx = min(used_slide_indices)
            if 0 <= first_used_idx < len(self.presentation.slides):
                if not hasattr(self, "_statement_table_slides"):
                    self._statement_table_slides = {}
                self._statement_table_slides[statement_type] = self.presentation.slides[first_used_idx]

        # Note: Unused slides will be removed at the end, after all content and tables are embedded
        # Store unused slides for later removal
        statement_slide_range = list(range(start_slide - 1, min(start_slide + 3, len(self.presentation.slides))))
        unused_slides = [idx for idx in statement_slide_range if idx not in used_slide_indices]
        if unused_slides:
            # Store for later removal - don't remove now
            if not hasattr(self, '_unused_slides_to_remove'):
                self._unused_slides_to_remove = []
            self._unused_slides_to_remove.extend(unused_slides)
            logger.info("Marked %s unused slides for %s for later removal: %s", len(unused_slides), statement_type, [idx + 1 for idx in unused_slides])
        logger.info(
            "PPTX stage apply_structured_data_to_slides[%s] took %.2fs across %s populated slides",
            statement_type,
            time.perf_counter() - stage_started_at,
            len(slides_content),
        )
    
    def _remove_slides(self, slide_indices):
        """Remove slides by indices (from backup method)"""
        # Sort in reverse order to maintain indices while removing
        for slide_idx in sorted(slide_indices, reverse=True):
            if slide_idx < len(self.presentation.slides):
                try:
                    # Use XML-based removal (from backup method)
                    xml_slides = self.presentation.slides._sldIdLst
                    slides = list(xml_slides)
                    
                    if slide_idx < len(slides):
                        # Get the slide element to remove
                        slide_element = slides[slide_idx]
                        # Remove relationship
                        rId = slide_element.rId
                        self.presentation.part.drop_rel(rId)
                        # Remove from XML
                        xml_slides.remove(slide_element)
                        logger.info("Removed slide %s", slide_idx + 1)
                    else:
                        logger.warning("Slide index %s out of range (only %s slides)", slide_idx, len(slides))
                except Exception as e:
                    logger.warning("Could not remove slide %s: %s", slide_idx + 1, e)
                    logger.debug(traceback.format_exc())
    
    def _set_cell_border(self, cell, border_position='top', color_rgb=None, width=Pt(1)):
        """Set cell border"""
        from pptx.oxml.xmlchemy import OxmlElement
        
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        
        # Map position to tag name
        tag_map = {'top': 'lnT', 'bottom': 'lnB', 'left': 'lnL', 'right': 'lnR'}
        tag_name = tag_map.get(border_position)
        if not tag_name:
            return
            
        # Check if line element exists
        ln = tcPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag_name}")
        if ln is None:
            ln = OxmlElement(f"a:{tag_name}")
            tcPr.append(ln)
            
        # Set properties
        ln.set('w', str(int(width)))
        ln.set('cap', 'flat')
        ln.set('cmpd', 'sng')
        ln.set('algn', 'ctr')

        # Calling this twice on the same side (e.g. a full-grid pass, then a
        # heavier total-row override) previously left BOTH the old and new
        # <a:solidFill>/<a:prstDash>/<a:round>/<a:headEnd>/<a:tailEnd>
        # children on `ln` -- append() never replaces, so the element ended
        # up with duplicates and PowerPoint's rendering of that is
        # undefined (in practice, whichever child renderers pick up first).
        # Clear any existing children before appending the new ones so a
        # second call genuinely overrides the first, not just adds to it.
        for child in list(ln):
            ln.remove(child)

        # Set color
        if color_rgb:
            solidFill = OxmlElement('a:solidFill')
            srgbClr = OxmlElement('a:srgbClr')
            # Convert RGBColor or tuple to hex string
            hex_color = "000000"
            if isinstance(color_rgb, str):
                hex_color = color_rgb.replace('#', '')
            elif isinstance(color_rgb, tuple) and len(color_rgb) == 3:
                hex_color = f"{color_rgb[0]:02x}{color_rgb[1]:02x}{color_rgb[2]:02x}"
            # If it's an RGBColor object, user should pass str or tuple for this low-level func
                
            srgbClr.set('val', hex_color)
            solidFill.append(srgbClr)
            ln.append(solidFill)
            
            prstDash = OxmlElement('a:prstDash')
            prstDash.set('val', 'solid')
            ln.append(prstDash)
            
            round_ = OxmlElement('a:round')
            ln.append(round_)
            
            headEnd = OxmlElement('a:headEnd')
            headEnd.set('type', 'none')
            headEnd.set('w', 'med')
            headEnd.set('len', 'med')
            ln.append(headEnd)
            
            tailEnd = OxmlElement('a:tailEnd')
            tailEnd.set('type', 'none')
            tailEnd.set('w', 'med')
            tailEnd.set('len', 'med')
            ln.append(tailEnd)

    def _fill_table_placeholder(
        self, shape, df, table_name: str = None, currency_unit: str = None, bounds: Dict[str, int] = None,
        mappings: Optional[Dict[str, Any]] = None, is_chinese_mode: bool = False,
    ):
        """Fill table placeholder with DataFrame data, preserving original formatting
        Args:
            shape: Table shape or placeholder
            df: DataFrame with data
            table_name: Name of the table (e.g., "示意性调整后资产负债表 - xxxx")
            currency_unit: Currency unit (e.g., "人民币千元" or "CNY'000") to replace "Description"
        """
        try:
            # Debug: Log DataFrame content
            logger.info("Filling table with DF shape: %s", df.shape)
            if not df.empty:
                logger.info("First row data: %s", df.iloc[0].to_dict())
                # Check if any data is non-zero
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    non_zero_count = (df[numeric_cols] != 0).sum().sum()
                    logger.info("Non-zero values in DF: %s", non_zero_count)
            
            # Find parent slide
            slide = None
            for s in self.presentation.slides:
                for shp in s.shapes:
                    if shp == shape:
                        slide = s
                        break
                if slide:
                    break
            
            if bounds is None:
                bounds = {
                    "left": shape.left,
                    "top": shape.top,
                    "width": shape.width,
                    "height": shape.height,
                }

            # Adjust position and size from resolved layout bounds
            try:
                shape.left = bounds["left"]
                shape.top = bounds["top"]
                shape.width = bounds["width"]
                shape.height = bounds["height"]
            except Exception as e:
                logger.warning("Could not adjust table position/width: %s", e)

            # Check if shape is a TablePlaceholder (textbox placeholder)
            from pptx.shapes.placeholder import TablePlaceholder
            
            table = None
            # Check if it's a TablePlaceholder (either a real pptx placeholder
            # of that type, OR a plain textbox/autoshape named "Table
            # Placeholder" used as a template bounds marker). isinstance()
            # returning False is NOT an exception -- the previous try/except
            # here never actually reached the name-based fallback, since
            # isinstance() simply returns False for a non-match rather than
            # raising, so a plain-shape placeholder (e.g. added directly via
            # python-pptx, not through the placeholder XML machinery) was
            # silently treated as "not a table placeholder" and fell through
            # to the broken "insert as plain text" error path below.
            is_table_placeholder = isinstance(shape, TablePlaceholder) or (
                'Table' in getattr(shape, 'name', '') and 'Placeholder' in getattr(shape, 'name', '')
            )
            
            if is_table_placeholder:
                # It's a table placeholder - insert a table into it
                logger.info("Found TablePlaceholder (%s), inserting table with %s rows and %s columns", shape.name if hasattr(shape, 'name') else 'unnamed', len(df), len(df.columns))
                try:
                    left = bounds["left"]
                    top = bounds["top"]
                    width = bounds["width"]
                    height = bounds["height"]
                    
                    # Find the slide containing this shape (already found above)
                    if slide:
                        # Remove the placeholder shape
                        sp = shape._element
                        slide.shapes._spTree.remove(sp)
                        
                        # Add new table at the same position
                        # Need: 1 row for title (if table_name), 1 for header, N for data
                        total_rows = len(df) + 2 if table_name else len(df) + 1
                        table_shape = slide.shapes.add_table(
                            rows=total_rows,
                            cols=len(df.columns),
                            left=left,
                            top=top,
                            width=width,
                            height=height
                        )
                        table = table_shape.table
                        logger.info("Inserted new table: %s rows, %s columns", len(table.rows), len(table.columns))
                except Exception as e:
                    logger.error("Could not insert table into placeholder: %s", e)
                    logger.debug(traceback.format_exc())
            elif hasattr(shape, 'table'):
                # Try to access existing table
                try:
                    table = shape.table
                    logger.info("Found existing table with %s rows and %s columns", len(table.rows), len(table.columns))
                except ValueError:
                    # Shape doesn't contain a table
                    logger.warning("Shape has table attribute but doesn't contain a table")
                    table = None
            
            if table:
                # Every cell's fill/font/border is set explicitly below, but a
                # freshly-created/template table shape defaults to
                # first_row=True, horz_banding=True -- PowerPoint's built-in
                # table-style theme then auto-applies its own "first row" and
                # alternating row-banding colors on top of (or in place of,
                # for rows this code doesn't explicitly touch) the direct
                # formatting, which is exactly the kind of stray theme
                # colour/stripe the plain reference table doesn't have.
                # Disable both so only the explicit per-cell styling renders.
                try:
                    table.first_row = False
                    table.horz_banding = False
                except Exception:
                    pass

                # Colors -- matches the company-format reference (real KPMG
                # deliverable page): navy #00338D fills the title band and
                # accents total-row borders; only the LAST column is
                # highlighted light blue (every row, header+data+totals);
                # every other cell stays unfilled (white bg, dark text) --
                # total rows are called out with border weight only, no fill.
                DARK_BLUE = RGBColor(0x00, 0x33, 0x8D)
                WHITE = RGBColor(255, 255, 255)
                BLACK = RGBColor(0, 0, 0)
                # Reference format (IMG_0035) calls out only the four
                # statement-level grand totals with a solid grey fill.
                GREY_TOTAL_FILL = RGBColor(0xD9, 0xD9, 0xD9)

                self._fit_table_columns(table, df)

                # Smaller/tighter than before across the board -- reference
                # format (IMG_0035) reads noticeably more compact than this
                # table previously rendered at.
                #
                # Category header rows (_insert_category_header_rows --
                # "Revenue"/"Expenses"/"Current assets"/etc.) don't count
                # toward this tier threshold: they're a single short label
                # with no figures, not a real data row, and letting them
                # count at full weight was dragging the WHOLE table (every
                # real data row too) down to a smaller font/row-height tier
                # than the actual data alone would need -- e.g. the IS
                # table's 16 real rows sit comfortably under the 20-row
                # threshold on their own, but 16 rows + 6 "Revenue"/
                # "Expenses" header rows = 22 pushed it into the tighter
                # tier for no real reason.
                _tier_keywords = list(
                    {'total', '合计', '总计', '小计', 'subtotal', 'sub-total', 'sub total'}
                    | set(SUMMARY_ACCOUNT_SKIP_KEYWORDS)
                )
                _category_header_row_count = 0
                for _, _tier_row in df.iterrows():
                    _tier_label = str(_tier_row.iloc[0]).strip()
                    if not _tier_label or any(kw in _tier_label.lower() for kw in _tier_keywords):
                        continue
                    if all(pd.isna(v) or (isinstance(v, str) and not v.strip()) for v in _tier_row.iloc[1:]):
                        _category_header_row_count += 1
                total_visible_rows = (len(df) - _category_header_row_count) + 1 + (1 if table_name else 0)
                if total_visible_rows >= 26:
                    data_font_size = Pt(6)
                    data_row_height = Inches(0.13)
                elif total_visible_rows >= 20:
                    data_font_size = Pt(6.5)
                    data_row_height = Inches(0.15)
                else:
                    data_font_size = Pt(7)
                    data_row_height = Inches(0.17)

                # Title band and column-header row scale WITH the data tier
                # instead of a fixed Pt(8)/Pt(7)/0.25" -- on a dense table
                # (total_visible_rows >= 26, data shrunk to 6pt/0.13") a
                # hardcoded 0.25"/Pt(7) header reads nearly 2x taller and a
                # full point bigger than every row beneath it, exactly the
                # "header row太大" mismatch a user screenshot flagged.
                header_font_size = Pt(data_font_size.pt + 1)
                header_row_height = Inches(data_row_height.inches + 0.03)
                title_font_size = Pt(data_font_size.pt + 2)
                title_row_height = Inches(data_row_height.inches + 0.05)

                # Hard clamp: the three fixed tiers above are picked from
                # ROW COUNT alone and don't know this table's actual
                # available height (varies by template/slide) -- a table
                # with unusually many rows for its placeholder could still
                # render taller than its bounds and spill onto the slide.
                # Scale every row height (and, floored at a legibility
                # minimum, font size) down uniformly so the whole table's
                # summed row heights never exceed bounds["height"].
                try:
                    _available_h_in = max(0.1, (bounds.get("height", 0) or 0) / 914400)
                    _needed_h_in = (
                        (title_row_height.inches if table_name else 0.0)
                        + header_row_height.inches
                        + len(df) * data_row_height.inches
                    )
                    if _needed_h_in > _available_h_in > 0:
                        _scale = _available_h_in / _needed_h_in
                        data_row_height = Inches(max(0.08, data_row_height.inches * _scale))
                        header_row_height = Inches(max(0.09, header_row_height.inches * _scale))
                        title_row_height = Inches(max(0.10, title_row_height.inches * _scale))
                        data_font_size = Pt(max(4.5, data_font_size.pt * _scale))
                        header_font_size = Pt(max(5.0, header_font_size.pt * _scale))
                        title_font_size = Pt(max(5.5, title_font_size.pt * _scale))
                except Exception:
                    pass

                # Add table name as first row if provided
                if table_name:
                    # Insert a new row at the top for table name
                    try:
                        # Ensure table has at least one row
                        if len(table.rows) == 0:
                            table.rows.add_row()
                            
                        name_row = table.rows[0]  # Use first row for name
                        name_row.height = title_row_height
                        
                        # Merge all cells in first row for table name
                        if len(table.columns) > 1:
                            name_row.cells[0].merge(name_row.cells[len(table.columns) - 1])
                        name_cell = name_row.cells[0]
                        name_cell.text = table_name
                        # Company-format reference (real KPMG deliverable page):
                        # the statement title sits in a navy-filled band spanning
                        # the full table width, white bold text, left-aligned --
                        # the date/header row directly beneath it is the one
                        # left unfilled (white bg, dark text), not the title.
                        if name_cell.text_frame.paragraphs:
                            p = name_cell.text_frame.paragraphs[0]
                            p.alignment = PP_ALIGN.LEFT
                            if p.runs:
                                run = p.runs[0]
                            else:
                                run = p.add_run()
                            run.font.name = 'Arial'
                            self._set_east_asian_typeface(run)
                            self._declare_run_language(run)
                            run.font.size = title_font_size
                            run.font.bold = True
                            run.font.color.rgb = WHITE

                            name_cell.fill.solid()
                            name_cell.fill.fore_color.rgb = DARK_BLUE

                        # Shift data down - we'll use rows starting from index 1
                        data_start_row = 1
                    except:
                        data_start_row = 0
                else:
                    data_start_row = 0
                
                # Fill header row with formatting
                max_cols = min(len(df.columns), len(table.columns))
                header_row_idx = data_start_row
                
                # Ensure header row exists
                if len(table.rows) <= header_row_idx:
                    table.rows.add_row()
                
                # Set header row height slightly taller for readability
                try:
                    table.rows[header_row_idx].height = header_row_height
                except:
                    pass
                    
                # TEMPORARILY REVERTED 2026-08-04 (was HEADER_ROW_BLUE fill
                # across every column, white text, white L/R borders) while
                # isolating a real blank-page bug: two real Chinese exports
                # in a row rendered COMPLETELY BLANK in real PowerPoint on
                # exactly the two slides this function's table lands on
                # (BS/IS overview "financials" pages), with abnormally slow
                # open, while python-pptx itself read every shape/cell back
                # intact from the same files. The Commentary-band deletion
                # was the first suspect and has since been replaced with an
                # in-place text swap (see _localize_label_bands) that keeps
                # every shape's structure untouched -- but the blank pages
                # persisted even after that fix, which rules the band out
                # and leaves this function's per-cell colour/border changes
                # (the only other thing ebf2179 touched on these same two
                # pages) as the remaining suspect. Reverted to the
                # known-safe pre-ebf2179 styling (white bg, black text,
                # only the last column tinted) as a controlled test -- if
                # blank pages persist even with this reverted too, the
                # cause is elsewhere entirely (not this session's colour
                # work), and this revert should itself be re-reverted once
                # that's confirmed.
                #
                # A touch more saturated than a bare tint (#DCE6F1) so the
                # highlighted column is unmistakable at a glance, not just
                # visible on close inspection -- a user photo of a fresh
                # export still read as "hasn't taken effect" at the lighter
                # shade.
                LIGHT_BLUE_HIGHLIGHT = RGBColor(0xBD, 0xD7, 0xEE)
                # A lighter blue than the navy title band directly above it,
                # matching the reference report's two-tone banner. Merging
                # them into one navy (tried first) removed the distinction the
                # reference draws. What made the header read as a SEPARATE
                # element before was not the shade but the black cell borders
                # drawn across it -- those are white hairlines now.
                HEADER_ROW_BLUE = RGBColor(0x1F, 0x4E, 0x96)
                _header_blue = bool(
                    self.pptx_settings.get("financial_table_header_blue", True)
                )

                for col_idx, col_name in enumerate(df.columns[:max_cols]):
                    if col_idx < len(table.columns):
                        cell = table.cell(header_row_idx, col_idx)
                        # Replace "Description" with currency unit if found
                        if currency_unit and (col_name.lower() == 'description' or '描述' in str(col_name) or '项目' in str(col_name)):
                            cell.text = currency_unit
                        else:
                            cell.text = str(col_name)
                        # Company-format reference: the date/header row sits
                        # directly under the navy title band with NO fill of
                        # its own (white bg, dark bold text) -- only the title
                        # row above it is filled navy. The first column's own
                        # header ("Description"/currency unit) is left-aligned
                        # like a row label; the date/period columns stay
                        # centered above their right-aligned numeric data.
                        if cell.text_frame.paragraphs:
                            p = cell.text_frame.paragraphs[0]
                            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.CENTER

                            if p.runs:
                                run = p.runs[0]
                            else:
                                run = p.add_run()

                            run.font.name = 'Arial'
                            self._set_east_asian_typeface(run)
                            self._declare_run_language(run)
                            run.font.size = header_font_size
                            run.font.bold = True
                            run.font.color.rgb = WHITE if _header_blue else BLACK
                            p.line_spacing = 1.0
                            self._apply_east_asian_line_breaking(p)

                        # Blue header band with white text, matching the
                        # reference report (IMG_0224): the date/period row
                        # sits under the navy title band on its own blue
                        # fill, not on white.
                        #
                        # READ THIS BEFORE TOUCHING IT. This exact fill is
                        # what ebf2179 applied and 93c41e4 reverted: two
                        # real Chinese exports came back with the BS and IS
                        # first pages rendering COMPLETELY BLANK in real
                        # PowerPoint (slow file-open, while python-pptx still
                        # saw every shape and all text intact), and those are
                        # exactly the two pages carrying this grid. It is not
                        # reproducible on a Mac -- no real PowerPoint, and a
                        # local template whose bands are plain text boxes
                        # rather than the native placeholders production uses.
                        # Re-applied here on the user's explicit sign-off
                        # against the reference photo, but behind
                        # pptx_settings["financial_table_header_blue"] so a
                        # recurrence is one config value to flip, not a code
                        # revert -- and so the answer stays a single known
                        # switch rather than another round of bisecting.
                        cell.fill.solid()
                        if _header_blue:
                            cell.fill.fore_color.rgb = HEADER_ROW_BLUE
                        else:
                            # Only the LAST column is highlighted light blue
                            # (matches the company-format reference's
                            # "adjusted figures" column) -- every other column
                            # is EXPLICITLY set to solid white, not left
                            # unset. Leaving fill untouched makes the cell
                            # "inherit" from the table's built-in style GUID
                            # (still referenced even with first_row/
                            # horz_banding disabled), which can be a themed/
                            # tinted colour -- confirmed via
                            # inspect_pptx_tables.py showing fill=inherited on
                            # every un-highlighted cell.
                            cell.fill.fore_color.rgb = (
                                LIGHT_BLUE_HIGHLIGHT if col_idx == max_cols - 1 else WHITE
                            )

                        # Vertical (column-separating) borders plus ONE
                        # bottom rule under the header row -- horizontal
                        # rules on every data row read as cluttered once the
                        # table has 20+ rows; the total/subtotal rows below
                        # get their own explicit top/bottom rule so those
                        # separators are still there where they matter.
                        # Black column separators drawn ACROSS a navy header
                        # band read as a grid pasted over the title, which is
                        # the "header 同藍色黐唔埋" complaint. Use white hairlines
                        # between the date columns instead -- visible against
                        # navy, invisible as a boundary -- and keep the black
                        # rule only where it does real work: under the band,
                        # separating it from the data.
                        _sep = "FFFFFF" if _header_blue else "000000"
                        for _side in ("left", "right"):
                            self._set_cell_border(cell, _side, color_rgb=_sep, width=Pt(0.5))
                        self._set_cell_border(cell, "bottom", color_rgb="000000", width=Pt(0.5))

                        try:
                            cell.margin_left = Inches(0.04)
                            cell.margin_right = Inches(0.04)
                            cell.margin_top = Inches(0.02)
                            cell.margin_bottom = Inches(0.02)
                        except Exception:
                            pass

                        logger.debug("Filled header cell %s: %s", col_idx, cell.text)
                
                # Fill data rows with formatting - show ALL rows (no limit)
                # Check if table has enough rows, if not, limit to available rows
                max_rows = len(df)  # Show all rows
                rows_needed = max_rows + data_start_row + 1  # +1 for header row
                available_rows = len(table.rows)
                
                if available_rows < rows_needed:
                    logger.warning("Table has %s rows but needs %s. Will only fill %s data rows.", available_rows, rows_needed, available_rows - data_start_row - 1)
                    max_rows = min(max_rows, available_rows - data_start_row - 1)
                    if max_rows < 0:
                        max_rows = 0
                
                logger.info("Table has %s rows available, will fill %s data rows", available_rows, max_rows)
                
                # Now fill all rows with Arial 9 font
                # Check for title, date, total, and subtotal rows to highlight
                logger.info("Filling %s data rows, starting at row index %s, table has %s rows", max_rows, header_row_idx + 1, len(table.rows))
                for row_idx in range(max_rows):
                    if row_idx >= len(df):
                        break
                    df_row = df.iloc[row_idx]
                    first_col_value = str(df_row.iloc[0]) if len(df_row) > 0 else ""

                    # A section header ("流动资产" / "非流动负债" / etc.) is a
                    # row the source Financials sheet gives a label but no
                    # figures at all -- blank on every numeric column, not
                    # just a zero (a genuine zero renders as "-" via
                    # _format_table_value, so it's distinguishable from a
                    # truly blank/missing value). "Blank" covers both
                    # pd.isna() (None/NaN) AND an empty/whitespace-only
                    # string -- a real extracted sheet's blank cell isn't
                    # guaranteed to come through as NaN specifically (could
                    # already be "" depending on how the extractor read it),
                    # and pd.isna("") is False, so a NaN-only check silently
                    # never fires on that variant. Reference format (real
                    # company deliverable, IMG_0035): these sit flush left
                    # with no indent, everything else between one and the
                    # next header/total is indented under it.
                    def _is_blank_cell(v) -> bool:
                        if pd.isna(v):
                            return True
                        if isinstance(v, str) and not v.strip():
                            return True
                        return False

                    is_category_header_row = bool(first_col_value.strip()) and all(
                        _is_blank_cell(df_row[col]) for col in df.columns[1:max_cols] if col in df_row.index
                    )

                    # Check if this is a title, total, or subtotal row
                    is_special_row = False
                    is_total_row = False
                    first_col_lower = first_col_value.lower()
                    # 'total'/合计 etc. cover BS lines (Total assets, Total
                    # current assets, ...); IS running-total lines (Gross
                    # profit, Operating profit, Net profit, ...) don't contain
                    # the word "total" at all, so pull in the canonical
                    # SUMMARY_ACCOUNT_SKIP_KEYWORDS list (already used
                    # elsewhere in the codebase for "is this a statement
                    # subtotal/total line, not a leaf account" detection)
                    # rather than hand-rolling a second copy of it here.
                    total_keywords = list(
                        {'total', '合计', '总计', '小计', 'subtotal', 'sub-total', 'sub total'}
                        | set(SUMMARY_ACCOUNT_SKIP_KEYWORDS)
                    )
                    # NOTE: previously also matched a 'date_keywords' list
                    # ('date'/'日期'/'年'/'月') into special_keywords AND a
                    # separate is_date_row flag that forced a fixed Pt(7)
                    # font. The actual date/period COLUMN HEADER row (e.g.
                    # "2023-12-31") is a different row entirely, handled
                    # above via header_row_idx -- this per-data-row loop only
                    # ever sees ACCOUNT LABELS, and '年'/'月' as bare
                    # substrings match ordinary Chinese account names like
                    # "一年内到期的非流动负债" ("non-current liabilities due
                    # within one year"), wrongly bolding them AND overriding
                    # their font size to Pt(7) regardless of the table's
                    # actual density tier (data_font_size could be 6/6.5/7pt)
                    # -- a real, visible font-size/weight mismatch against
                    # every other leaf row on the same table.
                    special_keywords = total_keywords + ['title', '标题']

                    if any(keyword in first_col_lower for keyword in special_keywords):
                        is_special_row = True

                    if any(keyword in first_col_lower for keyword in total_keywords):
                        is_total_row = True

                    # A blank-values row that ALSO matches a total keyword is
                    # a total row that simply has no figures yet (e.g. a
                    # subtotal for a section with no populated accounts),
                    # not a section header -- keep those two mutually
                    # exclusive so the total-row border/fill logic below
                    # still applies to it.
                    is_category_header_row = is_category_header_row and not is_total_row
                    if is_category_header_row:
                        is_special_row = True

                    # Statement-terminal lines (grand totals) get the heavier
                    # two-border/grey-fill treatment; everything else that only
                    # matches total_keywords (subsection subtotals like "Total
                    # current assets", or IS running subtotals like "Gross
                    # profit") gets a thin top border only, matching the
                    # company-format reference where those two tiers look
                    # visually distinct.
                    grand_total_keywords = [
                        'total assets', 'total liabilities', "total owners", "total owner's",
                        '总资产', '负债合计', '负债总计', '所有者权益合计',
                        '股东权益合计', '资产总计', 'net profit', 'net loss', '净利润', '净亏损',
                        'ebitda',
                    ]
                    # "Total equity attributable to owners of the Company" is a
                    # subtotal, not the statement-level grand total, even though
                    # it contains "total ... owners" as a substring (same trap
                    # applies to the Chinese "归属于母公司所有者权益合计" pattern) --
                    # exclude any label signalling it's scoped to a sub-group.
                    attributable_signal_keywords = ['attributable', '归属', '母公司']
                    is_attributable_subtotal = any(
                        keyword in first_col_lower for keyword in attributable_signal_keywords
                    )
                    is_grand_total_row = (
                        is_total_row
                        and not is_attributable_subtotal
                        and any(keyword in first_col_lower for keyword in grand_total_keywords)
                    )
                    
                    # Data row index = header_row_idx + 1 + row_idx
                    data_row_idx = header_row_idx + 1 + row_idx
                    if data_row_idx >= len(table.rows):
                        logger.warning("Data row index %s exceeds table rows %s, skipping", data_row_idx, len(table.rows))
                        break
                    
                    # Set data row height based on table density
                    try:
                        table.rows[data_row_idx].height = data_row_height
                    except:
                        pass
                    
                    # Log first row processing
                    if row_idx == 0:
                        logger.info("Processing first data row: %s", df_row.values[:3])

                    # Unit scaling policy for the financial table:
                    #   The extractor is called with multiply_values=False
                    #   (embed_financial_tables), so numeric values flow through
                    #   in the ORIGINAL source units declared by the workbook
                    #   header. If the header says CNY'000 / 人民币千元, values
                    #   already represent thousands and must NOT be multiplied.
                    #   Same for CNY'M / 人民币百万 (millions). The column
                    #   header shows the unit so the reader interprets them
                    #   correctly. Any accidental scaling here would double-count.
                    for col_idx, col_name in enumerate(df.columns[:max_cols]):
                        if col_idx >= len(table.columns):
                            break
                        cell = table.cell(data_row_idx, col_idx)

                        # Get value from DataFrame safely
                        value = df_row[col_name] if col_name in df_row.index else ""
                        text_val = self._format_table_value(value, is_numeric_column=col_idx > 0)

                        # Description column only -- the source Financials
                        # sheet's own row labels (e.g. "Cash at bank and on
                        # hand") stay whatever language that sheet was
                        # authored in, even when the REPORT is Chinese, since
                        # nothing upstream of this table translates them.
                        if col_idx == 0 and is_chinese_mode:
                            text_val = _translate_statement_row_label(text_val, mappings)

                        # Set text
                        cell.text = text_val
                        
                        # Log first cell value of first row
                        if row_idx == 0 and col_idx < 2:
                            logger.info("Setting cell (%s, %s) to: '%s'", data_row_idx, col_idx, text_val)
                        
                        # Apply formatting: Arial 7pt (reduced from 9pt) for all cells
                        # Note: Always access paragraphs[0] AFTER setting text
                        if not cell.text_frame.paragraphs:
                            cell.text_frame.add_paragraph()
                            
                        p = cell.text_frame.paragraphs[0]
                        if not p.runs:
                            p.add_run()
                            
                        # cell.text = text_val above already wrote the text into
                        # the run; setting run.text again was a redundant XML
                        # roundtrip. Just apply the font formatting.
                        for run in p.runs:
                            run.font.name = 'Arial'
                            self._set_east_asian_typeface(run)
                            self._declare_run_language(run)
                            run.font.size = data_font_size
                            try:
                                run.font.color.rgb = BLACK
                            except Exception:
                                pass
                            run.font.bold = is_special_row
                        try:
                            p.line_spacing = 1.0
                            self._apply_east_asian_line_breaking(p)
                        except Exception:
                            pass

                        # Give cells a small internal margin so text doesn't hug the border
                        try:
                            cell.margin_left = Inches(0.04)
                            cell.margin_right = Inches(0.04)
                            cell.margin_top = Inches(0.01)
                            cell.margin_bottom = Inches(0.01)
                        except Exception:
                            pass

                        # First column left-aligned, numeric columns right-aligned.
                        # Within the label column, a leaf line item (neither a
                        # section header nor a total/subtotal) is indented
                        # under whichever header it follows -- headers and
                        # every tier of total stay flush with the left margin
                        # (reference format, IMG_0035). Set via raw XML
                        # (_set_paragraph_left_indent) -- see that method's
                        # docstring for why paragraph.left_indent itself is a
                        # no-op in this python-pptx version.
                        try:
                            p.alignment = PP_ALIGN.LEFT if col_idx == 0 else PP_ALIGN.RIGHT
                            if col_idx == 0:
                                should_indent = not is_category_header_row and not is_total_row
                                self._set_paragraph_left_indent(p, int(Inches(0.12)) if should_indent else 0)
                        except Exception:
                            pass

                        # Vertical (column-separating) borders only -- a
                        # horizontal rule under every single data row reads
                        # as visually busy/cluttered once the table has 20+
                        # rows. Total/subtotal rows get their own explicit
                        # top (and, for grand totals, bottom) rule below,
                        # applied AFTER this, so those separators are still
                        # there exactly where they matter.
                        for _side in ("left", "right"):
                            self._set_cell_border(cell, _side, color_rgb="000000", width=Pt(0.5))

                        # Only the LAST column is highlighted light blue, on
                        # EVERY row including totals -- every other cell is
                        # EXPLICITLY set to solid white (not left unset --
                        # an untouched cell.fill inherits the table's built-in
                        # style GUID, which can render as a themed/tinted
                        # colour even with first_row/horz_banding disabled).
                        # The four statement-level grand totals (total
                        # assets/liabilities/equity, and liabilities+equity)
                        # are the one exception -- reference format (IMG_0035)
                        # calls those out with a solid grey fill across the
                        # whole row; every other total/subtotal tier stays
                        # white, called out by border weight only (below).
                        try:
                            cell.fill.solid()
                            if is_grand_total_row:
                                cell.fill.fore_color.rgb = GREY_TOTAL_FILL
                            else:
                                cell.fill.fore_color.rgb = (
                                    LIGHT_BLUE_HIGHLIGHT if col_idx == max_cols - 1 else WHITE
                                )
                        except Exception:
                            pass

                        # Thin top border on every total/subtotal row; grand
                        # totals additionally get a heavier bottom border,
                        # matching the reference's two-tier total styling.
                        if is_total_row:
                            try:
                                self._set_cell_border(cell, 'top', color_rgb="00338D", width=Pt(0.75))
                                if is_grand_total_row:
                                    self._set_cell_border(cell, 'bottom', color_rgb="00338D", width=Pt(2.25))
                            except Exception:
                                pass
                    
                    logger.debug("Filled table row %s (data_row_idx: %s, special: %s)", row_idx + 1, data_row_idx, is_special_row)
                
                logger.info("Updated table with Excel data (formatting preserved)")
            else:
                # If no table, this is an error - table placeholder should be a table shape
                logger.error("Table Placeholder is not a table shape! Cannot embed financial table.")
                logger.error("Shape type: %s, has_table: %s", type(shape), hasattr(shape, 'table'))
                logger.error("Shape name: %s", shape.name if hasattr(shape, 'name') else 'unnamed')
                # Check if shape has table attribute but it's None
                if hasattr(shape, 'table'):
                    logger.error("shape.table is: %s", shape.table)
                # Try to create a table representation in text frame as last resort
                if shape.has_text_frame:
                    shape.text_frame.clear()
                    # Convert DataFrame to formatted text table - show ALL rows
                    try:
                        # Show all rows, no limit
                        text_table = df.to_string(index=False)
                    except:
                        text_table = str(df)
                    
                    p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                    p.text = text_table
                    logger.warning("Added text table representation with all %s rows (%s chars) - NOT IDEAL, should be table format", len(df), len(text_table))
        except Exception as e:
            logger.error("Could not fill table placeholder: %s", e)
            logger.error(traceback.format_exc())
            # Fallback: add text representation - show ALL rows
            if shape.has_text_frame:
                shape.text_frame.clear()
                # Show all rows, not just first 10
                text_repr = df.to_string(index=False)
                p = shape.text_frame.paragraphs[0] if shape.text_frame.paragraphs else shape.text_frame.add_paragraph()
                p.text = text_repr
    
    def _detect_bullet_levels(self, text: str) -> List[Tuple[int, str]]:
        """
        Detect bullet levels (1-3) from commentary text
        Returns list of (level, text) tuples where level 0 = no bullet, 1-3 = bullet levels
        """
        lines = text.split('\n')
        bullet_lines = []
        
        for line in lines:
            stripped = line.strip()
            original_line = line
            
            # Detect bullet lines with '- ' prefix
            if original_line.lstrip().startswith('- '):
                # Calculate indentation level (based on spaces/tabs before the bullet)
                indent_spaces = len(original_line) - len(original_line.lstrip())
                
                # Determine bullet level based on indentation (2 spaces per level)
                level = min(3, (indent_spaces // 2) + 1)  # Cap at level 3
                
                # Clean and store bullet line
                clean_line = stripped[2:]  # Remove '- '
                
                # Special handling for level 3 bullets that contain a dash indicating sub-level
                if level == 3 and " - " in clean_line:
                    # Split at the first occurrence of " - "
                    parts = clean_line.split(" - ", 1)
                    if len(parts) > 1:
                        # Add level 3 content
                        bullet_lines.append((level, parts[0].strip()))
                        # Add continuation as level 3 (indented)
                        bullet_lines.append((level, parts[1].strip()))
                    else:
                        bullet_lines.append((level, clean_line))
                else:
                    bullet_lines.append((level, clean_line))
            elif stripped:
                # Regular content (no bullet) - level 0
                bullet_lines.append((0, stripped))
        
        return bullet_lines
    
    def _determine_slot_font_size(
        self,
        slot_accounts: List[Dict],
        shape,
        slot_name: str,
        statement_type: Optional[str] = None,
    ) -> int:
        """Deck-wide fixed size: 9pt Arial for every slot on every slide,
        regardless of language or content. Any per-slot adjustment here
        reintroduces size drift between slides."""
        return 9

    _CJK_RANGE = ("一", "鿿")

    def _declare_run_language(self, run) -> None:
        """Declare the run's LANGUAGE, which is what selects the
        line-breaking algorithm.

        Third attempt at the "。 starts a line" defect, and the first at this
        mechanism. The previous two targeted the RULE (paragraph eaLnBrk /
        hangingPunct) and the FONT (<a:ea> typeface); both are set correctly
        and neither changed the render. The missing piece is that our runs
        declared no language at all -- <a:rPr> carried no lang -- so
        PowerPoint fell back to the document default, and this template is an
        English one. Latin line-breaking has no 禁则处理, which is exactly the
        behaviour observed: a line beginning with 。

        lang is set from the run's own text so this is correct in an English
        deck too; altLang says what to do with East Asian characters inside
        an otherwise-Latin run, which is the common case here ("■ ", " - ",
        and a Chinese sentence are separate runs of one bullet).
        """
        try:
            text = run.text or ""
            has_cjk = any(self._CJK_RANGE[0] <= ch <= self._CJK_RANGE[1] for ch in text)
            rPr = run._r.get_or_add_rPr()
            rPr.set("lang", "zh-CN" if has_cjk else "en-US")
            rPr.set("altLang", "zh-CN")
        except Exception as exc:
            logger.debug("Could not declare run language: %s", exc)

    def _set_east_asian_typeface(self, run) -> None:
        """Declare the East Asian typeface on a run, alongside the Latin one.

        Every run this file creates sets font.name='Arial', which python-pptx
        writes as <a:latin>. Nothing declared <a:ea>, so PowerPoint had to
        pick a CJK fallback itself -- and, more importantly, may classify the
        run as Latin script and apply LATIN line-breaking to it. That is the
        most likely reason a real deck still broke a line before a 、 even
        after eaLnBrk/hangingPunct were set: the paragraph asked for East
        Asian rules while the run looked like Latin text.

        It also removes a real inconsistency: the packer MEASURES Chinese
        text with the configured CJK metrics (Microsoft YaHei by default) but
        was telling PowerPoint the font was Arial. <a:ea> applies only to
        East Asian characters, so declaring it cannot affect Latin text and
        is safe to set unconditionally.
        """
        try:
            from pptx.oxml.ns import qn
            family = self._packing_settings().get("font_family_chi") or "Microsoft YaHei"
            rPr = run._r.get_or_add_rPr()
            for existing in rPr.findall(qn('a:ea')):
                rPr.remove(existing)
            ea = rPr.makeelement(qn('a:ea'), {'typeface': family})
            # <a:ea> must follow <a:latin> in the CT_TextCharacterProperties
            # sequence; append after it when present, else at the end.
            latin = rPr.find(qn('a:latin'))
            if latin is not None:
                latin.addnext(ea)
            else:
                rPr.append(ea)
        except Exception as exc:
            logger.debug("Could not set East Asian typeface: %s", exc)

    @staticmethod
    def _apply_east_asian_line_breaking(paragraph) -> None:
        """Turn on East Asian line-breaking (禁则处理) and hanging punctuation
        for one paragraph.

        Without this a real deck put a full stop at the START of a line, and
        in the worst case a lone "。" on its own line under a paragraph that
        otherwise ended cleanly. Chinese typography forbids a line beginning
        with closing punctuation (。，）」etc.); the rule that prevents it is
        a PARAGRAPH property, and the template declares no <a:pPr> at all, so
        nothing was asserting it.

        eaLnBrk       -- apply East Asian line-break rules rather than Latin
                         ones. Our runs carry font.name='Arial' (a Latin
                         typeface) even for Chinese text, which is exactly
                         the case where PowerPoint may otherwise fall back to
                         Latin breaking.
        hangingPunct  -- let trailing punctuation hang past the right margin
                         instead of being pushed onto the next line, which is
                         what keeps "米。" together.

        Set explicitly rather than relied on as a schema default -- the
        observed render proves the default was not being applied here.
        """
        try:
            pPr = paragraph._p.get_or_add_pPr()
            pPr.set("eaLnBrk", "1")
            pPr.set("hangingPunct", "1")
        except Exception as exc:
            logger.debug("Could not set East Asian line-breaking: %s", exc)

    @staticmethod
    def _force_no_autofit(text_frame) -> None:
        """Set the text frame's bodyPr autofit to ``<a:noAutofit/>`` so
        PowerPoint never shrinks the text to fit the shape. The template
        ships with ``<a:spAutoFit/>`` (resize shape to fit text), which in
        some viewers falls back to shrinking the text when the shape can't
        grow. Forcing ``noAutofit`` keeps the text at the exact point size
        we set (9pt / 10pt); overflow is simply clipped at the shape edge."""
        try:
            from lxml import etree  # noqa: F401
            from pptx.oxml.ns import qn
            bodyPr = text_frame._txBody.bodyPr
            # Remove any existing autofit child (spAutoFit / normAutofit / noAutofit).
            for tag in ("a:spAutoFit", "a:normAutofit", "a:noAutofit"):
                for child in bodyPr.findall(qn(tag)):
                    bodyPr.remove(child)
            from pptx.oxml import parse_xml
            bodyPr.append(parse_xml(
                '<a:noAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"/>'
            ))
        except Exception as exc:
            logger.debug("Could not force noAutofit on text frame: %s", exc)

    # A slot the DP could only pack by relaxing past strict 1.0x capacity
    # (shape_height_utilization / the later 1.35x, 1.6x tiers) has no room
    # to actually hold its content at 9pt with noAutofit -- that combination
    # silently CLIPS the overflow at the shape edge, contradicting the DP's
    # own relax-factor comment ("PPT auto-fit can absorb that much
    # overflow"). 0.70 is a floor, not the typical case: the DP's own
    # relax tiers below its 10.0x always-feasible last resort top out at
    # 1.6x (needing ~63% scale), so this floor mainly guards against that
    # last-resort tier producing illegibly small text rather than clipping.
    _BOUNDED_AUTOFIT_MIN_SCALE = 0.70

    @classmethod
    def _apply_bounded_autofit(cls, text_frame, font_scale: float) -> None:
        """Set ``<a:normAutofit fontScale="..." lnSpcReduction="..."/>`` so
        PowerPoint actually shrinks text to fit instead of clipping it, for
        a slot the DP intentionally packed beyond strict 1.0x capacity.
        ``font_scale`` is capacity/used, clamped to
        ``_BOUNDED_AUTOFIT_MIN_SCALE`` so an extreme overflow (the DP's
        10.0x always-feasible last-resort tier) still clips rather than
        rendering illegibly small text -- a bounded, not unlimited, shrink.
        """
        scale = max(cls._BOUNDED_AUTOFIT_MIN_SCALE, min(1.0, float(font_scale)))
        font_pct = int(round(scale * 100000))
        # PowerPoint reduces line spacing somewhat less aggressively than
        # font size when auto-fitting -- half the font shrink, floored at 0.
        line_pct = int(round(max(0.0, (1.0 - scale) * 0.5) * 100000))
        try:
            from pptx.oxml.ns import qn
            from pptx.oxml import parse_xml
            bodyPr = text_frame._txBody.bodyPr
            for tag in ("a:spAutoFit", "a:normAutofit", "a:noAutofit"):
                for child in bodyPr.findall(qn(tag)):
                    bodyPr.remove(child)
            bodyPr.append(parse_xml(
                f'<a:normAutofit xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                f'fontScale="{font_pct}" lnSpcReduction="{line_pct}"/>'
            ))
        except Exception as exc:
            logger.debug("Could not apply bounded normAutofit on text frame: %s", exc)

    def _determine_slot_font_size_UNUSED(
        self,
        slot_accounts: List[Dict],
        shape,
        slot_name: str,
        statement_type: Optional[str] = None,
    ) -> int:
        """KEPT FOR REFERENCE — old shrink-to-fit logic (9→8→7pt)."""
        packing = self._packing_settings(statement_type)
        if not shape or not hasattr(shape, "height"):
            return 9

        is_chinese_slot = any(
            any("\u4e00" <= c <= "\u9fff" for c in str(a.get("commentary", "")))
            for a in slot_accounts
        )

        pillow_ok = self._pillow_fitting_enabled(packing)
        if pillow_ok:
            try:
                from fdd_utils.text_metrics import (
                    get_font,
                    line_height_pt as _line_h,
                    lines_that_fit,
                    text_box_from_shape,
                    wrap_text,
                )
                box = text_box_from_shape(shape)
                family = "Microsoft YaHei" if is_chinese_slot else "Arial"
                line_spacing = 0.95 if is_chinese_slot else 1.0
                for candidate_pt in (9, 8, 7):
                    font = get_font(family, candidate_pt, is_cjk=is_chinese_slot)
                    line_h = _line_h(font, line_spacing=line_spacing)
                    capacity = lines_that_fit(box.height_pt, line_h)
                    total_lines = 0
                    prev_cat = None
                    for acct in slot_accounts:
                        cat = acct.get("category", "")
                        if cat and cat != prev_cat:
                            total_lines += 1
                            prev_cat = cat
                        parts: List[str] = []
                        mapping_key = acct.get("mapping_key", acct.get("account_name", ""))
                        if mapping_key:
                            parts.append(str(mapping_key))
                        commentary = str(acct.get("commentary", ""))
                        if commentary:
                            parts.append(commentary)
                        joined = "\n".join(parts)
                        if joined:
                            total_lines += len(wrap_text(joined, font, box.width_pt))
                    if total_lines <= capacity:
                        return candidate_pt
                return 7
            except Exception:
                pass  # fall through to legacy

        for candidate_pt in (9, 8, 7):
            shape_height_pt = shape.height * 72 / 914400
            effective_height = shape_height_pt * float(packing.get("shape_height_utilization", 1.02))
            line_spacing = 0.95 if is_chinese_slot else 1.0
            line_height = (candidate_pt * line_spacing) + float(packing.get("line_height_padding_pt", 1.6))
            max_lines = int(effective_height / line_height)

            total_lines = 0
            prev_cat = None
            for acct in slot_accounts:
                cat = acct.get("category", "")
                if cat and cat != prev_cat:
                    total_lines += 1
                    prev_cat = cat
                commentary = str(acct.get("commentary", ""))
                is_chi = any("\u4e00" <= c <= "\u9fff" for c in commentary)
                base_cpl = self._estimate_chars_per_line(slot_name, is_chi, shape=shape, statement_type=statement_type)
                scale = 9.0 / candidate_pt
                cpl = max(16, int(base_cpl * scale))
                total_lines += 1  # key line
                for line in commentary.split("\n"):
                    if line.strip():
                        total_lines += max(1, (len(line) + cpl - 1) // cpl)

            if total_lines <= max_lines:
                return candidate_pt
        return 7

    @staticmethod
    def _build_clause_segments(
        commentary: str,
        clause_reviews: Optional[List[Dict[str, Any]]],
    ) -> Optional[List[Tuple[str, str]]]:
        """Split commentary into (text, category) segments using clause_reviews.

        Returns None if no clauses match. Falls back to a single 'data-backed'
        segment for any text not matched by any clause review (so unmatched
        prose stays black).
        """
        if not commentary or not clause_reviews:
            return None
        # Sort clauses by their position in the commentary
        positions: List[Tuple[int, int, str]] = []
        used_starts: set = set()
        for review in clause_reviews:
            clause_text = str(review.get("clause") or "").strip()
            if not clause_text:
                continue
            category = str(review.get("category") or ("data-backed" if review.get("supported") else "hallucination")).lower()
            search_from = 0
            # Find first non-overlapping occurrence
            while True:
                idx = commentary.find(clause_text, search_from)
                if idx == -1:
                    break
                if idx in used_starts:
                    search_from = idx + 1
                    continue
                used_starts.add(idx)
                positions.append((idx, idx + len(clause_text), category))
                break
        if not positions:
            return None
        positions.sort()
        # Merge overlaps by sorting and skipping fully-contained overlaps
        cleaned: List[Tuple[int, int, str]] = []
        for start, end, cat in positions:
            if cleaned and start < cleaned[-1][1]:
                continue
            cleaned.append((start, end, cat))
        # Build segments from start to end of commentary
        segments: List[Tuple[str, str]] = []
        cursor = 0
        for start, end, cat in cleaned:
            if start > cursor:
                segments.append((commentary[cursor:start], "data-backed"))
            segments.append((commentary[start:end], cat))
            cursor = end
        if cursor < len(commentary):
            segments.append((commentary[cursor:], "data-backed"))
        return segments

    @staticmethod
    def _category_to_rgb(category: str) -> Optional[Tuple[int, int, int]]:
        if category == "hallucination":
            return (200, 16, 46)  # red
        if category == "reasoning":
            return (213, 94, 0)  # orange
        return None

    def _add_runs_for_line(
        self,
        paragraph,
        line: str,
        clause_segments: Optional[List[Tuple[str, str]]],
        font_size_pt: int,
    ) -> None:
        """Add one or more runs to `paragraph` for `line`, applying clause colours when applicable."""
        from pptx.dml.color import RGBColor

        def _apply_run_format(run, color_rgb: Optional[Tuple[int, int, int]]):
            run.font.size = Pt(font_size_pt)
            run.font.name = 'Arial'
            self._set_east_asian_typeface(run)
            self._declare_run_language(run)
            run.font.bold = False
            try:
                if color_rgb is not None:
                    run.font.color.rgb = RGBColor(*color_rgb)
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
            except Exception:
                pass

        if not clause_segments:
            run = paragraph.add_run()
            run.text = line
            _apply_run_format(run, None)
            return

        # Find which segments overlap with this line. Since segments span the
        # whole commentary, we need to figure out where this line fits. For
        # simplicity, walk the segments while consuming this line's characters.
        remaining = line
        for segment_text, category in clause_segments:
            if not remaining:
                break
            if not segment_text:
                continue
            # Skip leading characters of the segment that aren't in the line
            # (e.g. the segment may start before this line begins).
            overlap_start = remaining.find(segment_text)
            if overlap_start == 0:
                run = paragraph.add_run()
                run.text = segment_text
                _apply_run_format(run, self._category_to_rgb(category))
                remaining = remaining[len(segment_text):]
            elif overlap_start > 0:
                # Plain prefix before this segment
                run = paragraph.add_run()
                run.text = remaining[:overlap_start]
                _apply_run_format(run, None)
                # Then the segment
                run = paragraph.add_run()
                run.text = segment_text
                _apply_run_format(run, self._category_to_rgb(category))
                remaining = remaining[overlap_start + len(segment_text):]
            # else: segment doesn't appear on this line, skip it
        if remaining:
            run = paragraph.add_run()
            run.text = remaining
            _apply_run_format(run, None)

    @staticmethod
    def _truncate_commentary_to_fit(commentary: str, max_chars: int) -> str:
        """Hard truncation at sentence boundary, with ellipsis."""
        if len(commentary) <= max_chars:
            return commentary
        truncated = commentary[:max_chars]
        # Try to cut at sentence boundary
        for end_char in (". ", "。", "! ", "？"):
            pos = truncated.rfind(end_char)
            if pos > max_chars * 0.5:
                return truncated[: pos + len(end_char)].rstrip()
        # Fall back to word boundary
        pos = truncated.rfind(" ", int(max_chars * 0.7))
        if pos > 0:
            return truncated[:pos].rstrip() + "..."
        return truncated.rstrip() + "..."

    _ORPHANABLE_END_PUNCT = "。．.；;"

    def _drop_orphan_trailing_punctuation(self, key_prefix: str, text: str, is_chinese: bool) -> str:
        """Delete a sentence-ending mark that is the only thing pushing a
        paragraph onto one more line.

        A lone "。" starting a line is wrong in Chinese typography and the
        project team will not accept it. The usual remedies are unavailable
        here: shrinking the font is explicitly ruled out, and the kinsoku
        controls PowerPoint should honour (eaLnBrk, hangingPunct, the run's
        own lang) were all set correctly across three attempts and changed
        nothing in the real render. So take the remedy the team named
        themselves -- "直接刪除句號 他不是那麼重要".

        Measured the way POWERPOINT breaks, not the way this repo does. Our
        own wrapper already implements kinsoku, so it hangs the mark on the
        previous line for free and a line-count comparison here can never
        detect the problem -- the first version of this check was written
        that way and provably never fired. Instead: wrap without the mark,
        then ask whether the mark still fits on the resulting last line. If
        it does not, PowerPoint (which does no hanging) puts it alone on the
        next line, and it goes.
        """
        from fdd_utils.text_metrics import get_measurer, text_box_from_shape

        stripped = (text or "").rstrip()
        if not stripped or stripped[-1] not in self._ORPHANABLE_END_PUNCT:
            return text
        shape = self._measurement_slot_shape()
        box = text_box_from_shape(shape)
        if box is None or box.width_pt <= 0:
            return text
        packing = self._packing_settings()
        measurer = get_measurer(
            self._measurer_family(is_chinese, packing),
            self._real_font_size_pt(is_chinese), is_cjk=is_chinese,
            line_spacing=self._real_line_spacing(is_chinese),
            metrics_path=self._resolve_font_metrics_path(is_chinese, packing),
        )
        body_width = max(1.0, box.width_pt - self._BULLET_HANGING_INDENT_PT)
        lines = measurer.wrap(
            key_prefix + stripped[:-1], body_width, first_line_width_pt=box.width_pt,
        )
        if not lines:
            return text
        if measurer.text_width_pt(lines[-1] + stripped[-1]) > body_width:
            return stripped[:-1]
        return text

    def _fill_text_main_bullets_with_category_and_key(self, text_frame, category: str, display_name: str,
                                                      commentary: str, is_chinese: bool, is_chinese_databook: bool = False,
                                                      needs_continuation: bool = False, font_size_pt: int = 9,
                                                      clause_reviews: Optional[List[Dict[str, Any]]] = None):
        """
        Fill textMainBullets shape with commentary formatted as:
        - Category as first level (dark blue Arial 9) - only if category is provided
        - Key name with filled round bullet + space + key name (black bold Arial 9) + "-" (not bold) + plain text
        - Indentation 0.15" with special hanging 0.15", spacing after 6pt
        - When clause_reviews is provided, non-data-backed clauses are coloured:
          orange (213, 94, 0) for 'reasoning', red (200, 16, 46) for 'hallucination'.
        """
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        # Before anything measures or renders this text, and before the
        # clause segments that carry hallucination highlighting are cut from
        # it, so every downstream consumer sees the same string.
        commentary = self._drop_orphan_trailing_punctuation(
            f"■ {display_name} - " if display_name else "", commentary, is_chinese,
        )

        clause_segments = self._build_clause_segments(commentary, clause_reviews) if clause_reviews else None

        # Add category as first level (if category exists and is not None)
        # Note: category is now handled at slide level, so this is only for individual calls
        if category:
            p_category = text_frame.add_paragraph()
            p_category.level = 0
            try:
                p_category.left_indent = Inches(0.21)
                p_category.first_line_indent = Inches(-0.19)
                p_category.space_before = Pt(0)
                p_category.space_after = Pt(0)
                p_category.line_spacing = 1.0
                self._apply_east_asian_line_breaking(p_category)
            except:
                pass
            
            run_category = p_category.add_run()
            # The parallel header site in apply_structured_data_to_slides
            # translates; this one never did, despite taking
            # is_chinese_databook as a parameter. Only the presentation-table
            # renderer passes a real category here (the other caller passes
            # None and writes its own header), so the symptom was narrow and
            # easy to miss: a Chinese deck whose table pages alone carried an
            # English section heading -- "Expenses" sitting above 税金及附加
            # and 营业成本, while every non-table page read 流动资产/营业收入.
            run_category.text = (
                translate_category_to_chinese(category) if is_chinese_databook else category
            )
            run_category.font.size = Pt(font_size_pt)
            run_category.font.name = 'Arial'
            self._set_east_asian_typeface(run_category)
            self._declare_run_language(run_category)
            run_category.font.bold = False
            try:
                run_category.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue
            except:
                pass
        
        # Add key name with grey char + space + key name (black bold) + "-" (not bold) + plain text
        p_key = text_frame.add_paragraph()
        p_key.level = 0  # No bullet level, we'll use grey character
        try:
            # Set formatting
            p_key.left_indent = Inches(0.15)  # 0.15" indent
            p_key.first_line_indent = Inches(-0.15)  # 0.15" special hanging
            p_key.space_before = Pt(0)
            p_key.space_after = Pt(3)  # Matches _PARA_SPACE_AFTER (cost estimator)
            p_key.line_spacing = 1.0
            self._apply_east_asian_line_breaking(p_key)
        except Exception as e:
            logger.warning("Could not set paragraph formatting: %s", e)
            pass
        
        # Grey char (U+25A0) + space
        run_bullet = p_key.add_run()
        run_bullet.text = '\u25A0 '  # U+25A0 (black square) + space
        run_bullet.font.size = Pt(font_size_pt)
        run_bullet.font.name = 'Arial'
        self._set_east_asian_typeface(run_bullet)
        self._declare_run_language(run_bullet)
        run_bullet.font.bold = False
        try:
            run_bullet.font.color.rgb = RGBColor(128, 128, 128)  # Grey
        except:
            pass

        # Key name (black bold)
        run_key = p_key.add_run()
        run_key.text = display_name
        run_key.font.size = Pt(font_size_pt)
        run_key.font.name = 'Arial'
        self._set_east_asian_typeface(run_key)
        self._declare_run_language(run_key)
        run_key.font.bold = True
        try:
            run_key.font.color.rgb = RGBColor(0, 0, 0)  # Black
        except:
            pass

        # "-" (not bold)
        run_dash = p_key.add_run()
        run_dash.text = " - "
        run_dash.font.size = Pt(font_size_pt)
        run_dash.font.name = 'Arial'
        self._set_east_asian_typeface(run_dash)
        self._declare_run_language(run_dash)
        run_dash.font.bold = False
        try:
            run_dash.font.color.rgb = RGBColor(0, 0, 0)  # Black
        except:
            pass
        
        # Plain text (commentary content)
        commentary_lines = commentary.split('\n')
        first_line_added = False
        for line_idx, line in enumerate(commentary_lines):
            line = line.strip()
            if not line:
                continue

            if not first_line_added:
                target_paragraph = p_key
                first_line_added = True
            else:
                # Subsequent lines as new paragraphs (indented continuation)
                p_text = text_frame.add_paragraph()
                p_text.level = 0  # No bullet for continuation
                try:
                    p_text.left_indent = Inches(0.15)  # 0.15" indent (same as key text)
                    p_text.first_line_indent = Inches(0)  # No hanging for continuation lines
                    p_text.space_before = Pt(0)
                    p_text.space_after = Pt(3)
                    p_text.line_spacing = 1.0
                    self._apply_east_asian_line_breaking(p_text)
                except:
                    pass
                target_paragraph = p_text

            self._add_runs_for_line(
                target_paragraph,
                line,
                clause_segments=clause_segments,
                font_size_pt=font_size_pt,
            )

        # Note: "(continued)" is now added to category header, not here
    
    def _fill_text_main_bullets_with_levels(self, text_frame, commentary: str, is_chinese: bool):
        """
        Fill textMainBullets shape with commentary using detailed line break logic
        and level 1-3 text handling with page breaks (legacy method, kept for compatibility)
        """
        from pptx.util import Inches
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        # Detect bullet levels
        bullet_lines = self._detect_bullet_levels(commentary)
        
        # Calculate max lines that can fit in the shape
        # Estimate based on shape height (conservative estimate)
        max_lines = 35  # Default conservative estimate
        
        lines_added = 0
        
        for level, text in bullet_lines:
            if not text.strip():
                continue
            
            # Check if we need a page break (if shape is getting full)
            # Note: Actual page breaks would require creating new slides, which is handled
            # at a higher level. Here we just ensure content fits.
            if lines_added >= max_lines:
                # Add continuation indicator
                p = text_frame.add_paragraph()
                p.level = 0
                run = p.add_run()
                run.text = "... (continued on next page)" if not is_chinese else "... (续下页)"
                run.font.size = get_font_size_for_text(run.text, force_chinese_mode=is_chinese)
                run.font.name = get_font_name_for_text(run.text)
                run.font.italic = True
                break
            
            # Create paragraph with appropriate level
            p = text_frame.add_paragraph()
            p.level = level  # Set bullet level (0-3)
            
            # Apply paragraph formatting based on level
            try:
                # Level 0 (no bullet) or Level 1 (main bullet)
                if level == 0 or level == 1:
                    p.left_indent = Inches(0.21)  # 0.21" indent before text
                    p.first_line_indent = Inches(-0.19)  # 0.19" special hanging
                    p.space_before = Pt(0)  # 0pt spacing before
                    p.space_after = Pt(0)  # 0pt spacing after
                    p.line_spacing = 1.0  # Single line spacing
                elif level == 2:
                    # Level 2 - more indented
                    p.left_indent = Inches(0.4)
                    p.first_line_indent = Inches(-0.19)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    p.line_spacing = 1.0
                    self._apply_east_asian_line_breaking(p)
                elif level == 3:
                    # Level 3 - most indented
                    p.left_indent = Inches(0.6)
                    p.first_line_indent = Inches(-0.19)
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    p.line_spacing = 1.0
                    self._apply_east_asian_line_breaking(p)
            except:
                pass  # Silently handle formatting errors
            
            # Add text with proper formatting
            run = p.add_run()
            run.text = text
            run.font.size = get_font_size_for_text(text, force_chinese_mode=is_chinese)
            run.font.name = get_font_name_for_text(text)
            
            # Apply level-specific formatting
            if level == 1:
                run.font.bold = True
                try:
                    run.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue for level 1
                except:
                    pass
            elif level == 0:
                # Regular text - no special formatting
                pass
            
            lines_added += 1
    
    def _validate_ai_summary(
        self,
        source_text: str,
        draft_summary: str,
        is_chinese: bool,
        ai_helper: Optional[Any] = None,
    ) -> Optional[str]:
        summary_settings = self._summary_settings()
        if not bool(summary_settings.get("enable_validation", True)):
            return draft_summary

        try:
            from fdd_utils.ai import AIClient

            model_type = ai_helper.model_type if ai_helper is not None else self._resolve_summary_model_type(is_chinese)
            if model_type == "local" and not bool(summary_settings.get("local_enable_validation", False)):
                logger.info("Skipping PPTX summary validation for local model; using draft summary directly")
                return draft_summary

            max_input_chars = int(summary_settings.get("max_input_chars", 1400))
            validation_max_tokens = int(summary_settings.get("validation_max_tokens", 90))
            max_numeric_sentences = int(summary_settings.get("max_numeric_sentences", 1))
            validation_timeout_seconds = float(summary_settings.get("validation_timeout_seconds", 25) or 25)
            # Sized to the real summary box, not a static config number --
            # see _summary_length_targets.
            target_chars_chi, target_words_eng = self._summary_length_targets(is_chinese)
            max_sentences_chi = int(summary_settings.get("max_sentences_chi", 4))
            max_sentences_eng = int(summary_settings.get("max_sentences_eng", 4))

            if is_chinese:
                prompt = f"""请校验以下PPT执行摘要草稿，使其适合作为财务PPT摘要框内容。

要求：
1. 只保留与原始评论一致的高层结论、趋势和核心驱动。
2. 控制在{max_sentences_chi}句话以内，长度约{target_chars_chi}字 —— 若草稿明显短于目标长度，请补充其他高层要点以达到目标，不要仅做压缩。
3. 最多保留{max_numeric_sentences}个数字或百分比，除非删除后会影响结论准确性。
4. 删除重复、堆叠金额和逐项罗列，但不得为了简短而牺牲目标长度或丢失核心趋势、驱动和结论。
5. 优先合并重复句、删去铺垫和次要背景，只保留最重要的业务含义。
6. 只输出最终摘要，不要解释。

原始评论：
{source_text[:max_input_chars]}

摘要草稿：
{draft_summary}"""
            else:
                prompt = f"""Validate the draft executive summary for a financial PPT summary box.

Requirements:
1. Keep only source-supported themes, trend, and core driver.
2. Limit the result to no more than {max_sentences_eng} sentences and about {target_words_eng} words —
   if the draft runs noticeably shorter than the target, add other high-level points to reach it
   rather than just compressing further.
3. Keep at most {max_numeric_sentences} number or percentage unless removing it would make the summary inaccurate.
4. Remove repeated phrasing, stacked figures, scene-setting language, and account-by-account detail,
   but do not sacrifice the target length or the key trend, driver, or conclusion just to be terse.
5. Output only the final validated summary paragraph.

Source commentary:
{source_text[:max_input_chars]}

Draft summary:
{draft_summary}"""

            ai_helper = ai_helper or AIClient(
                model_type=model_type,
                language='Chi' if is_chinese else 'Eng',
                model_name=self.model_name,
            )
            validation_max_retries = int(summary_settings.get("validation_max_retries", 2) or 2)
            response = self._call_with_timeout_retry(
                lambda: ai_helper.get_response(
                    user_prompt=prompt,
                    system_prompt=(
                        "You validate executive summaries for financial presentation slides. "
                        "Keep only source-supported, concise, presentation-ready conclusions."
                    ),
                    temperature=float(summary_settings.get("validation_temperature", 0.1) or 0.1),
                    max_tokens=validation_max_tokens,
                ),
                timeout_seconds=validation_timeout_seconds,
                max_retries=validation_max_retries,
                timeout_label="PPTX summary validation",
            )
            validated_summary = str((response or {}).get("content") or "").strip()
            if _looks_like_blocked_ai_content(validated_summary):
                logger.warning("PPTX summary validation returned blocked/network HTML content; using draft summary fallback")
                return draft_summary
            return validated_summary or draft_summary
        except Exception as exc:
            logger.warning("Could not validate AI summary: %s", exc)
            return draft_summary

    @classmethod
    def strip_table_detail_for_summary(cls, text: str, is_chinese: bool) -> str:
        """Truncates one account's commentary to its lead-in only, at the
        same "明细如下："/"the breakdown is set out below" handoff phrase
        _split_table_commentary already splits on for PPTX rendering --
        callers building a page/section-level summary blob (ui.py, both the
        single-file and batch paths) should use ONLY this lead-in, never
        the per-component "-"/"➢" bullets after it.

        Root cause this exists to prevent: the non-AI fallback summary
        (_generate_page_summary) picks each account's "first sentence" via
        a naive splitter that doesn't know "明细如下：" isn't a real sentence
        end -- for a table account's RAW (unsplit) commentary, that made
        its "first sentence" run straight through the handoff phrase and
        into the whole first "➢" bullet, producing a summary that visibly
        spliced in the table's own detail bullets and cut off mid-sentence
        (confirmed against a real Crescent export's coSummaryShape). Safe
        to call on every account unconditionally, table-bearing or not --
        _split_table_commentary already falls back to the whole text
        unchanged when the handoff phrase isn't present.
        """
        generator = cls.__new__(cls)
        lead_in, _post = generator._split_table_commentary(text, is_chinese)
        return lead_in

    @classmethod
    def generate_section_summary(
        cls,
        commentary: str,
        *,
        is_chinese: bool,
        language: str = "english",
        model_type: Optional[str] = None,
        model_name: Optional[str] = None,
    ) -> Optional[str]:
        """Top-level helper: generate one executive summary from concatenated
        commentary for a BS or IS section. Designed to be called from the UI
        during the AI commentary phase so the PPTX export becomes pure XML
        (no AI calls during export).

        Returns the summary string, or None if AI is disabled / fails.
        """
        try:
            generator = cls.__new__(cls)
            # Use the same config-merged settings the full generator uses,
            # otherwise the timeout/retry defaults from config.yml are
            # ignored (would fire at 10s instead of the configured value).
            generator.pptx_settings = _load_pptx_settings()
            generator.model_type = model_type
            generator.model_name = model_name
            generator.language = language
            result = generator._generate_ai_summary(commentary, commentary, is_chinese)
            if result is None:
                # AI timed out or unavailable — fall back to the rule-based
                # summary so coSummaryShape is never left blank.
                result = generator._generate_page_summary(commentary, is_chinese) or None
            return result
        except Exception as exc:
            logger.warning("generate_section_summary failed: %s", exc)
            return None

    def _generate_ai_summary(self, commentary: str, summary_source: str, is_chinese: bool) -> Optional[str]:
        """Generate and validate AI summary from page commentary."""
        try:
            from fdd_utils.ai import AIClient
            summary_settings = self._summary_settings()
            if not bool(summary_settings.get("enable_ai", True)):
                logger.info("PPTX summary AI disabled by config; using fallback summary")
                return None
            model_type = self._resolve_summary_model_type(is_chinese)
            max_input_chars = int(summary_settings.get("max_input_chars", 1600))
            max_tokens = int(summary_settings.get("max_tokens", 180))
            max_numeric_sentences = int(summary_settings.get("max_numeric_sentences", 1))
            # Use a shorter timeout for local models — they either answer fast
            # or they're not running; long waits just block the export.
            _is_local = str(model_type or "").lower() == "local"
            generation_timeout_seconds = float(
                summary_settings.get("local_generation_timeout_seconds", 10)
                if _is_local else
                summary_settings.get("generation_timeout_seconds", 20)
            )
            # Sized to the real summary box, not a static config number --
            # see _summary_length_targets.
            target_chars_chi, target_words_eng = self._summary_length_targets(is_chinese)
            max_sentences_chi = int(summary_settings.get("max_sentences_chi", 4))
            max_sentences_eng = int(summary_settings.get("max_sentences_eng", 4))
            source_text = str(commentary or summary_source or "").strip()
            if not source_text:
                return None

            if is_chinese:
                prompt = f"""请将以下财务评论改写成适合PPT摘要框的高层执行摘要。

目标长度：约{target_chars_chi}字，控制在{max_sentences_chi}句话以内 —— 请写满这个长度，不要明显短于目标。

要求：
1. 保留高层结论、趋势和核心驱动，可覆盖一个以上要点以达到目标长度。
2. 除非极其必要，最多保留{max_numeric_sentences}个数字或百分比。
3. 不要逐项复述账户，不要堆叠金额细节。
4. 语气要像管理层摘要，写成一个紧凑自然的短段落。
5. 优先删去次要说明、重复铺垫和账户层级细节，只保留最重要的业务结论、驱动和影响，但不要为了简短而牺牲目标长度。

原始内容：
{source_text[:max_input_chars]}"""
            else:
                prompt = f"""Write a short executive summary for a PPT summary box based on the following financial commentary.

Target length: about {target_words_eng} words, with no more than {max_sentences_eng} sentences —
write to fill this length, do not stop noticeably short of it.

Requirements:
1. Cover overall trend, key driver, and business implication — span more than one theme if needed to reach the target length.
2. Keep it high level and presentation-friendly.
3. Include at most {max_numeric_sentences} number or percentage unless absolutely necessary.
4. Do not list account-by-account detail or repeat many figures.
5. Write one compact management-style paragraph. Remove secondary detail, scene-setting language, and repeated wording, but do not sacrifice the target length just to be terse.
6. Remove secondary detail, scene-setting language, and repeated wording.

Original content:
{source_text[:max_input_chars]}"""

            ai_helper = AIClient(
                model_type=model_type,
                language='Chi' if is_chinese else 'Eng',
                model_name=self.model_name,
            )
            generation_max_retries = int(
                summary_settings.get("local_generation_max_retries", 1)
                if _is_local else
                summary_settings.get("generation_max_retries", 1)
            )
            response = self._call_with_timeout_retry(
                lambda: ai_helper.get_response(
                    user_prompt=prompt,
                    system_prompt=(
                        "You write concise executive summaries for financial presentation slides. "
                        "Prefer themes, drivers, and implications over detailed figures."
                    ),
                    temperature=float(summary_settings.get("generation_temperature", 0.2) or 0.2),
                    max_tokens=max_tokens,
                ),
                timeout_seconds=generation_timeout_seconds,
                max_retries=generation_max_retries,
                timeout_label="PPTX summary generation",
            )
            from fdd_utils.ai import strip_thinking
            summary = strip_thinking(str((response or {}).get("content") or "")).strip()
            if _looks_like_blocked_ai_content(summary):
                logger.warning(
                    "PPTX summary generation returned blocked/network HTML content; falling back to compact summary"
                )
                return None

            if summary:
                return self._validate_ai_summary(source_text, summary, is_chinese, ai_helper=ai_helper)
        except Exception as e:
            logger.warning("Could not generate AI summary: %s", e)
            logger.debug(traceback.format_exc())
        
        return None
    
    def _generate_page_summary(self, commentary: str, is_chinese: bool) -> str:
        """Fallback (non-AI) page summary.

        Sentences are taken in RANK ORDER across account paragraphs -- every
        account's opening sentence first, then every account's second, and so
        on -- until the measured character budget for the box is used up. The
        rank ordering keeps the original intent (the summary spans the whole
        page rather than covering only the first account); filling to the
        budget is what stops it under-running the box.

        The old version took exactly one sentence per account and stopped at
        a static max_sentences (4). Both limits starved the band on a real
        deck: a 7-account page showed 4 sentences / 168 chars against a ~276
        char budget, and a single-account page showed one sentence / 31
        chars -- reported as "exe sum 非常短". max_chars is derived from the
        real coSummaryShape width x target_lines, so it already expresses
        "as much as the box holds" far better than a hand-tuned sentence
        count, and now governs on its own.
        """
        if not commentary or not commentary.strip():
            return ""
        is_chinese_text = is_chinese or detect_chinese_text(commentary)
        _t_chi, _t_eng = self._summary_length_targets(is_chinese_text)
        max_chars = _t_chi if is_chinese_text else _t_eng * 6

        # Each account block is separated by "\n\n".
        blocks = [b.strip() for b in commentary.split("\n\n") if b.strip()]
        by_block = [_split_text_sentences(b, is_chinese_text) for b in blocks]
        by_block = [s for s in by_block if s]
        if not by_block:
            by_block = [_split_text_sentences(commentary, is_chinese_text)]
            by_block = [s for s in by_block if s]
            if not by_block:
                return ""

        sep = "" if is_chinese_text else " "
        picked: List[str] = []
        used = 0
        for rank in range(max(len(s) for s in by_block)):
            for sentences in by_block:
                if rank >= len(sentences):
                    continue
                candidate = sentences[rank].strip()
                if not candidate:
                    continue
                cost = len(candidate) + (len(sep) if picked else 0)
                if picked and used + cost > max_chars:
                    continue
                picked.append(candidate)
                used += cost

        if not picked:
            picked = [by_block[0][0]]

        summary = sep.join(picked).strip()
        if len(summary) > max_chars:
            summary = summary[:max_chars].rstrip(" ,;:/-") + "…"
        return summary.strip()

    def embed_financial_tables(
        self,
        excel_path: str,
        sheet_name: str,
        project_name: str,
        language: str,
        bs_is_results: Optional[Dict[str, Any]] = None,
        mappings: Optional[Dict[str, Any]] = None,
    ):
        """Embed financial tables: BS to page 1, IS to page 5"""
        try:
            import pandas as pd
            from fdd_utils.workbook import extract_balance_sheet_and_income_statement
            
            logger.info("Embedding financial tables from %s, sheet: %s", excel_path, sheet_name)

            # Only bail here if there's NEITHER a precomputed bs_is_results
            # NOR enough to extract one fresh -- excel_path/sheet_name being
            # blank is fine on its own when the caller already did the
            # extraction (e.g. from a roll-up workbook, or a synthesized
            # BS/IS with no literal Financials sheet at all).
            if bs_is_results is None and (not excel_path or not sheet_name):
                logger.warning("Missing excel_path (%s) or sheet_name (%s), and no precomputed BS/IS results -- skipping table embedding", excel_path, sheet_name)
                return
            
            # Use the precomputed results when available (they came from the
            # default extractor, multiply_values=True, so numeric values are
            # already in actual units). Only extract fresh if nothing was
            # passed in. Either way, at display time we rescale to the source
            # unit (CNY'000 / 人民币千元 / CNY'M / 人民币百万) so cells line up
            # with the header. This avoids the fragility of a second extract
            # that could silently return empty and lose the table.
            if bs_is_results is None:
                try:
                    logger.info("No precomputed BS/IS; extracting fresh")
                    bs_is_results = extract_balance_sheet_and_income_statement(
                        excel_path,
                        sheet_name,
                        debug=False,
                    )
                except Exception as exc:
                    logger.warning("Fresh BS/IS extraction failed: %s", exc)
                    return

            if not bs_is_results:
                logger.warning("No BS/IS data available for PPTX tables")
                return

            # Values as received have been multiplied by 1000 if the source
            # sheet declared CNY'000 / 人民币千元, and left as-is otherwise.
            # Rescale once we know the unit label (detected further below).
            values_pre_multiplied = True
            
            # Extract BS and IS DataFrames from results. Copy them — the rescale
            # block below divides numeric columns by 1000 IN PLACE, and
            # bs_is_results is the same object as session_state.bs_is_results,
            # which survives across re-exports (only cleared on a new file
            # upload). Without the copy, a second export in the same session
            # would divide already-divided values by 1000 again, showing
            # table figures 1000x too small.
            bs_df = bs_is_results.get('balance_sheet')
            is_df = bs_is_results.get('income_statement')
            bs_df = bs_df.copy() if bs_df is not None else None
            is_df = is_df.copy() if is_df is not None else None
            
            # Table titles follow the standard FDD phrasing regardless of what
            # the source Excel calls the sheet. Language-aware so Chinese decks
            # stay consistent with English decks.
            is_chinese_mode = str(language or "").strip().lower().startswith(("chi", "zh", "cn"))
            project_suffix = f" - {project_name}" if project_name else ""
            if is_chinese_mode:
                bs_table_name = f"示意性调整后资产负债表{project_suffix}"
                is_table_name = f"示意性调整后利润表{project_suffix}"
            else:
                bs_table_name = f"Indicative adjusted balance sheet{project_suffix}"
                is_table_name = f"Indicative adjusted income statement{project_suffix}"

            # Detect currency unit from the sheet header. Currency markers live
            # in the first 20 rows (table titles / unit row); reading the full
            # sheet via iterrows() was a ~1-3s hit on big workbooks. Cap to
            # nrows=20 and use vectorised astype(str) instead of iterrows.
            currency_unit = None
            try:
                excel_df = pd.read_excel(excel_path, sheet_name=sheet_name, header=None, nrows=20)
                blob = ' '.join(
                    excel_df.fillna('').astype(str).agg(' '.join, axis=1).tolist()
                )
                if '人民币百万' in blob:
                    currency_unit = '人民币百万'
                elif "CNY'M" in blob or 'CNY million' in blob or 'CNY mn' in blob.lower():
                    currency_unit = "CNY'M"
                elif '人民币千元' in blob:
                    currency_unit = '人民币千元'
                elif "CNY'000" in blob or "CNY 000" in blob:
                    currency_unit = "CNY'000"
            except Exception:
                pass

            # An English-labelled source databook (e.g. Kunshan) will only
            # ever have "CNY'000"/"CNY'M" markers to detect, even when the
            # REPORT is being generated in Chinese -- normalise the unit
            # label itself to match the report language, same as the table
            # title above, so a Chinese table never shows an English header.
            if is_chinese_mode:
                if currency_unit == "CNY'000":
                    currency_unit = "人民币千元"
                elif currency_unit == "CNY'M":
                    currency_unit = "人民币百万"

            logger.info("Extracted BS: %s, IS: %s", bs_df.shape if bs_df is not None else 'None', is_df.shape if is_df is not None else 'None')
            logger.info("Table names - BS: %s, IS: %s, Currency: %s", bs_table_name, is_table_name, currency_unit)

            # If the values came from the precomputed (multiply_values=True)
            # pipeline, rescale to the source unit so the cells match the
            # header. "CNY'000" / "人民币千元" → divide by 1000,
            # The workbook extractor multiplies by 1000 ONLY when the source
            # header is CNY'000 / 人民币千元 (it does not touch millions).
            # Divide by 1000 here so the displayed cells read as thousands
            # (matching the "CNY'000" header). For any other unit the
            # values pass through unchanged.
            if values_pre_multiplied and currency_unit and (
                "千" in currency_unit or "'000" in currency_unit or "000" in currency_unit
            ):
                logger.info("Rescaling values by 1/1000 to match unit %s", currency_unit)
                for _df in (bs_df, is_df):
                    if _df is None or _df.empty:
                        continue
                    for _col in _df.columns:
                        if pd.api.types.is_numeric_dtype(_df[_col]):
                            _df[_col] = _df[_col] / 1000.0
            
            # Target the ACTUAL first commentary slide of each statement (a slide
            # object recorded during apply_structured_data_to_slides), not a
            # hard-coded slides[0]/slides[4]. Commentary adds slides and unused
            # ones are removed, so fixed indices drift — and the BS table could
            # land on a slide that no longer corresponds to BS page 1. The slide
            # OBJECT survives that reshuffle because it is a used (kept) slide.
            tracked = getattr(self, "_statement_table_slides", {}) or {}

            # Embed BS table on the first BS commentary slide.
            bs_slide = tracked.get("BS")
            if bs_slide is None and len(self.presentation.slides) > 0:
                bs_slide = self.presentation.slides[0]  # fallback
            if bs_df is not None and not bs_df.empty and bs_slide is not None:
                logger.info("Embedding BS table on tracked slide (shapes: %s)",
                            [getattr(s, 'name', '?') for s in bs_slide.shapes])
                self._embed_statement_table(
                    bs_slide, bs_df, "BS",
                    table_name=bs_table_name, currency_unit=currency_unit,
                    mappings=mappings, is_chinese_mode=is_chinese_mode,
                )
            else:
                logger.warning(
                    "Skipping BS table — bs_df empty=%s, target slide=%s. If bs_df is "
                    "empty but the databook DOES have a balance sheet, the session's "
                    "bs_is_results is stale: re-run Process Data, then export.",
                    bs_df is None or getattr(bs_df, 'empty', True), bs_slide is not None,
                )

            # Embed IS table on the first IS commentary slide.
            is_slide = tracked.get("IS")
            if is_slide is None and len(self.presentation.slides) > 4:
                is_slide = self.presentation.slides[4]  # fallback
            if is_df is not None and not is_df.empty and is_slide is not None:
                logger.info("Embedding IS table on tracked slide (shapes: %s)",
                            [getattr(s, 'name', '?') for s in is_slide.shapes])
                self._embed_statement_table(
                    is_slide, is_df, "IS",
                    table_name=is_table_name, currency_unit=currency_unit,
                    mappings=mappings, is_chinese_mode=is_chinese_mode,
                )
            elif is_df is not None and not is_df.empty:
                logger.error("No target slide found for IS table (slides=%s)", len(self.presentation.slides))
                    
        except Exception as e:
            logger.error("Error embedding financial tables: %s", e)
            logger.error(traceback.format_exc())

    # Chinese text shown in the Commentary label band in Chinese mode --
    # a single word, easy to change if the team prefers another term.
    _COMMENTARY_LABEL_ZH = "评述"

    def _localize_label_bands(self):
        """Translate the template's static English "Commentary" label band
        (Text-commentary / _L / _R, present on every slide) to Chinese
        when the deliverable is Chinese. English mode is left untouched.

        History -- why this TRANSLATES instead of deleting: ebf2179
        DELETED these shapes outright (the real reference deck has no
        such band at all), but the next two real Chinese exports both
        came back with the two statement-first pages (BS/IS page 1)
        rendering COMPLETELY BLANK in real PowerPoint, with abnormally
        slow file-open -- while python-pptx still saw every shape and
        all text intact in the same files, and the continuation pages
        kept rendering fine. Those two blank pages are exactly the two
        where the production template's band is a native PLACEHOLDER
        shape (the _L/_R variants on continuation pages are plain text
        boxes there); deleting a placeholder from a slide whose layout
        still defines it is the prime suspect for PowerPoint's renderer
        choking. Not reproducible on this Mac (no real PowerPoint, and
        the local template uses plain text boxes throughout), so the
        deletion was replaced by this strictly-safer in-place text swap:
        shape, placeholder wiring, fill and geometry all stay untouched
        -- only the literal run text changes, run-by-run so the run's
        own font formatting is preserved. If the blank-page symptom
        somehow survives even this, the one remaining suspect from the
        same commit on those two pages is the grid-table header restyle
        in _fill_table_placeholder.

        Matched by SHAPE NAME, at save time, so _calculate_table_bounds'
        text-based label-anchor fallback (matches the literal word
        "commentary" during layout) has long since finished by the time
        the text changes.
        """
        if not getattr(self, "_is_chinese_mode", False):
            return
        label_names = {"Text-commentary", "Text-commentary_L", "Text-commentary_R"}
        for slide in self.presentation.slides:
            for shape in slide.shapes:
                if self._shape_name(shape) not in label_names:
                    continue
                if not getattr(shape, "has_text_frame", False):
                    continue
                try:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if (run.text or "").strip().lower() == "commentary":
                                run.text = self._COMMENTARY_LABEL_ZH
                except Exception as exc:
                    logger.debug("Could not localize Commentary label band: %s", exc)

    def save(self, output_path: str):
        """Save the presentation"""
        if not self.presentation:
            raise ValueError("No presentation loaded")

        self._localize_label_bands()

        # Ensure output directory exists. dirname() is "" for a bare filename
        # in the current directory, and makedirs("") raises FileNotFoundError
        # -- so any caller passing "out.pptx" rather than a path crashed here
        # after doing all the generation work.
        out_dir = os.path.dirname(output_path)
        if out_dir:
            os.makedirs(out_dir, exist_ok=True)

        self.presentation.save(output_path)
        logger.info("Presentation saved to: %s", output_path)
# --- end pptx/generation.py ---
